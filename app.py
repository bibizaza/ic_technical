# Custom momentum loading (commented out - using standard approach for all instruments)
"""
Streamlit application for technical dashboard and presentation generation.

This application allows users to upload data, configure year‑to‑date (YTD)
charts for various asset classes, perform technical analysis on the S&P 500
index (including a selectable assessment title and a table of scores) and
generate a customised PowerPoint presentation.  The app persists
configuration selections in the session state and leverages helper functions
from the ``technical_analysis.equity.spx`` module for chart creation and
PowerPoint editing.

Key modifications relative to the original application:

* The SPX “view” title is no longer automatically derived from the average
  of technical and momentum scores.  Instead, users can select a view
  (e.g., “Strongly Bullish”) via a dropdown.  The chosen view is
  prepended with “S&P 500:” and inserted into the PowerPoint slide.
* The Streamlit interface no longer displays an average gauge for the SPX
  scores.  Instead, a simple table shows the technical score, momentum
  score and their average (DMAS), helping users judge the trend.
* The selected view is stored in ``st.session_state["spx_selected_view"]``
  and passed to ``insert_spx_technical_assessment`` when generating the
  presentation.
"""

# =====================================================================
# IMPORTANT: Clean up sys.path BEFORE any imports
# =====================================================================
import sys
from pathlib import Path

# Remove any herculis-technical-score or herculis-assessment paths
# that might have been added by other modules or previous runs
_app_dir = str(Path(__file__).parent)
_paths_to_remove = [
    p for p in sys.path
    if p and ('herculis-technical-score' in p or 'herculis-assessment' in p or 'market_compass' in p)
]
for p in _paths_to_remove:
    while p in sys.path:
        try:
            sys.path.remove(p)
        except ValueError:
            break

# Also clear any cached imports of utils from wrong paths
if 'utils' in sys.modules:
    _utils_module_path = getattr(sys.modules['utils'], '__file__', '')
    if 'herculis-technical-score' in _utils_module_path:
        # Wrong utils module was loaded, remove it
        del sys.modules['utils']

# Ensure app directory is first in sys.path
if _app_dir in sys.path:
    sys.path.remove(_app_dir)
sys.path.insert(0, _app_dir)

import streamlit as st

# ---------------------------------------------------------------------
# Load API keys from Streamlit secrets
# ---------------------------------------------------------------------
try:
    # Load Anthropic API key for Claude subtitle generation
    if "anthropic" in st.secrets and "api_key" in st.secrets["anthropic"]:
        from market_compass.subtitle_generator import set_api_key
        set_api_key(st.secrets["anthropic"]["api_key"])
except Exception as e:
    pass  # Secrets not configured, will use pattern-based subtitles

# ---------------------------------------------------------------------
# Load historical DMAS values from history tracker into session state
# ---------------------------------------------------------------------
def _load_historical_dmas_to_session():
    """
    Load previous week's DMAS values from history tracker into session state.
    This enables week-over-week comparison in technical slides.

    Uses data_as_of date from calendar picker if available.
    Also stores days_gap and previous_date for proper change text formatting.
    """
    try:
        from market_compass.subtitle_generator.history_tracker import get_tracker
        import os
        tracker = get_tracker()

        # Debug: Show history file path and status
        print(f"[History] Looking for: {tracker.storage_path}")
        print(f"[History] File exists: {os.path.exists(tracker.storage_path)}")

        all_assets = tracker.get_all_assets()
        print(f"[History] Found {len(all_assets)} assets in history: {all_assets[:10]}{'...' if len(all_assets) > 10 else ''}")

        # Get the calendar-selected date if available
        current_date = st.session_state.get("data_as_of")
        print(f"[History] Using data_as_of date: {current_date}")

        # Map asset names (as stored in history) to ticker keys (as used in session state)
        ASSET_TO_TICKER_KEY = {
            # Equities
            "S&P 500": "spx",
            "SPX": "spx",
            "CSI 300": "csi",
            "CSI300": "csi",
            "Nikkei 225": "nikkei",
            "Nikkei": "nikkei",
            "TASI": "tasi",
            "Sensex": "sensex",
            "Dax": "dax",
            "DAX": "dax",
            "SMI": "smi",
            "Ibov": "ibov",
            "IBOV": "ibov",
            "Bovespa": "ibov",
            "Mexbol": "mexbol",
            "MEXBOL": "mexbol",
            # Precious Metals
            "Gold": "gold",
            "GOLD": "gold",
            "Silver": "silver",
            "SILVER": "silver",
            "Platinum": "platinum",
            "PLATINUM": "platinum",
            "Palladium": "palladium",
            "PALLADIUM": "palladium",
            # Energy & Industrial
            "Oil": "oil",
            "OIL": "oil",
            "WTI": "oil",
            "Copper": "copper",
            "COPPER": "copper",
            # Crypto
            "Bitcoin": "bitcoin",
            "BTC": "bitcoin",
            "Ethereum": "ethereum",
            "ETH": "ethereum",
            "Ripple": "ripple",
            "XRP": "ripple",
            "Solana": "solana",
            "SOL": "solana",
            "Binance": "binance",
            "BNB": "binance",
        }

        loaded_count = 0
        for asset_name, ticker_key in ASSET_TO_TICKER_KEY.items():
            # Get the previous week's data, passing current_date for proper lookup
            last_week = tracker.get_last_week(asset_name, current_date=current_date)

            if last_week is not None:
                # Load DMAS
                prev_dmas = last_week.get("dmas")
                if prev_dmas is not None:
                    session_key = f"{ticker_key}_last_week_avg"
                    # Only set if not already set (don't override transition sheet data)
                    if session_key not in st.session_state:
                        st.session_state[session_key] = float(prev_dmas)
                        loaded_count += 1
                        print(f"[History] Set {session_key} = {prev_dmas} (from {asset_name})")

                # Store days_gap for proper change text formatting
                days_gap = last_week.get("days_gap")
                if days_gap is not None:
                    gap_key = f"{ticker_key}_prev_days_gap"
                    st.session_state[gap_key] = days_gap
                    print(f"[History] Set {gap_key} = {days_gap}")

                # Store previous_date for display in change text
                prev_date = last_week.get("previous_date")
                if prev_date is not None:
                    date_key = f"{ticker_key}_prev_date"
                    st.session_state[date_key] = prev_date
                    print(f"[History] Set {date_key} = {prev_date}")

                # Load Technical score for trend arrow
                prev_tech = last_week.get("technical_score")
                if prev_tech is not None:
                    tech_key = f"{ticker_key}_last_week_tech"
                    if tech_key not in st.session_state:
                        st.session_state[tech_key] = float(prev_tech)
                        print(f"[History] Set {tech_key} = {prev_tech} (from {asset_name})")

                # Load Momentum score for trend arrow
                prev_mom = last_week.get("momentum_score")
                if prev_mom is not None:
                    mom_key = f"{ticker_key}_last_week_mom"
                    if mom_key not in st.session_state:
                        st.session_state[mom_key] = float(prev_mom)
                        print(f"[History] Set {mom_key} = {prev_mom} (from {asset_name})")

                # Load RSI for trend arrow
                prev_rsi = last_week.get("rsi")
                if prev_rsi is not None:
                    rsi_key = f"{ticker_key}_last_week_rsi"
                    if rsi_key not in st.session_state:
                        st.session_state[rsi_key] = float(prev_rsi)
                        print(f"[History] Set {rsi_key} = {prev_rsi} (from {asset_name})")

        if loaded_count > 0:
            print(f"[History] Loaded {loaded_count} previous DMAS values from history")
        else:
            print(f"[History] No previous DMAS values loaded (need at least 2 weeks of data per asset)")

    except ImportError:
        print("[History] History tracker not available (import error)")
    except Exception as e:
        import traceback
        print(f"[History] Warning: Could not load historical DMAS: {e}")
        traceback.print_exc()

# Load historical data on app startup
# NOTE: Moved to after calendar date picker to ensure correct date is used
# _load_historical_dmas_to_session()

# ---------------------------------------------------------------------
# Assessment and subtitle generation - LAZY LOADING
# ---------------------------------------------------------------------
# Don't import at module level to avoid sys.path pollution
# Import only when needed inside functions
SUBTITLE_GEN_AVAILABLE = True  # Assume available, will check when needed

# Define assessment options directly here to avoid import
ASSESSMENT_OPTIONS = [
    "Bearish",
    "Cautious",
    "Neutral",
    "Constructive",
    "Bullish",
]

# ---------------------------------------------------------------------
# Safe score helpers to avoid float(None) crashes in commodities/crypto
# ---------------------------------------------------------------------
from typing import Optional, Union

_Number = Union[int, float]

def _safe_avg_two(a: Optional[_Number], b: Optional[_Number], rounding: int = 1) -> Optional[float]:
    vals = []
    for v in (a, b):
        try:
            if v is not None:
                vals.append(float(v))
        except Exception:
            pass
    if not vals:
        return None
    return round(sum(vals) / len(vals), rounding)

def _fmt_score_cell(v: Optional[_Number]) -> str:
    try:
        return f"{int(round(float(v)))}"
    except Exception:
        return "N/A"

def _fmt_dmas_cell(v: Optional[_Number]) -> str:
    try:
        return f"{round(float(v), 1):.1f}"
    except Exception:
        return "N/A"
import pandas as pd
import numpy as np
import plotly.graph_objects as go
from sklearn.linear_model import LinearRegression
from io import BytesIO
from pptx import Presentation
import tempfile
from pathlib import Path

# Import SPX functions from the dedicated module.  The SPX module
# resides in ``technical_analysis/equity/spx.py`` and provides all helper
# functions for building charts, inserting data into slides, and
# computing scores.  Note that ``insert_spx_technical_assessment``
# accepts a manual description and ``insert_spx_source`` inserts the
# source footnote based on the selected price mode.

# --- FIX: Force reload of the spx module to ensure changes are always applied ---
from importlib import reload
try:
    import technical_analysis.equity.spx as _spx_module
    _spx_module = reload(_spx_module)
except ImportError:
    pass  # Let the next import handle the error if the module doesn't exist
# --- END FIX ---

from technical_analysis.equity.spx import (
    make_spx_figure,
    _get_spx_technical_score,
    _get_spx_momentum_score,
    generate_range_gauge_only_image,
    _compute_range_bounds as _compute_range_bounds_spx,
    # Technical Analysis v2 (Chart.js + Playwright)
    create_technical_analysis_v2_chart,
    insert_technical_analysis_v2_slide,
)

from importlib import reload
try:
    import funda_breadth.breadth_page as _breadth
except ModuleNotFoundError:
    import funda_breadth.breadth_page as _breadth
_breadth = reload(_breadth)

_load_breadth_page_data = _breadth._load_and_prepare
_style_breadth_page     = _breadth._apply_matrix_style
_debug_breadth_rows     = _breadth.debug_first_rows

# ─── Breadth‑picture helper (Excel → EMF → PPT) ───────────────────────
from pathlib import Path
import tempfile
from importlib import reload


# Import SMI functions from the dedicated module.  The SMI module resides
# in ``technical_analysis/equity/smi.py`` and provides helper functions
# for V2 chart generation (score/momentum retrieval, range computation).
try:
    from technical_analysis.equity.smi import (
        make_smi_figure,
        _get_smi_technical_score,
        _get_smi_momentum_score,
        _compute_range_bounds as _compute_range_bounds_smi,
    )
except Exception:
    # Define no-op stand-ins if the SMI module is unavailable
    def make_smi_figure(*args, **kwargs):
        return go.Figure()
    def _get_smi_technical_score(*args, **kwargs):
        return None
    def _get_smi_momentum_score(*args, **kwargs):
        return None

    # Fallback: if the SMI module is unavailable, fall back to the SPX range computation
    def _compute_range_bounds_smi(*args, **kwargs):  # type: ignore
        return _compute_range_bounds_spx(*args, **kwargs)

# Import IBOV functions from the dedicated module.  The IBOV module
# resides in ``technical_analysis/equity/ibov.py`` and provides helper
# functions for V2 chart generation (score/momentum retrieval, range computation).
try:
    from technical_analysis.equity.ibov import (
        make_ibov_figure,
        _get_ibov_technical_score,
        _get_ibov_momentum_score,
        _compute_range_bounds as _compute_range_bounds_ibov,
    )
except Exception:
    # Define no-op stand-ins if the IBOV module is unavailable
    def make_ibov_figure(*args, **kwargs):
        return go.Figure()
    def _get_ibov_technical_score(*args, **kwargs):
        return None
    def _get_ibov_momentum_score(*args, **kwargs):
        return None

    # Fallback: if the IBOV module is unavailable, fall back to the SPX range computation
    def _compute_range_bounds_ibov(*args, **kwargs):  # type: ignore
        return _compute_range_bounds_spx(*args, **kwargs)

# Import Mexbol functions from the dedicated module.  The Mexbol module
# resides in ``technical_analysis/equity/mexbol.py`` and provides helper
# functions for V2 chart generation (score/momentum retrieval, range computation).
try:
    from technical_analysis.equity.mexbol import (
        make_mexbol_figure,
        _get_mexbol_technical_score,
        _get_mexbol_momentum_score,
        _compute_range_bounds as _compute_range_bounds_mexbol,
    )
except Exception:
    # Define no-op stand-ins if the Mexbol module is unavailable
    def make_mexbol_figure(*args, **kwargs):
        return go.Figure()
    def _get_mexbol_technical_score(*args, **kwargs):
        return None
    def _get_mexbol_momentum_score(*args, **kwargs):
        return None

    # Fallback: if the Mexbol module is unavailable, fall back to the SPX range computation
    def _compute_range_bounds_mexbol(*args, **kwargs):  # type: ignore
        return _compute_range_bounds_spx(*args, **kwargs)

# Import Gold functions from the dedicated module.  The Gold module
# resides in ``technical_analysis/commodity/gold.py`` and provides helper
# functions for V2 chart generation (score/momentum retrieval, range computation).
try:
    from technical_analysis.commodity.gold import (
        make_gold_figure,
        _get_gold_technical_score,
        _get_gold_momentum_score,
        _compute_range_bounds as _compute_range_bounds_gold,
    )
except Exception:
    try:
        from gold import (
            make_gold_figure,
            _get_gold_technical_score,
            _get_gold_momentum_score,
            _compute_range_bounds as _compute_range_bounds_gold,
        )
    except Exception:
        # Define no-op stand-ins if the Gold module is unavailable
        def make_gold_figure(*args, **kwargs):
            return go.Figure()
        def _get_gold_technical_score(*args, **kwargs):
            return None
        def _get_gold_momentum_score(*args, **kwargs):
            return None
        def _compute_range_bounds_gold(*args, **kwargs):
            return _compute_range_bounds_spx(*args, **kwargs)

# Import Silver functions from the dedicated module.  The Silver module
# resides in ``technical_analysis/commodity/silver.py`` and provides helper
# functions for V2 chart generation (score/momentum retrieval, range computation).
try:
    from technical_analysis.commodity.silver import (
        make_silver_figure,
        _get_silver_technical_score,
        _get_silver_momentum_score,
        _compute_range_bounds as _compute_range_bounds_silver,
    )
except Exception:
    try:
        from silver import (
            make_silver_figure,
            _get_silver_technical_score,
            _get_silver_momentum_score,
            _compute_range_bounds as _compute_range_bounds_silver,
        )
    except Exception:
        # Define no-op stand-ins if the Silver module is unavailable
        def make_silver_figure(*args, **kwargs):
            return go.Figure()
        def _get_silver_technical_score(*args, **kwargs):
            return None
        def _get_silver_momentum_score(*args, **kwargs):
            return None
        def _compute_range_bounds_silver(*args, **kwargs):
            return _compute_range_bounds_spx(*args, **kwargs)

# Import Platinum functions from the dedicated module.  The Platinum module
# resides in ``technical_analysis/commodity/platinum.py`` and provides helper
# functions for V2 chart generation (score/momentum retrieval, range computation).
try:
    from technical_analysis.commodity.platinum import (
        make_platinum_figure,
        _get_platinum_technical_score,
        _get_platinum_momentum_score,
        _compute_range_bounds as _compute_range_bounds_platinum,
    )
except Exception:
    try:
        from platinum import (
            make_platinum_figure,
            _get_platinum_technical_score,
            _get_platinum_momentum_score,
            _compute_range_bounds as _compute_range_bounds_platinum,
        )
    except Exception:
        # Define no-op stand-ins if the Platinum module is unavailable
        def make_platinum_figure(*args, **kwargs):
            return go.Figure()
        def _get_platinum_technical_score(*args, **kwargs):
            return None
        def _get_platinum_momentum_score(*args, **kwargs):
            return None
        def _compute_range_bounds_platinum(*args, **kwargs):
            return _compute_range_bounds_spx(*args, **kwargs)

# Import Palladium functions from the dedicated module.  The Palladium module
# resides in ``technical_analysis/commodity/palladium.py`` and provides helper
# functions for V2 chart generation (score/momentum retrieval, range computation).
try:
    from technical_analysis.commodity.palladium import (
        make_palladium_figure,
        _get_palladium_technical_score,
        _get_palladium_momentum_score,
        _compute_range_bounds as _compute_range_bounds_palladium,
    )
except Exception:
    try:
        from palladium import (
            make_palladium_figure,
            _get_palladium_technical_score,
            _get_palladium_momentum_score,
            _compute_range_bounds as _compute_range_bounds_palladium,
        )
    except Exception:
        # Define no-op stand-ins if the Palladium module is unavailable
        def make_palladium_figure(*args, **kwargs):
            return go.Figure()
        def _get_palladium_technical_score(*args, **kwargs):
            return None
        def _get_palladium_momentum_score(*args, **kwargs):
            return None
        def _compute_range_bounds_palladium(*args, **kwargs):
            return _compute_range_bounds_spx(*args, **kwargs)

# Import Oil functions from the dedicated module.  The Oil module
# resides in ``technical_analysis/commodity/oil.py`` and provides helper
# functions for V2 chart generation (score/momentum retrieval, range computation).
try:
    from technical_analysis.commodity.oil import (
        make_oil_figure,
        _get_oil_technical_score,
        _get_oil_momentum_score,
        _compute_range_bounds as _compute_range_bounds_oil,
    )
except Exception:
    try:
        from oil import (
            make_oil_figure,
            _get_oil_technical_score,
            _get_oil_momentum_score,
            _compute_range_bounds as _compute_range_bounds_oil,
        )
    except Exception:
        # Define no-op stand-ins if the Oil module is unavailable
        def make_oil_figure(*args, **kwargs):
            return go.Figure()
        def _get_oil_technical_score(*args, **kwargs):
            return None
        def _get_oil_momentum_score(*args, **kwargs):
            return None
        def _compute_range_bounds_oil(*args, **kwargs):
            return _compute_range_bounds_spx(*args, **kwargs)

# Import Copper functions from the dedicated module.  The Copper module
# resides in ``technical_analysis/commodity/copper.py`` and provides helper
# functions for V2 chart generation (score/momentum retrieval, range computation).
try:
    from technical_analysis.commodity.copper import (
        make_copper_figure,
        _get_copper_technical_score,
        _get_copper_momentum_score,
        _compute_range_bounds as _compute_range_bounds_copper,
    )
except Exception:
    try:
        from copper import (
            make_copper_figure,
            _get_copper_technical_score,
            _get_copper_momentum_score,
            _compute_range_bounds as _compute_range_bounds_copper,
        )
    except Exception:
        # Define no-op stand-ins if the Copper module is unavailable
        def make_copper_figure(*args, **kwargs):
            return go.Figure()
        def _get_copper_technical_score(*args, **kwargs):
            return None
        def _get_copper_momentum_score(*args, **kwargs):
            return None
        def _compute_range_bounds_copper(*args, **kwargs):
            return _compute_range_bounds_spx(*args, **kwargs)

# Import Bitcoin functions from the dedicated module.  The Bitcoin module
# resides in ``technical_analysis/crypto/bitcoin.py`` and provides helper
# functions for V2 chart generation (score/momentum retrieval, range computation).
try:
    from technical_analysis.crypto.bitcoin import (
        make_bitcoin_figure,
        _get_bitcoin_technical_score,
        _get_bitcoin_momentum_score,
        _compute_range_bounds as _compute_range_bounds_bitcoin,
    )
except Exception:
    try:
        from bitcoin import (
            make_bitcoin_figure,
            _get_bitcoin_technical_score,
            _get_bitcoin_momentum_score,
            _compute_range_bounds as _compute_range_bounds_bitcoin,
        )
    except Exception:
        # Define no-op stand-ins if the Bitcoin module is unavailable
        def make_bitcoin_figure(*args, **kwargs):
            return go.Figure()
        def _get_bitcoin_technical_score(*args, **kwargs):
            return None
        def _get_bitcoin_momentum_score(*args, **kwargs):
            return None
        def _compute_range_bounds_bitcoin(*args, **kwargs):
            return _compute_range_bounds_spx(*args, **kwargs)

# Import Ethereum functions from the dedicated module.  The Ethereum module
# resides in ``technical_analysis/crypto/ethereum.py`` and provides helper
# functions for V2 chart generation (score/momentum retrieval, range computation).
try:
    from technical_analysis.crypto.ethereum import (
        make_ethereum_figure,
        _get_ethereum_technical_score,
        _get_ethereum_momentum_score,
        _compute_range_bounds as _compute_range_bounds_ethereum,
    )
except Exception:
    try:
        from ethereum import (
            make_ethereum_figure,
            _get_ethereum_technical_score,
            _get_ethereum_momentum_score,
            _compute_range_bounds as _compute_range_bounds_ethereum,
        )
    except Exception:
        # Define no-op stand-ins if the Ethereum module is unavailable
        def make_ethereum_figure(*args, **kwargs):
            return go.Figure()
        def _get_ethereum_technical_score(*args, **kwargs):
            return None
        def _get_ethereum_momentum_score(*args, **kwargs):
            return None
        def _compute_range_bounds_ethereum(*args, **kwargs):
            return _compute_range_bounds_spx(*args, **kwargs)

# Import Ripple functions from the dedicated module.  The Ripple module
# resides in ``technical_analysis/crypto/ripple.py`` and provides helper
# functions for V2 chart generation (score/momentum retrieval, range computation).
try:
    from technical_analysis.crypto.ripple import (
        make_ripple_figure,
        _get_ripple_technical_score,
        _get_ripple_momentum_score,
        _compute_range_bounds as _compute_range_bounds_ripple,
    )
except Exception:
    try:
        from ripple import (
            make_ripple_figure,
            _get_ripple_technical_score,
            _get_ripple_momentum_score,
            _compute_range_bounds as _compute_range_bounds_ripple,
        )
    except Exception:
        # Define no-op stand-ins if the Ripple module is unavailable
        def make_ripple_figure(*args, **kwargs):
            return go.Figure()
        def _get_ripple_technical_score(*args, **kwargs):
            return None
        def _get_ripple_momentum_score(*args, **kwargs):
            return None
        def _compute_range_bounds_ripple(*args, **kwargs):
            return _compute_range_bounds_spx(*args, **kwargs)

# Import Solana functions from the dedicated module.  The Solana module
# resides in ``technical_analysis/crypto/solana.py`` and provides helper
# functions for V2 chart generation (score/momentum retrieval, range computation).
try:
    from technical_analysis.crypto.solana import (
        make_solana_figure,
        _get_solana_technical_score,
        _get_solana_momentum_score,
        _compute_range_bounds as _compute_range_bounds_solana,
    )
except Exception:
    try:
        from solana import (
            make_solana_figure,
            _get_solana_technical_score,
            _get_solana_momentum_score,
            _compute_range_bounds as _compute_range_bounds_solana,
        )
    except Exception:
        # Define no-op stand-ins if the Solana module is unavailable
        def make_solana_figure(*args, **kwargs):
            return go.Figure()
        def _get_solana_technical_score(*args, **kwargs):
            return None
        def _get_solana_momentum_score(*args, **kwargs):
            return None
        def _compute_range_bounds_solana(*args, **kwargs):
            return _compute_range_bounds_spx(*args, **kwargs)

# Import Binance functions from the dedicated module.  The Binance module
# resides in ``technical_analysis/crypto/binance.py`` and provides helper
# functions for V2 chart generation (score/momentum retrieval, range computation).
try:
    from technical_analysis.crypto.binance import (
        make_binance_figure,
        _get_binance_technical_score,
        _get_binance_momentum_score,
        _compute_range_bounds as _compute_range_bounds_binance,
    )
except Exception:
    try:
        from binance import (
            make_binance_figure,
            _get_binance_technical_score,
            _get_binance_momentum_score,
            _compute_range_bounds as _compute_range_bounds_binance,
        )
    except Exception:
        # Define no-op stand-ins if the Binance module is unavailable
        def make_binance_figure(*args, **kwargs):
            return go.Figure()
        def _get_binance_technical_score(*args, **kwargs):
            return None
        def _get_binance_momentum_score(*args, **kwargs):
            return None
        def _compute_range_bounds_binance(*args, **kwargs):
            return _compute_range_bounds_spx(*args, **kwargs)

# Import CSI functions from the dedicated module.  The CSI module resides
# in ``technical_analysis/equity/csi.py`` and provides helper functions
# for V2 chart generation (score/momentum retrieval, range computation).
try:
    from technical_analysis.equity.csi import (
        make_csi_figure,
        _get_csi_technical_score,
        _get_csi_momentum_score,
        _compute_range_bounds as _compute_range_bounds_csi,
    )
except Exception:
    # Define no-op stand-ins if the CSI module is unavailable
    def make_csi_figure(*args, **kwargs):
        return go.Figure()
    def _get_csi_technical_score(*args, **kwargs):
        return None
    def _get_csi_momentum_score(*args, **kwargs):
        return None

    # Fallback: if the CSI module is unavailable, fall back to the SPX range computation
    def _compute_range_bounds_csi(*args, **kwargs):  # type: ignore
        return _compute_range_bounds_spx(*args, **kwargs)

# Import Nikkei functions from the dedicated module.  The Nikkei module
# resides in ``technical_analysis/equity/nikkei.py`` and provides helper
# functions for V2 chart generation (score/momentum retrieval, range computation).
try:
    from technical_analysis.equity.nikkei import (
        make_nikkei_figure,
        _get_nikkei_technical_score,
        _get_nikkei_momentum_score,
        _compute_range_bounds as _compute_range_bounds_nikkei,
    )
except Exception:
    # Define no-op stand-ins if the Nikkei module is unavailable
    def make_nikkei_figure(*args, **kwargs):
        return go.Figure()
    def _get_nikkei_technical_score(*args, **kwargs):
        return None
    def _get_nikkei_momentum_score(*args, **kwargs):
        return None

    # Fallback: if the Nikkei module is unavailable, fall back to the SPX range computation
    def _compute_range_bounds_nikkei(*args, **kwargs):  # type: ignore
        return _compute_range_bounds_spx(*args, **kwargs)

# Import TASI functions from the dedicated module.  The TASI module
# resides in ``technical_analysis/equity/tasi.py`` and provides helper
# functions for V2 chart generation (score/momentum retrieval, range computation).
try:
    from technical_analysis.equity.tasi import (
        make_tasi_figure,
        _get_tasi_technical_score,
        _get_tasi_momentum_score,
        _compute_range_bounds as _compute_range_bounds_tasi,
    )
except Exception:
    # Define no-op stand-ins if the TASI module is unavailable
    def make_tasi_figure(*args, **kwargs):
        return go.Figure()
    def _get_tasi_technical_score(*args, **kwargs):
        return None
    def _get_tasi_momentum_score(*args, **kwargs):
        return None

    # Fallback: if the TASI module is unavailable, fall back to the SPX range computation
    def _compute_range_bounds_tasi(*args, **kwargs):  # type: ignore
        return _compute_range_bounds_spx(*args, **kwargs)

# Import Sensex functions from the dedicated module.  The Sensex module
# resides in ``technical_analysis/equity/sensex.py`` and provides helper
# functions for V2 chart generation (score/momentum retrieval, range computation).
try:
    from technical_analysis.equity.sensex import (
        make_sensex_figure,
        _get_sensex_technical_score,
        _get_sensex_momentum_score,
        _compute_range_bounds as _compute_range_bounds_sensex,
    )
except Exception:
    # Define no-op stand-ins if the Sensex module is unavailable
    def make_sensex_figure(*args, **kwargs):
        return go.Figure()
    def _get_sensex_technical_score(*args, **kwargs):
        return None
    def _get_sensex_momentum_score(*args, **kwargs):
        return None

    # Fallback: if the Sensex module is unavailable, fall back to the SPX range computation
    def _compute_range_bounds_sensex(*args, **kwargs):  # type: ignore
        return _compute_range_bounds_spx(*args, **kwargs)

# Import DAX functions from the dedicated module.  The DAX module
# resides in ``technical_analysis/equity/dax.py`` and provides helper
# functions for V2 chart generation (score/momentum retrieval, range computation).
try:
    from technical_analysis.equity.dax import (
        make_dax_figure,
        _get_dax_technical_score,
        _get_dax_momentum_score,
        _compute_range_bounds as _compute_range_bounds_dax,
    )
except Exception:
    # Define no-op stand-ins if the DAX module is unavailable
    def make_dax_figure(*args, **kwargs):
        return go.Figure()
    def _get_dax_technical_score(*args, **kwargs):
        return None
    def _get_dax_momentum_score(*args, **kwargs):
        return None

    # Fallback: if the DAX module is unavailable, fall back to the SPX range computation
    def _compute_range_bounds_dax(*args, **kwargs):  # type: ignore
        return _compute_range_bounds_spx(*args, **kwargs)

# Import CSI functions from the dedicated module.  The CSI module resides
# in ``technical_analysis/equity/csi.py`` and provides helper functions
# for V2 chart generation (score/momentum retrieval, range computation).
try:
    from technical_analysis.equity.csi import (
        make_csi_figure,
        _get_csi_technical_score,
        _get_csi_momentum_score,
        _compute_range_bounds as _compute_range_bounds_csi,
    )
except Exception:
    # Define no-op stand-ins if the CSI module is unavailable
    def make_csi_figure(*args, **kwargs):
        return go.Figure()
    def _get_csi_technical_score(*args, **kwargs):
        return None
    def _get_csi_momentum_score(*args, **kwargs):
        return None

    # Fallback: if the CSI module is unavailable, fall back to the SPX range computation
    def _compute_range_bounds_csi(*args, **kwargs):  # type: ignore
        return _compute_range_bounds_spx(*args, **kwargs)

# Import Nikkei functions from the dedicated module.  The Nikkei module
# resides in ``technical_analysis/equity/nikkei.py`` and provides helper
# functions for V2 chart generation (score/momentum retrieval, range computation).
try:
    from technical_analysis.equity.nikkei import (
        make_nikkei_figure,
        _get_nikkei_technical_score,
        _get_nikkei_momentum_score,
        _compute_range_bounds as _compute_range_bounds_nikkei,
    )
except Exception:
    # Define no-op stand-ins if the Nikkei module is unavailable
    def make_nikkei_figure(*args, **kwargs):
        return go.Figure()
    def _get_nikkei_technical_score(*args, **kwargs):
        return None
    def _get_nikkei_momentum_score(*args, **kwargs):
        return None

    # Fallback: if the Nikkei module is unavailable, fall back to the SPX range computation
    def _compute_range_bounds_nikkei(*args, **kwargs):  # type: ignore
        return _compute_range_bounds_spx(*args, **kwargs)

# Import TASI functions from the dedicated module.  The TASI module
# resides in ``technical_analysis/equity/tasi.py`` and provides helper
# functions for V2 chart generation (score/momentum retrieval, range computation).
try:
    from technical_analysis.equity.tasi import (
        make_tasi_figure,
        _get_tasi_technical_score,
        _get_tasi_momentum_score,
        _compute_range_bounds as _compute_range_bounds_tasi,
    )
except Exception:
    # Define no-op stand-ins if the TASI module is unavailable
    def make_tasi_figure(*args, **kwargs):
        return go.Figure()
    def _get_tasi_technical_score(*args, **kwargs):
        return None
    def _get_tasi_momentum_score(*args, **kwargs):
        return None

    # Fallback: if the TASI module is unavailable, fall back to the SPX range computation
    def _compute_range_bounds_tasi(*args, **kwargs):  # type: ignore
        return _compute_range_bounds_spx(*args, **kwargs)

# Import Sensex functions from the dedicated module.  The Sensex module
# resides in ``technical_analysis/equity/sensex.py`` and provides helper
# functions for V2 chart generation (score/momentum retrieval, range computation).
try:
    from technical_analysis.equity.sensex import (
        make_sensex_figure,
        _get_sensex_technical_score,
        _get_sensex_momentum_score,
        _compute_range_bounds as _compute_range_bounds_sensex,
    )
except Exception:
    # Define no-op stand-ins if the Sensex module is unavailable
    def make_sensex_figure(*args, **kwargs):
        return go.Figure()
    def _get_sensex_technical_score(*args, **kwargs):
        return None
    def _get_sensex_momentum_score(*args, **kwargs):
        return None

    # Fallback: if the Sensex module is unavailable, fall back to the SPX range computation
    def _compute_range_bounds_sensex(*args, **kwargs):  # type: ignore
        return _compute_range_bounds_spx(*args, **kwargs)

# Import helper to adjust price data according to price mode.  The utils
# module resides at the project root (e.g. ``ic/utils.py``) so that it can
# be shared across technical analysis and performance modules.
from utils import adjust_prices_for_mode

# Import performance dashboard helpers (unchanged)
from performance.equity_perf import (
    create_weekly_performance_chart,
    create_historical_performance_table,
    insert_equity_performance_bar_slide,
    insert_equity_performance_histo_slide,
    create_equity_ytd_evolution_chart,
    insert_equity_ytd_evolution_slide,
    create_fx_impact_analysis_chart_eur,
    insert_fx_impact_analysis_slide_eur,
)

# Import FX performance functions
try:
    from performance.fx_perf import (
        create_weekly_performance_chart as create_weekly_fx_performance_chart,
        create_historical_performance_table as create_historical_fx_performance_table,
        insert_fx_performance_bar_slide,
        insert_fx_performance_histo_slide,
        create_weekly_html_performance_chart as create_weekly_fx_html_chart,
        insert_fx_weekly_html_slide,
        create_historical_html_performance_chart as create_historical_fx_html_chart,
        insert_fx_historical_html_slide,
    )
except Exception:
    # If FX module not available, define no-op placeholders
    def create_weekly_fx_performance_chart(*args, **kwargs):
        return (b"", None)
    def create_historical_fx_performance_table(*args, **kwargs):
        return (b"", None)
    def insert_fx_performance_bar_slide(prs, image_bytes, *args, **kwargs):
        return prs
    def insert_fx_performance_histo_slide(prs, image_bytes, *args, **kwargs):
        return prs
    def create_weekly_fx_html_chart(*args, **kwargs):
        return (b"", None)
    def insert_fx_weekly_html_slide(prs, image_bytes, *args, **kwargs):
        return prs
    def create_historical_fx_html_chart(*args, **kwargs):
        return (b"", None)
    def insert_fx_historical_html_slide(prs, image_bytes, *args, **kwargs):
        return prs

# Import Crypto performance functions
try:
    from performance.crypto_perf import (
        create_weekly_performance_chart as create_weekly_crypto_performance_chart,
        create_historical_performance_table as create_historical_crypto_performance_table,
        insert_crypto_performance_bar_slide,
        insert_crypto_performance_histo_slide,
        create_weekly_html_performance_chart as create_weekly_crypto_html_chart,
        insert_crypto_weekly_html_slide,
        create_historical_html_performance_chart as create_historical_crypto_html_chart,
        insert_crypto_historical_html_slide,
        create_crypto_ytd_evolution_chart,
        insert_crypto_ytd_evolution_slide,
    )
except Exception:
    # If Crypto module not available, define no-op placeholders
    def create_weekly_crypto_performance_chart(*args, **kwargs):
        return (b"", None)
    def create_historical_crypto_performance_table(*args, **kwargs):
        return (b"", None)
    def insert_crypto_performance_bar_slide(prs, image_bytes, *args, **kwargs):
        return prs
    def insert_crypto_performance_histo_slide(prs, image_bytes, *args, **kwargs):
        return prs
    def create_weekly_crypto_html_chart(*args, **kwargs):
        return (b"", None)
    def insert_crypto_weekly_html_slide(prs, image_bytes, *args, **kwargs):
        return prs
    def create_historical_crypto_html_chart(*args, **kwargs):
        return (b"", None)
    def insert_crypto_historical_html_slide(prs, image_bytes, *args, **kwargs):
        return prs

# Import Credit performance functions
try:
    from performance.corp_bonds_perf import (
        create_weekly_performance_chart as create_weekly_credit_performance_chart,
        insert_corp_bonds_performance_slide as insert_credit_performance_bar_slide,
        create_historical_performance_chart as create_historical_credit_performance_chart,
        insert_corp_bonds_historical_slide as insert_credit_performance_histo_slide,
    )
except Exception:
    # If Credit module not available, define no-op placeholders
    def create_weekly_credit_performance_chart(*args, **kwargs):  # type: ignore
        return (b"", None)
    def create_historical_credit_performance_chart(*args, **kwargs):  # type: ignore
        return (b"", None)
    def insert_credit_performance_bar_slide(prs, image_bytes, *args, **kwargs):  # type: ignore
        return prs
    def insert_credit_performance_histo_slide(prs, image_bytes, *args, **kwargs):  # type: ignore
        return prs

# Import Commodity performance functions
try:
    from performance.commodity_perf import (
        create_weekly_performance_chart as create_weekly_commodity_performance_chart,
        create_historical_performance_table as create_historical_commodity_performance_table,
        insert_commodity_performance_bar_slide,
        insert_commodity_performance_histo_slide,
        create_weekly_html_performance_chart as create_weekly_commodity_html_chart,
        insert_commodity_weekly_html_slide,
        create_historical_html_performance_chart as create_historical_commodity_html_chart,
        insert_commodity_historical_html_slide,
        create_commodity_ytd_evolution_chart,
        insert_commodity_ytd_evolution_slide,
    )
except Exception:
    def create_weekly_commodity_performance_chart(*args, **kwargs):  # type: ignore
        return (b"", None)
    def create_historical_commodity_performance_table(*args, **kwargs):  # type: ignore
        return (b"", None)
    def insert_commodity_performance_bar_slide(prs, image_bytes, *args, **kwargs):  # type: ignore
        return prs
    def insert_commodity_performance_histo_slide(prs, image_bytes, *args, **kwargs):  # type: ignore
        return prs
    def create_weekly_commodity_html_chart(*args, **kwargs):  # type: ignore
        return (b"", None)
    def insert_commodity_weekly_html_slide(prs, image_bytes, *args, **kwargs):  # type: ignore
        return prs
    def create_historical_commodity_html_chart(*args, **kwargs):  # type: ignore
        return (b"", None)
    def insert_commodity_historical_html_slide(prs, image_bytes, *args, **kwargs):  # type: ignore
        return prs
    
# Import Rates performance functions
try:
    from performance.rates_perf import (
        create_weekly_performance_chart as create_weekly_rates_performance_chart,
        create_historical_performance_table as create_historical_rates_performance_table,
        insert_rates_performance_bar_slide,
        insert_rates_performance_histo_slide,
    )
except Exception:
    # If Rates module is not available, define no-op placeholders
    def create_weekly_rates_performance_chart(*args, **kwargs):  # type: ignore
        return (b"", None)
    def create_historical_rates_performance_table(*args, **kwargs):  # type: ignore
        return (b"", None)
    def insert_rates_performance_bar_slide(prs, image_bytes, *args, **kwargs):  # type: ignore
        return prs
    def insert_rates_performance_histo_slide(prs, image_bytes, *args, **kwargs):  # type: ignore
        return prs

###############################################################################
# Synthetic data helpers (fallback when no Excel is loaded)
###############################################################################

def _create_synthetic_spx_series() -> pd.DataFrame:
    """Create a synthetic SPX price series for demonstration purposes."""
    end_date = pd.Timestamp.today().normalize()
    start_date = end_date - pd.Timedelta(days=730)
    dates = pd.date_range(start=start_date, end=end_date, freq="B")
    np.random.seed(42)
    returns = np.random.normal(loc=0, scale=0.01, size=len(dates))
    prices = 100 * np.exp(np.cumsum(returns))
    return pd.DataFrame({"Date": dates, "Price": prices})


def _add_moving_averages(df: pd.DataFrame) -> pd.DataFrame:
    """Add moving averages to a DataFrame with a Price column."""
    out = df.copy()
    for w in (50, 100, 200):
        out[f"MA_{w}"] = out["Price"].rolling(w, min_periods=1).mean()
    return out


def show_score_history(asset_name: str) -> None:
    """
    Display DMAS/Technical/Momentum history chart below the price chart.
    Shows score evolution over time with rating zones.
    """
    try:
        from market_compass.subtitle_generator.history_tracker import get_tracker
    except ImportError:
        return  # History tracker not available

    tracker = get_tracker()
    df = tracker.get_dataframe(asset_name)

    if df.empty or len(df) < 2:
        st.caption("📊 Score history will appear after 2+ weeks of data")
        return

    # Create chart
    fig = go.Figure()

    # DMAS line (thicker, primary)
    fig.add_trace(go.Scatter(
        x=df.index,
        y=df["dmas"],
        name="DMAS",
        line=dict(width=3, color="#1f77b4"),
        mode="lines+markers"
    ))

    # Technical line
    fig.add_trace(go.Scatter(
        x=df.index,
        y=df["technical_score"],
        name="Technical",
        line=dict(width=2, color="#2ca02c"),
        mode="lines"
    ))

    # Momentum line
    fig.add_trace(go.Scatter(
        x=df.index,
        y=df["momentum_score"],
        name="Momentum",
        line=dict(width=2, color="#ff7f0e"),
        mode="lines"
    ))

    # Rating zones (background colors)
    fig.add_hrect(y0=70, y1=100, fillcolor="green", opacity=0.08,
                  annotation_text="Bullish", annotation_position="top left")
    fig.add_hrect(y0=55, y1=70, fillcolor="lightgreen", opacity=0.08,
                  annotation_text="Constructive", annotation_position="top left")
    fig.add_hrect(y0=45, y1=55, fillcolor="gray", opacity=0.08,
                  annotation_text="Neutral", annotation_position="top left")
    fig.add_hrect(y0=30, y1=45, fillcolor="orange", opacity=0.08,
                  annotation_text="Cautious", annotation_position="top left")
    fig.add_hrect(y0=0, y1=30, fillcolor="red", opacity=0.08,
                  annotation_text="Bearish", annotation_position="top left")

    fig.update_layout(
        title="Score History",
        yaxis_title="Score",
        yaxis=dict(range=[0, 100]),
        height=300,
        margin=dict(t=40, b=20),
        legend=dict(orientation="h", yanchor="bottom", y=1.02),
        hovermode="x unified"
    )

    st.plotly_chart(fig, use_container_width=True)

    # Show context if available
    context = tracker.get_context_for_subtitle(asset_name)
    if context:
        st.caption(f"📈 {context}")


def _build_fallback_figure(
    df_full: pd.DataFrame, anchor_date: pd.Timestamp | None = None
) -> go.Figure:
    """
    Build a Plotly figure using synthetic data when no Excel file is loaded.
    """
    if df_full.empty:
        return go.Figure()

    today = df_full["Date"].max().normalize()
    # Determine the lookback window for the fallback chart based on the
    # currently selected analysis timeframe.  When running under
    # Streamlit the ``ta_timeframe_days`` key will be present in
    # ``st.session_state``; otherwise it falls back to three months
    # (90 days).  This ensures the synthetic fallback chart aligns
    # with the timeframe used for real data.
    try:
        lookback_days = int(st.session_state.get("ta_timeframe_days", 90))  # type: ignore
    except Exception:
        lookback_days = 90
    start = today - pd.Timedelta(days=lookback_days)
    df = df_full[df_full["Date"].between(start, today)].reset_index(drop=True)

    fig = go.Figure()
    fig.add_trace(
        go.Scatter(
            x=df["Date"],
            y=df["Price"],
            mode="lines",
            name="S&P 500 Price",
            line=dict(color="#153D64", width=2.5),
        )
    )
    fig.add_trace(
        go.Scatter(
            x=df["Date"],
            y=df.get("MA_50", df["Price"]),
            mode="lines",
            name="50-day MA",
            line=dict(color="#008000", width=1.5),
        )
    )
    fig.add_trace(
        go.Scatter(
            x=df["Date"],
            y=df.get("MA_100", df["Price"]),
            mode="lines",
            name="100-day MA",
            line=dict(color="#FFA500", width=1.5),
        )
    )
    fig.add_trace(
        go.Scatter(
            x=df["Date"],
            y=df.get("MA_200", df["Price"]),
            mode="lines",
            name="200-day MA",
            line=dict(color="#FF0000", width=1.5),
        )
    )

    hi, lo = df["Price"].max(), df["Price"].min()
    span = hi - lo
    for lvl in [hi, hi - 0.236 * span, hi - 0.382 * span, hi - 0.5 * span, hi - 0.618 * span, lo]:
        fig.add_hline(
            y=lvl, line=dict(color="grey", dash="dash", width=1), opacity=0.6
        )

    if anchor_date is not None:
        subset = df_full[df_full["Date"].between(anchor_date, today)].copy()
        if not subset.empty:
            X = subset["Date"].map(pd.Timestamp.toordinal).to_numpy().reshape(-1, 1)
            y_vals = subset["Price"].to_numpy()
            model = LinearRegression().fit(X, y_vals)
            trend = model.predict(X)
            resid = y_vals - trend
            upper = trend + resid.max()
            lower = trend + resid.min()
            uptrend = model.coef_[0] > 0
            lineclr = "green" if uptrend else "red"
            fillclr = "rgba(0,150,0,0.25)" if uptrend else "rgba(200,0,0,0.25)"
            fig.add_trace(
                go.Scatter(
                    x=subset["Date"],
                    y=upper,
                    mode="lines",
                    line=dict(color=lineclr, dash="dash"),
                    showlegend=False,
                )
            )
            fig.add_trace(
                go.Scatter(
                    x=subset["Date"],
                    y=lower,
                    mode="lines",
                    line=dict(color=lineclr, dash="dash"),
                    fill="tonexty",
                    fillcolor=fillclr,
                    showlegend=False,
                )
            )

    fig.update_layout(
        margin=dict(l=30, r=30, t=60, b=40),
        showlegend=True,
        legend=dict(
            orientation="h",
            yanchor="bottom",
            y=1.12,
            xanchor="center",
            x=0.5,
            font=dict(size=12),
        ),
        xaxis_title=None,
        yaxis_title=None,
        xaxis=dict(showgrid=False),
        yaxis=dict(showgrid=False, zeroline=False),
    )
    return fig


###############################################################################
# Streamlit configuration
###############################################################################


st.sidebar.title("Navigation")
page = st.sidebar.radio(
    "Select page", [
        "Upload", 
        "YTD Update", 
        "Technical Analysis", 
        "Market Breadth",
        "Generate Presentation",
    ]
)


def show_upload_page():
    """Handle file uploads for Excel and PowerPoint templates."""
    st.sidebar.header("Upload files")
    excel_file = st.sidebar.file_uploader(
        "Upload consolidated Excel file", type=["xlsx", "xlsm", "xls"], key="excel_upload"
    )
    if excel_file is not None:
        st.session_state["excel_file"] = excel_file

        # Load transition sheet to pre-populate fields
        try:
            import tempfile
            from pathlib import Path
            from transition_loader import read_transition_sheet, apply_transition_data_to_session_state

            # Save to temp file to read transition sheet
            with tempfile.NamedTemporaryFile(delete=False, suffix=".xlsx") as tmp:
                tmp.write(excel_file.getbuffer())
                tmp.flush()
                temp_path = Path(tmp.name)

            # Read transition data and apply to session state
            transition_data = read_transition_sheet(temp_path)
            apply_transition_data_to_session_state(transition_data, st.session_state)

            # Clean up temp file
            temp_path.unlink()
        except Exception as e:
            print(f"Warning: Could not load transition data: {e}")

    # =========================================================================
    # DATA AS OF DATE PICKER
    # =========================================================================
    if excel_file is not None:
        st.sidebar.markdown("---")
        st.sidebar.subheader("📅 Data Configuration")

        # Get max date from Excel data
        try:
            import tempfile
            with tempfile.NamedTemporaryFile(delete=False, suffix=".xlsx") as tmp:
                tmp.write(excel_file.getbuffer())
                tmp.flush()
                date_check_path = Path(tmp.name)

            # Read full sheet to let pandas properly parse dates (usecols breaks date parsing)
            df_dates = pd.read_excel(date_check_path, sheet_name="data_prices")
            df_dates = df_dates.drop(index=0)
            df_dates = df_dates[df_dates[df_dates.columns[0]] != "DATES"]
            df_dates["Date"] = pd.to_datetime(df_dates[df_dates.columns[0]], errors="coerce")
            df_dates = df_dates.dropna(subset=["Date"])
            max_date = df_dates["Date"].max().date()

            # Debug logging
            print(f"[Calendar] Date column type: {df_dates[df_dates.columns[0]].dtype}")
            print(f"[Calendar] Sample dates: {df_dates['Date'].tail(3).tolist()}")
            print(f"[Calendar] Max date: {max_date}")

            # Initialize session state if not set
            if "data_as_of" not in st.session_state:
                st.session_state["data_as_of"] = max_date

            # Date picker
            data_as_of = st.sidebar.date_input(
                "Data As Of",
                value=st.session_state.get("data_as_of", max_date),
                max_value=max_date,
                help=f"Latest date in Excel: {max_date.strftime('%d %b %Y')}. Adjust if markets were closed."
            )
            st.session_state["data_as_of"] = data_as_of

            # Load historical data NOW that we have the correct date
            # This ensures history lookup uses calendar date, not today's date
            _load_historical_dmas_to_session()

            # Show selected date info
            if data_as_of < max_date:
                st.sidebar.info(f"📆 Using data up to {data_as_of.strftime('%d %b %Y')}")

        except Exception as e:
            st.sidebar.warning(f"⚠️ Could not read dates from Excel: {str(e)[:50]}")

    pptx_file = st.sidebar.file_uploader(
        "Upload PowerPoint template", type=["pptx", "pptm"], key="ppt_upload"
    )
    if pptx_file is not None:
        st.session_state["pptx_file"] = pptx_file

    # =========================================================================
    # GLOBAL ANALYSIS BUTTON
    # =========================================================================
    if excel_file is not None:
        st.sidebar.markdown("---")
        st.sidebar.subheader("📊 Auto-Analysis")

        # Check if Claude API is available for better subtitles
        try:
            from assessment_integration import is_claude_available, CLAUDE_GEN_AVAILABLE
            claude_ready = CLAUDE_GEN_AVAILABLE and is_claude_available()
        except ImportError:
            claude_ready = False

        if claude_ready:
            st.sidebar.info("🤖 Claude API: Ready for unique subtitles")
        else:
            st.sidebar.warning("⚠️ Claude API not configured - using pattern-based subtitles")

        # Test CoinMarketCap API button
        if st.sidebar.button("🪙 Fetch Crypto Mkt Cap", help="Test CoinMarketCap API - fetches live market caps"):
            with st.spinner("Fetching crypto market caps..."):
                try:
                    from market_compass.technical_slide.crypto_data import fetch_crypto_market_caps
                    crypto_caps = fetch_crypto_market_caps()

                    if all(v == "—" for v in crypto_caps.values()):
                        st.sidebar.error("❌ API key not found. Add to .streamlit/secrets.toml")
                    else:
                        st.sidebar.success("✅ CoinMarketCap API working!")
                        for name, cap in crypto_caps.items():
                            st.sidebar.write(f"  {name}: **{cap}**")
                except Exception as e:
                    st.sidebar.error(f"❌ Error: {str(e)[:100]}")

        # Model selection for subtitle generation
        st.sidebar.markdown("---")
        st.sidebar.subheader("🤖 Model Selection")
        model_choice = st.sidebar.radio(
            "Claude Model:",
            options=["Haiku 3.5 (Current)", "Haiku 4.5 (New)"],
            index=0,
            help="Haiku 4.5 has better reasoning but costs ~25% more",
            key="model_choice"
        )
        # Map selection to model_key
        model_key = "haiku_35" if "3.5" in model_choice else "haiku_45"
        st.session_state["claude_model_key"] = model_key

        if st.sidebar.button("🚀 Run Full Analysis", type="primary", help="Compute technical scores, assessments, and subtitles for all assets"):
            with st.spinner("Running analysis for all assets..."):
                try:
                    # Import required modules
                    from technical_score_wrapper import compute_dmas_scores
                    from assessment_integration import (
                        generate_assessment_and_subtitle,
                        generate_claude_subtitles_batch,
                        is_claude_available,
                        CLAUDE_GEN_AVAILABLE,
                    )

                    # Fetch live crypto market caps from CoinMarketCap
                    try:
                        from market_compass.technical_slide.crypto_data import fetch_crypto_market_caps
                        crypto_caps = fetch_crypto_market_caps()
                        if not all(v == "—" for v in crypto_caps.values()):
                            st.session_state["crypto_market_caps"] = crypto_caps
                            st.sidebar.success("🪙 Crypto market caps fetched")
                    except Exception as e:
                        print(f"[CoinMarketCap] Error in full analysis: {e}")

                    # Save Excel to temp file
                    with tempfile.NamedTemporaryFile(delete=False, suffix=".xlsx") as tmp:
                        tmp.write(excel_file.getbuffer())
                        tmp.flush()
                        analysis_temp_path = Path(tmp.name)

                    # Read price data
                    df_prices = pd.read_excel(analysis_temp_path, sheet_name="data_prices")
                    df_prices = df_prices.drop(index=0)
                    df_prices = df_prices[df_prices[df_prices.columns[0]] != "DATES"]
                    df_prices["Date"] = pd.to_datetime(df_prices[df_prices.columns[0]], errors="coerce")

                    # Filter by "Data As Of" date if set
                    if "data_as_of" in st.session_state:
                        df_prices = df_prices[df_prices["Date"] <= pd.Timestamp(st.session_state["data_as_of"])]

                    # Calculate commodity market caps from reserves × spot price
                    try:
                        from market_compass.technical_slide.commodity_data import calculate_all_commodity_market_caps
                        commo_caps = calculate_all_commodity_market_caps(df_prices)
                        if not all(v == "—" for v in commo_caps.values()):
                            st.session_state["commodity_market_caps"] = commo_caps
                            st.sidebar.success("⚡ Commodity market caps calculated")
                    except Exception as e:
                        print(f"[Commodity] Error in full analysis: {e}")

                    # Asset mapping: ticker_key -> (Bloomberg ticker, display name)
                    asset_map = {
                        # Equity
                        "spx": ("SPX Index", "S&P 500"),
                        "csi": ("SHSZ300 Index", "CSI 300"),
                        "nikkei": ("NKY Index", "Nikkei 225"),
                        "tasi": ("SASEIDX Index", "TASI"),
                        "sensex": ("SENSEX Index", "Sensex"),
                        "dax": ("DAX Index", "Dax"),
                        "smi": ("SMI Index", "SMI"),
                        "ibov": ("IBOV Index", "IBOV"),
                        "mexbol": ("MEXBOL Index", "MEXBOL"),
                        # Commodity
                        "gold": ("GCA COMDTY", "Gold"),
                        "silver": ("SIA COMDTY", "Silver"),
                        "platinum": ("XPT COMDTY", "Platinum"),
                        "palladium": ("XPD CURNCY", "Palladium"),
                        "oil": ("CL1 COMDTY", "Oil"),
                        "copper": ("LP1 COMDTY", "Copper"),
                        # Crypto
                        "bitcoin": ("XBTUSD CURNCY", "Bitcoin"),
                        "ethereum": ("XETUSD CURNCY", "Ethereum"),
                        "ripple": ("XRPUSD CURNCY", "Ripple"),
                        "solana": ("XSOUSD CURNCY", "Solana"),
                        "binance": ("XBIUSD CURNCY", "Binance"),
                    }

                    results = []
                    assets_for_claude = []  # Collect assets for batch subtitle generation
                    prices_dict = {}  # Store prices for Claude batch generation
                    progress_bar = st.sidebar.progress(0)
                    status_text = st.sidebar.empty()

                    # Phase 1: Compute scores for all assets
                    for idx, (ticker_key, (bbg_ticker, display_name)) in enumerate(asset_map.items()):
                        try:
                            status_text.text(f"Computing scores for {display_name}...")

                            # Find matching column (case-insensitive, flexible matching)
                            matching_col = None
                            bbg_ticker_upper = bbg_ticker.upper()

                            for col in df_prices.columns:
                                if isinstance(col, str) and col.upper() == bbg_ticker_upper:
                                    matching_col = col
                                    break

                            if matching_col is None:
                                st.sidebar.warning(f"⚠️ {display_name}: No data found (looking for {bbg_ticker})")
                                continue

                            # Extract price series
                            prices_df = df_prices[["Date", matching_col]].copy()
                            prices_df.columns = ["Date", "Price"]
                            prices_df["Price"] = pd.to_numeric(prices_df["Price"], errors="coerce")
                            prices_df = prices_df.dropna(subset=["Date", "Price"]).sort_values("Date").reset_index(drop=True)

                            if len(prices_df) < 200:
                                st.sidebar.warning(f"⚠️ {display_name}: Insufficient data ({len(prices_df)} days)")
                                continue

                            prices = prices_df["Price"]
                            prices_dict[ticker_key] = prices

                            # Compute technical scores
                            scores = compute_dmas_scores(prices, ticker=bbg_ticker, excel_path=str(analysis_temp_path))
                            tech_score = scores["technical_score"]
                            mom_score = scores["momentum_score"]
                            dmas = scores["dmas"]
                            rsi = scores.get("rsi")

                            # Store scores in session state
                            st.session_state[f"{ticker_key}_tech_score"] = tech_score
                            st.session_state[f"{ticker_key}_mom_score"] = mom_score
                            st.session_state[f"{ticker_key}_dmas"] = dmas

                            # Get last week DMAS (if available from transition sheet)
                            dmas_prev = st.session_state.get(f"{ticker_key}_last_week_avg", dmas)

                            # Collect for batch subtitle generation
                            assets_for_claude.append({
                                "ticker_key": ticker_key,
                                "asset_name": display_name,
                                "dmas": dmas,
                                "technical_score": tech_score,
                                "momentum_score": mom_score,
                                "dmas_prev_week": dmas_prev,
                                "rsi": rsi,
                            })

                        except Exception as e:
                            st.sidebar.error(f"❌ {display_name}: {str(e)[:50]}")
                            import traceback
                            print(f"Error processing {display_name}: {traceback.format_exc()}")

                        # Update progress (0-50% for score computation)
                        progress_bar.progress((idx + 1) / len(asset_map) * 0.5)

                    # Phase 2: Generate subtitles (batch via Claude API if available)
                    use_claude = CLAUDE_GEN_AVAILABLE and is_claude_available()
                    if use_claude and assets_for_claude:
                        status_text.text("Generating subtitles via Claude API...")
                        # Pass data_as_of for history storage keyed by calendar date
                        data_as_of_str = None
                        if "data_as_of" in st.session_state:
                            data_as_of_str = st.session_state["data_as_of"].strftime("%Y-%m-%d")
                        # Get selected model from session state
                        selected_model_key = st.session_state.get("claude_model_key", "haiku_35")
                        subtitle_results = generate_claude_subtitles_batch(
                            assets_for_claude, prices_dict,
                            data_as_of=data_as_of_str,
                            model_key=selected_model_key
                        )
                    else:
                        # Fall back to individual generation
                        subtitle_results = {}
                        for asset in assets_for_claude:
                            ticker_key = asset["ticker_key"]
                            result = generate_assessment_and_subtitle(
                                ticker_key=ticker_key,
                                asset_name=asset["asset_name"],
                                prices=prices_dict.get(ticker_key),
                                dmas=asset["dmas"],
                                technical_score=asset["technical_score"],
                                momentum_score=asset["momentum_score"],
                                dmas_prev_week=asset["dmas_prev_week"],
                                subtitle_generator=st.session_state.get('subtitle_generator')
                            )
                            subtitle_results[ticker_key] = result

                    # Phase 2.5: Generate YTD recap subtitles for equity, commodities, crypto
                    status_text.text("Generating YTD recap subtitles...")
                    try:
                        from ytd_perf.equity_ytd import get_equity_ytd_series
                        from ytd_perf.commodity_ytd import get_commodity_ytd_series
                        from ytd_perf.crypto_ytd import get_crypto_ytd_series
                        from market_compass.subtitle_generator.claude_generator import generate_all_recaps

                        price_mode = st.session_state.get("price_mode", "Last Price")

                        # Define allowed commodities for YTD recap (only these 7)
                        YTD_RECAP_COMMODITIES = [
                            "SIA Comdty",   # Silver
                            "XPT Comdty",   # Platinum
                            "XPD Curncy",   # Palladium
                            "GCA Comdty",   # Gold
                            "CL1 Comdty",   # Oil (WTI)
                            "UXA1 Comdty",  # Uranium
                            "LP1 Comdty",   # Copper
                        ]

                        # Get YTD series for each asset class
                        # Note: Empty tickers list means include all tickers (fixed in ytd_perf modules)
                        eq_tickers = st.session_state.get("selected_eq_tickers", [])
                        cr_tickers = st.session_state.get("selected_cr_tickers", [])

                        df_eq = get_equity_ytd_series(str(analysis_temp_path), tickers=eq_tickers, price_mode=price_mode)
                        df_co = get_commodity_ytd_series(str(analysis_temp_path), tickers=YTD_RECAP_COMMODITIES, price_mode=price_mode)
                        df_cr = get_crypto_ytd_series(str(analysis_temp_path), tickers=cr_tickers, price_mode=price_mode)

                        # Helper to extract YTD performance from DataFrame
                        def extract_ytd_perf(df):
                            if df.empty:
                                return []
                            last_row = df.iloc[-1]
                            perf_data = []
                            for col in df.columns:
                                if col != 'Date':
                                    ytd_val = last_row[col]
                                    if pd.notna(ytd_val):
                                        perf_data.append({
                                            'asset': col,
                                            'ytd_pct': float(ytd_val),
                                            '1w_pct': 0.0,
                                            '1m_pct': 0.0
                                        })
                            return perf_data

                        # Build performance data and generate recaps
                        perf_data = {
                            'equity': extract_ytd_perf(df_eq),
                            'commodities': extract_ytd_perf(df_co),
                            'crypto': extract_ytd_perf(df_cr)
                        }

                        # Only generate if we have data
                        if any(perf_data.values()):
                            recaps = generate_all_recaps(perf_data)
                            if 'equity' in recaps:
                                st.session_state['eq_subtitle'] = recaps['equity']
                            if 'commodities' in recaps:
                                st.session_state['co_subtitle'] = recaps['commodities']
                            if 'crypto' in recaps:
                                st.session_state['cr_subtitle'] = recaps['crypto']
                    except Exception as e:
                        print(f"YTD recap generation failed: {e}")
                        # Non-fatal - continue with analysis

                    progress_bar.progress(0.75)

                    # Phase 3: Store results
                    status_text.text("Storing results...")
                    for asset in assets_for_claude:
                        ticker_key = asset["ticker_key"]
                        result = subtitle_results.get(ticker_key, {})

                        # Store results in session state
                        st.session_state[f"{ticker_key}_assessment"] = result.get("assessment", "Neutral")
                        st.session_state[f"{ticker_key}_subtitle"] = result.get("subtitle", "Technical analysis under review.")
                        st.session_state[f"{ticker_key}_selected_view"] = result.get("assessment", "Neutral")

                        results.append({
                            "Asset": asset["asset_name"],
                            "DMAS": f"{asset['dmas']:.0f}",
                            "Tech": f"{asset['technical_score']:.0f}",
                            "Mom": f"{asset['momentum_score']:.0f}",
                            "Assessment": result.get("assessment", "Neutral"),
                            "Subtitle": result.get("subtitle", "")[:50] + "..." if len(result.get("subtitle", "")) > 50 else result.get("subtitle", "")
                        })

                    progress_bar.progress(1.0)

                    # Clean up
                    analysis_temp_path.unlink()
                    progress_bar.empty()
                    status_text.empty()

                    # Show results
                    if results:
                        # Show which model was used
                        model_display = "Haiku 3.5" if selected_model_key == "haiku_35" else "Haiku 4.5"
                        st.sidebar.success(f"✅ Analysis complete! Processed {len(results)} assets")
                        st.sidebar.info(f"🤖 Model used: {model_display}")

                        # Show summary in main area
                        st.subheader("Analysis Results")
                        st.dataframe(pd.DataFrame(results), use_container_width=True)
                        st.info(f"📝 All scores, assessments, and subtitles have been auto-generated using **{model_display}**. Navigate to asset tabs to review and edit if needed.")
                    else:
                        st.sidebar.warning("⚠️ No assets were successfully processed")

                except ImportError as e:
                    st.sidebar.error(f"❌ Missing modules: {e}")
                    st.sidebar.info("Make sure herculis-technical-score and assessment modules are available")
                except Exception as e:
                    st.sidebar.error(f"❌ Analysis failed: {str(e)[:100]}")
                    import traceback
                    print(f"Global analysis error: {traceback.format_exc()}")

    st.sidebar.success("Files uploaded. Navigate to other pages to continue.")

    # Allow the user to choose between using the last recorded price (which may
    # be an intraday or current price) and the last close price (i.e. the
    # previous trading day's close).  The choice is stored in session state
    # and will affect how data is loaded and displayed elsewhere in the app.
    # Persist the selected price mode across pages.  Use the previously selected
    # value from session state (if any) to determine the default index.  If no
    # value has been stored yet, default to "Last Price".
    current_mode = st.session_state.get("price_mode", "Last Price")
    options = ["Last Price", "Last Close"]
    default_index = options.index(current_mode) if current_mode in options else 0
    price_mode = st.sidebar.radio(
        "Price mode",
        options=options,
        index=default_index,
        help=(
            "Select 'Last Close' to use the previous day's closing prices for all markets. "
            "Select 'Last Price' to use the most recent price in the data (which may be intraday)."
        ),
        key="price_mode_select",
    )
    st.session_state["price_mode"] = price_mode


def show_ytd_update_page():
    """Display YTD update charts and configuration."""
    st.sidebar.header("YTD Update")
    if "excel_file" not in st.session_state:
        st.sidebar.error("Please upload an Excel file on the Upload page first.")
        st.stop()

    # Lazy import heavy modules
    from ytd_perf.loader_update import load_data
    from utils import adjust_prices_for_mode
    from ytd_perf.equity_ytd import get_equity_ytd_series, create_equity_chart
    from ytd_perf.commodity_ytd import get_commodity_ytd_series, create_commodity_chart
    from ytd_perf.crypto_ytd import get_crypto_ytd_series, create_crypto_chart

    prices_df, params_df = load_data(st.session_state["excel_file"])
    # Determine whether to use the last price or the last close using
    # the centralised adjust_prices_for_mode helper.  This returns an
    # adjusted DataFrame and the effective date used for YTD calculations.
    used_date = None
    if not prices_df.empty:
        price_mode = st.session_state.get("price_mode", "Last Price")
        prices_df, used_date = adjust_prices_for_mode(prices_df, price_mode)
    # Display a caption indicating which date's prices are being used
    if used_date is not None:
        price_mode = st.session_state.get("price_mode", "Last Price")
        if price_mode == "Last Close":
            st.sidebar.caption(f"Prices as of {used_date.strftime('%d/%m/%Y')} close")
        else:
            st.sidebar.caption(f"Prices as of {used_date.strftime('%d/%m/%Y')}")

    # Equities configuration
    st.sidebar.subheader("Equities")
    eq_params = params_df[params_df["Asset Class"] == "Equity"]
    eq_name_to_ticker = {row["Name"]: row["Tickers"] for _, row in eq_params.iterrows()}
    eq_names_available = eq_params["Name"].tolist()
    default_eq = [
        name
        for name in [
            "Dax",
            "Ibov",
            "S&P 500",
            "Sensex",
            "SMI",
            "CSI 300",
            "Nikkei 225",
            "TASI",
            "Mexbol",
        ]
        if name in eq_names_available
    ]
    selected_eq_names = st.sidebar.multiselect(
        "Select equity indices",
        options=eq_names_available,
        default=st.session_state.get("selected_eq_names", default_eq),
        key="eq_indices",
    )
    st.session_state["selected_eq_names"] = selected_eq_names
    eq_tickers = [eq_name_to_ticker[name] for name in selected_eq_names]
    eq_subtitle = st.sidebar.text_input(
        "Equity subtitle", value=st.session_state.get("eq_subtitle", ""), key="eq_subtitle_input"
    )
    st.session_state["eq_subtitle"] = eq_subtitle

    # Commodities configuration
    st.sidebar.subheader("Commodities")
    co_params = params_df[params_df["Asset Class"] == "Commodity"]
    co_name_to_ticker = {row["Name"]: row["Tickers"] for _, row in co_params.iterrows()}
    co_names_available = co_params["Name"].tolist()
    default_co = [
        name
        for name in ["Gold", "Silver", "Oil (WTI)", "Platinum", "Copper", "Uranium","Palladium"]
        if name in co_names_available
    ]
    selected_co_names = st.sidebar.multiselect(
        "Select commodity indices",
        options=co_names_available,
        default=st.session_state.get("selected_co_names", default_co),
        key="co_indices",
    )
    st.session_state["selected_co_names"] = selected_co_names
    co_tickers = [co_name_to_ticker[name] for name in selected_co_names]
    co_subtitle = st.sidebar.text_input(
        "Commodity subtitle", value=st.session_state.get("co_subtitle", ""), key="co_subtitle_input"
    )
    st.session_state["co_subtitle"] = co_subtitle

    # Crypto configuration
    st.sidebar.subheader("Cryptocurrencies")
    cr_params = params_df[params_df["Asset Class"] == "Crypto"]
    cr_name_to_ticker = {row["Name"]: row["Tickers"] for _, row in cr_params.iterrows()}
    cr_names_available = cr_params["Name"].tolist()
    default_cr = [
        name
        for name in ["Ripple", "Bitcoin", "Binance", "Ethereum", "Solana","Ton"]
        if name in cr_names_available
    ]
    selected_cr_names = st.sidebar.multiselect(
        "Select crypto indices",
        options=cr_names_available,
        default=st.session_state.get("selected_cr_names", default_cr),
        key="cr_indices",
    )
    st.session_state["selected_cr_names"] = selected_cr_names
    cr_tickers = [cr_name_to_ticker[name] for name in selected_cr_names]
    cr_subtitle = st.sidebar.text_input(
        "Crypto subtitle", value=st.session_state.get("cr_subtitle", ""), key="cr_subtitle_input"
    )
    st.session_state["cr_subtitle"] = cr_subtitle

    # Persist selections
    st.session_state["selected_eq_tickers"] = eq_tickers
    st.session_state["selected_co_tickers"] = co_tickers
    st.session_state["selected_cr_tickers"] = cr_tickers

    st.header("YTD Performance Charts")

    # Compute YTD series for all asset classes
    price_mode = st.session_state.get("price_mode", "Last Price")
    with st.expander("Equity Chart", expanded=True):
        df_eq = get_equity_ytd_series(
            st.session_state["excel_file"], tickers=eq_tickers, price_mode=price_mode
        )
        st.pyplot(create_equity_chart(df_eq))
    with st.expander("Commodity Chart", expanded=False):
        df_co = get_commodity_ytd_series(
            st.session_state["excel_file"], tickers=co_tickers, price_mode=price_mode
        )
        st.pyplot(create_commodity_chart(df_co))
    with st.expander("Crypto Chart", expanded=False):
        df_cr = get_crypto_ytd_series(
            st.session_state["excel_file"], tickers=cr_tickers, price_mode=price_mode
        )
        st.pyplot(create_crypto_chart(df_cr))

    st.sidebar.success("Configure YTD charts, then go to 'Generate Presentation'.")


def show_technical_analysis_page():
    """Display the technical analysis interface for Equity (SPX) and other asset classes."""
    st.sidebar.header("Technical Analysis")
    asset_class = st.sidebar.radio(
        "Asset class", ["Equity", "Commodity", "Crypto"], index=0
    )

    # -------------------------------------------------------------------
    # Analysis timeframe selection
    # -------------------------------------------------------------------
    # Allow the user to choose the lookback horizon for all technical charts.
    # Provide multiple timeframe options: 3, 4, 5, and 6 months so users
    # can toggle between different analysis windows. The default is 4 months
    # as requested by management.
    timeframe_options: dict[str, int] = {
        "3 months": 90,
        "4 months": 120,
        "5 months": 150,
        "6 months": 180,
    }
    # Determine the default based on any previously stored selection; fall
    # back to 4 months if none is present.
    default_tf_label = st.session_state.get("ta_timeframe_label", "4 months")
    tf_labels = list(timeframe_options.keys())
    if default_tf_label not in tf_labels:
        default_tf_idx = 1  # Index 1 is "4 months"
    else:
        default_tf_idx = tf_labels.index(default_tf_label)
    selected_tf_label = st.sidebar.selectbox(
        "Analysis timeframe",
        options=tf_labels,
        index=default_tf_idx,
        key="ta_timeframe_select",
    )
    # Persist the selection and derive the numeric days value
    st.session_state["ta_timeframe_label"] = selected_tf_label
    st.session_state["ta_timeframe_days"] = timeframe_options[selected_tf_label]

    # Propagate the chosen timeframe into technical analysis modules that
    # support configurable lookback windows.  The Mexbol and Palladium
    # modules define a ``PLOT_LOOKBACK_DAYS`` constant which can be
    # overridden at runtime.  We attempt to set this attribute here.
    try:
        import technical_analysis.equity.mexbol as _mex_module  # type: ignore
        _mex_module.PLOT_LOOKBACK_DAYS = st.session_state["ta_timeframe_days"]
    except Exception:
        pass

    try:
        import technical_analysis.equity.csi as _csi_module  # same package as your CSI code
        if hasattr(_csi_module, "PLOT_LOOKBACK_DAYS"):
            _csi_module.PLOT_LOOKBACK_DAYS = st.session_state["ta_timeframe_days"]
    except Exception:
        pass
    
    try:
        import technical_analysis.equity.dax as _dax_module  # same package as your DAX code
        if hasattr(_dax_module, "PLOT_LOOKBACK_DAYS"):
            _dax_module.PLOT_LOOKBACK_DAYS = st.session_state["ta_timeframe_days"]
    except Exception:
        pass

    try:
        import technical_analysis.equity.ibov as _ibov_module  # same package as your IBOV code
        if hasattr(_ibov_module, "PLOT_LOOKBACK_DAYS"):
            _ibov_module.PLOT_LOOKBACK_DAYS = st.session_state["ta_timeframe_days"]
    except Exception:
        pass

    try:
        import technical_analysis.equity.nikkei as _nikkei_module  # same package as your NIKKEI code
        if hasattr(_nikkei_module, "PLOT_LOOKBACK_DAYS"):
            _nikkei_module.PLOT_LOOKBACK_DAYS = st.session_state["ta_timeframe_days"]
    except Exception:
        pass

    try:
        import technical_analysis.equity.sensex as _sensex_module  # same package as your SENSEX code
        if hasattr(_sensex_module, "PLOT_LOOKBACK_DAYS"):
            _sensex_module.PLOT_LOOKBACK_DAYS = st.session_state["ta_timeframe_days"]
    except Exception:
        pass

    try:
        import technical_analysis.equity.smi as _smi_module  # same package as your SMI code
        if hasattr(_smi_module, "PLOT_LOOKBACK_DAYS"):
            _smi_module.PLOT_LOOKBACK_DAYS = st.session_state["ta_timeframe_days"]
    except Exception:
        pass

    try:
        import technical_analysis.equity.spx as _spx_module  # same package as your SPX code
        if hasattr(_spx_module, "PLOT_LOOKBACK_DAYS"):
            _spx_module.PLOT_LOOKBACK_DAYS = st.session_state["ta_timeframe_days"]
    except Exception:
        pass

    try:
        import technical_analysis.equity.tasi as _tasi_module  # same package as your TASI code
        if hasattr(_tasi_module, "PLOT_LOOKBACK_DAYS"):
            _tasi_module.PLOT_LOOKBACK_DAYS = st.session_state["ta_timeframe_days"]
    except Exception:
        pass

    try:
        import technical_analysis.commodity.copper as _copper_module  # same package as your COPPER code
        if hasattr(_copper_module, "PLOT_LOOKBACK_DAYS"):
            _copper_module.PLOT_LOOKBACK_DAYS = st.session_state["ta_timeframe_days"]
    except Exception:
        pass

    try:
        import technical_analysis.commodity.gold as _gold_module  # same package as your GOLD code
        if hasattr(_gold_module, "PLOT_LOOKBACK_DAYS"):
            _gold_module.PLOT_LOOKBACK_DAYS = st.session_state["ta_timeframe_days"]
    except Exception:
        pass

    try:
        import technical_analysis.commodity.oil as _oil_module  # same package as your OIL code
        if hasattr(_oil_module, "PLOT_LOOKBACK_DAYS"):
            _oil_module.PLOT_LOOKBACK_DAYS = st.session_state["ta_timeframe_days"]
    except Exception:
        pass

    try:
        import technical_analysis.commodity.palladium as _palladium_module  # same package as your PALLADIUM code
        if hasattr(_palladium_module, "PLOT_LOOKBACK_DAYS"):
            _palladium_module.PLOT_LOOKBACK_DAYS = st.session_state["ta_timeframe_days"]
    except Exception:
        pass

    try:
        import technical_analysis.commodity.platinum as _platinum_module  # same package as your PLATINUM code
        if hasattr(_platinum_module, "PLOT_LOOKBACK_DAYS"):
            _platinum_module.PLOT_LOOKBACK_DAYS = st.session_state["ta_timeframe_days"]
    except Exception:
        pass

    try:
        import technical_analysis.commodity.silver as _silver_module  # same package as your SILVER code
        if hasattr(_silver_module, "PLOT_LOOKBACK_DAYS"):
            _silver_module.PLOT_LOOKBACK_DAYS = st.session_state["ta_timeframe_days"]
    except Exception:
        pass

    try:
        import technical_analysis.crypto.binance as _binance_module  # same package as your BINANCE code
        if hasattr(_binance_module, "PLOT_LOOKBACK_DAYS"):
            _binance_module.PLOT_LOOKBACK_DAYS = st.session_state["ta_timeframe_days"]
    except Exception:
        pass

    try:
        import technical_analysis.crypto.bitcoin as _bitcoin_module  # same package as your BITCOIN code
        if hasattr(_bitcoin_module, "PLOT_LOOKBACK_DAYS"):
            _bitcoin_module.PLOT_LOOKBACK_DAYS = st.session_state["ta_timeframe_days"]
    except Exception:
        pass

    try:
        import technical_analysis.crypto.ethereum as _ethereum_module  # same package as your ETHEREUM code
        if hasattr(_ethereum_module, "PLOT_LOOKBACK_DAYS"):
            _ethereum_module.PLOT_LOOKBACK_DAYS = st.session_state["ta_timeframe_days"]
    except Exception:
        pass

    try:
        import technical_analysis.crypto.ripple as _ripple_module  # same package as your RIPPLE code
        if hasattr(_ripple_module, "PLOT_LOOKBACK_DAYS"):
            _ripple_module.PLOT_LOOKBACK_DAYS = st.session_state["ta_timeframe_days"]
    except Exception:
        pass

    try:
        import technical_analysis.crypto.solana as _solana_module  # same package as your SOLANA code
        if hasattr(_solana_module, "PLOT_LOOKBACK_DAYS"):
            _solana_module.PLOT_LOOKBACK_DAYS = st.session_state["ta_timeframe_days"]
    except Exception:
        pass

    # Provide a clear channel button to reset the regression channel for both indices
    if st.sidebar.button("Clear channel", key="ta_clear_global"):
        # Remove stored anchors for all indices if present
        for key in [
            "spx_anchor",
            "csi_anchor",
            "nikkei_anchor",
            "tasi_anchor",
            "sensex_anchor",
            "dax_anchor",
            "smi_anchor",
            "ibov_anchor",
            "mexbol_anchor",
            # also clear anchors for commodity and crypto assets
            "gold_anchor",
            "silver_anchor",
            "platinum_anchor",
            "palladium_anchor",
            "oil_anchor",
            "copper_anchor",
            "bitcoin_anchor",
            "ethereum_anchor",
            "ripple_anchor",
            "solana_anchor",
            "binance_anchor",
        ]:
            if key in st.session_state:
                st.session_state.pop(key)
        st.rerun()

    excel_available = "excel_file" in st.session_state

    if asset_class == "Equity":
        # Allow the user to select which equity index they wish to analyse.  We
        # provide two options: S&P 500 and CSI 300.  The selection is stored
        # in session state to persist across reruns.
        # Provide index options.  Add Nikkei 225 alongside SPX and CSI.
        # Include SMI (Swiss Market Index) alongside existing indices
        # Include IBOV (Brazil Bovespa) alongside existing indices
        # Add Mexbol to the list of available equity indices
        index_options = ["S&P 500", "CSI 300", "Nikkei 225", "TASI", "Sensex", "Dax", "SMI", "Ibov", "Mexbol"]
        default_index = st.session_state.get("ta_equity_index", "S&P 500")
        selected_index = st.sidebar.selectbox(
            "Select equity index for technical analysis",
            options=index_options,
            index=index_options.index(default_index) if default_index in index_options else 0,
            key="ta_equity_index_select",
        )
        # Persist the selected index
        st.session_state["ta_equity_index"] = selected_index

        # Determine ticker and names based on the selected index
        # Determine ticker and label keys based on the selected index
        if selected_index == "S&P 500":
            ticker = "SPX Index"
            ticker_key = "spx"
            chart_title = "S&P 500 Technical Chart"
        elif selected_index == "CSI 300":
            ticker = "SHSZ300 Index"
            ticker_key = "csi"
            chart_title = "CSI 300 Technical Chart"
        elif selected_index == "Nikkei 225":
            ticker = "NKY Index"
            ticker_key = "nikkei"
            chart_title = "Nikkei 225 Technical Chart"
        elif selected_index == "TASI":
            ticker = "SASEIDX Index"
            ticker_key = "tasi"
            chart_title = "TASI Technical Chart"
        elif selected_index == "Sensex":
            ticker = "SENSEX Index"
            ticker_key = "sensex"
            chart_title = "Sensex Technical Chart"
        elif selected_index == "Dax":
            ticker = "DAX Index"
            ticker_key = "dax"
            chart_title = "DAX Technical Chart"
        elif selected_index == "SMI":
            ticker = "SMI Index"
            ticker_key = "smi"
            chart_title = "SMI Technical Chart"
        elif selected_index == "Ibov":
            ticker = "IBOV Index"
            ticker_key = "ibov"
            chart_title = "Ibov Technical Chart"
        elif selected_index == "Mexbol":
            # Mexbol index configuration
            ticker = "MEXBOL Index"
            ticker_key = "mexbol"
            chart_title = "Mexbol Technical Chart"
        else:
            # Default fallback (should not occur)
            ticker = "SPX Index"
            ticker_key = "spx"
            chart_title = "S&P 500 Technical Chart"

        # Load data for interactive chart (real or synthetic)
        if excel_available:
            with tempfile.NamedTemporaryFile(delete=False, suffix=".xlsx") as tmp:
                tmp.write(st.session_state["excel_file"].getbuffer())
                tmp.flush()
                temp_path = Path(tmp.name)
            df_prices = pd.read_excel(temp_path, sheet_name="data_prices")
            df_prices = df_prices.drop(index=0)
            df_prices = df_prices[df_prices[df_prices.columns[0]] != "DATES"]
            df_prices["Date"] = pd.to_datetime(
                df_prices[df_prices.columns[0]], errors="coerce"
            )
            # Filter by "Data As Of" date if set
            if "data_as_of" in st.session_state:
                df_prices = df_prices[df_prices["Date"] <= pd.Timestamp(st.session_state["data_as_of"])]
            df_prices["Price"] = pd.to_numeric(df_prices[ticker], errors="coerce")
            df_prices = df_prices.dropna(subset=["Date", "Price"]).sort_values("Date").reset_index(
                drop=True
            )
            # Adjust the prices according to the selected price mode using the helper.
            price_mode = st.session_state.get("price_mode", "Last Price")
            df_prices, used_date = adjust_prices_for_mode(df_prices, price_mode)
            df_full = df_prices.copy()
            # Store the used date for later caption display (per index)
            st.session_state[f"{ticker_key}_used_date"] = used_date
        else:
            # Use synthetic series only for SPX; for CSI default to SPX synthetic as fallback
            df_prices = _create_synthetic_spx_series()
            df_full = df_prices.copy()

        min_date = df_prices["Date"].min().date()
        max_date = df_prices["Date"].max().date()

        # Chart with controls in expander
        with st.expander(chart_title, expanded=True):
            # Display a caption indicating which date's prices are being used
            used_date = st.session_state.get(f"{ticker_key}_used_date")
            price_mode = st.session_state.get("price_mode", "Last Price")
            if used_date is not None:
                if price_mode == "Last Close":
                    st.caption(f"Prices as of {used_date.strftime('%d/%m/%Y')} close")
                else:
                    st.caption(f"Prices as of {used_date.strftime('%d/%m/%Y')}")
            # -------------------------------------------------------------------
            # Display technical and momentum scores first
            # -------------------------------------------------------------------
            st.subheader("Technical and momentum scores")
            tech_score = None
            mom_score = None
            if excel_available and df_prices is not None:
                try:
                    # Map display name to Bloomberg ticker
                    ticker_map = {
                        "S&P 500": "SPX Index",
                        "CSI 300": "SHSZ300 Index",
                        "Nikkei 225": "NKY Index",
                        "TASI": "SASEIDX Index",
                        "Sensex": "SENSEX Index",
                        "Dax": "DAX Index",
                        "SMI": "SMI Index",
                        "Ibov": "IBOV Index",
                        "Mexbol": "MEXBOL Index",
                        "Gold": "GCA COMDTY",
                        "Silver": "SIA COMDTY",
                        "Platinum": "XPT COMDTY",
                        "Palladium": "XPD CURNCY",
                        "Oil": "CL1 COMDTY",
                        "Copper": "LP1 COMDTY",
                        "Bitcoin": "XBTUSD CURNCY",
                        "Ethereum": "XETUSD CURNCY",
                        "Ripple": "XRPUSD CURNCY",
                        "Solana": "XSOUSD CURNCY",
                        "Binance": "XBIUSD CURNCY",
                    }

                    bbg_ticker = ticker_map.get(selected_index)
                    if bbg_ticker:
                        # Find matching column (case-insensitive)
                        matching_col = None
                        for col in df_prices.columns:
                            if isinstance(col, str) and col.upper() == bbg_ticker.upper():
                                matching_col = col
                                break

                        if matching_col and len(df_prices) >= 200:
                            # Extract prices
                            prices = df_prices[matching_col].dropna()
                            if len(prices) >= 200:
                                # Import and compute scores
                                from technical_score_wrapper import compute_dmas_scores
                                scores = compute_dmas_scores(prices, ticker=bbg_ticker, excel_path=str(temp_path))
                                tech_score = scores["technical_score"]
                                mom_score = scores["momentum_score"]

                except Exception as e:
                    print(f"Error computing scores for {selected_index}: {e}")
                    tech_score = None
                    mom_score = None

            # Prepare DMAS and table if both scores are available
            dmas = None
            if tech_score is not None or mom_score is not None:
                scores_data = {}
                if tech_score is not None:
                    scores_data["Technical Score"] = [float(tech_score)]
                if mom_score is not None:
                    scores_data["Momentum Score"] = [float(mom_score)]
                if tech_score is not None and mom_score is not None:
                    dmas = _safe_avg_two(tech_score, mom_score, rounding=1)
                    scores_data["Average (DMAS)"] = [dmas]
                st.table(pd.DataFrame(scores_data))
                # Provide an input for last week's average DMAS to be used in the gauge
                if selected_index == "S&P 500":
                    spx_last_week_input = st.number_input("Last week's average (DMAS)", min_value=0.0, max_value=100.0, value=st.session_state.get("spx_last_week_avg", 50.0), key="spx_last_week_avg_input")
                    st.session_state["spx_last_week_avg"] = spx_last_week_input
                elif selected_index == "CSI 300":
                    csi_last_week_input = st.number_input("Last week's average (DMAS)", min_value=0.0, max_value=100.0, value=st.session_state.get("csi_last_week_avg", 50.0), key="csi_last_week_avg_input")
                    st.session_state["csi_last_week_avg"] = csi_last_week_input
                elif selected_index == "Nikkei 225":
                    nikkei_last_week_input = st.number_input("Last week's average (DMAS)", min_value=0.0, max_value=100.0, value=st.session_state.get("nikkei_last_week_avg", 50.0), key="nikkei_last_week_avg_input")
                    st.session_state["nikkei_last_week_avg"] = nikkei_last_week_input
                elif selected_index == "TASI":
                    tasi_last_week_input = st.number_input("Last week's average (DMAS)", min_value=0.0, max_value=100.0, value=st.session_state.get("tasi_last_week_avg", 50.0), key="tasi_last_week_avg_input")
                    st.session_state["tasi_last_week_avg"] = tasi_last_week_input
                elif selected_index == "Sensex":
                    sensex_last_week_input = st.number_input("Last week's average (DMAS)", min_value=0.0, max_value=100.0, value=st.session_state.get("sensex_last_week_avg", 50.0), key="sensex_last_week_avg_input")
                    st.session_state["sensex_last_week_avg"] = sensex_last_week_input
                elif selected_index == "Dax":
                    dax_last_week_input = st.number_input("Last week's average (DMAS)", min_value=0.0, max_value=100.0, value=st.session_state.get("dax_last_week_avg", 50.0), key="dax_last_week_avg_input")
                    st.session_state["dax_last_week_avg"] = dax_last_week_input
                elif selected_index == "SMI":
                    smi_last_week_input = st.number_input("Last week's average (DMAS)", min_value=0.0, max_value=100.0, value=st.session_state.get("smi_last_week_avg", 50.0), key="smi_last_week_avg_input")
                    st.session_state["smi_last_week_avg"] = smi_last_week_input
                elif selected_index == "Ibov":
                    ibov_last_week_input = st.number_input("Last week's average (DMAS)", min_value=0.0, max_value=100.0, value=st.session_state.get("ibov_last_week_avg", 50.0), key="ibov_last_week_avg_input")
                    st.session_state["ibov_last_week_avg"] = ibov_last_week_input
                elif selected_index == "Mexbol":
                    mexbol_last_week_input = st.number_input("Last week's average (DMAS)", min_value=0.0, max_value=100.0, value=st.session_state.get("mexbol_last_week_avg", 50.0), key="mexbol_last_week_avg_input")
                    st.session_state["mexbol_last_week_avg"] = mexbol_last_week_input
            else:
                st.info("Neither Technical nor Momentum score could be calculated. Please check the Excel file.")
                st.info(
                    "Technical or momentum score not available in the uploaded Excel. "
                    "Please ensure sufficient price history and technical score data exist."
                )
            # -------------------------------------------------------------------
            # Show recent trading range (high/low) beneath the score table
            # -------------------------------------------------------------------
            try:
                # Compute trading range for the last 90 days based on implied volatility or realised volatility.
                current_price = df_full["Price"].iloc[-1] if not df_full.empty else None
                if current_price is not None and not np.isnan(current_price):
                    # Attempt to use implied volatility for the S&P 500 (VIX)
                    use_implied = False
                    vol_val = None
                    if selected_index == "S&P 500":
                        try:
                            df_vol = pd.read_excel(temp_path, sheet_name="data_prices")
                            df_vol = df_vol.drop(index=0)
                            df_vol = df_vol[df_vol[df_vol.columns[0]] != "DATES"]
                            df_vol["Date"] = pd.to_datetime(df_vol[df_vol.columns[0]], errors="coerce")
                            # Filter by "Data As Of" date if set
                            if "data_as_of" in st.session_state:
                                df_vol = df_vol[df_vol["Date"] <= pd.Timestamp(st.session_state["data_as_of"])]
                            if "VIX Index" in df_vol.columns:
                                df_vol["Price"] = pd.to_numeric(df_vol["VIX Index"], errors="coerce")
                                df_vol = df_vol.dropna(subset=["Date", "Price"]).sort_values("Date").reset_index(drop=True)[["Date", "Price"]]
                                pm = st.session_state.get("price_mode", "Last Price")
                                if adjust_prices_for_mode is not None:
                                    try:
                                        df_vol, _ = adjust_prices_for_mode(df_vol, pm)
                                    except Exception:
                                        pass
                                if not df_vol.empty:
                                    vol_val = float(df_vol["Price"].iloc[-1])
                                    use_implied = True
                        except Exception:
                            use_implied = False
                    elif selected_index == "SMI":
                        # Attempt to use implied volatility for SMI (VSMI1M)
                        try:
                            df_vol = pd.read_excel(temp_path, sheet_name="data_prices")
                            df_vol = df_vol.drop(index=0)
                            df_vol = df_vol[df_vol[df_vol.columns[0]] != "DATES"]
                            df_vol["Date"] = pd.to_datetime(df_vol[df_vol.columns[0]], errors="coerce")
                            # Filter by "Data As Of" date if set
                            if "data_as_of" in st.session_state:
                                df_vol = df_vol[df_vol["Date"] <= pd.Timestamp(st.session_state["data_as_of"])]
                            if "VSMI1M Index" in df_vol.columns:
                                df_vol["Price"] = pd.to_numeric(df_vol["VSMI1M Index"], errors="coerce")
                                df_vol = df_vol.dropna(subset=["Date", "Price"]).sort_values("Date").reset_index(drop=True)[["Date", "Price"]]
                                pm = st.session_state.get("price_mode", "Last Price")
                                if adjust_prices_for_mode is not None:
                                    try:
                                        df_vol, _ = adjust_prices_for_mode(df_vol, pm)
                                    except Exception:
                                        pass
                                if not df_vol.empty:
                                    vol_val = float(df_vol["Price"].iloc[-1])
                                    use_implied = True
                        except Exception:
                            use_implied = False
                    # Compute expected move from implied volatility if available
                    if use_implied and vol_val is not None:
                        expected_move = (current_price * (vol_val / 100.0)) / np.sqrt(52.0)
                        lower_bound = current_price - expected_move
                        upper_bound = current_price + expected_move
                        # Enforce minimum ±1 % band around current price
                        min_span = 0.02 * current_price
                        if (upper_bound - lower_bound) < min_span:
                            half = min_span / 2.0
                            lower_bound = current_price - half
                            upper_bound = current_price + half
                    else:
                        # Use realised-volatility based bounds based on the selected index
                        if selected_index == "S&P 500":
                            upper_bound, lower_bound = _compute_range_bounds_spx(df_full, lookback_days=90)
                        elif selected_index == "CSI 300":
                            upper_bound, lower_bound = _compute_range_bounds_csi(df_full, lookback_days=90)
                        elif selected_index == "Nikkei 225":
                            upper_bound, lower_bound = _compute_range_bounds_nikkei(df_full, lookback_days=90)
                        elif selected_index == "TASI":
                            upper_bound, lower_bound = _compute_range_bounds_tasi(df_full, lookback_days=90)
                        elif selected_index == "Sensex":
                            upper_bound, lower_bound = _compute_range_bounds_sensex(df_full, lookback_days=90)
                        elif selected_index == "Dax":
                            upper_bound, lower_bound = _compute_range_bounds_dax(df_full, lookback_days=90)
                        elif selected_index == "SMI":
                            upper_bound, lower_bound = _compute_range_bounds_smi(df_full, lookback_days=90)
                        elif selected_index == "Ibov":
                            upper_bound, lower_bound = _compute_range_bounds_ibov(df_full, lookback_days=90)
                        elif selected_index == "Mexbol":
                            upper_bound, lower_bound = _compute_range_bounds_mexbol(df_full, lookback_days=90)
                        else:
                            upper_bound, lower_bound = _compute_range_bounds_spx(df_full, lookback_days=90)
                    low_pct = (lower_bound - current_price) / current_price * 100.0
                    high_pct = (upper_bound - current_price) / current_price * 100.0
                    st.write(
                        f"Trading range (90d): Low {lower_bound:,.0f} ({low_pct:+.1f}%), "
                        f"High {upper_bound:,.0f} ({high_pct:+.1f}%)"
                    )
                else:
                    # No current price: use realised-volatility-based bounds
                    if selected_index == "S&P 500":
                        upper_bound, lower_bound = _compute_range_bounds_spx(df_full, lookback_days=90)
                    elif selected_index == "CSI 300":
                        upper_bound, lower_bound = _compute_range_bounds_csi(df_full, lookback_days=90)
                    elif selected_index == "Nikkei 225":
                        upper_bound, lower_bound = _compute_range_bounds_nikkei(df_full, lookback_days=90)
                    elif selected_index == "TASI":
                        upper_bound, lower_bound = _compute_range_bounds_tasi(df_full, lookback_days=90)
                    elif selected_index == "Sensex":
                        upper_bound, lower_bound = _compute_range_bounds_sensex(df_full, lookback_days=90)
                    elif selected_index == "Dax":
                        upper_bound, lower_bound = _compute_range_bounds_dax(df_full, lookback_days=90)
                    elif selected_index == "Ibov":
                        upper_bound, lower_bound = _compute_range_bounds_ibov(df_full, lookback_days=90)
                    elif selected_index == "Mexbol":
                        upper_bound, lower_bound = _compute_range_bounds_mexbol(df_full, lookback_days=90)
                    elif selected_index == "SMI":
                        upper_bound, lower_bound = _compute_range_bounds_smi(df_full, lookback_days=90)
                    elif selected_index == "Ibov":
                        upper_bound, lower_bound = _compute_range_bounds_ibov(df_full, lookback_days=90)
                    elif selected_index == "Mexbol":
                        upper_bound, lower_bound = _compute_range_bounds_mexbol(df_full, lookback_days=90)
                    else:
                        upper_bound, lower_bound = _compute_range_bounds_spx(df_full, lookback_days=90)
                    st.write(
                        f"Trading range (90d): Low {lower_bound:,.0f} – High {upper_bound:,.0f}"
                    )
            except Exception:
                pass

            # -------------------------------------------------------------------
            # Regression channel controls second
            # -------------------------------------------------------------------
            enable_channel = st.checkbox(
                "Enable regression channel",
                value=st.session_state.get(f"{ticker_key}_enable_channel", True),
                key=f"{ticker_key}_enable_channel",
            )

            anchor_ts = None
            if enable_channel:
                # When the regression channel is enabled, default the anchor to
                # the start of the selected analysis timeframe unless a
                # previous anchor has been stored in the session.  This
                # replaces the fixed 180‑day default with the user‑chosen
                # timeframe (e.g. 180 or 365 days).
                default_anchor = st.session_state.get(
                    f"{ticker_key}_anchor",
                    (max_date - pd.Timedelta(days=30))  # 1 month before last data,
                )
                # Ensure default is within valid range (convert to Timestamp for comparison)
                default_anchor = pd.Timestamp(default_anchor)
                default_anchor = max(pd.Timestamp(min_date), min(default_anchor, pd.Timestamp(max_date)))
                anchor_input = st.date_input(
                    "Select anchor date",
                    value=default_anchor,
                    min_value=min_date,
                    max_value=max_date,
                    key=f"{ticker_key}_anchor_date_input",
                )
                anchor_ts = pd.to_datetime(anchor_input)
                st.session_state[f"{ticker_key}_anchor"] = anchor_ts
            else:
                if f"{ticker_key}_anchor" in st.session_state:
                    st.session_state.pop(f"{ticker_key}_anchor")
                anchor_ts = None

            # -------------------------------------------------------------------
            # Assessment selection third
            # -------------------------------------------------------------------
            if tech_score is not None and mom_score is not None and dmas is not None:
                options = ASSESSMENT_OPTIONS
                def _default_index_from_dmas(val: float) -> int:
                    if val >= 70:
                        return options.index("Bullish")
                    elif val >= 55:
                        return options.index("Constructive")
                    elif val >= 45:
                        return options.index("Neutral")
                    elif val >= 30:
                        return options.index("Cautious")
                    else:
                        return options.index("Bearish")

                # Check if assessment was pre-set from transition sheet
                preset_assessment = st.session_state.get(f"{ticker_key}_assessment")
                if preset_assessment and preset_assessment in options:
                    default_idx = options.index(preset_assessment)
                else:
                    # Use automatic computation from DMAS
                    default_idx = _default_index_from_dmas(dmas)

                user_view = st.selectbox(
                    "Select your assessment",
                    options,
                    index=default_idx,
                    key=f"{ticker_key}_view_select",
                )
                st.session_state[f"{ticker_key}_selected_view"] = user_view
                st.caption(
                    "Your selection will override the automatically computed view in the presentation."
                )

            # -------------------------------------------------------------------
            # Subtitle input fourth
            # -------------------------------------------------------------------
            # Auto-generate subtitle button
            if SUBTITLE_GEN_AVAILABLE and tech_score is not None and mom_score is not None and dmas is not None:
                if st.button(f"🔄 Auto-generate {ticker_key.upper()} subtitle", key=f"{ticker_key}_autogen_subtitle"):
                    try:
                        # Get price data for structure analysis (if available)
                        try:
                            if excel_available:
                                df_price = pd.read_excel(temp_path, sheet_name="Price")
                                prices = df_price["Price"]
                            else:
                                prices = pd.Series([100] * 300)  # Fallback dummy data
                        except Exception:
                            prices = pd.Series([100] * 300)  # Fallback

                        # Get last week DMAS
                        dmas_prev = st.session_state.get(f"{ticker_key}_last_week_avg", dmas)

                        # Generate subtitle
                        result = generate_assessment_and_subtitle(
                            ticker_key=ticker_key,
                            asset_name=selected_index,
                            prices=prices,
                            dmas=dmas,
                            technical_score=tech_score,
                            momentum_score=mom_score,
                            dmas_prev_week=dmas_prev,
                            subtitle_generator=st.session_state.get('subtitle_generator')
                        )

                        # Update session state
                        st.session_state[f"{ticker_key}_subtitle"] = result["subtitle"]
                        st.success(f"Generated: {result['subtitle']}")

                    except Exception as e:
                        st.error(f"Could not auto-generate subtitle: {e}")

            subtitle_value = st.text_input(
                f"{ticker_key.upper()} subtitle" if selected_index == "S&P 500" else f"{ticker_key.upper()} subtitle",
                value=st.session_state.get(f"{ticker_key}_subtitle", ""),
                key=f"{ticker_key}_subtitle_input",
            )
            st.session_state[f"{ticker_key}_subtitle"] = subtitle_value

            # -------------------------------------------------------------------
            # Finally, build and show the interactive chart
            # -------------------------------------------------------------------
            if excel_available:
                pmode = st.session_state.get("price_mode", "Last Price")
                if selected_index == "S&P 500":
                    fig = make_spx_figure(temp_path, anchor_date=anchor_ts, price_mode=pmode)
                elif selected_index == "CSI 300":
                    fig = make_csi_figure(temp_path, anchor_date=anchor_ts, price_mode=pmode)
                elif selected_index == "Nikkei 225":
                    fig = make_nikkei_figure(temp_path, anchor_date=anchor_ts, price_mode=pmode)
                elif selected_index == "TASI":
                    fig = make_tasi_figure(temp_path, anchor_date=anchor_ts, price_mode=pmode)
                elif selected_index == "Sensex":
                    fig = make_sensex_figure(temp_path, anchor_date=anchor_ts, price_mode=pmode)
                elif selected_index == "Dax":
                    fig = make_dax_figure(temp_path, anchor_date=anchor_ts, price_mode=pmode)
                elif selected_index == "SMI":
                    fig = make_smi_figure(temp_path, anchor_date=anchor_ts, price_mode=pmode)
                elif selected_index == "Ibov":
                    fig = make_ibov_figure(temp_path, anchor_date=anchor_ts, price_mode=pmode)
                elif selected_index == "Mexbol":
                    fig = make_mexbol_figure(temp_path, anchor_date=anchor_ts, price_mode=pmode)
                else:
                    # default fallback: use SPX figure
                    fig = make_spx_figure(temp_path, anchor_date=anchor_ts, price_mode=pmode)
            else:
                df_ma = _add_moving_averages(df_full)
                fig = _build_fallback_figure(df_ma, anchor_date=anchor_ts)

            st.plotly_chart(fig, use_container_width=True)
            st.caption(
                "Use the controls above to enable and configure the regression channel. "
                "Green shading indicates an uptrend; red shading indicates a downtrend."
            )

            # Show score history below price chart
            show_score_history(selected_index)

    elif asset_class == "Commodity":
        # Delegate to the commodity technical analysis handler
        show_commodity_technical_analysis()
    elif asset_class == "Crypto":
        # Delegate to the crypto technical analysis handler
        show_crypto_technical_analysis()
    else:
        # Fallback for unsupported asset classes
        with st.expander(f"{asset_class} technical charts", expanded=False):
            st.info(f"{asset_class} technical analysis not implemented yet.")


def show_commodity_technical_analysis() -> None:
    """Render the technical analysis interface for commodity assets such as Gold.

    This function mirrors the equity technical analysis interface but is
    customised for commodity tickers.  Currently only Gold is supported.
    It handles data loading, score retrieval, DMAS computation, trading
    range estimation using an implied volatility index (XAUUSDV1M),
    regression channel controls, assessment selection, subtitle input and
    interactive chart rendering.  State is persisted in
    ``st.session_state`` to allow regeneration of the chart with the
    regression channel anchored at a user‑selected date.
    """
    # Identify whether an Excel file has been uploaded
    excel_available = "excel_file" in st.session_state

    # Commodity selection (Gold and Silver)
    # Include Gold, Silver, Platinum, Oil and Copper in the commodity options
    # Include Palladium in the list of supported commodities
    index_options = ["Gold", "Silver", "Platinum", "Palladium", "Oil", "Copper"]
    default_index = st.session_state.get("ta_commodity_index", "Gold")
    selected_index = st.sidebar.selectbox(
        "Select commodity for technical analysis",
        options=index_options,
        index=index_options.index(default_index) if default_index in index_options else 0,
        key="ta_commodity_index_select",
    )
    # Persist the selected commodity
    st.session_state["ta_commodity_index"] = selected_index

    # Determine ticker and keys based on selection
    if selected_index == "Gold":
        ticker = "GCA Comdty"
        ticker_key = "gold"
        chart_title = "Gold Technical Chart"
    elif selected_index == "Silver":
        ticker = "SIA Comdty"
        ticker_key = "silver"
        chart_title = "Silver Technical Chart"
    elif selected_index == "Platinum":
        ticker = "XPT Comdty"
        ticker_key = "platinum"
        chart_title = "Platinum Technical Chart"
    elif selected_index == "Palladium":
        ticker = "XPD Curncy"
        ticker_key = "palladium"
        chart_title = "Palladium Technical Chart"
    elif selected_index == "Oil":
        ticker = "CL1 Comdty"
        ticker_key = "oil"
        chart_title = "Oil Technical Chart"
    elif selected_index == "Copper":
        ticker = "LP1 Comdty"
        ticker_key = "copper"
        chart_title = "Copper Technical Chart"
    else:
        # Default back to Gold if an unknown commodity is selected
        ticker = "GCA Comdty"
        ticker_key = "gold"
        chart_title = f"{selected_index} Technical Chart"

    # Load price data (either from Excel or fallback synthetic)
    if excel_available:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".xlsx") as tmp:
            tmp.write(st.session_state["excel_file"].getbuffer())
            tmp.flush()
            temp_path = Path(tmp.name)
        # Read the data_prices sheet and tidy
        df_prices = pd.read_excel(temp_path, sheet_name="data_prices")
        df_prices = df_prices.drop(index=0)
        df_prices = df_prices[df_prices[df_prices.columns[0]] != "DATES"]
        df_prices["Date"] = pd.to_datetime(df_prices[df_prices.columns[0]], errors="coerce")
        # Filter by "Data As Of" date if set
        if "data_as_of" in st.session_state:
            df_prices = df_prices[df_prices["Date"] <= pd.Timestamp(st.session_state["data_as_of"])]
        df_prices["Price"] = pd.to_numeric(df_prices[ticker], errors="coerce")
        df_prices = df_prices.dropna(subset=["Date", "Price"]).sort_values("Date").reset_index(drop=True)
        # Adjust for price mode using the helper if available
        price_mode = st.session_state.get("price_mode", "Last Price")
        if adjust_prices_for_mode is not None and price_mode:
            try:
                df_prices, used_date = adjust_prices_for_mode(df_prices, price_mode)
            except Exception:
                used_date = None
        else:
            used_date = None
        df_full = df_prices.copy()
        st.session_state[f"{ticker_key}_used_date"] = used_date
    else:
        # Use a synthetic SPX series as a fallback when no Excel is provided
        df_prices = _create_synthetic_spx_series()
        df_full = df_prices.copy()
        used_date = None
    # Determine min and max dates for regression channel controls
    min_date = df_prices["Date"].min().date()
    max_date = df_prices["Date"].max().date()

    # Chart and controls
    with st.expander(chart_title, expanded=True):
        # Caption for date used
        used_date = st.session_state.get(f"{ticker_key}_used_date")
        price_mode = st.session_state.get("price_mode", "Last Price")
        if used_date is not None:
            if price_mode == "Last Close":
                st.caption(f"Prices as of {used_date.strftime('%d/%m/%Y')} close")
            else:
                st.caption(f"Prices as of {used_date.strftime('%d/%m/%Y')}")
        # -----------------------------------------------------------------
        # Technical and momentum scores
        # -----------------------------------------------------------------
        st.subheader("Technical and momentum scores")
        tech_score: Optional[float] = None
        mom_score: Optional[float] = None
        if excel_available and df_prices is not None:
            try:
                # Map display name to Bloomberg ticker
                ticker_map = {
                    "Gold": "GCA COMDTY",
                    "Silver": "SIA COMDTY",
                    "Platinum": "XPT COMDTY",
                    "Palladium": "XPD CURNCY",
                    "Oil": "CL1 COMDTY",
                    "Copper": "LP1 COMDTY",
                }

                bbg_ticker = ticker_map.get(selected_index)
                if bbg_ticker:
                    # Find matching column (case-insensitive)
                    matching_col = None
                    for col in df_prices.columns:
                        if isinstance(col, str) and col.upper() == bbg_ticker.upper():
                            matching_col = col
                            break

                    if matching_col and len(df_prices) >= 200:
                        # Extract prices
                        prices = df_prices[matching_col].dropna()
                        if len(prices) >= 200:
                            # Import and compute scores
                            from technical_score_wrapper import compute_dmas_scores
                            scores = compute_dmas_scores(prices, ticker=bbg_ticker, excel_path=str(temp_path))
                            tech_score = scores["technical_score"]
                            mom_score = scores["momentum_score"]

            except Exception as e:
                print(f"Error computing scores for {selected_index}: {e}")
                tech_score = None
                mom_score = None
        # Compute DMAS if scores are available
        dmas: Optional[float] = None
        if tech_score is not None or mom_score is not None:
            scores_data = {}
            if tech_score is not None:
                scores_data["Technical Score"] = [_fmt_score_cell(tech_score)]
            if mom_score is not None:
                scores_data["Momentum Score"] = [_fmt_score_cell(mom_score)]
            if tech_score is not None and mom_score is not None:
                dmas = _safe_avg_two(tech_score, mom_score, rounding=1)
                scores_data["Average (DMAS)"] = [_fmt_dmas_cell(dmas)]
            st.table(pd.DataFrame(scores_data))
            # Allow user to input last week's DMAS for the gauge. Use a key based on the commodity
            gauge_key = f"{ticker_key}_last_week_avg"
            gauge_input_key = f"{ticker_key}_last_week_avg_input"
            last_week_default = st.session_state.get(gauge_key, 50.0)
            last_week_input = st.number_input(
                "Last week's average (DMAS)",
                min_value=0.0,
                max_value=100.0,
                value=last_week_default,
                key=gauge_input_key,
            )
            st.session_state[gauge_key] = last_week_input
        else:
            st.info(
                "Technical or momentum score not available in the uploaded Excel. "
                "Please ensure sufficient price history and technical score data exist."
            )
        # -----------------------------------------------------------------
        # Trading range (90d) estimation
        # -----------------------------------------------------------------
        try:
            current_price = df_full["Price"].iloc[-1] if not df_full.empty else None
            if current_price is not None and not np.isnan(current_price):
                use_implied = False
                vol_val: Optional[float] = None
                # Attempt to use implied volatility via XAUUSDV1M BGN Curncy (Gold) or XAGUSDV1M BGN Curncy (Silver)
                vol_col_name = None
                if selected_index == "Gold":
                    vol_col_name = "XAUUSDV1M BGN Curncy"
                elif selected_index == "Silver":
                    vol_col_name = "XAGUSDV1M BGN Curncy"
                elif selected_index == "Platinum":
                    vol_col_name = "XPTUSDV1M BGN Curncy"
                elif selected_index == "Palladium":
                    vol_col_name = "XPDUSDV1M BGN Curncy"
                elif selected_index == "Oil":
                    # Oil implied volatility index column
                    vol_col_name = "WTI US 1M 50D VOL BVOL Equity"
                else:
                    # Copper implied volatility index column
                    vol_col_name = "LPR1 Index"
                if vol_col_name is not None:
                    try:
                        df_vol = pd.read_excel(temp_path, sheet_name="data_prices")
                        df_vol = df_vol.drop(index=0)
                        df_vol = df_vol[df_vol[df_vol.columns[0]] != "DATES"]
                        df_vol["Date"] = pd.to_datetime(df_vol[df_vol.columns[0]], errors="coerce")
                        # Filter by "Data As Of" date if set
                        if "data_as_of" in st.session_state:
                            df_vol = df_vol[df_vol["Date"] <= pd.Timestamp(st.session_state["data_as_of"])]
                        if vol_col_name in df_vol.columns:
                            df_vol["Price"] = pd.to_numeric(df_vol[vol_col_name], errors="coerce")
                            df_vol = df_vol.dropna(subset=["Date", "Price"]).sort_values("Date").reset_index(drop=True)[["Date", "Price"]]
                            pm = st.session_state.get("price_mode", "Last Price")
                            if adjust_prices_for_mode is not None:
                                try:
                                    df_vol, _ = adjust_prices_for_mode(df_vol, pm)
                                except Exception:
                                    pass
                            if not df_vol.empty:
                                vol_val = float(df_vol["Price"].iloc[-1])
                                use_implied = True
                    except Exception:
                        use_implied = False
                # If implied vol available, compute expected move
                if use_implied and vol_val is not None:
                    expected_move = (current_price * (vol_val / 100.0)) / np.sqrt(52.0)
                    lower_bound = current_price - expected_move
                    upper_bound = current_price + expected_move
                    min_span = 0.02 * current_price
                    if (upper_bound - lower_bound) < min_span:
                        half = min_span / 2.0
                        lower_bound = current_price - half
                        upper_bound = current_price + half
                else:
                    # Fallback to realised volatility when implied vol is unavailable
                    if selected_index == "Gold":
                        upper_bound, lower_bound = _compute_range_bounds_gold(df_full, lookback_days=90)
                    elif selected_index == "Silver":
                        upper_bound, lower_bound = _compute_range_bounds_silver(df_full, lookback_days=90)
                    elif selected_index == "Platinum":
                        upper_bound, lower_bound = _compute_range_bounds_platinum(df_full, lookback_days=90)
                    elif selected_index == "Palladium":
                        upper_bound, lower_bound = _compute_range_bounds_palladium(df_full, lookback_days=90)
                    elif selected_index == "Oil":
                        upper_bound, lower_bound = _compute_range_bounds_oil(df_full, lookback_days=90)
                    else:
                        # Copper (or any other commodity default)
                        upper_bound, lower_bound = _compute_range_bounds_copper(df_full, lookback_days=90)
                low_pct = (lower_bound - current_price) / current_price * 100.0
                high_pct = (upper_bound - current_price) / current_price * 100.0
                st.write(
                    f"Trading range (90d): Low {lower_bound:,.0f} ({low_pct:+.1f}%), "
                    f"High {upper_bound:,.0f} ({high_pct:+.1f}%)"
                )
            else:
                if selected_index == "Gold":
                    upper_bound, lower_bound = _compute_range_bounds_gold(df_full, lookback_days=90)
                elif selected_index == "Silver":
                    upper_bound, lower_bound = _compute_range_bounds_silver(df_full, lookback_days=90)
                elif selected_index == "Platinum":
                    upper_bound, lower_bound = _compute_range_bounds_platinum(df_full, lookback_days=90)
                elif selected_index == "Palladium":
                    upper_bound, lower_bound = _compute_range_bounds_palladium(df_full, lookback_days=90)
                elif selected_index == "Oil":
                    upper_bound, lower_bound = _compute_range_bounds_oil(df_full, lookback_days=90)
                else:
                    # Copper (or other commodity)
                    upper_bound, lower_bound = _compute_range_bounds_copper(df_full, lookback_days=90)
                st.write(
                    f"Trading range (90d): Low {lower_bound:,.0f} – High {upper_bound:,.0f}"
                )
        except Exception:
            pass
        # -----------------------------------------------------------------
        # Regression channel controls
        # -----------------------------------------------------------------
        enable_channel = st.checkbox(
            "Enable regression channel",
            value=st.session_state.get(f"{ticker_key}_enable_channel", True),
            key=f"{ticker_key}_enable_channel",
        )
        anchor_ts: Optional[pd.Timestamp] = None
        if enable_channel:
            # Default the anchor to the beginning of the selected
            # timeframe when no previous anchor is stored.  Uses
            # ``ta_timeframe_days`` instead of a fixed window.  When
            # ``ta_timeframe_days`` is not present (e.g. outside Streamlit),
            # a default of 90 days (three months) is used.
            default_anchor = st.session_state.get(
                f"{ticker_key}_anchor",
                (max_date - pd.Timedelta(days=30))  # 1 month before last data,
            )
            # Ensure default is within valid range (convert to Timestamp for comparison)
            default_anchor = pd.Timestamp(default_anchor)
            default_anchor = max(pd.Timestamp(min_date), min(default_anchor, pd.Timestamp(max_date)))
            anchor_input = st.date_input(
                "Select anchor date",
                value=default_anchor,
                min_value=min_date,
                max_value=max_date,
                key=f"{ticker_key}_anchor_date_input",
            )
            anchor_ts = pd.to_datetime(anchor_input)
            st.session_state[f"{ticker_key}_anchor"] = anchor_ts
        else:
            if f"{ticker_key}_anchor" in st.session_state:
                st.session_state.pop(f"{ticker_key}_anchor")
            anchor_ts = None
        # -----------------------------------------------------------------
        # Assessment selection
        # -----------------------------------------------------------------
        if tech_score is not None and mom_score is not None and dmas is not None:
            options = [
                "Strongly Bearish",
                "Bearish",
                "Slightly Bearish",
                "Neutral",
                "Slightly Bullish",
                "Bullish",
                "Strongly Bullish",
            ]
            def _default_index_from_dmas(val: float) -> int:
                if val >= 70:
                    return options.index("Bullish")
                elif val >= 55:
                    return options.index("Constructive")
                elif val >= 45:
                    return options.index("Neutral")
                elif val >= 30:
                    return options.index("Cautious")
                else:
                    return options.index("Bearish")
            preset_assessment = st.session_state.get(f"{ticker_key}_assessment")
            if preset_assessment and preset_assessment in options:
                default_idx = options.index(preset_assessment)
            else:
                # Use automatic computation from DMAS
                default_idx = _default_index_from_dmas(dmas)

            user_view = st.selectbox(
                "Select your assessment",
                options,
                index=default_idx,
                key=f"{ticker_key}_view_select",
            )
            st.session_state[f"{ticker_key}_selected_view"] = user_view
            st.caption(
                "Your selection will override the automatically computed view in the presentation."
            )
        # -----------------------------------------------------------------
        # Subtitle input
        # -----------------------------------------------------------------
        # Auto-generate subtitle button
        if SUBTITLE_GEN_AVAILABLE and tech_score is not None and mom_score is not None and dmas is not None:
            if st.button(f"🔄 Auto-generate {ticker_key.upper()} subtitle", key=f"{ticker_key}_autogen_subtitle"):
                try:
                    # Get price data for structure analysis (if available)
                    try:
                        if excel_available:
                            df_price = pd.read_excel(temp_path, sheet_name="Price")
                            prices = df_price["Price"]
                        else:
                            prices = pd.Series([100] * 300)  # Fallback dummy data
                    except Exception:
                        prices = pd.Series([100] * 300)  # Fallback

                    # Get last week DMAS
                    dmas_prev = st.session_state.get(f"{ticker_key}_last_week_avg", dmas)

                    # Generate subtitle
                    result = generate_assessment_and_subtitle(
                        ticker_key=ticker_key,
                        asset_name=selected_index,
                        prices=prices,
                        dmas=dmas,
                        technical_score=tech_score,
                        momentum_score=mom_score,
                        dmas_prev_week=dmas_prev,
                        subtitle_generator=st.session_state.get('subtitle_generator')
                    )

                    # Update session state
                    st.session_state[f"{ticker_key}_subtitle"] = result["subtitle"]
                    st.success(f"Generated: {result['subtitle']}")

                except Exception as e:
                    st.error(f"Could not auto-generate subtitle: {e}")

        subtitle_value = st.text_input(
            f"{ticker_key.upper()} subtitle",
            value=st.session_state.get(f"{ticker_key}_subtitle", ""),
            key=f"{ticker_key}_subtitle_input",
        )
        st.session_state[f"{ticker_key}_subtitle"] = subtitle_value
        # -----------------------------------------------------------------
        # Interactive chart
        # -----------------------------------------------------------------
        if excel_available:
            pmode = st.session_state.get("price_mode", "Last Price")
            if selected_index == "Gold":
                fig = make_gold_figure(temp_path, anchor_date=anchor_ts, price_mode=pmode)
            elif selected_index == "Silver":
                fig = make_silver_figure(temp_path, anchor_date=anchor_ts, price_mode=pmode)
            elif selected_index == "Platinum":
                fig = make_platinum_figure(temp_path, anchor_date=anchor_ts, price_mode=pmode)
            elif selected_index == "Palladium":
                fig = make_palladium_figure(temp_path, anchor_date=anchor_ts, price_mode=pmode)
            elif selected_index == "Oil":
                fig = make_oil_figure(temp_path, anchor_date=anchor_ts, price_mode=pmode)
            elif selected_index == "Copper":
                fig = make_copper_figure(temp_path, anchor_date=anchor_ts, price_mode=pmode)
            else:
                # Fallback: show an empty figure if unknown commodity
                fig = go.Figure()
        else:
            # Fallback: compute simple MA and regression channel on synthetic data
            from technical_analysis.equity.spx import _add_moving_averages, _build_fallback_figure  # type: ignore
            df_ma = _add_moving_averages(df_full)
            fig = _build_fallback_figure(df_ma, anchor_date=anchor_ts)
        st.plotly_chart(fig, use_container_width=True)
        st.caption(
            "Use the controls above to enable and configure the regression channel. "
            "Green shading indicates an uptrend; red shading indicates a downtrend."
        )

        # Show score history below price chart
        show_score_history(selected_index)


def show_crypto_technical_analysis() -> None:
    """Render the technical analysis interface for crypto assets such as Bitcoin.

    This function closely mirrors the commodity technical analysis interface
    but is tailored for crypto tickers.  At present only Bitcoin is
    supported.  It handles data loading from the uploaded Excel, score
    retrieval, DMAS computation, trading range estimation using the
    Bitcoin volatility index (BVXS Index) with a realised volatility
    fallback, regression channel controls, assessment selection, subtitle
    input and interactive chart rendering.  State is persisted in
    ``st.session_state`` so that the regression channel can be anchored
    at a user‑selected date across reruns.
    """
    # Identify whether an Excel file has been uploaded
    excel_available = "excel_file" in st.session_state

    # Allow selection of supported crypto assets
    index_options = ["Bitcoin", "Ethereum", "Ripple", "Solana", "Binance"]
    default_index = st.session_state.get("ta_crypto_index", "Bitcoin") if st.session_state.get("ta_crypto_index") in index_options else "Bitcoin"
    selected_index = st.sidebar.selectbox(
        "Select crypto for technical analysis",
        options=index_options,
        index=index_options.index(default_index) if default_index in index_options else 0,
        key="ta_crypto_index_select",
    )
    # Persist selection
    st.session_state["ta_crypto_index"] = selected_index

    # Determine ticker, key, chart title, volatility column and helper functions based on selected crypto
    if selected_index == "Bitcoin":
        ticker = "XBTUSD Curncy"
        ticker_key = "bitcoin"
        chart_title = "Bitcoin Technical Chart"
        vol_col_name = "BVXS Index"
        get_tech_score = _get_bitcoin_technical_score
        get_mom_score = _get_bitcoin_momentum_score
        compute_range_fallback = _compute_range_bounds_bitcoin
        make_figure_func = make_bitcoin_figure
    elif selected_index == "Ethereum":
        ticker = "XETUSD Curncy"
        ticker_key = "ethereum"
        chart_title = "Ethereum Technical Chart"
        # Placeholder for implied volatility column; not yet available
        vol_col_name = "XETUSDV1M BGN Curncy"
        get_tech_score = _get_ethereum_technical_score
        get_mom_score = _get_ethereum_momentum_score
        compute_range_fallback = _compute_range_bounds_ethereum
        make_figure_func = make_ethereum_figure
    elif selected_index == "Ripple":
        ticker = "XRPUSD Curncy"
        ticker_key = "ripple"
        chart_title = "Ripple Technical Chart"
        # Placeholder for implied volatility column; not yet available
        vol_col_name = "XRPUSDV1M BGN Curncy"
        get_tech_score = _get_ripple_technical_score
        get_mom_score = _get_ripple_momentum_score
        compute_range_fallback = _compute_range_bounds_ripple
        make_figure_func = make_ripple_figure
    elif selected_index == "Solana":
        ticker = "XSOUSD Curncy"
        ticker_key = "solana"
        chart_title = "Solana Technical Chart"
        # Placeholder for implied volatility column; not yet available
        vol_col_name = "XSOUSDV1M BGN Curncy"
        get_tech_score = _get_solana_technical_score
        get_mom_score = _get_solana_momentum_score
        compute_range_fallback = _compute_range_bounds_solana
        make_figure_func = make_solana_figure
    elif selected_index == "Binance":
        ticker = "XBIUSD Curncy"
        ticker_key = "binance"
        chart_title = "Binance Technical Chart"
        # Placeholder for implied volatility column; not yet available
        vol_col_name = "XBIUSDV1M BGN Curncy"
        get_tech_score = _get_binance_technical_score
        get_mom_score = _get_binance_momentum_score
        compute_range_fallback = _compute_range_bounds_binance
        make_figure_func = make_binance_figure
    else:
        # Default to Bitcoin if unknown selection
        ticker = "XBTUSD Curncy"
        ticker_key = "bitcoin"
        chart_title = f"{selected_index} Technical Chart"
        vol_col_name = "BVXS Index"
        get_tech_score = _get_bitcoin_technical_score
        get_mom_score = _get_bitcoin_momentum_score
        compute_range_fallback = _compute_range_bounds_bitcoin
        make_figure_func = make_bitcoin_figure

    # Load price data
    if excel_available:
        # Save uploaded Excel to a temporary path
        with tempfile.NamedTemporaryFile(delete=False, suffix=".xlsx") as tmp:
            tmp.write(st.session_state["excel_file"].getbuffer())
            tmp.flush()
            temp_path = Path(tmp.name)
        # Read data_prices and clean
        df_prices = pd.read_excel(temp_path, sheet_name="data_prices")
        df_prices = df_prices.drop(index=0)
        df_prices = df_prices[df_prices[df_prices.columns[0]] != "DATES"]
        df_prices["Date"] = pd.to_datetime(df_prices[df_prices.columns[0]], errors="coerce")
        # Filter by "Data As Of" date if set
        if "data_as_of" in st.session_state:
            df_prices = df_prices[df_prices["Date"] <= pd.Timestamp(st.session_state["data_as_of"])]
        df_prices["Price"] = pd.to_numeric(df_prices[ticker], errors="coerce")
        df_prices = df_prices.dropna(subset=["Date", "Price"]).sort_values("Date").reset_index(drop=True)
        # Adjust prices for selected price mode
        price_mode = st.session_state.get("price_mode", "Last Price")
        used_date = None
        if adjust_prices_for_mode is not None and price_mode:
            try:
                df_prices, used_date = adjust_prices_for_mode(df_prices, price_mode)
            except Exception:
                used_date = None
        st.session_state[f"{ticker_key}_used_date"] = used_date
        df_full = df_prices.copy()
    else:
        # Fallback: synthetic SPX series (not ideal but ensures a chart)
        df_prices = _create_synthetic_spx_series()
        df_full = df_prices.copy()
        used_date = None

    # Determine date range for channel controls
    min_date = df_prices["Date"].min().date()
    max_date = df_prices["Date"].max().date()

    # Chart and controls
    with st.expander(chart_title, expanded=True):
        # Caption with used date
        used_date = st.session_state.get(f"{ticker_key}_used_date")
        price_mode = st.session_state.get("price_mode", "Last Price")
        if used_date is not None:
            if price_mode == "Last Close":
                st.caption(f"Prices as of {used_date.strftime('%d/%m/%Y')} close")
            else:
                st.caption(f"Prices as of {used_date.strftime('%d/%m/%Y')}")
        # -----------------------------------------------------------------
        # Technical and momentum scores
        # -----------------------------------------------------------------
        st.subheader("Technical and momentum scores")
        tech_score: Optional[float] = None
        mom_score: Optional[float] = None
        if excel_available and df_prices is not None:
            try:
                # Map display name to Bloomberg ticker
                ticker_map = {
                    "Bitcoin": "XBTUSD CURNCY",
                    "Ethereum": "XETUSD CURNCY",
                    "Ripple": "XRPUSD CURNCY",
                    "Solana": "XSOUSD CURNCY",
                    "Binance": "XBIUSD CURNCY",
                }

                bbg_ticker = ticker_map.get(selected_index)
                if bbg_ticker:
                    # Find matching column (case-insensitive)
                    matching_col = None
                    for col in df_prices.columns:
                        if isinstance(col, str) and col.upper() == bbg_ticker.upper():
                            matching_col = col
                            break

                    if matching_col and len(df_prices) >= 200:
                        # Extract prices
                        prices = df_prices[matching_col].dropna()
                        if len(prices) >= 200:
                            # Import and compute scores
                            from technical_score_wrapper import compute_dmas_scores
                            scores = compute_dmas_scores(prices, ticker=bbg_ticker, excel_path=str(temp_path))
                            tech_score = scores["technical_score"]
                            mom_score = scores["momentum_score"]

            except Exception as e:
                print(f"Error computing scores for {selected_index}: {e}")
                tech_score = None
                mom_score = None
        # Compute DMAS if available
        dmas: Optional[float] = None
        if tech_score is not None or mom_score is not None:
            scores_data = {}
            if tech_score is not None:
                scores_data["Technical Score"] = [_fmt_score_cell(tech_score)]
            if mom_score is not None:
                scores_data["Momentum Score"] = [_fmt_score_cell(mom_score)]
            if tech_score is not None and mom_score is not None:
                dmas = _safe_avg_two(tech_score, mom_score, rounding=1)
                scores_data["Average (DMAS)"] = [_fmt_dmas_cell(dmas)]
            st.table(pd.DataFrame(scores_data))
            # Last week's DMAS input
            gauge_key = f"{ticker_key}_last_week_avg"
            gauge_input_key = f"{ticker_key}_last_week_avg_input"
            last_week_default = st.session_state.get(gauge_key, 50.0)
            last_week_input = st.number_input(
                "Last week's average (DMAS)",
                min_value=0.0,
                max_value=100.0,
                value=last_week_default,
                key=gauge_input_key,
            )
            st.session_state[gauge_key] = last_week_input
        else:
            st.info(
                "Technical or momentum score not available in the uploaded Excel. "
                "Please ensure sufficient price history and technical score data exist."
            )
        # -----------------------------------------------------------------
        # Trading range estimation (90d)
        # -----------------------------------------------------------------
        try:
            current_price = df_full["Price"].iloc[-1] if not df_full.empty else None
            if current_price is not None and not np.isnan(current_price):
                use_implied = False
                vol_val: Optional[float] = None
                # Already determined vol_col_name based on selected crypto
                # Attempt to load implied volatility from the Excel file
                if excel_available:
                    try:
                        df_vol = pd.read_excel(temp_path, sheet_name="data_prices")
                        df_vol = df_vol.drop(index=0)
                        df_vol = df_vol[df_vol[df_vol.columns[0]] != "DATES"]
                        df_vol["Date"] = pd.to_datetime(df_vol[df_vol.columns[0]], errors="coerce")
                        # Filter by "Data As Of" date if set
                        if "data_as_of" in st.session_state:
                            df_vol = df_vol[df_vol["Date"] <= pd.Timestamp(st.session_state["data_as_of"])]
                        if vol_col_name in df_vol.columns:
                            df_vol["Price"] = pd.to_numeric(df_vol[vol_col_name], errors="coerce")
                            df_vol = df_vol.dropna(subset=["Date", "Price"]).sort_values("Date").reset_index(drop=True)[[
                                "Date",
                                "Price",
                            ]]
                            pm = st.session_state.get("price_mode", "Last Price")
                            if adjust_prices_for_mode is not None:
                                try:
                                    df_vol, _ = adjust_prices_for_mode(df_vol, pm)
                                except Exception:
                                    pass
                            if not df_vol.empty:
                                vol_val = float(df_vol["Price"].iloc[-1])
                                use_implied = True
                    except Exception:
                        use_implied = False
                if use_implied and vol_val is not None:
                    expected_move = (current_price * (vol_val / 100.0)) / np.sqrt(52.0)
                    lower_bound = current_price - expected_move
                    upper_bound = current_price + expected_move
                    min_span = 0.02 * current_price
                    if (upper_bound - lower_bound) < min_span:
                        half = min_span / 2.0
                        lower_bound = current_price - half
                        upper_bound = current_price + half
                else:
                    # Fallback to realised volatility
                    # Use crypto-specific realised volatility fallback
                    upper_bound, lower_bound = compute_range_fallback(df_full, lookback_days=90)
                low_pct = (lower_bound - current_price) / current_price * 100.0
                high_pct = (upper_bound - current_price) / current_price * 100.0
                st.write(
                    f"Trading range (90d): Low {lower_bound:,.0f} ({low_pct:+.1f}%), "
                    f"High {upper_bound:,.0f} ({high_pct:+.1f}%)"
                )
            else:
                # Use crypto-specific realised volatility fallback
                upper_bound, lower_bound = compute_range_fallback(df_full, lookback_days=90)
                st.write(
                    f"Trading range (90d): Low {lower_bound:,.0f} – High {upper_bound:,.0f}"
                )
        except Exception:
            pass
        # -----------------------------------------------------------------
        # Regression channel controls
        # -----------------------------------------------------------------
        enable_channel = st.checkbox(
            "Enable regression channel",
            value=st.session_state.get(f"{ticker_key}_enable_channel", True),
            key=f"{ticker_key}_enable_channel",
        )
        anchor_ts: Optional[pd.Timestamp] = None
        if enable_channel:
            # Default anchor uses the selected timeframe rather than a
            # fixed window when none is stored in session.  When
            # ``ta_timeframe_days`` is not available a default of 90 days
            # (three months) is used.
            default_anchor = st.session_state.get(
                f"{ticker_key}_anchor",
                (max_date - pd.Timedelta(days=30))  # 1 month before last data,
            )
            # Ensure default is within valid range (convert to Timestamp for comparison)
            default_anchor = pd.Timestamp(default_anchor)
            default_anchor = max(pd.Timestamp(min_date), min(default_anchor, pd.Timestamp(max_date)))
            anchor_input = st.date_input(
                "Select anchor date",
                value=default_anchor,
                min_value=min_date,
                max_value=max_date,
                key=f"{ticker_key}_anchor_date_input",
            )
            anchor_ts = pd.to_datetime(anchor_input)
            st.session_state[f"{ticker_key}_anchor"] = anchor_ts
        else:
            if f"{ticker_key}_anchor" in st.session_state:
                st.session_state.pop(f"{ticker_key}_anchor")
            anchor_ts = None
        # -----------------------------------------------------------------
        # Assessment selection
        # -----------------------------------------------------------------
        if tech_score is not None and mom_score is not None and dmas is not None:
            options = [
                "Strongly Bearish",
                "Bearish",
                "Slightly Bearish",
                "Neutral",
                "Slightly Bullish",
                "Bullish",
                "Strongly Bullish",
            ]
            def _default_index_from_dmas(val: float) -> int:
                if val >= 70:
                    return options.index("Bullish")
                elif val >= 55:
                    return options.index("Constructive")
                elif val >= 45:
                    return options.index("Neutral")
                elif val >= 30:
                    return options.index("Cautious")
                else:
                    return options.index("Bearish")
            preset_assessment = st.session_state.get(f"{ticker_key}_assessment")
            if preset_assessment and preset_assessment in options:
                default_idx = options.index(preset_assessment)
            else:
                # Use automatic computation from DMAS
                default_idx = _default_index_from_dmas(dmas)

            user_view = st.selectbox(
                "Select your assessment",
                options,
                index=default_idx,
                key=f"{ticker_key}_view_select",
            )
            st.session_state[f"{ticker_key}_selected_view"] = user_view
            st.caption(
                "Your selection will override the automatically computed view in the presentation."
            )
        # -----------------------------------------------------------------
        # Subtitle input
        # -----------------------------------------------------------------
        # Auto-generate subtitle button
        if SUBTITLE_GEN_AVAILABLE and tech_score is not None and mom_score is not None and dmas is not None:
            if st.button(f"🔄 Auto-generate {ticker_key.upper()} subtitle", key=f"{ticker_key}_autogen_subtitle"):
                try:
                    # Get price data for structure analysis (if available)
                    try:
                        if excel_available:
                            df_price = pd.read_excel(temp_path, sheet_name="Price")
                            prices = df_price["Price"]
                        else:
                            prices = pd.Series([100] * 300)  # Fallback dummy data
                    except Exception:
                        prices = pd.Series([100] * 300)  # Fallback

                    # Get last week DMAS
                    dmas_prev = st.session_state.get(f"{ticker_key}_last_week_avg", dmas)

                    # Generate subtitle
                    result = generate_assessment_and_subtitle(
                        ticker_key=ticker_key,
                        asset_name=selected_index,
                        prices=prices,
                        dmas=dmas,
                        technical_score=tech_score,
                        momentum_score=mom_score,
                        dmas_prev_week=dmas_prev,
                        subtitle_generator=st.session_state.get('subtitle_generator')
                    )

                    # Update session state
                    st.session_state[f"{ticker_key}_subtitle"] = result["subtitle"]
                    st.success(f"Generated: {result['subtitle']}")

                except Exception as e:
                    st.error(f"Could not auto-generate subtitle: {e}")

        subtitle_value = st.text_input(
            f"{ticker_key.upper()} subtitle",
            value=st.session_state.get(f"{ticker_key}_subtitle", ""),
            key=f"{ticker_key}_subtitle_input",
        )
        st.session_state[f"{ticker_key}_subtitle"] = subtitle_value
        # -----------------------------------------------------------------
        # Interactive chart
        # -----------------------------------------------------------------
        if excel_available:
            pmode = st.session_state.get("price_mode", "Last Price")
            fig = make_figure_func(temp_path, anchor_date=anchor_ts, price_mode=pmode)
        else:
            # Build fallback figure on synthetic data
            from technical_analysis.equity.spx import _add_moving_averages, _build_fallback_figure  # type: ignore
            df_ma = _add_moving_averages(df_full)
            fig = _build_fallback_figure(df_ma, anchor_date=anchor_ts)
        st.plotly_chart(fig, use_container_width=True)
        st.caption(
            "Use the controls above to enable and configure the regression channel. "
            "Green shading indicates an uptrend; red shading indicates a downtrend."
        )

        # Show score history below price chart
        show_score_history(selected_index)


def show_market_breadth_page() -> None:
    st.header("Market Breadth")

    if "excel_file" not in st.session_state:
        st.error("Upload an Excel file first (Upload page).")
        return

    import tempfile, pathlib
    with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as tmp:
        tmp.write(st.session_state["excel_file"].getbuffer())
        xl_path = pathlib.Path(tmp.name)

    df = _load_breadth_page_data(xl_path)
    if df.empty:
        st.warning("No breadth data found in **bql_formula** (columns AB–AE).")
        return

    st.dataframe(
        _style_breadth_page(df),
        use_container_width=True,
        height=min(600, 50 + 25 * len(df)),
    )

    with st.expander("Debug – first parsed rows"):
        st.write(_debug_breadth_rows(xl_path))


def show_generate_presentation_page():
    """Generate a customised PowerPoint presentation based on user selections."""
    st.sidebar.header("Generate Presentation")
    if "excel_file" not in st.session_state or "pptx_file" not in st.session_state:
        st.sidebar.error(
            "Please upload both an Excel file and a PowerPoint template in the Upload page."
        )
        st.stop()

    st.sidebar.write("### Summary of selections")
    st.sidebar.write("Equities:", st.session_state.get("selected_eq_names", []))
    st.sidebar.write("Commodities:", st.session_state.get("selected_co_names", []))
    st.sidebar.write("Cryptos:", st.session_state.get("selected_cr_names", []))
    # Display FX pairs being analysed (fixed list) for user awareness
    st.sidebar.write(
        "FX:",
        [
            "DXY",
            "EUR/USD",
            "EUR/CHF",
            "EUR/GBP",
            "EUR/JPY",
            "EUR/AUD",
            "EUR/CAD",
            "EUR/BRL",
            "EUR/RUB",
            "EUR/ZAR",
            "EUR/MXN",
        ],
    )
    # Display rates tickers being analysed (fixed list)
    st.sidebar.write(
        "Rates:",
        [
            "US - 2Y",
            "US - 10Y",
            "US - 30Y",
            "EUR - 2Y",
            "EUR - 10Y",
            "EUR - 30Y",
            "CN - 2Y",
            "CN - 10Y",
            "CN - 30Y",
            "JP - 2Y",
            "JP - 10Y",
            "JP - 30Y",
        ],
    )
    # Display credit indices being analysed (fixed list) for user awareness
    st.sidebar.write(
        "Credit:",
        [
            "USD - IG",
            "USD - HY",
            "EUR - IG",
            "EUR - HY",
            "Asia (ex JP) - IG",
            "Asia - HY",
            "EM - IG",
            "EM - HY",
        ],
    )

    if st.sidebar.button("Generate updated PPTX", key="gen_ppt_button"):
        import time

        # Initialize progress tracking
        total_steps = 32  # Total number of major operations
        current_step = 0
        start_time = time.time()

        # Create progress bar and status text placeholders
        progress_bar = st.progress(0)
        status_text = st.empty()
        time_text = st.empty()

        def update_progress(step_name):
            """Update progress bar and estimated time remaining."""
            nonlocal current_step
            current_step += 1
            progress = min(current_step / total_steps, 1.0)  # Cap at 100%
            progress_bar.progress(progress)

            # Calculate elapsed and estimated remaining time
            elapsed = time.time() - start_time
            if current_step > 0:
                avg_time_per_step = elapsed / current_step
                remaining_steps = total_steps - current_step
                estimated_remaining = avg_time_per_step * remaining_steps

                # Format time display
                if estimated_remaining < 60:
                    time_str = f"{int(estimated_remaining)}s remaining"
                else:
                    mins = int(estimated_remaining // 60)
                    secs = int(estimated_remaining % 60)
                    time_str = f"{mins}m {secs}s remaining"

                time_text.text(f"⏱️ {time_str} (Step {current_step}/{total_steps})")

            status_text.text(f"🔄 {step_name}")

        # Write the uploaded PPTX to a temporary file so that python-pptx
        # can read it reliably.  Also write the uploaded Excel file to a
        # temporary XLSX path so that multiple reads do not exhaust the
        # underlying file-like object.  The Excel path is reused for
        # inserting charts and scores throughout the presentation.
        update_progress("Loading PowerPoint template...")
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp_input:
            tmp_input.write(st.session_state["pptx_file"].getbuffer())
            tmp_input.flush()
            prs = Presentation(tmp_input.name)

        # Persist the Excel to a temporary path to avoid file pointer
        # exhaustion when pandas reads multiple sheets.  Without this,
        # repeated reads from the UploadedFile can yield empty DataFrames.
        update_progress("Loading Excel data...")
        with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as tmp_xls:
            tmp_xls.write(st.session_state["excel_file"].getbuffer())
            tmp_xls.flush()
            excel_path_for_ppt = Path(tmp_xls.name)

        # Clear Excel cache to ensure clean state for this generation run
        # This is part of Phase 1 optimization - Excel sheets will be cached
        # during generation to avoid redundant reads across 20+ instruments
        from technical_analysis.common_helpers import clear_excel_cache
        clear_excel_cache()

        # ----------------------------------------------------------------------
        # Automatically update the date on the first slide and prepare filename
        #
        # Compute the current date in Europe/Zurich, format it for the
        # DataIC textbox and build a stamp for the output filename.  A
        # helper function is defined to locate the "DataIC" shape on
        # the first slide and replace its text while preserving
        # existing styling (size, colour, bold/italic).  See user
        # requirement: the first slide should show today's date in the
        # form "September 10, 2025" automatically when generating
        # presentations.  The filename will use the same date in
        # DDMMYYYY format.

        import pandas as pd  # local import to avoid polluting module scope

        def _set_text_in_named_textbox(prs_obj, shape_name: str, text: str):
            """Replace text in a named textbox while keeping the first run's style."""
            for slide in prs_obj.slides:
                for shape in slide.shapes:
                    if getattr(shape, "name", "") == shape_name and getattr(shape, "has_text_frame", False):
                        p = shape.text_frame.paragraphs[0]
                        if p.runs:
                            r0 = p.runs[0]
                            size = r0.font.size
                            color = r0.font.color
                            rgb = getattr(color, "rgb", None)
                            theme_color = getattr(color, "theme_color", None)
                            brightness = getattr(color, "brightness", None)
                            bold = r0.font.bold
                            italic = r0.font.italic
                            shape.text_frame.clear()
                            pr = shape.text_frame.paragraphs[0]
                            new_run = pr.add_run()
                            new_run.text = text
                            if size:
                                new_run.font.size = size
                            try:
                                if rgb:
                                    new_run.font.color.rgb = rgb
                                elif theme_color:
                                    new_run.font.color.theme_color = theme_color
                                    if brightness is not None:
                                        new_run.font.color.brightness = brightness
                            except Exception:
                                pass
                            if bold is not None:
                                new_run.font.bold = bold
                            if italic is not None:
                                new_run.font.italic = italic
                        else:
                            shape.text_frame.text = text
                        return prs_obj
            return prs_obj

        # Determine today's date in Europe/Zurich timezone
        ts = pd.Timestamp.now(tz="Europe/Zurich")
        human_date = f"{ts.strftime('%B')} {ts.day}, {ts.year}"
        stamp_ddmmyyyy = f"{ts.day:02d}{ts.month:02d}{ts.year}"
        # Try to replace a placeholder [DataIC] within any text run on the first slide.
        # The placeholder is replaced with the formatted date while preserving the
        # run's font attributes (size, colour, bold and italic).  If no
        # run with text exactly "[DataIC]" is found, fall back to replacing
        # the entire shape named "DataIC" using the helper above.
        _date_set = False
        for slide in prs.slides:
            if _date_set:
                break
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                tf = shape.text_frame
                for para in tf.paragraphs:
                    for run in para.runs:
                        if run.text.strip() == "[DataIC]":
                            # Preserve the existing font styling
                            size = run.font.size
                            color = run.font.color
                            rgb = getattr(color, "rgb", None) if color else None
                            theme_color = getattr(color, "theme_color", None) if color else None
                            brightness = getattr(color, "brightness", None) if color else None
                            bold = run.font.bold
                            italic = run.font.italic
                            # Replace the text
                            run.text = human_date
                            # Reapply styling
                            if size:
                                run.font.size = size
                            try:
                                if rgb:
                                    run.font.color.rgb = rgb
                                elif theme_color:
                                    run.font.color.theme_color = theme_color
                                    if brightness is not None:
                                        run.font.color.brightness = brightness
                            except Exception:
                                pass
                            if bold is not None:
                                run.font.bold = bold
                            if italic is not None:
                                run.font.italic = italic
                            _date_set = True
                            break
                    if _date_set:
                        break
                if _date_set:
                    break
        # Fallback: if placeholder not found, replace shape named "DataIC"
        if not _date_set:
            prs = _set_text_in_named_textbox(prs, "DataIC", human_date)

        # Determine which equity index was selected for technical analysis (not used here since we insert all indices)
        selected_index = st.session_state.get("ta_equity_index", "S&P 500")

        # Retrieve anchors for SPX, CSI, Nikkei, TASI, Sensex and DAX slides
        spx_anchor_dt = st.session_state.get("spx_anchor")
        csi_anchor_dt = st.session_state.get("csi_anchor")
        nikkei_anchor_dt = st.session_state.get("nikkei_anchor")
        tasi_anchor_dt = st.session_state.get("tasi_anchor")
        sensex_anchor_dt = st.session_state.get("sensex_anchor")
        dax_anchor_dt = st.session_state.get("dax_anchor")
        smi_anchor_dt = st.session_state.get("smi_anchor")
        ibov_anchor_dt = st.session_state.get("ibov_anchor")
        mexbol_anchor_dt = st.session_state.get("mexbol_anchor")
        # Anchor for Gold regression channel (commodity)
        gold_anchor_dt = st.session_state.get("gold_anchor")
        # Anchor for Silver regression channel (commodity)
        silver_anchor_dt = st.session_state.get("silver_anchor")
        # Anchor for Platinum regression channel (commodity)
        platinum_anchor_dt = st.session_state.get("platinum_anchor")
        # Anchor for Palladium regression channel (commodity)
        palladium_anchor_dt = st.session_state.get("palladium_anchor")
        # Anchor for Oil regression channel (commodity)
        oil_anchor_dt = st.session_state.get("oil_anchor")
        # Anchor for Copper regression channel (commodity)
        copper_anchor_dt = st.session_state.get("copper_anchor")
        # Anchor for Bitcoin regression channel (crypto)
        bitcoin_anchor_dt = st.session_state.get("bitcoin_anchor")
        # Anchor for Ethereum regression channel (crypto)
        ethereum_anchor_dt = st.session_state.get("ethereum_anchor")
        # Anchor for Ripple regression channel (crypto)
        ripple_anchor_dt = st.session_state.get("ripple_anchor")
        # Anchor for Solana regression channel (crypto)
        solana_anchor_dt = st.session_state.get("solana_anchor")
        # Anchor for Binance regression channel (crypto)
        binance_anchor_dt = st.session_state.get("binance_anchor")

        # Common price mode
        pmode = st.session_state.get("price_mode", "Last Price")

        # ===================================================================
        # PHASE 2 OPTIMIZATION: DISABLED (overhead > benefit on fast systems)
        # ===================================================================
        # Testing showed that parallel chart generation adds overhead that exceeds
        # the benefit on fast systems (Mac M4). The bottleneck is likely PowerPoint
        # insertion, not chart generation. Keeping Phase 1 (Excel caching) which
        # provides consistent 20-30% speedup with zero overhead.
        #
        # The parallel code is preserved below but commented out. It may be useful
        # for slower systems where chart generation is the bottleneck.
        # ===================================================================

        # DISABLED: Parallel chart generation (adds 10s overhead on Mac M4)
        # from technical_analysis.common_helpers import enable_image_cache, disable_image_cache
        # from technical_analysis.prewarm_all_instrument_charts import prewarm_all_instrument_charts
        # enable_image_cache()
        # ... (parallel generation code) ...

        # Continue with sequential PPT generation (uses Phase 1 Excel caching)
        # ===================================================================

        # ------------------------------------------------------------------
        # Insert SPX Technical Analysis v2 chart (Chart.js + Playwright)
        # ------------------------------------------------------------------
        try:
            update_progress("Processing S&P 500 Technical Analysis...")
            # Get DMAS scores from session state (uses individual keys like "spx_dmas")
            spx_dmas = st.session_state.get("spx_dmas", 50)
            spx_dmas_prev = st.session_state.get("spx_last_week_avg", spx_dmas)
            spx_tech = _get_spx_technical_score(excel_path_for_ppt)
            spx_momentum = _get_spx_momentum_score(excel_path_for_ppt)
            print(f"[Tech V2] SPX DMAS: {spx_dmas}, Prev Week: {spx_dmas_prev}, Tech: {spx_tech}, Mom: {spx_momentum}")

            # Get previous week Technical/Momentum/RSI scores from history
            spx_tech_prev = st.session_state.get("spx_last_week_tech", None)
            spx_mom_prev = st.session_state.get("spx_last_week_mom", None)
            spx_rsi_prev = st.session_state.get("spx_last_week_rsi", None)
            print(f"[Tech V2] Prev week scores - Tech: {spx_tech_prev}, Mom: {spx_mom_prev}, RSI: {spx_rsi_prev}")

            # Get gap information for change text formatting
            spx_days_gap = st.session_state.get("spx_prev_days_gap", None)
            spx_prev_date = st.session_state.get("spx_prev_date", None)

            # Compute used date for SPX source footnote
            try:
                import pandas as pd
                df_prices = pd.read_excel(excel_path_for_ppt, sheet_name="data_prices")
                df_prices = df_prices.drop(index=0)
                df_prices = df_prices[df_prices[df_prices.columns[0]] != "DATES"]
                df_prices["Date"] = pd.to_datetime(df_prices[df_prices.columns[0]], errors="coerce")
                # Filter by "Data As Of" date if set
                if "data_as_of" in st.session_state:
                    df_prices = df_prices[df_prices["Date"] <= pd.Timestamp(st.session_state["data_as_of"])]
                df_prices["Price"] = pd.to_numeric(df_prices["SPX Index"], errors="coerce")
                df_prices = df_prices.dropna(subset=["Date", "Price"]).sort_values("Date").reset_index(drop=True)[
                    ["Date", "Price"]
                ]
                df_adj, used_date_spx = adjust_prices_for_mode(df_prices, pmode)
            except Exception:
                used_date_spx = None

            v2_bytes, v2_date = create_technical_analysis_v2_chart(
                excel_path_for_ppt,
                ticker="SPX Index",
                price_mode=pmode,
                dmas_score=int(spx_dmas),
                dmas_prev_week=int(spx_dmas_prev),
                technical_score=spx_tech,
                technical_prev_week=spx_tech_prev,
                momentum_score=spx_momentum,
                momentum_prev_week=spx_mom_prev,
                rsi_prev_week=spx_rsi_prev,
                days_gap=spx_days_gap,
                previous_date=spx_prev_date,
            )
            # Get the same view and subtitle
            v2_view_text = st.session_state.get("spx_selected_view")
            # Prepend index name if not already present
            if v2_view_text and not v2_view_text.lower().startswith("s&p 500"):
                v2_view_text = f"S&P 500: {v2_view_text}"
            v2_subtitle = st.session_state.get("spx_subtitle", "")

            prs = insert_technical_analysis_v2_slide(
                prs,
                v2_bytes,
                used_date=used_date_spx,
                price_mode=pmode,
                placeholder_name="spx_v2",
                view_text=v2_view_text,
                subtitle_text=v2_subtitle,
            )
        except Exception as e:
            print(f"[Tech V2] SPX v2 chart error: {e}")
            import traceback
            traceback.print_exc()

        # ------------------------------------------------------------------
        # Insert CSI Technical Analysis v2 chart (Chart.js + Playwright)
        # ------------------------------------------------------------------
        try:
            update_progress("Processing CSI 300 Technical Analysis...")
            # Get DMAS scores from session state
            csi_dmas = st.session_state.get("csi_dmas", 50)
            csi_dmas_prev = st.session_state.get("csi_last_week_avg", csi_dmas)
            csi_tech = _get_csi_technical_score(excel_path_for_ppt)
            csi_momentum = _get_csi_momentum_score(excel_path_for_ppt)
            print(f"[Tech V2] CSI DMAS: {csi_dmas}, Prev Week: {csi_dmas_prev}, Tech: {csi_tech}, Mom: {csi_momentum}")

            # Get previous week Technical/Momentum/RSI scores from history
            csi_tech_prev = st.session_state.get("csi_last_week_tech", None)
            csi_mom_prev = st.session_state.get("csi_last_week_mom", None)
            csi_rsi_prev = st.session_state.get("csi_last_week_rsi", None)
            print(f"[Tech V2] CSI Prev week scores - Tech: {csi_tech_prev}, Mom: {csi_mom_prev}, RSI: {csi_rsi_prev}")

            # Get gap information for change text formatting
            csi_days_gap = st.session_state.get("csi_prev_days_gap", None)
            csi_prev_date = st.session_state.get("csi_prev_date", None)

            # Compute used date for CSI source footnote
            try:
                import pandas as pd
                df_prices_csi = pd.read_excel(excel_path_for_ppt, sheet_name="data_prices")
                df_prices_csi = df_prices_csi.drop(index=0)
                df_prices_csi = df_prices_csi[df_prices_csi[df_prices_csi.columns[0]] != "DATES"]
                df_prices_csi["Date"] = pd.to_datetime(df_prices_csi[df_prices_csi.columns[0]], errors="coerce")
                # Filter by "Data As Of" date if set
                if "data_as_of" in st.session_state:
                    df_prices_csi = df_prices_csi[df_prices_csi["Date"] <= pd.Timestamp(st.session_state["data_as_of"])]
                df_prices_csi["Price"] = pd.to_numeric(df_prices_csi["SHSZ300 Index"], errors="coerce")
                df_prices_csi = df_prices_csi.dropna(subset=["Date", "Price"]).sort_values("Date").reset_index(drop=True)[
                    ["Date", "Price"]
                ]
                df_adj_csi, used_date_csi = adjust_prices_for_mode(df_prices_csi, pmode)
            except Exception:
                used_date_csi = None

            v2_bytes_csi, v2_date_csi = create_technical_analysis_v2_chart(
                excel_path_for_ppt,
                ticker="SHSZ300 Index",
                price_mode=pmode,
                dmas_score=int(csi_dmas),
                dmas_prev_week=int(csi_dmas_prev),
                technical_score=csi_tech,
                technical_prev_week=csi_tech_prev,
                momentum_score=csi_momentum,
                momentum_prev_week=csi_mom_prev,
                rsi_prev_week=csi_rsi_prev,
                days_gap=csi_days_gap,
                previous_date=csi_prev_date,
            )
            # Get the view and subtitle
            v2_view_text_csi = st.session_state.get("csi_selected_view")
            # Prepend index name if not already present
            if v2_view_text_csi and not v2_view_text_csi.lower().startswith("csi"):
                v2_view_text_csi = f"CSI 300: {v2_view_text_csi}"
            v2_subtitle_csi = st.session_state.get("csi_subtitle", "")

            prs = insert_technical_analysis_v2_slide(
                prs,
                v2_bytes_csi,
                used_date=used_date_csi,
                price_mode=pmode,
                placeholder_name="csi_v2",
                view_text=v2_view_text_csi,
                subtitle_text=v2_subtitle_csi,
            )
        except Exception as e:
            print(f"[Tech V2] CSI v2 chart error: {e}")
            import traceback
            traceback.print_exc()

        # ------------------------------------------------------------------
        # Insert Nikkei Technical Analysis v2 chart (Chart.js + Playwright)
        # ------------------------------------------------------------------
        try:
            update_progress("Processing Nikkei 225 Technical Analysis...")
            # Get DMAS scores from session state
            nikkei_dmas = st.session_state.get("nikkei_dmas", 50)
            nikkei_dmas_prev = st.session_state.get("nikkei_last_week_avg", nikkei_dmas)
            nikkei_tech = _get_nikkei_technical_score(excel_path_for_ppt)
            nikkei_momentum = _get_nikkei_momentum_score(excel_path_for_ppt)
            print(f"[Tech V2] Nikkei DMAS: {nikkei_dmas}, Prev Week: {nikkei_dmas_prev}, Tech: {nikkei_tech}, Mom: {nikkei_momentum}")

            # Get previous week Technical/Momentum/RSI scores from history
            nikkei_tech_prev = st.session_state.get("nikkei_last_week_tech", None)
            nikkei_mom_prev = st.session_state.get("nikkei_last_week_mom", None)
            nikkei_rsi_prev = st.session_state.get("nikkei_last_week_rsi", None)
            print(f"[Tech V2] Nikkei Prev week scores - Tech: {nikkei_tech_prev}, Mom: {nikkei_mom_prev}, RSI: {nikkei_rsi_prev}")

            # Get gap information for change text formatting
            nikkei_days_gap = st.session_state.get("nikkei_prev_days_gap", None)
            nikkei_prev_date = st.session_state.get("nikkei_prev_date", None)

            # Compute used date for Nikkei source footnote
            try:
                import pandas as pd
                df_prices_nikkei = pd.read_excel(excel_path_for_ppt, sheet_name="data_prices")
                df_prices_nikkei = df_prices_nikkei.drop(index=0)
                df_prices_nikkei = df_prices_nikkei[df_prices_nikkei[df_prices_nikkei.columns[0]] != "DATES"]
                df_prices_nikkei["Date"] = pd.to_datetime(df_prices_nikkei[df_prices_nikkei.columns[0]], errors="coerce")
                # Filter by "Data As Of" date if set
                if "data_as_of" in st.session_state:
                    df_prices_nikkei = df_prices_nikkei[df_prices_nikkei["Date"] <= pd.Timestamp(st.session_state["data_as_of"])]
                df_prices_nikkei["Price"] = pd.to_numeric(df_prices_nikkei["NKY Index"], errors="coerce")
                df_prices_nikkei = df_prices_nikkei.dropna(subset=["Date", "Price"]).sort_values("Date").reset_index(drop=True)[
                    ["Date", "Price"]
                ]
                df_adj_nikkei, used_date_nikkei = adjust_prices_for_mode(df_prices_nikkei, pmode)
            except Exception:
                used_date_nikkei = None

            v2_bytes_nikkei, v2_date_nikkei = create_technical_analysis_v2_chart(
                excel_path_for_ppt,
                ticker="NKY Index",
                price_mode=pmode,
                dmas_score=int(nikkei_dmas),
                dmas_prev_week=int(nikkei_dmas_prev),
                technical_score=nikkei_tech,
                technical_prev_week=nikkei_tech_prev,
                momentum_score=nikkei_momentum,
                momentum_prev_week=nikkei_mom_prev,
                rsi_prev_week=nikkei_rsi_prev,
                days_gap=nikkei_days_gap,
                previous_date=nikkei_prev_date,
            )
            # Get the view and subtitle
            v2_view_text_nikkei = st.session_state.get("nikkei_selected_view")
            # Prepend index name if not already present
            if v2_view_text_nikkei and not v2_view_text_nikkei.lower().startswith("nikkei"):
                v2_view_text_nikkei = f"Nikkei 225: {v2_view_text_nikkei}"
            v2_subtitle_nikkei = st.session_state.get("nikkei_subtitle", "")

            prs = insert_technical_analysis_v2_slide(
                prs,
                v2_bytes_nikkei,
                used_date=used_date_nikkei,
                price_mode=pmode,
                placeholder_name="nikkei_v2",
                view_text=v2_view_text_nikkei,
                subtitle_text=v2_subtitle_nikkei,
            )
        except Exception as e:
            print(f"[Tech V2] Nikkei v2 chart error: {e}")
            import traceback
            traceback.print_exc()

        # ------------------------------------------------------------------
        # Insert TASI Technical Analysis v2 chart (Chart.js + Playwright)
        # ------------------------------------------------------------------
        try:
            update_progress("Processing TASI Technical Analysis...")
            # Get DMAS scores from session state
            tasi_dmas = st.session_state.get("tasi_dmas", 50)
            tasi_dmas_prev = st.session_state.get("tasi_last_week_avg", tasi_dmas)
            tasi_tech = _get_tasi_technical_score(excel_path_for_ppt)
            tasi_momentum = _get_tasi_momentum_score(excel_path_for_ppt)
            print(f"[Tech V2] TASI DMAS: {tasi_dmas}, Prev Week: {tasi_dmas_prev}, Tech: {tasi_tech}, Mom: {tasi_momentum}")

            # Get previous week Technical/Momentum/RSI scores from history
            tasi_tech_prev = st.session_state.get("tasi_last_week_tech", None)
            tasi_mom_prev = st.session_state.get("tasi_last_week_mom", None)
            tasi_rsi_prev = st.session_state.get("tasi_last_week_rsi", None)
            print(f"[Tech V2] TASI Prev week scores - Tech: {tasi_tech_prev}, Mom: {tasi_mom_prev}, RSI: {tasi_rsi_prev}")

            # Get gap information for change text formatting
            tasi_days_gap = st.session_state.get("tasi_prev_days_gap", None)
            tasi_prev_date = st.session_state.get("tasi_prev_date", None)

            # Compute used date for TASI source footnote
            try:
                import pandas as pd
                df_prices_tasi = pd.read_excel(excel_path_for_ppt, sheet_name="data_prices")
                df_prices_tasi = df_prices_tasi.drop(index=0)
                df_prices_tasi = df_prices_tasi[df_prices_tasi[df_prices_tasi.columns[0]] != "DATES"]
                df_prices_tasi["Date"] = pd.to_datetime(df_prices_tasi[df_prices_tasi.columns[0]], errors="coerce")
                # Filter by "Data As Of" date if set
                if "data_as_of" in st.session_state:
                    df_prices_tasi = df_prices_tasi[df_prices_tasi["Date"] <= pd.Timestamp(st.session_state["data_as_of"])]
                df_prices_tasi["Price"] = pd.to_numeric(df_prices_tasi["SASEIDX Index"], errors="coerce")
                df_prices_tasi = df_prices_tasi.dropna(subset=["Date", "Price"]).sort_values("Date").reset_index(drop=True)[
                    ["Date", "Price"]
                ]
                df_adj_tasi, used_date_tasi = adjust_prices_for_mode(df_prices_tasi, pmode)
            except Exception:
                used_date_tasi = None

            v2_bytes_tasi, v2_date_tasi = create_technical_analysis_v2_chart(
                excel_path_for_ppt,
                ticker="SASEIDX Index",
                price_mode=pmode,
                dmas_score=int(tasi_dmas),
                dmas_prev_week=int(tasi_dmas_prev),
                technical_score=tasi_tech,
                technical_prev_week=tasi_tech_prev,
                momentum_score=tasi_momentum,
                momentum_prev_week=tasi_mom_prev,
                rsi_prev_week=tasi_rsi_prev,
                days_gap=tasi_days_gap,
                previous_date=tasi_prev_date,
            )
            # Get the view and subtitle
            v2_view_text_tasi = st.session_state.get("tasi_selected_view")
            # Prepend index name if not already present
            if v2_view_text_tasi and not v2_view_text_tasi.lower().startswith("tasi"):
                v2_view_text_tasi = f"TASI: {v2_view_text_tasi}"
            v2_subtitle_tasi = st.session_state.get("tasi_subtitle", "")

            prs = insert_technical_analysis_v2_slide(
                prs,
                v2_bytes_tasi,
                used_date=used_date_tasi,
                price_mode=pmode,
                placeholder_name="tasi_v2",
                view_text=v2_view_text_tasi,
                subtitle_text=v2_subtitle_tasi,
            )
        except Exception as e:
            print(f"[Tech V2] TASI v2 chart error: {e}")
            import traceback
            traceback.print_exc()

        # ------------------------------------------------------------------
        # Insert Sensex Technical Analysis v2 chart (Chart.js + Playwright)
        # ------------------------------------------------------------------
        try:
            update_progress("Processing Sensex Technical Analysis...")
            # Get DMAS scores from session state
            sensex_dmas = st.session_state.get("sensex_dmas", 50)
            sensex_dmas_prev = st.session_state.get("sensex_last_week_avg", sensex_dmas)
            sensex_tech = _get_sensex_technical_score(excel_path_for_ppt)
            sensex_momentum = _get_sensex_momentum_score(excel_path_for_ppt)
            print(f"[Tech V2] Sensex DMAS: {sensex_dmas}, Prev Week: {sensex_dmas_prev}, Tech: {sensex_tech}, Mom: {sensex_momentum}")

            # Get previous week Technical/Momentum/RSI scores from history
            sensex_tech_prev = st.session_state.get("sensex_last_week_tech", None)
            sensex_mom_prev = st.session_state.get("sensex_last_week_mom", None)
            sensex_rsi_prev = st.session_state.get("sensex_last_week_rsi", None)
            print(f"[Tech V2] Sensex Prev week scores - Tech: {sensex_tech_prev}, Mom: {sensex_mom_prev}, RSI: {sensex_rsi_prev}")

            # Get gap information for change text formatting
            sensex_days_gap = st.session_state.get("sensex_prev_days_gap", None)
            sensex_prev_date = st.session_state.get("sensex_prev_date", None)

            # Compute used date for Sensex source footnote
            try:
                import pandas as pd
                df_prices_sensex = pd.read_excel(excel_path_for_ppt, sheet_name="data_prices")
                df_prices_sensex = df_prices_sensex.drop(index=0)
                df_prices_sensex = df_prices_sensex[df_prices_sensex[df_prices_sensex.columns[0]] != "DATES"]
                df_prices_sensex["Date"] = pd.to_datetime(df_prices_sensex[df_prices_sensex.columns[0]], errors="coerce")
                # Filter by "Data As Of" date if set
                if "data_as_of" in st.session_state:
                    df_prices_sensex = df_prices_sensex[df_prices_sensex["Date"] <= pd.Timestamp(st.session_state["data_as_of"])]
                df_prices_sensex["Price"] = pd.to_numeric(df_prices_sensex["SENSEX Index"], errors="coerce")
                df_prices_sensex = df_prices_sensex.dropna(subset=["Date", "Price"]).sort_values("Date").reset_index(drop=True)[
                    ["Date", "Price"]
                ]
                df_adj_sensex, used_date_sensex = adjust_prices_for_mode(df_prices_sensex, pmode)
            except Exception:
                used_date_sensex = None

            v2_bytes_sensex, v2_date_sensex = create_technical_analysis_v2_chart(
                excel_path_for_ppt,
                ticker="SENSEX Index",
                price_mode=pmode,
                dmas_score=int(sensex_dmas),
                dmas_prev_week=int(sensex_dmas_prev),
                technical_score=sensex_tech,
                technical_prev_week=sensex_tech_prev,
                momentum_score=sensex_momentum,
                momentum_prev_week=sensex_mom_prev,
                rsi_prev_week=sensex_rsi_prev,
                days_gap=sensex_days_gap,
                previous_date=sensex_prev_date,
            )
            # Get the view and subtitle
            v2_view_text_sensex = st.session_state.get("sensex_selected_view")
            # Prepend index name if not already present
            if v2_view_text_sensex and not v2_view_text_sensex.lower().startswith("sensex"):
                v2_view_text_sensex = f"Sensex: {v2_view_text_sensex}"
            v2_subtitle_sensex = st.session_state.get("sensex_subtitle", "")

            prs = insert_technical_analysis_v2_slide(
                prs,
                v2_bytes_sensex,
                used_date=used_date_sensex,
                price_mode=pmode,
                placeholder_name="sensex_v2",
                view_text=v2_view_text_sensex,
                subtitle_text=v2_subtitle_sensex,
            )
        except Exception as e:
            print(f"[Tech V2] Sensex v2 chart error: {e}")
            import traceback
            traceback.print_exc()

        # ------------------------------------------------------------------
        # Insert DAX Technical Analysis v2 chart (Chart.js + Playwright)
        # ------------------------------------------------------------------
        try:
            update_progress("Processing DAX Technical Analysis...")
            # Get DMAS scores from session state
            dax_dmas = st.session_state.get("dax_dmas", 50)
            dax_dmas_prev = st.session_state.get("dax_last_week_avg", dax_dmas)
            dax_tech = _get_dax_technical_score(excel_path_for_ppt)
            dax_momentum = _get_dax_momentum_score(excel_path_for_ppt)
            print(f"[Tech V2] DAX DMAS: {dax_dmas}, Prev Week: {dax_dmas_prev}, Tech: {dax_tech}, Mom: {dax_momentum}")

            # Get previous week Technical/Momentum/RSI scores from history
            dax_tech_prev = st.session_state.get("dax_last_week_tech", None)
            dax_mom_prev = st.session_state.get("dax_last_week_mom", None)
            dax_rsi_prev = st.session_state.get("dax_last_week_rsi", None)
            print(f"[Tech V2] DAX Prev week scores - Tech: {dax_tech_prev}, Mom: {dax_mom_prev}, RSI: {dax_rsi_prev}")

            # Get gap information for change text formatting
            dax_days_gap = st.session_state.get("dax_prev_days_gap", None)
            dax_prev_date = st.session_state.get("dax_prev_date", None)

            # Compute used date for DAX source footnote
            try:
                import pandas as pd
                df_prices_dax = pd.read_excel(excel_path_for_ppt, sheet_name="data_prices")
                df_prices_dax = df_prices_dax.drop(index=0)
                df_prices_dax = df_prices_dax[df_prices_dax[df_prices_dax.columns[0]] != "DATES"]
                df_prices_dax["Date"] = pd.to_datetime(df_prices_dax[df_prices_dax.columns[0]], errors="coerce")
                # Filter by "Data As Of" date if set
                if "data_as_of" in st.session_state:
                    df_prices_dax = df_prices_dax[df_prices_dax["Date"] <= pd.Timestamp(st.session_state["data_as_of"])]
                df_prices_dax["Price"] = pd.to_numeric(df_prices_dax["DAX Index"], errors="coerce")
                df_prices_dax = df_prices_dax.dropna(subset=["Date", "Price"]).sort_values("Date").reset_index(drop=True)[
                    ["Date", "Price"]
                ]
                df_adj_dax, used_date_dax = adjust_prices_for_mode(df_prices_dax, pmode)
            except Exception:
                used_date_dax = None

            v2_bytes_dax, v2_date_dax = create_technical_analysis_v2_chart(
                excel_path_for_ppt,
                ticker="DAX Index",
                price_mode=pmode,
                dmas_score=int(dax_dmas),
                dmas_prev_week=int(dax_dmas_prev),
                technical_score=dax_tech,
                technical_prev_week=dax_tech_prev,
                momentum_score=dax_momentum,
                momentum_prev_week=dax_mom_prev,
                rsi_prev_week=dax_rsi_prev,
                days_gap=dax_days_gap,
                previous_date=dax_prev_date,
            )
            # Get the view and subtitle
            v2_view_text_dax = st.session_state.get("dax_selected_view")
            # Prepend index name if not already present
            if v2_view_text_dax and not v2_view_text_dax.lower().startswith("dax"):
                v2_view_text_dax = f"DAX: {v2_view_text_dax}"
            v2_subtitle_dax = st.session_state.get("dax_subtitle", "")

            prs = insert_technical_analysis_v2_slide(
                prs,
                v2_bytes_dax,
                used_date=used_date_dax,
                price_mode=pmode,
                placeholder_name="dax_v2",
                view_text=v2_view_text_dax,
                subtitle_text=v2_subtitle_dax,
            )
        except Exception as e:
            print(f"[Tech V2] DAX v2 chart error: {e}")
            import traceback
            traceback.print_exc()

        # ------------------------------------------------------------------
        # Insert SMI Technical Analysis v2 chart (Chart.js + Playwright)
        # ------------------------------------------------------------------
        try:
            update_progress("Processing SMI Technical Analysis...")
            # Get DMAS scores from session state
            smi_dmas = st.session_state.get("smi_dmas", 50)
            smi_dmas_prev = st.session_state.get("smi_last_week_avg", smi_dmas)
            smi_tech = _get_smi_technical_score(excel_path_for_ppt)
            smi_momentum = _get_smi_momentum_score(excel_path_for_ppt)
            print(f"[Tech V2] SMI DMAS: {smi_dmas}, Prev Week: {smi_dmas_prev}, Tech: {smi_tech}, Mom: {smi_momentum}")

            # Get previous week Technical/Momentum/RSI scores from history
            smi_tech_prev = st.session_state.get("smi_last_week_tech", None)
            smi_mom_prev = st.session_state.get("smi_last_week_mom", None)
            smi_rsi_prev = st.session_state.get("smi_last_week_rsi", None)
            print(f"[Tech V2] SMI Prev week scores - Tech: {smi_tech_prev}, Mom: {smi_mom_prev}, RSI: {smi_rsi_prev}")

            # Get gap information for change text formatting
            smi_days_gap = st.session_state.get("smi_prev_days_gap", None)
            smi_prev_date = st.session_state.get("smi_prev_date", None)

            # Compute used date for SMI source footnote
            try:
                import pandas as pd
                df_prices_smi = pd.read_excel(excel_path_for_ppt, sheet_name="data_prices")
                df_prices_smi = df_prices_smi.drop(index=0)
                df_prices_smi = df_prices_smi[df_prices_smi[df_prices_smi.columns[0]] != "DATES"]
                df_prices_smi["Date"] = pd.to_datetime(df_prices_smi[df_prices_smi.columns[0]], errors="coerce")
                # Filter by "Data As Of" date if set
                if "data_as_of" in st.session_state:
                    df_prices_smi = df_prices_smi[df_prices_smi["Date"] <= pd.Timestamp(st.session_state["data_as_of"])]
                df_prices_smi["Price"] = pd.to_numeric(df_prices_smi["SMI Index"], errors="coerce")
                df_prices_smi = df_prices_smi.dropna(subset=["Date", "Price"]).sort_values("Date").reset_index(drop=True)[
                    ["Date", "Price"]
                ]
                df_adj_smi, used_date_smi = adjust_prices_for_mode(df_prices_smi, pmode)
            except Exception:
                used_date_smi = None

            v2_bytes_smi, v2_date_smi = create_technical_analysis_v2_chart(
                excel_path_for_ppt,
                ticker="SMI Index",
                price_mode=pmode,
                dmas_score=int(smi_dmas),
                dmas_prev_week=int(smi_dmas_prev),
                technical_score=smi_tech,
                technical_prev_week=smi_tech_prev,
                momentum_score=smi_momentum,
                momentum_prev_week=smi_mom_prev,
                rsi_prev_week=smi_rsi_prev,
                days_gap=smi_days_gap,
                previous_date=smi_prev_date,
            )
            # Get the view and subtitle
            v2_view_text_smi = st.session_state.get("smi_selected_view")
            # Prepend index name if not already present
            if v2_view_text_smi and not v2_view_text_smi.lower().startswith("smi"):
                v2_view_text_smi = f"SMI: {v2_view_text_smi}"
            v2_subtitle_smi = st.session_state.get("smi_subtitle", "")

            prs = insert_technical_analysis_v2_slide(
                prs,
                v2_bytes_smi,
                used_date=used_date_smi,
                price_mode=pmode,
                placeholder_name="smi_v2",
                view_text=v2_view_text_smi,
                subtitle_text=v2_subtitle_smi,
            )
        except Exception as e:
            print(f"[Tech V2] SMI v2 chart error: {e}")
            import traceback
            traceback.print_exc()

        # ------------------------------------------------------------------
        # Insert IBOV Technical Analysis v2 chart (Chart.js + Playwright)
        # ------------------------------------------------------------------
        try:
            update_progress("Processing IBOV Technical Analysis...")
            # Get DMAS scores from session state
            ibov_dmas = st.session_state.get("ibov_dmas", 50)
            ibov_dmas_prev = st.session_state.get("ibov_last_week_avg", ibov_dmas)
            ibov_tech = _get_ibov_technical_score(excel_path_for_ppt)
            ibov_momentum = _get_ibov_momentum_score(excel_path_for_ppt)
            print(f"[Tech V2] IBOV DMAS: {ibov_dmas}, Prev Week: {ibov_dmas_prev}, Tech: {ibov_tech}, Mom: {ibov_momentum}")

            # Get previous week Technical/Momentum/RSI scores from history
            ibov_tech_prev = st.session_state.get("ibov_last_week_tech", None)
            ibov_mom_prev = st.session_state.get("ibov_last_week_mom", None)
            ibov_rsi_prev = st.session_state.get("ibov_last_week_rsi", None)
            print(f"[Tech V2] IBOV Prev week scores - Tech: {ibov_tech_prev}, Mom: {ibov_mom_prev}, RSI: {ibov_rsi_prev}")

            # Get gap information for change text formatting
            ibov_days_gap = st.session_state.get("ibov_prev_days_gap", None)
            ibov_prev_date = st.session_state.get("ibov_prev_date", None)

            # Compute used date for IBOV source footnote
            try:
                import pandas as pd
                df_prices_ibov = pd.read_excel(excel_path_for_ppt, sheet_name="data_prices")
                df_prices_ibov = df_prices_ibov.drop(index=0)
                df_prices_ibov = df_prices_ibov[df_prices_ibov[df_prices_ibov.columns[0]] != "DATES"]
                df_prices_ibov["Date"] = pd.to_datetime(df_prices_ibov[df_prices_ibov.columns[0]], errors="coerce")
                # Filter by "Data As Of" date if set
                if "data_as_of" in st.session_state:
                    df_prices_ibov = df_prices_ibov[df_prices_ibov["Date"] <= pd.Timestamp(st.session_state["data_as_of"])]
                df_prices_ibov["Price"] = pd.to_numeric(df_prices_ibov["IBOV Index"], errors="coerce")
                df_prices_ibov = df_prices_ibov.dropna(subset=["Date", "Price"]).sort_values("Date").reset_index(drop=True)[
                    ["Date", "Price"]
                ]
                df_adj_ibov, used_date_ibov = adjust_prices_for_mode(df_prices_ibov, pmode)
            except Exception:
                used_date_ibov = None

            v2_bytes_ibov, v2_date_ibov = create_technical_analysis_v2_chart(
                excel_path_for_ppt,
                ticker="IBOV Index",
                price_mode=pmode,
                dmas_score=int(ibov_dmas),
                dmas_prev_week=int(ibov_dmas_prev),
                technical_score=ibov_tech,
                technical_prev_week=ibov_tech_prev,
                momentum_score=ibov_momentum,
                momentum_prev_week=ibov_mom_prev,
                rsi_prev_week=ibov_rsi_prev,
                days_gap=ibov_days_gap,
                previous_date=ibov_prev_date,
            )
            # Get the view and subtitle
            v2_view_text_ibov = st.session_state.get("ibov_selected_view")
            # Prepend index name if not already present
            if v2_view_text_ibov and not v2_view_text_ibov.lower().startswith("ibov"):
                v2_view_text_ibov = f"Bovespa: {v2_view_text_ibov}"
            v2_subtitle_ibov = st.session_state.get("ibov_subtitle", "")

            prs = insert_technical_analysis_v2_slide(
                prs,
                v2_bytes_ibov,
                used_date=used_date_ibov,
                price_mode=pmode,
                placeholder_name="ibov_v2",
                view_text=v2_view_text_ibov,
                subtitle_text=v2_subtitle_ibov,
            )
        except Exception as e:
            print(f"[Tech V2] IBOV v2 chart error: {e}")
            import traceback
            traceback.print_exc()

        # ------------------------------------------------------------------
        # Insert Mexbol Technical Analysis v2 chart (Chart.js + Playwright)
        # ------------------------------------------------------------------
        try:
            update_progress("Processing Mexbol Technical Analysis...")
            # Get DMAS scores from session state
            mexbol_dmas = st.session_state.get("mexbol_dmas", 50)
            mexbol_dmas_prev = st.session_state.get("mexbol_last_week_avg", mexbol_dmas)
            mexbol_tech = _get_mexbol_technical_score(excel_path_for_ppt)
            mexbol_momentum = _get_mexbol_momentum_score(excel_path_for_ppt)
            print(f"[Tech V2] Mexbol DMAS: {mexbol_dmas}, Prev Week: {mexbol_dmas_prev}, Tech: {mexbol_tech}, Mom: {mexbol_momentum}")

            # Get previous week Technical/Momentum/RSI scores from history
            mexbol_tech_prev = st.session_state.get("mexbol_last_week_tech", None)
            mexbol_mom_prev = st.session_state.get("mexbol_last_week_mom", None)
            mexbol_rsi_prev = st.session_state.get("mexbol_last_week_rsi", None)
            print(f"[Tech V2] Mexbol Prev week scores - Tech: {mexbol_tech_prev}, Mom: {mexbol_mom_prev}, RSI: {mexbol_rsi_prev}")

            # Get gap information for change text formatting
            mexbol_days_gap = st.session_state.get("mexbol_prev_days_gap", None)
            mexbol_prev_date = st.session_state.get("mexbol_prev_date", None)

            # Compute used date for Mexbol source footnote
            try:
                import pandas as pd
                df_prices_mexbol = pd.read_excel(excel_path_for_ppt, sheet_name="data_prices")
                df_prices_mexbol = df_prices_mexbol.drop(index=0)
                df_prices_mexbol = df_prices_mexbol[df_prices_mexbol[df_prices_mexbol.columns[0]] != "DATES"]
                df_prices_mexbol["Date"] = pd.to_datetime(df_prices_mexbol[df_prices_mexbol.columns[0]], errors="coerce")
                # Filter by "Data As Of" date if set
                if "data_as_of" in st.session_state:
                    df_prices_mexbol = df_prices_mexbol[df_prices_mexbol["Date"] <= pd.Timestamp(st.session_state["data_as_of"])]
                df_prices_mexbol["Price"] = pd.to_numeric(df_prices_mexbol["MEXBOL Index"], errors="coerce")
                df_prices_mexbol = df_prices_mexbol.dropna(subset=["Date", "Price"]).sort_values("Date").reset_index(drop=True)[
                    ["Date", "Price"]
                ]
                df_adj_mexbol, used_date_mexbol = adjust_prices_for_mode(df_prices_mexbol, pmode)
            except Exception:
                used_date_mexbol = None

            v2_bytes_mexbol, v2_date_mexbol = create_technical_analysis_v2_chart(
                excel_path_for_ppt,
                ticker="MEXBOL Index",
                price_mode=pmode,
                dmas_score=int(mexbol_dmas),
                dmas_prev_week=int(mexbol_dmas_prev),
                technical_score=mexbol_tech,
                technical_prev_week=mexbol_tech_prev,
                momentum_score=mexbol_momentum,
                momentum_prev_week=mexbol_mom_prev,
                rsi_prev_week=mexbol_rsi_prev,
                days_gap=mexbol_days_gap,
                previous_date=mexbol_prev_date,
            )
            # Get the view and subtitle
            v2_view_text_mexbol = st.session_state.get("mexbol_selected_view")
            # Prepend index name if not already present
            if v2_view_text_mexbol and not v2_view_text_mexbol.lower().startswith("mexbol"):
                v2_view_text_mexbol = f"Mexbol: {v2_view_text_mexbol}"
            v2_subtitle_mexbol = st.session_state.get("mexbol_subtitle", "")

            prs = insert_technical_analysis_v2_slide(
                prs,
                v2_bytes_mexbol,
                used_date=used_date_mexbol,
                price_mode=pmode,
                placeholder_name="mexbol_v2",
                view_text=v2_view_text_mexbol,
                subtitle_text=v2_subtitle_mexbol,
            )
        except Exception as e:
            print(f"[Tech V2] Mexbol v2 chart error: {e}")
            import traceback
            traceback.print_exc()

        # ------------------------------------------------------------------
        # Insert Gold Technical Analysis v2 chart (Chart.js + Playwright)
        # ------------------------------------------------------------------
        try:
            update_progress("Processing Gold Technical Analysis...")
            # Get DMAS scores from session state
            gold_dmas = st.session_state.get("gold_dmas", 50)
            gold_dmas_prev = st.session_state.get("gold_last_week_avg", gold_dmas)
            gold_tech = _get_gold_technical_score(excel_path_for_ppt)
            gold_momentum = _get_gold_momentum_score(excel_path_for_ppt)
            print(f"[Tech V2] Gold DMAS: {gold_dmas}, Prev Week: {gold_dmas_prev}, Tech: {gold_tech}, Mom: {gold_momentum}")

            # Get previous week Technical/Momentum/RSI scores from history
            gold_tech_prev = st.session_state.get("gold_last_week_tech", None)
            gold_mom_prev = st.session_state.get("gold_last_week_mom", None)
            gold_rsi_prev = st.session_state.get("gold_last_week_rsi", None)
            print(f"[Tech V2] Gold Prev week scores - Tech: {gold_tech_prev}, Mom: {gold_mom_prev}, RSI: {gold_rsi_prev}")

            # Get gap information for change text formatting
            gold_days_gap = st.session_state.get("gold_prev_days_gap", None)
            gold_prev_date = st.session_state.get("gold_prev_date", None)

            # Compute used date for Gold source footnote
            try:
                import pandas as pd
                df_prices_gold = pd.read_excel(excel_path_for_ppt, sheet_name="data_prices")
                df_prices_gold = df_prices_gold.drop(index=0)
                df_prices_gold = df_prices_gold[df_prices_gold[df_prices_gold.columns[0]] != "DATES"]
                df_prices_gold["Date"] = pd.to_datetime(df_prices_gold[df_prices_gold.columns[0]], errors="coerce")
                # Filter by "Data As Of" date if set
                if "data_as_of" in st.session_state:
                    df_prices_gold = df_prices_gold[df_prices_gold["Date"] <= pd.Timestamp(st.session_state["data_as_of"])]
                df_prices_gold["Price"] = pd.to_numeric(df_prices_gold["GCA Comdty"], errors="coerce")
                df_prices_gold = df_prices_gold.dropna(subset=["Date", "Price"]).sort_values("Date").reset_index(drop=True)[
                    ["Date", "Price"]
                ]
                df_adj_gold, used_date_gold = adjust_prices_for_mode(df_prices_gold, pmode)
            except Exception:
                used_date_gold = None

            v2_bytes_gold, v2_date_gold = create_technical_analysis_v2_chart(
                excel_path_for_ppt,
                ticker="GCA Comdty",
                price_mode=pmode,
                dmas_score=int(gold_dmas),
                dmas_prev_week=int(gold_dmas_prev),
                technical_score=gold_tech,
                technical_prev_week=gold_tech_prev,
                momentum_score=gold_momentum,
                momentum_prev_week=gold_mom_prev,
                rsi_prev_week=gold_rsi_prev,
                days_gap=gold_days_gap,
                previous_date=gold_prev_date,
            )
            # Get the view and subtitle
            v2_view_text_gold = st.session_state.get("gold_selected_view")
            # Prepend commodity name if not already present
            if v2_view_text_gold and not v2_view_text_gold.lower().startswith("gold"):
                v2_view_text_gold = f"Gold: {v2_view_text_gold}"
            v2_subtitle_gold = st.session_state.get("gold_subtitle", "")

            prs = insert_technical_analysis_v2_slide(
                prs,
                v2_bytes_gold,
                used_date=used_date_gold,
                price_mode=pmode,
                placeholder_name="gold_v2",
                view_text=v2_view_text_gold,
                subtitle_text=v2_subtitle_gold,
            )
        except Exception as e:
            print(f"[Tech V2] Gold v2 chart error: {e}")
            import traceback
            traceback.print_exc()

        # ------------------------------------------------------------------
        # Insert Silver Technical Analysis v2 chart (Chart.js + Playwright)
        # ------------------------------------------------------------------
        try:
            update_progress("Processing Silver Technical Analysis...")
            # Get DMAS scores from session state
            silver_dmas = st.session_state.get("silver_dmas", 50)
            silver_dmas_prev = st.session_state.get("silver_last_week_avg", silver_dmas)
            silver_tech = _get_silver_technical_score(excel_path_for_ppt)
            silver_momentum = _get_silver_momentum_score(excel_path_for_ppt)
            print(f"[Tech V2] Silver DMAS: {silver_dmas}, Prev Week: {silver_dmas_prev}, Tech: {silver_tech}, Mom: {silver_momentum}")

            # Get previous week Technical/Momentum/RSI scores from history
            silver_tech_prev = st.session_state.get("silver_last_week_tech", None)
            silver_mom_prev = st.session_state.get("silver_last_week_mom", None)
            silver_rsi_prev = st.session_state.get("silver_last_week_rsi", None)
            print(f"[Tech V2] Silver Prev week scores - Tech: {silver_tech_prev}, Mom: {silver_mom_prev}, RSI: {silver_rsi_prev}")

            # Get gap information for change text formatting
            silver_days_gap = st.session_state.get("silver_prev_days_gap", None)
            silver_prev_date = st.session_state.get("silver_prev_date", None)

            # Compute used date for Silver source footnote
            try:
                import pandas as pd
                df_prices_silver = pd.read_excel(excel_path_for_ppt, sheet_name="data_prices")
                df_prices_silver = df_prices_silver.drop(index=0)
                df_prices_silver = df_prices_silver[df_prices_silver[df_prices_silver.columns[0]] != "DATES"]
                df_prices_silver["Date"] = pd.to_datetime(df_prices_silver[df_prices_silver.columns[0]], errors="coerce")
                # Filter by "Data As Of" date if set
                if "data_as_of" in st.session_state:
                    df_prices_silver = df_prices_silver[df_prices_silver["Date"] <= pd.Timestamp(st.session_state["data_as_of"])]
                df_prices_silver["Price"] = pd.to_numeric(df_prices_silver["SIA Comdty"], errors="coerce")
                df_prices_silver = df_prices_silver.dropna(subset=["Date", "Price"]).sort_values("Date").reset_index(drop=True)[
                    ["Date", "Price"]
                ]
                df_adj_silver, used_date_silver = adjust_prices_for_mode(df_prices_silver, pmode)
            except Exception:
                used_date_silver = None

            v2_bytes_silver, v2_date_silver = create_technical_analysis_v2_chart(
                excel_path_for_ppt,
                ticker="SIA Comdty",
                price_mode=pmode,
                dmas_score=int(silver_dmas),
                dmas_prev_week=int(silver_dmas_prev),
                technical_score=silver_tech,
                technical_prev_week=silver_tech_prev,
                momentum_score=silver_momentum,
                momentum_prev_week=silver_mom_prev,
                rsi_prev_week=silver_rsi_prev,
                days_gap=silver_days_gap,
                previous_date=silver_prev_date,
            )
            # Get the view and subtitle
            v2_view_text_silver = st.session_state.get("silver_selected_view")
            # Prepend commodity name if not already present
            if v2_view_text_silver and not v2_view_text_silver.lower().startswith("silver"):
                v2_view_text_silver = f"Silver: {v2_view_text_silver}"
            v2_subtitle_silver = st.session_state.get("silver_subtitle", "")

            prs = insert_technical_analysis_v2_slide(
                prs,
                v2_bytes_silver,
                used_date=used_date_silver,
                price_mode=pmode,
                placeholder_name="silver_v2",
                view_text=v2_view_text_silver,
                subtitle_text=v2_subtitle_silver,
            )
        except Exception as e:
            print(f"[Tech V2] Silver v2 chart error: {e}")
            import traceback
            traceback.print_exc()

        # ------------------------------------------------------------------
        # Insert Platinum Technical Analysis v2 chart (Chart.js + Playwright)
        # ------------------------------------------------------------------
        try:
            update_progress("Processing Platinum Technical Analysis...")
            # Get DMAS scores from session state
            platinum_dmas = st.session_state.get("platinum_dmas", 50)
            platinum_dmas_prev = st.session_state.get("platinum_last_week_avg", platinum_dmas)
            platinum_tech = _get_platinum_technical_score(excel_path_for_ppt)
            platinum_momentum = _get_platinum_momentum_score(excel_path_for_ppt)
            print(f"[Tech V2] Platinum DMAS: {platinum_dmas}, Prev Week: {platinum_dmas_prev}, Tech: {platinum_tech}, Mom: {platinum_momentum}")

            # Get previous week Technical/Momentum/RSI scores from history
            platinum_tech_prev = st.session_state.get("platinum_last_week_tech", None)
            platinum_mom_prev = st.session_state.get("platinum_last_week_mom", None)
            platinum_rsi_prev = st.session_state.get("platinum_last_week_rsi", None)
            print(f"[Tech V2] Platinum Prev week scores - Tech: {platinum_tech_prev}, Mom: {platinum_mom_prev}, RSI: {platinum_rsi_prev}")

            # Get gap information for change text formatting
            platinum_days_gap = st.session_state.get("platinum_prev_days_gap", None)
            platinum_prev_date = st.session_state.get("platinum_prev_date", None)

            # Compute used date for Platinum source footnote
            try:
                import pandas as pd
                df_prices_platinum = pd.read_excel(excel_path_for_ppt, sheet_name="data_prices")
                df_prices_platinum = df_prices_platinum.drop(index=0)
                df_prices_platinum = df_prices_platinum[df_prices_platinum[df_prices_platinum.columns[0]] != "DATES"]
                df_prices_platinum["Date"] = pd.to_datetime(df_prices_platinum[df_prices_platinum.columns[0]], errors="coerce")
                # Filter by "Data As Of" date if set
                if "data_as_of" in st.session_state:
                    df_prices_platinum = df_prices_platinum[df_prices_platinum["Date"] <= pd.Timestamp(st.session_state["data_as_of"])]
                df_prices_platinum["Price"] = pd.to_numeric(df_prices_platinum["XPT Comdty"], errors="coerce")
                df_prices_platinum = df_prices_platinum.dropna(subset=["Date", "Price"]).sort_values("Date").reset_index(drop=True)[
                    ["Date", "Price"]
                ]
                df_adj_platinum, used_date_platinum = adjust_prices_for_mode(df_prices_platinum, pmode)
            except Exception:
                used_date_platinum = None

            v2_bytes_platinum, v2_date_platinum = create_technical_analysis_v2_chart(
                excel_path_for_ppt,
                ticker="XPT Comdty",
                price_mode=pmode,
                dmas_score=int(platinum_dmas),
                dmas_prev_week=int(platinum_dmas_prev),
                technical_score=platinum_tech,
                technical_prev_week=platinum_tech_prev,
                momentum_score=platinum_momentum,
                momentum_prev_week=platinum_mom_prev,
                rsi_prev_week=platinum_rsi_prev,
                days_gap=platinum_days_gap,
                previous_date=platinum_prev_date,
            )
            # Get the view and subtitle
            v2_view_text_platinum = st.session_state.get("platinum_selected_view")
            # Prepend commodity name if not already present
            if v2_view_text_platinum and not v2_view_text_platinum.lower().startswith("platinum"):
                v2_view_text_platinum = f"Platinum: {v2_view_text_platinum}"
            v2_subtitle_platinum = st.session_state.get("platinum_subtitle", "")

            prs = insert_technical_analysis_v2_slide(
                prs,
                v2_bytes_platinum,
                used_date=used_date_platinum,
                price_mode=pmode,
                placeholder_name="platinum_v2",
                view_text=v2_view_text_platinum,
                subtitle_text=v2_subtitle_platinum,
            )
        except Exception as e:
            print(f"[Tech V2] Platinum v2 chart error: {e}")
            import traceback
            traceback.print_exc()

        # ------------------------------------------------------------------
        # Insert Palladium Technical Analysis v2 chart (Chart.js + Playwright)
        # ------------------------------------------------------------------
        try:
            update_progress("Processing Palladium Technical Analysis...")
            # Get DMAS scores from session state
            palladium_dmas = st.session_state.get("palladium_dmas", 50)
            palladium_dmas_prev = st.session_state.get("palladium_last_week_avg", palladium_dmas)
            palladium_tech = _get_palladium_technical_score(excel_path_for_ppt)
            palladium_momentum = _get_palladium_momentum_score(excel_path_for_ppt)
            print(f"[Tech V2] Palladium DMAS: {palladium_dmas}, Prev Week: {palladium_dmas_prev}, Tech: {palladium_tech}, Mom: {palladium_momentum}")

            # Get previous week Technical/Momentum/RSI scores from history
            palladium_tech_prev = st.session_state.get("palladium_last_week_tech", None)
            palladium_mom_prev = st.session_state.get("palladium_last_week_mom", None)
            palladium_rsi_prev = st.session_state.get("palladium_last_week_rsi", None)
            print(f"[Tech V2] Palladium Prev week scores - Tech: {palladium_tech_prev}, Mom: {palladium_mom_prev}, RSI: {palladium_rsi_prev}")

            # Get gap information for change text formatting
            palladium_days_gap = st.session_state.get("palladium_prev_days_gap", None)
            palladium_prev_date = st.session_state.get("palladium_prev_date", None)

            # Compute used date for Palladium source footnote
            try:
                import pandas as pd
                df_prices_palladium = pd.read_excel(excel_path_for_ppt, sheet_name="data_prices")
                df_prices_palladium = df_prices_palladium.drop(index=0)
                df_prices_palladium = df_prices_palladium[df_prices_palladium[df_prices_palladium.columns[0]] != "DATES"]
                df_prices_palladium["Date"] = pd.to_datetime(df_prices_palladium[df_prices_palladium.columns[0]], errors="coerce")
                # Filter by "Data As Of" date if set
                if "data_as_of" in st.session_state:
                    df_prices_palladium = df_prices_palladium[df_prices_palladium["Date"] <= pd.Timestamp(st.session_state["data_as_of"])]
                df_prices_palladium["Price"] = pd.to_numeric(df_prices_palladium["XPD Curncy"], errors="coerce")
                df_prices_palladium = df_prices_palladium.dropna(subset=["Date", "Price"]).sort_values("Date").reset_index(drop=True)[
                    ["Date", "Price"]
                ]
                df_adj_palladium, used_date_palladium = adjust_prices_for_mode(df_prices_palladium, pmode)
            except Exception:
                used_date_palladium = None

            v2_bytes_palladium, v2_date_palladium = create_technical_analysis_v2_chart(
                excel_path_for_ppt,
                ticker="XPD Curncy",
                price_mode=pmode,
                dmas_score=int(palladium_dmas),
                dmas_prev_week=int(palladium_dmas_prev),
                technical_score=palladium_tech,
                technical_prev_week=palladium_tech_prev,
                momentum_score=palladium_momentum,
                momentum_prev_week=palladium_mom_prev,
                rsi_prev_week=palladium_rsi_prev,
                days_gap=palladium_days_gap,
                previous_date=palladium_prev_date,
            )
            # Get the view and subtitle
            v2_view_text_palladium = st.session_state.get("palladium_selected_view")
            # Prepend commodity name if not already present
            if v2_view_text_palladium and not v2_view_text_palladium.lower().startswith("palladium"):
                v2_view_text_palladium = f"Palladium: {v2_view_text_palladium}"
            v2_subtitle_palladium = st.session_state.get("palladium_subtitle", "")

            prs = insert_technical_analysis_v2_slide(
                prs,
                v2_bytes_palladium,
                used_date=used_date_palladium,
                price_mode=pmode,
                placeholder_name="palladium_v2",
                view_text=v2_view_text_palladium,
                subtitle_text=v2_subtitle_palladium,
            )
        except Exception as e:
            print(f"[Tech V2] Palladium v2 chart error: {e}")
            import traceback
            traceback.print_exc()

        # ------------------------------------------------------------------
        # Insert Oil Technical Analysis v2 chart (Chart.js + Playwright)
        # ------------------------------------------------------------------
        try:
            update_progress("Processing Oil Technical Analysis...")
            # Get DMAS scores from session state
            oil_dmas = st.session_state.get("oil_dmas", 50)
            oil_dmas_prev = st.session_state.get("oil_last_week_avg", oil_dmas)
            oil_tech = _get_oil_technical_score(excel_path_for_ppt)
            oil_momentum = _get_oil_momentum_score(excel_path_for_ppt)
            print(f"[Tech V2] Oil DMAS: {oil_dmas}, Prev Week: {oil_dmas_prev}, Tech: {oil_tech}, Mom: {oil_momentum}")

            # Get previous week Technical/Momentum/RSI scores from history
            oil_tech_prev = st.session_state.get("oil_last_week_tech", None)
            oil_mom_prev = st.session_state.get("oil_last_week_mom", None)
            oil_rsi_prev = st.session_state.get("oil_last_week_rsi", None)
            print(f"[Tech V2] Oil Prev week scores - Tech: {oil_tech_prev}, Mom: {oil_mom_prev}, RSI: {oil_rsi_prev}")

            # Get gap information for change text formatting
            oil_days_gap = st.session_state.get("oil_prev_days_gap", None)
            oil_prev_date = st.session_state.get("oil_prev_date", None)

            # Compute used date for Oil source footnote
            try:
                import pandas as pd
                df_prices_oil = pd.read_excel(excel_path_for_ppt, sheet_name="data_prices")
                df_prices_oil = df_prices_oil.drop(index=0)
                df_prices_oil = df_prices_oil[df_prices_oil[df_prices_oil.columns[0]] != "DATES"]
                df_prices_oil["Date"] = pd.to_datetime(df_prices_oil[df_prices_oil.columns[0]], errors="coerce")
                # Filter by "Data As Of" date if set
                if "data_as_of" in st.session_state:
                    df_prices_oil = df_prices_oil[df_prices_oil["Date"] <= pd.Timestamp(st.session_state["data_as_of"])]
                df_prices_oil["Price"] = pd.to_numeric(df_prices_oil["CL1 Comdty"], errors="coerce")
                df_prices_oil = df_prices_oil.dropna(subset=["Date", "Price"]).sort_values("Date").reset_index(drop=True)[
                    ["Date", "Price"]
                ]
                df_adj_oil, used_date_oil = adjust_prices_for_mode(df_prices_oil, pmode)
            except Exception:
                used_date_oil = None

            v2_bytes_oil, v2_date_oil = create_technical_analysis_v2_chart(
                excel_path_for_ppt,
                ticker="CL1 Comdty",
                price_mode=pmode,
                dmas_score=int(oil_dmas),
                dmas_prev_week=int(oil_dmas_prev),
                technical_score=oil_tech,
                technical_prev_week=oil_tech_prev,
                momentum_score=oil_momentum,
                momentum_prev_week=oil_mom_prev,
                rsi_prev_week=oil_rsi_prev,
                days_gap=oil_days_gap,
                previous_date=oil_prev_date,
            )
            # Get the view and subtitle
            v2_view_text_oil = st.session_state.get("oil_selected_view")
            # Prepend commodity name if not already present
            if v2_view_text_oil and not v2_view_text_oil.lower().startswith("oil"):
                v2_view_text_oil = f"Oil: {v2_view_text_oil}"
            v2_subtitle_oil = st.session_state.get("oil_subtitle", "")

            prs = insert_technical_analysis_v2_slide(
                prs,
                v2_bytes_oil,
                used_date=used_date_oil,
                price_mode=pmode,
                placeholder_name="oil_v2",
                view_text=v2_view_text_oil,
                subtitle_text=v2_subtitle_oil,
            )
        except Exception as e:
            print(f"[Tech V2] Oil v2 chart error: {e}")
            import traceback
            traceback.print_exc()

        # ------------------------------------------------------------------
        # Insert Copper Technical Analysis v2 chart (Chart.js + Playwright)
        # ------------------------------------------------------------------
        try:
            update_progress("Processing Copper Technical Analysis...")
            # Get DMAS scores from session state
            copper_dmas = st.session_state.get("copper_dmas", 50)
            copper_dmas_prev = st.session_state.get("copper_last_week_avg", copper_dmas)
            copper_tech = _get_copper_technical_score(excel_path_for_ppt)
            copper_momentum = _get_copper_momentum_score(excel_path_for_ppt)
            print(f"[Tech V2] Copper DMAS: {copper_dmas}, Prev Week: {copper_dmas_prev}, Tech: {copper_tech}, Mom: {copper_momentum}")

            # Get previous week Technical/Momentum/RSI scores from history
            copper_tech_prev = st.session_state.get("copper_last_week_tech", None)
            copper_mom_prev = st.session_state.get("copper_last_week_mom", None)
            copper_rsi_prev = st.session_state.get("copper_last_week_rsi", None)
            print(f"[Tech V2] Copper Prev week scores - Tech: {copper_tech_prev}, Mom: {copper_mom_prev}, RSI: {copper_rsi_prev}")

            # Get gap information for change text formatting
            copper_days_gap = st.session_state.get("copper_prev_days_gap", None)
            copper_prev_date = st.session_state.get("copper_prev_date", None)

            # Compute used date for Copper source footnote
            try:
                import pandas as pd
                df_prices_copper = pd.read_excel(excel_path_for_ppt, sheet_name="data_prices")
                df_prices_copper = df_prices_copper.drop(index=0)
                df_prices_copper = df_prices_copper[df_prices_copper[df_prices_copper.columns[0]] != "DATES"]
                df_prices_copper["Date"] = pd.to_datetime(df_prices_copper[df_prices_copper.columns[0]], errors="coerce")
                # Filter by "Data As Of" date if set
                if "data_as_of" in st.session_state:
                    df_prices_copper = df_prices_copper[df_prices_copper["Date"] <= pd.Timestamp(st.session_state["data_as_of"])]
                df_prices_copper["Price"] = pd.to_numeric(df_prices_copper["LP1 Comdty"], errors="coerce")
                df_prices_copper = df_prices_copper.dropna(subset=["Date", "Price"]).sort_values("Date").reset_index(drop=True)[
                    ["Date", "Price"]
                ]
                df_adj_copper, used_date_copper = adjust_prices_for_mode(df_prices_copper, pmode)
            except Exception:
                used_date_copper = None

            v2_bytes_copper, v2_date_copper = create_technical_analysis_v2_chart(
                excel_path_for_ppt,
                ticker="LP1 Comdty",
                price_mode=pmode,
                dmas_score=int(copper_dmas),
                dmas_prev_week=int(copper_dmas_prev),
                technical_score=copper_tech,
                technical_prev_week=copper_tech_prev,
                momentum_score=copper_momentum,
                momentum_prev_week=copper_mom_prev,
                rsi_prev_week=copper_rsi_prev,
                days_gap=copper_days_gap,
                previous_date=copper_prev_date,
            )
            # Get the view and subtitle
            v2_view_text_copper = st.session_state.get("copper_selected_view")
            # Prepend commodity name if not already present
            if v2_view_text_copper and not v2_view_text_copper.lower().startswith("copper"):
                v2_view_text_copper = f"Copper: {v2_view_text_copper}"
            v2_subtitle_copper = st.session_state.get("copper_subtitle", "")

            prs = insert_technical_analysis_v2_slide(
                prs,
                v2_bytes_copper,
                used_date=used_date_copper,
                price_mode=pmode,
                placeholder_name="copper_v2",
                view_text=v2_view_text_copper,
                subtitle_text=v2_subtitle_copper,
            )
        except Exception as e:
            print(f"[Tech V2] Copper v2 chart error: {e}")
            import traceback
            traceback.print_exc()

        # ------------------------------------------------------------------
        # Insert Bitcoin Technical Analysis v2 chart (Chart.js + Playwright)
        # ------------------------------------------------------------------
        try:
            update_progress("Processing Bitcoin Technical Analysis...")
            # Get DMAS scores from session state
            bitcoin_dmas = st.session_state.get("bitcoin_dmas", 50)
            bitcoin_dmas_prev = st.session_state.get("bitcoin_last_week_avg", bitcoin_dmas)
            bitcoin_tech = _get_bitcoin_technical_score(excel_path_for_ppt)
            bitcoin_momentum = _get_bitcoin_momentum_score(excel_path_for_ppt)
            print(f"[Tech V2] Bitcoin DMAS: {bitcoin_dmas}, Prev Week: {bitcoin_dmas_prev}, Tech: {bitcoin_tech}, Mom: {bitcoin_momentum}")

            # Get previous week Technical/Momentum/RSI scores from history
            bitcoin_tech_prev = st.session_state.get("bitcoin_last_week_tech", None)
            bitcoin_mom_prev = st.session_state.get("bitcoin_last_week_mom", None)
            bitcoin_rsi_prev = st.session_state.get("bitcoin_last_week_rsi", None)
            print(f"[Tech V2] Bitcoin Prev week scores - Tech: {bitcoin_tech_prev}, Mom: {bitcoin_mom_prev}, RSI: {bitcoin_rsi_prev}")

            # Get gap information for change text formatting
            bitcoin_days_gap = st.session_state.get("bitcoin_prev_days_gap", None)
            bitcoin_prev_date = st.session_state.get("bitcoin_prev_date", None)

            # Compute used date for Bitcoin source footnote
            try:
                import pandas as pd
                df_prices_bitcoin = pd.read_excel(excel_path_for_ppt, sheet_name="data_prices")
                df_prices_bitcoin = df_prices_bitcoin.drop(index=0)
                df_prices_bitcoin = df_prices_bitcoin[df_prices_bitcoin[df_prices_bitcoin.columns[0]] != "DATES"]
                df_prices_bitcoin["Date"] = pd.to_datetime(df_prices_bitcoin[df_prices_bitcoin.columns[0]], errors="coerce")
                # Filter by "Data As Of" date if set
                if "data_as_of" in st.session_state:
                    df_prices_bitcoin = df_prices_bitcoin[df_prices_bitcoin["Date"] <= pd.Timestamp(st.session_state["data_as_of"])]
                df_prices_bitcoin["Price"] = pd.to_numeric(df_prices_bitcoin["XBTUSD Curncy"], errors="coerce")
                df_prices_bitcoin = df_prices_bitcoin.dropna(subset=["Date", "Price"]).sort_values("Date").reset_index(drop=True)[
                    ["Date", "Price"]
                ]
                df_adj_bitcoin, used_date_bitcoin = adjust_prices_for_mode(df_prices_bitcoin, pmode)
            except Exception:
                used_date_bitcoin = None

            v2_bytes_bitcoin, v2_date_bitcoin = create_technical_analysis_v2_chart(
                excel_path_for_ppt,
                ticker="XBTUSD Curncy",
                price_mode=pmode,
                dmas_score=int(bitcoin_dmas),
                dmas_prev_week=int(bitcoin_dmas_prev),
                technical_score=bitcoin_tech,
                technical_prev_week=bitcoin_tech_prev,
                momentum_score=bitcoin_momentum,
                momentum_prev_week=bitcoin_mom_prev,
                rsi_prev_week=bitcoin_rsi_prev,
                days_gap=bitcoin_days_gap,
                previous_date=bitcoin_prev_date,
            )
            # Get the view and subtitle
            v2_view_text_bitcoin = st.session_state.get("bitcoin_selected_view")
            # Prepend crypto name if not already present
            if v2_view_text_bitcoin and not v2_view_text_bitcoin.lower().startswith("bitcoin"):
                v2_view_text_bitcoin = f"Bitcoin: {v2_view_text_bitcoin}"
            v2_subtitle_bitcoin = st.session_state.get("bitcoin_subtitle", "")

            prs = insert_technical_analysis_v2_slide(
                prs,
                v2_bytes_bitcoin,
                used_date=used_date_bitcoin,
                price_mode=pmode,
                placeholder_name="bitcoin_v2",
                view_text=v2_view_text_bitcoin,
                subtitle_text=v2_subtitle_bitcoin,
            )
        except Exception as e:
            print(f"[Tech V2] Bitcoin v2 chart error: {e}")
            import traceback
            traceback.print_exc()

        # ------------------------------------------------------------------
        # Insert Ethereum Technical Analysis v2 chart (Chart.js + Playwright)
        # ------------------------------------------------------------------
        try:
            update_progress("Processing Ethereum Technical Analysis...")
            # Get DMAS scores from session state
            ethereum_dmas = st.session_state.get("ethereum_dmas", 50)
            ethereum_dmas_prev = st.session_state.get("ethereum_last_week_avg", ethereum_dmas)
            ethereum_tech = _get_ethereum_technical_score(excel_path_for_ppt)
            ethereum_momentum = _get_ethereum_momentum_score(excel_path_for_ppt)
            print(f"[Tech V2] Ethereum DMAS: {ethereum_dmas}, Prev Week: {ethereum_dmas_prev}, Tech: {ethereum_tech}, Mom: {ethereum_momentum}")

            # Get previous week Technical/Momentum/RSI scores from history
            ethereum_tech_prev = st.session_state.get("ethereum_last_week_tech", None)
            ethereum_mom_prev = st.session_state.get("ethereum_last_week_mom", None)
            ethereum_rsi_prev = st.session_state.get("ethereum_last_week_rsi", None)
            print(f"[Tech V2] Ethereum Prev week scores - Tech: {ethereum_tech_prev}, Mom: {ethereum_mom_prev}, RSI: {ethereum_rsi_prev}")

            # Get gap information for change text formatting
            ethereum_days_gap = st.session_state.get("ethereum_prev_days_gap", None)
            ethereum_prev_date = st.session_state.get("ethereum_prev_date", None)

            # Compute used date for Ethereum source footnote
            try:
                import pandas as pd
                df_prices_ethereum = pd.read_excel(excel_path_for_ppt, sheet_name="data_prices")
                df_prices_ethereum = df_prices_ethereum.drop(index=0)
                df_prices_ethereum = df_prices_ethereum[df_prices_ethereum[df_prices_ethereum.columns[0]] != "DATES"]
                df_prices_ethereum["Date"] = pd.to_datetime(df_prices_ethereum[df_prices_ethereum.columns[0]], errors="coerce")
                # Filter by "Data As Of" date if set
                if "data_as_of" in st.session_state:
                    df_prices_ethereum = df_prices_ethereum[df_prices_ethereum["Date"] <= pd.Timestamp(st.session_state["data_as_of"])]
                df_prices_ethereum["Price"] = pd.to_numeric(df_prices_ethereum["XETUSD Curncy"], errors="coerce")
                df_prices_ethereum = df_prices_ethereum.dropna(subset=["Date", "Price"]).sort_values("Date").reset_index(drop=True)[
                    ["Date", "Price"]
                ]
                df_adj_ethereum, used_date_ethereum = adjust_prices_for_mode(df_prices_ethereum, pmode)
            except Exception:
                used_date_ethereum = None

            v2_bytes_ethereum, v2_date_ethereum = create_technical_analysis_v2_chart(
                excel_path_for_ppt,
                ticker="XETUSD Curncy",
                price_mode=pmode,
                dmas_score=int(ethereum_dmas),
                dmas_prev_week=int(ethereum_dmas_prev),
                technical_score=ethereum_tech,
                technical_prev_week=ethereum_tech_prev,
                momentum_score=ethereum_momentum,
                momentum_prev_week=ethereum_mom_prev,
                rsi_prev_week=ethereum_rsi_prev,
                days_gap=ethereum_days_gap,
                previous_date=ethereum_prev_date,
            )
            # Get the view and subtitle
            v2_view_text_ethereum = st.session_state.get("ethereum_selected_view")
            # Prepend crypto name if not already present
            if v2_view_text_ethereum and not v2_view_text_ethereum.lower().startswith("ethereum"):
                v2_view_text_ethereum = f"Ethereum: {v2_view_text_ethereum}"
            v2_subtitle_ethereum = st.session_state.get("ethereum_subtitle", "")

            prs = insert_technical_analysis_v2_slide(
                prs,
                v2_bytes_ethereum,
                used_date=used_date_ethereum,
                price_mode=pmode,
                placeholder_name="ethereum_v2",
                view_text=v2_view_text_ethereum,
                subtitle_text=v2_subtitle_ethereum,
            )
        except Exception as e:
            print(f"[Tech V2] Ethereum v2 chart error: {e}")
            import traceback
            traceback.print_exc()

        # ------------------------------------------------------------------
        # Insert Ripple Technical Analysis v2 chart (Chart.js + Playwright)
        # ------------------------------------------------------------------
        try:
            update_progress("Processing Ripple Technical Analysis...")
            # Get DMAS scores from session state
            ripple_dmas = st.session_state.get("ripple_dmas", 50)
            ripple_dmas_prev = st.session_state.get("ripple_last_week_avg", ripple_dmas)
            ripple_tech = _get_ripple_technical_score(excel_path_for_ppt)
            ripple_momentum = _get_ripple_momentum_score(excel_path_for_ppt)
            print(f"[Tech V2] Ripple DMAS: {ripple_dmas}, Prev Week: {ripple_dmas_prev}, Tech: {ripple_tech}, Mom: {ripple_momentum}")

            # Get previous week Technical/Momentum/RSI scores from history
            ripple_tech_prev = st.session_state.get("ripple_last_week_tech", None)
            ripple_mom_prev = st.session_state.get("ripple_last_week_mom", None)
            ripple_rsi_prev = st.session_state.get("ripple_last_week_rsi", None)
            print(f"[Tech V2] Ripple Prev week scores - Tech: {ripple_tech_prev}, Mom: {ripple_mom_prev}, RSI: {ripple_rsi_prev}")

            # Get gap information for change text formatting
            ripple_days_gap = st.session_state.get("ripple_prev_days_gap", None)
            ripple_prev_date = st.session_state.get("ripple_prev_date", None)

            # Compute used date for Ripple source footnote
            try:
                import pandas as pd
                df_prices_ripple = pd.read_excel(excel_path_for_ppt, sheet_name="data_prices")
                df_prices_ripple = df_prices_ripple.drop(index=0)
                df_prices_ripple = df_prices_ripple[df_prices_ripple[df_prices_ripple.columns[0]] != "DATES"]
                df_prices_ripple["Date"] = pd.to_datetime(df_prices_ripple[df_prices_ripple.columns[0]], errors="coerce")
                # Filter by "Data As Of" date if set
                if "data_as_of" in st.session_state:
                    df_prices_ripple = df_prices_ripple[df_prices_ripple["Date"] <= pd.Timestamp(st.session_state["data_as_of"])]
                df_prices_ripple["Price"] = pd.to_numeric(df_prices_ripple["XRPUSD Curncy"], errors="coerce")
                df_prices_ripple = df_prices_ripple.dropna(subset=["Date", "Price"]).sort_values("Date").reset_index(drop=True)[
                    ["Date", "Price"]
                ]
                df_adj_ripple, used_date_ripple = adjust_prices_for_mode(df_prices_ripple, pmode)
            except Exception:
                used_date_ripple = None

            v2_bytes_ripple, v2_date_ripple = create_technical_analysis_v2_chart(
                excel_path_for_ppt,
                ticker="XRPUSD Curncy",
                price_mode=pmode,
                dmas_score=int(ripple_dmas),
                dmas_prev_week=int(ripple_dmas_prev),
                technical_score=ripple_tech,
                technical_prev_week=ripple_tech_prev,
                momentum_score=ripple_momentum,
                momentum_prev_week=ripple_mom_prev,
                rsi_prev_week=ripple_rsi_prev,
                days_gap=ripple_days_gap,
                previous_date=ripple_prev_date,
            )
            # Get the view and subtitle
            v2_view_text_ripple = st.session_state.get("ripple_selected_view")
            # Prepend crypto name if not already present
            if v2_view_text_ripple and not v2_view_text_ripple.lower().startswith("ripple"):
                v2_view_text_ripple = f"Ripple: {v2_view_text_ripple}"
            v2_subtitle_ripple = st.session_state.get("ripple_subtitle", "")

            prs = insert_technical_analysis_v2_slide(
                prs,
                v2_bytes_ripple,
                used_date=used_date_ripple,
                price_mode=pmode,
                placeholder_name="ripple_v2",
                view_text=v2_view_text_ripple,
                subtitle_text=v2_subtitle_ripple,
            )
        except Exception as e:
            print(f"[Tech V2] Ripple v2 chart error: {e}")
            import traceback
            traceback.print_exc()

        # ------------------------------------------------------------------
        # Insert Solana Technical Analysis v2 chart (Chart.js + Playwright)
        # ------------------------------------------------------------------
        try:
            update_progress("Processing Solana Technical Analysis...")
            # Get DMAS scores from session state
            solana_dmas = st.session_state.get("solana_dmas", 50)
            solana_dmas_prev = st.session_state.get("solana_last_week_avg", solana_dmas)
            solana_tech = _get_solana_technical_score(excel_path_for_ppt)
            solana_momentum = _get_solana_momentum_score(excel_path_for_ppt)
            print(f"[Tech V2] Solana DMAS: {solana_dmas}, Prev Week: {solana_dmas_prev}, Tech: {solana_tech}, Mom: {solana_momentum}")

            # Get previous week Technical/Momentum/RSI scores from history
            solana_tech_prev = st.session_state.get("solana_last_week_tech", None)
            solana_mom_prev = st.session_state.get("solana_last_week_mom", None)
            solana_rsi_prev = st.session_state.get("solana_last_week_rsi", None)
            print(f"[Tech V2] Solana Prev week scores - Tech: {solana_tech_prev}, Mom: {solana_mom_prev}, RSI: {solana_rsi_prev}")

            # Get gap information for change text formatting
            solana_days_gap = st.session_state.get("solana_prev_days_gap", None)
            solana_prev_date = st.session_state.get("solana_prev_date", None)

            # Compute used date for Solana source footnote
            try:
                import pandas as pd
                df_prices_solana = pd.read_excel(excel_path_for_ppt, sheet_name="data_prices")
                df_prices_solana = df_prices_solana.drop(index=0)
                df_prices_solana = df_prices_solana[df_prices_solana[df_prices_solana.columns[0]] != "DATES"]
                df_prices_solana["Date"] = pd.to_datetime(df_prices_solana[df_prices_solana.columns[0]], errors="coerce")
                # Filter by "Data As Of" date if set
                if "data_as_of" in st.session_state:
                    df_prices_solana = df_prices_solana[df_prices_solana["Date"] <= pd.Timestamp(st.session_state["data_as_of"])]
                df_prices_solana["Price"] = pd.to_numeric(df_prices_solana["XSOUSD Curncy"], errors="coerce")
                df_prices_solana = df_prices_solana.dropna(subset=["Date", "Price"]).sort_values("Date").reset_index(drop=True)[
                    ["Date", "Price"]
                ]
                df_adj_solana, used_date_solana = adjust_prices_for_mode(df_prices_solana, pmode)
            except Exception:
                used_date_solana = None

            v2_bytes_solana, v2_date_solana = create_technical_analysis_v2_chart(
                excel_path_for_ppt,
                ticker="XSOUSD Curncy",
                price_mode=pmode,
                dmas_score=int(solana_dmas),
                dmas_prev_week=int(solana_dmas_prev),
                technical_score=solana_tech,
                technical_prev_week=solana_tech_prev,
                momentum_score=solana_momentum,
                momentum_prev_week=solana_mom_prev,
                rsi_prev_week=solana_rsi_prev,
                days_gap=solana_days_gap,
                previous_date=solana_prev_date,
            )
            # Get the view and subtitle
            v2_view_text_solana = st.session_state.get("solana_selected_view")
            # Prepend crypto name if not already present
            if v2_view_text_solana and not v2_view_text_solana.lower().startswith("solana"):
                v2_view_text_solana = f"Solana: {v2_view_text_solana}"
            v2_subtitle_solana = st.session_state.get("solana_subtitle", "")

            prs = insert_technical_analysis_v2_slide(
                prs,
                v2_bytes_solana,
                used_date=used_date_solana,
                price_mode=pmode,
                placeholder_name="solana_v2",
                view_text=v2_view_text_solana,
                subtitle_text=v2_subtitle_solana,
            )
        except Exception as e:
            print(f"[Tech V2] Solana v2 chart error: {e}")
            import traceback
            traceback.print_exc()

        # ------------------------------------------------------------------
        # Insert Binance Technical Analysis v2 chart (Chart.js + Playwright)
        # ------------------------------------------------------------------
        try:
            update_progress("Processing Binance Technical Analysis...")
            # Get DMAS scores from session state
            binance_dmas = st.session_state.get("binance_dmas", 50)
            binance_dmas_prev = st.session_state.get("binance_last_week_avg", binance_dmas)
            binance_tech = _get_binance_technical_score(excel_path_for_ppt)
            binance_momentum = _get_binance_momentum_score(excel_path_for_ppt)
            print(f"[Tech V2] Binance DMAS: {binance_dmas}, Prev Week: {binance_dmas_prev}, Tech: {binance_tech}, Mom: {binance_momentum}")

            # Get previous week Technical/Momentum/RSI scores from history
            binance_tech_prev = st.session_state.get("binance_last_week_tech", None)
            binance_mom_prev = st.session_state.get("binance_last_week_mom", None)
            binance_rsi_prev = st.session_state.get("binance_last_week_rsi", None)
            print(f"[Tech V2] Binance Prev week scores - Tech: {binance_tech_prev}, Mom: {binance_mom_prev}, RSI: {binance_rsi_prev}")

            # Get gap information for change text formatting
            binance_days_gap = st.session_state.get("binance_prev_days_gap", None)
            binance_prev_date = st.session_state.get("binance_prev_date", None)

            # Compute used date for Binance source footnote
            try:
                import pandas as pd
                df_prices_binance = pd.read_excel(excel_path_for_ppt, sheet_name="data_prices")
                df_prices_binance = df_prices_binance.drop(index=0)
                df_prices_binance = df_prices_binance[df_prices_binance[df_prices_binance.columns[0]] != "DATES"]
                df_prices_binance["Date"] = pd.to_datetime(df_prices_binance[df_prices_binance.columns[0]], errors="coerce")
                # Filter by "Data As Of" date if set
                if "data_as_of" in st.session_state:
                    df_prices_binance = df_prices_binance[df_prices_binance["Date"] <= pd.Timestamp(st.session_state["data_as_of"])]
                df_prices_binance["Price"] = pd.to_numeric(df_prices_binance["XBIUSD Curncy"], errors="coerce")
                df_prices_binance = df_prices_binance.dropna(subset=["Date", "Price"]).sort_values("Date").reset_index(drop=True)[
                    ["Date", "Price"]
                ]
                df_adj_binance, used_date_binance = adjust_prices_for_mode(df_prices_binance, pmode)
            except Exception:
                used_date_binance = None

            v2_bytes_binance, v2_date_binance = create_technical_analysis_v2_chart(
                excel_path_for_ppt,
                ticker="XBIUSD Curncy",
                price_mode=pmode,
                dmas_score=int(binance_dmas),
                dmas_prev_week=int(binance_dmas_prev),
                technical_score=binance_tech,
                technical_prev_week=binance_tech_prev,
                momentum_score=binance_momentum,
                momentum_prev_week=binance_mom_prev,
                rsi_prev_week=binance_rsi_prev,
                days_gap=binance_days_gap,
                previous_date=binance_prev_date,
            )
            # Get the view and subtitle
            v2_view_text_binance = st.session_state.get("binance_selected_view")
            # Prepend crypto name if not already present
            if v2_view_text_binance and not v2_view_text_binance.lower().startswith("binance"):
                v2_view_text_binance = f"Binance: {v2_view_text_binance}"
            v2_subtitle_binance = st.session_state.get("binance_subtitle", "")

            prs = insert_technical_analysis_v2_slide(
                prs,
                v2_bytes_binance,
                used_date=used_date_binance,
                price_mode=pmode,
                placeholder_name="binance_v2",
                view_text=v2_view_text_binance,
                subtitle_text=v2_subtitle_binance,
            )
        except Exception as e:
            print(f"[Tech V2] Binance v2 chart error: {e}")
            import traceback
            traceback.print_exc()

        # When CSI 300 is the selected index, the technical analysis slides
        # for CSI have already been inserted in the branch above.  Avoid
        # inserting CSI slides again here.  Likewise, when SPX is selected,
        # CSI slides are not inserted at all.  This prevents duplicate
        # insertion of CSI slides that could override SPX content or leave
        # placeholders empty.

        # ------------------------------------------------------------------
        # Insert Equity performance charts
        # ------------------------------------------------------------------
        try:
            # Generate the weekly performance bar chart with price-mode adjustment
            bar_bytes, perf_used_date = create_weekly_performance_chart(
                excel_path_for_ppt,
                price_mode=st.session_state.get("price_mode", "Last Price"),
            )
            prs = insert_equity_performance_bar_slide(
                prs,
                bar_bytes,
                used_date=perf_used_date,
                price_mode=st.session_state.get("price_mode", "Last Price"),
                left_cm=3.47,
                top_cm=5.28,
                width_cm=17.31,
                height_cm=10,
            )
            # Generate the historical performance heatmap with price-mode adjustment
            histo_bytes, histo_used_date = create_historical_performance_table(
                excel_path_for_ppt,
                price_mode=st.session_state.get("price_mode", "Last Price"),
            )
            prs = insert_equity_performance_histo_slide(
                prs,
                histo_bytes,
                used_date=histo_used_date,
                price_mode=st.session_state.get("price_mode", "Last Price"),
                left_cm=2.16,
                top_cm=4.70,
                width_cm=19.43,
                height_cm=10.61,
            )

            # Generate the Equity YTD Evolution chart (Chart.js line chart)
            ytd_evo_bytes, ytd_evo_date = create_equity_ytd_evolution_chart(
                excel_path_for_ppt,
                price_mode=st.session_state.get("price_mode", "Last Price"),
            )
            prs = insert_equity_ytd_evolution_slide(
                prs,
                ytd_evo_bytes,
                used_date=ytd_evo_date,
                price_mode=st.session_state.get("price_mode", "Last Price"),
                subtitle=st.session_state.get("eq_subtitle"),
            )

            # Generate the FX Impact Analysis chart for EUR investors
            fx_impact_bytes, fx_impact_date = create_fx_impact_analysis_chart_eur(
                excel_path_for_ppt,
                price_mode=st.session_state.get("price_mode", "Last Price"),
            )
            prs = insert_fx_impact_analysis_slide_eur(
                prs,
                fx_impact_bytes,
                used_date=fx_impact_date,
                price_mode=st.session_state.get("price_mode", "Last Price"),
            )
        except Exception as e:
            print(f"Equity performance charts error: {e}")

        # ------------------------------------------------------------------
        # Insert FX performance charts
        # ------------------------------------------------------------------
        try:
            update_progress("Processing FX performance charts...")
            # Generate the weekly FX performance bar chart (HTML-based)
            fx_bar_bytes, fx_used_date = create_weekly_fx_html_chart(
                excel_path_for_ppt,
                price_mode=st.session_state.get("price_mode", "Last Price"),
            )
            prs = insert_fx_weekly_html_slide(
                prs,
                fx_bar_bytes,
                used_date=fx_used_date,
                price_mode=st.session_state.get("price_mode", "Last Price"),
            )

            # Generate the FX historical performance heatmap with price-mode adjustment (HTML-based)
            fx_histo_bytes, fx_used_date2 = create_historical_fx_html_chart(
                excel_path_for_ppt,
                price_mode=st.session_state.get("price_mode", "Last Price"),
            )
            prs = insert_fx_historical_html_slide(
                prs,
                fx_histo_bytes,
                used_date=fx_used_date2,
                price_mode=st.session_state.get("price_mode", "Last Price"),
            )
        except Exception as e:
            print(f"FX performance charts error: {e}")

        # ------------------------------------------------------------------
        # Insert cryptocurrency performance charts
        # ------------------------------------------------------------------
        try:
            update_progress("Processing Crypto performance charts...")
            # Generate the weekly crypto performance bar chart (HTML-based)
            crypto_bar_bytes, crypto_used_date = create_weekly_crypto_html_chart(
                excel_path_for_ppt,
                price_mode=st.session_state.get("price_mode", "Last Price"),
            )
            prs = insert_crypto_weekly_html_slide(
                prs,
                crypto_bar_bytes,
                used_date=crypto_used_date,
                price_mode=st.session_state.get("price_mode", "Last Price"),
            )

            # Generate the cryptocurrency historical performance heatmap (HTML-based) with price-mode adjustment
            crypto_histo_bytes, crypto_used_date2 = create_historical_crypto_html_chart(
                excel_path_for_ppt,
                price_mode=st.session_state.get("price_mode", "Last Price"),
            )
            prs = insert_crypto_historical_html_slide(
                prs,
                crypto_histo_bytes,
                used_date=crypto_used_date2,
                price_mode=st.session_state.get("price_mode", "Last Price"),
            )

            # Generate the Crypto YTD Evolution chart (Chart.js line chart)
            crypto_ytd_evo_bytes, crypto_ytd_evo_date = create_crypto_ytd_evolution_chart(
                excel_path_for_ppt,
                price_mode=st.session_state.get("price_mode", "Last Price"),
            )
            prs = insert_crypto_ytd_evolution_slide(
                prs,
                crypto_ytd_evo_bytes,
                used_date=crypto_ytd_evo_date,
                price_mode=st.session_state.get("price_mode", "Last Price"),
                subtitle=st.session_state.get("cr_subtitle"),
            )
        except Exception as e:
            print(f"Crypto performance charts error: {e}")

        # ------------------------------------------------------------------
        # Insert Rates performance charts
        # ------------------------------------------------------------------
        try:
            update_progress("Processing Rates performance charts...")
            # Generate the weekly rates performance bar chart with price-mode adjustment
            rates_bar_bytes, rates_used_date = create_weekly_rates_performance_chart(
                excel_path_for_ppt,
                price_mode=st.session_state.get("price_mode", "Last Price"),
            )
            prs = insert_rates_performance_bar_slide(
                prs,
                rates_bar_bytes,
                used_date=rates_used_date,
                price_mode=st.session_state.get("price_mode", "Last Price"),
                left_cm=3.35,
                top_cm=4.6,
                width_cm=17.02,
            )

            # Generate the rates historical performance heatmap with price-mode adjustment
            rates_histo_bytes, rates_used_date2 = create_historical_rates_performance_table(
                excel_path_for_ppt,
                price_mode=st.session_state.get("price_mode", "Last Price"),
            )
            prs = insert_rates_performance_histo_slide(
                prs,
                rates_histo_bytes,
                used_date=rates_used_date2,
                price_mode=st.session_state.get("price_mode", "Last Price"),
                left_cm=3.35,
                top_cm=4.6,
                width_cm=17.02,
            )
        except Exception as e:
            print(f"Rates performance charts error: {e}")

        # ------------------------------------------------------------------
        # Insert Credit performance charts
        # ------------------------------------------------------------------
        try:
            update_progress("Processing Credit performance charts...")
            # Generate the weekly credit performance bar chart with price-mode adjustment
            credit_bar_bytes, credit_used_date = create_weekly_credit_performance_chart(
                excel_path_for_ppt,
                price_mode=st.session_state.get("price_mode", "Last Price"),
            )
            prs = insert_credit_performance_bar_slide(
                prs,
                credit_bar_bytes,
                used_date=credit_used_date,
                price_mode=st.session_state.get("price_mode", "Last Price"),
                left_cm=1.63,
                top_cm=4.73,
                width_cm=22.48,
                height_cm=10.61,
            )

            # Generate the credit historical performance heatmap with price-mode adjustment
            credit_histo_bytes, credit_used_date2 = create_historical_credit_performance_chart(
                excel_path_for_ppt,
                price_mode=st.session_state.get("price_mode", "Last Price"),
            )
            prs = insert_credit_performance_histo_slide(
                prs,
                credit_histo_bytes,
                used_date=credit_used_date2,
                price_mode=st.session_state.get("price_mode", "Last Price"),
            )
        except Exception as e:
            print(f"Credit performance charts error: {e}")

        # ------------------------------------------------------------------
        # Insert Commodity performance charts
        # ------------------------------------------------------------------
        try:
            update_progress("Processing Commodity performance charts...")
            # Generate the weekly commodity performance bar chart (HTML-based)
            commo_bar_bytes, commo_used_date = create_weekly_commodity_html_chart(
                excel_path_for_ppt,
                price_mode=st.session_state.get("price_mode", "Last Price"),
            )
            prs = insert_commodity_weekly_html_slide(
                prs,
                commo_bar_bytes,
                used_date=commo_used_date,
                price_mode=st.session_state.get("price_mode", "Last Price"),
            )

            # Generate the commodity historical performance heatmap (HTML-based)
            commo_histo_bytes, commo_used_date2 = create_historical_commodity_html_chart(
                excel_path_for_ppt,
                price_mode=st.session_state.get("price_mode", "Last Price"),
            )
            prs = insert_commodity_historical_html_slide(
                prs,
                commo_histo_bytes,
                used_date=commo_used_date2,
                price_mode=st.session_state.get("price_mode", "Last Price"),
            )

            # Generate the Commodity YTD Evolution chart (Chart.js line chart)
            commo_ytd_evo_bytes, commo_ytd_evo_date = create_commodity_ytd_evolution_chart(
                excel_path_for_ppt,
                price_mode=st.session_state.get("price_mode", "Last Price"),
            )
            prs = insert_commodity_ytd_evolution_slide(
                prs,
                commo_ytd_evo_bytes,
                used_date=commo_ytd_evo_date,
                price_mode=st.session_state.get("price_mode", "Last Price"),
                subtitle=st.session_state.get("co_subtitle"),
            )
        except Exception as e:
            print(f"Commodity performance charts error: {e}")

        # ------------------------------------------------------------------
        # Technical Analysis In A Nutshell slide
        # ------------------------------------------------------------------
        update_progress("Generating Technical Analysis slide...")
        try:
            from market_compass.technical_slide import (
                prepare_slide_data,
                insert_technical_analysis_slide,
            )
            import pandas as pd

            # Read price data for the slide
            df_prices_tech = pd.read_excel(excel_path_for_ppt, sheet_name="data_prices")
            df_prices_tech = df_prices_tech.drop(index=0)
            df_prices_tech = df_prices_tech[
                df_prices_tech[df_prices_tech.columns[0]] != "DATES"
            ]
            df_prices_tech["Date"] = pd.to_datetime(
                df_prices_tech[df_prices_tech.columns[0]], errors="coerce"
            )
            df_prices_tech = df_prices_tech.dropna(subset=["Date"]).sort_values(
                "Date"
            ).reset_index(drop=True)

            print(f"[Technical Nutshell] Price data loaded: {len(df_prices_tech)} rows, columns: {list(df_prices_tech.columns)[:10]}...")

            # Adjust for price mode
            df_prices_tech_adj, tech_used_date = adjust_prices_for_mode(
                df_prices_tech, pmode
            )

            # Collect DMAS scores from session state
            dmas_scores = {
                # Equity
                "spx": st.session_state.get("spx_dmas", 50),
                "csi": st.session_state.get("csi_dmas", 50),
                "nikkei": st.session_state.get("nikkei_dmas", 50),
                "tasi": st.session_state.get("tasi_dmas", 50),
                "sensex": st.session_state.get("sensex_dmas", 50),
                "dax": st.session_state.get("dax_dmas", 50),
                "smi": st.session_state.get("smi_dmas", 50),
                "ibov": st.session_state.get("ibov_dmas", 50),
                "mexbol": st.session_state.get("mexbol_dmas", 50),
                # Commodity
                "gold": st.session_state.get("gold_dmas", 50),
                "silver": st.session_state.get("silver_dmas", 50),
                "platinum": st.session_state.get("platinum_dmas", 50),
                "palladium": st.session_state.get("palladium_dmas", 50),
                "oil": st.session_state.get("oil_dmas", 50),
                "copper": st.session_state.get("copper_dmas", 50),
                # Crypto
                "bitcoin": st.session_state.get("bitcoin_dmas", 50),
                "ethereum": st.session_state.get("ethereum_dmas", 50),
                "ripple": st.session_state.get("ripple_dmas", 50),
                "solana": st.session_state.get("solana_dmas", 50),
                "binance": st.session_state.get("binance_dmas", 50),
            }

            print(f"[Technical Nutshell] DMAS scores collected: {dmas_scores}")

            # Prepare slide data
            rows = prepare_slide_data(
                df_prices_tech_adj,
                dmas_scores,
                str(excel_path_for_ppt),
                price_mode=pmode,
            )

            print(f"[Technical Nutshell] Prepared {len(rows)} asset rows")

            # Insert the slide (finds "technical_nutshell" placeholder)
            if rows:
                insert_technical_analysis_slide(
                    prs,
                    rows,
                    placeholder_name="technical_nutshell",
                    used_date=tech_used_date,
                    price_mode=pmode,
                )
                print(f"[Technical Nutshell] Slide generated with {len(rows)} assets")
            else:
                print("[Technical Nutshell] Warning: No assets found for slide")

        except Exception as e:
            print(f"Technical Analysis slide error: {e}")
            import traceback
            traceback.print_exc()

        # ------------------------------------------------------------------
        # Market Breadth slide (full-width)
        # ------------------------------------------------------------------
        update_progress("Generating Market Breadth slide...")
        try:
            from market_compass.breadth_slide import generate_breadth_slide

            prs = generate_breadth_slide(
                prs,
                excel_path=str(excel_path_for_ppt),
                slide_name="slide_breadth",
            )
            print("[Market Breadth] Slide generation complete")

        except Exception as e:
            print(f"Market Breadth slide error: {e}")
            import traceback
            traceback.print_exc()

        # ------------------------------------------------------------------
        # Fundamental Analysis slide (full-width)
        # ------------------------------------------------------------------
        update_progress("Generating Fundamental Analysis slide...")
        try:
            from market_compass.fundamental_slide import generate_fundamental_slide

            prs = generate_fundamental_slide(
                prs,
                excel_path=str(excel_path_for_ppt),
                slide_name="slide_fundamentals",
            )
            print("[Fundamental Analysis] Slide generation complete")

        except Exception as e:
            print(f"Fundamental Analysis slide error: {e}")
            import traceback
            traceback.print_exc()

        # ------------------------------------------------------------------
        # Enforce Calibri font on all table shapes
        # ------------------------------------------------------------------
        # The weekly performance tables and other tables inserted into the
        # presentation are created by external helper functions.  Those functions
        # may not explicitly set a font for each run in the table.  To ensure
        # visual consistency with the template (which uses Calibri), iterate
        # through all tables in the presentation and set every run's font name
        # to "Calibri".  If Calibri is unavailable on the host system, the
        # fallback fonts defined in the template will still apply.  The helper
        # functions below are defined in-line so they have access to the
        # surrounding scope.
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from pptx.util import Pt

        def _force_textframe_calibri(tf, size_pt: int = 11):
            """Set all runs in a text frame to Calibri with the given size."""
            if not tf:
                return
            for p in tf.paragraphs:
                for r in p.runs:
                    r.font.name = "Calibri"
                    # Preserve the original size if not provided
                    if size_pt is not None:
                        r.font.size = Pt(size_pt)

        def force_table_calibri(table, size_pt: int = 11):
            """Set Calibri for every cell of a python‑pptx table."""
            for row in table.rows:
                for cell in row.cells:
                    if hasattr(cell, "text_frame") and cell.text_frame is not None:
                        _force_textframe_calibri(cell.text_frame, size_pt=size_pt)

        def force_all_tables_calibri(prs, size_pt: int = 11):
            """Iterate through all tables in the presentation and enforce Calibri."""
            try:
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                            # Table shapes have a .table property exposing row/col API
                            try:
                                force_table_calibri(shape.table, size_pt=size_pt)
                            except Exception:
                                pass
            except Exception:
                # Never break presentation generation if font enforcement fails
                pass

        # Apply the font correction to all tables after all slides have been added
        force_all_tables_calibri(prs, size_pt=11)

        def disable_image_compression(prs):
            """Disable automatic image compression in PowerPoint.

            Sets compression state to 'none' for all embedded images (blips)
            to preserve full resolution when the presentation is opened.
            """
            DRAWING_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
            try:
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, '_element'):
                            # Find blip elements (images)
                            for blip in shape._element.iter(f'{DRAWING_NS}blip'):
                                # Set compression state to 'none'
                                blip.set(f'{DRAWING_NS}cstate', 'none')
            except Exception:
                # Never break presentation generation if compression disable fails
                pass

        # Disable image compression to preserve chart quality
        disable_image_compression(prs)

        update_progress("Finalizing presentation...")
        out_stream = BytesIO()
        prs.save(out_stream)
        out_stream.seek(0)
        updated_bytes = out_stream.getvalue()

        # Always generate a macro‑free PowerPoint (.pptx).  Converting a
        # macro‑enabled template (.pptm) to .pptx removes any embedded VBA
        # projects and prevents runtime errors when opening the file.  The
        # MIME type for .pptx files is used for all downloads.
        # Use the computed stamp_ddmmyyyy to construct the output filename.
        # This yields names like "03092025_Herculis_Partners_Macro_Update.pptx".
        fname = f"{stamp_ddmmyyyy}_Herculis_Partners_Technical_Update.pptx"
        mime = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

        # PHASE 2: DISABLED (no image cache to clean up)
        # disable_image_cache()

        # Complete the progress
        progress_bar.progress(1.0)
        elapsed_total = time.time() - start_time
        if elapsed_total < 60:
            time_str = f"{int(elapsed_total)}s"
        else:
            mins = int(elapsed_total // 60)
            secs = int(elapsed_total % 60)
            time_str = f"{mins}m {secs}s"
        status_text.text(f"✅ Presentation generated successfully in {time_str}! 🚀 (Phase 2 parallel acceleration active)")
        time_text.text("")

        st.sidebar.success("Updated presentation created successfully.")
        st.sidebar.download_button(
            "Download updated presentation",
            data=updated_bytes,
            file_name=fname,
            mime=mime,
            key="download_ppt_button",
        )
    st.write("Click the button in the sidebar to generate your updated presentation.")


# -----------------------------------------------------------------------------
# Main navigation dispatch
# -----------------------------------------------------------------------------
if page == "Upload":
    show_upload_page()
elif page == "YTD Update":
    show_ytd_update_page()
elif page == "Technical Analysis":
    show_technical_analysis_page()
elif page == "Market Breadth":
    show_market_breadth_page()
elif page == "Generate Presentation":
    show_generate_presentation_page()