"""
IC Technical Presentation Engine - Headless Generation Logic

This module contains the core presentation generation logic extracted from app.py.
It is designed to work without Streamlit - all state is passed as a dictionary.

Usage:
    from ic_engine import generate_presentation

    state = {...}  # Populated with scores, subtitles, etc.
    pptx_bytes, filename = generate_presentation(
        prices_path="/path/to/master_prices.csv",
        excel_path="/path/to/ic_file.xlsx",
        template_path="/path/to/template.pptx",
        state=state,
        progress_callback=print
    )
"""

import time
import tempfile
from io import BytesIO
from pathlib import Path
from typing import Dict, Any, Callable, Optional, Tuple
from datetime import datetime

import pandas as pd
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt

# Import helper functions
from utils import adjust_prices_for_mode
from technical_analysis.common_helpers import clear_excel_cache

# Import CSV data loader
from data_loader import load_prices_from_csv, get_price_series, create_temp_excel_from_csv

# Import chart generation functions
from technical_analysis.equity.spx import (
    create_technical_analysis_v2_chart,
    insert_technical_analysis_v2_slide,
    _get_spx_technical_score,
    _get_spx_momentum_score,
)

# Import score getters for all instruments
from technical_analysis.equity.csi import _get_csi_technical_score, _get_csi_momentum_score
from technical_analysis.equity.nikkei import _get_nikkei_technical_score, _get_nikkei_momentum_score
from technical_analysis.equity.tasi import _get_tasi_technical_score, _get_tasi_momentum_score
from technical_analysis.equity.sensex import _get_sensex_technical_score, _get_sensex_momentum_score
from technical_analysis.equity.dax import _get_dax_technical_score, _get_dax_momentum_score
from technical_analysis.equity.smi import _get_smi_technical_score, _get_smi_momentum_score
from technical_analysis.equity.ibov import _get_ibov_technical_score, _get_ibov_momentum_score
from technical_analysis.equity.mexbol import _get_mexbol_technical_score, _get_mexbol_momentum_score

from technical_analysis.commodity.gold import _get_gold_technical_score, _get_gold_momentum_score
from technical_analysis.commodity.silver import _get_silver_technical_score, _get_silver_momentum_score
from technical_analysis.commodity.platinum import _get_platinum_technical_score, _get_platinum_momentum_score
from technical_analysis.commodity.palladium import _get_palladium_technical_score, _get_palladium_momentum_score
from technical_analysis.commodity.oil import _get_oil_technical_score, _get_oil_momentum_score
from technical_analysis.commodity.copper import _get_copper_technical_score, _get_copper_momentum_score

from technical_analysis.crypto.bitcoin import _get_bitcoin_technical_score, _get_bitcoin_momentum_score
from technical_analysis.crypto.ethereum import _get_ethereum_technical_score, _get_ethereum_momentum_score
from technical_analysis.crypto.ripple import _get_ripple_technical_score, _get_ripple_momentum_score
from technical_analysis.crypto.solana import _get_solana_technical_score, _get_solana_momentum_score
from technical_analysis.crypto.binance import _get_binance_technical_score, _get_binance_momentum_score

# Import performance chart functions
from performance.equity_perf import (
    create_weekly_performance_chart,
    create_historical_performance_table,
    insert_equity_performance_bar_slide,
    insert_equity_performance_histo_slide,
    create_equity_ytd_evolution_chart,
    insert_equity_ytd_evolution_slide,
    create_fx_impact_analysis_chart_eur,
    insert_fx_impact_analysis_slide_eur,
    create_fx_impact_analysis_chart_chf,
    insert_fx_impact_analysis_slide_chf,
)

from performance.fx_perf import (
    create_weekly_html_performance_chart as create_weekly_fx_html_chart,
    insert_fx_weekly_html_slide,
    create_historical_html_performance_chart as create_historical_fx_html_chart,
    insert_fx_historical_html_slide,
)

from performance.crypto_perf import (
    create_weekly_html_performance_chart as create_weekly_crypto_html_chart,
    insert_crypto_weekly_html_slide,
    create_historical_html_performance_chart as create_historical_crypto_html_chart,
    insert_crypto_historical_html_slide,
    create_crypto_ytd_evolution_chart,
    insert_crypto_ytd_evolution_slide,
)

from performance.rates_perf import (
    create_weekly_performance_chart as create_weekly_rates_performance_chart,
    create_historical_performance_table as create_historical_rates_performance_table,
    insert_rates_performance_bar_slide,
    insert_rates_performance_histo_slide,
)

from performance.corp_bonds_perf import (
    create_weekly_performance_chart as create_weekly_credit_performance_chart,
    insert_corp_bonds_performance_slide as insert_credit_performance_bar_slide,
    create_historical_performance_chart as create_historical_credit_performance_chart,
    insert_corp_bonds_historical_slide as insert_credit_performance_histo_slide,
)

from performance.commodity_perf import (
    create_weekly_html_performance_chart as create_weekly_commodity_html_chart,
    insert_commodity_weekly_html_slide,
    create_historical_html_performance_chart as create_historical_commodity_html_chart,
    insert_commodity_historical_html_slide,
    create_commodity_ytd_evolution_chart,
    insert_commodity_ytd_evolution_slide,
)


# ==============================================================================
# INSTRUMENT CONFIGURATION
# ==============================================================================

# Each instrument is defined by:
# - ticker_key: Internal key for state dict (e.g., "spx")
# - display_name: Human-readable name (e.g., "S&P 500")
# - bloomberg_ticker: Ticker for chart generation (e.g., "SPX Index")
# - price_column: Column name in data_prices sheet (e.g., "SPX Index")
# - placeholder_name: Slide placeholder name (e.g., "spx_v2")
# - get_tech_score: Function to get technical score
# - get_mom_score: Function to get momentum score

INSTRUMENTS = [
    # Equity
    {
        "ticker_key": "spx",
        "display_name": "S&P 500",
        "bloomberg_ticker": "SPX Index",
        "price_column": "SPX Index",
        "placeholder_name": "spx_v2",
        "get_tech_score": _get_spx_technical_score,
        "get_mom_score": _get_spx_momentum_score,
    },
    {
        "ticker_key": "csi",
        "display_name": "CSI 300",
        "bloomberg_ticker": "SHSZ300 Index",
        "price_column": "SHSZ300 Index",
        "placeholder_name": "csi_v2",
        "get_tech_score": _get_csi_technical_score,
        "get_mom_score": _get_csi_momentum_score,
    },
    {
        "ticker_key": "nikkei",
        "display_name": "Nikkei 225",
        "bloomberg_ticker": "NKY Index",
        "price_column": "NKY Index",
        "placeholder_name": "nikkei_v2",
        "get_tech_score": _get_nikkei_technical_score,
        "get_mom_score": _get_nikkei_momentum_score,
    },
    {
        "ticker_key": "tasi",
        "display_name": "TASI",
        "bloomberg_ticker": "SASEIDX Index",
        "price_column": "SASEIDX Index",
        "placeholder_name": "tasi_v2",
        "get_tech_score": _get_tasi_technical_score,
        "get_mom_score": _get_tasi_momentum_score,
    },
    {
        "ticker_key": "sensex",
        "display_name": "Sensex",
        "bloomberg_ticker": "SENSEX Index",
        "price_column": "SENSEX Index",
        "placeholder_name": "sensex_v2",
        "get_tech_score": _get_sensex_technical_score,
        "get_mom_score": _get_sensex_momentum_score,
    },
    {
        "ticker_key": "dax",
        "display_name": "DAX",
        "bloomberg_ticker": "DAX Index",
        "price_column": "DAX Index",
        "placeholder_name": "dax_v2",
        "get_tech_score": _get_dax_technical_score,
        "get_mom_score": _get_dax_momentum_score,
    },
    {
        "ticker_key": "smi",
        "display_name": "SMI",
        "bloomberg_ticker": "SMI Index",
        "price_column": "SMI Index",
        "placeholder_name": "smi_v2",
        "get_tech_score": _get_smi_technical_score,
        "get_mom_score": _get_smi_momentum_score,
    },
    {
        "ticker_key": "ibov",
        "display_name": "Ibovespa",
        "bloomberg_ticker": "IBOV Index",
        "price_column": "IBOV Index",
        "placeholder_name": "ibov_v2",
        "get_tech_score": _get_ibov_technical_score,
        "get_mom_score": _get_ibov_momentum_score,
    },
    {
        "ticker_key": "mexbol",
        "display_name": "MEXBOL",
        "bloomberg_ticker": "MEXBOL Index",
        "price_column": "MEXBOL Index",
        "placeholder_name": "mexbol_v2",
        "get_tech_score": _get_mexbol_technical_score,
        "get_mom_score": _get_mexbol_momentum_score,
    },
    # Commodities
    {
        "ticker_key": "gold",
        "display_name": "Gold",
        "bloomberg_ticker": "GCA Comdty",
        "price_column": "GCA Comdty",
        "placeholder_name": "gold_v2",
        "get_tech_score": _get_gold_technical_score,
        "get_mom_score": _get_gold_momentum_score,
    },
    {
        "ticker_key": "silver",
        "display_name": "Silver",
        "bloomberg_ticker": "SIA Comdty",
        "price_column": "SIA Comdty",
        "placeholder_name": "silver_v2",
        "get_tech_score": _get_silver_technical_score,
        "get_mom_score": _get_silver_momentum_score,
    },
    {
        "ticker_key": "platinum",
        "display_name": "Platinum",
        "bloomberg_ticker": "XPT Comdty",
        "price_column": "XPT Comdty",
        "placeholder_name": "platinum_v2",
        "get_tech_score": _get_platinum_technical_score,
        "get_mom_score": _get_platinum_momentum_score,
    },
    {
        "ticker_key": "palladium",
        "display_name": "Palladium",
        "bloomberg_ticker": "XPD Curncy",
        "price_column": "XPD Curncy",
        "placeholder_name": "palladium_v2",
        "get_tech_score": _get_palladium_technical_score,
        "get_mom_score": _get_palladium_momentum_score,
    },
    {
        "ticker_key": "oil",
        "display_name": "Oil",
        "bloomberg_ticker": "CL1 Comdty",
        "price_column": "CL1 Comdty",
        "placeholder_name": "oil_v2",
        "get_tech_score": _get_oil_technical_score,
        "get_mom_score": _get_oil_momentum_score,
    },
    {
        "ticker_key": "copper",
        "display_name": "Copper",
        "bloomberg_ticker": "LP1 Comdty",
        "price_column": "LP1 Comdty",
        "placeholder_name": "copper_v2",
        "get_tech_score": _get_copper_technical_score,
        "get_mom_score": _get_copper_momentum_score,
    },
    # Crypto
    {
        "ticker_key": "bitcoin",
        "display_name": "Bitcoin",
        "bloomberg_ticker": "XBTUSD Curncy",
        "price_column": "XBTUSD Curncy",
        "placeholder_name": "bitcoin_v2",
        "get_tech_score": _get_bitcoin_technical_score,
        "get_mom_score": _get_bitcoin_momentum_score,
    },
    {
        "ticker_key": "ethereum",
        "display_name": "Ethereum",
        "bloomberg_ticker": "XETUSD Curncy",
        "price_column": "XETUSD Curncy",
        "placeholder_name": "ethereum_v2",
        "get_tech_score": _get_ethereum_technical_score,
        "get_mom_score": _get_ethereum_momentum_score,
    },
    {
        "ticker_key": "ripple",
        "display_name": "Ripple",
        "bloomberg_ticker": "XRPUSD Curncy",
        "price_column": "XRPUSD Curncy",
        "placeholder_name": "ripple_v2",
        "get_tech_score": _get_ripple_technical_score,
        "get_mom_score": _get_ripple_momentum_score,
    },
    {
        "ticker_key": "solana",
        "display_name": "Solana",
        "bloomberg_ticker": "XSOUSD Curncy",
        "price_column": "XSOUSD Curncy",
        "placeholder_name": "solana_v2",
        "get_tech_score": _get_solana_technical_score,
        "get_mom_score": _get_solana_momentum_score,
    },
    {
        "ticker_key": "binance",
        "display_name": "Binance",
        "bloomberg_ticker": "XBIUSD Curncy",
        "price_column": "XBIUSD Curncy",
        "placeholder_name": "binance_v2",
        "get_tech_score": _get_binance_technical_score,
        "get_mom_score": _get_binance_momentum_score,
    },
]


# ==============================================================================
# HELPER FUNCTIONS
# ==============================================================================

def _update_date_placeholder(prs, human_date: str) -> None:
    """Replace [DataIC] placeholder with formatted date."""
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            tf = shape.text_frame
            for para in tf.paragraphs:
                for run in para.runs:
                    if run.text.strip() == "[DataIC]":
                        # Preserve styling
                        size = run.font.size
                        color = run.font.color
                        rgb = getattr(color, "rgb", None) if color else None
                        bold = run.font.bold
                        italic = run.font.italic
                        # Replace text
                        run.text = human_date
                        # Reapply styling
                        if size:
                            run.font.size = size
                        try:
                            if rgb:
                                run.font.color.rgb = rgb
                        except Exception:
                            pass
                        if bold is not None:
                            run.font.bold = bold
                        if italic is not None:
                            run.font.italic = italic
                        return


def _force_textframe_calibri(tf, size_pt: int = 11):
    """Set all runs in a text frame to Calibri."""
    if not tf:
        return
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.name = "Calibri"
            if size_pt is not None:
                r.font.size = Pt(size_pt)


def _force_all_tables_calibri(prs, size_pt: int = 11):
    """Enforce Calibri font on all tables."""
    try:
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    try:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                if hasattr(cell, "text_frame") and cell.text_frame:
                                    _force_textframe_calibri(cell.text_frame, size_pt)
                    except Exception:
                        pass
    except Exception:
        pass


def _disable_image_compression(prs):
    """Disable image compression to preserve chart quality."""
    DRAWING_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    try:
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, '_element'):
                    for blip in shape._element.iter(f'{DRAWING_NS}blip'):
                        blip.set(f'{DRAWING_NS}cstate', 'none')
    except Exception:
        pass


def _load_prices_dataframe(prices_path: Path, data_as_of: Optional[pd.Timestamp] = None) -> pd.DataFrame:
    """Load and prepare prices dataframe from CSV."""
    # Convert Timestamp to date for the loader
    data_as_of_date = data_as_of.date() if data_as_of is not None else None
    return load_prices_from_csv(prices_path, data_as_of_date)


def _get_used_date_for_instrument(
    df_prices: pd.DataFrame,
    price_column: str,
    price_mode: str
) -> Optional[pd.Timestamp]:
    """Get the used date for a specific instrument."""
    try:
        df = df_prices[["Date", price_column]].copy()
        df.columns = ["Date", "Price"]
        df["Price"] = pd.to_numeric(df["Price"], errors="coerce")
        df = df.dropna(subset=["Date", "Price"]).sort_values("Date").reset_index(drop=True)
        _, used_date = adjust_prices_for_mode(df, price_mode)
        return used_date
    except Exception:
        return None


# ==============================================================================
# MAIN GENERATION FUNCTION
# ==============================================================================

def generate_presentation(
    prices_path: Path,
    excel_path: Path,
    template_path: Path,
    state: Dict[str, Any],
    progress_callback: Optional[Callable[[str], None]] = None,
) -> Tuple[bytes, str, Dict[str, dict], Dict[str, int]]:
    """
    Generate the IC Technical presentation headlessly.

    Parameters
    ----------
    prices_path : Path
        Path to master_prices.csv with daily price data.
    excel_path : Path
        Path to ic_file.xlsx with parameters, mars_score, transition sheets.
    template_path : Path
        Path to the PowerPoint template.
    state : dict
        State dictionary containing all scores, subtitles, and configuration.
        Must include keys like:
        - data_as_of: Optional[date] - Date to filter data
        - price_mode: str - "Last Price" or other
        - {ticker}_dmas: float - DMAS score for each instrument
        - {ticker}_last_week_avg: float - Previous week DMAS
        - {ticker}_selected_view: str - Assessment text
        - {ticker}_subtitle: str - Subtitle text
        - eq_subtitle, co_subtitle, cr_subtitle: str - YTD recap subtitles
    progress_callback : callable, optional
        Function to call with progress messages.

    Returns
    -------
    tuple[bytes, str, dict, dict]
        PowerPoint file bytes, suggested filename, breadth ranks dict, and fundamental ranks dict.
        Breadth ranks: {"U.S.": {"rank": 1, "pct_both": 82}, ...}
        Fundamental ranks: {"U.S.": 3, ...}
    """
    start_time = time.time()

    def progress(msg: str):
        if progress_callback:
            progress_callback(msg)
        print(f"[Engine] {msg}")

    progress("Loading PowerPoint template...")
    prs = Presentation(str(template_path))

    # Clear Excel cache for clean state
    clear_excel_cache()

    # Get configuration from state
    data_as_of = state.get("data_as_of")
    if data_as_of is not None:
        data_as_of = pd.Timestamp(data_as_of)
    price_mode = state.get("price_mode", "Last Price")

    # Compute date for slide and filename
    ts = pd.Timestamp.now(tz="Europe/Zurich")
    human_date = f"{ts.strftime('%B')} {ts.day}, {ts.year}"
    stamp_ddmmyyyy = f"{ts.day:02d}{ts.month:02d}{ts.year}"

    # Update date placeholder
    progress("Updating date placeholder...")
    _update_date_placeholder(prs, human_date)

    # Load prices dataframe once from CSV
    progress("Loading price data from CSV...")
    df_prices = _load_prices_dataframe(prices_path, data_as_of)

    # Create temporary Excel file with data_prices sheet from CSV
    # This is needed because chart functions in technical_analysis/ expect Excel format
    progress("Creating temporary Excel file from CSV data...")
    data_as_of_date = data_as_of.date() if data_as_of is not None else None
    temp_excel_path = create_temp_excel_from_csv(prices_path, excel_path, data_as_of_date)
    progress(f"Temporary Excel created at: {temp_excel_path}")

    # Use temp Excel for chart generation (has data_prices from CSV + other sheets from source)
    chart_excel_path = temp_excel_path

    # ==========================================================================
    # TECHNICAL ANALYSIS SLIDES
    # ==========================================================================

    total_instruments = len(INSTRUMENTS)
    for idx, inst in enumerate(INSTRUMENTS):
        ticker_key = inst["ticker_key"]
        display_name = inst["display_name"]

        progress(f"Processing {display_name} ({idx + 1}/{total_instruments})...")

        try:
            # Get scores from state
            dmas = state.get(f"{ticker_key}_dmas", 50)
            dmas_prev = state.get(f"{ticker_key}_last_week_avg", dmas)

            # Get computed scores from Excel (use temp Excel with data_prices from CSV)
            tech_score = inst["get_tech_score"](chart_excel_path)
            mom_score = inst["get_mom_score"](chart_excel_path)

            # Get previous week scores
            tech_prev = state.get(f"{ticker_key}_last_week_tech")
            mom_prev = state.get(f"{ticker_key}_last_week_mom")
            rsi_prev = state.get(f"{ticker_key}_last_week_rsi")

            # Get gap info
            days_gap = state.get(f"{ticker_key}_prev_days_gap")
            prev_date = state.get(f"{ticker_key}_prev_date")

            # Get used date
            used_date = _get_used_date_for_instrument(
                df_prices, inst["price_column"], price_mode
            )

            # Generate chart (use temp Excel with data_prices from CSV)
            chart_bytes, _ = create_technical_analysis_v2_chart(
                chart_excel_path,
                ticker=inst["bloomberg_ticker"],
                price_mode=price_mode,
                dmas_score=int(dmas),
                dmas_prev_week=int(dmas_prev),
                technical_score=tech_score,
                technical_prev_week=tech_prev,
                momentum_score=mom_score,
                momentum_prev_week=mom_prev,
                rsi_prev_week=rsi_prev,
                days_gap=days_gap,
                previous_date=prev_date,
            )

            # Get view and subtitle
            view_text = state.get(f"{ticker_key}_selected_view", "")
            if view_text and not view_text.lower().startswith(display_name.lower()):
                view_text = f"{display_name}: {view_text}"
            subtitle_text = state.get(f"{ticker_key}_subtitle", "")

            # Insert slide
            prs = insert_technical_analysis_v2_slide(
                prs,
                chart_bytes,
                used_date=used_date,
                price_mode=price_mode,
                placeholder_name=inst["placeholder_name"],
                view_text=view_text,
                subtitle_text=subtitle_text,
            )

        except Exception as e:
            print(f"[Engine] Error processing {display_name}: {e}")
            import traceback
            traceback.print_exc()

    # ==========================================================================
    # PERFORMANCE SLIDES
    # ==========================================================================

    progress("Generating Equity performance slides...")
    try:
        # Weekly bar chart
        bar_bytes, perf_used_date = create_weekly_performance_chart(chart_excel_path, price_mode=price_mode)
        prs = insert_equity_performance_bar_slide(
            prs, bar_bytes, used_date=perf_used_date, price_mode=price_mode,
            left_cm=3.47, top_cm=5.28, width_cm=17.31, height_cm=10
        )

        # Historical heatmap
        histo_bytes, histo_used_date = create_historical_performance_table(chart_excel_path, price_mode=price_mode)
        prs = insert_equity_performance_histo_slide(
            prs, histo_bytes, used_date=histo_used_date, price_mode=price_mode,
            left_cm=2.16, top_cm=4.70, width_cm=19.43, height_cm=10.61
        )

        # YTD evolution
        ytd_bytes, ytd_date = create_equity_ytd_evolution_chart(chart_excel_path, price_mode=price_mode)
        prs = insert_equity_ytd_evolution_slide(
            prs, ytd_bytes, used_date=ytd_date, price_mode=price_mode,
            subtitle=state.get("eq_subtitle")
        )

        # FX impact EUR
        fx_eur_bytes, fx_eur_date = create_fx_impact_analysis_chart_eur(chart_excel_path, price_mode=price_mode)
        prs = insert_fx_impact_analysis_slide_eur(prs, fx_eur_bytes, used_date=fx_eur_date, price_mode=price_mode)

        # FX impact CHF
        fx_chf_bytes, fx_chf_date = create_fx_impact_analysis_chart_chf(chart_excel_path, price_mode=price_mode)
        prs = insert_fx_impact_analysis_slide_chf(prs, fx_chf_bytes, used_date=fx_chf_date, price_mode=price_mode)
    except Exception as e:
        print(f"[Engine] Equity performance error: {e}")

    progress("Generating FX performance slides...")
    try:
        fx_bar_bytes, fx_date = create_weekly_fx_html_chart(chart_excel_path, price_mode=price_mode)
        prs = insert_fx_weekly_html_slide(prs, fx_bar_bytes, used_date=fx_date, price_mode=price_mode)

        fx_histo_bytes, fx_date2 = create_historical_fx_html_chart(chart_excel_path, price_mode=price_mode)
        prs = insert_fx_historical_html_slide(prs, fx_histo_bytes, used_date=fx_date2, price_mode=price_mode)
    except Exception as e:
        print(f"[Engine] FX performance error: {e}")

    progress("Generating Crypto performance slides...")
    try:
        crypto_bar_bytes, crypto_date = create_weekly_crypto_html_chart(chart_excel_path, price_mode=price_mode)
        prs = insert_crypto_weekly_html_slide(prs, crypto_bar_bytes, used_date=crypto_date, price_mode=price_mode)

        crypto_histo_bytes, crypto_date2 = create_historical_crypto_html_chart(chart_excel_path, price_mode=price_mode)
        prs = insert_crypto_historical_html_slide(prs, crypto_histo_bytes, used_date=crypto_date2, price_mode=price_mode)

        crypto_ytd_bytes, crypto_ytd_date = create_crypto_ytd_evolution_chart(chart_excel_path, price_mode=price_mode)
        prs = insert_crypto_ytd_evolution_slide(
            prs, crypto_ytd_bytes, used_date=crypto_ytd_date, price_mode=price_mode,
            subtitle=state.get("cr_subtitle")
        )
    except Exception as e:
        print(f"[Engine] Crypto performance error: {e}")

    progress("Generating Rates performance slides...")
    try:
        rates_bar_bytes, rates_date = create_weekly_rates_performance_chart(chart_excel_path, price_mode=price_mode)
        prs = insert_rates_performance_bar_slide(
            prs, rates_bar_bytes, used_date=rates_date, price_mode=price_mode,
            left_cm=3.35, top_cm=4.6, width_cm=17.02
        )

        rates_histo_bytes, rates_date2 = create_historical_rates_performance_table(chart_excel_path, price_mode=price_mode)
        prs = insert_rates_performance_histo_slide(
            prs, rates_histo_bytes, used_date=rates_date2, price_mode=price_mode,
            left_cm=3.35, top_cm=4.6, width_cm=17.02
        )
    except Exception as e:
        print(f"[Engine] Rates performance error: {e}")

    progress("Generating Credit performance slides...")
    try:
        credit_bar_bytes, credit_date = create_weekly_credit_performance_chart(chart_excel_path, price_mode=price_mode)
        prs = insert_credit_performance_bar_slide(
            prs, credit_bar_bytes, used_date=credit_date, price_mode=price_mode,
            left_cm=1.63, top_cm=4.73, width_cm=22.48, height_cm=10.61
        )

        credit_histo_bytes, credit_date2 = create_historical_credit_performance_chart(chart_excel_path, price_mode=price_mode)
        prs = insert_credit_performance_histo_slide(
            prs, credit_histo_bytes, used_date=credit_date2, price_mode=price_mode
        )
    except Exception as e:
        print(f"[Engine] Credit performance error: {e}")

    progress("Generating Commodity performance slides...")
    try:
        commo_bar_bytes, commo_date = create_weekly_commodity_html_chart(chart_excel_path, price_mode=price_mode)
        prs = insert_commodity_weekly_html_slide(prs, commo_bar_bytes, used_date=commo_date, price_mode=price_mode)

        commo_histo_bytes, commo_date2 = create_historical_commodity_html_chart(chart_excel_path, price_mode=price_mode)
        prs = insert_commodity_historical_html_slide(prs, commo_histo_bytes, used_date=commo_date2, price_mode=price_mode)

        commo_ytd_bytes, commo_ytd_date = create_commodity_ytd_evolution_chart(chart_excel_path, price_mode=price_mode)
        prs = insert_commodity_ytd_evolution_slide(
            prs, commo_ytd_bytes, used_date=commo_ytd_date, price_mode=price_mode,
            subtitle=state.get("co_subtitle")
        )
    except Exception as e:
        print(f"[Engine] Commodity performance error: {e}")

    # ==========================================================================
    # SUMMARY SLIDES
    # ==========================================================================

    progress("Generating Technical Analysis summary slide...")
    try:
        from market_compass.technical_slide import prepare_slide_data, insert_technical_analysis_slide

        # Adjust prices for mode
        df_prices_adj, tech_used_date = adjust_prices_for_mode(df_prices, price_mode)

        # Collect DMAS scores
        dmas_scores = {inst["ticker_key"]: state.get(f"{inst['ticker_key']}_dmas", 50) for inst in INSTRUMENTS}

        # Prepare and insert slide
        rows = prepare_slide_data(df_prices_adj, dmas_scores, str(chart_excel_path), price_mode=price_mode)
        if rows:
            insert_technical_analysis_slide(
                prs, rows, placeholder_name="technical_nutshell",
                used_date=tech_used_date, price_mode=price_mode
            )
    except Exception as e:
        print(f"[Engine] Technical summary error: {e}")
        import traceback
        traceback.print_exc()

    progress("Generating Market Breadth slide...")
    breadth_ranks = {}
    try:
        from market_compass.breadth_slide import generate_breadth_slide
        prs, breadth_ranks = generate_breadth_slide(prs, excel_path=str(chart_excel_path), slide_name="slide_breadth")
    except Exception as e:
        print(f"[Engine] Breadth slide error: {e}")

    progress("Generating Fundamental Analysis slide...")
    fundamental_ranks = {}
    try:
        from market_compass.fundamental_slide import generate_fundamental_slide
        prs, fundamental_ranks = generate_fundamental_slide(prs, excel_path=str(chart_excel_path), slide_name="slide_fundamentals")
    except Exception as e:
        print(f"[Engine] Fundamental slide error: {e}")

    # ==========================================================================
    # FINALIZE
    # ==========================================================================

    progress("Finalizing presentation...")
    _force_all_tables_calibri(prs, size_pt=11)
    _disable_image_compression(prs)

    # Save to bytes
    out_stream = BytesIO()
    prs.save(out_stream)
    out_stream.seek(0)
    pptx_bytes = out_stream.getvalue()

    # Generate filename
    filename = f"{stamp_ddmmyyyy}_Herculis_Partners_Technical_Update.pptx"

    elapsed = time.time() - start_time
    if elapsed < 60:
        time_str = f"{int(elapsed)}s"
    else:
        time_str = f"{int(elapsed // 60)}m {int(elapsed % 60)}s"

    progress(f"Presentation generated in {time_str}")

    # Cleanup temporary Excel file
    try:
        import os
        if temp_excel_path.exists():
            os.remove(temp_excel_path)
            progress("Cleaned up temporary Excel file")
    except Exception as e:
        print(f"[Engine] Warning: Could not clean up temp file: {e}")

    return pptx_bytes, filename, breadth_ranks, fundamental_ranks


# ==============================================================================
# EXPORTS
# ==============================================================================

__all__ = [
    "generate_presentation",
    "INSTRUMENTS",
]
