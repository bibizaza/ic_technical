"""Combined Global Indices slide generator - Breadth & Fundamental tables."""

from dataclasses import dataclass
from typing import List, Dict
import tempfile
from pathlib import Path

import pandas as pd
import numpy as np
from jinja2 import Template
from html2image import Html2Image
from pptx import Presentation
from pptx.util import Cm

from .html_template import GLOBAL_INDICES_HTML_TEMPLATE


# =============================================================================
# RESOLUTION SETTINGS
# =============================================================================

SCALE_FACTOR = 4
BASE_WIDTH = 903
BASE_HEIGHT = 315  # Compact table fits in smaller height
COMBINED_WIDTH_PX = BASE_WIDTH * SCALE_FACTOR   # 2709
COMBINED_HEIGHT_PX = BASE_HEIGHT * SCALE_FACTOR  # 945

# PowerPoint placement (cm)
COMBINED_LEFT_CM = 0.91
COMBINED_TOP_CM = 7.29
COMBINED_WIDTH_CM = 24.03
COMBINED_HEIGHT_CM = 8.36


# =============================================================================
# INDEX NAME MAPPING
# =============================================================================

INDEX_NAME_MAP = {
    "SPX Index": "U.S.",
    "SHSZ300 Index": "China",
    "NKY Index": "Japan",
    "SASEIDX Index": "Saudi Arabia",
    "SENSEX Index": "India",
    "DAX Index": "Germany",
    "SMI Index": "Switzerland",
    "MEXBOL Index": "Mexico",
    "IBOV Index": "Brazil",
}


# =============================================================================
# FACTOR METRICS CONFIGURATION (for Fundamental)
# =============================================================================

FACTOR_METRICS = {
    "Relative_Valuation": {
        "trailing": ["#PE", "#PB", "#PS", "#EV_EBIT", "#EV_EBITDA", "#EV_SALES"],
        "forward": ["#Forward_PE", "#Forward_PB", "#Forward_PS", "#Forward_EV_EBIT", "#Forward_EV_EBITDA", "#Forward_EV_SALES"],
        "zscore": ["#Z_PER", "#Z_Pb", "#Z_PS", "#Z_EV_EBIT", "#Z_EV_EBITDA", "#Z_EV_SALES"],
        "direction": "lower_is_better",
        "weight_trailing": 0.375,
        "weight_forward": 0.375,
        "weight_zscore": 0.25,
    },
    "Growth": {
        "trailing": ["#Growth_EPS"],
        "forward": ["#Forward_EPS"],
        "zscore": ["#Z_EPS"],
        "direction": "higher_is_better",
        "weight_trailing": 0.375,
        "weight_forward": 0.375,
        "weight_zscore": 0.25,
    },
    "Profitability": {
        "trailing": ["#Gross_margin", "#oper_margin", "#ebitda_margin", "#net_prof_margin"],
        "forward": ["#Forward_Gross_margin", "#Forward_oper_margin", "#Forward_ebitda_margin", "#Forward_prof_margin"],
        "zscore": ["#Z_gross_margin", "#Z_oper_margin", "#Z_ebitda_margin", "#Z_prof_margin"],
        "direction": "higher_is_better",
        "weight_trailing": 0.375,
        "weight_forward": 0.375,
        "weight_zscore": 0.25,
    },
    "Quality": {
        "trailing": ["#ROE", "#ROA", "#ROC"],
        "forward": ["#Forward_ROE", "#Forward_ROA", "#Forward_ROC"],
        "zscore": ["#Z_ROE", "#Z_ROA", "#Z_ROC"],
        "direction": "higher_is_better",
        "weight_trailing": 0.375,
        "weight_forward": 0.375,
        "weight_zscore": 0.25,
    },
    "Leverage": {
        "trailing": ["#tot_debt_ebitda", "#net_debt_ebitda", "#tot_debt_asset", "#tot_debt_equity"],
        "forward": ["#Forward_tot_debt_ebitda", "#Forward_net_debt_ebitda", "#Forward_tot_debt_asset", "#Forward_tot_debt_equity"],
        "zscore": ["#Z_tot_debt_ebitda", "#Z_net_debt_ebitda", "#Z_tot_debt_asset", "#Z_tot_debt_equity"],
        "direction": "lower_is_better",
        "weight_trailing": 0.375,
        "weight_forward": 0.375,
        "weight_zscore": 0.25,
    },
    "Dividend": {
        "trailing": ["#Div_Yield"],
        "forward": [],
        "zscore": [],
        "direction": "higher_is_better",
        "weight_trailing": 1.0,
        "weight_forward": 0.0,
        "weight_zscore": 0.0,
    },
}


# =============================================================================
# DATA STRUCTURES
# =============================================================================

@dataclass
class BreadthRow:
    """Data for a single row in the Breadth Rank table."""
    index_name: str
    rank: int
    pct_both: int   # % above both MAs
    pct_20d: int    # % above 20D MA
    pct_50d: int    # % above 50D MA


@dataclass
class FundamentalRow:
    """Data for a single row in the Fundamental Rank table."""
    index_name: str
    rank: int
    rv: int
    growth: int
    profitability: int
    quality: int
    leverage: int
    dividend: int


# =============================================================================
# BREADTH DATA PREPARATION
# =============================================================================

def prepare_breadth_data(excel_path: str) -> List[BreadthRow]:
    """
    Prepare Breadth Rank data from Excel sheet.

    Parameters
    ----------
    excel_path : str
        Path to Excel file

    Returns
    -------
    List[BreadthRow]
        List of breadth data rows, sorted by % above both MAs
    """
    try:
        df = pd.read_excel(excel_path, sheet_name="helper_breadth")
        print(f"[Global Indices] Loaded helper_breadth with {len(df)} rows")

        cols = df.columns.tolist()
        if len(cols) < 10:
            print(f"[Global Indices] Warning: Expected at least 10 columns, got {len(cols)}")
            return []

        ticker_col = cols[2]  # Column C
        pct_both_col = cols[7]  # Column H
        pct_20d_col = cols[8]  # Column I
        pct_50d_col = cols[9]  # Column J

        rows = []
        for _, row in df.iterrows():
            ticker = str(row[ticker_col]).strip() if pd.notna(row[ticker_col]) else ""

            if not ticker or ticker.lower() in ["ticker", "index", ""]:
                continue

            display_name = INDEX_NAME_MAP.get(ticker, ticker.replace(" Index", ""))

            try:
                pct_both = int(round(float(row[pct_both_col]) * 100)) if pd.notna(row[pct_both_col]) else 0
                pct_20d = int(round(float(row[pct_20d_col]) * 100)) if pd.notna(row[pct_20d_col]) else 0
                pct_50d = int(round(float(row[pct_50d_col]) * 100)) if pd.notna(row[pct_50d_col]) else 0
            except (ValueError, TypeError):
                pct_both = int(row[pct_both_col]) if pd.notna(row[pct_both_col]) else 0
                pct_20d = int(row[pct_20d_col]) if pd.notna(row[pct_20d_col]) else 0
                pct_50d = int(row[pct_50d_col]) if pd.notna(row[pct_50d_col]) else 0

            pct_both = max(0, min(100, pct_both))
            pct_20d = max(0, min(100, pct_20d))
            pct_50d = max(0, min(100, pct_50d))

            rows.append(BreadthRow(
                index_name=display_name,
                rank=0,
                pct_both=pct_both,
                pct_20d=pct_20d,
                pct_50d=pct_50d,
            ))

        rows.sort(key=lambda r: r.pct_both, reverse=True)
        for i, row in enumerate(rows):
            row.rank = i + 1

        print(f"[Global Indices] Prepared {len(rows)} breadth rows")
        return rows

    except Exception as e:
        print(f"[Global Indices] Error reading breadth data: {e}")
        import traceback
        traceback.print_exc()
        return []


# =============================================================================
# FUNDAMENTAL RANKING ALGORITHM
# =============================================================================

def _compute_factor_score(df: pd.DataFrame, factor_config: dict) -> pd.Series:
    """Compute weighted composite score for a factor."""
    scores = pd.Series(0.0, index=df.index)

    for component in ["trailing", "forward", "zscore"]:
        metrics = factor_config.get(component, [])
        weight = factor_config.get(f"weight_{component}", 0)

        if not metrics or weight == 0:
            continue

        available_metrics = [m for m in metrics if m in df.columns]
        if not available_metrics:
            continue

        component_ranks = []
        for metric in available_metrics:
            values = df[metric]
            if values.isna().all():
                continue

            if factor_config["direction"] == "lower_is_better":
                ranks = values.rank(method="min", ascending=True)
            else:
                ranks = values.rank(method="min", ascending=False)
            component_ranks.append(ranks)

        if component_ranks:
            avg_ranks = pd.concat(component_ranks, axis=1).mean(axis=1)
            scores += avg_ranks * weight

    return scores


def compute_fundamental_ranks(excel_path: str) -> List[FundamentalRow]:
    """
    Compute fundamental ranks from Excel file.

    Parameters
    ----------
    excel_path : str
        Path to Excel file with data_fundamental sheet

    Returns
    -------
    List[FundamentalRow]
        List of fundamental data rows, sorted by overall rank
    """
    try:
        df = pd.read_excel(excel_path, sheet_name="data_fundamental")
        print(f"[Global Indices] Loaded data_fundamental with {len(df)} rows")

        id_col = None
        for col in df.columns:
            # Skip non-string columns (e.g., NaN from empty Excel headers)
            if isinstance(col, str) and col.lower() in ["id", "index", "ticker"]:
                id_col = col
                break
        if id_col is None:
            id_col = df.columns[0]

        df = df.set_index(id_col)

        valid_indices = [idx for idx in df.index if idx in INDEX_NAME_MAP]
        if not valid_indices:
            print(f"[Global Indices] No valid indices found")
            return []

        df = df.loc[valid_indices]

        factor_scores = {}
        for factor_name, config in FACTOR_METRICS.items():
            factor_scores[factor_name] = _compute_factor_score(df, config)

        factor_ranks = {}
        for factor_name, scores in factor_scores.items():
            if scores.isna().all():
                factor_ranks[factor_name] = pd.Series(5, index=scores.index)
            else:
                factor_ranks[factor_name] = scores.rank(method="min", ascending=True).fillna(5).astype(int)

        overall_scores = sum(factor_ranks.values())
        overall_ranks = overall_scores.rank(method="min", ascending=True).astype(int)

        results = []
        for idx in df.index:
            results.append(FundamentalRow(
                index_name=INDEX_NAME_MAP.get(idx, idx),
                rank=int(overall_ranks[idx]),
                rv=int(factor_ranks["Relative_Valuation"][idx]),
                growth=int(factor_ranks["Growth"][idx]),
                profitability=int(factor_ranks["Profitability"][idx]),
                quality=int(factor_ranks["Quality"][idx]),
                leverage=int(factor_ranks["Leverage"][idx]),
                dividend=int(factor_ranks["Dividend"][idx]),
            ))

        results.sort(key=lambda x: x.rank)

        print(f"[Global Indices] Computed {len(results)} fundamental rows")
        return results

    except Exception as e:
        print(f"[Global Indices] Error computing fundamental ranks: {e}")
        import traceback
        traceback.print_exc()
        return []


# =============================================================================
# COLOR HELPERS
# =============================================================================

def get_pct_class(value: int) -> str:
    """Breadth: percentage color class."""
    if value >= 55:
        return "high"
    elif value >= 45:
        return "med-high"
    elif value >= 35:
        return "med"
    elif value >= 25:
        return "med-low"
    else:
        return "low"


def get_rank_class(rank: int) -> str:
    """Fundamental: rank color class (1-9)."""
    if rank <= 3:
        return "good"
    elif rank <= 6:
        return "neutral"
    else:
        return "bad"


# =============================================================================
# HTML GENERATION
# =============================================================================

def _prepare_breadth_for_template(rows: List[BreadthRow]) -> List[dict]:
    """Prepare breadth rows for Jinja template."""
    return [{
        "index_name": r.index_name,
        "rank": r.rank,
        "pct_both": r.pct_both,
        "pct_both_class": get_pct_class(r.pct_both),
        "pct_20d": r.pct_20d,
        "pct_20d_class": get_pct_class(r.pct_20d),
        "pct_50d": r.pct_50d,
        "pct_50d_class": get_pct_class(r.pct_50d),
    } for r in rows]


def _prepare_fundamental_for_template(rows: List[FundamentalRow]) -> List[dict]:
    """Prepare fundamental rows for Jinja template."""
    return [{
        "index_name": r.index_name,
        "rank": r.rank,
        "rv": r.rv,
        "rv_class": get_rank_class(r.rv),
        "growth": r.growth,
        "growth_class": get_rank_class(r.growth),
        "profitability": r.profitability,
        "profitability_class": get_rank_class(r.profitability),
        "quality": r.quality,
        "quality_class": get_rank_class(r.quality),
        "leverage": r.leverage,
        "leverage_class": get_rank_class(r.leverage),
        "dividend": r.dividend,
        "dividend_class": get_rank_class(r.dividend),
    } for r in rows]


def _generate_html(breadth_rows: List[BreadthRow], fundamental_rows: List[FundamentalRow]) -> str:
    """Generate combined HTML from both data sets."""
    template = Template(GLOBAL_INDICES_HTML_TEMPLATE)
    return template.render(
        breadth_rows=_prepare_breadth_for_template(breadth_rows),
        fundamental_rows=_prepare_fundamental_for_template(fundamental_rows),
        width=COMBINED_WIDTH_PX,
        height=COMBINED_HEIGHT_PX,
        scale=SCALE_FACTOR,
    )


def _html_to_png(html: str, output_path: str) -> str:
    """Convert HTML to PNG image."""
    with tempfile.TemporaryDirectory() as tmpdir:
        hti = Html2Image(
            output_path=tmpdir,
            size=(COMBINED_WIDTH_PX, COMBINED_HEIGHT_PX)
        )
        hti.screenshot(html_str=html, save_as="global_indices.png")

        import shutil
        shutil.move(str(Path(tmpdir) / "global_indices.png"), output_path)

    return output_path


# =============================================================================
# SLIDE INSERTION
# =============================================================================

def insert_global_indices(
    prs: Presentation,
    breadth_rows: List[BreadthRow],
    fundamental_rows: List[FundamentalRow],
    slide_title: str = "Global Indices Insight"
) -> Presentation:
    """
    Insert combined tables PNG into PowerPoint slide.

    Parameters
    ----------
    prs : Presentation
        PowerPoint presentation object
    breadth_rows : List[BreadthRow]
        List of breadth data rows
    fundamental_rows : List[FundamentalRow]
        List of fundamental data rows
    slide_title : str
        Title text to search for to find the correct slide

    Returns
    -------
    Presentation
        Modified presentation
    """
    if not breadth_rows and not fundamental_rows:
        print("[Global Indices] No data to display")
        return prs

    print(f"[Global Indices] Generating combined table...")
    print(f"[Global Indices] Resolution: {COMBINED_WIDTH_PX}x{COMBINED_HEIGHT_PX}px (3x)")
    print(f"[Global Indices] Breadth rows: {len(breadth_rows)}, Fundamental rows: {len(fundamental_rows)}")

    # Generate HTML
    html = _generate_html(breadth_rows, fundamental_rows)

    # Convert to PNG
    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as f:
        img_path = f.name

    _html_to_png(html, img_path)

    # Find slide by title text
    target_slide = None
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                if slide_title.lower() in shape.text.lower():
                    target_slide = slide
                    print(f"[Global Indices] Found slide '{slide_title}' at index {slide_idx + 1}")
                    break
        if target_slide:
            break

    if not target_slide:
        print(f"[Global Indices] No slide with title '{slide_title}' found")
        Path(img_path).unlink(missing_ok=True)
        return prs

    # Insert image at fixed position
    pic = target_slide.shapes.add_picture(
        img_path,
        left=Cm(COMBINED_LEFT_CM),
        top=Cm(COMBINED_TOP_CM),
        width=Cm(COMBINED_WIDTH_CM),
        height=Cm(COMBINED_HEIGHT_CM)
    )

    # Send to back (behind other elements like footnote)
    spTree = target_slide.shapes._spTree
    sp = pic._element
    spTree.remove(sp)
    spTree.insert(2, sp)  # Index 2 = back (0 and 1 are reserved)

    # Cleanup
    Path(img_path).unlink(missing_ok=True)

    print(f"[Global Indices] Combined table inserted at ({COMBINED_LEFT_CM}, {COMBINED_TOP_CM}) cm")
    print(f"[Global Indices] Size: {COMBINED_WIDTH_CM} x {COMBINED_HEIGHT_CM} cm")

    return prs


def generate_global_indices_slide(
    prs: Presentation,
    excel_path: str,
    slide_title: str = "Global Indices Insight"
) -> Presentation:
    """
    Generate combined Global Indices slide from Excel data.

    Parameters
    ----------
    prs : Presentation
        PowerPoint presentation object
    excel_path : str
        Path to Excel file with helper_breadth and data_fundamental sheets
    slide_title : str
        Title text to search for to find the correct slide

    Returns
    -------
    Presentation
        Modified presentation
    """
    # Prepare both data sets
    breadth_rows = prepare_breadth_data(excel_path)
    fundamental_rows = compute_fundamental_ranks(excel_path)

    if not breadth_rows and not fundamental_rows:
        print("[Global Indices] No data found in Excel sheets")
        return prs

    # Insert combined image into slide
    return insert_global_indices(prs, breadth_rows, fundamental_rows, slide_title)
