"""Technical Analysis slide generator - HIGH RESOLUTION.

Renders the 3 tables (Equity, Commodity, Crypto) as a high-resolution image
and inserts into the existing slide at FIXED position and size.

Uses 2x scale factor for crisp output when PowerPoint scales down.
"""

import tempfile
import shutil
from pathlib import Path
from typing import List, Optional
from datetime import datetime

from jinja2 import Template
from html2image import Html2Image
from pptx import Presentation
from pptx.util import Cm

from .html_template import TABLES_HTML_TEMPLATE
from .data_prep import AssetRow
from helpers.flag_utils import get_flag_html


# ============================================================
# RESOLUTION SETTINGS
# ============================================================

# Scale factor for crisp output (3x = ultra crisp)
SCALE_FACTOR = 4

# Base dimensions (logical size)
BASE_WIDTH = 1200
BASE_HEIGHT = 650  # Increased to fit all rows

# Actual render size (scaled up for quality)
IMAGE_WIDTH_PX = BASE_WIDTH * SCALE_FACTOR   # 2400
IMAGE_HEIGHT_PX = BASE_HEIGHT * SCALE_FACTOR  # 1300

# PowerPoint placement (cm)
PPTX_LEFT = 0.91
PPTX_TOP = 4.98
PPTX_WIDTH = 24.03
PPTX_HEIGHT = 11.74


def _get_rsi_class(rsi: int) -> str:
    """Get CSS class for RSI value."""
    if rsi > 70:
        return "rsi-overbought"
    elif rsi < 30:
        return "rsi-oversold"
    return "neutral"


def _get_ma_class(vs_50d: float) -> str:
    """Get CSS class for vs 50d MA value."""
    return "positive" if vs_50d >= 0 else "negative"


def _get_dmas_class(dmas: int) -> str:
    """Get CSS class for DMAS value."""
    if dmas >= 55:
        return "positive"
    elif dmas < 45:
        return "negative"
    return "neutral"


def _prepare_row(row: AssetRow) -> dict:
    """Convert AssetRow to template dict."""
    dmas = int(row.dmas)
    rsi = int(row.rsi)
    vs_50d = row.vs_50d_ma

    # Generate flag HTML for crypto logos (size 18 for table cells)
    flag_html = ""
    if hasattr(row, 'flag') and row.flag:
        flag_html = get_flag_html(row.flag, size=18)

    return {
        "name": row.name,
        "market_cap": row.market_cap,
        "rsi": rsi,
        "rsi_class": _get_rsi_class(rsi),
        "vs_50d_ma_fmt": f"{vs_50d:+.1f}%",
        "ma_class": _get_ma_class(vs_50d),
        "dmas": dmas,
        "dmas_class": _get_dmas_class(dmas),
        "outlook": row.outlook,
        "outlook_lower": row.outlook.lower(),
        "flag_html": flag_html,
    }


def _generate_tables_html(rows: List[AssetRow]) -> str:
    """Generate HTML for tables with scaled dimensions."""
    equity = [_prepare_row(r) for r in rows if r.asset_class == "equity"]
    commodity = [_prepare_row(r) for r in rows if r.asset_class == "commodities"]
    crypto = [_prepare_row(r) for r in rows if r.asset_class == "crypto"]

    print(f"[Technical Nutshell] Prepared rows - Equity: {len(equity)}, Commodity: {len(commodity)}, Crypto: {len(crypto)}")

    template = Template(TABLES_HTML_TEMPLATE)
    return template.render(
        equity_rows=equity,
        commodity_rows=commodity,
        crypto_rows=crypto,
        width=IMAGE_WIDTH_PX,
        height=IMAGE_HEIGHT_PX,
        scale=SCALE_FACTOR,
    )


def _html_to_png(html: str, output_path: str) -> str:
    """Convert HTML to high-resolution PNG."""
    with tempfile.TemporaryDirectory() as tmpdir:
        hti = Html2Image(
            output_path=tmpdir,
            size=(IMAGE_WIDTH_PX, IMAGE_HEIGHT_PX)
        )
        hti.screenshot(html_str=html, save_as="tables.png")

        shutil.move(str(Path(tmpdir) / "tables.png"), output_path)

    return output_path


def insert_technical_analysis_slide(
    prs: Presentation,
    rows: List[AssetRow],
    placeholder_name: str = "technical_nutshell",
    used_date: Optional[datetime] = None,
    price_mode: str = "Last Price"
) -> Presentation:
    """
    Generate high-resolution tables image and insert into slide.

    Uses 2x scale factor for crisp output when PowerPoint scales down.

    Parameters
    ----------
    prs : Presentation
        The PowerPoint presentation object
    rows : List[AssetRow]
        List of asset rows with data
    placeholder_name : str
        Name of the placeholder shape to find and remove
    used_date : datetime, optional
        Date to display (not used in tables-only mode)
    price_mode : str
        Price mode (not used in tables-only mode)

    Returns
    -------
    Presentation
        The modified presentation
    """
    print(f"[Technical Nutshell] Generating HTML tables with {len(rows)} assets...")
    print(f"[Technical Nutshell] Resolution: {IMAGE_WIDTH_PX}x{IMAGE_HEIGHT_PX}px ({SCALE_FACTOR}x scale)")

    # Generate HTML
    html = _generate_tables_html(rows)

    # Convert to PNG
    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as f:
        img_path = f.name

    print(f"[Technical Nutshell] Rendering image...")
    _html_to_png(html, img_path)

    # Find slide with placeholder
    target_slide = None

    print(f"[Technical Nutshell] Searching for placeholder '{placeholder_name}'...")

    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            name = getattr(shape, "name", "")
            if placeholder_name.lower() in name.lower():
                target_slide = slide
                print(f"[Technical Nutshell] Found on slide {slide_idx + 1}")
                # Remove placeholder shape
                sp = shape._element
                sp.getparent().remove(sp)
                break
        if target_slide:
            break

    if not target_slide:
        print(f"[Technical Nutshell] No placeholder found, creating new slide")
        target_slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Insert image at fixed position
    target_slide.shapes.add_picture(
        img_path,
        left=Cm(PPTX_LEFT),
        top=Cm(PPTX_TOP),
        width=Cm(PPTX_WIDTH),
        height=Cm(PPTX_HEIGHT)
    )

    # Cleanup
    Path(img_path).unlink(missing_ok=True)

    print(f"[Technical Nutshell] ✅ High-res tables inserted at ({PPTX_LEFT}, {PPTX_TOP}) cm")

    return prs


# Backwards compatibility
def generate_technical_analysis_slide(
    prs: Presentation,
    rows: List[AssetRow],
    used_date: Optional[datetime] = None,
    price_mode: str = "Last Price"
) -> Presentation:
    """Generate slide (backwards compatibility wrapper)."""
    return insert_technical_analysis_slide(prs, rows, "technical_nutshell", used_date, price_mode)
