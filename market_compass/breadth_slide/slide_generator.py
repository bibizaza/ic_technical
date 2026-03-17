"""Composite Breadth Score slide generator.

Renders a 6-column table (Index, Rank, Composite, Trend, Conviction, Sentiment)
from breadth records computed by pipeline/breadth.py, then inserts into PPTX.
"""

from typing import List, Dict, Tuple, Optional
import tempfile
from pathlib import Path

from jinja2 import Template
from pptx import Presentation
from pptx.util import Cm

from .html_template import BREADTH_HTML_TEMPLATE
from helpers.html_to_image import render_html_to_image


# =============================================================================
# RESOLUTION SETTINGS (matches fundamental_slide exactly)
# =============================================================================

SCALE_FACTOR = 4
BREADTH_BASE_WIDTH = 750
BREADTH_BASE_HEIGHT = 360
BREADTH_WIDTH_PX = BREADTH_BASE_WIDTH * SCALE_FACTOR   # 3000
BREADTH_HEIGHT_PX = BREADTH_BASE_HEIGHT * SCALE_FACTOR  # 1440

# PowerPoint placement (cm) — matches fundamental_slide
BREADTH_LEFT_CM = 2.7
BREADTH_TOP_CM = 7.25
BREADTH_WIDTH_CM = 20.0
BREADTH_HEIGHT_CM = 9.5


# =============================================================================
# INDEX NAME + FLAG MAPPING — same 9 markets as fundamental/IC universe
# =============================================================================

INDEX_NAME_MAP = {
    # Keyed by display name (from config/tickers.yaml breadth_indices)
    "S&P 500":   ("\U0001F1FA\U0001F1F8", "U.S."),
    "SMI":       ("\U0001F1E8\U0001F1ED", "Switzerland"),
    "DAX":       ("\U0001F1E9\U0001F1EA", "Germany"),
    "Nikkei 225":("\U0001F1EF\U0001F1F5", "Japan"),
    "CSI 300":   ("\U0001F1E8\U0001F1F3", "China"),
    "TASI":      ("\U0001F1F8\U0001F1E6", "Saudi Arabia"),
    "Sensex":    ("\U0001F1EE\U0001F1F3", "India"),
    "MEXBOL":    ("\U0001F1F2\U0001F1FD", "Mexico"),
    "IBOV":      ("\U0001F1E7\U0001F1F7", "Brazil"),
}


# =============================================================================
# COLOR CLASSIFICATION
# =============================================================================

def _color_class(value: float) -> str:
    """Classify value into green/amber/red for CSS styling."""
    if value >= 55:
        return "green"
    elif value >= 35:
        return "amber"
    else:
        return "red"


# =============================================================================
# DATA PREPARATION (from draft_state breadth records)
# =============================================================================

def _prepare_rows_from_records(breadth_records: list) -> list:
    """
    Convert breadth records (from pipeline/breadth.py via draft_state.json)
    into template-ready row dicts.

    Each record has: name, composite, trend, conviction, sentiment, rank
    """
    rows = []
    for rec in breadth_records:
        ticker = rec["name"]
        mapping = INDEX_NAME_MAP.get(ticker)
        if not mapping:
            continue

        flag, display_name = mapping
        composite = float(rec["composite"])
        trend = float(rec["trend"])
        # conviction (was "momentum" in old records)
        conviction = float(rec.get("conviction", rec.get("momentum", 50)))
        # sentiment (was "skew" / "extension" in old records)
        sentiment = float(rec.get("sentiment", rec.get("skew", rec.get("extension", 50))))

        rows.append({
            "rank": int(rec["rank"]),
            "flag": flag,
            "name": display_name,
            "composite": int(round(composite)),
            "composite_class": _color_class(composite),
            "trend": min(100, max(0, trend)),
            "trend_int": int(round(trend)),
            "trend_class": _color_class(trend),
            "conviction": min(100, max(0, conviction)),
            "conviction_int": int(round(conviction)),
            "conviction_class": _color_class(conviction),
            "sentiment": min(100, max(0, sentiment)),
            "sentiment_int": int(round(sentiment)),
            "sentiment_class": _color_class(sentiment),
        })

    rows.sort(key=lambda r: r["rank"])
    return rows


# =============================================================================
# HTML GENERATION + RENDERING
# =============================================================================

def _generate_html(rows: list) -> str:
    """Generate HTML from prepared row data."""
    template = Template(BREADTH_HTML_TEMPLATE)
    return template.render(
        rows=rows,
        width=BREADTH_WIDTH_PX,
        height=BREADTH_HEIGHT_PX,
        scale=SCALE_FACTOR,
    )


def _html_to_png(html: str, output_path: str) -> str:
    """Convert HTML to PNG image using Playwright."""
    print(f"[Composite Breadth] Rendering image ({BREADTH_WIDTH_PX}x{BREADTH_HEIGHT_PX}px)...")
    return render_html_to_image(
        html_content=html,
        output_path=output_path,
        size=(BREADTH_WIDTH_PX, BREADTH_HEIGHT_PX),
        filename="breadth_composite.png",
        device_scale_factor=1,
    )


# =============================================================================
# SLIDE INSERTION
# =============================================================================

def insert_composite_breadth(
    prs: Presentation,
    rows: list,
    slide_name: str = "slide_breadth",
) -> Presentation:
    """Insert Composite Breadth Score table into PowerPoint slide."""
    if not rows:
        print("[Composite Breadth] No data to display")
        return prs

    print(f"[Composite Breadth] Generating table with {len(rows)} rows...")

    html = _generate_html(rows)

    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as f:
        img_path = f.name

    _html_to_png(html, img_path)

    # Find slide by shape name
    target_slide = None
    slide_name_lower = slide_name.lower().strip()

    print(f"[Composite Breadth] Searching for slide with shape named '{slide_name}'...")
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, "name") and shape.name:
                if shape.name.lower().strip() == slide_name_lower:
                    target_slide = slide
                    print(f"[Composite Breadth] Found slide at index {slide_idx + 1}")
                    break
        if target_slide:
            break

    if not target_slide:
        print(f"[Composite Breadth] No slide with shape name '{slide_name}' found")
        Path(img_path).unlink(missing_ok=True)
        return prs

    # Insert image
    picture = target_slide.shapes.add_picture(
        img_path,
        left=Cm(BREADTH_LEFT_CM),
        top=Cm(BREADTH_TOP_CM),
        width=Cm(BREADTH_WIDTH_CM),
        height=Cm(BREADTH_HEIGHT_CM),
    )

    # Send to back
    spTree = target_slide.shapes._spTree
    sp = picture._element
    spTree.remove(sp)
    spTree.insert(2, sp)

    Path(img_path).unlink(missing_ok=True)

    print(f"[Composite Breadth] Table inserted at ({BREADTH_LEFT_CM}, {BREADTH_TOP_CM}) cm")
    return prs


# =============================================================================
# PUBLIC API: called from assemble.py
# =============================================================================

def generate_composite_breadth_slide(
    prs: Presentation,
    breadth_records: list,
    slide_name: str = "slide_breadth",
) -> Presentation:
    """
    Generate Composite Breadth Score slide from breadth records.

    Parameters
    ----------
    prs : Presentation
        PowerPoint presentation object
    breadth_records : list
        List of dicts with keys: name, composite, trend, conviction, sentiment, rank
        (output of pipeline/breadth.py compute_composite_breadth)
    slide_name : str
        Shape name to find the target slide

    Returns
    -------
    Presentation
        Modified presentation
    """
    rows = _prepare_rows_from_records(breadth_records)

    if not rows:
        print("[Composite Breadth] No breadth data available")
        return prs

    for row in rows:
        print(f"  {row['rank']}. {row['name']}: Composite={row['composite']}, "
              f"Trend={row['trend_int']}, Conv={row['conviction_int']}, Sent={row['sentiment_int']}")

    return insert_composite_breadth(prs, rows, slide_name)


# =============================================================================
# BACKWARDS COMPAT: keep old API working (used by existing code paths)
# =============================================================================

# Re-export old names for any code that imports them
from dataclasses import dataclass

@dataclass
class BreadthRow:
    """Legacy data structure (kept for backwards compatibility)."""
    index_name: str
    flag_code: str
    rank: int
    pct_both: int
    pct_20d: int
    pct_50d: int


def prepare_breadth_data(excel_path: str) -> list:
    """Legacy function — returns empty list, use generate_composite_breadth_slide instead."""
    print("[Breadth Rank] Warning: prepare_breadth_data is deprecated, use composite breadth")
    return []


def insert_breadth_rank(prs, rows, slide_name="slide_breadth"):
    """Legacy function — no-op."""
    return prs


def generate_breadth_slide(prs, excel_path="", slide_name="slide_breadth"):
    """Legacy function — returns unchanged prs."""
    print("[Breadth Rank] Warning: generate_breadth_slide is deprecated")
    return prs, {}
