"""Weekly Performance chart generator for all asset classes."""

from dataclasses import dataclass
from typing import List, Optional
import tempfile
from pathlib import Path

from jinja2 import Template
from html2image import Html2Image
from pptx import Presentation
from pptx.util import Cm

from .html_template import WEEKLY_PERFORMANCE_HTML_TEMPLATE


# =============================================================================
# CONSTANTS
# =============================================================================

SCALE_FACTOR = 3

# Index name and flag mapping for Equity
# Flag codes are ISO 3166-1 alpha-2 for flagcdn.com
EQUITY_INDEX_MAP = {
    "SPX Index": {"name": "S&P 500", "flag": "us"},
    "DAX Index": {"name": "Dax", "flag": "de"},
    "SMI Index": {"name": "SMI", "flag": "ch"},
    "NKY Index": {"name": "Nikkei 225", "flag": "jp"},
    "SHSZ300 Index": {"name": "CSI 300", "flag": "cn"},
    "SENSEX Index": {"name": "Sensex", "flag": "in"},
    "IBOV Index": {"name": "Bovespa", "flag": "br"},
    "MEXBOL Index": {"name": "Mexbol", "flag": "mx"},
    "SASEIDX Index": {"name": "TASI", "flag": "sa"},
}


# =============================================================================
# DATA STRUCTURES
# =============================================================================

@dataclass
class PerformanceRow:
    """Data for a single row in the Weekly Performance chart."""
    name: str           # Display name (e.g., "Dax", "S&P 500", "Gold")
    value: float        # Weekly return as decimal (e.g., 0.006 for +0.6%)
    flag: str           # Emoji flag or icon (e.g., flag emoji, medal emoji)


# =============================================================================
# HELPER FUNCTIONS
# =============================================================================

def prepare_performance_rows(
    rows: List[PerformanceRow],
    max_abs_value: float = None
) -> List[dict]:
    """
    Prepare rows for template rendering.
    Rows should be pre-sorted by value (descending).
    """
    if not rows:
        return []

    # Calculate max for bar scaling
    if max_abs_value is None:
        max_abs_value = max(abs(r.value) for r in rows)

    # Ensure minimum scale
    max_abs_value = max(max_abs_value, 0.01)  # At least 1%

    prepared = []
    for i, row in enumerate(rows):
        # Determine highlight class
        highlight_class = ""
        if i == 0 and row.value > 0:
            highlight_class = "top-performer"
        elif i == len(rows) - 1 and row.value < 0:
            highlight_class = "worst-performer"

        # Bar width as percentage of half the chart (since center is 50%)
        bar_width = abs(row.value) / max_abs_value * 50
        bar_width = min(bar_width, 48)  # Cap at 48% to leave margin

        # Value classes
        if row.value > 0:
            bar_class = "positive"
            value_class = "positive"
            formatted_value = f"+{row.value * 100:.1f}%"
        elif row.value < 0:
            bar_class = "negative"
            value_class = "negative"
            formatted_value = f"{row.value * 100:.1f}%"
        else:
            bar_class = ""
            value_class = "zero"
            formatted_value = "0.0%"

        prepared.append({
            "name": row.name,
            "flag": row.flag,
            "value": row.value,
            "highlight_class": highlight_class,
            "bar_class": bar_class,
            "bar_width": bar_width,
            "value_class": value_class,
            "formatted_value": formatted_value,
        })

    return prepared


def calculate_scale(max_abs_value: float) -> dict:
    """Calculate nice scale values."""
    # Round up to nice number
    if max_abs_value <= 0.01:
        scale_max = 1
    elif max_abs_value <= 0.02:
        scale_max = 2
    elif max_abs_value <= 0.05:
        scale_max = 5
    elif max_abs_value <= 0.10:
        scale_max = 10
    else:
        scale_max = int(max_abs_value * 100) + 5

    return {
        "scale_min": f"-{scale_max}%",
        "scale_mid_low": f"-{scale_max // 2}%",
        "scale_mid_high": f"+{scale_max // 2}%",
        "scale_max": f"+{scale_max}%",
    }


# =============================================================================
# PNG GENERATION
# =============================================================================

def generate_weekly_performance_png(
    rows: List[PerformanceRow],
    output_path: str,
    width_px: int = 2400,
    height_px: int = 1050,
) -> str:
    """Generate Weekly Performance PNG."""

    # Sort by value descending
    sorted_rows = sorted(rows, key=lambda x: x.value, reverse=True)

    # Calculate scale
    max_abs = max(abs(r.value) for r in sorted_rows) if sorted_rows else 0.02
    scale_values = calculate_scale(max_abs)

    # Prepare data
    prepared_rows = prepare_performance_rows(sorted_rows, max_abs)

    # Generate HTML
    template = Template(WEEKLY_PERFORMANCE_HTML_TEMPLATE)
    html = template.render(
        rows=prepared_rows,
        width=width_px,
        height=height_px,
        scale=SCALE_FACTOR,
        **scale_values,
    )

    # Convert to PNG
    with tempfile.TemporaryDirectory() as tmpdir:
        hti = Html2Image(output_path=tmpdir, size=(width_px, height_px))
        hti.screenshot(html_str=html, save_as="weekly_perf.png")

        import shutil
        shutil.move(str(Path(tmpdir) / "weekly_perf.png"), output_path)

    return output_path


# =============================================================================
# SLIDE INSERTION
# =============================================================================

def insert_weekly_performance(
    prs: Presentation,
    rows: List[PerformanceRow],
    slide_title: str,
    left_cm: float,
    top_cm: float,
    width_cm: float,
    height_cm: float,
) -> Presentation:
    """
    Insert Weekly Performance chart into slide.

    Parameters
    ----------
    prs : Presentation
        PowerPoint presentation object
    rows : List[PerformanceRow]
        List of performance data rows
    slide_title : str
        Title text to search for to find the correct slide
    left_cm, top_cm : float
        Position in centimeters
    width_cm, height_cm : float
        Size in centimeters

    Returns
    -------
    Presentation
        Modified presentation
    """
    if not rows:
        print(f"[Weekly Performance] No data provided")
        return prs

    # Calculate pixel dimensions (37.8 px per cm at 96 DPI)
    width_px = int(width_cm * 37.8 * SCALE_FACTOR)
    height_px = int(height_cm * 37.8 * SCALE_FACTOR)

    print(f"[Weekly Performance] Generating chart...")
    print(f"[Weekly Performance] Resolution: {width_px}x{height_px}px")

    # Generate PNG
    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as f:
        img_path = f.name

    generate_weekly_performance_png(rows, img_path, width_px, height_px)

    # Find slide
    target_slide = None
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                if slide_title.lower() in shape.text.lower():
                    target_slide = slide
                    print(f"[Weekly Performance] Found slide '{slide_title}' at index {slide_idx + 1}")
                    break
        if target_slide:
            break

    if not target_slide:
        print(f"[Weekly Performance] ERROR: Slide '{slide_title}' not found")
        Path(img_path).unlink(missing_ok=True)
        return prs

    # Insert picture
    pic = target_slide.shapes.add_picture(
        img_path,
        left=Cm(left_cm),
        top=Cm(top_cm),
        width=Cm(width_cm),
        height=Cm(height_cm)
    )

    # Send to back (behind other elements like footnote)
    spTree = target_slide.shapes._spTree
    sp = pic._element
    spTree.remove(sp)
    spTree.insert(2, sp)  # Index 2 = back (0 and 1 are reserved)

    # Cleanup
    Path(img_path).unlink(missing_ok=True)

    print(f"[Weekly Performance] Chart inserted at ({left_cm}, {top_cm}) cm")
    print(f"[Weekly Performance] Size: {width_cm} x {height_cm} cm")

    return prs
