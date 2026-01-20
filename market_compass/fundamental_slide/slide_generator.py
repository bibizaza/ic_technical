"""Fundamental Rank slide generator."""

from dataclasses import dataclass
from typing import List, Dict, Optional
import tempfile
from pathlib import Path

import pandas as pd
import numpy as np
from jinja2 import Template
from html2image import Html2Image
from pptx import Presentation
from pptx.util import Cm

from .html_template import FUNDAMENTAL_HTML_TEMPLATE
from helpers.flag_utils import get_flag_html


# =============================================================================
# RESOLUTION SETTINGS (Full-width standalone slide)
# =============================================================================

SCALE_FACTOR = 4
FUNDAMENTAL_BASE_WIDTH = 750  # Wider for full-width display
FUNDAMENTAL_BASE_HEIGHT = 331
FUNDAMENTAL_WIDTH_PX = FUNDAMENTAL_BASE_WIDTH * SCALE_FACTOR   # 3000
FUNDAMENTAL_HEIGHT_PX = FUNDAMENTAL_BASE_HEIGHT * SCALE_FACTOR  # 1324

# PowerPoint placement (cm) - same as Breadth slide
FUNDAMENTAL_LEFT_CM = 2.7    # Horizontal position
FUNDAMENTAL_TOP_CM = 7.25    # Vertical position (same as Breadth)
FUNDAMENTAL_WIDTH_CM = 20.0  # Full width
FUNDAMENTAL_HEIGHT_CM = 9.5  # Same height as Breadth


# =============================================================================
# INDEX NAME MAPPING
# =============================================================================

# Map tickers to (display_name, flag_code)
INDEX_NAME_MAP = {
    "SPX Index": ("U.S.", "us"),
    "SHSZ300 Index": ("China", "cn"),
    "NKY Index": ("Japan", "jp"),
    "SASEIDX Index": ("Saudi Arabia", "sa"),
    "SENSEX Index": ("India", "in"),
    "DAX Index": ("Germany", "de"),
    "SMI Index": ("Switzerland", "ch"),
    "MEXBOL Index": ("Mexico", "mx"),
    "IBOV Index": ("Brazil", "br"),
}


# =============================================================================
# FACTOR METRICS CONFIGURATION
# =============================================================================

FACTOR_METRICS = {
    "Relative_Valuation": {
        "trailing": ["#PE", "#PB", "#PS", "#EV_EBIT", "#EV_EBITDA", "#EV_SALES"],
        "forward": ["#Forward_PE", "#Forward_PB", "#Forward_PS", "#Forward_EV_EBIT", "#Forward_EV_EBITDA", "#Forward_EV_SALES"],
        "zscore": ["#Z_PER", "#Z_Pb", "#Z_PS", "#Z_EV_EBIT", "#Z_EV_EBITDA", "#Z_EV_SALES"],
        "direction": "lower_is_better",
        "weight_trailing": 0.375,
        "weight_forward": 0.375,
        "weight_zscore": 0.25,
    },
    "Growth": {
        "trailing": ["#Growth_EPS"],
        "forward": ["#Forward_EPS"],
        "zscore": ["#Z_EPS"],
        "direction": "higher_is_better",
        "weight_trailing": 0.375,
        "weight_forward": 0.375,
        "weight_zscore": 0.25,
    },
    "Profitability": {
        "trailing": ["#Gross_margin", "#oper_margin", "#ebitda_margin", "#net_prof_margin"],
        "forward": ["#Forward_Gross_margin", "#Forward_oper_margin", "#Forward_ebitda_margin", "#Forward_prof_margin"],
        "zscore": ["#Z_gross_margin", "#Z_oper_margin", "#Z_ebitda_margin", "#Z_prof_margin"],
        "direction": "higher_is_better",
        "weight_trailing": 0.375,
        "weight_forward": 0.375,
        "weight_zscore": 0.25,
    },
    "Quality": {
        "trailing": ["#ROE", "#ROA", "#ROC"],
        "forward": ["#Forward_ROE", "#Forward_ROA", "#Forward_ROC"],
        "zscore": ["#Z_ROE", "#Z_ROA", "#Z_ROC"],
        "direction": "higher_is_better",
        "weight_trailing": 0.375,
        "weight_forward": 0.375,
        "weight_zscore": 0.25,
    },
    "Leverage": {
        "trailing": ["#tot_debt_ebitda", "#net_debt_ebitda", "#tot_debt_asset", "#tot_debt_equity"],
        "forward": ["#Forward_tot_debt_ebitda", "#Forward_net_debt_ebitda", "#Forward_tot_debt_asset", "#Forward_tot_debt_equity"],
        "zscore": ["#Z_tot_debt_ebitda", "#Z_net_debt_ebitda", "#Z_tot_debt_asset", "#Z_tot_debt_equity"],
        "direction": "lower_is_better",
        "weight_trailing": 0.375,
        "weight_forward": 0.375,
        "weight_zscore": 0.25,
    },
    "Dividend": {
        "trailing": ["#Div_Yield"],
        "forward": [],
        "zscore": [],
        "direction": "higher_is_better",
        "weight_trailing": 1.0,
        "weight_forward": 0.0,
        "weight_zscore": 0.0,
    },
}


# =============================================================================
# DATA STRUCTURES
# =============================================================================

@dataclass
class FundamentalRow:
    """Data for a single row in the Fundamental Rank table."""
    index_name: str
    flag_code: str  # ISO country code for flag
    rank: int
    rv: int
    growth: int
    profitability: int
    quality: int
    leverage: int
    dividend: int


# =============================================================================
# RANKING ALGORITHM
# =============================================================================

def _compute_factor_score(df: pd.DataFrame, factor_config: dict) -> pd.Series:
    """
    Compute weighted composite score for a factor.

    For each metric category (trailing, forward, zscore):
    1. Rank each metric 1-9 (handle direction)
    2. Average the ranks within the category
    3. Apply weight
    4. Sum weighted averages
    """
    scores = pd.Series(0.0, index=df.index)

    for component in ["trailing", "forward", "zscore"]:
        metrics = factor_config.get(component, [])
        weight = factor_config.get(f"weight_{component}", 0)

        if not metrics or weight == 0:
            continue

        # Get available metrics
        available_metrics = [m for m in metrics if m in df.columns]
        if not available_metrics:
            continue

        # Rank each metric
        component_ranks = []
        for metric in available_metrics:
            values = df[metric]

            if values.isna().all():
                continue

            # Rank based on direction
            if factor_config["direction"] == "lower_is_better":
                ranks = values.rank(method="min", ascending=True)
            else:
                ranks = values.rank(method="min", ascending=False)

            component_ranks.append(ranks)

        if component_ranks:
            avg_ranks = pd.concat(component_ranks, axis=1).mean(axis=1)
            scores += avg_ranks * weight

    return scores


def compute_fundamental_ranks(excel_path: str) -> List[FundamentalRow]:
    """
    Compute fundamental ranks from Excel file.

    Parameters
    ----------
    excel_path : str
        Path to Excel file with data_fundamental sheet

    Returns
    -------
    List[FundamentalRow]
        List of fundamental data rows, sorted by overall rank
    """
    try:
        # Read data
        df = pd.read_excel(excel_path, sheet_name="data_fundamental")
        print(f"[Fundamental Rank] Loaded sheet with {len(df)} rows")

        # Try to identify the ID column
        id_col = None
        for col in df.columns:
            if col.lower() in ["id", "index", "ticker"]:
                id_col = col
                break

        if id_col is None:
            # Use first column as ID
            id_col = df.columns[0]

        df = df.set_index(id_col)
        print(f"[Fundamental Rank] Using '{id_col}' as index, indices: {list(df.index)}")

        # Filter to only our indices
        valid_indices = [idx for idx in df.index if idx in INDEX_NAME_MAP]
        if not valid_indices:
            print(f"[Fundamental Rank] No valid indices found. Available: {list(df.index)[:10]}")
            return []

        df = df.loc[valid_indices]
        print(f"[Fundamental Rank] Processing {len(df)} indices")

        # Compute score for each factor
        factor_scores = {}
        for factor_name, config in FACTOR_METRICS.items():
            factor_scores[factor_name] = _compute_factor_score(df, config)

        # Convert scores to ranks (1-9) for each factor
        factor_ranks = {}
        for factor_name, scores in factor_scores.items():
            # Handle all NaN case
            if scores.isna().all():
                factor_ranks[factor_name] = pd.Series(5, index=scores.index)  # Neutral
            else:
                factor_ranks[factor_name] = scores.rank(method="min", ascending=True).fillna(5).astype(int)

        # Compute overall score (sum of all factor ranks)
        overall_scores = sum(factor_ranks.values())
        overall_ranks = overall_scores.rank(method="min", ascending=True).astype(int)

        # Build result rows
        results = []
        for idx in df.index:
            # Get display name and flag code
            name_flag = INDEX_NAME_MAP.get(idx)
            if name_flag:
                display_name, flag_code = name_flag
            else:
                display_name = idx
                flag_code = ""

            results.append(FundamentalRow(
                index_name=display_name,
                flag_code=flag_code,
                rank=int(overall_ranks[idx]),
                rv=int(factor_ranks["Relative_Valuation"][idx]),
                growth=int(factor_ranks["Growth"][idx]),
                profitability=int(factor_ranks["Profitability"][idx]),
                quality=int(factor_ranks["Quality"][idx]),
                leverage=int(factor_ranks["Leverage"][idx]),
                dividend=int(factor_ranks["Dividend"][idx]),
            ))

        # Sort by overall rank
        results.sort(key=lambda x: x.rank)

        print(f"[Fundamental Rank] Computed ranks for {len(results)} indices")
        for row in results:
            print(f"  {row.rank}. {row.index_name}: RV={row.rv}, Grw={row.growth}, Prof={row.profitability}, Qual={row.quality}, Lev={row.leverage}, Div={row.dividend}")

        return results

    except Exception as e:
        print(f"[Fundamental Rank] Error computing ranks: {e}")
        import traceback
        traceback.print_exc()
        return []


# =============================================================================
# HTML GENERATION
# =============================================================================

def _get_rank_class(rank: int) -> str:
    """Get CSS class for rank value (1-9)."""
    if rank <= 3:
        return "good"
    elif rank <= 6:
        return "neutral"
    else:
        return "bad"


def _prepare_row_for_template(row: FundamentalRow) -> dict:
    """Prepare row data for Jinja template."""
    # Generate flag HTML (size 22 to match other slides)
    flag_html = get_flag_html(row.flag_code, size=22) if row.flag_code else ""

    return {
        "index_name": row.index_name,
        "flag_html": flag_html,
        "rank": row.rank,
        "rv": row.rv,
        "rv_class": _get_rank_class(row.rv),
        "growth": row.growth,
        "growth_class": _get_rank_class(row.growth),
        "profitability": row.profitability,
        "profitability_class": _get_rank_class(row.profitability),
        "quality": row.quality,
        "quality_class": _get_rank_class(row.quality),
        "leverage": row.leverage,
        "leverage_class": _get_rank_class(row.leverage),
        "dividend": row.dividend,
        "dividend_class": _get_rank_class(row.dividend),
    }


def _generate_html(rows: List[FundamentalRow]) -> str:
    """Generate HTML from data."""
    template = Template(FUNDAMENTAL_HTML_TEMPLATE)
    return template.render(
        rows=[_prepare_row_for_template(r) for r in rows],
        width=FUNDAMENTAL_WIDTH_PX,
        height=FUNDAMENTAL_HEIGHT_PX,
        scale=SCALE_FACTOR,
    )


def _html_to_png(html: str, output_path: str) -> str:
    """Convert HTML to PNG image."""
    with tempfile.TemporaryDirectory() as tmpdir:
        hti = Html2Image(
            output_path=tmpdir,
            size=(FUNDAMENTAL_WIDTH_PX, FUNDAMENTAL_HEIGHT_PX)
        )
        hti.screenshot(html_str=html, save_as="fundamental.png")

        import shutil
        shutil.move(str(Path(tmpdir) / "fundamental.png"), output_path)

    return output_path


# =============================================================================
# SLIDE INSERTION
# =============================================================================

def insert_fundamental_rank(
    prs: Presentation,
    rows: List[FundamentalRow],
    slide_title: str = "Fundamental Analysis"
) -> Presentation:
    """
    Insert Fundamental Rank table into PowerPoint slide.

    Parameters
    ----------
    prs : Presentation
        PowerPoint presentation object
    rows : List[FundamentalRow]
        List of fundamental data rows
    slide_title : str
        Title text to search for to find the correct slide

    Returns
    -------
    Presentation
        Modified presentation
    """
    if not rows:
        print("[Fundamental Rank] No data to display")
        return prs

    print(f"[Fundamental Rank] Generating table with {len(rows)} rows...")
    print(f"[Fundamental Rank] Resolution: {FUNDAMENTAL_WIDTH_PX}x{FUNDAMENTAL_HEIGHT_PX}px (3x)")

    # Generate HTML
    html = _generate_html(rows)

    # Convert to PNG
    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as f:
        img_path = f.name

    _html_to_png(html, img_path)

    # Find slide by title text
    target_slide = None

    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                if slide_title.lower() in shape.text.lower():
                    target_slide = slide
                    print(f"[Fundamental Rank] Found slide with '{slide_title}' at index {slide_idx + 1}")
                    break
        if target_slide:
            break

    if not target_slide:
        print(f"[Fundamental Rank] No slide with title '{slide_title}' found")
        Path(img_path).unlink(missing_ok=True)
        return prs

    # Insert image at FIXED position
    target_slide.shapes.add_picture(
        img_path,
        left=Cm(FUNDAMENTAL_LEFT_CM),
        top=Cm(FUNDAMENTAL_TOP_CM),
        width=Cm(FUNDAMENTAL_WIDTH_CM),
        height=Cm(FUNDAMENTAL_HEIGHT_CM)
    )

    # Cleanup
    Path(img_path).unlink(missing_ok=True)

    print(f"[Fundamental Rank] ✅ Table inserted at ({FUNDAMENTAL_LEFT_CM}, {FUNDAMENTAL_TOP_CM}) cm")

    return prs


def generate_fundamental_slide(
    prs: Presentation,
    excel_path: str,
    slide_title: str = "Fundamental Analysis"
) -> Presentation:
    """
    Generate Fundamental Rank slide from Excel data.

    Parameters
    ----------
    prs : Presentation
        PowerPoint presentation object
    excel_path : str
        Path to Excel file with data_fundamental sheet
    slide_title : str
        Title text to search for to find the correct slide

    Returns
    -------
    Presentation
        Modified presentation
    """
    # Compute ranks
    rows = compute_fundamental_ranks(excel_path)

    if not rows:
        print("[Fundamental Rank] No data found in data_fundamental sheet")
        return prs

    # Insert into slide
    return insert_fundamental_rank(prs, rows, slide_title)
