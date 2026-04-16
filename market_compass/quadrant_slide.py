"""Quadrant chart slide: Breadth Rank vs Fundamental Rank scatter plot with WoW arrows."""

from __future__ import annotations

import json
import re
import tempfile
from pathlib import Path
from typing import Dict, Optional, Tuple

import pandas as pd
from pptx import Presentation
from pptx.util import Cm

from helpers.html_to_image import render_html_to_image

# ── Chart dimensions (1050×500 at 4x) ────────────────────────────────────
CHART_WIDTH = 1050
CHART_HEIGHT = 500
SCALE_FACTOR = 4
RENDER_WIDTH = CHART_WIDTH * SCALE_FACTOR    # 4200
RENDER_HEIGHT = CHART_HEIGHT * SCALE_FACTOR  # 2000

# ── PPTX placement ───────────────────────────────────────────────────────
LEFT_CM = 2.20
TOP_CM = 4.94
WIDTH_CM = 21.0
HEIGHT_CM = 10.0

# ── 9 equity indices with chart styling ──────────────────────────────────
QUADRANT_INDICES = [
    {"name": "IBOV",       "label": "Ibovespa",   "flag": "BR", "color": "#2EA043"},
    {"name": "Nikkei 225", "label": "Nikkei 225", "flag": "JP", "color": "#7B2D8E"},
    {"name": "SMI",        "label": "SMI",        "flag": "CH", "color": "#636363"},
    {"name": "MEXBOL",     "label": "Mexbol",     "flag": "MX", "color": "#0288D1"},
    {"name": "S&P 500",    "label": "S&P 500",    "flag": "US", "color": "#1A237E"},
    {"name": "Sensex",     "label": "Sensex",     "flag": "IN", "color": "#E65100"},
    {"name": "DAX",        "label": "Dax",        "flag": "DE", "color": "#D32F2F"},
    {"name": "CSI 300",    "label": "CSI 300",    "flag": "CN", "color": "#E64A19"},
    {"name": "TASI",       "label": "Tasi",       "flag": "SA", "color": "#2E7D32"},
]

# Map from fundamental slide display names → instrument names
_FUND_DISPLAY_TO_INSTRUMENT = {
    "U.S.": "S&P 500",
    "China": "CSI 300",
    "Japan": "Nikkei 225",
    "Saudi Arabia": "TASI",
    "India": "Sensex",
    "Germany": "DAX",
    "Switzerland": "SMI",
    "Brazil": "IBOV",
    "Mexico": "MEXBOL",
}


def _load_prev_ranks(
    history_path: str, ic_date: str
) -> Dict[str, Tuple[int, int]]:
    """Load previous week's (breadth_rank, fundamental_rank) from history.json.

    Returns dict keyed by instrument name → (breadth_rank, fund_rank).
    """
    prev: Dict[str, Tuple[int, int]] = {}
    path = Path(history_path)
    if not path.exists():
        return prev

    with open(path) as f:
        history = json.load(f)

    ic_ts = pd.Timestamp(ic_date).date()

    for name, entries in history.items():
        sorted_entries = sorted(
            entries, key=lambda x: x.get("date", ""), reverse=True
        )
        for entry in sorted_entries:
            try:
                edate = pd.Timestamp(entry["date"]).date()
            except Exception:
                continue
            if edate < ic_ts:
                br = entry.get("breadth_rank")
                fr = entry.get("fundamental_rank")
                if br is not None and fr is not None:
                    prev[name] = (int(br), int(fr))
                break

    return prev


def _flag_data_uri(code: str) -> str:
    """Return a base64 data URI for the flag PNG, or empty string if not found."""
    import base64
    flag_path = Path(__file__).resolve().parent.parent / "assets" / "flags" / f"{code.lower()}.png"
    if not flag_path.exists():
        return ""
    b64 = base64.b64encode(flag_path.read_bytes()).decode("ascii")
    return f"data:image/png;base64,{b64}"


def _build_data_block(
    breadth_ranks: Dict[str, int],
    fundamental_ranks: Dict[str, int],
    prev_ranks: Dict[str, Tuple[int, int]],
) -> str:
    """Build the JavaScript `const indices = [...]` block for injection."""
    lines = []
    for idx in QUADRANT_INDICES:
        name = idx["name"]
        br = breadth_ranks.get(name, 5)
        fr = fundamental_ranks.get(name, 5)
        pbr, pfr = prev_ranks.get(name, (br, fr))  # default: no arrow
        flag_uri = _flag_data_uri(idx["flag"])

        lines.append(
            f"  {{ label: '{idx['label']}', flag: '{idx['flag']}', "
            f"flagImg: '{flag_uri}', "
            f"prev: [{pbr}, {pfr}], now: [{br}, {fr}], color: '{idx['color']}' }}"
        )

    block = "const indices = [\n" + ",\n".join(lines) + "\n];\n"
    # Pre-load flag images so they're ready for canvas drawImage
    block += """
// Pre-load flag images
const flagImages = {};
let flagsLoaded = 0;
indices.forEach(idx => {
  if (idx.flagImg) {
    const img = new Image();
    img.onload = () => { flagsLoaded++; };
    img.onerror = () => { flagsLoaded++; };
    img.src = idx.flagImg;
    flagImages[idx.flag] = img;
  } else {
    flagsLoaded++;
  }
});
"""
    return block


def generate_quadrant_slide(
    prs: Presentation,
    breadth_ranks: Dict[str, int],
    fundamental_ranks: Dict[str, int],
    history_path: str,
    ic_date: str,
    slide_name: str = "slide_quadrant",
) -> Presentation:
    """Render quadrant chart and insert into the PPTX slide.

    Parameters
    ----------
    prs : Presentation
        PowerPoint presentation being assembled.
    breadth_ranks : dict
        Instrument name → breadth rank (1-9), matching the breadth table.
    fundamental_ranks : dict
        Instrument name → fundamental rank (1-9), matching the fundamental table.
    history_path : str
        Path to history.json for WoW arrows.
    ic_date : str
        Current IC date (YYYY-MM-DD).
    slide_name : str
        Shape name to locate the target slide.
    """
    print(f"[Quadrant] Generating quadrant chart...")

    # Load previous ranks from history
    prev_ranks = _load_prev_ranks(history_path, ic_date)
    if prev_ranks:
        print(f"[Quadrant] Loaded previous ranks for {len(prev_ranks)} indices")
    else:
        print(f"[Quadrant] No previous ranks — arrows will be hidden")

    # Read HTML template
    template_path = Path(__file__).parent / "templates" / "quadrant_chart.html"
    html = template_path.read_text(encoding="utf-8")

    # Build and inject data block
    data_block = _build_data_block(breadth_ranks, fundamental_ranks, prev_ranks)
    html = re.sub(
        r"// ========== DATA.*?// ========== END DATA ==========",
        f"// ========== DATA — injected by assemble ==========\n{data_block}\n// ========== END DATA ==========",
        html,
        flags=re.DOTALL,
    )

    # Render to PNG via Playwright
    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as f:
        img_path = f.name

    print(f"[Quadrant] Rendering at {RENDER_WIDTH}x{RENDER_HEIGHT}px (4x scale)...")
    render_html_to_image(
        html_content=html,
        output_path=img_path,
        size=(CHART_WIDTH, CHART_HEIGHT),
        filename="quadrant.png",
        device_scale_factor=SCALE_FACTOR,
    )

    # Find target slide by shape name
    target_slide = None
    slide_name_lower = slide_name.lower().strip()

    print(f"[Quadrant] Searching for shape '{slide_name}'...")
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, "name") and shape.name:
                if shape.name.lower().strip() == slide_name_lower:
                    target_slide = slide
                    print(f"[Quadrant] Found slide at index {slide_idx + 1}")
                    break
        if target_slide:
            break

    if not target_slide:
        print(f"[Quadrant] ERROR: No slide with shape '{slide_name}' found")
        Path(img_path).unlink(missing_ok=True)
        return prs

    # Insert image
    picture = target_slide.shapes.add_picture(
        img_path,
        left=Cm(LEFT_CM),
        top=Cm(TOP_CM),
        width=Cm(WIDTH_CM),
        height=Cm(HEIGHT_CM),
    )

    # Send to back
    spTree = target_slide.shapes._spTree
    sp = picture._element
    spTree.remove(sp)
    spTree.insert(2, sp)

    Path(img_path).unlink(missing_ok=True)

    # Log positions
    for idx_cfg in QUADRANT_INDICES:
        name = idx_cfg["name"]
        br = breadth_ranks.get(name, "?")
        fr = fundamental_ranks.get(name, "?")
        print(f"[Quadrant]   {name}: breadth={br}, fundamental={fr}")

    print(f"[Quadrant] OK: Chart inserted at ({LEFT_CM}, {TOP_CM}) cm")
    return prs
