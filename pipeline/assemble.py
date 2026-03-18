"""
Stage: assemble

Reads complete draft_state.json (with subtitles filled by Claude Code)
and builds the final PowerPoint presentation by calling the same battle-tested
functions from ic_engine.py / technical_analysis / performance / market_compass.
"""
from __future__ import annotations

import json
import logging
import os
from datetime import datetime
from io import BytesIO
from pathlib import Path
from typing import Optional

import pandas as pd
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE

log = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Instrument → placeholder mapping (mirrors ic_engine.INSTRUMENTS)
# ---------------------------------------------------------------------------
INSTRUMENT_CONFIG = {
    "S&P 500":   {"ticker_key": "spx",       "bbg_ticker": "SPX Index",       "placeholder": "spx_v2"},
    "CSI 300":   {"ticker_key": "csi",       "bbg_ticker": "SHSZ300 Index",   "placeholder": "csi_v2"},
    "Nikkei 225":{"ticker_key": "nikkei",    "bbg_ticker": "NKY Index",       "placeholder": "nikkei_v2"},
    "TASI":      {"ticker_key": "tasi",      "bbg_ticker": "SASEIDX Index",   "placeholder": "tasi_v2"},
    "Sensex":    {"ticker_key": "sensex",    "bbg_ticker": "SENSEX Index",    "placeholder": "sensex_v2"},
    "DAX":       {"ticker_key": "dax",       "bbg_ticker": "DAX Index",       "placeholder": "dax_v2"},
    "SMI":       {"ticker_key": "smi",       "bbg_ticker": "SMI Index",       "placeholder": "smi_v2"},
    "IBOV":      {"ticker_key": "ibov",      "bbg_ticker": "IBOV Index",      "placeholder": "ibov_v2"},
    "MEXBOL":    {"ticker_key": "mexbol",    "bbg_ticker": "MEXBOL Index",    "placeholder": "mexbol_v2"},
    "Gold":      {"ticker_key": "gold",      "bbg_ticker": "GCA Comdty",      "placeholder": "gold_v2"},
    "Silver":    {"ticker_key": "silver",    "bbg_ticker": "SIA Comdty",      "placeholder": "silver_v2"},
    "Platinum":  {"ticker_key": "platinum",  "bbg_ticker": "XPT Comdty",      "placeholder": "platinum_v2"},
    "Palladium": {"ticker_key": "palladium", "bbg_ticker": "XPD Curncy",      "placeholder": "palladium_v2"},
    "Oil":       {"ticker_key": "oil",       "bbg_ticker": "CL1 Comdty",      "placeholder": "oil_v2"},
    "Copper":    {"ticker_key": "copper",    "bbg_ticker": "LP1 Comdty",      "placeholder": "copper_v2"},
    "Bitcoin":   {"ticker_key": "bitcoin",   "bbg_ticker": "XBTUSD Curncy",   "placeholder": "bitcoin_v2"},
    "Ethereum":  {"ticker_key": "ethereum",  "bbg_ticker": "XETUSD Curncy",   "placeholder": "ethereum_v2"},
    "Ripple":    {"ticker_key": "ripple",    "bbg_ticker": "XRPUSD Curncy",   "placeholder": "ripple_v2"},
    "Solana":    {"ticker_key": "solana",    "bbg_ticker": "XSOUSD Curncy",   "placeholder": "solana_v2"},
    "Binance":   {"ticker_key": "binance",   "bbg_ticker": "XBIUSD Curncy",   "placeholder": "binance_v2"},
}


def _load_draft(draft_path: str) -> dict:
    with open(draft_path) as f:
        return json.load(f)


def run_assemble(
    draft_path: str = "draft_state.json",
    template_path: Optional[str] = None,
    output_path: Optional[str] = None,
    history_path: str = "market_compass/data/history.json",
    config_path: str = "config/tickers.yaml",
) -> str:
    """
    Build the final PowerPoint from draft_state.json.

    Uses the same functions as ic_engine.py:
    - insert_technical_analysis_v2_slide for chart + subtitle + view insertion
    - Performance chart functions for equity, FX, crypto, rates, credit, commodity
    - Market compass slides (technical summary, breadth, fundamentals)

    Returns: path to the output PPTX file.
    """
    draft = _load_draft(draft_path)
    ic_date = draft["date"]
    instruments = draft["instruments"]
    ytd_subtitles = draft.get("ytd_subtitles", {})

    log.info("Assembling presentation for date: %s", ic_date)

    # Validate subtitles
    missing_subtitles = [nm for nm, d in instruments.items() if not d.get("subtitle")]
    if missing_subtitles:
        log.warning("Missing subtitles: %s", ", ".join(missing_subtitles))

    # Resolve paths
    dropbox_path = os.environ.get(
        "IC_DROPBOX_PATH",
        "/Users/larazanella/Library/CloudStorage/Dropbox/Tools_In_Construction/ic",
    )
    if template_path is None:
        template_path = str(Path(dropbox_path) / "shadow_template.pptx")
    if output_path is None:
        date_str = ic_date.replace("-", "")
        output_path = str(Path(dropbox_path) / f"Market_Compass_{date_str}.pptx")

    master_csv = str(Path(dropbox_path) / "master_prices.csv")
    excel_path = str(Path(dropbox_path) / "ic_file.xlsx")

    if not Path(template_path).exists():
        raise FileNotFoundError(f"Template not found: {template_path}")

    # Load template
    prs = Presentation(template_path)
    log.info("Loaded template: %s (%d slides)", template_path, len(prs.slides))

    # Create temp Excel for chart/performance functions
    from data_loader import create_temp_excel_from_csv, load_prices_from_csv
    from technical_analysis.common_helpers import clear_excel_cache

    clear_excel_cache()

    data_as_of = pd.Timestamp(ic_date).date()
    temp_excel = create_temp_excel_from_csv(
        Path(master_csv), Path(excel_path), data_as_of
    )
    chart_excel_path = str(temp_excel)
    log.info("Created temp Excel: %s", temp_excel)

    # Load prices DataFrame for summary slides
    df_prices = load_prices_from_csv(Path(master_csv), data_as_of)

    # =====================================================================
    # 1. Update [DataIC] date placeholder
    # =====================================================================
    _update_date_placeholder(prs, ic_date)

    # =====================================================================
    # 2. Technical analysis slides (20 instruments)
    # =====================================================================
    from technical_analysis.equity.spx import (
        create_technical_analysis_v2_chart,
        insert_technical_analysis_v2_slide,
    )

    # Load previous week's scores from history.json for WoW delta and arrows
    prev_dmas:  dict[str, int] = {}
    prev_tech:  dict[str, int] = {}
    prev_mom:   dict[str, int] = {}
    prev_rsi:   dict[str, int] = {}
    try:
        hist_path = Path(history_path)
        if hist_path.exists():
            with open(hist_path) as _hf:
                _hist = json.load(_hf)
            _ic_ts = pd.Timestamp(ic_date).date()
            for _name, _entries in _hist.items():
                _sorted = sorted(_entries, key=lambda x: x.get("date", ""), reverse=True)
                for _e in _sorted:
                    try:
                        _edate = pd.Timestamp(_e["date"]).date()
                    except Exception:
                        continue
                    if _edate < _ic_ts and _e.get("dmas") is not None:
                        prev_dmas[_name] = int(_e["dmas"])
                        if _e.get("technical_score") is not None:
                            prev_tech[_name] = int(_e["technical_score"])
                        if _e.get("momentum_score") is not None:
                            prev_mom[_name] = int(_e["momentum_score"])
                        if _e.get("rsi") is not None:
                            prev_rsi[_name] = int(_e["rsi"])
                        break
            log.info("Loaded WoW scores for %d instruments from history.json", len(prev_dmas))
    except Exception as _e:
        log.warning("Could not load history.json for WoW deltas: %s", _e)

    total = len(instruments)
    for idx, (name, data) in enumerate(instruments.items()):
        if name not in INSTRUMENT_CONFIG:
            log.warning("Unknown instrument %s — skipping", name)
            continue

        cfg = INSTRUMENT_CONFIG[name]
        log.info("Processing %s (%d/%d)...", name, idx + 1, total)

        try:
            # Generate chart via the V2 chart function (reads from temp Excel)
            chart_bytes, used_date = create_technical_analysis_v2_chart(
                chart_excel_path,
                ticker=cfg["bbg_ticker"],
                price_mode="Last Price",
                dmas_score=data.get("dmas"),
                dmas_prev_week=prev_dmas.get(name),
                technical_score=data.get("technical"),
                technical_prev_week=prev_tech.get(name),
                momentum_score=data.get("momentum"),
                momentum_prev_week=prev_mom.get(name),
                rsi_prev_week=prev_rsi.get(name),
            )

            if not chart_bytes:
                log.warning("No chart generated for %s", name)
                continue

            # Build view text from rating
            rating = data.get("rating", "")
            view_text = f"{name}: {rating}" if rating else ""

            # Subtitle
            subtitle_text = data.get("subtitle") or ""

            # Insert into slide using the battle-tested V2 function
            prs = insert_technical_analysis_v2_slide(
                prs,
                chart_bytes,
                used_date=used_date,
                price_mode="Last Price",
                placeholder_name=cfg["placeholder"],
                view_text=view_text,
                subtitle_text=subtitle_text,
            )

            log.info("Inserted slide for %s", name)

        except Exception as e:
            log.warning("Failed to process %s: %s", name, e)

    # =====================================================================
    # 3. Performance slides
    # =====================================================================
    _insert_performance_slides(prs, chart_excel_path, ytd_subtitles)

    # =====================================================================
    # 4. Summary slides (technical nutshell, breadth, fundamentals)
    # =====================================================================
    breadth_records = draft.get("breadth", [])
    breadth_ranks, fundamental_ranks = _insert_summary_slides(
        prs, df_prices, instruments, chart_excel_path, breadth_records,
        ic_date=ic_date, history_path=history_path,
    )

    # =====================================================================
    # 5. Finalize
    # =====================================================================
    _force_all_tables_calibri(prs)
    _disable_image_compression(prs)

    prs.save(output_path)
    log.info("Saved presentation: %s", output_path)

    # Clean up temp Excel
    try:
        os.remove(temp_excel)
    except Exception:
        pass

    # Update history.json (includes breadth/fundamental ranks for quadrant WoW)
    _update_history(history_path, ic_date, instruments, breadth_ranks, fundamental_ranks)

    print(f"\n✓ Presentation assembled: {output_path}")
    print(f"  Date: {ic_date}")
    print(f"  Instruments: {len(instruments)}")
    return output_path


# =========================================================================
# Helper: date placeholder
# =========================================================================

def _update_date_placeholder(prs, ic_date: str) -> None:
    """Replace [DataIC] placeholder with formatted date."""
    ts = pd.Timestamp(ic_date)
    human_date = f"{ts.strftime('%B')} {ts.day}, {ts.year}"

    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.strip() == "[DataIC]":
                        size = run.font.size
                        bold = run.font.bold
                        italic = run.font.italic
                        run.text = human_date
                        if size:
                            run.font.size = size
                        if bold is not None:
                            run.font.bold = bold
                        if italic is not None:
                            run.font.italic = italic
                        return


# =========================================================================
# Helper: performance slides
# =========================================================================

def _insert_performance_slides(prs, chart_excel_path: str, ytd_subtitles: dict = None) -> None:
    """Insert all performance slides (equity, FX, crypto, rates, credit, commodity)."""
    if ytd_subtitles is None:
        ytd_subtitles = {}
    price_mode = "Last Price"

    # --- Equity ---
    try:
        from performance.equity_perf import (
            create_weekly_performance_chart,
            create_historical_performance_table,
            insert_equity_performance_bar_slide,
            insert_equity_performance_histo_slide,
            create_equity_ytd_evolution_chart,
            insert_equity_ytd_evolution_slide,
            create_fx_impact_analysis_chart_eur,
            insert_fx_impact_analysis_slide_eur,
            create_fx_impact_analysis_chart_chf,
            insert_fx_impact_analysis_slide_chf,
        )
        bar_bytes, date1 = create_weekly_performance_chart(chart_excel_path, price_mode=price_mode)
        prs = insert_equity_performance_bar_slide(
            prs, bar_bytes, used_date=date1, price_mode=price_mode,
            left_cm=3.47, top_cm=5.28, width_cm=17.31, height_cm=10
        )
        histo_bytes, date2 = create_historical_performance_table(chart_excel_path, price_mode=price_mode)
        prs = insert_equity_performance_histo_slide(
            prs, histo_bytes, used_date=date2, price_mode=price_mode,
            left_cm=2.16, top_cm=4.70, width_cm=19.43, height_cm=10.61
        )
        ytd_bytes, date3 = create_equity_ytd_evolution_chart(chart_excel_path, price_mode=price_mode)
        prs = insert_equity_ytd_evolution_slide(prs, ytd_bytes, used_date=date3, price_mode=price_mode, subtitle=ytd_subtitles.get("equity"))

        fx_eur_bytes, date4 = create_fx_impact_analysis_chart_eur(chart_excel_path, price_mode=price_mode)
        prs = insert_fx_impact_analysis_slide_eur(prs, fx_eur_bytes, used_date=date4, price_mode=price_mode)
        fx_chf_bytes, date5 = create_fx_impact_analysis_chart_chf(chart_excel_path, price_mode=price_mode)
        prs = insert_fx_impact_analysis_slide_chf(prs, fx_chf_bytes, used_date=date5, price_mode=price_mode)
        log.info("Equity performance slides inserted")
    except Exception as e:
        log.warning("Equity performance error: %s", e)

    # --- FX ---
    try:
        from performance.fx_perf import (
            create_weekly_html_performance_chart as create_weekly_fx_html_chart,
            insert_fx_weekly_html_slide,
            create_historical_html_performance_chart as create_historical_fx_html_chart,
            insert_fx_historical_html_slide,
        )
        fx_bar, fxd1 = create_weekly_fx_html_chart(chart_excel_path, price_mode=price_mode)
        prs = insert_fx_weekly_html_slide(prs, fx_bar, used_date=fxd1, price_mode=price_mode)
        fx_histo, fxd2 = create_historical_fx_html_chart(chart_excel_path, price_mode=price_mode)
        prs = insert_fx_historical_html_slide(prs, fx_histo, used_date=fxd2, price_mode=price_mode)
        log.info("FX performance slides inserted")
    except Exception as e:
        log.warning("FX performance error: %s", e)

    # --- Crypto ---
    try:
        from performance.crypto_perf import (
            create_weekly_html_performance_chart as create_weekly_crypto_html_chart,
            insert_crypto_weekly_html_slide,
            create_historical_html_performance_chart as create_historical_crypto_html_chart,
            insert_crypto_historical_html_slide,
            create_crypto_ytd_evolution_chart,
            insert_crypto_ytd_evolution_slide,
        )
        cr_bar, crd1 = create_weekly_crypto_html_chart(chart_excel_path, price_mode=price_mode)
        prs = insert_crypto_weekly_html_slide(prs, cr_bar, used_date=crd1, price_mode=price_mode)
        cr_histo, crd2 = create_historical_crypto_html_chart(chart_excel_path, price_mode=price_mode)
        prs = insert_crypto_historical_html_slide(prs, cr_histo, used_date=crd2, price_mode=price_mode)
        cr_ytd, crd3 = create_crypto_ytd_evolution_chart(chart_excel_path, price_mode=price_mode)
        prs = insert_crypto_ytd_evolution_slide(prs, cr_ytd, used_date=crd3, price_mode=price_mode, subtitle=ytd_subtitles.get("crypto"))
        log.info("Crypto performance slides inserted")
    except Exception as e:
        log.warning("Crypto performance error: %s", e)

    # --- Rates ---
    try:
        from performance.rates_perf import (
            create_weekly_performance_chart as create_weekly_rates_chart,
            create_historical_performance_table as create_historical_rates_table,
            insert_rates_performance_bar_slide,
            insert_rates_performance_histo_slide,
        )
        rt_bar, rtd1 = create_weekly_rates_chart(chart_excel_path, price_mode=price_mode)
        prs = insert_rates_performance_bar_slide(
            prs, rt_bar, used_date=rtd1, price_mode=price_mode,
            left_cm=3.35, top_cm=4.6, width_cm=17.02
        )
        rt_histo, rtd2 = create_historical_rates_table(chart_excel_path, price_mode=price_mode)
        prs = insert_rates_performance_histo_slide(
            prs, rt_histo, used_date=rtd2, price_mode=price_mode,
            left_cm=3.35, top_cm=4.6, width_cm=17.02
        )
        log.info("Rates performance slides inserted")
    except Exception as e:
        log.warning("Rates performance error: %s", e)

    # --- Credit ---
    try:
        from performance.corp_bonds_perf import (
            create_weekly_performance_chart as create_weekly_credit_chart,
            insert_corp_bonds_performance_slide as insert_credit_bar,
            create_historical_performance_chart as create_historical_credit_chart,
            insert_corp_bonds_historical_slide as insert_credit_histo,
        )
        cd_bar, cdd1 = create_weekly_credit_chart(chart_excel_path, price_mode=price_mode)
        prs = insert_credit_bar(
            prs, cd_bar, used_date=cdd1, price_mode=price_mode,
            left_cm=1.63, top_cm=4.73, width_cm=22.48, height_cm=10.61
        )
        cd_histo, cdd2 = create_historical_credit_chart(chart_excel_path, price_mode=price_mode)
        prs = insert_credit_histo(prs, cd_histo, used_date=cdd2, price_mode=price_mode)
        log.info("Credit performance slides inserted")
    except Exception as e:
        log.warning("Credit performance error: %s", e)

    # --- Commodity ---
    try:
        from performance.commodity_perf import (
            create_weekly_html_performance_chart as create_weekly_commo_chart,
            insert_commodity_weekly_html_slide,
            create_historical_html_performance_chart as create_historical_commo_chart,
            insert_commodity_historical_html_slide,
            create_commodity_ytd_evolution_chart,
            insert_commodity_ytd_evolution_slide,
        )
        co_bar, cod1 = create_weekly_commo_chart(chart_excel_path, price_mode=price_mode)
        prs = insert_commodity_weekly_html_slide(prs, co_bar, used_date=cod1, price_mode=price_mode)
        co_histo, cod2 = create_historical_commo_chart(chart_excel_path, price_mode=price_mode)
        prs = insert_commodity_historical_html_slide(prs, co_histo, used_date=cod2, price_mode=price_mode)
        co_ytd, cod3 = create_commodity_ytd_evolution_chart(chart_excel_path, price_mode=price_mode)
        prs = insert_commodity_ytd_evolution_slide(prs, co_ytd, used_date=cod3, price_mode=price_mode, subtitle=ytd_subtitles.get("commodity"))
        log.info("Commodity performance slides inserted")
    except Exception as e:
        log.warning("Commodity performance error: %s", e)


# =========================================================================
# Helper: summary slides
# =========================================================================

def _insert_summary_slides(
    prs, df_prices, instruments, chart_excel_path,
    breadth_records=None, ic_date: str = "", history_path: str = "",
) -> tuple[dict[str, int], dict[str, int]]:
    """Insert technical summary, breadth, fundamental, and quadrant slides.

    Returns (breadth_ranks, fundamental_ranks) for history tracking.
    """
    from utils import adjust_prices_for_mode

    # Technical Analysis summary
    try:
        from market_compass.technical_slide import prepare_slide_data, insert_technical_analysis_slide

        df_prices_adj, tech_used_date = adjust_prices_for_mode(df_prices, "Last Price")
        dmas_scores = {}
        for name, data in instruments.items():
            cfg = INSTRUMENT_CONFIG.get(name)
            if cfg:
                dmas_scores[cfg["ticker_key"]] = data.get("dmas", 50)

        rows = prepare_slide_data(df_prices_adj, dmas_scores, chart_excel_path, price_mode="Last Price")
        if rows:
            insert_technical_analysis_slide(
                prs, rows, placeholder_name="technical_nutshell",
                used_date=tech_used_date, price_mode="Last Price"
            )
        log.info("Technical summary slide inserted")
    except Exception as e:
        log.warning("Technical summary error: %s", e)

    # Breadth slide (new composite breadth table)
    breadth_ranks: dict[str, int] = {}
    try:
        from market_compass.breadth_slide import generate_composite_breadth_slide
        prs = generate_composite_breadth_slide(
            prs, breadth_records=breadth_records or [], slide_name="slide_breadth"
        )
        # Extract breadth ranks from records (matches the table)
        for rec in (breadth_records or []):
            breadth_ranks[rec["name"]] = int(rec["rank"])
        log.info("Breadth slide inserted")
    except Exception as e:
        log.warning("Breadth slide error: %s", e)

    # Fundamental slide
    fundamental_ranks: dict[str, int] = {}
    try:
        from market_compass.fundamental_slide import generate_fundamental_slide
        from market_compass.quadrant_slide import _FUND_DISPLAY_TO_INSTRUMENT
        prs, fund_display_ranks = generate_fundamental_slide(prs, excel_path=chart_excel_path, slide_name="slide_fundamentals")
        # Map display names (U.S., Japan, ...) → instrument names (S&P 500, Nikkei 225, ...)
        for display_name, rank in fund_display_ranks.items():
            instr_name = _FUND_DISPLAY_TO_INSTRUMENT.get(display_name)
            if instr_name:
                fundamental_ranks[instr_name] = int(rank)
        log.info("Fundamental slide inserted")
    except Exception as e:
        log.warning("Fundamental slide error: %s", e)

    # Quadrant slide (breadth rank vs fundamental rank scatter)
    try:
        from market_compass.quadrant_slide import generate_quadrant_slide
        prs = generate_quadrant_slide(
            prs,
            breadth_ranks=breadth_ranks,
            fundamental_ranks=fundamental_ranks,
            history_path=history_path,
            ic_date=ic_date,
        )
        log.info("Quadrant slide inserted")
    except Exception as e:
        log.warning("Quadrant slide error: %s", e)

    return breadth_ranks, fundamental_ranks


# =========================================================================
# Helpers: finalization (from ic_engine.py)
# =========================================================================

def _force_all_tables_calibri(prs, size_pt: int = 11):
    """Enforce Calibri font on all tables."""
    try:
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if hasattr(cell, "text_frame") and cell.text_frame:
                                for p in cell.text_frame.paragraphs:
                                    for r in p.runs:
                                        r.font.name = "Calibri"
                                        r.font.size = Pt(size_pt)
    except Exception:
        pass


def _disable_image_compression(prs):
    """Disable image compression to preserve chart quality."""
    DRAWING_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    try:
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, '_element'):
                    for blip in shape._element.iter(f'{DRAWING_NS}blip'):
                        blip.set(f'{DRAWING_NS}cstate', 'none')
    except Exception:
        pass


# =========================================================================
# Helper: history.json
# =========================================================================

def _update_history(
    history_path: str, ic_date: str, instruments: dict,
    breadth_ranks: dict | None = None, fundamental_ranks: dict | None = None,
) -> None:
    """Update history.json with current week's scores and ranks."""
    path = Path(history_path)
    if path.exists():
        with open(path) as f:
            history = json.load(f)
    else:
        history = {}

    if breadth_ranks is None:
        breadth_ranks = {}
    if fundamental_ranks is None:
        fundamental_ranks = {}

    for name, data in instruments.items():
        if name not in history:
            history[name] = []

        entry = {
            "date":               ic_date,
            "dmas":               data.get("dmas"),
            "technical_score":    data.get("technical"),
            "momentum_score":     data.get("momentum"),
            "price_vs_50ma_pct":  _parse_pct(data.get("vs_50d", "0%")),
            "price_vs_100ma_pct": _parse_pct(data.get("vs_100d", "0%")),
            "price_vs_200ma_pct": _parse_pct(data.get("vs_200d", "0%")),
            "rating":             data.get("rating"),
        }
        # Store ranks for quadrant WoW arrows (only for equity indices)
        if name in breadth_ranks:
            entry["breadth_rank"] = breadth_ranks[name]
        if name in fundamental_ranks:
            entry["fundamental_rank"] = fundamental_ranks[name]

        # Remove duplicate entries for the same date
        history[name] = [h for h in history[name] if h.get("date") != ic_date]
        history[name].append(entry)
        # Keep last 52 weeks
        history[name] = sorted(history[name], key=lambda x: x["date"])[-52:]

    with open(path, "w") as f:
        json.dump(history, f, indent=2)

    log.info("Updated history.json (%d instruments)", len(instruments))


def _parse_pct(pct_str: str) -> float:
    """Parse '+1.5%' or '-2.3%' to float."""
    try:
        return float(pct_str.replace("%", "").replace("+", ""))
    except (ValueError, AttributeError):
        return 0.0
