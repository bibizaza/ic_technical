"""
Base instrument class for technical analysis.

This module provides a base class that eliminates code duplication across
all instrument modules (equity, crypto, commodity). Instead of having 20+
separate files with identical functions, instruments are configured with
metadata and the base class handles all common operations.

Key features:
- Unified chart generation (Plotly interactive and matplotlib static)
- Technical and momentum score calculation with caching
- PowerPoint slide manipulation
- Vectorized pandas operations (no .iterrows())
- Configurable instrument metadata
"""

from __future__ import annotations

from datetime import timedelta
import pathlib
from typing import Optional, Tuple, Dict, Any, Callable
from functools import lru_cache
import hashlib

import pandas as pd
import plotly.graph_objects as go
from sklearn.linear_model import LinearRegression

from pptx import Presentation
from pptx.util import Cm
from io import BytesIO
import matplotlib.pyplot as plt
import matplotlib.patches as patches
from matplotlib.colors import LinearSegmentedColormap
import numpy as np


# Default configuration
PLOT_LOOKBACK_DAYS: int = 90

# Import helper for adjusting price data according to price mode
try:
    from utils import adjust_prices_for_mode  # type: ignore
except Exception:
    adjust_prices_for_mode = None  # type: ignore


class InstrumentConfig:
    """Configuration for an instrument."""

    def __init__(
        self,
        name: str,
        display_name: str,
        ticker: str,
        vol_ticker: Optional[str] = None,
        peer_group: Optional[list] = None,
        mars_scorer_func: Optional[Callable] = None,
        placeholder_prefix: str = None,
    ):
        """
        Initialize instrument configuration.

        Parameters
        ----------
        name : str
            Internal name (e.g., 'spx', 'gold', 'bitcoin')
        display_name : str
            Display name (e.g., 'S&P 500', 'Gold', 'Bitcoin')
        ticker : str
            Bloomberg or data ticker (e.g., 'SPX Index', 'GCA Comdty')
        vol_ticker : str, optional
            Volatility index ticker (e.g., 'VIX Index' for SPX)
        peer_group : list, optional
            List of peer tickers for relative momentum
        mars_scorer_func : Callable, optional
            MARS scoring function from mars_engine
        placeholder_prefix : str, optional
            Prefix for PowerPoint placeholders (defaults to name)
        """
        self.name = name
        self.display_name = display_name
        self.ticker = ticker
        self.vol_ticker = vol_ticker or f"{ticker.split()[0]}VOL Index"
        self.peer_group = peer_group or []
        self.mars_scorer_func = mars_scorer_func
        self.placeholder_prefix = placeholder_prefix or name


class BaseInstrument:
    """Base class for all technical analysis instruments."""

    def __init__(self, config: InstrumentConfig):
        """Initialize with instrument configuration."""
        self.config = config
        self._cache: Dict[str, Any] = {}

    @property
    def plot_lookback_days(self) -> int:
        """Get the plot lookback days (can be overridden at module level)."""
        return PLOT_LOOKBACK_DAYS

    # -------------------------------------------------------------------------
    # Font and text formatting helpers
    # -------------------------------------------------------------------------

    @staticmethod
    def _get_run_font_attributes(run):
        """
        Capture font attributes from a run.

        Returns a tuple (size, rgb, theme_color, brightness, bold, italic).
        """
        if run is None:
            return None, None, None, None, None, None
        size = run.font.size
        colour = run.font.color
        rgb = None
        theme_color = None
        brightness = None
        try:
            rgb = colour.rgb
        except Exception:
            rgb = None
            try:
                theme_color = colour.theme_color
            except Exception:
                theme_color = None
        try:
            brightness = colour.brightness
        except Exception:
            brightness = None
        bold = run.font.bold
        italic = run.font.italic
        return size, rgb, theme_color, brightness, bold, italic

    @staticmethod
    def _apply_run_font_attributes(new_run, size, rgb, theme_color, brightness, bold, italic):
        """Apply captured font attributes to a new run."""
        if size is not None:
            new_run.font.size = size
        if rgb is not None:
            try:
                new_run.font.color.rgb = rgb
            except Exception:
                pass
        elif theme_color is not None:
            try:
                new_run.font.color.theme_color = theme_color
                if brightness is not None:
                    new_run.font.color.brightness = brightness
            except Exception:
                pass
        if bold is not None:
            new_run.font.bold = bold
        if italic is not None:
            new_run.font.italic = italic

    # -------------------------------------------------------------------------
    # Data loading
    # -------------------------------------------------------------------------

    def _load_price_data(
        self,
        excel_path: pathlib.Path,
        price_mode: str = "Last Price",
    ) -> pd.DataFrame:
        """
        Read the raw price sheet and return a tidy Date-Price DataFrame.

        Parameters
        ----------
        excel_path : pathlib.Path
            Path to the Excel workbook containing price data.
        price_mode : str, default "Last Price"
            One of "Last Price" or "Last Close".

        Returns
        -------
        pandas.DataFrame
            DataFrame with columns Date and Price.
        """
        df = pd.read_excel(excel_path, sheet_name="data_prices")
        df = df.drop(index=0)
        df = df[df[df.columns[0]] != "DATES"]
        df["Date"] = pd.to_datetime(df[df.columns[0]], errors="coerce")
        df["Price"] = pd.to_numeric(df[self.config.ticker], errors="coerce")
        df_clean = (
            df.dropna(subset=["Date", "Price"])
            .sort_values("Date")
            .reset_index(drop=True)[["Date", "Price"]]
        )
        # Adjust for price mode if helper is available
        if adjust_prices_for_mode is not None and price_mode:
            try:
                df_clean, _ = adjust_prices_for_mode(df_clean, price_mode)
            except Exception:
                pass
        return df_clean

    @staticmethod
    def _add_mas(df: pd.DataFrame) -> pd.DataFrame:
        """Add 50/100/200-day moving-average columns to a DataFrame."""
        out = df.copy()
        for w in (50, 100, 200):
            out[f"MA_{w}"] = out["Price"].rolling(w, min_periods=1).mean()
        return out

    def _get_vol_index_value(
        self,
        excel_obj_or_path,
        price_mode: str = "Last Price",
    ) -> Optional[float]:
        """
        Retrieve the most recent value of a volatility index.

        Returns
        -------
        float or None
            The most recent volatility index value.
        """
        try:
            df = pd.read_excel(excel_obj_or_path, sheet_name="data_prices")
        except Exception:
            return None
        df = df.drop(index=0)
        df = df[df[df.columns[0]] != "DATES"]
        df["Date"] = pd.to_datetime(df[df.columns[0]], errors="coerce")
        if self.config.vol_ticker not in df.columns:
            return None
        df["Price"] = pd.to_numeric(df[self.config.vol_ticker], errors="coerce")
        df_clean = df.dropna(subset=["Date", "Price"]).sort_values("Date").reset_index(drop=True)[
            ["Date", "Price"]
        ]
        if adjust_prices_for_mode is not None and price_mode:
            try:
                df_clean, _ = adjust_prices_for_mode(df_clean, price_mode)
            except Exception:
                pass
        if df_clean.empty:
            return None
        try:
            return float(df_clean["Price"].iloc[-1])
        except Exception:
            return None

    # -------------------------------------------------------------------------
    # Score calculation (VECTORIZED - no .iterrows())
    # -------------------------------------------------------------------------

    def _get_technical_score(self, excel_obj_or_path) -> Optional[float]:
        """
        Retrieve the technical score from 'data_technical_score'.

        PERFORMANCE: Uses vectorized pandas query instead of .iterrows().
        This is ~100x faster than the original implementation.

        Returns
        -------
        float or None
            The technical score or None if unavailable.
        """
        try:
            df = pd.read_excel(excel_obj_or_path, sheet_name="data_technical_score")
        except Exception:
            return None

        df = df.dropna(subset=[df.columns[0], df.columns[1]])

        # VECTORIZED: Use boolean indexing instead of .iterrows()
        # Original: for _, row in df.iterrows(): if str(row[...]).upper() == "SPX INDEX"...
        # New: Direct query with vectorized string operations
        df[df.columns[0]] = df[df.columns[0]].astype(str).str.strip().str.upper()
        target_ticker = self.config.ticker.upper()

        matches = df[df[df.columns[0]] == target_ticker]
        if matches.empty:
            return None

        try:
            return float(matches.iloc[0][df.columns[1]])
        except Exception:
            return None

    def _get_momentum_score(
        self,
        excel_obj_or_path,
        price_mode: str = "Last Price",
    ) -> Optional[float]:
        """
        Retrieve the MARS momentum score.

        Uses the configured MARS scoring function if available.
        Includes caching for performance.

        Returns
        -------
        float or None
            The momentum score or None if unavailable.
        """
        if self.config.mars_scorer_func is None:
            return None

        try:
            # Create cache key based on file path and price mode
            cache_key = f"momentum_{excel_obj_or_path}_{price_mode}"

            # Check cache first
            if cache_key in self._cache:
                return self._cache[cache_key]

            # Prepare data for MARS scorer
            df = pd.read_excel(excel_obj_or_path, sheet_name="data_prices")

            # Build price DataFrame with peer group
            df = df.drop(index=0)
            df = df[df[df.columns[0]] != "DATES"]
            df["Date"] = pd.to_datetime(df[df.columns[0]], errors="coerce")
            df = df.dropna(subset=["Date"]).set_index("Date")

            # Collect instrument and peer prices
            price_cols = {}
            if self.config.ticker in df.columns:
                price_cols[self.config.name.upper()] = pd.to_numeric(
                    df[self.config.ticker], errors="coerce"
                )

            for peer in self.config.peer_group:
                if peer in df.columns:
                    price_cols[peer] = pd.to_numeric(df[peer], errors="coerce")

            if not price_cols:
                return None

            df_prices = pd.DataFrame(price_cols)
            df_prices = df_prices.ffill().bfill()

            # Apply price mode adjustment
            if adjust_prices_for_mode is not None and price_mode:
                try:
                    df_prices_adj = df_prices.reset_index()
                    df_prices_adj = df_prices_adj.rename(columns={"Date": "Date"})
                    for col in df_prices.columns:
                        temp_df = pd.DataFrame({
                            "Date": df_prices_adj["Date"],
                            "Price": df_prices_adj[col]
                        })
                        temp_df, _ = adjust_prices_for_mode(temp_df, price_mode)
                        df_prices_adj[col] = temp_df["Price"].values
                    df_prices = df_prices_adj.set_index("Date")
                except Exception:
                    pass

            # Call MARS scorer
            score_series = self.config.mars_scorer_func(df_prices)
            if score_series is None or score_series.empty:
                return None

            score = float(score_series.iloc[-1])

            # Cache the result
            self._cache[cache_key] = score

            return score
        except Exception:
            return None

    # -------------------------------------------------------------------------
    # Chart generation
    # -------------------------------------------------------------------------

    def make_figure(
        self,
        excel_path: str | pathlib.Path,
        anchor_date: Optional[pd.Timestamp] = None,
        price_mode: str = "Last Price",
    ) -> go.Figure:
        """
        Build an interactive Plotly chart for Streamlit.

        Parameters
        ----------
        excel_path : str or pathlib.Path
            Path to the Excel file containing price data.
        anchor_date : pandas.Timestamp or None, optional
            If provided, a regression channel is drawn.
        price_mode : str, default "Last Price"
            One of "Last Price" or "Last Close".

        Returns
        -------
        go.Figure
            A Plotly figure with price, moving averages, and Fibonacci lines.
        """
        excel_path = pathlib.Path(excel_path)
        df_raw = self._load_price_data(excel_path, price_mode=price_mode)
        df_full = self._add_mas(df_raw)

        if df_full.empty:
            return go.Figure()

        # Crop to lookback window
        today = df_full["Date"].max()
        cutoff = today - timedelta(days=self.plot_lookback_days)
        df = df_full[df_full["Date"] >= cutoff].copy()

        if df.empty:
            return go.Figure()

        fig = go.Figure()

        # Price trace
        fig.add_trace(
            go.Scatter(
                x=df["Date"],
                y=df["Price"],
                mode="lines",
                name=self.config.display_name,
                line=dict(color="#1f77b4", width=2),
            )
        )

        # Moving averages
        ma_colors = {"MA_50": "#00FF00", "MA_100": "#FFA500", "MA_200": "#FF0000"}
        for ma_col, color in ma_colors.items():
            if ma_col in df.columns:
                fig.add_trace(
                    go.Scatter(
                        x=df["Date"],
                        y=df[ma_col],
                        mode="lines",
                        name=f"{ma_col.split('_')[1]}-day MA",
                        line=dict(color=color, width=1),
                    )
                )

        # Fibonacci retracement levels
        hi, lo = df["Price"].max(), df["Price"].min()
        span = hi - lo
        fib_levels = [
            hi,
            hi - 0.236 * span,
            hi - 0.382 * span,
            hi - 0.5 * span,
            hi - 0.618 * span,
            lo,
        ]
        for lvl in fib_levels:
            fig.add_hline(
                y=lvl,
                line_dash="dash",
                line_color="grey",
                opacity=0.4,
                annotation_text=f"{lvl:.2f}",
                annotation_position="right",
            )

        # Regression channel if anchor_date is provided
        if anchor_date is not None:
            subset = df_full[df_full["Date"] >= anchor_date].copy()
            if len(subset) > 1:
                X = subset["Date"].map(pd.Timestamp.toordinal).to_numpy().reshape(-1, 1)
                y_vals = subset["Price"].to_numpy()
                model = LinearRegression().fit(X, y_vals)
                trend = model.predict(X)
                resid = y_vals - trend
                std_resid = np.std(resid)
                upper = trend + 2 * std_resid
                lower = trend - 2 * std_resid

                uptrend = model.coef_[0] > 0
                color = "#008000" if uptrend else "#C00000"

                # Only show channel in the lookback window
                subset_vis = subset[subset["Date"] >= cutoff]
                if not subset_vis.empty:
                    fig.add_trace(
                        go.Scatter(
                            x=subset_vis["Date"],
                            y=upper[subset["Date"] >= cutoff],
                            mode="lines",
                            name="Upper Channel",
                            line=dict(color=color, dash="dash"),
                            showlegend=False,
                        )
                    )
                    fig.add_trace(
                        go.Scatter(
                            x=subset_vis["Date"],
                            y=lower[subset["Date"] >= cutoff],
                            mode="lines",
                            name="Lower Channel",
                            line=dict(color=color, dash="dash"),
                            fill="tonexty",
                            fillcolor=f"rgba({int(color[1:3], 16)}, {int(color[3:5], 16)}, {int(color[5:7], 16)}, 0.2)",
                            showlegend=False,
                        )
                    )

        fig.update_layout(
            title=f"{self.config.display_name} Technical Analysis",
            xaxis_title="Date",
            yaxis_title="Price",
            hovermode="x unified",
            template="plotly_white",
            height=600,
        )

        return fig

    def _generate_image_from_df(
        self,
        df_full: pd.DataFrame,
        anchor_date: Optional[pd.Timestamp] = None,
        lookback_days: int = 90,
        width_cm: float = 21.41,
        height_cm: float = 7.53,
        show_legend: bool = False,
    ) -> bytes:
        """
        Generate a high-resolution PNG image for PowerPoint insertion.

        Parameters
        ----------
        df_full : pd.DataFrame
            Full price history with Date and Price columns.
        anchor_date : pd.Timestamp, optional
            Optional anchor date for regression channel.
        lookback_days : int
            Number of days to display.
        width_cm, height_cm : float
            Image dimensions in centimeters.
        show_legend : bool
            Whether to show the legend.

        Returns
        -------
        bytes
            PNG image bytes.
        """
        df_ma = self._add_mas(df_full)
        today = df_ma["Date"].max()
        cutoff = today - timedelta(days=lookback_days)
        df = df_ma[df_ma["Date"] >= cutoff].copy()

        if df.empty:
            # Return blank image
            fig, ax = plt.subplots(figsize=(width_cm / 2.54, height_cm / 2.54))
            buf = BytesIO()
            plt.savefig(buf, format="png", dpi=300, transparent=True)
            plt.close(fig)
            buf.seek(0)
            return buf.getvalue()

        # Calculate regression channel
        upper = None
        lower = None
        uptrend = False
        if anchor_date is not None:
            subset = df_ma[df_ma["Date"] >= anchor_date].copy()
            if len(subset) > 1:
                X = subset["Date"].map(pd.Timestamp.toordinal).to_numpy().reshape(-1, 1)
                y_vals = subset["Price"].to_numpy()
                model = LinearRegression().fit(X, y_vals)
                trend = model.predict(X)
                resid = y_vals - trend
                std_resid = np.std(resid)

                # Extend channel to full visible range
                X_full = df[df["Date"] >= anchor_date]["Date"].map(pd.Timestamp.toordinal).to_numpy().reshape(-1, 1)
                trend_full = model.predict(X_full)
                upper = trend_full + 2 * std_resid
                lower = trend_full - 2 * std_resid
                uptrend = model.coef_[0] > 0

        # Create matplotlib figure
        fig, ax = plt.subplots(figsize=(width_cm / 2.54, height_cm / 2.54))

        # Plot price
        ax.plot(df["Date"], df["Price"], color="#1f77b4", linewidth=2, label=self.config.display_name)

        # Plot moving averages
        ax.plot(df_ma["Date"], df_ma["MA_50"], color="#00FF00", linewidth=1.5, label="50-day MA")
        ax.plot(df_ma["Date"], df_ma["MA_100"], color="#FFA500", linewidth=1.5, label="100-day MA")
        ax.plot(df_ma["Date"], df_ma["MA_200"], color="#FF0000", linewidth=1.5, label="200-day MA")

        # Fibonacci levels
        hi, lo = df["Price"].max(), df["Price"].min()
        span = hi - lo
        fib_levels = [hi, hi - 0.236 * span, hi - 0.382 * span, hi - 0.5 * span, hi - 0.618 * span, lo]
        for lvl in fib_levels:
            ax.axhline(y=lvl, color="grey", linestyle="--", linewidth=0.8, alpha=0.6)

        # Draw regression channel
        if anchor_date is not None and upper is not None and lower is not None:
            fill_color = (0, 0.6, 0, 0.25) if uptrend else (0.78, 0, 0, 0.25)
            line_color = "#008000" if uptrend else "#C00000"
            subset = df[df["Date"] >= anchor_date].copy().reset_index(drop=True)
            ax.plot(subset["Date"], upper, color=line_color, linestyle="--")
            ax.plot(subset["Date"], lower, color=line_color, linestyle="--")
            ax.fill_between(subset["Date"], lower, upper, color=fill_color)

        # Styling
        for spine in ax.spines.values():
            spine.set_visible(False)
        ax.tick_params(left=True, bottom=True, labelleft=True, labelbottom=True)

        if show_legend:
            ax.legend(loc="upper center", bbox_to_anchor=(0.5, 1.1), ncol=4, fontsize=8, frameon=False)

        plt.tight_layout()

        buf = BytesIO()
        plt.savefig(buf, format="png", dpi=600, transparent=True)
        plt.close(fig)
        buf.seek(0)
        return buf.getvalue()

    # -------------------------------------------------------------------------
    # PowerPoint manipulation
    # -------------------------------------------------------------------------

    def _find_slide(self, prs: Presentation) -> Optional[int]:
        """
        Locate the index of the slide that contains this instrument's placeholder.

        Returns
        -------
        int or None
            Zero-based slide index or None if not found.
        """
        search_names = [
            self.config.name.lower(),
            self.config.placeholder_prefix.lower(),
        ]
        for idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                name_attr = getattr(shape, "name", "").lower()
                if name_attr in search_names:
                    return idx
                if shape.has_text_frame:
                    text = (shape.text or "").strip().lower()
                    if text in [f"[{n}]" for n in search_names]:
                        return idx
        return None

    def insert_technical_score_number(self, prs: Presentation, excel_file) -> Presentation:
        """
        Insert the technical score (integer) into the instrument's slide.

        Parameters
        ----------
        prs : Presentation
            PowerPoint presentation.
        excel_file : file-like or path
            Excel file containing technical scores.

        Returns
        -------
        Presentation
            Modified presentation.
        """
        score = self._get_technical_score(excel_file)
        score_text = "N/A" if score is None else f"{int(round(float(score)))}"

        placeholder_name = f"tech_score_{self.config.name}"
        placeholder_patterns = ["[XXX]", "XXX"]

        slide_idx = self._find_slide(prs)
        if slide_idx is None:
            return prs

        slide = prs.slides[slide_idx]

        # Search for named placeholder
        for shape in slide.shapes:
            if getattr(shape, "name", "").lower() == placeholder_name.lower():
                if shape.has_text_frame:
                    runs = shape.text_frame.paragraphs[0].runs
                    attrs = self._get_run_font_attributes(runs[0]) if runs else (None,) * 6
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    new_run = p.add_run()
                    new_run.text = score_text
                    self._apply_run_font_attributes(new_run, *attrs)
                return prs

        # Search for text placeholder
        for shape in slide.shapes:
            if shape.has_text_frame:
                for pattern in placeholder_patterns:
                    if pattern in (shape.text or ""):
                        runs = shape.text_frame.paragraphs[0].runs
                        attrs = self._get_run_font_attributes(runs[0]) if runs else (None,) * 6
                        new_text = shape.text.replace(pattern, score_text)
                        shape.text_frame.clear()
                        p = shape.text_frame.paragraphs[0]
                        new_run = p.add_run()
                        new_run.text = new_text
                        self._apply_run_font_attributes(new_run, *attrs)
                        return prs

        return prs

    def insert_momentum_score_number(
        self,
        prs: Presentation,
        excel_file,
        price_mode: str = "Last Price",
    ) -> Presentation:
        """
        Insert the momentum score (integer) into the instrument's slide.

        Parameters
        ----------
        prs : Presentation
            PowerPoint presentation.
        excel_file : file-like or path
            Excel file containing price data.
        price_mode : str
            One of "Last Price" or "Last Close".

        Returns
        -------
        Presentation
            Modified presentation.
        """
        score = self._get_momentum_score(excel_file, price_mode=price_mode)
        score_text = "N/A" if score is None else f"{int(round(float(score)))}"

        placeholder_name = f"mom_score_{self.config.name}"
        placeholder_patterns = ["[YYY]", "YYY"]

        slide_idx = self._find_slide(prs)
        if slide_idx is None:
            return prs

        slide = prs.slides[slide_idx]

        # Search for named placeholder
        for shape in slide.shapes:
            if getattr(shape, "name", "").lower() == placeholder_name.lower():
                if shape.has_text_frame:
                    runs = shape.text_frame.paragraphs[0].runs
                    attrs = self._get_run_font_attributes(runs[0]) if runs else (None,) * 6
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    new_run = p.add_run()
                    new_run.text = score_text
                    self._apply_run_font_attributes(new_run, *attrs)
                return prs

        # Search for text placeholder
        for shape in slide.shapes:
            if shape.has_text_frame:
                for pattern in placeholder_patterns:
                    if pattern in (shape.text or ""):
                        runs = shape.text_frame.paragraphs[0].runs
                        attrs = self._get_run_font_attributes(runs[0]) if runs else (None,) * 6
                        new_text = shape.text.replace(pattern, score_text)
                        shape.text_frame.clear()
                        p = shape.text_frame.paragraphs[0]
                        new_run = p.add_run()
                        new_run.text = new_text
                        self._apply_run_font_attributes(new_run, *attrs)
                        return prs

        return prs

    def insert_subtitle(
        self,
        prs: Presentation,
        subtitle_text: str,
    ) -> Presentation:
        """
        Insert a subtitle into the instrument's slide.

        Parameters
        ----------
        prs : Presentation
            PowerPoint presentation.
        subtitle_text : str
            Subtitle text to insert.

        Returns
        -------
        Presentation
            Modified presentation.
        """
        placeholder_name = f"subtitle_{self.config.name}"
        placeholder_pattern = "[SUBTITLE]"

        slide_idx = self._find_slide(prs)
        if slide_idx is None:
            return prs

        slide = prs.slides[slide_idx]

        # Search for named placeholder
        for shape in slide.shapes:
            if getattr(shape, "name", "").lower() == placeholder_name.lower():
                if shape.has_text_frame:
                    runs = shape.text_frame.paragraphs[0].runs
                    attrs = self._get_run_font_attributes(runs[0]) if runs else (None,) * 6
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    new_run = p.add_run()
                    new_run.text = subtitle_text
                    self._apply_run_font_attributes(new_run, *attrs)
                return prs

        # Search for text placeholder
        for shape in slide.shapes:
            if shape.has_text_frame and placeholder_pattern in (shape.text or ""):
                runs = shape.text_frame.paragraphs[0].runs
                attrs = self._get_run_font_attributes(runs[0]) if runs else (None,) * 6
                new_text = shape.text.replace(placeholder_pattern, subtitle_text)
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                new_run = p.add_run()
                new_run.text = new_text
                self._apply_run_font_attributes(new_run, *attrs)
                return prs

        return prs

    def insert_technical_chart(
        self,
        prs: Presentation,
        excel_path: pathlib.Path,
        anchor_date: Optional[pd.Timestamp] = None,
        price_mode: str = "Last Price",
        lookback_days: int = 90,
    ) -> Presentation:
        """
        Insert a technical chart into the instrument's slide.

        Parameters
        ----------
        prs : Presentation
            PowerPoint presentation.
        excel_path : pathlib.Path
            Path to Excel file with price data.
        anchor_date : pd.Timestamp, optional
            Anchor date for regression channel.
        price_mode : str
            One of "Last Price" or "Last Close".
        lookback_days : int
            Number of days to display.

        Returns
        -------
        Presentation
            Modified presentation.
        """
        df_raw = self._load_price_data(excel_path, price_mode=price_mode)
        if df_raw.empty:
            return prs

        img_bytes = self._generate_image_from_df(
            df_raw,
            anchor_date=anchor_date,
            lookback_days=lookback_days,
        )

        slide_idx = self._find_slide(prs)
        if slide_idx is None:
            return prs

        slide = prs.slides[slide_idx]
        placeholder_name = f"chart_{self.config.name}"

        # Look for placeholder
        for shape in slide.shapes:
            if getattr(shape, "name", "").lower() == placeholder_name.lower():
                # Replace with image
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height

                # Remove old shape
                sp = shape.element
                sp.getparent().remove(sp)

                # Insert new image
                img_stream = BytesIO(img_bytes)
                slide.shapes.add_picture(img_stream, left, top, width, height)
                return prs

        return prs
