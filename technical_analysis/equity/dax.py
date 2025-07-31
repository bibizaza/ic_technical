"""
Utility functions for DAX technical analysis and high‑resolution export.

This module provides tools to build interactive and static charts for the
DAX index, calculate and insert technical and momentum scores into
PowerPoint presentations, generate horizontal and vertical gauges that
visualise the average of the technical and momentum scores, as well as
contextual trading ranges (higher and lower range bounds).  Functions
fall back to sensible defaults when placeholders are not found.

Key functions include:

* ``make_dax_figure`` – interactive Plotly chart for Streamlit.
* ``insert_dax_technical_chart`` – insert a static DAX chart into a PPTX.
* ``insert_dax_technical_score_number`` – insert the technical score (integer).
* ``insert_dax_momentum_score_number`` – insert the momentum score (integer).
* ``insert_dax_subtitle`` – insert a user‑defined subtitle into the DAX slide.
* ``generate_average_gauge_image`` – create a horizontal gauge image.
* ``insert_dax_average_gauge`` – insert the gauge into a PPT slide.
* ``insert_dax_technical_assessment`` – insert a descriptive “view” text.
* ``generate_range_gauge_chart_image`` – create a combined price chart with
  a vertical range gauge on the right hand side, including a horizontal line
  connecting the last price to the gauge.  This function is used by
  ``insert_dax_technical_chart_with_range``.
* ``insert_dax_technical_chart_with_range`` – insert the DAX technical
  analysis chart with the higher/lower range gauge into the PPT.

The range gauge illustrates the recent trading range for the DAX.
Instead of using the absolute high and low closes of the last 90 days,
the bounds are estimated from recent volatility.  Whenever possible the
code looks up the forward‑looking volatility index (V1X) and computes a
1‑week expected move as ``(current_price × (V1X / 100)) / sqrt(52)``.
The upper and lower bounds are the current price plus and minus that
expected move.  If the volatility index is unavailable, the code falls
back to using realised volatility: it computes the standard deviation of
30‑day daily returns, annualises it and converts it to a 1‑week move
using the same formula.  As a last resort when neither volatility
measure can be computed the bounds default to ±2 % of the current price.
A minimum ±1 % band is enforced to avoid overlapping annotations when
volatility is extremely low.  A horizontal line continues the last
price through to the gauge so that viewers can quickly assess how far
the index lies from its typical volatility band.  This method
provides context on potential upside and downside moves based on
recent price variability rather than simply the most recent highs and
lows.
"""

from __future__ import annotations

from datetime import timedelta
import pathlib
from typing import Optional, Tuple

import pandas as pd
import plotly.graph_objects as go
from sklearn.linear_model import LinearRegression

from pptx import Presentation
from pptx.util import Cm
from pptx.dml.color import RGBColor
from io import BytesIO
import matplotlib.pyplot as plt
import matplotlib.patches as patches
from matplotlib.colors import LinearSegmentedColormap
import numpy as np

# Import helper for adjusting price data according to price mode.  The utils
# module must reside at the project root.  It is deliberately not imported
# from ``technical_analysis.utils`` so that this module can be reused when
# nested within ``ic`` and ``utils.py`` sits at the project root.
try:
    from utils import adjust_prices_for_mode  # type: ignore
except Exception:
    # If the utils module is not available, define a no-op fallback.  This
    # preserves compatibility with environments where price mode is not used.
    adjust_prices_for_mode = None  # type: ignore

###############################################################################
# Internal helpers
###############################################################################

def _get_run_font_attributes(run):
    """Capture font attributes from a run.

    Returns a tuple ``(size, rgb, theme_color, brightness, bold, italic)``.
    The colour information includes either the RGB value if explicitly
    defined, or the theme colour and brightness for a scheme colour.  If
    colour information is not available, ``rgb`` and ``theme_color`` are
    ``None``.  Bold and italic attributes are preserved as provided.
    """
    if run is None:
        return None, None, None, None, None, None
    size = run.font.size
    colour = run.font.color
    rgb = None
    theme_color = None
    brightness = None
    # Try to capture an explicit RGB value
    try:
        rgb = colour.rgb
    except Exception:
        rgb = None
        # If no RGB value, attempt to capture a theme colour
        try:
            theme_color = colour.theme_color
        except Exception:
            theme_color = None
    # Capture brightness adjustment if available
    try:
        brightness = colour.brightness
    except Exception:
        brightness = None
    bold = run.font.bold
    italic = run.font.italic
    return size, rgb, theme_color, brightness, bold, italic


def _apply_run_font_attributes(new_run, size, rgb, theme_color, brightness, bold, italic):
    """Apply captured font attributes to a new run.

    Parameters
    ----------
    new_run : pptx.text.run.Run
        The run to which attributes should be applied.
    size : pptx.util.Length or None
        The font size to apply.
    rgb : pptx.dml.color.RGBColor or None
        The explicit RGB colour value to apply.
    theme_color : MSO_THEME_COLOR or None
        The theme colour value to apply if no RGB colour is defined.
    brightness : float or None
        Brightness adjustment for the colour, if any.
    bold : bool or None
        Whether the font should be bold.
    italic : bool or None
        Whether the font should be italic.
    """
    if size is not None:
        new_run.font.size = size
    # Apply colour: prefer explicit RGB, otherwise theme colour
    if rgb is not None:
        try:
            new_run.font.color.rgb = rgb
        except Exception:
            pass
    elif theme_color is not None:
        try:
            new_run.font.color.theme_color = theme_color
            if brightness is not None:
                new_run.font.color.brightness = brightness
        except Exception:
            pass
    # Apply bold and italic
    if bold is not None:
        new_run.font.bold = bold
    if italic is not None:
        new_run.font.italic = italic


def _load_price_data(
    excel_path: pathlib.Path,
    ticker: str = "DAX Index",
    price_mode: str = "Last Price",
) -> pd.DataFrame:
    """
    Read the raw price sheet and return a tidy Date‑Price DataFrame.

    Parameters
    ----------
    excel_path : pathlib.Path
        Path to the Excel workbook containing price data.
    ticker : str, default "DAX Index"
        Column name corresponding to the desired ticker in the Excel sheet.
    price_mode : str, default "Last Price"
        One of "Last Price" or "Last Close".  If ``adjust_prices_for_mode``
        is available and the mode is "Last Close", rows with the last
        recorded date (if equal to today's date) will be dropped.

    Returns
    -------
    pandas.DataFrame
        DataFrame with columns ``Date`` and ``Price``.  The data are
        sorted by date and any rows with missing values are removed.
    """
    df = pd.read_excel(excel_path, sheet_name="data_prices")
    df = df.drop(index=0)
    df = df[df[df.columns[0]] != "DATES"]
    df["Date"] = pd.to_datetime(df[df.columns[0]], errors="coerce")
    df["Price"] = pd.to_numeric(df[ticker], errors="coerce")
    df_clean = (
        df.dropna(subset=["Date", "Price"]).sort_values("Date").reset_index(drop=True)[
            ["Date", "Price"]
        ]
    )
    # Adjust for price mode if helper is available
    if adjust_prices_for_mode is not None and price_mode:
        try:
            df_clean, _ = adjust_prices_for_mode(df_clean, price_mode)
        except Exception:
            # If adjustment fails, silently fall back to unadjusted data
            pass
    return df_clean


def _add_mas(df: pd.DataFrame) -> pd.DataFrame:
    """Add 50/100/200‑day moving‑average columns to a DataFrame."""
    out = df.copy()
    for w in (50, 100, 200):
        out[f"MA_{w}"] = out["Price"].rolling(w, min_periods=1).mean()
    return out

def _get_vol_index_value(
    excel_obj_or_path,
    price_mode: str = "Last Price",
    vol_ticker: str = "V1X Index",
) -> Optional[float]:
    """
    Retrieve the most recent value of a volatility index (e.g. V1X) from
    the ``data_prices`` sheet.  If ``price_mode`` is ``"Last Close"``,
    the most recent date is dropped if it matches today's date.  The
    returned value is the last available entry after price‑mode adjustment.

    Parameters
    ----------
    excel_obj_or_path : file‑like or path
        The Excel workbook containing price data.
    price_mode : str, default "Last Price"
        One of "Last Price" or "Last Close".  When set to "Last Close"
        rows corresponding to the most recent date (if equal to today's
        date) will be excluded before taking the last value.
    vol_ticker : str, default "V1X Index"
        Column name in the ``data_prices`` sheet corresponding to the
        volatility index whose level should be used.

    Returns
    -------
    float or None
        The most recent volatility index value, or ``None`` if it cannot
        be read or converted to a float.
    """
    try:
        df = pd.read_excel(excel_obj_or_path, sheet_name="data_prices")
    except Exception:
        return None
    # Drop first row and metadata rows
    df = df.drop(index=0)
    df = df[df[df.columns[0]] != "DATES"]
    # Parse dates and the volatility index column
    df["Date"] = pd.to_datetime(df[df.columns[0]], errors="coerce")
    # Ensure the volatility index column exists
    if vol_ticker not in df.columns:
        return None
    df["Price"] = pd.to_numeric(df[vol_ticker], errors="coerce")
    df_clean = df.dropna(subset=["Date", "Price"]).sort_values("Date").reset_index(drop=True)[
        ["Date", "Price"]
    ]
    # Apply price mode adjustment if possible
    if adjust_prices_for_mode is not None and price_mode:
        try:
            df_clean, _ = adjust_prices_for_mode(df_clean, price_mode)
        except Exception:
            pass
    if df_clean.empty:
        return None
    try:
        return float(df_clean["Price"].iloc[-1])
    except Exception:
        return None


###############################################################################
# Plotly interactive chart for Streamlit
###############################################################################

def make_dax_figure(
    excel_path: str | pathlib.Path,
    anchor_date: Optional[pd.Timestamp] = None,
    price_mode: str = "Last Price",
) -> go.Figure:
    """
    Build an interactive DAX chart for Streamlit.

    Parameters
    ----------
    excel_path : str or pathlib.Path
        Path to the Excel file containing DAX price data.
    anchor_date : pandas.Timestamp or None, optional
        If provided, a regression channel is drawn from ``anchor_date`` to the
        latest date.
    price_mode : str, default "Last Price"
        One of "Last Price" or "Last Close".  When set to "Last Close"
        and if the ``utils.adjust_prices_for_mode`` helper is available,
        rows corresponding to the most recent date (if equal to today's
        date) are excluded from the chart.  This allows the chart to
        represent the prior day's closing prices rather than an
        intraday or current price.

    Returns
    -------
    go.Figure
        A Plotly figure with price, moving averages, Fibonacci lines and
        an optional regression channel.
    """
    excel_path = pathlib.Path(excel_path)
    # Load data and adjust according to the price mode
    df_raw = _load_price_data(excel_path, "DAX Index", price_mode=price_mode)
    df_full = _add_mas(df_raw)

    if df_full.empty:
        return go.Figure()

    today = df_full["Date"].max().normalize()
    start = today - timedelta(days=365)
    df = df_full[df_full["Date"].between(start, today)].reset_index(drop=True)

    if df.empty:
        return go.Figure()

    last_price = df["Price"].iloc[-1]
    last_price_str = f"{last_price:,.2f}"

    fig = go.Figure()
    fig.add_trace(
        go.Scatter(
            x=df["Date"],
            y=df["Price"],
            mode="lines",
            name=f"DAX Price (last: {last_price_str})",
            line=dict(color="#153D64", width=2.5),
        )
    )
    fig.add_trace(
        go.Scatter(
            x=df["Date"],
            y=df["MA_50"],
            mode="lines",
            name="50‑day MA",
            line=dict(color="#008000", width=1.5),
        )
    )
    fig.add_trace(
        go.Scatter(
            x=df["Date"],
            y=df["MA_100"],
            mode="lines",
            name="100‑day MA",
            line=dict(color="#FFA500", width=1.5),
        )
    )
    fig.add_trace(
        go.Scatter(
            x=df["Date"],
            y=df["MA_200"],
            mode="lines",
            name="200‑day MA",
            line=dict(color="#FF0000", width=1.5),
        )
    )

    hi, lo = df["Price"].max(), df["Price"].min()
    span = hi - lo
    for lvl in [
        hi,
        hi - 0.236 * span,
        hi - 0.382 * span,
        hi - 0.5 * span,
        hi - 0.618 * span,
        lo,
    ]:
        fig.add_hline(
            y=lvl, line=dict(color="grey", dash="dash", width=1), opacity=0.6
        )

    if anchor_date is not None:
        per = df_full[df_full["Date"].between(anchor_date, today)].copy()
        if not per.empty:
            X = per["Date"].map(pd.Timestamp.toordinal).to_numpy().reshape(-1, 1)
            y_vals = per["Price"].to_numpy()
            model = LinearRegression().fit(X, y_vals)
            trend = model.predict(X)
            resid = y_vals - trend
            upper = trend + resid.max()
            lower = trend + resid.min()

            uptrend = model.coef_[0] > 0
            lineclr = "green" if uptrend else "red"
            fillclr = "rgba(0,150,0,0.25)" if uptrend else "rgba(200,0,0,0.25)"

            fig.add_trace(
                go.Scatter(
                    x=per["Date"],
                    y=upper,
                    mode="lines",
                    line=dict(color=lineclr, dash="dash"),
                    showlegend=False,
                )
            )
            fig.add_trace(
                go.Scatter(
                    x=per["Date"],
                    y=lower,
                    mode="lines",
                    line=dict(color=lineclr, dash="dash"),
                    fill="tonexty",
                    fillcolor=fillclr,
                    showlegend=False,
                )
            )

    fig.update_layout(
        margin=dict(l=30, r=30, t=60, b=40),
        showlegend=True,
        legend=dict(
            orientation="h",
            yanchor="bottom",
            y=1.12,
            xanchor="center",
            x=0.5,
            font=dict(size=12),
        ),
        xaxis_title=None,
        yaxis_title=None,
        xaxis=dict(showgrid=False),
        yaxis=dict(showgrid=False, zeroline=False),
    )
    return fig


###############################################################################
# High‑resolution chart export (PNG)
###############################################################################

def _generate_dax_image_from_df(
    df_full: pd.DataFrame,
    anchor_date: Optional[pd.Timestamp],
    width_cm: float = 21.41,
    height_cm: float = 7.53,
) -> bytes:
    """
    Create a high‑resolution (dpi=300) transparent PNG chart from the DataFrame.
    Includes price, moving averages, Fibonacci lines and optional regression channel.
    """
    today = df_full["Date"].max().normalize()
    start = today - timedelta(days=365)
    df = df_full[df_full["Date"].between(start, today)].reset_index(drop=True)

    df_ma = df.copy()
    for w in (50, 100, 200):
        df_ma[f"MA_{w}"] = df_ma["Price"].rolling(w, min_periods=1).mean()

    uptrend = False
    upper = lower = None
    if anchor_date is not None:
        subset = df_full[df_full["Date"].between(anchor_date, today)].copy()
        if not subset.empty:
            X = subset["Date"].map(pd.Timestamp.toordinal).to_numpy().reshape(-1, 1)
            y_vals = subset["Price"].to_numpy()
            model = LinearRegression().fit(X, y_vals)
            trend = model.predict(X)
            resid = y_vals - trend
            uptrend = model.coef_[0] > 0
            upper = trend + resid.max()
            lower = trend + resid.min()

    last_price = df["Price"].iloc[-1]
    last_price_str = f"{last_price:,.2f}"

    fig_width_in, fig_height_in = width_cm / 2.54, height_cm / 2.54
    plt.style.use("default")
    fig, ax = plt.subplots(figsize=(fig_width_in, fig_height_in))

    ax.plot(
        df["Date"],
        df["Price"],
        color="#153D64",
        linewidth=2.5,
        label=f"DAX Price (last: {last_price_str})",
    )
    ax.plot(
        df_ma["Date"],
        df_ma["MA_50"],
        color="#008000",
        linewidth=1.5,
        label="50‑day MA",
    )
    ax.plot(
        df_ma["Date"],
        df_ma["MA_100"],
        color="#FFA500",
        linewidth=1.5,
        label="100‑day MA",
    )
    ax.plot(
        df_ma["Date"],
        df_ma["MA_200"],
        color="#FF0000",
        linewidth=1.5,
        label="200‑day MA",
    )

    hi, lo = df["Price"].max(), df["Price"].min()
    span = hi - lo
    fib_levels = [
        hi,
        hi - 0.236 * span,
        hi - 0.382 * span,
        hi - 0.5 * span,
        hi - 0.618 * span,
        lo,
    ]
    for lvl in fib_levels:
        ax.axhline(y=lvl, color="grey", linestyle="--", linewidth=0.8, alpha=0.6)

    if anchor_date is not None and upper is not None and lower is not None:
        fill_color = (0, 0.6, 0, 0.25) if uptrend else (0.78, 0, 0, 0.25)
        line_color = "#008000" if uptrend else "#C00000"
        subset = (
            df_full[df_full["Date"].between(anchor_date, today)].copy().reset_index(drop=True)
        )
        ax.plot(subset["Date"], upper, color=line_color, linestyle="--")
        ax.plot(subset["Date"], lower, color=line_color, linestyle="--")
        ax.fill_between(subset["Date"], lower, upper, color=fill_color)

    for spine in ax.spines.values():
        spine.set_visible(False)
    ax.tick_params(left=True, bottom=True, labelleft=True, labelbottom=True)

    ax.legend(
        loc="upper center", bbox_to_anchor=(0.5, 1.1), ncol=4, fontsize=8, frameon=False
    )
    plt.tight_layout()

    buf = BytesIO()
    plt.savefig(buf, format="png", dpi=600, transparent=True)
    plt.close(fig)
    buf.seek(0)
    return buf.getvalue()


###############################################################################
# Score helpers
###############################################################################

def _get_dax_technical_score(excel_obj_or_path) -> Optional[float]:
    """
    Retrieve the technical score for DAX from 'data_technical_score' (col A, B).
    Returns None if the sheet or score is unavailable.
    """
    try:
        df = pd.read_excel(excel_obj_or_path, sheet_name="data_technical_score")
    except Exception:
        return None
    df = df.dropna(subset=[df.columns[0], df.columns[1]])
    for _, row in df.iterrows():
        if str(row[df.columns[0]]).strip().upper() == "DAX INDEX":
            try:
                return float(row[df.columns[1]])
            except Exception:
                return None
    return None


def _find_dax_slide(prs: Presentation) -> Optional[int]:
    """Locate the index of the slide that contains the DAX placeholder.

    This helper searches for a slide containing a shape named ``dax`` or
    whose text is exactly ``[dax]`` (case‑insensitive).  It returns the
    zero‑based slide index or ``None`` if no such slide exists.
    """
    for idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            name_attr = getattr(shape, "name", "").lower()
            if name_attr == "dax":
                return idx
            if shape.has_text_frame:
                if (shape.text or "").strip().lower() == "[dax]":
                    return idx
    return None


def insert_dax_technical_score_number(prs: Presentation, excel_file) -> Presentation:
    """
    Insert the DAX technical score (integer) into the DAX slide.

    This function looks for a shape named ``tech_score_dax`` on the slide
    identified by the ``dax`` placeholder.  If not found, it searches for
    placeholders ``[XXX]`` or ``XXX`` within that slide.  Formatting from
    the original placeholder run is preserved.  Other slides are not
    modified, avoiding accidental replacement of CSI placeholders.
    """
    score = _get_dax_technical_score(excel_file)
    score_text = "N/A" if score is None else f"{int(round(float(score)))}"

    placeholder_name = "tech_score_dax"
    placeholder_patterns = ["[XXX]", "XXX"]

    dax_idx = _find_dax_slide(prs)
    if dax_idx is None:
        # No DAX slide found; return unmodified
        return prs
    slide = prs.slides[dax_idx]
    # First search for a shape named exactly as the placeholder
    for shape in slide.shapes:
        if getattr(shape, "name", "").lower() == placeholder_name:
            if shape.has_text_frame:
                runs = shape.text_frame.paragraphs[0].runs
                attrs = _get_run_font_attributes(runs[0]) if runs else (None, None, None, None, None, None)
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                new_run = p.add_run()
                new_run.text = score_text
                _apply_run_font_attributes(new_run, *attrs)
            return prs
    # Otherwise, search for textual placeholders within shapes on the DAX slide
    for shape in slide.shapes:
        if shape.has_text_frame:
            for pattern in placeholder_patterns:
                if pattern in (shape.text or ""):
                    runs = shape.text_frame.paragraphs[0].runs
                    attrs = _get_run_font_attributes(runs[0]) if runs else (None, None, None, None, None, None)
                    new_text = shape.text.replace(pattern, score_text)
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    new_run = p.add_run()
                    new_run.text = new_text
                    _apply_run_font_attributes(new_run, *attrs)
                    return prs
    return prs


###############################################################################
# Call‑out range helpers and insertion
###############################################################################

def generate_range_callout_chart_image(
    df_full: pd.DataFrame,
    anchor_date: Optional[pd.Timestamp] = None,
    lookback_days: int = 90,
    width_cm: float = 21.41,
    height_cm: float = 7.53,
    callout_width_cm: float = 3.5,
    *,
    vol_index_value: Optional[float] = None,
) -> bytes:
    """
    Create a PNG image of the DAX price chart with a textual call‑out on the
    right summarising the recent trading range.  The call‑out lists the
    higher and lower range values (with ±% changes relative to the last
    price) and draws small coloured markers aligned with those levels on
    the y‑axis.  This design preserves the full chart width and avoids
    overlapping the price plot with additional graphics.

    Parameters
    ----------
    df_full : pandas.DataFrame
        Full DAX price history with 'Date' and 'Price' columns.
    anchor_date : pandas.Timestamp or None, optional
        Optional anchor date for a regression channel; if provided, the
        channel is drawn on the price chart.
    lookback_days : int, default 90
        The lookback window for computing the high and low bounds.
    width_cm : float, default 21.41
        Overall width of the output image in centimetres.  This should
        correspond to the template placeholder width.
    height_cm : float, default 7.53
        Height of the output image in centimetres.
    callout_width_cm : float, default 3.5
        Width of the call‑out area on the right where the range summary
        appears.  The remaining width is used for the chart.

    Returns
    -------
    bytes
        PNG image bytes with transparency.
    """
    import matplotlib.pyplot as plt
    import matplotlib.patches as mpatches

    if df_full.empty:
        return b""

    # Restrict to the last year of data for plotting
    today = df_full["Date"].max().normalize()
    start = today - timedelta(days=365)
    df = df_full[df_full["Date"].between(start, today)].reset_index(drop=True)

    # Calculate moving averages on the 1‑year subset
    df_ma = _add_mas(df)

    # Optional regression channel
    uptrend = False
    upper_channel = lower_channel = None
    if anchor_date is not None:
        subset_full = df_full[df_full["Date"].between(anchor_date, today)].copy()
        if not subset_full.empty:
            X = subset_full["Date"].map(pd.Timestamp.toordinal).to_numpy().reshape(-1, 1)
            y_vals = subset_full["Price"].to_numpy()
            model = LinearRegression().fit(X, y_vals)
            trend = model.predict(X)
            resid = y_vals - trend
            uptrend = model.coef_[0] > 0
            upper_channel = trend + resid.max()
            lower_channel = trend + resid.min()

    # Compute high/low bounds and current price.  If an implied volatility
    # value is provided (e.g. the V1X level), use it to estimate the
    # expected one‑week move.  The expected move is computed as
    # ``last_price × (vol_index_value/100) / sqrt(52)``.  Otherwise
    # fall back to the realised‑volatility‑based bounds returned by
    # ``_compute_range_bounds``.
    last_price = df["Price"].iloc[-1]
    if vol_index_value is not None and last_price and not np.isnan(last_price):
        expected_move = (last_price * (vol_index_value / 100.0)) / np.sqrt(52.0)
        upper_bound = last_price + expected_move
        lower_bound = last_price - expected_move
    else:
        upper_bound, lower_bound = _compute_range_bounds(df_full, lookback_days=lookback_days)
    # Enforce a minimum total range (e.g. ±1 % of the current price) to avoid overlapping text.
    min_range_pct = 0.02  # 2% total band → ±1% around the current price
    if last_price and not np.isnan(last_price):
        range_span_pct = (upper_bound - lower_bound) / last_price if last_price else 0.0
        if range_span_pct < min_range_pct:
            half_span = (min_range_pct * last_price) / 2.0
            upper_bound = last_price + half_span
            lower_bound = last_price - half_span
        # Recompute percentage differences after adjusting range
        up_pct = (upper_bound - last_price) / last_price * 100.0
        down_pct = (last_price - lower_bound) / last_price * 100.0
    else:
        # Handle missing last_price gracefully
        up_pct = 0.0
        down_pct = 0.0

    # Determine y‑axis limits: ensure the axis includes the entire trading
    # range and the observed price range.  We add a small margin so the
    # labels and markers do not overlap the top or bottom edges.
    hi = df["Price"].max()
    lo = df["Price"].min()
    y_max = max(hi, upper_bound) * 1.02
    y_min = min(lo, lower_bound) * 0.98

    # Compute widths for chart and call‑out.  Reserve a small margin on the
    # left of the chart to ensure that y‑axis tick labels and the legend
    # remain visible when the image is inserted into a PowerPoint slide.
    callout_width_cm = min(callout_width_cm, width_cm)
    chart_width_cm = max(width_cm - callout_width_cm, 0.0)
    fig_w_in, fig_h_in = width_cm / 2.54, height_cm / 2.54
    fig = plt.figure(figsize=(fig_w_in, fig_h_in))

    # Relative widths as fractions of the full figure width
    chart_rel_width = chart_width_cm / width_cm if width_cm > 0 else 0.0
    callout_rel_width = callout_width_cm / width_cm if width_cm > 0 else 0.0

    # Define a margin fraction for the left side of the chart.  Without
    # this margin, tick labels and the legend can be clipped when the
    # combined image is saved at high DPI.  Use up to 4% of the figure
    # width or 10% of the chart portion, whichever is smaller.
    margin_rel = min(0.04, 0.10 * chart_rel_width)

    # Axes for chart and call‑out; share the y‑axis so that the call‑out
    # markers align with the same price levels as the chart.  The chart
    # occupies the left portion of the figure starting at margin_rel; the
    # call‑out uses the remaining width starting at chart_rel_width.
    ax_chart = fig.add_axes([margin_rel, 0.0, chart_rel_width - margin_rel, 1.0])
    # Create a separate y‑axis for the call‑out so that hiding its ticks
    # does not remove the ticks from the main chart.  We will manually
    # synchronise the y‑limits below.
    ax_callout = fig.add_axes([chart_rel_width, 0.0, callout_rel_width, 1.0])

    # Set y‑limits before plotting so that shared axes align properly
    ax_chart.set_ylim(y_min, y_max)
    ax_callout.set_ylim(y_min, y_max)

    # Plot price and moving averages on the main chart
    ax_chart.plot(df["Date"], df["Price"], color="#153D64", linewidth=2.5,
                  label=f"DAX Price (last: {last_price:,.2f})")
    ax_chart.plot(df_ma["Date"], df_ma["MA_50"], color="#008000", linewidth=1.5, label="50‑day MA")
    ax_chart.plot(df_ma["Date"], df_ma["MA_100"], color="#FFA500", linewidth=1.5, label="100‑day MA")
    ax_chart.plot(df_ma["Date"], df_ma["MA_200"], color="#FF0000", linewidth=1.5, label="200‑day MA")
    # Fibonacci levels on the subset
    sub_hi, sub_lo = df["Price"].max(), df["Price"].min()
    sub_span = sub_hi - sub_lo
    for lvl in [sub_hi, sub_hi - 0.236 * sub_span, sub_hi - 0.382 * sub_span,
                sub_hi - 0.5 * sub_span, sub_hi - 0.618 * sub_span, sub_lo]:
        ax_chart.axhline(lvl, color="grey", linestyle="--", linewidth=0.8, alpha=0.6)

    # Regression channel shading
    if anchor_date is not None and upper_channel is not None and lower_channel is not None:
        subset = df_full[df_full["Date"].between(anchor_date, today)].copy().reset_index(drop=True)
        fill_color = (0, 0.6, 0, 0.25) if uptrend else (0.78, 0, 0, 0.25)
        line_color = "#008000" if uptrend else "#C00000"
        ax_chart.plot(subset["Date"], upper_channel, color=line_color, linestyle="--")
        ax_chart.plot(subset["Date"], lower_channel, color=line_color, linestyle="--")
        ax_chart.fill_between(subset["Date"], lower_channel, upper_channel, color=fill_color)

    # Style the main chart: remove spines and configure ticks.  We set the
    # y‑axis tick length to zero so that the small horizontal tick marks
    # next to the axis labels are not visible, while keeping the labels
    # themselves.  The x‑axis ticks retain their default length.
    for spine in ax_chart.spines.values():
        spine.set_visible(False)
    ax_chart.tick_params(axis="y", which="both", length=0)
    ax_chart.tick_params(axis="x", which="both", length=2)
    ax_chart.tick_params(left=True, bottom=True, labelleft=True, labelbottom=True)
    # Legend: place the legend just above the main chart, aligned to the
    # left so that it does not overlap the call‑out panel.  Use a
    # multi‑column layout to fit all entries on a single line.  The
    # bounding box is anchored slightly above the axes (y=1.05).
    ax_chart.legend(loc="upper left", bbox_to_anchor=(0.0, 1.05), ncol=4,
                    fontsize=8, frameon=False)

    # Configure call‑out axis: remove ticks and spines; set background white
    ax_callout.set_xlim(0, 1)
    ax_callout.set_xticks([])
    ax_callout.set_yticks([])
    for spine in ax_callout.spines.values():
        spine.set_visible(False)
    ax_callout.set_facecolor("white")

    # Determine x positions for markers and text in relative coordinates.
    # Place the markers near the left of the call‑out area and the text
    # closer to the left so that the numeric portions align on their
    # left edges.  Using ``ha='left'`` keeps all values aligned to the
    # same left margin while the middle line still aligns with the price
    # axis via ``va='center'`` and symmetrical blank lines.
    marker_start_x = 0.02
    marker_end_x = 0.08
    text_x = 0.15

    # Draw small horizontal bars as markers aligned with the high/low bounds
    ax_callout.hlines(upper_bound, xmin=marker_start_x, xmax=marker_end_x,
                      colors="#009951", linewidth=2, transform=ax_callout.transData)
    ax_callout.hlines(lower_bound, xmin=marker_start_x, xmax=marker_end_x,
                      colors="#C00000", linewidth=2, transform=ax_callout.transData)

    # Helper to format values with apostrophes for thousands separators
    def _fmt(val: float) -> str:
        try:
            return f"{val:,.0f}".replace(",", "'")
        except Exception:
            return f"{val:.0f}"

    # Compose label strings with percentage differences.  The index level and
    # percentage are shown together on one line to minimise overlap.  The
    # "Higher Range" label appears above its number, while the "Lower Range"
    # label appears below its number.
    # Compose label strings with percentage differences.  We construct
    # multi‑line strings with symmetrical blank lines so that when
    # ``va='center'`` is used, the index/percentage line (the middle line)
    # aligns exactly with the price level on the y‑axis.  For the upper
    # bound, place "Higher Range" above the value and a blank line below.
    # For the lower bound, place a blank line above the value and
    # "Lower Range" below.  This results in three lines for each label
    # block, ensuring the middle line (index and percent) sits on the
    # specified y‑coordinate.
    upper_text = (
        f"Higher Range\n"
        f"{_fmt(upper_bound)} (+{up_pct:.1f}%)\n"
        f""
    )
    lower_text = (
        f"\n"
        f"{_fmt(lower_bound)} (-{down_pct:.1f}%)\n"
        f"Lower Range"
    )

    # Add the text labels at the appropriate y positions.  Using
    # ``va='center'`` ensures that the middle line (index and percent)
    # aligns with the price level, because there is one line above and
    # one line below.  We align the text to the right so that the plus
    # and minus signs line up neatly.
    ax_callout.text(text_x, upper_bound, upper_text, color="#009951",
                    ha="left", va="center", fontsize=8, fontweight='bold',
                    transform=ax_callout.transData)
    ax_callout.text(text_x, lower_bound, lower_text, color="#C00000",
                    ha="left", va="center", fontsize=8, fontweight='bold',
                    transform=ax_callout.transData)

    # Export to transparent PNG.  Use bbox_inches='tight' so that the
    # entire figure (including legends and tick labels) is saved without
    # cropping.  A small padding is added to provide breathing room.
    buf = BytesIO()
    fig.savefig(buf, format="png", dpi=600, transparent=True,
                bbox_inches="tight", pad_inches=0.05)
    plt.close(fig)
    buf.seek(0)
    return buf.getvalue()


def insert_dax_technical_chart_with_callout(
    prs: Presentation,
    excel_file,
    anchor_date: Optional[pd.Timestamp] = None,
    lookback_days: int = 90,
    price_mode: str = "Last Price",
) -> Presentation:
    """
    Insert the DAX technical analysis chart with the trading range call‑out
    into the PowerPoint.  This function mirrors the behaviour of
    ``insert_dax_technical_chart_with_range`` but uses the call‑out style to
    display the high and low bounds instead of a vertical gauge.

    The image is placed at the fixed coordinates (0.93 cm left, 4.40 cm top)
    with dimensions 21.41 cm wide by 7.53 cm high, matching the template.

    Parameters
    ----------
    prs : Presentation
        The PowerPoint presentation to modify.
    excel_file : file‑like object or path
        Excel workbook containing DAX price data.
    anchor_date : pandas.Timestamp or None, optional
        Optional anchor date for a regression channel.
    lookback_days : int, default 90
        Lookback window for computing the high/low range.

    Returns
    -------
    Presentation
        The presentation with the updated slide.
    """
    # Load the price data from the Excel file
    try:
        df_full = _load_price_data_from_obj(excel_file, "DAX Index", price_mode=price_mode)
    except Exception:
        df_full = _load_price_data(pathlib.Path(excel_file), "DAX Index", price_mode=price_mode)

    # Determine the implied volatility index value (V1X) from the Excel file
    # so that the expected one‑week trading range can be estimated.  If the
    # volatility index cannot be read, ``None`` is returned and the range
    # will fall back to an ATR‑based estimate.
    vol_val = _get_vol_index_value(excel_file, price_mode=price_mode, vol_ticker="V1X Index")
    # Generate the image with the call‑out.  Use an extended width of
    # 25.0 cm while keeping the height at 7.3 cm.  Pass the volatility index
    # value to ``generate_range_callout_chart_image`` so that the range
    # calculation can use the implied volatility if available.
    img_bytes = generate_range_callout_chart_image(
        df_full,
        anchor_date=anchor_date,
        lookback_days=lookback_days,
        width_cm=25.0,
        height_cm=7.3,
        vol_index_value=vol_val,
    )

    # Locate the slide containing the 'dax' placeholder or text
    target_slide = None
    for slide in prs.slides:
        for shape in slide.shapes:
            name_attr = getattr(shape, "name", "").lower()
            if name_attr == "dax":
                target_slide = slide
                break
            if shape.has_text_frame:
                if (shape.text or "").strip().lower() == "[dax]":
                    target_slide = slide
                    break
        if target_slide:
            break
    if target_slide is None:
        target_slide = prs.slides[min(11, len(prs.slides) - 1)]

    # Insert the image at the requested coordinates.  The dimensions 25 cm
    # wide and 7.3 cm high and position (0.93 cm, 4.80 cm) come from the
    # template.
    left = Cm(0.93)
    top = Cm(4.80)
    width = Cm(25.0)
    height = Cm(7.3)
    stream = BytesIO(img_bytes)
    # Add the picture and bring it to the front.  In some templates,
    # additional shapes (e.g. a placeholder gauge) may overlap the chart.
    # Removing and reinserting the picture element near the start of the
    # shape tree ensures the chart remains visible above other content.
    picture = target_slide.shapes.add_picture(stream, left, top, width=width, height=height)
    try:
        sp_tree = target_slide.shapes._spTree
        # Remove the element and reinsert at position 1 (after background)
        sp_tree.remove(picture._element)
        sp_tree.insert(1, picture._element)
    except Exception:
        # Fallback: leave the picture at the end of the shape list
        pass
    return prs


def _get_dax_momentum_score(excel_obj_or_path) -> Optional[float]:
    """Return DAX momentum score, mapping letter grades to numeric if needed."""
    try:
        df = pd.read_excel(excel_obj_or_path, sheet_name="data_trend_rating")
    except Exception:
        return None
    # find DAX row
    mask = df.iloc[:, 0].astype(str).str.strip().str.upper() == "DAX INDEX"
    if not mask.any():
        return None
    row = df.loc[mask].iloc[0]
    # try to convert the existing value to float
    try:
        return float(row.iloc[3])
    except Exception:
        pass
    # fall back to mapping letter rating to numeric using parameters sheet
    rating = str(row.iloc[2]).strip().upper()  # 'Current' column
    mapping = {"A": 100.0, "B": 70.0, "C": 40.0, "D": 0.0}
    # optionally lookup in 'parameters' sheet for customised mapping
    try:
        params = pd.read_excel(excel_obj_or_path, sheet_name="parameters")
        dax_param = params[params["Tickers"].astype(str).str.upper() == "DAX INDEX"]
        if not dax_param.empty and "Unnamed: 8" in dax_param:
            return float(dax_param["Unnamed: 8"].dropna().iloc[0])
    except Exception:
        pass
    return mapping.get(rating)


def insert_dax_momentum_score_number(prs: Presentation, excel_file) -> Presentation:
    """
    Insert the DAX momentum score (integer) into the DAX slide.

    The momentum score is inserted into a shape named ``mom_score_dax`` on
    the DAX slide.  If that shape is not found, any ``XXX`` or ``[XXX]``
    placeholder within the DAX slide is replaced instead.  This avoids
    inadvertently replacing placeholders on other slides.
    """
    score = _get_dax_momentum_score(excel_file)
    score_text = "N/A" if score is None else f"{int(round(float(score)))}"

    placeholder_name = "mom_score_dax"
    placeholder_patterns = ["[XXX]", "XXX"]

    dax_idx = _find_dax_slide(prs)
    if dax_idx is None:
        return prs
    slide = prs.slides[dax_idx]
    # Attempt to replace the named placeholder first
    for shape in slide.shapes:
        if getattr(shape, "name", "").lower() == placeholder_name:
            if shape.has_text_frame:
                runs = shape.text_frame.paragraphs[0].runs
                attrs = _get_run_font_attributes(runs[0]) if runs else (None, None, None, None, None, None)
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                new_run = p.add_run()
                new_run.text = score_text
                _apply_run_font_attributes(new_run, *attrs)
            return prs
    # Otherwise, replace placeholder patterns on the DAX slide only
    for shape in slide.shapes:
        if shape.has_text_frame:
            for pattern in placeholder_patterns:
                if pattern in (shape.text or ""):
                    runs = shape.text_frame.paragraphs[0].runs
                    attrs = _get_run_font_attributes(runs[0]) if runs else (None, None, None, None, None, None)
                    new_text = shape.text.replace(pattern, score_text)
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    new_run = p.add_run()
                    new_run.text = new_text
                    _apply_run_font_attributes(new_run, *attrs)
                    return prs
    return prs


###############################################################################
# Chart insertion
###############################################################################

def insert_dax_technical_chart(
    prs: Presentation,
    excel_file,
    anchor_date: Optional[pd.Timestamp] = None,
    price_mode: str = "Last Price",
) -> Presentation:
    """
    Insert the DAX technical‑analysis chart into the PPT.

    We only use the textbox named ``dax`` (or containing “[dax]”) to locate
    the correct slide; the chart itself is always pasted at the fixed
    coordinates (0.93 cm left, 4.39 cm top, 21.41 cm wide, 7.53 cm high).
    """
    # Load data and generate image
    try:
        df_full = _load_price_data_from_obj(excel_file, "DAX Index", price_mode=price_mode)
    except Exception:
        df_full = _load_price_data(pathlib.Path(excel_file), "DAX Index", price_mode=price_mode)
    img_bytes = _generate_dax_image_from_df(df_full, anchor_date)

    # Find the slide containing the 'dax' placeholder
    target_slide = None
    for slide in prs.slides:
        for shape in slide.shapes:
            name_attr = getattr(shape, "name", "").lower()
            if name_attr == "dax":
                target_slide = slide
                break
            if shape.has_text_frame:
                if (shape.text or "").strip().lower() == "[dax]":
                    target_slide = slide
                    break
        if target_slide:
            break

    if target_slide is None:
        target_slide = prs.slides[min(11, len(prs.slides) - 1)]

    # Always paste the chart at fixed coordinates
    left = Cm(0.93)
    top = Cm(4.39)
    width = Cm(21.41)
    height = Cm(7.53)
    stream = BytesIO(img_bytes)
    target_slide.shapes.add_picture(stream, left, top, width=width, height=height)
    return prs


###############################################################################
# Subtitle insertion
###############################################################################

def insert_dax_subtitle(prs: Presentation, subtitle: str) -> Presentation:
    """
    Replace the DAX subtitle placeholder with the provided text.

    Only the slide identified by the ``dax`` placeholder is modified.  A
    shape named ``dax_text`` takes precedence; if it does not exist
    within the DAX slide, any occurrences of ``XXX`` or ``[XXX]`` on
    that slide are replaced instead.  Formatting of the original run is
    preserved.
    """
    placeholder_name = "dax_text"
    placeholder_patterns = ["[XXX]", "XXX"]
    subtitle_text = subtitle or ""
    # Determine the slide index for DAX and bail early if not found
    dax_idx = _find_dax_slide(prs)
    if dax_idx is None:
        return prs
    slide = prs.slides[dax_idx]
    # 1) Attempt to update a shape whose name exactly matches the placeholder
    for shape in slide.shapes:
        if getattr(shape, "name", "").lower() == placeholder_name:
            if shape.has_text_frame:
                runs = shape.text_frame.paragraphs[0].runs
                # capture existing font attributes if available
                attrs = _get_run_font_attributes(runs[0]) if runs else (None, None, None, None, None, None)
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                new_run = p.add_run()
                new_run.text = subtitle_text
                # apply original attributes (size, weight, italics)
                _apply_run_font_attributes(new_run, *attrs)
                # Ensure text is visible: if the original run had no RGB or was white, set to black
                try:
                    rgb = attrs[1]
                    if rgb is None or str(rgb).lower() == 'ffffff':
                        new_run.font.color.rgb = RGBColor(0, 0, 0)
                except Exception:
                    # If attrs cannot be inspected, default to black
                    new_run.font.color.rgb = RGBColor(0, 0, 0)
            return prs
    # 2) Attempt to update any shape whose name contains "text" (e.g. reused SPX placeholders)
    for shape in slide.shapes:
        name_attr = getattr(shape, "name", "").lower()
        if "text" in name_attr and shape.has_text_frame:
            runs = shape.text_frame.paragraphs[0].runs
            attrs = _get_run_font_attributes(runs[0]) if runs else (None, None, None, None, None, None)
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            new_run = p.add_run()
            new_run.text = subtitle_text
            _apply_run_font_attributes(new_run, *attrs)
            # Force text colour to black if original colour is white or undefined
            try:
                rgb = attrs[1]
                if rgb is None or str(rgb).lower() == 'ffffff':
                    new_run.font.color.rgb = RGBColor(0, 0, 0)
            except Exception:
                new_run.font.color.rgb = RGBColor(0, 0, 0)
            return prs
    # 3) Otherwise, replace placeholder patterns within the DAX slide
    for shape in slide.shapes:
        if shape.has_text_frame:
            for pattern in placeholder_patterns:
                if pattern in (shape.text or ""):
                    runs = shape.text_frame.paragraphs[0].runs
                    attrs = _get_run_font_attributes(runs[0]) if runs else (None, None, None, None, None, None)
                    new_text = shape.text.replace(pattern, subtitle_text)
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    new_run = p.add_run()
                    new_run.text = new_text
                    _apply_run_font_attributes(new_run, *attrs)
                    # Force the text colour to black if original colour is white or undefined
                    try:
                        rgb = attrs[1]
                        if rgb is None or str(rgb).lower() == 'ffffff':
                            new_run.font.color.rgb = RGBColor(0, 0, 0)
                    except Exception:
                        new_run.font.color.rgb = RGBColor(0, 0, 0)
                    return prs
    # If no placeholder found, leave slide unchanged
    return prs


###############################################################################
# Colour interpolation for gauge
###############################################################################

def _interpolate_color(value: float) -> Tuple[float, float, float]:
    """
    Interpolate from red→yellow→green for a 0–100 value.  Pure red at 0,
    bright yellow at 40 and rich green at 70.
    """
    red = (1.0, 0.0, 0.0)
    yellow = (1.0, 204 / 255, 0.0)
    green = (0.0, 153 / 255, 81 / 255)
    if value <= 40:
        t = value / 40.0
        return tuple(red[i] + t * (yellow[i] - red[i]) for i in range(3))
    elif value <= 70:
        t = (value - 40) / 30.0
        return tuple(yellow[i] + t * (green[i] - yellow[i]) for i in range(3))
    return green


def generate_average_gauge_image(
    tech_score: float,
    mom_score: float,
    last_week_avg: float,
    date_text: str | None = None,
    last_label_text: str = "Last Week",
    width_cm: float = 15.15,
    height_cm: float = 3.13,
) -> bytes:
    """
    Create a horizontal gauge with a red→yellow→green gradient, marking the
    average of technical and momentum scores against last week’s average.

    This implementation mirrors the SPX gauge: it draws a gradient bar
    spanning 0–100 and places two triangular markers with associated
    numeric labels and optional text annotations.  The current average
    (computed from the technical and momentum scores) is shown on top;
    the previous week’s average is shown below.  Colours for the
    markers are interpolated from the gradient.
    """
    def clamp100(x: float) -> float:
        return max(0.0, min(100.0, float(x)))

    curr = (clamp100(tech_score) + clamp100(mom_score)) / 2.0
    prev = clamp100(last_week_avg)

    cmap = LinearSegmentedColormap.from_list(
        "gauge_gradient", ["#FF0000", "#FFCC00", "#009951"], N=256
    )

    fig_w, fig_h = width_cm / 2.54, height_cm / 2.54
    plt.style.use("default")
    fig, ax = plt.subplots(figsize=(fig_w, fig_h))

    gradient = np.linspace(0, 1, 500).reshape(1, -1)
    bar_thickness = 0.4
    bar_bottom_y = -bar_thickness / 2.0
    bar_top_y = bar_thickness / 2.0
    ax.imshow(
        gradient,
        extent=[0, 100, bar_bottom_y, bar_top_y],
        aspect="auto",
        cmap=cmap,
        origin="lower",
    )

    # Marker dimensions and spacing
    marker_width = 3.0
    marker_height = 0.15
    gap = 0.10
    number_space = 0.25
    top_label_offset = 0.40
    bottom_label_offset = 0.40

    # Y positions for current (top) marker and labels
    top_apex_y = bar_top_y + gap
    top_base_y = top_apex_y + marker_height
    top_number_y = top_base_y + number_space
    top_label_y = top_number_y + top_label_offset

    # Y positions for previous (bottom) marker and labels
    bottom_apex_y = bar_bottom_y - gap
    bottom_base_y = bottom_apex_y - marker_height
    bottom_number_y = bottom_base_y - number_space
    bottom_label_y = bottom_number_y - bottom_label_offset

    curr_colour = _interpolate_color(curr)
    prev_colour = _interpolate_color(prev)

    # Draw triangles and numbers
    ax.add_patch(
        patches.Polygon(
            [
                (curr - marker_width / 2, top_base_y),
                (curr + marker_width / 2, top_base_y),
                (curr, top_apex_y),
            ],
            color=curr_colour,
        )
    )
    ax.add_patch(
        patches.Polygon(
            [
                (prev - marker_width / 2, bottom_base_y),
                (prev + marker_width / 2, bottom_base_y),
                (prev, bottom_apex_y),
            ],
            color=prev_colour,
        )
    )
    ax.text(
        curr,
        top_number_y,
        f"{curr:.0f}",
        color=curr_colour,
        ha="center",
        va="center",
        fontsize=8,
        fontweight="bold",
    )
    ax.text(
        prev,
        bottom_number_y,
        f"{prev:.0f}",
        color=prev_colour,
        ha="center",
        va="center",
        fontsize=8,
        fontweight="bold",
    )

    if date_text:
        ax.text(
            curr,
            top_label_y,
            date_text,
            color="#0063B0",
            ha="center",
            va="center",
            fontsize=7,
            fontweight="bold",
        )
    ax.text(
        prev,
        bottom_label_y,
        last_label_text,
        color="#133C74",
        ha="center",
        va="center",
        fontsize=7,
        fontweight="bold",
    )

    ax.set_xlim(0, 100)
    ax.set_ylim(bottom_label_y - 0.35, top_label_y + 0.35)
    ax.axis("off")

    buf = BytesIO()
    plt.savefig(buf, format="png", dpi=600, transparent=True)
    plt.close(fig)
    return buf.getvalue()


def insert_dax_average_gauge(
    prs: Presentation, excel_file, last_week_avg: float
) -> Presentation:
    """
    Insert the DAX average gauge into the DAX slide.

    The gauge shows the average of the technical and momentum scores and
    last week's average.  It is inserted into a shape named
    ``gauge_dax`` within the DAX slide.  If such a shape is not found
    within the DAX slide, placeholders ``[GAUGE]``, ``GAUGE`` or
    ``gauge_dax`` on the DAX slide are used instead.  If neither is
    present, the gauge is placed at a default position below the chart
    on the DAX slide.  Other slides remain untouched.
    """
    # Compute scores
    tech_score = _get_dax_technical_score(excel_file)
    mom_score = _get_dax_momentum_score(excel_file)
    if tech_score is None or mom_score is None:
        return prs
    # Generate the gauge image using the same dimensions as SPX
    try:
        gauge_bytes = generate_average_gauge_image(
            tech_score,
            mom_score,
            last_week_avg,
            date_text="Last",
            last_label_text="Previous Week",
            width_cm=15.15,
            height_cm=3.13,
        )
    except Exception:
        return prs
    # Identify DAX slide
    dax_idx = _find_dax_slide(prs)
    if dax_idx is None:
        return prs
    slide = prs.slides[dax_idx]
    # Define placeholder names and patterns
    placeholder_name = "gauge_dax"
    placeholder_patterns = ["[GAUGE]", "GAUGE", "gauge_dax"]
    # 1) Search for a shape named exactly gauge_dax (case-insensitive)
    for shape in slide.shapes:
        if getattr(shape, "name", "").lower() == placeholder_name:
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            # clear any existing text in the placeholder
            if shape.has_text_frame:
                try:
                    shape.text = ""
                except Exception:
                    pass
                try:
                    shape.text_frame.clear()
                except Exception:
                    pass
            # insert the gauge image scaled to the placeholder
            stream = BytesIO(gauge_bytes)
            slide.shapes.add_picture(stream, left, top, width=width, height=height)
            return prs
    # 2) Search for textual gauge placeholders within shapes on the DAX slide
    for shape in slide.shapes:
        if shape.has_text_frame:
            for pattern in placeholder_patterns:
                text_lower = (shape.text or "").lower()
                if pattern.lower() in text_lower:
                    left, top, width, height = shape.left, shape.top, shape.width, shape.height
                    # remove the placeholder text
                    try:
                        shape.text = shape.text.replace(pattern, "")
                    except Exception:
                        pass
                    # insert the gauge image scaled to the placeholder
                    stream = BytesIO(gauge_bytes)
                    slide.shapes.add_picture(stream, left, top, width=width, height=height)
                    return prs
    # 3) Additional fallback: search for any shape whose name contains 'gauge'
    for shape in slide.shapes:
        name_attr = getattr(shape, "name", "").lower()
        if "gauge" in name_attr:
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            if shape.has_text_frame:
                try:
                    shape.text = ""
                except Exception:
                    pass
                try:
                    shape.text_frame.clear()
                except Exception:
                    pass
            stream = BytesIO(gauge_bytes)
            slide.shapes.add_picture(stream, left, top, width=width, height=height)
            return prs
    # 4) Default fallback: insert below the chart using template coordinates
    left = Cm(8.97)
    top = Cm(12.13)
    width = Cm(15.15)
    height = Cm(3.13)
    stream = BytesIO(gauge_bytes)
    slide.shapes.add_picture(stream, left, top, width=width, height=height)
    return prs


###############################################################################
# Technical assessment insertion
###############################################################################

def insert_dax_technical_assessment(
    prs: Presentation,
    excel_file,
    manual_desc: Optional[str] = None,
) -> Presentation:
    """
    Insert a descriptive assessment text into the DAX slide.

    The assessment is written into a shape named ``dax_view`` on the DAX
    slide.  If no such shape exists, the function replaces any
    occurrences of ``[dax_view]`` or ``dax_view`` in text on that slide.
    A manual description may be provided; if not, the function computes
    the view from the average of the technical and momentum scores.

    Other slides remain unmodified.
    """
    # Determine the description to insert
    if manual_desc is not None and isinstance(manual_desc, str):
        desc = manual_desc.strip()
        if desc and not desc.lower().startswith("dax"):
            desc = f"DAX: {desc}"
    else:
        tech_score = _get_dax_technical_score(excel_file)
        mom_score = _get_dax_momentum_score(excel_file)
        if tech_score is None or mom_score is None:
            return prs
        avg = (float(tech_score) + float(mom_score)) / 2.0
        if avg >= 80:
            desc = "DAX: Strongly Bullish"
        elif avg >= 70:
            desc = "DAX: Bullish"
        elif avg >= 60:
            desc = "DAX: Slightly Bullish"
        elif avg >= 40:
            desc = "DAX: Neutral"
        elif avg >= 30:
            desc = "DAX: Slightly Bearish"
        elif avg >= 20:
            desc = "DAX: Bearish"
        else:
            desc = "DAX: Strongly Bearish"

    target_name = "dax_view"
    placeholder_patterns = ["[dax_view]", "dax_view"]

    dax_idx = _find_dax_slide(prs)
    if dax_idx is None:
        return prs
    slide = prs.slides[dax_idx]
    # Try to locate a shape by name on the DAX slide
    for shape in slide.shapes:
        name_attr = getattr(shape, "name", "")
        if name_attr and name_attr.lower() == target_name:
            if shape.has_text_frame:
                runs = shape.text_frame.paragraphs[0].runs
                attrs = _get_run_font_attributes(runs[0]) if runs else (None, None, None, None, None, None)
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                new_run = p.add_run()
                new_run.text = desc
                _apply_run_font_attributes(new_run, *attrs)
                # Ensure visibility: if original colour is missing or white, set black
                try:
                    rgb = attrs[1]
                    if rgb is None or str(rgb).lower() == 'ffffff':
                        new_run.font.color.rgb = RGBColor(0, 0, 0)
                except Exception:
                    new_run.font.color.rgb = RGBColor(0, 0, 0)
            return prs
    # Otherwise, replace placeholder patterns on the DAX slide
    for shape in slide.shapes:
        if shape.has_text_frame:
            for pattern in placeholder_patterns:
                if pattern.lower() in (shape.text or "").lower():
                    runs = shape.text_frame.paragraphs[0].runs
                    attrs = _get_run_font_attributes(runs[0]) if runs else (None, None, None, None, None, None)
                    try:
                        new_text = shape.text.replace(pattern, desc)
                    except Exception:
                        new_text = desc
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    new_run = p.add_run()
                    new_run.text = new_text
                    _apply_run_font_attributes(new_run, *attrs)
                    # Ensure visible colour
                    try:
                        rgb = attrs[1]
                        if rgb is None or str(rgb).lower() == 'ffffff':
                            new_run.font.color.rgb = RGBColor(0, 0, 0)
                    except Exception:
                        new_run.font.color.rgb = RGBColor(0, 0, 0)
                    return prs
    return prs


###############################################################################
# Source footnote insertion
###############################################################################

def insert_dax_source(
    prs: Presentation,
    used_date: Optional[pd.Timestamp],
    price_mode: str,
) -> Presentation:
    """
    Insert the source footnote into a shape named 'dax_source' (or
    containing '[dax_source]').  The footnote text depends on the selected
    price mode.  For example:

      * Last Close  → "Source: Bloomberg, Herculis Group, Data as of 29/07/2025 Close"
      * Last Price  → "Source: Bloomberg, Herculis Group, Data as of 29/07/2025"

    Parameters
    ----------
    prs : Presentation
        The PowerPoint presentation to modify.
    used_date : pandas.Timestamp or None
        The date that should appear in the footnote.  If ``None``, no
        changes are made.
    price_mode : str
        Either 'Last Price' or 'Last Close'.  Determines whether the
        suffix " Close" is appended to the date.

    Returns
    -------
    Presentation
        The modified presentation.
    """
    if used_date is None:
        return prs
    try:
        date_str = used_date.strftime("%d/%m/%Y")
    except Exception:
        return prs
    suffix = " Close" if str(price_mode).lower() == "last close" else ""
    source_text = f"Source: Bloomberg, Herculis Group, Data as of {date_str}{suffix}"
    placeholder_name = "dax_source"
    placeholder_patterns = ["[dax_source]", "dax_source"]
    # Restrict insertion to the DAX slide only
    dax_idx = _find_dax_slide(prs)
    if dax_idx is None:
        return prs
    slide = prs.slides[dax_idx]
    # Case 1: replace a shape named exactly as the placeholder
    for shape in slide.shapes:
        name_attr = getattr(shape, "name", "")
        if name_attr and name_attr.lower() == placeholder_name:
            if shape.has_text_frame:
                runs = shape.text_frame.paragraphs[0].runs
                attrs = _get_run_font_attributes(runs[0]) if runs else (None, None, None, None, None, None)
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                new_run = p.add_run()
                new_run.text = source_text
                _apply_run_font_attributes(new_run, *attrs)
            return prs
    # Case 2: replace occurrences of the placeholder pattern in text on the DAX slide
    for shape in slide.shapes:
        if shape.has_text_frame:
            for pattern in placeholder_patterns:
                if pattern.lower() in (shape.text or "").lower():
                    runs = shape.text_frame.paragraphs[0].runs
                    attrs = _get_run_font_attributes(runs[0]) if runs else (None, None, None, None, None, None)
                    try:
                        new_text = (shape.text or "").replace(pattern, source_text)
                    except Exception:
                        new_text = source_text
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    new_run = p.add_run()
                    new_run.text = new_text
                    _apply_run_font_attributes(new_run, *attrs)
                    return prs
    return prs


###############################################################################
# Range gauge helpers and insertion
###############################################################################

def _compute_range_bounds(
    df_full: pd.DataFrame, lookback_days: int = 90
) -> Tuple[float, float]:
    """
    Compute fallback high and low range bounds for the DAX using
    realised volatility.

    This helper is used when an implied volatility index (e.g. V1X) is
    unavailable.  It computes the annualised realised volatility over a
    30‑session window by taking the standard deviation of daily
    percentage returns, multiplying by ``sqrt(252)`` and converting to
    a percentage.  The resulting 1‑week expected move is
    ``(current_price × (realised_vol / 100)) / sqrt(52)``.  The upper
    and lower bounds are the current price plus and minus this
    expected move.  If realised volatility cannot be computed or is
    zero, the function falls back to a ±2 % band around the current
    price.

    Parameters
    ----------
    df_full : pandas.DataFrame
        DataFrame containing at least 'Date' and 'Price' columns,
        sorted by date ascending.
    lookback_days : int, optional
        Number of trading days used to compute the approximate true range
        if realised volatility is unavailable.  Currently unused but
        retained for API compatibility.

    Returns
    -------
    Tuple[float, float]
        A two‑tuple ``(upper_bound, lower_bound)`` representing the
        current closing price plus and minus the realised volatility
        based expected move, or ±2 % of the current price if no
        volatility can be computed.
    """
    if df_full.empty:
        return (np.nan, np.nan)
    current_price = df_full["Price"].iloc[-1]
    # Attempt to compute 30‑day realised volatility (annualised) as a fallback.  Use
    # the last 30 trading days of closing prices to compute daily returns.
    # If the realised volatility can be computed, convert it into a 1‑week
    # expected move.  Otherwise fall back to a ±2 % band.
    try:
        # At least 2 data points are needed for pct_change; ensure there are
        # enough rows (we use min to handle shorter histories gracefully).
        lookback = 30
        window_prices = df_full["Price"].tail(lookback)
        # Compute daily percentage returns
        rets = window_prices.pct_change().dropna()
        # Standard deviation of daily returns
        std_daily = rets.std()
        if std_daily is not None and not np.isnan(std_daily) and std_daily > 0:
            # Annualise the standard deviation (multiply by sqrt(252)) and convert to %
            realised_vol = std_daily * np.sqrt(252.0) * 100.0
            # Convert to 1‑week expected move by dividing by sqrt(52)
            expected_move = (current_price * (realised_vol / 100.0)) / np.sqrt(52.0)
            upper_bound = current_price + expected_move
            lower_bound = current_price - expected_move
            return (float(upper_bound), float(lower_bound))
    except Exception:
        pass
    # Fallback: ±2 % of the current price
    return (float(current_price * 1.02), float(current_price * 0.98))


def generate_range_gauge_chart_image(
    df_full: pd.DataFrame,
    anchor_date: Optional[pd.Timestamp] = None,
    lookback_days: int = 90,
    width_cm: float = 21.41,
    height_cm: float = 7.53,
    chart_width_cm: float = None,
    gauge_width_cm: float = 4.0,
    *,
    vol_index_value: Optional[float] = None,
) -> bytes:
    """
    Create a PNG image of the DAX price chart with a vertical range gauge
    appended on the right.  The gauge shows a green–to–red gradient between
    recent high and support levels, with labels for the upper and lower
    bounds.  A horizontal line continues the last price into the gauge so
    that viewers can assess relative positioning.  This function is used by
    ``insert_dax_technical_chart_with_range``.

    Parameters
    ----------
    df_full : pandas.DataFrame
        Full DAX price history.
    anchor_date : pandas.Timestamp or None, optional
        Optional anchor date for a regression channel.
    lookback_days : int, default 90
        Lookback window for computing the high/low range.
    width_cm : float, default 21.41
        Total width of the output image in centimetres.
    height_cm : float, default 7.53
        Height of the output image in centimetres.
    chart_width_cm : float, optional
        Width of the price chart portion; if None, uses width_cm − gauge_width_cm.
    gauge_width_cm : float, default 4.0
        Width of the gauge portion.
    vol_index_value : float or None, optional
        Optional implied volatility index value used to estimate expected
        move; if None, uses realised volatility from the price data.

    Returns
    -------
    bytes
        PNG image bytes with transparency.
    """
    import matplotlib.pyplot as plt
    import matplotlib.patches as mpatches

    if df_full.empty:
        return b""

    # Use the last year of data for plotting the price chart
    today = df_full["Date"].max().normalize()
    start = today - timedelta(days=365)
    df = df_full[df_full["Date"].between(start, today)].reset_index(drop=True)
    df_ma = _add_mas(df)

    # Optional regression channel for trend shading
    uptrend = False
    upper_channel = lower_channel = None
    if anchor_date is not None:
        subset_full = df_full[df_full["Date"].between(anchor_date, today)].copy()
        if not subset_full.empty:
            X = subset_full["Date"].map(pd.Timestamp.toordinal).to_numpy().reshape(-1, 1)
            y_vals = subset_full["Price"].to_numpy()
            model = LinearRegression().fit(X, y_vals)
            trend = model.predict(X)
            resid = y_vals - trend
            uptrend = model.coef_[0] > 0
            upper_channel = trend + resid.max()
            lower_channel = trend + resid.min()

    # Last price and range bounds
    last_price = df["Price"].iloc[-1]
    if vol_index_value is not None and last_price and not np.isnan(last_price):
        expected_move = (last_price * (vol_index_value / 100.0)) / np.sqrt(52.0)
        upper_bound = last_price + expected_move
        lower_bound = last_price - expected_move
    else:
        upper_bound, lower_bound = _compute_range_bounds(df_full, lookback_days=lookback_days)
    # Minimum band of ±1% around current price
    min_range_pct = 0.02
    if last_price and not np.isnan(last_price):
        range_span_pct = (upper_bound - lower_bound) / last_price if last_price else 0.0
        if range_span_pct < min_range_pct:
            half_span = (min_range_pct * last_price) / 2.0
            upper_bound = last_price + half_span
            lower_bound = last_price - half_span
        up_pct = (upper_bound - last_price) / last_price * 100.0
        down_pct = (last_price - lower_bound) / last_price * 100.0
    else:
        up_pct = down_pct = 0.0

    # Determine axis limits
    hi = df["Price"].max()
    lo = df["Price"].min()
    y_max = max(hi, upper_bound) * 1.02
    y_min = min(lo, lower_bound) * 0.98

    # Compute chart and gauge widths
    if chart_width_cm is None:
        chart_width_cm = max(width_cm - gauge_width_cm, 0.0)
    fig_w_in, fig_h_in = width_cm / 2.54, height_cm / 2.54
    fig = plt.figure(figsize=(fig_w_in, fig_h_in))

    chart_rel_width = chart_width_cm / width_cm if width_cm > 0 else 0.0
    gauge_rel_width = gauge_width_cm / width_cm if width_cm > 0 else 0.0
    margin_rel = min(0.04, 0.10 * chart_rel_width)

    ax_chart = fig.add_axes([margin_rel, 0.0, chart_rel_width - margin_rel, 1.0])
    ax_gauge = fig.add_axes([chart_rel_width, 0.0, gauge_rel_width, 1.0])

    # Set y-limits
    ax_chart.set_ylim(y_min, y_max)
    ax_gauge.set_ylim(y_min, y_max)

    # Plot price and moving averages
    ax_chart.plot(df["Date"], df["Price"], color="#153D64", linewidth=2.5,
                  label=f"DAX Price (last: {last_price:,.2f})")
    ax_chart.plot(df_ma["Date"], df_ma["MA_50"], color="#008000", linewidth=1.5, label="50‑day MA")
    ax_chart.plot(df_ma["Date"], df_ma["MA_100"], color="#FFA500", linewidth=1.5, label="100‑day MA")
    ax_chart.plot(df_ma["Date"], df_ma["MA_200"], color="#FF0000", linewidth=1.5, label="200‑day MA")
    # Fibonacci lines on subset
    sub_hi, sub_lo = df["Price"].max(), df["Price"].min()
    sub_span = sub_hi - sub_lo
    for lvl in [sub_hi, sub_hi - 0.236 * sub_span, sub_hi - 0.382 * sub_span,
                sub_hi - 0.5 * sub_span, sub_hi - 0.618 * sub_span, sub_lo]:
        ax_chart.axhline(lvl, color="grey", linestyle="--", linewidth=0.8, alpha=0.6)

    # Regression channel shading
    if anchor_date is not None and upper_channel is not None and lower_channel is not None:
        subset = df_full[df_full["Date"].between(anchor_date, today)].copy().reset_index(drop=True)
        fill_color = (0, 0.6, 0, 0.25) if uptrend else (0.78, 0, 0, 0.25)
        line_color = "#008000" if uptrend else "#C00000"
        ax_chart.plot(subset["Date"], upper_channel, color=line_color, linestyle="--")
        ax_chart.plot(subset["Date"], lower_channel, color=line_color, linestyle="--")
        ax_chart.fill_between(subset["Date"], lower_channel, upper_channel, color=fill_color)

    # Style chart axes
    for spine in ax_chart.spines.values():
        spine.set_visible(False)
    ax_chart.tick_params(axis="y", which="both", length=0)
    ax_chart.tick_params(axis="x", which="both", length=2)
    ax_chart.tick_params(left=True, bottom=True, labelleft=True, labelbottom=True)
    ax_chart.legend(loc="upper left", bbox_to_anchor=(0.0, 1.05), ncol=4,
                    fontsize=8, frameon=False)

    # Draw gauge gradient
    n_steps = 1000
    y_vals = np.linspace(lower_bound, upper_bound, n_steps)
    # Map each y to a color: green at top, red at bottom
    colors = plt.cm.RdYlGn((y_vals - lower_bound) / (upper_bound - lower_bound))
    for y, c in zip(y_vals, colors):
        ax_gauge.axhline(y, color=c, linewidth=2)
    # Draw current price line across gauge
    ax_gauge.hlines(last_price, xmin=0, xmax=1, colors="black", linewidth=1.5)
    # Add markers and labels for high/low
    ax_gauge.hlines(upper_bound, xmin=0, xmax=0.3, colors="#009951", linewidth=2)
    ax_gauge.hlines(lower_bound, xmin=0, xmax=0.3, colors="#C00000", linewidth=2)
    # Remove gauge axes ticks and spines
    ax_gauge.set_xticks([])
    ax_gauge.set_yticks([])
    for spine in ax_gauge.spines.values():
        spine.set_visible(False)

    # Add labels on gauge
    ax_gauge.text(0.5, upper_bound, f"{upper_bound:,.0f} (+{up_pct:.1f}%)", ha="left", va="center",
                  color="#009951", fontsize=7, transform=ax_gauge.transData)
    ax_gauge.text(0.5, lower_bound, f"{lower_bound:,.0f} (-{down_pct:.1f}%)", ha="left", va="center",
                  color="#C00000", fontsize=7, transform=ax_gauge.transData)

    buf = BytesIO()
    fig.savefig(buf, format="png", dpi=600, transparent=True,
                bbox_inches="tight", pad_inches=0.05)
    plt.close(fig)
    buf.seek(0)
    return buf.getvalue()


def insert_dax_technical_chart_with_range(
    prs: Presentation,
    excel_file,
    anchor_date: Optional[pd.Timestamp] = None,
    lookback_days: int = 90,
    price_mode: str = "Last Price",
) -> Presentation:
    """
    Insert the DAX technical analysis chart with the trading range gauge
    into the PowerPoint.  The range gauge appears on the right of the
    chart, with a gradient from green (high) to red (low) and labels
    denoting the upper and lower bounds and their percentage moves.

    Parameters
    ----------
    prs : Presentation
        The PowerPoint presentation to modify.
    excel_file : file-like object or path
        Excel workbook containing DAX price data.
    anchor_date : pandas.Timestamp or None, optional
        Optional anchor date for a regression channel.
    lookback_days : int, default 90
        Lookback window for computing the high/low range.
    price_mode : str, default "Last Price"
        Determines whether the most recent date is excluded for Last Close mode.

    Returns
    -------
    Presentation
        The modified presentation.
    """
    # Load price data
    try:
        df_full = _load_price_data_from_obj(excel_file, "DAX Index", price_mode=price_mode)
    except Exception:
        df_full = _load_price_data(pathlib.Path(excel_file), "DAX Index", price_mode=price_mode)

    # Retrieve implied volatility index (V1X) value from the file
    vol_val = _get_vol_index_value(excel_file, price_mode=price_mode, vol_ticker="V1X Index")
    # Generate combined chart and gauge image
    img_bytes = generate_range_gauge_chart_image(
        df_full,
        anchor_date=anchor_date,
        lookback_days=lookback_days,
        width_cm=25.0,
        height_cm=7.3,
        chart_width_cm=21.41,
        gauge_width_cm=3.59,
        vol_index_value=vol_val,
    )

    # Locate the DAX slide
    target_slide = None
    for slide in prs.slides:
        for shape in slide.shapes:
            name_attr = getattr(shape, "name", "").lower()
            if name_attr == "dax":
                target_slide = slide
                break
            if shape.has_text_frame:
                if (shape.text or "").strip().lower() == "[dax]":
                    target_slide = slide
                    break
        if target_slide:
            break
    if target_slide is None:
        target_slide = prs.slides[min(11, len(prs.slides) - 1)]

    # Insert image at specified coordinates (0.93 cm left, 4.80 cm top)
    left = Cm(0.93)
    top = Cm(4.80)
    width = Cm(25.0)
    height = Cm(7.3)
    stream = BytesIO(img_bytes)
    picture = target_slide.shapes.add_picture(stream, left, top, width=width, height=height)
    try:
        sp_tree = target_slide.shapes._spTree
        sp_tree.remove(picture._element)
        sp_tree.insert(1, picture._element)
    except Exception:
        pass
    return prs


def _load_price_data_from_obj(
    excel_obj,
    ticker: str = "DAX Index",
    price_mode: str = "Last Price",
) -> pd.DataFrame:
    """
    Load price data from a file-like object and return a tidy DataFrame.

    Parameters
    ----------
    excel_obj : file-like
        File-like object representing an Excel workbook containing a
        ``data_prices`` sheet.
    ticker : str, default "DAX Index"
        Column name corresponding to the desired ticker in the Excel sheet.
    price_mode : str, default "Last Price"
        One of "Last Price" or "Last Close".  If ``adjust_prices_for_mode``
        is available and the mode is "Last Close", rows corresponding to
        the most recent date (if equal to today's date) will be dropped.

    Returns
    -------
    pandas.DataFrame
        DataFrame with columns ``Date`` and ``Price``.  The data are
        sorted by date and any rows with missing values are removed.
    """
    df = pd.read_excel(excel_obj, sheet_name="data_prices")
    df = df.drop(index=0)
    df = df[df[df.columns[0]] != "DATES"]
    df["Date"] = pd.to_datetime(df[df.columns[0]], errors="coerce")
    df["Price"] = pd.to_numeric(df[ticker], errors="coerce")
    df_clean = (
        df.dropna(subset=["Date", "Price"]).sort_values("Date").reset_index(drop=True)[
            ["Date", "Price"]
        ]
    )
    # Adjust for price mode if helper is available
    if adjust_prices_for_mode is not None and price_mode:
        try:
            df_clean, _ = adjust_prices_for_mode(df_clean, price_mode)
        except Exception:
            pass
    return df_clean