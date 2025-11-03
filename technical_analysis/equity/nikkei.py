"""
Utility functions for NIKKEI technical analysis and high‑resolution export.

This module provides tools to build interactive and static charts for the
NIKKEI index, calculate and insert technical and momentum scores into
PowerPoint presentations, generate horizontal and vertical gauges that
visualise the average of the technical and momentum scores, as well as
contextual trading ranges (higher and lower range bounds).  Functions
fall back to sensible defaults when placeholders are not found.

Key functions include:

* ``make_nikkei_figure`` – interactive Plotly chart for Streamlit.
* ``insert_nikkei_technical_chart`` – insert a static NIKKEI chart into a PPTX.
* ``insert_nikkei_technical_score_number`` – insert the technical score (integer).
* ``insert_nikkei_momentum_score_number`` – insert the momentum score (integer).
* ``insert_nikkei_subtitle`` – insert a user‑defined subtitle into the NIKKEI slide.
* ``generate_average_gauge_image`` – create a horizontal gauge image.
* ``insert_nikkei_average_gauge`` – insert the gauge into a PPT slide.
* ``insert_nikkei_technical_assessment`` – insert a descriptive “view” text.
* ``generate_range_gauge_chart_image`` – create a combined price chart with
  a vertical range gauge on the right hand side, including a horizontal line
  connecting the last price to the gauge.  This function is used by
  ``insert_nikkei_technical_chart_with_range``.
* ``insert_nikkei_technical_chart_with_range`` – insert the NIKKEI technical
  analysis chart with the higher/lower range gauge into the PPT.

The range gauge illustrates the recent trading range for the NIKKEI.
Instead of using the absolute high and low closes of the last 90 days,
the bounds are estimated from recent volatility.  Whenever possible the
code looks up the forward‑looking volatility index (VIX) and computes a
1‑week expected move as ``(current_price × (VIX / 100)) / sqrt(52)``.
The upper and lower bounds are the current price plus and minus that
expected move.  If the volatility index is unavailable, the code falls
back to using realised volatility: it computes the standard deviation of
30‑day daily returns, annualises it and converts it to a 1‑week move
using the same formula.  As a last resort when neither volatility
measure can be computed the bounds default to ±2 % of the current price.
A minimum ±1 % band is enforced to avoid overlapping annotations when
volatility is extremely low.  A horizontal line continues the last
price through to the gauge so that viewers can quickly assess how far
the index lies from its typical volatility band.  This method
provides context on potential upside and downside moves based on
recent price variability rather than simply the most recent highs and
lows.
"""

from __future__ import annotations

from datetime import timedelta
import pathlib
from typing import Optional, Tuple

import pandas as pd
import plotly.graph_objects as go
from sklearn.linear_model import LinearRegression

from pptx import Presentation

# Import common helpers (eliminates code duplication)
from technical_analysis.common_helpers import (
    _get_run_font_attributes,
    _apply_run_font_attributes,
    _add_mas,
    _get_technical_score_generic,
    _get_momentum_score_generic,
    _interpolate_color,
    _load_price_data_from_obj,
    _compute_range_bounds,
    generate_range_callout_chart_image,
)
from technical_analysis.powerpoint_utils import (
    find_slide_by_placeholder,
    insert_score_number,
    insert_subtitle,
)
from pptx.util import Cm
from io import BytesIO
import matplotlib.pyplot as plt
import matplotlib.patches as patches
from matplotlib.colors import LinearSegmentedColormap
import numpy as np

# Import helper for adjusting price data according to price mode.  The utils
# module must reside at the project root.  It is deliberately not imported
# from ``technical_analysis.utils`` so that this module can be reused when
# nested within ``ic`` and ``utils.py`` sits at the project root.
try:
    from utils import adjust_prices_for_mode  # type: ignore
except Exception:
    # If the utils module is not available, define a no-op fallback.  This
    # preserves compatibility with environments where price mode is not used.
    adjust_prices_for_mode = None  # type: ignore

# Default lookback window (in days) for plotting.  The app can override
# this value at runtime by setting the module-level ``PLOT_LOOKBACK_DAYS``
# attribute.  We use 90 days (approximately 3 months) by default to
# align with the updated requirement from management.  When the user
# selects a different timeframe (e.g. 6 months), ``app.py`` will
# temporarily override this constant to 180 days.
PLOT_LOOKBACK_DAYS: int = 90

###############################################################################
# Internal helpers
###############################################################################





def _load_price_data(
    excel_path: pathlib.Path,
    ticker: str = "NKY Index",
    price_mode: str = "Last Price",
) -> pd.DataFrame:
    """
    Read the raw price sheet and return a tidy Date‑Price DataFrame.

    Parameters
    ----------
    excel_path : pathlib.Path
        Path to the Excel workbook containing price data.
    ticker : str, default "NKY Index"
        Column name corresponding to the desired ticker in the Excel sheet.
    price_mode : str, default "Last Price"
        One of "Last Price" or "Last Close".  If ``adjust_prices_for_mode``
        is available and the mode is "Last Close", rows with the last
        recorded date (if equal to today's date) will be dropped.

    Returns
    -------
    pandas.DataFrame
        DataFrame with columns ``Date`` and ``Price``.  The data are
        sorted by date and any rows with missing values are removed.
    """
    df = pd.read_excel(excel_path, sheet_name="data_prices")
    df = df.drop(index=0)
    df = df[df[df.columns[0]] != "DATES"]
    df["Date"] = pd.to_datetime(df[df.columns[0]], errors="coerce")
    df["Price"] = pd.to_numeric(df[ticker], errors="coerce")
    df_clean = (
        df.dropna(subset=["Date", "Price"]).sort_values("Date").reset_index(drop=True)[
            ["Date", "Price"]
        ]
    )
    # Adjust for price mode if helper is available
    if adjust_prices_for_mode is not None and price_mode:
        try:
            df_clean, _ = adjust_prices_for_mode(df_clean, price_mode)
        except Exception:
            # If adjustment fails, silently fall back to unadjusted data
            pass
    return df_clean




def _get_vol_index_value(
    excel_obj_or_path,
    price_mode: str = "Last Price",
    vol_ticker: str = "VIX Index",
) -> Optional[float]:
    """
    Retrieve the most recent value of a volatility index (e.g. VIX) from
    the ``data_prices`` sheet.  If ``price_mode`` is ``"Last Close"``,
    the most recent date is dropped if it matches today's date.  The
    returned value is the last available entry after price‑mode adjustment.

    Parameters
    ----------
    excel_obj_or_path : file‑like or path
        The Excel workbook containing price data.
    price_mode : str, default "Last Price"
        One of "Last Price" or "Last Close".  When set to "Last Close"
        rows corresponding to the most recent date (if equal to today's
        date) will be excluded before taking the last value.
    vol_ticker : str, default "VIX Index"
        Column name in the ``data_prices`` sheet corresponding to the
        volatility index whose level should be used.

    Returns
    -------
    float or None
        The most recent volatility index value, or ``None`` if it cannot
        be read or converted to a float.
    """
    try:
        df = pd.read_excel(excel_obj_or_path, sheet_name="data_prices")
    except Exception:
        return None
    # Drop first row and metadata rows
    df = df.drop(index=0)
    df = df[df[df.columns[0]] != "DATES"]
    # Parse dates and the volatility index column
    df["Date"] = pd.to_datetime(df[df.columns[0]], errors="coerce")
    # Ensure the volatility index column exists
    if vol_ticker not in df.columns:
        return None
    df["Price"] = pd.to_numeric(df[vol_ticker], errors="coerce")
    df_clean = df.dropna(subset=["Date", "Price"]).sort_values("Date").reset_index(drop=True)[
        ["Date", "Price"]
    ]
    # Apply price mode adjustment if possible
    if adjust_prices_for_mode is not None and price_mode:
        try:
            df_clean, _ = adjust_prices_for_mode(df_clean, price_mode)
        except Exception:
            pass
    if df_clean.empty:
        return None
    try:
        return float(df_clean["Price"].iloc[-1])
    except Exception:
        return None


###############################################################################
# Plotly interactive chart for Streamlit
###############################################################################

def make_nikkei_figure(
    excel_path: str | pathlib.Path,
    anchor_date: Optional[pd.Timestamp] = None,
    price_mode: str = "Last Price",
) -> go.Figure:
    """
    Build an interactive NIKKEI chart for Streamlit.

    Parameters
    ----------
    excel_path : str or pathlib.Path
        Path to the Excel file containing NIKKEI price data.
    anchor_date : pandas.Timestamp or None, optional
        If provided, a regression channel is drawn from ``anchor_date`` to the
        latest date.
    price_mode : str, default "Last Price"
        One of "Last Price" or "Last Close".  When set to "Last Close"
        and if the ``utils.adjust_prices_for_mode`` helper is available,
        rows corresponding to the most recent date (if equal to today's
        date) are excluded from the chart.  This allows the chart to
        represent the prior day's closing prices rather than an
        intraday or current price.

    Returns
    -------
    go.Figure
        A Plotly figure with price, moving averages, Fibonacci lines and
        an optional regression channel.
    """
    excel_path = pathlib.Path(excel_path)
    # Load data and adjust according to the price mode
    df_raw = _load_price_data(excel_path, "NKY Index", price_mode=price_mode)
    df_full = _add_mas(df_raw)

    if df_full.empty:
        return go.Figure()

    today = df_full["Date"].max().normalize()
    start = today - timedelta(days=PLOT_LOOKBACK_DAYS)
    df = df_full[df_full["Date"].between(start, today)].reset_index(drop=True)

    if df.empty:
        return go.Figure()

    last_price = df["Price"].iloc[-1]
    last_price_str = f"{last_price:,.2f}"

    fig = go.Figure()
    fig.add_trace(
        go.Scatter(
            x=df["Date"],
            y=df["Price"],
            mode="lines",
            name=f"NIKKEI Price (last: {last_price_str})",
            line=dict(color="#153D64", width=2.5),
        )
    )
    fig.add_trace(
        go.Scatter(
            x=df["Date"],
            y=df["MA_50"],
            mode="lines",
            name="50‑day MA",
            line=dict(color="#008000", width=1.5),
        )
    )
    fig.add_trace(
        go.Scatter(
            x=df["Date"],
            y=df["MA_100"],
            mode="lines",
            name="100‑day MA",
            line=dict(color="#FFA500", width=1.5),
        )
    )
    fig.add_trace(
        go.Scatter(
            x=df["Date"],
            y=df["MA_200"],
            mode="lines",
            name="200‑day MA",
            line=dict(color="#FF0000", width=1.5),
        )
    )

    hi, lo = df["Price"].max(), df["Price"].min()
    span = hi - lo
    for lvl in [
        hi,
        hi - 0.236 * span,
        hi - 0.382 * span,
        hi - 0.5 * span,
        hi - 0.618 * span,
        lo,
    ]:
        fig.add_hline(
            y=lvl, line=dict(color="grey", dash="dash", width=1), opacity=0.6
        )

    if anchor_date is not None:
        per = df_full[df_full["Date"].between(anchor_date, today)].copy()
        if not per.empty:
            X = per["Date"].map(pd.Timestamp.toordinal).to_numpy().reshape(-1, 1)
            y_vals = per["Price"].to_numpy()
            model = LinearRegression().fit(X, y_vals)
            trend = model.predict(X)
            resid = y_vals - trend
            upper = trend + resid.max()
            lower = trend + resid.min()

            uptrend = model.coef_[0] > 0
            lineclr = "green" if uptrend else "red"
            fillclr = "rgba(0,150,0,0.25)" if uptrend else "rgba(200,0,0,0.25)"

            fig.add_trace(
                go.Scatter(
                    x=per["Date"],
                    y=upper,
                    mode="lines",
                    line=dict(color=lineclr, dash="dash"),
                    showlegend=False,
                )
            )
            fig.add_trace(
                go.Scatter(
                    x=per["Date"],
                    y=lower,
                    mode="lines",
                    line=dict(color=lineclr, dash="dash"),
                    fill="tonexty",
                    fillcolor=fillclr,
                    showlegend=False,
                )
            )

    fig.update_layout(
        margin=dict(l=30, r=30, t=60, b=40),
        showlegend=True,
        legend=dict(
            orientation="h",
            yanchor="bottom",
            y=1.12,
            xanchor="center",
            x=0.5,
            font=dict(size=12),
        ),
        xaxis_title=None,
        yaxis_title=None,
        xaxis=dict(showgrid=False),
        yaxis=dict(showgrid=False, zeroline=False),
    )
    return fig


###############################################################################
# High‑resolution chart export (PNG)
###############################################################################

def _generate_nikkei_image_from_df(
    df_full: pd.DataFrame,
    anchor_date: Optional[pd.Timestamp],
    width_cm: float = 21.41,
    height_cm: float = 7.53,
) -> bytes:
    """
    Create a high‑resolution (dpi=300) transparent PNG chart from the DataFrame.
    Includes price, moving averages, Fibonacci lines and optional regression channel.
    """
    today = df_full["Date"].max().normalize()
    start = today - timedelta(days=PLOT_LOOKBACK_DAYS)
    df = df_full[df_full["Date"].between(start, today)].reset_index(drop=True)

    df_ma = df.copy()
    for w in (50, 100, 200):
        df_ma[f"MA_{w}"] = df_ma["Price"].rolling(w, min_periods=1).mean()

    uptrend = False
    upper = lower = None
    if anchor_date is not None:
        subset = df_full[df_full["Date"].between(anchor_date, today)].copy()
        if not subset.empty:
            X = subset["Date"].map(pd.Timestamp.toordinal).to_numpy().reshape(-1, 1)
            y_vals = subset["Price"].to_numpy()
            model = LinearRegression().fit(X, y_vals)
            trend = model.predict(X)
            resid = y_vals - trend
            uptrend = model.coef_[0] > 0
            upper = trend + resid.max()
            lower = trend + resid.min()

    last_price = df["Price"].iloc[-1]
    last_price_str = f"{last_price:,.2f}"

    fig_width_in, fig_height_in = width_cm / 2.54, height_cm / 2.54
    plt.style.use("default")
    fig, ax = plt.subplots(figsize=(fig_width_in, fig_height_in))

    ax.plot(
        df["Date"],
        df["Price"],
        color="#153D64",
        linewidth=2.5,
        label=f"NIKKEI Price (last: {last_price_str})",
    )
    ax.plot(
        df_ma["Date"],
        df_ma["MA_50"],
        color="#008000",
        linewidth=1.5,
        label="50‑day MA",
    )
    ax.plot(
        df_ma["Date"],
        df_ma["MA_100"],
        color="#FFA500",
        linewidth=1.5,
        label="100‑day MA",
    )
    ax.plot(
        df_ma["Date"],
        df_ma["MA_200"],
        color="#FF0000",
        linewidth=1.5,
        label="200‑day MA",
    )

    hi, lo = df["Price"].max(), df["Price"].min()
    span = hi - lo
    fib_levels = [
        hi,
        hi - 0.236 * span,
        hi - 0.382 * span,
        hi - 0.5 * span,
        hi - 0.618 * span,
        lo,
    ]
    for lvl in fib_levels:
        ax.axhline(y=lvl, color="grey", linestyle="--", linewidth=0.8, alpha=0.6)

    if anchor_date is not None and upper is not None and lower is not None:
        fill_color = (0, 0.6, 0, 0.25) if uptrend else (0.78, 0, 0, 0.25)
        line_color = "#008000" if uptrend else "#C00000"
        subset = (
            df_full[df_full["Date"].between(anchor_date, today)].copy().reset_index(drop=True)
        )
        ax.plot(subset["Date"], upper, color=line_color, linestyle="--")
        ax.plot(subset["Date"], lower, color=line_color, linestyle="--")
        ax.fill_between(subset["Date"], lower, upper, color=fill_color)

    for spine in ax.spines.values():
        spine.set_visible(False)
    ax.tick_params(left=True, bottom=True, labelleft=True, labelbottom=True)

    ax.legend(
        loc="upper center", bbox_to_anchor=(0.5, 1.1), ncol=4, fontsize=8, frameon=False
    )
    plt.tight_layout()

    buf = BytesIO()
    plt.savefig(buf, format="png", dpi=600, transparent=True)
    plt.close(fig)
    buf.seek(0)
    return buf.getvalue()


###############################################################################
# Score helpers
###############################################################################

def _get_nikkei_technical_score(excel_obj_or_path) -> Optional[float]:
    """
    Retrieve the technical score for NIKKEI.
    Uses common helper with instrument-specific ticker.
    """
    return _get_technical_score_generic(excel_obj_or_path, "NKY INDEX")



def _find_nikkei_slide(prs: Presentation) -> Optional[int]:
    """Find the NIKKEI slide by placeholder."""
    return find_slide_by_placeholder(prs, "nikkei")



def insert_nikkei_technical_score_number(prs: Presentation, excel_file) -> Presentation:
    """Insert the NIKKEI technical score into the slide."""
    score = _get_nikkei_technical_score(excel_file)
    return insert_score_number(prs, score, "nikkei", "tech_score")


###############################################################################
# Call‑out range helpers and insertion
###############################################################################



def insert_nikkei_technical_chart_with_callout(
    prs: Presentation,
    excel_file,
    anchor_date: Optional[pd.Timestamp] = None,
    lookback_days: int = 90,
    price_mode: str = "Last Price",
) -> Presentation:
    """
    Insert the NIKKEI technical analysis chart with the trading range call‑out
    into the PowerPoint.  This function mirrors the behaviour of
    ``insert_nikkei_technical_chart_with_range`` but uses the call‑out style to
    display the high and low bounds instead of a vertical gauge.

    The image is placed at the fixed coordinates (0.93 cm left, 5.46 cm top)
    with dimensions 24.2 cm wide by 6.52 cm high.  These values match those
    used on the NIKKEI slide and leave room above for a separate legend on the
    PowerPoint slide.  When inserting into the presentation the legend is
    suppressed in the image itself so that it can be added manually.

    Parameters
    ----------
    prs : Presentation
        The PowerPoint presentation to modify.
    excel_file : file‑like object or path
        Excel workbook containing NIKKEI price data.
    anchor_date : pandas.Timestamp or None, optional
        Optional anchor date for a regression channel.
    lookback_days : int, default 90
        Lookback window for computing the high/low range.

    Returns
    -------
    Presentation
        The presentation with the updated slide.
    """
    # Load the price data from the Excel file
    try:
        df_full = _load_price_data_from_obj(excel_file, "NKY Index", price_mode=price_mode)
    except Exception:
        df_full = _load_price_data(pathlib.Path(excel_file), "NKY Index", price_mode=price_mode)

    # For the NIKKEI index there is no commonly used implied volatility index.
    # Do not attempt to read a volatility index from the Excel file.  The
    # range calculation in ``generate_range_callout_chart_image`` will
    # automatically fall back to realised volatility when the
    # ``vol_index_value`` is ``None``.
    vol_val = None
    # Generate the image with the call‑out.  Use a width of 24.2 cm and a
    # height of 6.52 cm (matching the NIKKEI template) so that there is
    # sufficient space above the chart for an external legend.  Pass
    # ``show_legend=False`` to suppress the internal legend on the figure.
    img_bytes = generate_range_callout_chart_image(
        df_full,
        anchor_date=anchor_date,
        lookback_days=lookback_days,
        width_cm=24.2,
        height_cm=6.52,
        vol_index_value=vol_val,
        show_legend=False,
    )

    # Locate the slide containing the 'nikkei' placeholder or text
    target_slide = None
    for slide in prs.slides:
        for shape in slide.shapes:
            name_attr = getattr(shape, "name", "").lower()
            if name_attr == "nikkei":
                target_slide = slide
                break
            if shape.has_text_frame:
                if (shape.text or "").strip().lower() == "[nikkei]":
                    target_slide = slide
                    break
        if target_slide:
            break
    if target_slide is None:
        target_slide = prs.slides[min(11, len(prs.slides) - 1)]

    # Insert the image at the requested coordinates.  The dimensions
    # 24.2 cm wide and 6.52 cm high and position (0.93 cm, 5.46 cm)
    # mirror those used on the NIKKEI slide.  These values leave room
    # above for a separate legend, which can be added later.  Add the
    # picture and bring it to the front so that it is not obscured by
    # other shapes (e.g. a placeholder gauge).
    left = Cm(0.93)
    top = Cm(5.46)
    width = Cm(24.2)
    height = Cm(6.52)
    stream = BytesIO(img_bytes)
    picture = target_slide.shapes.add_picture(stream, left, top, width=width, height=height)
    try:
        sp_tree = target_slide.shapes._spTree
        # Remove and reinsert near the front (after background)
        sp_tree.remove(picture._element)
        sp_tree.insert(1, picture._element)
    except Exception:
        # Fallback: leave the picture at the end of the shape list
        pass

    # Replace the last‑price placeholder on the NIKKEI slide.  Compute the
    # most recent price and format it with two decimal places; fall back
    # to 'N/A' if unavailable.  The placeholder may be a shape named
    # ``last_price_nikkei`` or text containing ``[last_price_nikkei]`` or
    # ``last_price_nikkei``.  Font attributes are preserved.
    last_price = None
    if df_full is not None and not df_full.empty:
        try:
            last_price = float(df_full["Price"].iloc[-1])
        except Exception:
            last_price = None
    last_str = f"(last: {last_price:,.2f})" if last_price is not None else "(last: N/A)"
    placeholder_name = "last_price_nikkei"
    placeholder_patterns = ["[last_price_nikkei]", "last_price_nikkei"]
    replaced = False
    for shp in target_slide.shapes:
        # Match by shape name
        if getattr(shp, "name", "").lower() == placeholder_name:
            if shp.has_text_frame:
                runs = shp.text_frame.paragraphs[0].runs
                attrs = _get_run_font_attributes(runs[0]) if runs else (
                    None,
                    None,
                    None,
                    None,
                    None,
                    None,
                )
                shp.text_frame.clear()
                p = shp.text_frame.paragraphs[0]
                new_run = p.add_run()
                new_run.text = last_str
                _apply_run_font_attributes(new_run, *attrs)
            replaced = True
            break
        # Match placeholder patterns within the text
        if shp.has_text_frame:
            original_text = shp.text or ""
            for pattern in placeholder_patterns:
                if pattern in original_text:
                    runs = shp.text_frame.paragraphs[0].runs
                    attrs = _get_run_font_attributes(runs[0]) if runs else (
                        None,
                        None,
                        None,
                        None,
                        None,
                        None,
                    )
                    new_text = original_text.replace(pattern, last_str)
                    shp.text_frame.clear()
                    p = shp.text_frame.paragraphs[0]
                    new_run = p.add_run()
                    new_run.text = new_text
                    _apply_run_font_attributes(new_run, *attrs)
                    replaced = True
                    break
        if replaced:
            break
    return prs


def _get_nikkei_momentum_score(excel_obj_or_path) -> Optional[float]:
    """
    Retrieve the momentum score for NIKKEI.
    Uses common helper with instrument-specific ticker.
    """
    return _get_momentum_score_generic(excel_obj_or_path, "NKY INDEX")



def insert_nikkei_momentum_score_number(prs: Presentation, excel_file) -> Presentation:
    """Insert the NIKKEI momentum score into the slide."""
    score = _get_nikkei_momentum_score(excel_file)
    return insert_score_number(prs, score, "nikkei", "momentum_score")


###############################################################################
# Chart insertion
###############################################################################

def insert_nikkei_technical_chart(
    prs: Presentation,
    excel_file,
    anchor_date: Optional[pd.Timestamp] = None,
    price_mode: str = "Last Price",
) -> Presentation:
    """
    Insert the NIKKEI technical‑analysis chart into the PPT.

    We only use the textbox named ``nikkei`` (or containing “[nikkei]”) to locate
    the correct slide; the chart itself is always pasted at the fixed
    coordinates (0.93 cm left, 4.39 cm top, 21.41 cm wide, 7.53 cm high).
    """
    # Load data and generate image
    try:
        df_full = _load_price_data_from_obj(excel_file, "NKY Index", price_mode=price_mode)
    except Exception:
        df_full = _load_price_data(pathlib.Path(excel_file), "NKY Index", price_mode=price_mode)
    img_bytes = _generate_nikkei_image_from_df(df_full, anchor_date)

    # Find the slide containing the 'nikkei' placeholder
    target_slide = None
    for slide in prs.slides:
        for shape in slide.shapes:
            name_attr = getattr(shape, "name", "").lower()
            if name_attr == "nikkei":
                target_slide = slide
                break
            if shape.has_text_frame:
                if (shape.text or "").strip().lower() == "[nikkei]":
                    target_slide = slide
                    break
        if target_slide:
            break

    if target_slide is None:
        target_slide = prs.slides[min(11, len(prs.slides) - 1)]

    # Always paste the chart at fixed coordinates
    left = Cm(0.93)
    top = Cm(4.39)
    width = Cm(21.41)
    height = Cm(7.53)
    stream = BytesIO(img_bytes)
    target_slide.shapes.add_picture(stream, left, top, width=width, height=height)
    return prs


###############################################################################
# Subtitle insertion
###############################################################################

def insert_nikkei_subtitle(prs: Presentation, subtitle: str) -> Presentation:
    """Insert subtitle into the NIKKEI slide."""
    return insert_subtitle(prs, subtitle, "nikkei")


###############################################################################
# Colour interpolation for gauge
###############################################################################





def insert_nikkei_average_gauge(
    prs: Presentation, excel_file, last_week_avg: float
) -> Presentation:
    """
    Insert the NIKKEI average gauge into the NIKKEI slide.

    The gauge shows the average of the technical and momentum scores and
    last week's average.  It is inserted into a shape named
    ``gauge_nikkei`` within the NIKKEI slide.  If such a shape is not found
    within the NIKKEI slide, placeholders ``[GAUGE]``, ``GAUGE`` or
    ``gauge_nikkei`` on the NIKKEI slide are used instead.  If neither is
    present, the gauge is placed at a default position below the chart
    on the NIKKEI slide.  Other slides remain untouched.
    """
    tech_score = _get_nikkei_technical_score(excel_file)
    mom_score = _get_nikkei_momentum_score(excel_file)
    if tech_score is None or mom_score is None:
        return prs
    try:
        gauge_bytes = generate_average_gauge_image(
            tech_score,
            mom_score,
            last_week_avg,
            date_text="Current Week",
            last_label_text="Previous Week",
            width_cm=15.15,
            height_cm=3.13,
        )
    except Exception:
        return prs
    # Identify NIKKEI slide
    nikkei_idx = _find_nikkei_slide(prs)
    if nikkei_idx is None:
        return prs
    slide = prs.slides[nikkei_idx]
    placeholder_name = "gauge_nikkei"
    placeholder_patterns = ["[GAUGE]", "GAUGE", "gauge_nikkei"]
    # Search for named gauge placeholder first
    for shape in slide.shapes:
        if getattr(shape, "name", "").lower() == placeholder_name:
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            if shape.has_text_frame:
                shape.text = ""
            stream = BytesIO(gauge_bytes)
            slide.shapes.add_picture(stream, left, top, width=width, height=height)
            return prs
    # Then search for textual gauge placeholders on the NIKKEI slide
    for shape in slide.shapes:
        if shape.has_text_frame:
            for pattern in placeholder_patterns:
                if pattern.lower() in (shape.text or "").lower():
                    left, top, width, height = shape.left, shape.top, shape.width, shape.height
                    shape.text = shape.text.replace(pattern, "")
                    stream = BytesIO(gauge_bytes)
                    slide.shapes.add_picture(stream, left, top, width=width, height=height)
                    return prs
    # Fallback: insert below the chart within the NIKKEI slide using template coordinates
    left = Cm(8.97)
    top = Cm(12.13)
    width = Cm(15.15)
    height = Cm(3.13)
    stream = BytesIO(gauge_bytes)
    slide.shapes.add_picture(stream, left, top, width=width, height=height)
    return prs


###############################################################################
# Technical assessment insertion
###############################################################################

def insert_nikkei_technical_assessment(
    prs: Presentation,
    excel_file,
    manual_desc: Optional[str] = None,
) -> Presentation:
    """
    Insert a descriptive assessment text into the NIKKEI slide.

    The assessment is written into a shape named ``nikkei_view`` on the NIKKEI
    slide.  If no such shape exists, the function replaces any
    occurrences of ``[nikkei_view]`` or ``nikkei_view`` in text on that slide.
    A manual description may be provided; if not, the function computes
    the view from the average of the technical and momentum scores.

    Other slides remain unmodified.
    """
    # Determine the description to insert
    if manual_desc is not None and isinstance(manual_desc, str):
        desc = manual_desc.strip()
        if desc and not desc.lower().startswith("s&p 500"):
            desc = f"NIKKEI: {desc}"
    else:
        tech_score = _get_nikkei_technical_score(excel_file)
        mom_score = _get_nikkei_momentum_score(excel_file)
        if tech_score is None or mom_score is None:
            return prs
        avg = (float(tech_score) + float(mom_score)) / 2.0
        if avg >= 80:
            desc = "NIKKEI: Strongly Bullish"
        elif avg >= 70:
            desc = "NIKKEI: Bullish"
        elif avg >= 60:
            desc = "NIKKEI: Slightly Bullish"
        elif avg >= 40:
            desc = "NIKKEI: Neutral"
        elif avg >= 30:
            desc = "NIKKEI: Slightly Bearish"
        elif avg >= 20:
            desc = "NIKKEI: Bearish"
        else:
            desc = "NIKKEI: Strongly Bearish"

    target_name = "nikkei_view"
    placeholder_patterns = ["[nikkei_view]", "nikkei_view"]

    nikkei_idx = _find_nikkei_slide(prs)
    if nikkei_idx is None:
        return prs
    slide = prs.slides[nikkei_idx]
    # Try to locate a shape by name on the NIKKEI slide
    for shape in slide.shapes:
        name_attr = getattr(shape, "name", "")
        if name_attr and name_attr.lower() == target_name:
            if shape.has_text_frame:
                runs = shape.text_frame.paragraphs[0].runs
                attrs = _get_run_font_attributes(runs[0]) if runs else (None, None, None, None, None, None)
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                new_run = p.add_run()
                new_run.text = desc
                _apply_run_font_attributes(new_run, *attrs)
            return prs
    # Otherwise, replace placeholder patterns on the NIKKEI slide
    for shape in slide.shapes:
        if shape.has_text_frame:
            for pattern in placeholder_patterns:
                if pattern.lower() in (shape.text or "").lower():
                    runs = shape.text_frame.paragraphs[0].runs
                    attrs = _get_run_font_attributes(runs[0]) if runs else (None, None, None, None, None, None)
                    try:
                        new_text = shape.text.replace(pattern, desc)
                    except Exception:
                        new_text = desc
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    new_run = p.add_run()
                    new_run.text = new_text
                    _apply_run_font_attributes(new_run, *attrs)
                    return prs
    return prs


###############################################################################
# Source footnote insertion
###############################################################################

def insert_nikkei_source(
    prs: Presentation,
    used_date: Optional[pd.Timestamp],
    price_mode: str,
) -> Presentation:
    """
    Insert the source footnote into a shape named 'nikkei_source' (or
    containing '[nikkei_source]').  The footnote text depends on the selected
    price mode.  For example:

      * Last Close  → "Source: Bloomberg, Herculis Group, Data as of 29/07/2025 Close"
      * Last Price  → "Source: Bloomberg, Herculis Group, Data as of 29/07/2025"

    Parameters
    ----------
    prs : Presentation
        The PowerPoint presentation to modify.
    used_date : pandas.Timestamp or None
        The date that should appear in the footnote.  If ``None``, no
        changes are made.
    price_mode : str
        Either 'Last Price' or 'Last Close'.  Determines whether the
        suffix " Close" is appended to the date.

    Returns
    -------
    Presentation
        The modified presentation.
    """
    if used_date is None:
        return prs
    try:
        date_str = used_date.strftime("%d/%m/%Y")
    except Exception:
        return prs
    suffix = " Close" if str(price_mode).lower() == "last close" else ""
    source_text = f"Source: Bloomberg, Herculis Group, Data as of {date_str}{suffix}"
    placeholder_name = "nikkei_source"
    placeholder_patterns = ["[nikkei_source]", "nikkei_source"]
    # Restrict insertion to the NIKKEI slide only
    nikkei_idx = _find_nikkei_slide(prs)
    if nikkei_idx is None:
        return prs
    slide = prs.slides[nikkei_idx]
    # Case 1: replace a shape named exactly as the placeholder
    for shape in slide.shapes:
        name_attr = getattr(shape, "name", "")
        if name_attr and name_attr.lower() == placeholder_name:
            if shape.has_text_frame:
                runs = shape.text_frame.paragraphs[0].runs
                attrs = _get_run_font_attributes(runs[0]) if runs else (None, None, None, None, None, None)
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                new_run = p.add_run()
                new_run.text = source_text
                _apply_run_font_attributes(new_run, *attrs)
            return prs
    # Case 2: replace occurrences of the placeholder pattern in text on the NIKKEI slide
    for shape in slide.shapes:
        if shape.has_text_frame:
            for pattern in placeholder_patterns:
                if pattern.lower() in (shape.text or "").lower():
                    runs = shape.text_frame.paragraphs[0].runs
                    attrs = _get_run_font_attributes(runs[0]) if runs else (None, None, None, None, None, None)
                    try:
                        new_text = (shape.text or "").replace(pattern, source_text)
                    except Exception:
                        new_text = source_text
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    new_run = p.add_run()
                    new_run.text = new_text
                    _apply_run_font_attributes(new_run, *attrs)
                    return prs
    return prs


###############################################################################
# Range gauge helpers and insertion
###############################################################################





def insert_nikkei_technical_chart_with_range(
    prs: Presentation,
    excel_file,
    anchor_date: Optional[pd.Timestamp] = None,
    lookback_days: int = 90,
    price_mode: str = "Last Price",
) -> Presentation:
    """
    Insert the NIKKEI technical analysis chart with the vertical range gauge into the PPT.

    This function behaves similarly to ``insert_nikkei_technical_chart`` but uses
    ``generate_range_gauge_chart_image`` to draw a combined chart and gauge.
    It attempts to find a shape named 'nikkei' or containing '[nikkei]' to locate the
    slide for insertion.  The image is placed at fixed coordinates matching the
    original template (0.93 cm left, 4.39 cm top, 21.41 cm wide, 7.53 cm high).

    Parameters
    ----------
    prs : Presentation
        The PowerPoint presentation into which the chart should be inserted.
    excel_file : file‑like object or path
        Excel workbook containing NIKKEI price data.
    anchor_date : pandas.Timestamp or None, optional
        Optional anchor date for the regression channel.
    lookback_days : int, default 90
        Lookback window for computing the high and low range bounds.

    Returns
    -------
    Presentation
        The presentation with the updated slide.
    """
    # Load data
    try:
        df_full = _load_price_data_from_obj(excel_file, "NKY Index", price_mode=price_mode)
    except Exception:
        df_full = _load_price_data(pathlib.Path(excel_file), "NKY Index", price_mode=price_mode)
    # Determine the implied volatility index value (VIX) from the Excel file
    # so that the expected one‑week trading range can be estimated.  If the
    # volatility index cannot be read, ``None`` is returned and the range
    # will fall back to an ATR‑based estimate.
    vol_val = _get_vol_index_value(excel_file, price_mode=price_mode, vol_ticker="VIX Index")
    img_bytes = generate_range_gauge_chart_image(
        df_full,
        anchor_date=anchor_date,
        lookback_days=lookback_days,
        vol_index_value=vol_val,
    )

    # Locate target slide
    target_slide = None
    for slide in prs.slides:
        for shape in slide.shapes:
            name_attr = getattr(shape, "name", "").lower()
            if name_attr == "nikkei":
                target_slide = slide
                break
            if shape.has_text_frame:
                if (shape.text or "").strip().lower() == "[nikkei]":
                    target_slide = slide
                    break
        if target_slide:
            break
    if target_slide is None:
        target_slide = prs.slides[min(11, len(prs.slides) - 1)]

    # Position and dimensions tailored to the original placeholder size.
    # The NIKKEI slide in the template allocates ~21.41 cm for the chart area
    # and reserves the remaining width for the chart title, subtitle and
    # margins.  We therefore insert the combined chart‑and‑gauge image
    # using the original dimensions (21.41 cm × 7.53 cm) and rely on the
    # gauge function to include the gauge within that width.  This avoids
    # cropping the chart when the image is inserted into the slide.
    left = Cm(0.93)
    top = Cm(4.40)
    width = Cm(21.41)
    height = Cm(7.53)
    stream = BytesIO(img_bytes)
    target_slide.shapes.add_picture(stream, left, top, width=width, height=height)
    return prs