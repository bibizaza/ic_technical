"""
Utility functions for S&P 500 technical analysis and high‑resolution export.

This module provides tools to build interactive and static charts for the
S&P 500 index, calculate and insert technical and momentum scores into
PowerPoint presentations, generate horizontal and vertical gauges that
visualise the average of the technical and momentum scores, as well as
contextual trading ranges (higher and lower range bounds).  Functions
fall back to sensible defaults when placeholders are not found.

Key functions include:

* ``make_spx_figure`` – interactive Plotly chart for Streamlit.
* ``insert_spx_technical_chart`` – insert a static SPX chart into a PPTX.
* ``insert_spx_technical_score_number`` – insert the technical score (integer).
* ``insert_spx_momentum_score_number`` – insert the momentum score (integer).
* ``insert_spx_subtitle`` – insert a user‑defined subtitle into the SPX slide.
* ``generate_average_gauge_image`` – create a horizontal gauge image.
* ``insert_spx_average_gauge`` – insert the gauge into a PPT slide.
* ``insert_spx_technical_assessment`` – insert a descriptive “view” text.
* ``generate_range_gauge_chart_image`` – create a combined price chart with
  a vertical range gauge on the right hand side, including a horizontal line
  connecting the last price to the gauge.  This function is used by
  ``insert_spx_technical_chart_with_range``.
* ``insert_spx_technical_chart_with_range`` – insert the SPX technical
  analysis chart with the higher/lower range gauge into the PPT.

The range gauge illustrates the recent trading range for the S&P 500.
Instead of using the absolute high and low closes of the last 90 days,
the bounds are estimated from recent volatility.  Whenever possible the
code looks up the forward‑looking volatility index (VIX) and computes a
1‑week expected move as ``(current_price × (VIX / 100)) / sqrt(52)``.
The upper and lower bounds are the current price plus and minus that
expected move.  If the volatility index is unavailable, the code falls
back to using realised volatility: it computes the standard deviation of
30‑day daily returns, annualises it and converts it to a 1‑week move
using the same formula.  As a last resort when neither volatility
measure can be computed the bounds default to ±2 % of the current price.
A minimum ±1 % band is enforced to avoid overlapping annotations when
volatility is extremely low.  A horizontal line continues the last
price through to the gauge so that viewers can quickly assess how far
the index lies from its typical volatility band.  This method
provides context on potential upside and downside moves based on
recent price variability rather than simply the most recent highs and
lows.
"""

from __future__ import annotations

# === Windows Playwright Fix - MUST BE EARLY ===
import sys
import asyncio
if sys.platform == 'win32':
    asyncio.set_event_loop_policy(asyncio.WindowsProactorEventLoopPolicy())

try:
    import nest_asyncio
    nest_asyncio.apply()
except ImportError:
    pass
# === End Fix ===

from datetime import timedelta
import pathlib
from typing import Optional, Tuple

import pandas as pd
import plotly.graph_objects as go
from sklearn.linear_model import LinearRegression
import streamlit as st

from pptx import Presentation
from pptx.util import Cm
from io import BytesIO
import matplotlib.pyplot as plt
import matplotlib.dates as mdates
import matplotlib.patches as patches
from matplotlib.colors import LinearSegmentedColormap
import numpy as np

# Import common helpers (eliminates code duplication)
from technical_analysis.common_helpers import (
    _get_run_font_attributes,
    _apply_run_font_attributes,
    _add_mas,
    _get_technical_score_generic,
    _get_momentum_score_generic,
    _interpolate_color,
    _load_price_data_from_obj,
    _load_price_data_generic,
    _compute_range_bounds,
    generate_average_gauge_image,
    generate_range_gauge_only_image,
    generate_range_gauge_chart_image,
    generate_range_callout_chart_image,
)
from technical_analysis.powerpoint_utils import (
    find_slide_by_placeholder,
    insert_score_number,
    insert_chart_image,
    insert_subtitle,
    insert_technical_assessment,
    insert_source,
)

# Import MARS momentum scoring engine (no longer used for calculation - kept for reference)
# from mars_engine import (
#     generate_spx_score_history,
#     load_prices_for_mars,
# )
# MARS scores are now read directly from mars_score sheet in Excel (see _compute_spx_mars_score_cached)

def _load_price_data(
    excel_path: pathlib.Path,
    ticker: str = "SPX Index",
    price_mode: str = "Last Price",
) -> pd.DataFrame:
    """
    Read the raw price sheet and return a tidy Date‑Price DataFrame.

    This is a wrapper around _load_price_data_generic with the instrument-specific
    default ticker.

    Parameters
    ----------
    excel_path : pathlib.Path
        Path to the Excel workbook containing price data.
    ticker : str, default "SPX Index"
        Column name corresponding to the desired ticker in the Excel sheet.
    price_mode : str, default "Last Price"
        One of "Last Price" or "Last Close".  If ``adjust_prices_for_mode``
        is available and the mode is "Last Close", rows with the last
        recorded date (if equal to today's date) will be dropped.

    Returns
    -------
    pandas.DataFrame
        DataFrame with columns ``Date`` and ``Price``.  The data are
        sorted by date and any rows with missing values are removed.
    """
    return _load_price_data_generic(excel_path, ticker, price_mode)

# ---------------------------------------------------------------------------

# ---------------------------------------------------------------------------
# Configuration
#
# ``PLOT_LOOKBACK_DAYS`` controls the default window of historical data used
# when drawing charts.  It is expressed in trading days and set by default
# to 90 (approximately 3 months) to match the desired timeframe.  The
# Streamlit application may override this value at runtime by assigning
# a new integer to this module-level constant before invoking any chart
# functions.  When computing moving averages for the static charts, the
# rolling windows are always computed on the full price history and then
# cropped to the lookback window to avoid artificially shortening longer
# moving averages when displaying a shorter slice of data.
PLOT_LOOKBACK_DAYS: int = 90

# Import helper for adjusting price data according to price mode.  The utils
# module must reside at the project root.  It is deliberately not imported
# from ``technical_analysis.utils`` so that this module can be reused when
# nested within ``ic`` and ``utils.py`` sits at the project root.
try:
    from utils import adjust_prices_for_mode  # type: ignore
except Exception:
    # If the utils module is not available, define a no-op fallback.  This
    # preserves compatibility with environments where price mode is not used.
    adjust_prices_for_mode = None  # type: ignore

###############################################################################
# Internal helpers
###############################################################################

def _get_vol_index_value(
    excel_obj_or_path,
    price_mode: str = "Last Price",
    vol_ticker: str = "VIX Index",
) -> Optional[float]:
    """
    Retrieve the most recent value of a volatility index (e.g. VIX) from
    the ``data_prices`` sheet.  If ``price_mode`` is ``"Last Close"``,
    the most recent date is dropped if it matches today's date.  The
    returned value is the last available entry after price‑mode adjustment.

    Parameters
    ----------
    excel_obj_or_path : file‑like or path
        The Excel workbook containing price data.
    price_mode : str, default "Last Price"
        One of "Last Price" or "Last Close".  When set to "Last Close"
        rows corresponding to the most recent date (if equal to today's
        date) will be excluded before taking the last value.
    vol_ticker : str, default "VIX Index"
        Column name in the ``data_prices`` sheet corresponding to the
        volatility index whose level should be used.

    Returns
    -------
    float or None
        The most recent volatility index value, or ``None`` if it cannot
        be read or converted to a float.
    """
    try:
        df = pd.read_excel(excel_obj_or_path, sheet_name="data_prices")
    except Exception:
        return None

# ---------------------------------------------------------------------------
# Momentum score helpers
#
def make_spx_figure(
    excel_path: str | pathlib.Path,
    anchor_date: Optional[pd.Timestamp] = None,
    price_mode: str = "Last Price",
) -> go.Figure:
    """
    Build an interactive SPX chart for Streamlit.

    Parameters
    ----------
    excel_path : str or pathlib.Path
        Path to the Excel file containing SPX price data.
    anchor_date : pandas.Timestamp or None, optional
        If provided, a regression channel is drawn from ``anchor_date`` to the
        latest date.
    price_mode : str, default "Last Price"
        One of "Last Price" or "Last Close".  When set to "Last Close"
        and if the ``utils.adjust_prices_for_mode`` helper is available,
        rows corresponding to the most recent date (if equal to today's
        date) are excluded from the chart.  This allows the chart to
        represent the prior day's closing prices rather than an
        intraday or current price.

    Returns
    -------
    go.Figure
        A Plotly figure with price, moving averages, Fibonacci lines and
        an optional regression channel.
    """
    excel_path = pathlib.Path(excel_path)
    # Load data and adjust according to the price mode
    df_raw = _load_price_data(excel_path, "SPX Index", price_mode=price_mode)
    df_full = _add_mas(df_raw)

    if df_full.empty:
        return go.Figure()

    today = df_full["Date"].max().normalize()
    # Restrict the chart to the configured lookback window (e.g. last 90 days).
    # ``PLOT_LOOKBACK_DAYS`` defaults to 90 but may be overridden at runtime.
    start = today - timedelta(days=PLOT_LOOKBACK_DAYS)
    df = df_full[df_full["Date"].between(start, today)].reset_index(drop=True)

    if df.empty:
        return go.Figure()

    last_price = df["Price"].iloc[-1]
    last_price_str = f"{last_price:,.2f}"

    fig = go.Figure()
    fig.add_trace(
        go.Scatter(
            x=df["Date"],
            y=df["Price"],
            mode="lines",
            name=f"S&P 500 Price (last: {last_price_str})",
            line=dict(color="#153D64", width=2.5),
        )
    )
    fig.add_trace(
        go.Scatter(
            x=df["Date"],
            y=df["MA_50"],
            mode="lines",
            name="50‑day MA",
            line=dict(color="#008000", width=1.5),
        )
    )
    fig.add_trace(
        go.Scatter(
            x=df["Date"],
            y=df["MA_100"],
            mode="lines",
            name="100‑day MA",
            line=dict(color="#FFA500", width=1.5),
        )
    )
    fig.add_trace(
        go.Scatter(
            x=df["Date"],
            y=df["MA_200"],
            mode="lines",
            name="200‑day MA",
            line=dict(color="#FF0000", width=1.5),
        )
    )

    hi, lo = df["Price"].max(), df["Price"].min()
    span = hi - lo
    for lvl in [
        hi,
        hi - 0.236 * span,
        hi - 0.382 * span,
        hi - 0.5 * span,
        hi - 0.618 * span,
        lo,
    ]:
        fig.add_hline(
            y=lvl, line=dict(color="grey", dash="dash", width=1), opacity=0.6
        )

    if anchor_date is not None:
        # Limit regression channel to the same date range as the chart (start to today)
        # This prevents the X-axis from extending beyond PLOT_LOOKBACK_DAYS
        channel_start = max(anchor_date, start)  # Don't go before chart start
        per = df_full[df_full["Date"].between(channel_start, today)].copy()
        if not per.empty:
            X = per["Date"].map(pd.Timestamp.toordinal).to_numpy().reshape(-1, 1)
            y_vals = per["Price"].to_numpy()
            model = LinearRegression().fit(X, y_vals)
            trend = model.predict(X)
            resid = y_vals - trend
            upper = trend + resid.max()
            lower = trend + resid.min()

            uptrend = model.coef_[0] > 0
            lineclr = "green" if uptrend else "red"
            fillclr = "rgba(0,150,0,0.25)" if uptrend else "rgba(200,0,0,0.25)"

            fig.add_trace(
                go.Scatter(
                    x=per["Date"],
                    y=upper,
                    mode="lines",
                    line=dict(color=lineclr, dash="dash"),
                    showlegend=False,
                )
            )
            fig.add_trace(
                go.Scatter(
                    x=per["Date"],
                    y=lower,
                    mode="lines",
                    line=dict(color=lineclr, dash="dash"),
                    fill="tonexty",
                    fillcolor=fillclr,
                    showlegend=False,
                )
            )

    fig.update_layout(
        margin=dict(l=30, r=30, t=60, b=40),
        showlegend=True,
        legend=dict(
            orientation="h",
            yanchor="bottom",
            y=1.12,
            xanchor="center",
            x=0.5,
            font=dict(size=12),
        ),
        xaxis_title=None,
        yaxis_title=None,
        xaxis=dict(showgrid=False),
        yaxis=dict(showgrid=False, zeroline=False),
    )
    return fig

###############################################################################
# High‑resolution chart export (PNG)
###############################################################################

def _generate_spx_image_from_df(
    df_full: pd.DataFrame,
    anchor_date: Optional[pd.Timestamp],
    width_cm: float = 21.41,
    height_cm: float = 7.53,
) -> bytes:
    """
    Create a high‑resolution (dpi=300) transparent PNG chart from the DataFrame.
    Includes price, moving averages, Fibonacci lines and optional regression channel.
    """
    # Set font to Calibri for all chart text
    plt.rcParams['font.family'] = 'Calibri'
    plt.rcParams['font.sans-serif'] = ['Calibri']

    today = df_full["Date"].max().normalize()
    # Compute the lookback start based on the configurable window
    start = today - timedelta(days=PLOT_LOOKBACK_DAYS)
    # Slice the price history to the configured lookback window
    df = df_full[df_full["Date"].between(start, today)].reset_index(drop=True)

    # Compute moving averages on the full dataset once, then slice to the
    # same window.  This prevents shorter lookback windows from truncating
    # the rolling windows for the 50-, 100- and 200-day moving averages.
    df_ma_full = _add_mas(df_full)
    df_ma = df_ma_full[df_ma_full["Date"].between(start, today)].reset_index(drop=True)

    uptrend = False
    upper = lower = None
    if anchor_date is not None:
        subset = df_full[df_full["Date"].between(anchor_date, today)].copy()
        if not subset.empty:
            X = subset["Date"].map(pd.Timestamp.toordinal).to_numpy().reshape(-1, 1)
            y_vals = subset["Price"].to_numpy()
            model = LinearRegression().fit(X, y_vals)
            trend = model.predict(X)
            resid = y_vals - trend
            uptrend = model.coef_[0] > 0
            upper = trend + resid.max()
            lower = trend + resid.min()

    last_price = df["Price"].iloc[-1]
    last_price_str = f"{last_price:,.2f}"

    fig_width_in, fig_height_in = width_cm / 2.54, height_cm / 2.54
    plt.style.use("default")
    fig, ax = plt.subplots(figsize=(fig_width_in, fig_height_in))

    ax.plot(
        df["Date"],
        df["Price"],
        color="#153D64",
        linewidth=2.5,
        label=f"S&P 500 Price (last: {last_price_str})",
    )
    ax.plot(
        df_ma["Date"],
        df_ma["MA_50"],
        color="#008000",
        linewidth=1.5,
        label="50‑day MA",
    )
    ax.plot(
        df_ma["Date"],
        df_ma["MA_100"],
        color="#FFA500",
        linewidth=1.5,
        label="100‑day MA",
    )
    ax.plot(
        df_ma["Date"],
        df_ma["MA_200"],
        color="#FF0000",
        linewidth=1.5,
        label="200‑day MA",
    )

    hi, lo = df["Price"].max(), df["Price"].min()
    span = hi - lo
    fib_levels = [
        hi,
        hi - 0.236 * span,
        hi - 0.382 * span,
        hi - 0.5 * span,
        hi - 0.618 * span,
        lo,
    ]
    for lvl in fib_levels:
        ax.axhline(y=lvl, color="grey", linestyle="--", linewidth=0.8, alpha=0.6)

    if anchor_date is not None and upper is not None and lower is not None:
        fill_color = (0, 0.6, 0, 0.25) if uptrend else (0.78, 0, 0, 0.25)
        line_color = "#008000" if uptrend else "#C00000"
        subset = (
            df_full[df_full["Date"].between(anchor_date, today)].copy().reset_index(drop=True)
        )
        ax.plot(subset["Date"], upper, color=line_color, linestyle="--")
        ax.plot(subset["Date"], lower, color=line_color, linestyle="--")
        ax.fill_between(subset["Date"], lower, upper, color=fill_color)

    for spine in ax.spines.values():
        spine.set_visible(False)
    ax.tick_params(left=True, bottom=True, labelleft=True, labelbottom=True)

    # Format x-axis dates as "Aug-01" to save space
    ax.xaxis.set_major_formatter(mdates.DateFormatter('%b-%d'))
    fig.autofmt_xdate()  # Rotate date labels for better readability

    ax.legend(
        loc="upper center", bbox_to_anchor=(0.5, 1.1), ncol=4, fontsize=8, frameon=False
    )
    plt.tight_layout()

    buf = BytesIO()
    plt.savefig(buf, format="png", dpi=600, transparent=True)
    plt.close(fig)
    buf.seek(0)
    return buf.getvalue()

###############################################################################
# Score helpers
###############################################################################

def _get_spx_technical_score(excel_obj_or_path) -> Optional[float]:
    """
    Retrieve the technical score for S&P 500.
    Uses common helper with SPX-specific ticker.
    """
    return _get_technical_score_generic(excel_obj_or_path, "SPX INDEX")

def _find_spx_slide(prs: Presentation) -> Optional[int]:
    """Find the S&P 500 slide by placeholder."""
    return find_slide_by_placeholder(prs, "spx")

def insert_spx_technical_score_number(prs: Presentation, excel_file) -> Presentation:
    """Insert the S&P 500 technical score into the slide."""
    score = _get_spx_technical_score(excel_file)
    return insert_score_number(prs, score, "spx", "tech_score")

###############################################################################
# Call‑out range helpers and insertion
###############################################################################

def insert_spx_technical_chart_with_callout(
    prs: Presentation,
    excel_file,
    anchor_date: Optional[pd.Timestamp] = None,
    lookback_days: int = 90,
    price_mode: str = "Last Price",
) -> Presentation:
    """
    Insert the SPX technical analysis chart with the trading range call‑out
    into the PowerPoint.  This function mirrors the behaviour of
    ``insert_spx_technical_chart_with_range`` but uses the call‑out style to
    display the high and low bounds instead of a vertical gauge.

    The image is placed at the fixed coordinates (0.93 cm left, 5.46 cm top)
    with dimensions 24.2 cm wide by 6.52 cm high.  These values match those
    used on the IBOV slide and leave room above for a separate legend on the
    PowerPoint slide.  When inserting into the presentation the legend is
    suppressed in the image itself so that it can be added manually.

    Parameters
    ----------
    prs : Presentation
        The PowerPoint presentation to modify.
    excel_file : file‑like object or path
        Excel workbook containing SPX price data.
    anchor_date : pandas.Timestamp or None, optional
        Optional anchor date for a regression channel.
    lookback_days : int, default 90
        Lookback window for computing the high/low range.

    Returns
    -------
    Presentation
        The presentation with the updated slide.
    """
    # Load the price data from the Excel file
    try:
        df_full = _load_price_data_from_obj(excel_file, "SPX Index", price_mode=price_mode)
    except Exception:
        df_full = _load_price_data(pathlib.Path(excel_file), "SPX Index", price_mode=price_mode)

    # Determine the implied volatility index value (VIX) from the Excel file
    # so that the expected one‑week trading range can be estimated.  If the
    # volatility index cannot be read, ``None`` is returned and the range
    # will fall back to an ATR‑based estimate.
    vol_val = _get_vol_index_value(excel_file, price_mode=price_mode, vol_ticker="VIX Index")
    # Generate the image with the call‑out.  Use a width of 24.2 cm and a
    # height of 6.52 cm (matching the IBOV template) so that there is
    # sufficient space above the chart for an external legend.  Pass
    # ``show_legend=False`` to suppress the internal legend on the figure and
    # provide the volatility index value so that the range calculation can
    # use the implied volatility if available.
    img_bytes = generate_range_callout_chart_image(
        df_full,
        anchor_date=anchor_date,
        lookback_days=lookback_days,
        width_cm=24.2,
        height_cm=6.52,
        vol_index_value=vol_val,
        show_legend=False,
        cache_key="spx_main_callout",  # Phase 2: Use cached chart if available
    )

    # Skip chart insertion if no image was generated (empty data)
    if not img_bytes:
        print("Warning: SPX chart image could not be generated (empty data)")
        return prs

    # Locate the slide containing the 'spx' placeholder or text
    target_slide = None
    for slide in prs.slides:
        for shape in slide.shapes:
            name_attr = getattr(shape, "name", "").lower()
            if name_attr == "spx":
                target_slide = slide
                break
            if shape.has_text_frame:
                if (shape.text or "").strip().lower() == "[spx]":
                    target_slide = slide
                    break
        if target_slide:
            break
    if target_slide is None:
        target_slide = prs.slides[min(11, len(prs.slides) - 1)]

    # Insert the image at the requested coordinates.  The dimensions
    # 24.2 cm wide and 6.52 cm high and position (0.93 cm, 5.46 cm)
    # mirror those used on the IBOV slide.  These values leave room
    # above for a separate legend, which can be added later.  Add the
    # picture and bring it to the front so that it is not obscured by
    # other shapes (e.g. a placeholder gauge).
    left = Cm(0.93)
    top = Cm(5.46)
    width = Cm(24.2)
    height = Cm(6.52)
    stream = BytesIO(img_bytes)
    picture = target_slide.shapes.add_picture(stream, left, top, width=width, height=height)
    try:
        sp_tree = target_slide.shapes._spTree
        # Remove and reinsert near the front (after background)
        sp_tree.remove(picture._element)
        sp_tree.insert(1, picture._element)
    except Exception:
        # Fallback: leave the picture at the end of the shape list
        pass

    # Replace the last‑price placeholder on the SPX slide.  Compute the
    # most recent price and format it with two decimal places; fall back
    # to 'N/A' if unavailable.  The placeholder may be a shape named
    # ``last_price_spx`` or text containing ``[last_price_spx]`` or
    # ``last_price_spx``.  Font attributes are preserved.
    last_price = None
    if df_full is not None and not df_full.empty:
        try:
            last_price = float(df_full["Price"].iloc[-1])
        except Exception:
            last_price = None
    last_str = f"(last: {last_price:,.2f})" if last_price is not None else "(last: N/A)"
    placeholder_name = "last_price_spx"
    placeholder_patterns = ["[last_price_spx]", "last_price_spx"]
    replaced = False
    for shp in target_slide.shapes:
        # Match by shape name
        if getattr(shp, "name", "").lower() == placeholder_name:
            if shp.has_text_frame:
                runs = shp.text_frame.paragraphs[0].runs
                attrs = _get_run_font_attributes(runs[0]) if runs else (
                    None,
                    None,
                    None,
                    None,
                    None,
                    None,
                )
                shp.text_frame.clear()
                p = shp.text_frame.paragraphs[0]
                new_run = p.add_run()
                new_run.text = last_str
                _apply_run_font_attributes(new_run, *attrs)
            replaced = True
            break
        # Match placeholder patterns within the text
        if shp.has_text_frame:
            original_text = shp.text or ""
            for pattern in placeholder_patterns:
                if pattern in original_text:
                    runs = shp.text_frame.paragraphs[0].runs
                    attrs = _get_run_font_attributes(runs[0]) if runs else (
                        None,
                        None,
                        None,
                        None,
                        None,
                        None,
                    )
                    new_text = original_text.replace(pattern, last_str)
                    shp.text_frame.clear()
                    p = shp.text_frame.paragraphs[0]
                    new_run = p.add_run()
                    new_run.text = new_text
                    _apply_run_font_attributes(new_run, *attrs)
                    replaced = True
                    break
        if replaced:
            break
    return prs

@st.cache_data(show_spinner=False)
def _compute_spx_mars_score_cached(excel_path: str) -> Optional[float]:
    """
    Read pre-computed MARS momentum score for SPX from the mars_score sheet.

    The mars_score sheet should contain pre-computed scores from the standalone
    MARS application. This ensures 100% consistency with the MARS app and
    provides instant performance (no calculation needed).

    Parameters
    ----------
    excel_path : str
        Path to Excel file (must be string for caching)

    Returns
    -------
    float or None
        Latest MARS momentum score (0-100), or None if not found
    """
    try:
        from mars_engine.data_loader import load_mars_scores

        # Load all pre-computed MARS scores
        mars_scores = load_mars_scores(excel_path)

        # Look for SPX score (try multiple ticker name variations)
        for ticker in ["SPX", "SPX Index", "S&P 500"]:
            if ticker in mars_scores:
                return float(mars_scores[ticker])

        print("Warning: SPX not found in mars_score sheet")
        return None

    except Exception as e:
        print(f"Warning: Could not read MARS score for SPX: {e}")
        return None

def _get_spx_momentum_score(excel_obj_or_path) -> Optional[float]:
    """
    Retrieve the momentum score for S&P 500 from the mars_score Excel sheet.

    Uses the generic momentum score function to read pre-computed MARS scores.
    """
    return _get_momentum_score_generic(excel_obj_or_path, "SPX INDEX")

def insert_spx_momentum_score_number(prs: Presentation, excel_file) -> Presentation:
    """Insert the S&P 500 momentum score into the slide."""
    score = _get_spx_momentum_score(excel_file)
    return insert_score_number(prs, score, "spx", "momentum_score")

###############################################################################
# Chart insertion
###############################################################################

def insert_spx_technical_chart(
    prs: Presentation,
    excel_file,
    anchor_date: Optional[pd.Timestamp] = None,
    price_mode: str = "Last Price",
) -> Presentation:
    """
    Insert the SPX technical‑analysis chart into the PPT.

    We only use the textbox named ``spx`` (or containing “[spx]”) to locate
    the correct slide; the chart itself is always pasted at the fixed
    coordinates (0.93 cm left, 4.39 cm top, 21.41 cm wide, 7.53 cm high).
    """
    # Load data and generate image
    try:
        df_full = _load_price_data_from_obj(excel_file, "SPX Index", price_mode=price_mode)
    except Exception:
        df_full = _load_price_data(pathlib.Path(excel_file), "SPX Index", price_mode=price_mode)
    img_bytes = _generate_spx_image_from_df(df_full, anchor_date)

    # Find the slide containing the 'spx' placeholder
    target_slide = None
    for slide in prs.slides:
        for shape in slide.shapes:
            name_attr = getattr(shape, "name", "").lower()
            if name_attr == "spx":
                target_slide = slide
                break
            if shape.has_text_frame:
                if (shape.text or "").strip().lower() == "[spx]":
                    target_slide = slide
                    break
        if target_slide:
            break

    if target_slide is None:
        target_slide = prs.slides[min(11, len(prs.slides) - 1)]

    # Always paste the chart at fixed coordinates
    left = Cm(0.93)
    top = Cm(4.39)
    width = Cm(21.41)
    height = Cm(7.53)
    stream = BytesIO(img_bytes)
    target_slide.shapes.add_picture(stream, left, top, width=width, height=height)
    return prs

###############################################################################
# Subtitle insertion
###############################################################################

def insert_spx_subtitle(prs: Presentation, subtitle: str) -> Presentation:
    """Insert subtitle into the S&P 500 slide."""
    return insert_subtitle(prs, subtitle, "spx")

###############################################################################
# Colour interpolation for gauge
###############################################################################

def insert_spx_average_gauge(
    prs: Presentation, excel_file, last_week_avg: float
) -> Presentation:
    """
    Insert the SPX average gauge into the SPX slide.

    The gauge shows the average of the technical and momentum scores and
    last week's average.  It is inserted into a shape named
    ``gauge_spx`` within the SPX slide.  If such a shape is not found
    within the SPX slide, placeholders ``[GAUGE]``, ``GAUGE`` or
    ``gauge_spx`` on the SPX slide are used instead.  If neither is
    present, the gauge is placed at a default position below the chart
    on the SPX slide.  Other slides remain untouched.
    """
    tech_score = _get_spx_technical_score(excel_file)
    mom_score = _get_spx_momentum_score(excel_file)
    if tech_score is None or mom_score is None:
        return prs
    try:
        gauge_bytes = generate_average_gauge_image(
            tech_score,
            mom_score,
            last_week_avg,
            date_text="Current Week",
            last_label_text="Previous Week",
            width_cm=15.15,
            height_cm=3.13,
            cache_key="spx_avg_gauge",  # Phase 2: Use cached gauge if available
        )
    except Exception:
        return prs
    # Identify SPX slide
    spx_idx = _find_spx_slide(prs)
    if spx_idx is None:
        return prs
    slide = prs.slides[spx_idx]
    placeholder_name = "gauge_spx"
    placeholder_patterns = ["[GAUGE]", "GAUGE", "gauge_spx"]
    # Search for named gauge placeholder first
    for shape in slide.shapes:
        if getattr(shape, "name", "").lower() == placeholder_name:
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            if shape.has_text_frame:
                shape.text = ""
            stream = BytesIO(gauge_bytes)
            slide.shapes.add_picture(stream, left, top, width=width, height=height)
            return prs
    # Then search for textual gauge placeholders on the SPX slide
    for shape in slide.shapes:
        if shape.has_text_frame:
            for pattern in placeholder_patterns:
                if pattern.lower() in (shape.text or "").lower():
                    left, top, width, height = shape.left, shape.top, shape.width, shape.height
                    shape.text = shape.text.replace(pattern, "")
                    stream = BytesIO(gauge_bytes)
                    slide.shapes.add_picture(stream, left, top, width=width, height=height)
                    return prs
    # Fallback: insert below the chart within the SPX slide using template coordinates
    left = Cm(8.97)
    top = Cm(12.13)
    width = Cm(15.15)
    height = Cm(3.13)
    stream = BytesIO(gauge_bytes)
    slide.shapes.add_picture(stream, left, top, width=width, height=height)
    return prs

###############################################################################
# Technical assessment insertion
###############################################################################

def insert_spx_technical_assessment(
    prs: Presentation,
    excel_file,
    manual_desc: Optional[str] = None,
) -> Presentation:
    """
    Insert a descriptive assessment text into the SPX slide.

    The assessment is written into a shape named ``spx_view`` on the SPX
    slide.  If no such shape exists, the function replaces any
    occurrences of ``[spx_view]`` or ``spx_view`` in text on that slide.
    A manual description may be provided; if not, the function computes
    the view from the average of the technical and momentum scores.

    Other slides remain unmodified.
    """
    # Determine the description to insert
    if manual_desc is not None and isinstance(manual_desc, str):
        desc = manual_desc.strip()
        if desc and not desc.lower().startswith("s&p 500"):
            desc = f"S&P 500: {desc}"
    else:
        tech_score = _get_spx_technical_score(excel_file)
        mom_score = _get_spx_momentum_score(excel_file)
        if tech_score is None or mom_score is None:
            return prs
        avg = (float(tech_score) + float(mom_score)) / 2.0
        if avg >= 80:
            desc = "S&P 500: Strongly Bullish"
        elif avg >= 70:
            desc = "S&P 500: Bullish"
        elif avg >= 60:
            desc = "S&P 500: Slightly Bullish"
        elif avg >= 40:
            desc = "S&P 500: Neutral"
        elif avg >= 30:
            desc = "S&P 500: Slightly Bearish"
        elif avg >= 20:
            desc = "S&P 500: Bearish"
        else:
            desc = "S&P 500: Strongly Bearish"

    target_name = "spx_view"
    placeholder_patterns = ["[spx_view]", "spx_view"]

    spx_idx = _find_spx_slide(prs)
    if spx_idx is None:
        return prs
    slide = prs.slides[spx_idx]
    # Try to locate a shape by name on the SPX slide
    for shape in slide.shapes:
        name_attr = getattr(shape, "name", "")
        if name_attr and name_attr.lower() == target_name:
            if shape.has_text_frame:
                runs = shape.text_frame.paragraphs[0].runs
                attrs = _get_run_font_attributes(runs[0]) if runs else (None, None, None, None, None, None)
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                new_run = p.add_run()
                new_run.text = desc
                _apply_run_font_attributes(new_run, *attrs)
            return prs
    # Otherwise, replace placeholder patterns on the SPX slide
    for shape in slide.shapes:
        if shape.has_text_frame:
            for pattern in placeholder_patterns:
                if pattern.lower() in (shape.text or "").lower():
                    runs = shape.text_frame.paragraphs[0].runs
                    attrs = _get_run_font_attributes(runs[0]) if runs else (None, None, None, None, None, None)
                    try:
                        new_text = shape.text.replace(pattern, desc)
                    except Exception:
                        new_text = desc
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    new_run = p.add_run()
                    new_run.text = new_text
                    _apply_run_font_attributes(new_run, *attrs)
                    return prs
    return prs

###############################################################################
# Source footnote insertion
###############################################################################

def insert_spx_source(
    prs: Presentation,
    used_date: Optional[pd.Timestamp],
    price_mode: str,
) -> Presentation:
    """
    Insert the source footnote into a shape named 'spx_source' (or
    containing '[spx_source]').  The footnote text depends on the selected
    price mode.  For example:

      * Last Close  → "Source: Bloomberg, Herculis Group, Data as of 29/07/2025 Close"
      * Last Price  → "Source: Bloomberg, Herculis Group, Data as of 29/07/2025"

    Parameters
    ----------
    prs : Presentation
        The PowerPoint presentation to modify.
    used_date : pandas.Timestamp or None
        The date that should appear in the footnote.  If ``None``, no
        changes are made.
    price_mode : str
        Either 'Last Price' or 'Last Close'.  Determines whether the
        suffix " Close" is appended to the date.

    Returns
    -------
    Presentation
        The modified presentation.
    """
    if used_date is None:
        return prs
    try:
        date_str = used_date.strftime("%d/%m/%Y")
    except Exception:
        return prs
    suffix = " Close" if str(price_mode).lower() == "last close" else ""
    source_text = f"Source: Bloomberg, Herculis Group, Data as of {date_str}{suffix}"
    placeholder_name = "spx_source"
    placeholder_patterns = ["[spx_source]", "spx_source"]
    # Restrict insertion to the SPX slide only
    spx_idx = _find_spx_slide(prs)
    if spx_idx is None:
        return prs
    slide = prs.slides[spx_idx]
    # Case 1: replace a shape named exactly as the placeholder
    for shape in slide.shapes:
        name_attr = getattr(shape, "name", "")
        if name_attr and name_attr.lower() == placeholder_name:
            if shape.has_text_frame:
                runs = shape.text_frame.paragraphs[0].runs
                attrs = _get_run_font_attributes(runs[0]) if runs else (None, None, None, None, None, None)
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                new_run = p.add_run()
                new_run.text = source_text
                _apply_run_font_attributes(new_run, *attrs)
            return prs
    # Case 2: replace occurrences of the placeholder pattern in text on the SPX slide
    for shape in slide.shapes:
        if shape.has_text_frame:
            for pattern in placeholder_patterns:
                if pattern.lower() in (shape.text or "").lower():
                    runs = shape.text_frame.paragraphs[0].runs
                    attrs = _get_run_font_attributes(runs[0]) if runs else (None, None, None, None, None, None)
                    try:
                        new_text = (shape.text or "").replace(pattern, source_text)
                    except Exception:
                        new_text = source_text
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    new_run = p.add_run()
                    new_run.text = new_text
                    _apply_run_font_attributes(new_run, *attrs)
                    return prs
    return prs

###############################################################################
# Range gauge helpers and insertion
###############################################################################

def insert_spx_technical_chart_with_range(
    prs: Presentation,
    excel_file,
    anchor_date: Optional[pd.Timestamp] = None,
    lookback_days: int = 90,
    price_mode: str = "Last Price",
) -> Presentation:
    """
    Insert the SPX technical analysis chart with the vertical range gauge into the PPT.

    This function behaves similarly to ``insert_spx_technical_chart`` but uses
    ``generate_range_gauge_chart_image`` to draw a combined chart and gauge.
    It attempts to find a shape named 'spx' or containing '[spx]' to locate the
    slide for insertion.  The image is placed at fixed coordinates matching the
    original template (0.93 cm left, 4.39 cm top, 21.41 cm wide, 7.53 cm high).

    Parameters
    ----------
    prs : Presentation
        The PowerPoint presentation into which the chart should be inserted.
    excel_file : file‑like object or path
        Excel workbook containing SPX price data.
    anchor_date : pandas.Timestamp or None, optional
        Optional anchor date for the regression channel.
    lookback_days : int, default 90
        Lookback window for computing the high and low range bounds.

    Returns
    -------
    Presentation
        The presentation with the updated slide.
    """
    # Load data
    try:
        df_full = _load_price_data_from_obj(excel_file, "SPX Index", price_mode=price_mode)
    except Exception:
        df_full = _load_price_data(pathlib.Path(excel_file), "SPX Index", price_mode=price_mode)
    # Determine the implied volatility index value (VIX) from the Excel file
    # so that the expected one‑week trading range can be estimated.  If the
    # volatility index cannot be read, ``None`` is returned and the range
    # will fall back to an ATR‑based estimate.
    vol_val = _get_vol_index_value(excel_file, price_mode=price_mode, vol_ticker="VIX Index")
    img_bytes = generate_range_gauge_chart_image(
        df_full,
        anchor_date=anchor_date,
        lookback_days=lookback_days,
        vol_index_value=vol_val,
    )

    # Locate target slide
    target_slide = None
    for slide in prs.slides:
        for shape in slide.shapes:
            name_attr = getattr(shape, "name", "").lower()
            if name_attr == "spx":
                target_slide = slide
                break
            if shape.has_text_frame:
                if (shape.text or "").strip().lower() == "[spx]":
                    target_slide = slide
                    break
        if target_slide:
            break
    if target_slide is None:
        target_slide = prs.slides[min(11, len(prs.slides) - 1)]

    # Position and dimensions tailored to the original placeholder size.
    # The SPX slide in the template allocates ~21.41 cm for the chart area
    # and reserves the remaining width for the chart title, subtitle and
    # margins.  We therefore insert the combined chart‑and‑gauge image
    # using the original dimensions (21.41 cm × 7.53 cm) and rely on the
    # gauge function to include the gauge within that width.  This avoids
    # cropping the chart when the image is inserted into the slide.
    left = Cm(0.93)
    top = Cm(4.40)
    width = Cm(21.41)
    height = Cm(7.53)
    stream = BytesIO(img_bytes)
    target_slide.shapes.add_picture(stream, left, top, width=width, height=height)
    return prs

# =============================================================================
# TECHNICAL ANALYSIS CHART V2 - Chart.js + Playwright
# =============================================================================

import json
from jinja2 import Environment
from playwright.sync_api import sync_playwright

# Chart dimensions for v2 - HTML at base size, Playwright scales up
# Base dimensions for HTML body (smaller = sharper when scaled)
TECH_V2_BASE_WIDTH = 950
TECH_V2_BASE_HEIGHT = 420  # Reduced for better fit
TECH_V2_DEVICE_SCALE = 4   # Playwright device scale factor for high-res output
TECH_V2_HTML_SCALE = 1     # Scale factor for HTML elements (1 = base size)
TECH_V2_PNG_WIDTH_PX = TECH_V2_BASE_WIDTH * TECH_V2_DEVICE_SCALE   # 3800
TECH_V2_PNG_HEIGHT_PX = TECH_V2_BASE_HEIGHT * TECH_V2_DEVICE_SCALE  # 1680
TECH_V2_LOOKBACK_DAYS = 85  # 4 months of trading days

def _compute_rsi(prices: pd.Series, period: int = 14) -> pd.Series:
    """Compute RSI for a price series."""
    delta = prices.diff()
    gain = delta.where(delta > 0, 0.0)
    loss = (-delta).where(delta < 0, 0.0)
    avg_gain = gain.rolling(window=period, min_periods=period).mean()
    avg_loss = loss.rolling(window=period, min_periods=period).mean()
    rs = avg_gain / avg_loss
    rsi = 100 - (100 / (1 + rs))
    return rsi

def _compute_fibonacci_levels(high: float, low: float) -> list:
    """Compute Fibonacci retracement levels."""
    diff = high - low
    levels = [
        low,                         # 0%
        low + diff * 0.236,          # 23.6%
        low + diff * 0.382,          # 38.2%
        low + diff * 0.5,            # 50%
        low + diff * 0.618,          # 61.8%
        low + diff * 0.786,          # 78.6%
        high,                        # 100%
    ]
    return [round(l, 2) for l in levels]

def _get_score_status(score: int) -> tuple:
    """Get status text and color based on score."""
    if score >= 81:
        return "Strong", "#22C55E"
    elif score >= 61:
        return "Good", "#84CC16"
    elif score >= 41:
        return "Neutral", "#EAB308"
    elif score >= 21:
        return "Weak", "#F97316"
    else:
        return "Very Weak", "#EF4444"

def _get_rsi_interpretation(rsi: float) -> tuple:
    """Get RSI interpretation text and color."""
    if rsi >= 70:
        return "Overbought", "#EF4444", "Caution: Potential pullback"
    elif rsi <= 30:
        return "Oversold", "#22C55E", "Opportunity: Potential bounce"
    elif rsi >= 50:
        return "Neutral", "#EAB308", "Room to run"
    else:
        return "Neutral", "#EAB308", "Watching support"

def create_technical_analysis_v2_chart(
    excel_path,
    ticker: str = "SPX Index",
    *,
    price_mode: str = "Last Price",
    dmas_score: int = None,
    dmas_prev_week: int = None,
    technical_score: int = None,
    technical_prev_week: int = None,
    momentum_score: int = None,
    momentum_prev_week: int = None,
    lookback_days: int = TECH_V2_LOOKBACK_DAYS,
) -> tuple:
    """Generate Technical Analysis v2 chart using Chart.js + Playwright.

    Parameters
    ----------
    excel_path : str or Path
        Path to Excel file with price data.
    ticker : str
        Bloomberg ticker for the instrument.
    price_mode : str
        "Last Price" or "Last Close".
    dmas_score : int, optional
        Current DMAS score (0-100). If None, calculated from technical + momentum.
    dmas_prev_week : int, optional
        Previous week's DMAS score for change calculation.
    technical_score : int, optional
        Technical component score. If None, calculated from price data.
    technical_prev_week : int, optional
        Previous week's Technical score for trend arrow.
    momentum_score : int, optional
        Momentum component score. If None, uses placeholder.
    momentum_prev_week : int, optional
        Previous week's Momentum score for trend arrow.
    lookback_days : int
        Number of trading days for price chart.

    Returns
    -------
    tuple
        (PNG bytes, effective date used)
    """
    from market_compass.weekly_performance.html_template import TECHNICAL_ANALYSIS_V2_HTML_TEMPLATE

    # Load price data
    try:
        df = _load_price_data_from_obj(excel_path, ticker, price_mode=price_mode)
    except Exception:
        df = _load_price_data(pathlib.Path(excel_path), ticker, price_mode=price_mode)

    if df.empty:
        print(f"[Tech V2] No data for {ticker}")
        return None, None

    # Get the last N days
    df = df.tail(lookback_days + 200)  # Extra for MA calculation

    # Calculate moving averages on full data
    df["MA50"] = df["Price"].rolling(window=50, min_periods=1).mean()
    df["MA100"] = df["Price"].rolling(window=100, min_periods=1).mean()
    df["MA200"] = df["Price"].rolling(window=200, min_periods=1).mean()
    df["RSI"] = _compute_rsi(df["Price"], 14)

    # Trim to lookback window
    df = df.tail(lookback_days)

    used_date = df["Date"].max()

    # Prepare data for chart
    price_labels = df["Date"].dt.strftime("%b-%d").tolist()
    price_data = df["Price"].round(2).tolist()
    ma50_data = df["MA50"].round(2).tolist()
    ma100_data = df["MA100"].round(2).tolist()
    ma200_data = df["MA200"].round(2).tolist()
    rsi_data = df["RSI"].round(1).tolist()

    # Last price - convert to native Python float for JSON serialization
    last_price = float(df["Price"].iloc[-1])
    last_price_str = f"{last_price:,.2f}"

    # Fibonacci levels
    period_high = df["Price"].max()
    period_low = df["Price"].min()
    fib_levels = _compute_fibonacci_levels(period_high, period_low)

    # Trading range (use volatility-based if available)
    vol_val = _get_vol_index_value(excel_path, price_mode=price_mode, vol_ticker="VIX Index")
    if vol_val:
        expected_move = (last_price * (vol_val / 100)) / (52 ** 0.5)
        higher_range = float(round(last_price + expected_move, 0))
        lower_range = float(round(last_price - expected_move, 0))
    else:
        # Fallback: 2% range
        higher_range = float(round(last_price * 1.02, 0))
        lower_range = float(round(last_price * 0.98, 0))

    higher_range_pct = f"+{((higher_range / last_price - 1) * 100):.1f}%"
    lower_range_pct = f"{((lower_range / last_price - 1) * 100):.1f}%"

    # Y-axis bounds: First Fibonacci below (floor) to Higher Trading Range (ceiling)
    # This focuses the chart on relevant price action
    # Asymmetric buffers: tighter on top, more room on bottom
    Y_MAX_BUFFER_PCT = 0.005  # 0.5% buffer above Higher Trading Range
    Y_MIN_BUFFER_PCT = 0.015  # 1.5% buffer below First Fibonacci

    # Y-MAX: Higher Trading Range (volatility-based) + buffer
    price_y_max = higher_range * (1 + Y_MAX_BUFFER_PCT)

    # Y-MIN: First Fibonacci below current price - buffer
    fib_below = [f for f in fib_levels if f < last_price]
    first_fib_below = max(fib_below) if fib_below else min(fib_levels)
    price_y_min = first_fib_below * (1 - Y_MIN_BUFFER_PCT)

    # Safety: extend if actual price data exceeds bounds
    price_min = float(df["Price"].min())
    price_max = float(df["Price"].max())

    if price_min < price_y_min:
        price_y_min = price_min * 0.995  # 0.5% below min price
    if price_max > price_y_max:
        price_y_max = price_max * 1.005  # 0.5% above max price

    # Convert to native Python floats and round for clean axis labels
    price_y_min = float(round(price_y_min, 0))
    price_y_max = float(round(price_y_max, 0))

    # Debug logging
    print(f"[Tech V2] First Fib below: {first_fib_below:.2f}, Higher Range: {higher_range:.2f}")
    print(f"[Tech V2] Y-axis bounds: {price_y_min:.0f} - {price_y_max:.0f}")

    # RSI current
    rsi_current = int(round(df["RSI"].iloc[-1], 0)) if not pd.isna(df["RSI"].iloc[-1]) else 50
    rsi_interpretation, rsi_color, rsi_context = _get_rsi_interpretation(rsi_current)

    # DMAS scores - ensure all are integers (no decimals)
    if technical_score is None:
        technical_score = _get_technical_score_generic(df, "SPX")
    technical_score = int(round(technical_score)) if technical_score is not None else 50

    if momentum_score is None:
        momentum_score = 50  # Default
    momentum_score = int(round(momentum_score))

    if dmas_score is None:
        dmas_score = int(round((technical_score + momentum_score) / 2))
    else:
        dmas_score = int(round(dmas_score))

    if dmas_prev_week is None:
        dmas_prev_week = dmas_score  # No change
    else:
        dmas_prev_week = int(round(dmas_prev_week))

    dmas_change = dmas_score - dmas_prev_week
    if dmas_change > 0:
        dmas_change_text = f"▲ +{dmas_change} WoW"
        dmas_change_color = "#22C55E"
    elif dmas_change < 0:
        dmas_change_text = f"▼ {dmas_change} WoW"
        dmas_change_color = "#EF4444"
    else:
        dmas_change_text = "— Unchanged WoW"
        dmas_change_color = "#9CA3AF"

    technical_status, technical_color = _get_score_status(technical_score)
    momentum_status, momentum_color = _get_score_status(momentum_score)
    dmas_status, dmas_color = _get_score_status(dmas_score)

    # Calculate trend arrows for Technical and Momentum
    # Technical trend
    if technical_prev_week is None:
        technical_prev_week = technical_score
    else:
        technical_prev_week = int(round(technical_prev_week))

    if technical_score > technical_prev_week:
        technical_trend = "▲"
        technical_trend_color = "#22C55E"  # Green
    elif technical_score < technical_prev_week:
        technical_trend = "▼"
        technical_trend_color = "#EF4444"  # Red
    else:
        technical_trend = "—"
        technical_trend_color = "#9CA3AF"  # Gray

    # Momentum trend
    if momentum_prev_week is None:
        momentum_prev_week = momentum_score
    else:
        momentum_prev_week = int(round(momentum_prev_week))

    if momentum_score > momentum_prev_week:
        momentum_trend = "▲"
        momentum_trend_color = "#22C55E"  # Green
    elif momentum_score < momentum_prev_week:
        momentum_trend = "▼"
        momentum_trend_color = "#EF4444"  # Red
    else:
        momentum_trend = "—"
        momentum_trend_color = "#9CA3AF"  # Gray

    # Debug logging
    print(f"[Tech V2] Data points: {len(price_data)}, RSI current: {rsi_current}")
    print(f"[Tech V2] Scores - DMAS: {dmas_score}, Technical: {technical_score}, Momentum: {momentum_score}")

    # Render HTML template
    env = Environment()
    env.filters['tojson'] = json.dumps
    template = env.from_string(TECHNICAL_ANALYSIS_V2_HTML_TEMPLATE)

    html_content = template.render(
        width=TECH_V2_BASE_WIDTH,   # Base dimensions - Playwright scales up
        height=TECH_V2_BASE_HEIGHT,
        scale=TECH_V2_HTML_SCALE,   # 1 = no CSS scaling, Playwright does the scaling
        # Price chart data
        price_labels=price_labels,
        price_data=price_data,
        ma50_data=ma50_data,
        ma100_data=ma100_data,
        ma200_data=ma200_data,
        show_ma50=True,   # Always show - Y-axis clips distant MAs naturally
        show_ma100=True,
        show_ma200=True,
        fib_levels=fib_levels,
        price_y_min=price_y_min,
        price_y_max=price_y_max,
        last_price=last_price_str,
        higher_range=higher_range,
        lower_range=lower_range,
        higher_range_pct=higher_range_pct,
        lower_range_pct=lower_range_pct,
        # RSI chart data
        rsi_labels=price_labels,
        rsi_data=rsi_data,
        rsi_current=int(rsi_current),
        rsi_color=rsi_color,
        rsi_interpretation=rsi_interpretation,
        rsi_context=rsi_context,
        # DMAS panel
        dmas_score=dmas_score,
        dmas_color=dmas_color,
        dmas_change_text=dmas_change_text,
        dmas_change_color=dmas_change_color,
        technical_score=technical_score,
        technical_color=technical_color,
        technical_status=technical_status,
        technical_trend=technical_trend,
        technical_trend_color=technical_trend_color,
        momentum_score=momentum_score,
        momentum_color=momentum_color,
        momentum_status=momentum_status,
        momentum_trend=momentum_trend,
        momentum_trend_color=momentum_trend_color,
    )

    # Debug: Save HTML for inspection
    print(f"[Tech V2] HTML body: {TECH_V2_BASE_WIDTH}×{TECH_V2_BASE_HEIGHT}px, Playwright scale: {TECH_V2_DEVICE_SCALE}x -> {TECH_V2_PNG_WIDTH_PX}×{TECH_V2_PNG_HEIGHT_PX}px output")
    try:
        import tempfile
        import os
        debug_html_path = os.path.join(tempfile.gettempdir(), "tech_v2_debug.html")
        with open(debug_html_path, "w") as f:
            f.write(html_content)
        print(f"[Tech V2] Debug HTML saved to: {debug_html_path}")
    except Exception as e:
        print(f"[Tech V2] Could not save debug HTML: {e}")

    # Render with Playwright
    png_bytes = None
    try:
        with sync_playwright() as p:
            browser = p.chromium.launch(headless=True)
            # Use base dimensions for viewport, device_scale_factor for high-res output
            page = browser.new_page(
                viewport={
                    'width': TECH_V2_BASE_WIDTH,
                    'height': TECH_V2_BASE_HEIGHT
                },
                device_scale_factor=TECH_V2_DEVICE_SCALE
            )

            # Set content and wait for network idle
            page.set_content(html_content, wait_until='networkidle')

            # Wait for Chart.js to load and render
            try:
                # First wait for Chart.js to be available
                page.wait_for_function("typeof Chart !== 'undefined'", timeout=10000)
                print("[Tech V2] Chart.js library loaded")

                # Then wait for chart ready signal
                page.wait_for_selector('body[data-chart-ready="true"]', timeout=15000)
                print("[Tech V2] Chart.js rendering complete")
            except Exception as wait_err:
                print(f"[Tech V2] Wait timeout, using fallback: {wait_err}")
                # Fallback: just wait 5 seconds
                page.wait_for_timeout(5000)

            # Take screenshot
            png_bytes = page.screenshot()
            print(f"[Tech V2] Screenshot taken: {len(png_bytes)} bytes")
            browser.close()
    except Exception as e:
        print(f"[Tech V2] Playwright error: {e}")
        import traceback
        traceback.print_exc()

    return png_bytes, used_date

def insert_technical_analysis_v2_slide(
    prs: Presentation,
    image_bytes: bytes,
    used_date=None,
    price_mode: str = "Last Price",
    *,
    placeholder_name: str = "spx_v2",
    left_cm: float = 1.13,
    top_cm: float = 4.8,      # Slightly lower to avoid subtitle
    width_cm: float = 23.67,
    height_cm: float = 10.5,  # Reduced from 11.5 for better fit
) -> Presentation:
    """Insert Technical Analysis v2 chart into PowerPoint.

    Parameters
    ----------
    prs : Presentation
        PowerPoint presentation to modify.
    image_bytes : bytes
        PNG chart image data.
    used_date : Timestamp, optional
        Date for source footnote.
    price_mode : str
        "Last Price" or "Last Close".
    placeholder_name : str
        Name of placeholder shape to find target slide.
    left_cm, top_cm, width_cm : float
        Chart position and size.

    Returns
    -------
    Presentation
        Modified presentation.
    """
    if not image_bytes:
        return prs

    # Find target slide
    target_slide = None
    for slide in prs.slides:
        for shape in slide.shapes:
            name_attr = getattr(shape, "name", "").lower()
            if name_attr == placeholder_name.lower():
                target_slide = slide
                break
            if shape.has_text_frame:
                text_lower = (shape.text or "").strip().lower()
                if text_lower == f"[{placeholder_name.lower()}]":
                    target_slide = slide
                    break
        if target_slide:
            break

    if target_slide is None:
        print(f"[Tech V2] Slide with placeholder '{placeholder_name}' not found")
        return prs

    # Insert chart image
    stream = BytesIO(image_bytes)
    picture = target_slide.shapes.add_picture(
        stream, Cm(left_cm), Cm(top_cm), width=Cm(width_cm), height=Cm(height_cm)
    )

    # Send to back
    spTree = target_slide.shapes._spTree
    pic_element = picture._element
    spTree.remove(pic_element)
    spTree.insert(2, pic_element)

    print(f"[DEBUG] Inserting image: {width_cm} × {height_cm} cm at ({left_cm}, {top_cm})")

    # Insert source footnote
    if used_date is not None:
        date_str = used_date.strftime("%d/%m/%Y")
        suffix = " Close" if price_mode.lower() == "last close" else ""
        source_text = f"Source: Bloomberg, Herculis Group, Data as of {date_str}{suffix}"

        source_placeholder = f"{placeholder_name}_source"
        source_patterns = [f"[{source_placeholder}]", source_placeholder]

        for shape in target_slide.shapes:
            name_attr = getattr(shape, "name", "")
            if name_attr and name_attr.lower() == source_placeholder.lower():
                if shape.has_text_frame:
                    runs = shape.text_frame.paragraphs[0].runs
                    attrs = _get_run_font_attributes(runs[0]) if runs else (None, None, None, None, None, None)
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    new_run = p.add_run()
                    new_run.text = source_text
                    _apply_run_font_attributes(new_run, *attrs)
                break
            if shape.has_text_frame:
                for pattern in source_patterns:
                    if pattern.lower() in (shape.text or "").lower():
                        runs = shape.text_frame.paragraphs[0].runs
                        attrs = _get_run_font_attributes(runs[0]) if runs else (None, None, None, None, None, None)
                        shape.text_frame.clear()
                        p = shape.text_frame.paragraphs[0]
                        new_run = p.add_run()
                        new_run.text = source_text
                        _apply_run_font_attributes(new_run, *attrs)
                        break
                else:
                    continue
                break

    return prs