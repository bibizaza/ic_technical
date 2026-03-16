"""
Generic PowerPoint manipulation utilities for technical analysis.

This module provides generic functions for inserting scores, charts, gauges, and other
elements into PowerPoint presentations. All instrument-specific modules use these
generic functions by providing instrument-specific parameters.

This eliminates ~10,000 lines of duplicate code across 20 instrument files.
"""

from typing import Optional, Callable
from pptx import Presentation
from pptx.util import Cm
from io import BytesIO

from technical_analysis.common_helpers import (
    _get_run_font_attributes,
    _apply_run_font_attributes,
)


def find_slide_by_placeholder(prs: Presentation, placeholder_text: str) -> Optional[int]:
    """
    Find the index of a slide containing a specific placeholder text.

    This function searches for slides with exact placeholder matches to avoid
    false positives (e.g., YTD equity slide containing "Nikkei 225" in a legend
    should not match when searching for the dedicated "nikkei" slide).

    Parameters
    ----------
    prs : Presentation
        PowerPoint presentation object.
    placeholder_text : str
        Text to search for (e.g., "spx", "gold", "bitcoin").

    Returns
    -------
    int or None
        Index of the slide containing the placeholder, or None if not found.

    Matching Rules
    --------------
    1. Shape name exactly matches placeholder_text (case-insensitive)
    2. Shape text exactly matches "[placeholder_text]" (case-insensitive)
    3. Shape text exactly matches "placeholder_text" as standalone word
    """
    target_lower = placeholder_text.lower()
    bracket_placeholder = f"[{target_lower}]"

    for idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            # Rule 1: Check if shape name matches exactly
            shape_name = getattr(shape, "name", "").lower()
            if shape_name == target_lower:
                return idx

            # Rule 2 & 3: Check text content for exact matches
            if shape.has_text_frame:
                text = (shape.text or "").strip().lower()
                # Match [placeholder] or exact placeholder text
                if text == bracket_placeholder or text == target_lower:
                    return idx

    return None


def insert_score_number(
    prs: Presentation,
    score: Optional[float],
    instrument_name: str,
    score_type: str,  # 'tech_score' or 'momentum_score'
    placeholder_patterns: list = None,
) -> Presentation:
    """
    Insert a technical or momentum score into a PowerPoint slide.

    Generic version that works for all instruments.

    Parameters
    ----------
    prs : Presentation
        PowerPoint presentation object.
    score : float or None
        The score to insert.
    instrument_name : str
        Instrument identifier (e.g., "spx", "gold", "bitcoin").
    score_type : str
        Type of score: 'tech_score' or 'momentum_score'.
    placeholder_patterns : list, optional
        List of placeholder patterns to search for (default: ["[XXX]", "XXX"]).

    Returns
    -------
    Presentation
        Modified presentation object.
    """
    if placeholder_patterns is None:
        placeholder_patterns = ["[XXX]", "XXX"]

    score_text = "N/A" if score is None else f"{int(round(float(score)))}"
    placeholder_name = f"{score_type}_{instrument_name}"

    slide_idx = find_slide_by_placeholder(prs, instrument_name)
    if slide_idx is None:
        return prs

    slide = prs.slides[slide_idx]

    # First search for a shape named exactly as the placeholder
    for shape in slide.shapes:
        if getattr(shape, "name", "").lower() == placeholder_name:
            if shape.has_text_frame:
                runs = shape.text_frame.paragraphs[0].runs
                attrs = _get_run_font_attributes(runs[0]) if runs else (None, None, None, None, None, None)
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                new_run = p.add_run()
                new_run.text = score_text
                _apply_run_font_attributes(new_run, *attrs)
            return prs

    # Otherwise, search for textual placeholders within shapes
    for shape in slide.shapes:
        if shape.has_text_frame:
            for pattern in placeholder_patterns:
                if pattern in (shape.text or ""):
                    runs = shape.text_frame.paragraphs[0].runs
                    attrs = _get_run_font_attributes(runs[0]) if runs else (None, None, None, None, None, None)
                    new_text = shape.text.replace(pattern, score_text)
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    new_run = p.add_run()
                    new_run.text = new_text
                    _apply_run_font_attributes(new_run, *attrs)
                    return prs

    return prs


def insert_chart_image(
    prs: Presentation,
    image_bytes: BytesIO,
    instrument_name: str,
    placeholder_name: str,
    left_cm: float = 1.66,
    top_cm: float = 7.11,
    width_cm: float = 21.41,
    height_cm: float = 7.53,
) -> Presentation:
    """
    Insert a chart image into a PowerPoint slide.

    Generic version that works for all instruments.

    Parameters
    ----------
    prs : Presentation
        PowerPoint presentation object.
    image_bytes : BytesIO
        Image data as bytes.
    instrument_name : str
        Instrument identifier (e.g., "spx", "gold", "bitcoin").
    placeholder_name : str
        Name of the placeholder shape to replace.
    left_cm, top_cm, width_cm, height_cm : float
        Position and dimensions of the image in centimeters.

    Returns
    -------
    Presentation
        Modified presentation object.
    """
    slide_idx = find_slide_by_placeholder(prs, instrument_name)
    if slide_idx is None:
        return prs

    slide = prs.slides[slide_idx]

    # Look for the placeholder shape
    for shape in slide.shapes:
        if getattr(shape, "name", "").lower() == placeholder_name.lower():
            # Remove the placeholder
            sp = shape.element
            sp.getparent().remove(sp)
            break

    # Insert the image
    slide.shapes.add_picture(
        image_bytes,
        Cm(left_cm),
        Cm(top_cm),
        width=Cm(width_cm),
        height=Cm(height_cm),
    )

    return prs


def insert_subtitle(
    prs: Presentation,
    subtitle: str,
    instrument_name: str,
    placeholder_name: str = None,
) -> Presentation:
    """
    Insert a subtitle into a PowerPoint slide.

    Generic version that works for all instruments.

    Parameters
    ----------
    prs : Presentation
        PowerPoint presentation object.
    subtitle : str
        Subtitle text to insert.
    instrument_name : str
        Instrument identifier (e.g., "spx", "gold", "bitcoin").
    placeholder_name : str, optional
        Name of the placeholder shape (default: "{instrument_name}_text").

    Returns
    -------
    Presentation
        Modified presentation object.
    """
    if placeholder_name is None:
        placeholder_name = f"{instrument_name}_text"

    subtitle_text = subtitle or ""
    placeholder_patterns = ["[XXX]", "XXX"]

    slide_idx = find_slide_by_placeholder(prs, instrument_name)
    if slide_idx is None:
        return prs

    slide = prs.slides[slide_idx]

    # Search for the placeholder shape by name
    for shape in slide.shapes:
        if getattr(shape, "name", "").lower() == placeholder_name.lower():
            if shape.has_text_frame:
                runs = shape.text_frame.paragraphs[0].runs
                attrs = _get_run_font_attributes(runs[0]) if runs else (None, None, None, None, None, None)
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                new_run = p.add_run()
                new_run.text = subtitle_text
                _apply_run_font_attributes(new_run, *attrs)
            return prs

    # Otherwise search for placeholder text patterns
    for shape in slide.shapes:
        if shape.has_text_frame:
            for pattern in placeholder_patterns:
                if pattern in (shape.text or ""):
                    runs = shape.text_frame.paragraphs[0].runs
                    attrs = _get_run_font_attributes(runs[0]) if runs else (None, None, None, None, None, None)
                    new_text = shape.text.replace(pattern, subtitle_text)
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    new_run = p.add_run()
                    new_run.text = new_text
                    _apply_run_font_attributes(new_run, *attrs)
                    return prs

    return prs


def insert_technical_assessment(
    prs: Presentation,
    view_text: str,
    instrument_name: str,
    placeholder_name: str = None,
) -> Presentation:
    """
    Insert a technical assessment/view text into a PowerPoint slide.

    Generic version that works for all instruments.

    Parameters
    ----------
    prs : Presentation
        PowerPoint presentation object.
    view_text : str
        Assessment text to insert.
    instrument_name : str
        Instrument identifier (e.g., "spx", "gold", "bitcoin").
    placeholder_name : str, optional
        Name of the placeholder shape (default: "view_{instrument_name}").

    Returns
    -------
    Presentation
        Modified presentation object.
    """
    if placeholder_name is None:
        placeholder_name = f"view_{instrument_name}"

    slide_idx = find_slide_by_placeholder(prs, instrument_name)
    if slide_idx is None:
        return prs

    slide = prs.slides[slide_idx]

    # Search for the placeholder shape
    for shape in slide.shapes:
        if getattr(shape, "name", "").lower() == placeholder_name.lower():
            if shape.has_text_frame:
                runs = shape.text_frame.paragraphs[0].runs
                attrs = _get_run_font_attributes(runs[0]) if runs else (None, None, None, None, None, None)
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                new_run = p.add_run()
                new_run.text = view_text
                _apply_run_font_attributes(new_run, *attrs)
            return prs

    # Otherwise search for placeholder text patterns
    for shape in slide.shapes:
        if shape.has_text_frame:
            for pattern in ["[view]", "view", "[View]", "View"]:
                if pattern in (shape.text or ""):
                    runs = shape.text_frame.paragraphs[0].runs
                    attrs = _get_run_font_attributes(runs[0]) if runs else (None, None, None, None, None, None)
                    new_text = shape.text.replace(pattern, view_text)
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    new_run = p.add_run()
                    new_run.text = new_text
                    _apply_run_font_attributes(new_run, *attrs)
                    return prs

    return prs


def insert_source(
    prs: Presentation,
    instrument_name: str,
    source_text: str = "Bloomberg",
    placeholder_name: str = None,
) -> Presentation:
    """
    Insert a source attribution into a PowerPoint slide.

    Generic version that works for all instruments.

    Parameters
    ----------
    prs : Presentation
        PowerPoint presentation object.
    instrument_name : str
        Instrument identifier (e.g., "spx", "gold", "bitcoin").
    source_text : str, optional
        Source text (default: "Bloomberg").
    placeholder_name : str, optional
        Name of the placeholder shape (default: "source_{instrument_name}").

    Returns
    -------
    Presentation
        Modified presentation object.
    """
    if placeholder_name is None:
        placeholder_name = f"source_{instrument_name}"

    slide_idx = find_slide_by_placeholder(prs, instrument_name)
    if slide_idx is None:
        return prs

    slide = prs.slides[slide_idx]

    # Search for the placeholder shape
    for shape in slide.shapes:
        if getattr(shape, "name", "").lower() == placeholder_name.lower():
            if shape.has_text_frame:
                runs = shape.text_frame.paragraphs[0].runs
                attrs = _get_run_font_attributes(runs[0]) if runs else (None, None, None, None, None, None)
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                new_run = p.add_run()
                new_run.text = source_text
                _apply_run_font_attributes(new_run, *attrs)
            return prs

    # Otherwise search for placeholder text
    for shape in slide.shapes:
        if shape.has_text_frame:
            for pattern in ["[source]", "[Source]", "source", "Source"]:
                if pattern in (shape.text or ""):
                    runs = shape.text_frame.paragraphs[0].runs
                    attrs = _get_run_font_attributes(runs[0]) if runs else (None, None, None, None, None, None)
                    new_text = shape.text.replace(pattern, source_text)
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    new_run = p.add_run()
                    new_run.text = new_text
                    _apply_run_font_attributes(new_run, *attrs)
                    return prs

    return prs
