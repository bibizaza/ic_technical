"""Cryptocurrency YTD performance chart generation and insertion.

This module computes year‑to‑date (YTD) performance for a set of
cryptocurrencies, generates a YTD performance chart with connectors
and bold labels, and inserts it into a slide identified by a
placeholder named ``ytd_crypto_perf``.  The subtitle for the chart
is written into a textbox named ``ytd_crypto_subtitle`` by replacing
a ``XXX`` placeholder while preserving formatting.
"""

from __future__ import annotations

from datetime import datetime
import os
import tempfile
from typing import List, Optional

import matplotlib.dates as mdates
import matplotlib.pyplot as plt
import pandas as pd
from pptx import Presentation
from pptx.util import Cm

from .loader_update import load_data

# Colour definitions for crypto (RGB → hex)
CRYPTO_COLOURS = {
    "Ripple": "#A9D18E",    # 169,209,142
    "XRPUSD Curncy": "#A9D18E",
    "Bitcoin": "#203864",   # 32,56,100
    "XBTUSD Curncy": "#203864",
    "Binance": "#BF9000",   # 191,144,0
    "XBIUSD LUKK Curncy": "#BF9000",
    "Ethereum": "#00B0F0",  # 0,176,240
    "XETUSD Curncy": "#00B0F0",
    "Solana": "#FFC000",    # 255,192,0
    "XSOUSD Curncy": "#FFC000",
}


def get_crypto_ytd_series(file_path: str, tickers: Optional[List[str]] = None) -> pd.DataFrame:
    """
    Compute YTD performance for crypto tickers.  If ``tickers`` is
    ``None``, all crypto tickers from the parameters sheet are
    included.  YTD is calculated from 1 January of the current year.

    Returns a DataFrame with ``Date`` and one column per crypto name.
    """
    prices_df, params_df = load_data(file_path)
    now = datetime.now()
    year_start = datetime(now.year, 1, 1)
    current_year_df = prices_df[prices_df["Date"] >= year_start].reset_index(drop=True)
    crypto_params = params_df[params_df["Asset Class"] == "Crypto"].copy()
    if tickers is not None:
        crypto_params = crypto_params[crypto_params["Tickers"].isin(tickers)]
    result = pd.DataFrame()
    result["Date"] = current_year_df["Date"]
    for _, row in crypto_params.iterrows():
        ticker = str(row["Tickers"]).strip()
        name = str(row["Name"]).strip()
        if ticker not in current_year_df.columns:
            continue
        series = current_year_df[ticker].astype(float)
        base_series = series.dropna()
        if base_series.empty:
            continue
        base_val = base_series.iloc[0]
        if base_val == 0 or pd.isna(base_val):
            continue
        ytd = (series / base_val - 1.0) * 100.0
        result[name] = ytd.reset_index(drop=True)
    return result


def create_crypto_chart(df: pd.DataFrame) -> plt.Figure:
    """Create a crypto YTD chart with connectors and bold labels."""
    fig, ax = plt.subplots(figsize=(10, 5.5))
    for col in df.columns:
        if col == "Date":
            continue
        colour = CRYPTO_COLOURS.get(col, None)
        ax.plot(df["Date"], df[col], color=colour, linewidth=2)
    # Determine y-range and sort by final values
    y_min = df[[c for c in df.columns if c != "Date"]].min().min()
    y_max = df[[c for c in df.columns if c != "Date"]].max().max()
    y_range = y_max - y_min if y_max > y_min else 1.0
    series_cols = [c for c in df.columns if c != "Date"]
    sorted_cols = sorted(series_cols, key=lambda c: df[c].iloc[-1], reverse=True)
    # Annotate connectors and labels
    for idx, col in enumerate(sorted_cols):
        colour = CRYPTO_COLOURS.get(col, None)
        last_x = df["Date"].iloc[-1]
        last_y = df[col].iloc[-1]
        offset_y = -idx * 0.02 * y_range
        target_y = last_y + offset_y
        perf_text = f"{last_y:+.1f}%"
        annotation = f"{col}: {perf_text}"
        x_offset = pd.Timedelta(days=5)
        ax.plot([last_x, last_x + x_offset], [last_y, target_y], color="#BFBFBF", linewidth=0.5)
        ax.text(
            last_x + x_offset,
            target_y,
            annotation,
            color=colour,
            fontsize=8,
            fontweight="bold",
            va="center",
            ha="left",
        )
    # Style chart
    ax.set_title("YTD Performance of Cryptocurrencies (%)", fontsize=14, color="#0A1F44")
    ax.set_xlabel("")
    ax.set_ylabel("Performance")
    ax.xaxis.set_major_locator(mdates.MonthLocator())
    ax.xaxis.set_major_formatter(mdates.DateFormatter('%b'))
    fig.autofmt_xdate()
    ax.axhline(0, color='gray', linewidth=0.8, linestyle='--')
    ax.grid(False)
    for spine in ax.spines.values():
        spine.set_visible(False)
    fig.tight_layout()
    return fig


def insert_crypto_chart(
    prs: Presentation,
    file_path: str,
    subtitle: str = "",
    tickers: Optional[List[str]] = None,
    *,
    left_cm: float = 1.87,
    top_cm: float = 5.49,
    width_cm: float = 20.64,
    height_cm: float = 9.57,
) -> Presentation:
    """Insert a crypto YTD chart into the slide identified by ``ytd_crypto_perf``.

    The function searches for a shape named ``ytd_crypto_perf`` (or
    containing ``[ytd_crypto_perf]``) to locate the slide.  The chart is
    then inserted at the specified coordinates.  The subtitle is
    written into ``ytd_crypto_subtitle`` by replacing a ``XXX``
    placeholder, preserving the original run formatting.  If no
    placeholder slide is found, it falls back to slide 26 (index 25).

    Parameters
    ----------
    prs : Presentation
        PowerPoint presentation to modify.
    file_path : str
        Path to the Excel workbook.
    subtitle : str, optional
        Subtitle text to insert.
    tickers : list of str, optional
        Crypto tickers to include; if omitted, all cryptos are used.
    left_cm, top_cm, width_cm, height_cm : float
        Position and dimensions for the chart insertion, in centimetres.

    Returns
    -------
    Presentation
        The modified presentation.
    """
    # Compute series and figure
    df_crypto = get_crypto_ytd_series(file_path, tickers)
    fig = create_crypto_chart(df_crypto)
    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp_png:
        fig.savefig(tmp_png.name, dpi=200)
        chart_path = tmp_png.name
    try:
        # Locate slide using placeholder
        target_slide = None
        for slide in prs.slides:
            for shape in slide.shapes:
                name_attr = getattr(shape, "name", "").lower()
                if name_attr == "ytd_crypto_perf":
                    target_slide = slide
                    break
                if shape.has_text_frame:
                    text_lower = (shape.text or "").strip().lower()
                    if text_lower == "[ytd_crypto_perf]":
                        target_slide = slide
                        break
            if target_slide:
                break
        # Fallback to slide 26 (index 25)
        if target_slide is None:
            target_slide = prs.slides[min(25, len(prs.slides) - 1)]
        # Insert subtitle
        if subtitle:
            for shape in target_slide.shapes:
                name_attr = getattr(shape, "name", "")
                if name_attr and name_attr.lower() == "ytd_crypto_subtitle" and shape.has_text_frame:
                    tf = shape.text_frame
                    for paragraph in tf.paragraphs:
                        for run in paragraph.runs:
                            if "XXX" in run.text:
                                orig_font = run.font
                                run.text = run.text.replace("XXX", subtitle)
                                run.font.name = orig_font.name
                                run.font.size = orig_font.size
                                run.font.bold = orig_font.bold
                                run.font.italic = orig_font.italic
                                run.font.color.rgb = orig_font.color.rgb
                                break
                    break
        # Convert coordinates to EMU
        left = Cm(left_cm)
        top = Cm(top_cm)
        width = Cm(width_cm)
        height = Cm(height_cm)
        # Insert the chart
        picture = target_slide.shapes.add_picture(chart_path, left, top, width, height)
        # Bring picture to the front
        sp_tree = target_slide.shapes._spTree
        sp_tree.remove(picture._element)
        sp_tree.insert(1, picture._element)
        return prs
    finally:
        os.remove(chart_path)