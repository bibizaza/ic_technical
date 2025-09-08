"""Equity YTD performance chart generation and insertion.

This module computes year‑to‑date (YTD) performance for a set of
equity indices, builds a matplotlib chart with connectors and
annotations, and provides a helper to insert the chart, its subtitle
and a data‑source footnote into a PowerPoint slide.  The insertion
helper locates the appropriate slide via a placeholder named
``ytd_eq_perf`` and inserts the chart at user‑specified coordinates.
The subtitle is written into a separate textbox named
``ytd_eq_subtitle`` using formatting preserved from a ``XXX``
placeholder.  The data source footnote is inserted into a textbox
named ``ytd_eq_source`` (or containing ``[ytd_eq_source]``) and
reflects the effective price mode (``Last Price`` or ``Last Close``)
selected by the user.

Functions
---------

``get_equity_ytd_series``
    Compute percentage change from the start of the current year for
    each equity ticker using the selected price mode.
``create_equity_chart``
    Build a YTD line chart with connectors and annotations.
``insert_equity_chart``
    Insert the chart, subtitle and source footnote into a slide
    identified by placeholders.
"""

from __future__ import annotations

import os
import tempfile
from datetime import datetime
from typing import List, Optional, Tuple

import matplotlib.dates as mdates
import matplotlib.pyplot as plt
import pandas as pd
from pptx import Presentation
from pptx.util import Cm

from utils import adjust_prices_for_mode

try:
    # Reuse formatting helpers from the SPX module to capture and
    # reapply font attributes (size, colour, theme colour, brightness,
    # bold and italic).  These are considered internal but can be
    # imported since they reside in the same repository.  Should this
    # import fail, fallback helpers defined locally will be used.
    from technical_analysis.equity.spx import (
        _get_run_font_attributes as _capture_font_attrs,
        _apply_run_font_attributes as _apply_font_attrs,
    )
except Exception:
    def _capture_font_attrs(run) -> Tuple:
        """Fallback font attribute extractor.  Captures only size, rgb
        and basic bold/italic settings.  Theme colours and brightness
        adjustments are ignored in this fallback implementation."""
        if run is None:
            return (None, None, None, None, None, None)
        size = run.font.size
        colour = run.font.color
        rgb = None
        try:
            rgb = colour.rgb
        except Exception:
            rgb = None
        bold = run.font.bold
        italic = run.font.italic
        return (size, rgb, None, None, bold, italic)

    def _apply_font_attrs(new_run, size, rgb, theme_color, brightness, bold, italic) -> None:
        """Fallback font attribute applier.  Applies size, explicit RGB
        colour, bold and italic attributes to the provided run."""
        if size is not None:
            new_run.font.size = size
        if rgb is not None:
            try:
                new_run.font.color.rgb = rgb
            except Exception:
                pass
        if bold is not None:
            new_run.font.bold = bold
        if italic is not None:
            new_run.font.italic = italic


# Colour mapping for equities
EQUITY_COLOURS = {
    "S&P 500": "#203864",
    "SPX Index": "#203864",
    "Shenzen CSI 300": "#00B0F0",
    "CSI 300": "#00B0F0",
    "SHSZ300 Index": "#00B0F0",
    "Dax": "#FF0000",
    "DAX": "#FF0000",
    "DAX Index": "#FF0000",
    "Ibov": "#7030A0",
    "IBOV": "#7030A0",
    "IBOV Index": "#7030A0",
    "Sensex": "#70AD47",
    "Sensex Index": "#70AD47",
    "SENSEX Index": "#70AD47",
    "TASI": "#A6A6A6",
    "SASEIDX Index": "#A6A6A6",
    "Nikkei 225": "#FFC000",
    "Nikkei": "#FFC000",
    "NKY Index": "#FFC000",
    "SMI": "#C55A11",
    "SMI Index": "#C55A11",
    "Mexbol":"#00DAB0",
    "MEXBOL":"#00DAB0",
    "MEXBOL INDEX":"#00DAB0",
}


def _load_prices_and_params(file_path: str, price_mode: str) -> Tuple[pd.DataFrame, pd.DataFrame, Optional[pd.Timestamp]]:
    """Load price and parameter data from the Excel workbook, applying the
    specified price mode adjustment to the price data.

    Parameters
    ----------
    file_path : str
        Path to the Excel workbook containing a ``data_prices`` sheet and
        a ``parameters`` sheet.
    price_mode : str
        Either ``"Last Price"`` or ``"Last Close"``.  When set to
        ``"Last Close"``, rows with the most recent date (if equal to
        today's date) will be dropped prior to computing performance.

    Returns
    -------
    tuple
        A three‑tuple ``(prices_df, params_df, used_date)`` where
        ``prices_df`` contains columns ``Date`` and one column per ticker,
        ``params_df`` contains the parameter table, and ``used_date`` is
        the maximum date in the adjusted ``prices_df`` or ``None`` if no
        dates are available.
    """
    # Read the parameters sheet
    params_df = pd.read_excel(file_path, sheet_name="parameters")

    # Read the raw price sheet
    df = pd.read_excel(file_path, sheet_name="data_prices")
    # First row contains text headers (e.g. DATES, #price).  Drop it.
    df = df.drop(index=0)
    # The first column holds dates; drop any rows labelled "DATES" just in case
    df = df[df[df.columns[0]] != "DATES"]
    df.loc[:, "Date"] = pd.to_datetime(df[df.columns[0]], errors="coerce")
    # Remove the original first column
    df = df.drop(columns=[df.columns[0]])
    # Convert all price columns to numeric
    for col in df.columns:
        if col != "Date":
            df[col] = pd.to_numeric(df[col], errors="coerce")
    # Drop rows where Date is NaT
    df = df.dropna(subset=["Date"]).reset_index(drop=True)
    # Apply price mode adjustment
    prices_df, used_date = adjust_prices_for_mode(df, price_mode)
    return prices_df, params_df, used_date


def get_equity_ytd_series(
    file_path: str,
    tickers: Optional[List[str]] = None,
    *,
    price_mode: str = "Last Price",
) -> pd.DataFrame:
    """
    Compute YTD performance for equity tickers using the selected price mode.

    This function returns only the YTD performance DataFrame.  The
    effective date used for the computation (i.e. the last date after
    applying the price mode adjustment) can be obtained separately
    via ``_load_prices_and_params``.

    Parameters
    ----------
    file_path : str
        Path to the Excel workbook containing price and parameter sheets.
    tickers : list of str, optional
        Specific tickers to include; if omitted, all equities are used.
    price_mode : str, default ``"Last Price"``
        Either ``"Last Price"`` or ``"Last Close"``.  When set to
        ``"Last Close"``, rows with the most recent date (if equal to
        today's date) will be dropped prior to computing performance.

    Returns
    -------
    DataFrame
        A DataFrame containing a ``Date`` column and one column per
        equity name, holding percentage changes from the start of the
        current year.
    """
    prices_df, params_df, _ = _load_prices_and_params(file_path, price_mode)
    now = datetime.now()
    year_start = datetime(now.year, 1, 1)
    # Filter data from the start of the year
    current_year_df = prices_df[prices_df["Date"] >= year_start].reset_index(drop=True)
    # Filter parameter rows to equities
    eq_params = params_df[params_df["Asset Class"] == "Equity"].copy()
    if tickers is not None:
        eq_params = eq_params[eq_params["Tickers"].isin(tickers)]
    result = pd.DataFrame()
    result["Date"] = current_year_df["Date"].reset_index(drop=True)
    for _, row in eq_params.iterrows():
        ticker = str(row["Tickers"]).strip()
        name = str(row["Name"]).strip()
        if ticker not in current_year_df.columns:
            continue
        series = current_year_df[ticker].astype(float)
        base_series = series.dropna()
        if base_series.empty:
            continue
        base_val = base_series.iloc[0]
        if base_val == 0 or pd.isna(base_val):
            continue
        ytd = (series / base_val - 1.0) * 100.0
        result[name] = ytd.reset_index(drop=True)
    return result


def create_equity_chart(df: pd.DataFrame) -> plt.Figure:
    """Create a YTD equity chart with connectors and bold labels.

    The chart plots one line per equity index.  At the right edge,
    connectors lead to bold labels indicating the final value of each
    series.  Colours are assigned according to ``EQUITY_COLOURS``.

    Parameters
    ----------
    df : DataFrame
        The DataFrame returned by ``get_equity_ytd_series``.

    Returns
    -------
    Figure
        A matplotlib Figure object ready for saving.
    """
    fig, ax = plt.subplots(figsize=(10, 5.5))
    # Plot each equity line
    for col in df.columns:
        if col == "Date":
            continue
        colour = EQUITY_COLOURS.get(col, None)
        ax.plot(df["Date"], df[col], color=colour, linewidth=2)
    # Determine y‑range and sort by final values
    y_min = df[[c for c in df.columns if c != "Date"]].min().min()
    y_max = df[[c for c in df.columns if c != "Date"]].max().max()
    y_range = y_max - y_min if y_max > y_min else 1.0
    series_cols = [c for c in df.columns if c != "Date"]
    sorted_cols = sorted(series_cols, key=lambda c: df[c].iloc[-1], reverse=True)
    # Annotate with connectors
    for idx, col in enumerate(sorted_cols):
        colour = EQUITY_COLOURS.get(col, None)
        last_x = df["Date"].iloc[-1]
        last_y = df[col].iloc[-1]
        offset_y = -idx * 0.02 * y_range
        target_y = last_y + offset_y
        perf_text = f"{last_y:+.1f}%"
        annotation = f"{col}: {perf_text}"
        x_offset = pd.Timedelta(days=5)
        ax.plot([last_x, last_x + x_offset], [last_y, target_y], color="#BFBFBF", linewidth=0.5)
        ax.text(
            last_x + x_offset,
            target_y,
            annotation,
            color=colour,
            fontsize=8,
            fontweight="bold",
            va="center",
            ha="left",
        )
    # Style axes
    ax.set_title("YTD Performance of Equity Indices (%)", fontsize=14, color="#0A1F44")
    ax.set_xlabel("")
    ax.set_ylabel("Performance")
    ax.xaxis.set_major_locator(mdates.MonthLocator())
    ax.xaxis.set_major_formatter(mdates.DateFormatter('%b'))
    fig.autofmt_xdate()
    ax.axhline(0, color='gray', linewidth=0.8, linestyle='--')
    ax.grid(False)
    for spine in ax.spines.values():
        spine.set_visible(False)
    fig.tight_layout()
    return fig


def insert_equity_chart(
    prs: Presentation,
    file_path: str,
    subtitle: str = "",
    tickers: Optional[List[str]] = None,
    *,
    price_mode: str = "Last Price",
    left_cm: float = 1.87,
    top_cm: float = 5.49,
    width_cm: float = 20.64,
    height_cm: float = 9.57,
) -> Presentation:
    """Insert the YTD equity chart, subtitle and source footnote into a slide.

    This helper searches for a shape named ``ytd_eq_perf`` or containing
    ``[ytd_eq_perf]`` to locate the correct slide.  It then inserts
    the chart image at the specified coordinates.  The subtitle is
    written into the shape named ``ytd_eq_subtitle`` by replacing the
    placeholder ``XXX`` or ``[ytd_eq_subtitle]`` while preserving the
    original formatting of the text run.  Finally, a source footnote
    indicating the data source and date is inserted into a shape named
    ``ytd_eq_source`` (or containing ``[ytd_eq_source]``).  If the
    placeholders are not found, the corresponding elements are not
    inserted or replaced.

    Parameters
    ----------
    prs : Presentation
        PowerPoint presentation to modify.
    file_path : str
        Path to the Excel workbook.
    subtitle : str, optional
        Subtitle text to insert in place of ``XXX``.
    tickers : list of str, optional
        Equity tickers to include; if omitted, all equities are
        included.
    price_mode : str, default ``"Last Price"``
        Either ``"Last Price"`` or ``"Last Close"``.  Determines how
        price data is adjusted prior to computing YTD performance and
        affects the date displayed in the source footnote.
    left_cm, top_cm, width_cm, height_cm : float
        Coordinates and dimensions (in centimetres) for the chart
        placement.

    Returns
    -------
    Presentation
        The modified presentation.
    """
    # Compute YTD data
    df_eq = get_equity_ytd_series(file_path, tickers, price_mode=price_mode)
    # Determine the date used for the data: reload price data with the
    # specified price mode and extract the maximum date.
    _, _, used_date = _load_prices_and_params(file_path, price_mode)
    # Build the chart
    fig = create_equity_chart(df_eq)
    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp_png:
        fig.savefig(tmp_png.name, dpi=200)
        chart_path = tmp_png.name
    try:
        # Locate slide by ytd_eq_perf placeholder
        target_slide = None
        for slide in prs.slides:
            for shape in slide.shapes:
                name_attr = getattr(shape, "name", "").lower()
                if name_attr == "ytd_eq_perf":
                    target_slide = slide
                    break
                if shape.has_text_frame:
                    text_lower = (shape.text or "").strip().lower()
                    if text_lower == "[ytd_eq_perf]":
                        target_slide = slide
                        break
            if target_slide:
                break
        # If the placeholder slide is not found, skip inserting the chart.
        if target_slide is None:
            return prs
        # ------------------------------------------------------------------
        # Insert subtitle into the 'ytd_eq_subtitle' placeholder
        # ------------------------------------------------------------------
        if subtitle:
            for shape in target_slide.shapes:
                name_attr = getattr(shape, "name", "")
                if name_attr and name_attr.lower() == "ytd_eq_subtitle" and shape.has_text_frame:
                    tf = shape.text_frame
                    # We'll modify only the first paragraph's first run for simplicity
                    paragraph = tf.paragraphs[0]
                    runs = paragraph.runs
                    attrs = _capture_font_attrs(runs[0]) if runs else (None, None, None, None, None, None)
                    # Determine replacement: replace tokens "XXX" or "[ytd_eq_subtitle]"
                    original_text = "".join(run.text for run in runs) if runs else ""
                    new_text = original_text
                    if "XXX" in new_text:
                        new_text = new_text.replace("XXX", subtitle)
                    elif "[ytd_eq_subtitle]" in new_text:
                        new_text = new_text.replace("[ytd_eq_subtitle]", subtitle)
                    else:
                        new_text = subtitle
                    # Clear and insert new run
                    tf.clear()
                    p = tf.paragraphs[0]
                    new_run = p.add_run()
                    new_run.text = new_text
                    _apply_font_attrs(new_run, *attrs)
                    break
        # ------------------------------------------------------------------
        # Insert the chart image
        # ------------------------------------------------------------------
        left = Cm(left_cm)
        top = Cm(top_cm)
        width = Cm(width_cm)
        height = Cm(height_cm)
        picture = target_slide.shapes.add_picture(chart_path, left, top, width, height)
        # Move the picture to the front (just after the slide's background)
        sp_tree = target_slide.shapes._spTree
        sp_tree.remove(picture._element)
        sp_tree.insert(1, picture._element)
        # ------------------------------------------------------------------
        # Insert data source footnote
        # ------------------------------------------------------------------
        if used_date is not None:
            date_str = used_date.strftime("%d/%m/%Y")
            suffix = " Close" if price_mode.lower() == "last close" else ""
            source_text = f"Source: Bloomberg, Herculis Group, Data as of {date_str}{suffix}"
            placeholder_name = "ytd_eq_source"
            placeholder_patterns = ["[ytd_eq_source]", "ytd_eq_source"]
            inserted = False
            for shape in target_slide.shapes:
                name_attr = getattr(shape, "name", "")
                if name_attr and name_attr.lower() == placeholder_name:
                    if shape.has_text_frame:
                        runs = shape.text_frame.paragraphs[0].runs
                        attrs = _capture_font_attrs(runs[0]) if runs else (None, None, None, None, None, None)
                        shape.text_frame.clear()
                        p = shape.text_frame.paragraphs[0]
                        new_run = p.add_run()
                        new_run.text = source_text
                        _apply_font_attrs(new_run, *attrs)
                    inserted = True
                    break
                if shape.has_text_frame:
                    for pattern in placeholder_patterns:
                        if pattern.lower() in (shape.text or "").lower():
                            runs = shape.text_frame.paragraphs[0].runs
                            attrs = _capture_font_attrs(runs[0]) if runs else (None, None, None, None, None, None)
                            try:
                                new_text = shape.text.replace(pattern, source_text)
                            except Exception:
                                new_text = source_text
                            shape.text_frame.clear()
                            p = shape.text_frame.paragraphs[0]
                            new_run = p.add_run()
                            new_run.text = new_text
                            _apply_font_attrs(new_run, *attrs)
                            inserted = True
                            break
                    if inserted:
                        break
        return prs
    finally:
        # Always remove the temporary image file
        if os.path.exists(chart_path):
            try:
                os.remove(chart_path)
            except Exception:
                pass