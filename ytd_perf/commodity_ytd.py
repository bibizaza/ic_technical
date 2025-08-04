"""Commodity YTD performance chart generation and insertion.

This module computes year‑to‑date (YTD) performance for commodity
indices, creates a YTD performance chart with labelled connectors and
bold annotations, and inserts the chart into a slide identified by
a placeholder named ``ytd_commo_perf``.  The subtitle for the chart
is placed into a separate textbox named ``ytd_commo_subtitle`` by
replacing a ``XXX`` placeholder while preserving the original
formatting.  A source footnote is also inserted into a textbox named
``ytd_commo_source`` (or containing ``[ytd_commo_source]``) and
reflects the effective price mode (``Last Price`` or ``Last Close``).

Functions
---------

``get_commodity_ytd_series``
    Compute percentage change from the start of the current year for
    each commodity ticker using the selected price mode.
``create_commodity_chart``
    Build a YTD line chart with connectors and annotations.
``insert_commodity_chart``
    Insert the chart, subtitle and source footnote into a slide
    identified by placeholders.
"""

from __future__ import annotations

from datetime import datetime
import os
import tempfile
from typing import List, Optional, Tuple

import matplotlib.dates as mdates
import matplotlib.pyplot as plt
import pandas as pd
from pptx import Presentation
from pptx.util import Cm

from .loader_update import load_data
from utils import adjust_prices_for_mode

try:
    # Import formatting helpers from the SPX module to preserve
    # font attributes when replacing placeholder text.
    from technical_analysis.equity.spx import (
        _get_run_font_attributes as _capture_font_attrs,
        _apply_run_font_attributes as _apply_font_attrs,
    )
except Exception:
    def _capture_font_attrs(run):
        if run is None:
            return None, None, None, None, None, None
        size = run.font.size
        colour = run.font.color
        rgb = None
        try:
            rgb = colour.rgb
        except Exception:
            rgb = None
        return size, rgb, None, None, run.font.bold, run.font.italic
    def _apply_font_attrs(new_run, size, rgb, theme_color, brightness, bold, italic):
        if size is not None:
            new_run.font.size = size
        if rgb is not None:
            try:
                new_run.font.color.rgb = rgb
            except Exception:
                pass
        if bold is not None:
            new_run.font.bold = bold
        if italic is not None:
            new_run.font.italic = italic


# Commodity colours
COMMO_COLOURS = {
    "Gold": "#BF9000",
    "Silver": "#A6A6A6",
    "Oil (WTI)": "#3B3838",
    "Platinum": "#7F7F7F",
    "Copper": "#C55A11",
    "Uranium": "#7F6000",
}


def get_commodity_ytd_series(
    file_path: str,
    tickers: Optional[List[str]] = None,
    *,
    price_mode: str = "Last Price",
) -> pd.DataFrame:
    """
    Compute YTD performance for commodity tickers using the selected price mode.

    If ``tickers`` is ``None``, all commodity tickers are included.
    YTD is calculated from 1 January of the current year using price data
    adjusted according to ``price_mode``.

    Parameters
    ----------
    file_path : str
        Path to the Excel workbook.
    tickers : list of str, optional
        Specific commodity tickers to include; if omitted, all commodities
        are included.
    price_mode : str, default ``"Last Price"``
        Either ``"Last Price"`` or ``"Last Close"``.  When set to
        ``"Last Close"``, rows with the most recent date (if equal to
        today's date) will be dropped prior to computing performance.

    Returns
    -------
    DataFrame
        A DataFrame with a ``Date`` column and one column per
        commodity, containing percentage changes from the start of
        the current year.
    """
    # Load raw price and parameter data
    prices_df, params_df = load_data(file_path)
    # Apply price mode adjustment
    prices_df, _ = adjust_prices_for_mode(prices_df, price_mode)
    now = datetime.now()
    year_start = datetime(now.year, 1, 1)
    current_year_df = prices_df[prices_df["Date"] >= year_start].reset_index(drop=True)
    commo_params = params_df[params_df["Asset Class"] == "Commodity"].copy()
    if tickers is not None:
        commo_params = commo_params[commo_params["Tickers"].isin(tickers)]
    result = pd.DataFrame()
    result["Date"] = current_year_df["Date"]
    for _, row in commo_params.iterrows():
        ticker = str(row["Tickers"]).strip()
        name = str(row["Name"]).strip()
        if ticker not in current_year_df.columns:
            continue
        series = current_year_df[ticker].astype(float)
        base_series = series.dropna()
        if base_series.empty:
            continue
        base_val = base_series.iloc[0]
        if base_val == 0 or pd.isna(base_val):
            continue
        ytd = (series / base_val - 1.0) * 100.0
        result[name] = ytd.reset_index(drop=True)
    return result


def create_commodity_chart(df: pd.DataFrame) -> plt.Figure:
    """Create a commodity YTD chart with connectors and bold labels.

    Parameters
    ----------
    df : DataFrame
        The DataFrame returned by ``get_commodity_ytd_series``.

    Returns
    -------
    Figure
        A matplotlib Figure object representing the chart.
    """
    fig, ax = plt.subplots(figsize=(10, 5.5))
    for col in df.columns:
        if col == "Date":
            continue
        colour = COMMO_COLOURS.get(col, None)
        ax.plot(df["Date"], df[col], color=colour, linewidth=2)
    # Determine y-range and sort by final values
    y_min = df[[c for c in df.columns if c != "Date"]].min().min()
    y_max = df[[c for c in df.columns if c != "Date"]].max().max()
    y_range = y_max - y_min if y_max > y_min else 1.0
    series_cols = [c for c in df.columns if c != "Date"]
    sorted_cols = sorted(series_cols, key=lambda c: df[c].iloc[-1], reverse=True)
    # Annotate with connectors
    for idx, col in enumerate(sorted_cols):
        colour = COMMO_COLOURS.get(col, None)
        last_x = df["Date"].iloc[-1]
        last_y = df[col].iloc[-1]
        offset_y = -idx * 0.02 * y_range
        target_y = last_y + offset_y
        perf_text = f"{last_y:+.1f}%"
        annotation = f"{col}: {perf_text}"
        x_offset = pd.Timedelta(days=5)
        ax.plot([last_x, last_x + x_offset], [last_y, target_y], color="#BFBFBF", linewidth=0.5)
        ax.text(
            last_x + x_offset,
            target_y,
            annotation,
            color=colour,
            fontsize=8,
            fontweight="bold",
            va="center",
            ha="left",
        )
    # Style chart
    ax.set_title("YTD Performance of Commodities (%)", fontsize=14, color="#0A1F44")
    ax.set_xlabel("")
    ax.set_ylabel("Performance")
    ax.xaxis.set_major_locator(mdates.MonthLocator())
    ax.xaxis.set_major_formatter(mdates.DateFormatter('%b'))
    fig.autofmt_xdate()
    ax.axhline(0, color='gray', linewidth=0.8, linestyle='--')
    ax.grid(False)
    for spine in ax.spines.values():
        spine.set_visible(False)
    fig.tight_layout()
    return fig


def insert_commodity_chart(
    prs: Presentation,
    file_path: str,
    subtitle: str = "",
    tickers: Optional[List[str]] = None,
    *,
    price_mode: str = "Last Price",
    left_cm: float = 1.87,
    top_cm: float = 5.49,
    width_cm: float = 20.64,
    height_cm: float = 9.57,
) -> Presentation:
    """Insert a commodity YTD chart, subtitle and source footnote into the appropriate slide.

    This function searches for a slide containing a shape named
    ``ytd_commo_perf`` (or containing ``[ytd_commo_perf]``).  Once found,
    it inserts the chart at the specified coordinates on that slide.
    The subtitle is inserted into the shape named ``ytd_commo_subtitle``
    by replacing a ``XXX`` placeholder while preserving formatting.
    A source footnote is written into a shape named ``ytd_commo_source``
    (or containing ``[ytd_commo_source]``) using the same formatting as
    the placeholder.  If no suitable slide is found, the function
    falls back to slide 20 (index 19).

    Parameters
    ----------
    prs : Presentation
        The PowerPoint presentation.
    file_path : str
        Path to the Excel workbook.
    subtitle : str, optional
        Text to replace the ``XXX`` placeholder in the subtitle box.
    tickers : list of str, optional
        Commodity tickers to include; if omitted, all commodities are used.
    price_mode : str, default ``"Last Price"``
        Either ``"Last Price"`` or ``"Last Close"``.  Determines how
        price data is adjusted prior to computing YTD performance and
        affects the date displayed in the source footnote.
    left_cm, top_cm, width_cm, height_cm : float
        Coordinates and size for inserting the chart, in centimetres.

    Returns
    -------
    Presentation
        The modified presentation.
    """
    # Compute series and figure
    df_commo = get_commodity_ytd_series(file_path, tickers, price_mode=price_mode)
    fig = create_commodity_chart(df_commo)
    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp_png:
        fig.savefig(tmp_png.name, dpi=200)
        chart_path = tmp_png.name
    try:
        # Locate slide by placeholder name
        target_slide = None
        for slide in prs.slides:
            for shape in slide.shapes:
                name_attr = getattr(shape, "name", "").lower()
                if name_attr == "ytd_commo_perf":
                    target_slide = slide
                    break
                if shape.has_text_frame:
                    text_lower = (shape.text or "").strip().lower()
                    if text_lower == "[ytd_commo_perf]":
                        target_slide = slide
                        break
            if target_slide:
                break
        # Fallback to slide 20 (index 19) if placeholder not found
        if target_slide is None:
            target_slide = prs.slides[min(19, len(prs.slides) - 1)]
        # Insert subtitle into designated textbox
        if subtitle:
            for shape in target_slide.shapes:
                name_attr = getattr(shape, "name", "")
                if name_attr and name_attr.lower() == "ytd_commo_subtitle" and shape.has_text_frame:
                    tf = shape.text_frame
                    paragraph = tf.paragraphs[0]
                    runs = paragraph.runs
                    attrs = _capture_font_attrs(runs[0]) if runs else (None, None, None, None, None, None)
                    original_text = "".join(run.text for run in runs) if runs else ""
                    new_text = original_text
                    if "XXX" in new_text:
                        new_text = new_text.replace("XXX", subtitle)
                    elif "[ytd_commo_subtitle]" in new_text:
                        new_text = new_text.replace("[ytd_commo_subtitle]", subtitle)
                    else:
                        new_text = subtitle
                    tf.clear()
                    p = tf.paragraphs[0]
                    new_run = p.add_run()
                    new_run.text = new_text
                    _apply_font_attrs(new_run, *attrs)
                    break
        # Convert coordinates to EMU and insert the chart picture
        left = Cm(left_cm)
        top = Cm(top_cm)
        width = Cm(width_cm)
        height = Cm(height_cm)
        picture = target_slide.shapes.add_picture(chart_path, left, top, width, height)
        # Bring the picture to the front
        sp_tree = target_slide.shapes._spTree
        sp_tree.remove(picture._element)
        sp_tree.insert(1, picture._element)
        # Insert data source footnote
        # Determine used date after adjusting for price mode
        prices_df, _ = load_data(file_path)
        prices_df, used_date = adjust_prices_for_mode(prices_df, price_mode)
        if used_date is not None:
            date_str = used_date.strftime("%d/%m/%Y")
            suffix = " Close" if price_mode.lower() == "last close" else ""
            source_text = f"Source: Bloomberg, Herculis Group, Data as of {date_str}{suffix}"
            placeholder_name = "ytd_commo_source"
            placeholder_patterns = ["[ytd_commo_source]", "ytd_commo_source"]
            inserted = False
            for shape in target_slide.shapes:
                name_attr2 = getattr(shape, "name", "")
                if name_attr2 and name_attr2.lower() == placeholder_name:
                    if shape.has_text_frame:
                        runs2 = shape.text_frame.paragraphs[0].runs
                        attrs2 = _capture_font_attrs(runs2[0]) if runs2 else (None, None, None, None, None, None)
                        shape.text_frame.clear()
                        p2 = shape.text_frame.paragraphs[0]
                        new_run2 = p2.add_run()
                        new_run2.text = source_text
                        _apply_font_attrs(new_run2, *attrs2)
                    inserted = True
                    break
                if shape.has_text_frame:
                    for pattern in placeholder_patterns:
                        if pattern.lower() in (shape.text or "").lower():
                            runs2 = shape.text_frame.paragraphs[0].runs
                            attrs2 = _capture_font_attrs(runs2[0]) if runs2 else (None, None, None, None, None, None)
                            try:
                                new_text2 = shape.text.replace(pattern, source_text)
                            except Exception:
                                new_text2 = source_text
                            shape.text_frame.clear()
                            p2 = shape.text_frame.paragraphs[0]
                            new_run2 = p2.add_run()
                            new_run2.text = new_text2
                            _apply_font_attrs(new_run2, *attrs2)
                            inserted = True
                            break
                    if inserted:
                        break
        return prs
    finally:
        os.remove(chart_path)