"""
Streamlit application for generating YTD performance charts and inserting them
into a PowerPoint template (equity indices on slide 11, crypto indices on slide 26,
and commodity indices on slide 20).

Run with:
    streamlit run app.py
"""

import streamlit as st
import pandas as pd
import matplotlib.pyplot as plt
import matplotlib.dates as mdates
from io import BytesIO
import tempfile
import os
from pptx import Presentation
from pptx.util import Inches
import zipfile


def extract_ytd_series_from_df(raw: pd.DataFrame) -> pd.DataFrame:
    """
    Extract YTD percentage series by detecting non‑blank names and selecting columns
    whose base row contains a numeric value (e.g. 100). Works for YTD_Perf, Crypto_Update
    and Commodities_Update sheets.
    """
    names_row = 8
    base_row = 10
    date_col = 1
    start_data_row = base_row + 1
    dates = raw.iloc[start_data_row:, date_col].reset_index(drop=True)
    names = raw.iloc[names_row]
    base_vals = raw.iloc[base_row]
    series = {}
    for col_idx, name in enumerate(names):
        if col_idx == date_col:
            continue
        if pd.notna(name) and str(name).strip() != "" and pd.api.types.is_number(base_vals[col_idx]):
            diff_series = raw.iloc[start_data_row:, col_idx]
            series[str(name)] = diff_series.reset_index(drop=True)
    result = pd.DataFrame({"Date": dates})
    for name, s in series.items():
        result[name] = s
    return result


def create_ytd_chart_figure(df: pd.DataFrame, colour_map: dict = None):
    """
    Return a Matplotlib figure plotting YTD percentage series with annotations.
    Optionally accepts a custom colour map keyed by column names.
    """
    default_map = {
        "DAX Index": "#E21818",
        "IBOV Index": "#800080",
        "SMI Index": "#F4A460",
        "Sensex Index": "#006400",
        "SHSZ300 Index": "#1E90FF",
        "SPX Index": "#2E64FE",
        "SASEIDX Index": "#708090",
        "NKY Index": "#DAA520",
    }
    colour_map = colour_map or default_map
    fig, ax = plt.subplots(figsize=(10, 5.5))
    y_min, y_max = float('inf'), float('-inf')
    for col in df.columns:
        if col == "Date":
            continue
        y_min = min(y_min, df[col].min())
        y_max = max(y_max, df[col].max())
    y_range = y_max - y_min if y_max > y_min else 1.0
    annotations = []
    for col in df.columns:
        if col == "Date":
            continue
        colour = colour_map.get(col, None)
        label = col.replace(" Index", "")
        ax.plot(df["Date"], df[col], color=colour, linewidth=2)
        last_x = df["Date"].iloc[-1]
        last_y = df[col].iloc[-1]
        perf = last_y
        annotations.append((label, last_x, last_y, colour, perf))
    annotations.sort(key=lambda t: t[2], reverse=True)
    for idx, (label, last_x, last_y, colour, perf) in enumerate(annotations):
        offset_y = -idx * 0.02 * y_range
        perf_text = f"{perf:+.1f}%"
        annotation_text = f"{label}: {perf_text}"
        ax.annotate(
            annotation_text,
            xy=(last_x, last_y),
            xytext=(last_x + pd.Timedelta(days=5), last_y + offset_y),
            textcoords="data",
            color=colour,
            fontsize=8,
            va="center",
            ha="left",
            arrowprops=dict(arrowstyle="-", color=colour, lw=0.5),
        )
    ax.set_title("YTD Performance of Indices (%)", fontsize=14, color="#0A1F44")
    ax.set_xlabel("")
    ax.set_ylabel("Performance")
    ax.xaxis.set_major_locator(mdates.MonthLocator())
    ax.xaxis.set_major_formatter(mdates.DateFormatter('%b'))
    fig.autofmt_xdate()
    ax.axhline(0, color='gray', linewidth=0.8, linestyle='--')
    ax.grid(False)
    for spine in ax.spines.values():
        spine.set_visible(False)
    fig.tight_layout()
    return fig


def insert_chart_into_equity_slide(prs: Presentation, fig, subtitle_text: str = "") -> Presentation:
    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp_png:
        fig.savefig(tmp_png.name, dpi=200)
        chart_path = tmp_png.name
    try:
        slide = prs.slides[10]
        if subtitle_text:
            for shape in slide.shapes:
                if shape.name == "YTD_EQ_Perf" and shape.has_text_frame:
                    tf = shape.text_frame
                    for paragraph in tf.paragraphs:
                        for run in paragraph.runs:
                            if "XXX" in run.text:
                                original_font = run.font
                                run.text = run.text.replace("XXX", subtitle_text)
                                run.font.name = original_font.name
                                run.font.size = original_font.size
                                run.font.bold = original_font.bold
                                run.font.italic = original_font.italic
                                run.font.color.rgb = original_font.color.rgb
                    break
        left = Inches(1.87 / 2.54)
        top = Inches(5.49 / 2.54)
        width = Inches(20.64 / 2.54)
        height = Inches(9.57 / 2.54)
        picture = slide.shapes.add_picture(chart_path, left, top, width=width, height=height)
        sp_tree = slide.shapes._spTree
        sp_tree.remove(picture._element)
        sp_tree.insert(1, picture._element)
        return prs
    finally:
        os.remove(chart_path)


def insert_chart_into_crypto_slide(prs: Presentation, fig, subtitle_text: str = "") -> Presentation:
    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp_png:
        fig.savefig(tmp_png.name, dpi=200)
        chart_path = tmp_png.name
    try:
        slide_idx = 25
        for idx, slide in enumerate(prs.slides):
            if any(shp.name == "Crypto_YTD" for shp in slide.shapes):
                slide_idx = idx
                break
        slide = prs.slides[slide_idx]
        updated = False
        if subtitle_text:
            for shape in slide.shapes:
                if shape.name == "Crypto_YTD_Subtitle" and shape.has_text_frame:
                    tf = shape.text_frame
                    for paragraph in tf.paragraphs:
                        for run in paragraph.runs:
                            if "XXX" in run.text:
                                original_font = run.font
                                run.text = run.text.replace("XXX", subtitle_text)
                                run.font.name = original_font.name
                                run.font.size = original_font.size
                                run.font.bold = original_font.bold
                                run.font.italic = original_font.italic
                                run.font.color.rgb = original_font.color.rgb
                                updated = True
                                break
                    if updated:
                        break
            if not updated:
                for shape in slide.shapes:
                    if shape.has_text_frame and "Crypto Update" in shape.text:
                        tf = shape.text_frame
                        first_run = tf.paragraphs[0].runs[0]
                        new_para = tf.add_paragraph()
                        new_run = new_para.add_run()
                        new_run.text = subtitle_text
                        new_run.font.name = first_run.font.name
                        new_run.font.size = first_run.font.size
                        new_run.font.bold = first_run.font.bold
                        new_run.font.italic = first_run.font.italic
                        new_run.font.color.rgb = first_run.font.color.rgb
                        break
        left = Inches(1.87 / 2.54)
        top = Inches(5.49 / 2.54)
        width = Inches(20.64 / 2.54)
        height = Inches(9.57 / 2.54)
        picture = slide.shapes.add_picture(chart_path, left, top, width=width, height=height)
        sp_tree = slide.shapes._spTree
        sp_tree.remove(picture._element)
        sp_tree.insert(1, picture._element)
        return prs
    finally:
        os.remove(chart_path)


def insert_chart_into_commodity_slide(prs: Presentation, fig, subtitle_text: str = "") -> Presentation:
    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp_png:
        fig.savefig(tmp_png.name, dpi=200)
        chart_path = tmp_png.name
    try:
        slide_idx = 19
        for idx, slide in enumerate(prs.slides):
            if any(shp.name == "YTD_Commo" for shp in slide.shapes):
                slide_idx = idx
                break
        slide = prs.slides[slide_idx]
        updated = False
        if subtitle_text:
            for shape in slide.shapes:
                if shape.name == "YTD_Commo_Subtitle" and shape.has_text_frame:
                    tf = shape.text_frame
                    for paragraph in tf.paragraphs:
                        for run in paragraph.runs:
                            if "XXX" in run.text:
                                original_font = run.font
                                run.text = run.text.replace("XXX", subtitle_text)
                                run.font.name = original_font.name
                                run.font.size = original_font.size
                                run.font.bold = original_font.bold
                                run.font.italic = original_font.italic
                                run.font.color.rgb = original_font.color.rgb
                                updated = True
                                break
                    if updated:
                        break
            if not updated:
                for shape in slide.shapes:
                    if shape.has_text_frame and "Commodity Update" in shape.text:
                        tf = shape.text_frame
                        first_run = tf.paragraphs[0].runs[0]
                        new_para = tf.add_paragraph()
                        new_run = new_para.add_run()
                        new_run.text = subtitle_text
                        new_run.font.name = first_run.font.name
                        new_run.font.size = first_run.font.size
                        new_run.font.bold = first_run.font.bold
                        new_run.font.italic = first_run.font.italic
                        new_run.font.color.rgb = first_run.font.color.rgb
                        break
        left = Inches(1.87 / 2.54)
        top = Inches(5.49 / 2.54)
        width = Inches(20.64 / 2.54)
        height = Inches(9.57 / 2.54)
        picture = slide.shapes.add_picture(chart_path, left, top, width=width, height=height)
        sp_tree = slide.shapes._spTree
        sp_tree.remove(picture._element)
        sp_tree.insert(1, picture._element)
        return prs
    finally:
        os.remove(chart_path)


def merge_vba(original_ppt_bytes: bytes, new_pptx_bytes: bytes) -> bytes:
    orig_io = BytesIO(original_ppt_bytes)
    new_io = BytesIO(new_pptx_bytes)
    with zipfile.ZipFile(orig_io, 'r') as orig_zip:
        vba_name = 'ppt/vbaProject.bin'
        if vba_name not in orig_zip.namelist():
            return new_pptx_bytes
        vba_data = orig_zip.read(vba_name)
    with zipfile.ZipFile(new_io, 'a') as new_zip:
        new_zip.writestr(vba_name, vba_data)
    new_io.seek(0)
    return new_io.read()


# Streamlit UI
st.set_page_config(page_title="Market, Crypto & Commodities YTD Chart Updater", layout="wide")
st.title("Market, Crypto & Commodities YTD Chart Updater")
st.write(
    "Upload your Excel file containing the **YTD_Perf**, **Crypto_Update** and **Commodities_Update** sheets and your PPTM/PPTX "
    "template. The app will generate charts for equity indices (slide 11), crypto indices (slide 26) and commodities (slide 20), "
    "update subtitles and insert them into the appropriate slides."
)

excel_file = st.file_uploader("Upload Global_IC Excel file", type=["xlsm", "xlsx", "xls"])
pptx_file = st.file_uploader("Upload PowerPoint template (PPTM/PPTX)", type=["pptm", "pptx"])

if excel_file is not None:
    try:
        raw_df_ytd = pd.read_excel(excel_file, sheet_name="YTD_Perf", header=None)
        df_ytd = extract_ytd_series_from_df(raw_df_ytd)
        raw_df_crypto = pd.read_excel(excel_file, sheet_name="Crypto_Update", header=None)
        df_crypto = extract_ytd_series_from_df(raw_df_crypto)
        raw_df_commo = pd.read_excel(excel_file, sheet_name="Commodities_Update", header=None)
        df_commo = extract_ytd_series_from_df(raw_df_commo)
        # Rename crypto tickers to names
        rename_map_crypto = {
            "XRPUSD Curncy": "Ripple",
            "XBTUSD Curncy": "Bitcoin",
            "XETUSD Curncy": "Ethereum",
            "XSOUSD Curncy": "Solana",
            "XBIUSD LUKK Curncy": "Binance",
        }
        df_crypto = df_crypto.rename(columns=rename_map_crypto)
        # Rename commodity tickers to names if known
        rename_map_commo = {
            "GC1 Comdty": "Gold",
            "SI1 Comdty": "Silver",
            "CL1 Comdty": "Oil (WTI)",
            "PL1 Comdty": "Platinum",
            "HG1 Comdty": "Copper",
            "UXA1 Comdty": "Uranium",
            "LP1 Curncy": "Platinum",  # fallback if ticker differs
        }
        df_commo = df_commo.rename(columns=rename_map_commo)
        st.success("Excel file loaded successfully.")
        # Show previews
        st.subheader("Preview of extracted YTD series (Equities)")
        st.dataframe(df_ytd.head())
        st.subheader("Preview of extracted YTD series (Crypto)")
        st.dataframe(df_crypto.head())
        st.subheader("Preview of extracted YTD series (Commodities)")
        st.dataframe(df_commo.head())
        # Colour maps
        crypto_colours = {
            "Ripple": "#5eb05a",     # green
            "Bitcoin": "#103a77",    # dark blue
            "Ethereum": "#007fc0",   # medium blue
            "Solana": "#f4b840",     # yellow/orange
            "Binance": "#d5282f",    # red
        }
        commo_colours = {
            "Gold": "#c69c3a",       # gold
            "Silver": "#c0c0c0",     # silver/gray
            "Oil (WTI)": "#333333",  # dark gray/black
            "Platinum": "#a7a7a7",   # light gray
            "Copper": "#d66e2f",     # copper/orange
            "Uranium": "#818118",    # olive green
        }
        fig_ytd = create_ytd_chart_figure(df_ytd)
        fig_crypto = create_ytd_chart_figure(df_crypto, colour_map=crypto_colours)
        fig_commo = create_ytd_chart_figure(df_commo, colour_map=commo_colours)
        st.subheader("Equity YTD Performance Chart")
        st.pyplot(fig_ytd)
        st.subheader("Crypto YTD Performance Chart")
        st.pyplot(fig_crypto)
        st.subheader("Commodities YTD Performance Chart")
        st.pyplot(fig_commo)
        if pptx_file is not None:
            subtitle_eq = st.text_input("Subtitle for the equity chart", "", key="subtitle_eq")
            subtitle_crypto = st.text_input("Subtitle for the crypto chart", "", key="subtitle_crypto")
            subtitle_commo = st.text_input("Subtitle for the commodities chart", "", key="subtitle_commo")
            if st.button("Generate updated PPT"):
                template_bytes = pptx_file.getbuffer().tobytes()
                with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp_input:
                    tmp_input.write(template_bytes)
                    tmp_input.flush()
                    prs = Presentation(tmp_input.name)
                prs = insert_chart_into_equity_slide(prs, fig_ytd, subtitle_text=subtitle_eq)
                prs = insert_chart_into_crypto_slide(prs, fig_crypto, subtitle_text=subtitle_crypto)
                prs = insert_chart_into_commodity_slide(prs, fig_commo, subtitle_text=subtitle_commo)
                out_stream = BytesIO()
                prs.save(out_stream)
                new_pptx_bytes = out_stream.getvalue()
                merged_bytes = merge_vba(template_bytes, new_pptx_bytes)
                if merged_bytes != new_pptx_bytes:
                    file_name = "updated_presentation.pptm"
                    mime = "application/vnd.ms-powerpoint.presentation.macroEnabled.12"
                else:
                    file_name = "updated_presentation.pptx"
                    mime = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                st.success("Updated presentation created successfully.")
                st.download_button(
                    label="Download updated presentation",
                    data=merged_bytes,
                    file_name=file_name,
                    mime=mime,
                )
    except Exception as e:
        st.error(f"Error processing Excel file: {e}")
        st.stop()
