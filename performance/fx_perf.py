"""FX performance dashboard generation with price‐mode awareness and source footnotes.

This module produces charts summarising recent performance for a selection
of foreign‐exchange (FX) pairs and indices.  Users can choose to base
calculations on either the most recent intraday price ("Last Price") or
the previous day's closing price ("Last Close").  The resulting
figures include a weekly returns bar chart and a heatmap of longer
horizons.  When the charts are inserted into a PowerPoint
presentation, a source footnote is added to the designated text
boxes on each slide, reflecting the effective date used in the
computations and the chosen price mode.  Text formatting from the
placeholders is preserved.

Key functions
-------------

* ``create_weekly_performance_chart`` – Build the weekly bar chart and
  return both the image bytes and the effective date used for
  computation.
* ``create_historical_performance_table`` – Build the heatmap of
  returns for multiple horizons and return the image bytes and the
  effective date.
* ``insert_fx_performance_bar_slide`` – Insert the weekly bar chart
  and source footnote into its slide.
* ``insert_fx_performance_histo_slide`` – Insert the historical
  performance heatmap and source footnote into its slide.

Usage example::

    from performance.fx_perf import (
        create_weekly_performance_chart,
        create_historical_performance_table,
        insert_fx_performance_bar_slide,
        insert_fx_performance_histo_slide,
    )

    # Generate charts with price-mode awareness
    bar_bytes, used_date = create_weekly_performance_chart(excel_path, price_mode="Last Close")
    histo_bytes, _ = create_historical_performance_table(excel_path, price_mode="Last Close")
    # Insert into PPT, supplying used_date and price_mode for source footnote
    prs = insert_fx_performance_bar_slide(prs, bar_bytes, used_date, price_mode)
    prs = insert_fx_performance_histo_slide(prs, histo_bytes, used_date, price_mode)

"""

from __future__ import annotations

import io
import pathlib
from typing import Dict, List, Tuple, Union, Optional

import matplotlib.pyplot as plt
import matplotlib.colors as mcolors
import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.util import Cm
from helpers.flag_utils import get_flag_html
from jinja2 import Template
from html2image import Html2Image

from utils import adjust_prices_for_mode
from market_compass.weekly_performance.html_template import (
    CURRENCY_WEEKLY_HTML_TEMPLATE,
    CURRENCY_HISTORICAL_HTML_TEMPLATE,
)

try:
    # Reuse font attribute helpers from the SPX module if available
    from technical_analysis.equity.spx import (
        _get_run_font_attributes as _capture_font_attrs,  # type: ignore
        _apply_run_font_attributes as _apply_font_attrs,  # type: ignore
    )
except Exception:
    # Fallback helpers: only support basic size and RGB; ignore theme colours
    def _capture_font_attrs(run):  # type: ignore
        if run is None:
            return (None, None, None, None, None, None)
        size = run.font.size
        colour = run.font.color
        rgb = None
        try:
            rgb = colour.rgb
        except Exception:
            rgb = None
        return (size, rgb, None, None, run.font.bold, run.font.italic)

    def _apply_font_attrs(new_run, size, rgb, theme_color, brightness, bold, italic):  # type: ignore
        if size is not None:
            new_run.font.size = size
        if rgb is not None:
            try:
                new_run.font.color.rgb = rgb
            except Exception:
                pass
        if bold is not None:
            new_run.font.bold = bold
        if italic is not None:
            new_run.font.italic = italic

###############################################################################
# Data loading and return calculations
###############################################################################

def _load_price_data(excel_path: Union[str, pathlib.Path], tickers: List[str]) -> pd.DataFrame:
    """Load a DataFrame of dates and prices for the specified FX tickers.

    Parameters
    ----------
    excel_path : str or pathlib.Path
        Path to the Excel workbook containing a sheet named ``data_prices``.
    tickers : list of str
        Column names in the sheet corresponding to the desired FX pairs.

    Returns
    -------
    DataFrame
        A tidy DataFrame with columns ``Date`` and one column per ticker,
        sorted by date and with numeric prices.
    """
    df = pd.read_excel(excel_path, sheet_name="data_prices")
    # Drop the first row which usually contains metadata and rows with the
    # sentinel ``DATES`` value.
    df = df.drop(index=0)
    df = df[df[df.columns[0]] != "DATES"].copy()
    df = df.rename(columns={df.columns[0]: "Date"})
    df["Date"] = pd.to_datetime(df["Date"], errors="coerce")
    # Restrict to requested tickers and coerce to numeric
    out = df[["Date"] + tickers].copy()
    for t in tickers:
        out[t] = pd.to_numeric(out[t], errors="coerce")
    return out.dropna(subset=["Date"]).sort_values("Date").reset_index(drop=True)


def _compute_horizon_returns(
    df: pd.DataFrame,
    ticker: str,
    today: pd.Timestamp,
    days: int,
) -> float:
    """Compute the percentage return of a series over the given lookback.

    This helper finds the last available price on or before ``today`` minus
    ``days`` and compares it to the most recent price.  If insufficient
    history exists, NaN is returned.

    Parameters
    ----------
    df : DataFrame
        A tidy DataFrame with ``Date`` and price columns.
    ticker : str
        Column name for which to compute the return.
    today : Timestamp
        Anchor date; typically the last date in the dataset.
    days : int
        Lookback period in days.

    Returns
    -------
    float
        Percent return over the lookback horizon, or NaN if not computable.
    """
    past_date = today - pd.Timedelta(days=days)
    past_series = df.loc[df["Date"] <= past_date, ticker]
    if len(past_series) == 0:
        return float("nan")
    past_price = past_series.iloc[-1]
    current_price = df[ticker].iloc[-1]
    if past_price == 0 or pd.isna(past_price):
        return float("nan")
    return (current_price - past_price) / past_price * 100.0


def _build_returns_table(
    df: pd.DataFrame,
    mapping: Dict[str, str],
) -> pd.DataFrame:
    """Calculate performance metrics for each FX ticker.

    Returns a DataFrame indexed by the ticker code with the following
    columns: ``Name``, ``1W``, ``1M``, ``3M``, ``6M``, ``12M`` and ``YTD``.

    Parameters
    ----------
    df : DataFrame
        A tidy DataFrame as returned by ``_load_price_data``.
    mapping : dict
        Mapping from ticker codes to human‑friendly names.

    Returns
    -------
    DataFrame
        Performance table.
    """
    today = df["Date"].max()
    results: Dict[str, Dict[str, float]] = {}
    for ticker in mapping:
        results[ticker] = {
            "1W": _compute_horizon_returns(df, ticker, today, 7),
            "1M": _compute_horizon_returns(df, ticker, today, 30),
            "3M": _compute_horizon_returns(df, ticker, today, 90),
            "6M": _compute_horizon_returns(df, ticker, today, 180),
            "12M": _compute_horizon_returns(df, ticker, today, 365),
        }
        # Year‑to‑date: look back to the start of the calendar year
        start_of_year = pd.Timestamp(year=today.year, month=1, day=1)
        past_series = df.loc[df["Date"] <= start_of_year, ticker]
        if len(past_series) > 0:
            past_price = past_series.iloc[-1]
            current_price = df[ticker].iloc[-1]
            if past_price and not pd.isna(past_price):
                results[ticker]["YTD"] = (current_price - past_price) / past_price * 100.0
            else:
                results[ticker]["YTD"] = float("nan")
        else:
            results[ticker]["YTD"] = float("nan")
    table = pd.DataFrame.from_dict(results, orient="index")
    table["Name"] = table.index.map(mapping.get)
    return table[["Name", "1W", "1M", "3M", "6M", "12M", "YTD"]]


###############################################################################
# Figure generation functions
###############################################################################

def create_weekly_performance_chart(
    excel_path: Union[str, pathlib.Path],
    ticker_mapping: Dict[str, str] | None = None,
    *,
    width: float = 14.0,
    height: float = 5.0,
    price_mode: str = "Last Price",
) -> Tuple[bytes, Optional[pd.Timestamp]]:
    """Generate a bar chart of 1‑week FX returns with price‑mode adjustment.

    This helper reads the specified Excel file, optionally adjusts the price
    history according to the selected ``price_mode`` (``"Last Price"`` vs
    ``"Last Close"``), computes 1‑week returns for the configured FX
    tickers, sorts them in descending order and constructs a horizontal
    bar chart.  The effective date used for the calculations is returned
    along with the PNG data.

    Parameters
    ----------
    excel_path : str or pathlib.Path
        Path to the Excel workbook containing the price series.
    ticker_mapping : dict, optional
        Mapping from ticker codes to display names.  If not provided,
        sensible defaults are used.
    width, height : float
        Dimensions of the generated figure in inches.
    price_mode : str, default ``"Last Price"``
        One of ``"Last Price"`` or ``"Last Close"``.  Determines whether
        intraday prices are used or the most recent row is dropped if it
        corresponds to today's date.

    Returns
    -------
    tuple
        A two‑tuple ``(image_bytes, used_date)`` where ``image_bytes`` is
        the PNG data for the bar chart and ``used_date`` is the effective
        date in the adjusted price series.
    """
    default_mapping = {
        "DXY Curncy": "Dollar Index",
        "EURUSD Curncy": "EUR/USD",
        "EURCHF Curncy": "EUR/CHF",
        "EURJPY Curncy": "EUR/JPY",
        "EURBRL Curncy": "EUR/BRL",
        "EURMXN Curncy": "EUR/MXN",
    }
    mapping = ticker_mapping or default_mapping
    tickers = list(mapping.keys())
    # Load data and adjust according to price mode
    df = _load_price_data(excel_path, tickers)
    df_adj, used_date = adjust_prices_for_mode(df, price_mode)
    perf = _build_returns_table(df_adj, mapping)
    # Build bar chart data sorted by performance descending
    bar_df = perf[["Name", "1W"]].dropna().sort_values("1W", ascending=False).reset_index(drop=True)
    # Create figure and axis
    fig, ax = plt.subplots(figsize=(width, height))
    bar_colors = ["#2E8B57" if x > 0 else "#C70039" for x in bar_df["1W"]]
    bars = ax.barh(bar_df["Name"], bar_df["1W"], color=bar_colors)
    # Add labels
    for bar in bars:
        width_val = bar.get_width()
        x_pos = width_val + (0.1 if width_val > 0 else -0.1)
        ax.text(
            x_pos,
            bar.get_y() + bar.get_height() / 2.0,
            f"{width_val:.1f}%",
            va="center",
            ha="left" if width_val > 0 else "right",
            fontweight="bold",
            color="black",
        )
    # Style chart
    ax.axvline(0.0, color="grey", linewidth=0.8, linestyle="--")
    ax.set_xticks([])
    ax.set_yticks(range(len(bar_df)))
    ax.set_yticklabels(bar_df["Name"], fontsize=10)
    ax.invert_yaxis()
    for spine in ax.spines.values():
        spine.set_visible(False)
    ax.tick_params(axis="y", length=0)
    # Provide extra horizontal space so bars do not overlap names
    min_val = min(bar_df["1W"].min(), 0)
    max_val = bar_df["1W"].max()
    margin = (max_val - min_val) * 0.2
    ax.set_xlim(min_val - margin, max_val + margin)
    fig.subplots_adjust(left=0.4)
    # Export to PNG
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=300, bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf.getvalue(), used_date


def create_historical_performance_table(
    excel_path: Union[str, pathlib.Path],
    ticker_mapping: Dict[str, str] | None = None,
    *,
    width: float = 14.0,
    height: float = 6.0,
    price_mode: str = "Last Price",
) -> Tuple[bytes, Optional[pd.Timestamp]]:
    """Generate a heatmap table of FX returns with price‑mode adjustment.

    This helper reads the price data, optionally adjusts it according to
    ``price_mode``, computes returns for multiple horizons, sorts the
    table by YTD performance (descending) and constructs a heatmap where
    each column is coloured independently.  The effective date used
    for the computations is returned along with the PNG data.

    Parameters
    ----------
    excel_path : str or pathlib.Path
        Path to the Excel workbook containing the price series.
    ticker_mapping : dict, optional
        Mapping from ticker codes to display names.  If not provided,
        sensible defaults are used.
    width, height : float
        Dimensions of the generated figure in inches.
    price_mode : str, default ``"Last Price"``
        Either ``"Last Price"`` or ``"Last Close"``.  Determines how
        price data is adjusted prior to computing returns and which
        effective date appears in the source footnote.

    Returns
    -------
    tuple
        A two‑tuple ``(image_bytes, used_date)`` where ``image_bytes`` is
        the PNG data for the heatmap and ``used_date`` is the effective
        date in the adjusted price series.
    """
    default_mapping = {
        "DXY Curncy": "Dollar Index",
        "EURUSD Curncy": "EUR/USD",
        "EURCHF Curncy": "EUR/CHF",
        "EURJPY Curncy": "EUR/JPY",
        "EURBRL Curncy": "EUR/BRL",
        "EURMXN Curncy": "EUR/MXN",
    }
    mapping = ticker_mapping or default_mapping
    tickers = list(mapping.keys())
    # Load data and adjust according to price mode
    df = _load_price_data(excel_path, tickers)
    df_adj, used_date = adjust_prices_for_mode(df, price_mode)
    perf = _build_returns_table(df_adj, mapping)
    heat_df = perf[["Name", "YTD", "1M", "3M", "6M", "12M"]].dropna().sort_values("YTD", ascending=False).reset_index(drop=True)
    n_rows = len(heat_df)
    cols = ["YTD", "1M", "3M", "6M", "12M"]
    color_array = np.zeros((n_rows, len(cols), 4))
    for j, col in enumerate(cols):
        col_vals = heat_df[col].astype(float).values
        pos_vals = col_vals[col_vals > 0]
        neg_vals = col_vals[col_vals < 0]
        max_pos = pos_vals.max() if len(pos_vals) > 0 else 0.0
        min_neg = neg_vals.min() if len(neg_vals) > 0 else 0.0
        for i, val in enumerate(col_vals):
            color_array[i, j] = _colour_for_value(float(val), max_pos, min_neg)
    # Create figure and axis
    fig, ax = plt.subplots(figsize=(width, height))
    ax.imshow(color_array, aspect="auto")
    # Add text
    for i in range(n_rows):
        for j, col in enumerate(cols):
            val = heat_df.iloc[i][col]
            ax.text(
                j,
                i,
                f"{val:.1f}%",
                ha="center",
                va="center",
                fontsize=9,
                fontweight="bold",
                color="black",
            )
    ax.set_xticks(range(len(cols)))
    ax.set_xticklabels(cols, fontsize=12, fontweight="bold")
    ax.xaxis.set_ticks_position("top")
    ax.xaxis.set_label_position("top")
    ax.set_yticks(range(n_rows))
    ax.set_yticklabels(heat_df["Name"], fontsize=9)
    ax.tick_params(axis="y", length=0)
    ax.tick_params(axis="x", length=0)
    for spine in ax.spines.values():
        spine.set_visible(False)
    fig.subplots_adjust(left=0.5)
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=300, bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf.getvalue(), used_date


###############################################################################
# Colour helper
###############################################################################

def _colour_for_value(val: float, max_pos: float, min_neg: float) -> Tuple[float, float, float, float]:
    """Return an RGBA colour for a single value using independent scales.

    Positive values produce a gradient from white to green; negative
    values from white to red.  Zeros (or undefined scales) return
    white.  The input ``max_pos`` should be the maximum positive value
    in the column, and ``min_neg`` the minimum (most negative) value.
    """
    white = np.array([1.0, 1.0, 1.0, 1.0])
    red = np.array(mcolors.to_rgba("#C70039"))
    green = np.array(mcolors.to_rgba("#2E8B57"))
    if val > 0 and max_pos > 0:
        ratio = float(val) / max_pos
        return tuple(white + ratio * (green - white))
    if val < 0 and min_neg < 0:
        ratio = float(val) / min_neg  # both negative → positive ratio
        return tuple(white + ratio * (red - white))
    return tuple(white)


###############################################################################
# Slide insertion helpers
###############################################################################

def _insert_dashboard_to_placeholder(
    prs: Presentation,
    image_bytes: bytes,
    placeholder_names: List[str],
    *,
    left_cm: float,
    top_cm: float,
    width_cm: float,
    height_cm: float,
    used_date: Optional[pd.Timestamp],
    price_mode: str,
    source_placeholder_names: List[str],
) -> Presentation:
    """Helper to insert an image and source footnote into a slide identified by placeholders.

    This internal helper searches for shapes whose names match any of
    ``placeholder_names`` or whose text contains any of those names in
    square brackets.  Once found, it uses the slide but ignores the
    shape's dimensions, inserting the provided image at the specified
    coordinates and size.  It also searches for a text box named
    ``source_placeholder_names`` (or containing the pattern) and
    replaces its contents with a source footnote based on ``used_date``
    and ``price_mode`` while preserving the original formatting.  If no
    matching placeholder is found, the image is inserted on the last
    slide.  Footnote insertion is skipped if ``used_date`` is ``None``.

    Parameters
    ----------
    prs : Presentation
        The presentation into which the image should be inserted.
    image_bytes : bytes
        The PNG data to insert.
    placeholder_names : list of str
        Names of shapes to look for (lowercased).  The function also
        recognises text placeholders like ``[fx_perf_1week]``.
    left_cm, top_cm, width_cm, height_cm : float
        Desired position and size in centimetres.
    used_date : pandas.Timestamp or None
        The effective date for the source footnote.  If ``None`` the
        source is not inserted.
    price_mode : str
        Either ``"Last Price"`` or ``"Last Close"``, used to suffix
        ``" Close"`` in the footnote.
    source_placeholder_names : list of str
        Names or patterns (without brackets) for the source text box.

    Returns
    -------
    Presentation
        The modified presentation.
    """
    target_slide = None
    placeholder_box = None
    # Normalise names for comparison
    name_candidates = [n.lower() for n in placeholder_names]
    pattern_candidates = [f"[{n}]" for n in name_candidates]
    for slide in prs.slides:
        for shape in slide.shapes:
            name_attr = getattr(shape, "name", "").lower()
            if name_attr in name_candidates:
                target_slide = slide
                placeholder_box = shape
                break
            if shape.has_text_frame:
                text_lower = (shape.text or "").strip().lower()
                if text_lower in [p.lower() for p in pattern_candidates]:
                    target_slide = slide
                    placeholder_box = shape
                    break
        if target_slide:
            break
    if target_slide is None:
        target_slide = prs.slides[min(11, len(prs.slides) - 1)]
    # Do not modify or remove the placeholder box; it serves only to locate
    # the slide.  Leave its text intact so the placeholder remains in
    # the template for future reference.
    left = Cm(left_cm)
    top = Cm(top_cm)
    width = Cm(width_cm)
    height = Cm(height_cm)
    stream = io.BytesIO(image_bytes)
    target_slide.shapes.add_picture(stream, left, top, width=width, height=height)
    # Insert source footnote if a date is available
    if used_date is not None:
        date_str = used_date.strftime("%d/%m/%Y")
        source_text = f"Source: Bloomberg, Herculis Group. Data as of {date_str}"
        # Look for source placeholder
        source_candidates = [n.lower() for n in source_placeholder_names]
        source_patterns = [f"[{n}]" for n in source_candidates]
        for shape in target_slide.shapes:
            name_attr2 = getattr(shape, "name", "").lower()
            if name_attr2 in source_candidates:
                if shape.has_text_frame:
                    runs = shape.text_frame.paragraphs[0].runs
                    attrs = _capture_font_attrs(runs[0]) if runs else (None, None, None, None, None, None)
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    new_run = p.add_run()
                    new_run.text = source_text
                    _apply_font_attrs(new_run, *attrs)
                break
            if shape.has_text_frame:
                for pattern in source_patterns:
                    if pattern.lower() in (shape.text or "").lower():
                        runs = shape.text_frame.paragraphs[0].runs
                        attrs = _capture_font_attrs(runs[0]) if runs else (None, None, None, None, None, None)
                        try:
                            new_text = shape.text.replace(pattern, source_text)
                        except Exception:
                            new_text = source_text
                        shape.text_frame.clear()
                        p = shape.text_frame.paragraphs[0]
                        new_run = p.add_run()
                        new_run.text = new_text
                        _apply_font_attrs(new_run, *attrs)
                        break
                else:
                    continue
                break
    return prs


def insert_fx_performance_bar_slide(
    prs: Presentation,
    image_bytes: bytes,
    used_date: Optional[pd.Timestamp] = None,
    price_mode: str = "Last Price",
    *,
    left_cm: float = 0.0,
    top_cm: float = 3.22,
    width_cm: float = 25.0,
    height_cm: float = 11.66,
) -> Presentation:
    """Insert the weekly FX performance bar chart and source footnote into its designated slide.

    The slide is identified by a shape named ``fx_perf_1week`` or
    ``fx_perf_1w`` (or containing ``[fx_perf_1week]`` etc.).  The chart
    is inserted at the specified coordinates regardless of the
    placeholder's original size.  A source footnote is written into
    ``fx_1w_source`` or a text box containing ``[fx_1w_source]`` reflecting
    the ``used_date`` and ``price_mode``.

    Parameters
    ----------
    prs : Presentation
        The PowerPoint presentation to modify.
    image_bytes : bytes
        PNG data for the weekly bar chart.
    used_date : pandas.Timestamp or None, optional
        Effective date used for performance calculations.  If ``None``,
        the source footnote is not inserted.
    price_mode : str, default ``"Last Price"``
        Either ``"Last Price"`` or ``"Last Close"``.  Determines the
        suffix appended to the date in the source footnote.
    left_cm, top_cm, width_cm, height_cm : float
        Position and size for inserting the chart.

    Returns
    -------
    Presentation
        The modified presentation.
    """
    return _insert_dashboard_to_placeholder(
        prs,
        image_bytes,
        placeholder_names=["fx_perf_1w", "fx_perf_1week"],
        left_cm=left_cm,
        top_cm=top_cm,
        width_cm=width_cm,
        height_cm=height_cm,
        used_date=used_date,
        price_mode=price_mode,
        source_placeholder_names=["fx_1w_source"],
    )


def insert_fx_performance_histo_slide(
    prs: Presentation,
    image_bytes: bytes,
    used_date: Optional[pd.Timestamp] = None,
    price_mode: str = "Last Price",
    *,
    left_cm: float = 0.0,
    top_cm: float = 3.22,
    width_cm: float = 25.0,
    height_cm: float = 11.66,
) -> Presentation:
    """Insert the FX historical performance heatmap and source footnote into its designated slide.

    The slide is identified by a shape named ``fx_perf_histo`` (or
    containing ``[fx_perf_histo]``).  The heatmap is inserted at the
    specified coordinates.  A source footnote is written into
    ``fx_1w_source2`` or a text box containing ``[fx_1w_source2]`` based
    on the provided ``used_date`` and ``price_mode``.

    Parameters
    ----------
    prs : Presentation
        The PowerPoint presentation to modify.
    image_bytes : bytes
        PNG data for the heatmap.
    used_date : pandas.Timestamp or None, optional
        Effective date used for performance calculations.  If ``None``,
        the source footnote is not inserted.
    price_mode : str, default ``"Last Price"``
        Either ``"Last Price"`` or ``"Last Close"``.  Determines the
        suffix appended to the date in the source footnote.
    left_cm, top_cm, width_cm, height_cm : float
        Position and size for inserting the heatmap.

    Returns
    -------
    Presentation
        The modified presentation.
    """
    return _insert_dashboard_to_placeholder(
        prs,
        image_bytes,
        placeholder_names=["fx_perf_histo"],
        left_cm=left_cm,
        top_cm=top_cm,
        width_cm=width_cm,
        height_cm=height_cm,
        used_date=used_date,
        price_mode=price_mode,
        source_placeholder_names=["fx_1w_source2"],
    )


# =============================================================================
# HTML-BASED CURRENCY WEEKLY PERFORMANCE CHART
# =============================================================================

# Currency configuration with flag country codes (for flagcdn.com)
CURRENCY_CONFIG = [
    {"ticker": "DXY Curncy", "name": "Dollar Index", "flag": "us"},
    {"ticker": "EURCHF Curncy", "name": "EUR/CHF", "flag": "ch"},
    {"ticker": "EURJPY Curncy", "name": "EUR/JPY", "flag": "jp"},
    {"ticker": "EURMXN Curncy", "name": "EUR/MXN", "flag": "mx"},
    {"ticker": "EURUSD Curncy", "name": "EUR/USD", "flag": "us"},
    {"ticker": "EURBRL Curncy", "name": "EUR/BRL", "flag": "br"},
]

# Chart dimensions (hardcoded)
FX_PNG_WIDTH_PX = 1963
FX_PNG_HEIGHT_PX = 1134
FX_PPT_WIDTH_CM = 17.31
FX_PPT_HEIGHT_CM = 10.0
FX_PPT_LEFT_CM = 3.35
FX_PPT_TOP_CM = 5.5
FX_HTML_SCALE = 3


def _format_fx_percentage(value: float) -> str:
    """Format percentage with sign."""
    if value >= 0:
        return f"+{value:.1f}%"
    else:
        return f"{value:.1f}%"


def create_weekly_html_performance_chart(
    excel_path: Union[str, pathlib.Path],
    *,
    price_mode: str = "Last Price",
) -> Tuple[bytes, Optional[pd.Timestamp]]:
    """Generate HTML-based weekly currency performance bar chart.

    Returns:
        Tuple of (PNG bytes, effective date used for computation)
    """
    # Get all tickers from configuration
    tickers = [c["ticker"] for c in CURRENCY_CONFIG]

    # Load and adjust price data
    df = _load_price_data(excel_path, tickers)
    df_adj, used_date = adjust_prices_for_mode(df, price_mode)
    today = df_adj["Date"].max()

    # Compute 1W returns for all currencies
    returns_dict = {}
    for ticker in tickers:
        try:
            ret_val = _compute_horizon_returns(df_adj, ticker, today, 7)
            returns_dict[ticker] = ret_val
        except Exception:
            returns_dict[ticker] = float("nan")

    # Build row data with returns
    rows_data = []
    for config in CURRENCY_CONFIG:
        ticker = config["ticker"]
        value = returns_dict.get(ticker, 0.0)
        if pd.isna(value):
            value = 0.0
        rows_data.append({
            "name": config["name"],
            "flag": config["flag"],
            "flag_html": get_flag_html(config["flag"]),
            "value": value,
        })

    # Sort by value descending (best performers first)
    rows_data.sort(key=lambda x: x["value"], reverse=True)

    # Calculate max absolute value for bar scaling (minimum 1%)
    valid_values = [abs(r["value"]) for r in rows_data if r["value"] != 0]
    max_abs_value = max(valid_values) if valid_values else 1.0
    max_abs_value = max(max_abs_value, 1.0)  # At least 1%

    # Prepare rows for template with highlight classes
    prepared_rows = []
    for i, row in enumerate(rows_data):
        value = row["value"]
        bar_width = abs(value) / max_abs_value * 48

        # Determine highlight class
        if i == 0:
            highlight_class = "top-performer"
        elif i == len(rows_data) - 1:
            highlight_class = "worst-performer"
        else:
            highlight_class = ""

        prepared_rows.append({
            "name": row["name"],
            "flag": row["flag"],
            "flag_html": get_flag_html(row["flag"]),
            "value": value,
            "bar_class": "positive" if value >= 0 else "negative",
            "bar_width": bar_width,
            "value_class": "positive" if value >= 0 else "negative",
            "formatted_value": _format_fx_percentage(value),
            "highlight_class": highlight_class,
        })

    # Calculate scale values
    scale_values = {
        "scale_min": f"-{max_abs_value:.1f}%",
        "scale_mid_low": f"-{max_abs_value/2:.1f}%",
        "scale_mid_high": f"+{max_abs_value/2:.1f}%",
        "scale_max": f"+{max_abs_value:.1f}%",
    }

    # Render HTML template
    template = Template(CURRENCY_WEEKLY_HTML_TEMPLATE)
    html_content = template.render(
        scale=FX_HTML_SCALE,
        width=FX_PNG_WIDTH_PX,
        height=FX_PNG_HEIGHT_PX,
        rows=prepared_rows,
        **scale_values,
    )

    # Convert HTML to PNG
    hti = Html2Image(size=(FX_PNG_WIDTH_PX, FX_PNG_HEIGHT_PX))
    hti.screenshot(html_str=html_content, save_as="fx_weekly_temp.png")

    # Read the generated PNG
    with open("fx_weekly_temp.png", "rb") as f:
        png_bytes = f.read()

    # Clean up temp file
    import os
    try:
        os.remove("fx_weekly_temp.png")
    except Exception:
        pass

    return png_bytes, used_date


def insert_fx_weekly_html_slide(
    prs: Presentation,
    image_bytes: bytes,
    used_date: Optional[pd.Timestamp] = None,
    price_mode: str = "Last Price",
) -> Presentation:
    """Insert the HTML-based currency weekly performance chart into PowerPoint."""
    if not image_bytes:
        return prs

    # Find target slide by placeholder name
    target_slide = None
    placeholder_names = ["fx_perf_1w", "fx_perf_1week"]
    name_candidates = [n.lower() for n in placeholder_names]
    pattern_candidates = [f"[{n}]" for n in name_candidates]

    for slide in prs.slides:
        for shape in slide.shapes:
            name_attr = getattr(shape, "name", "").lower()
            if name_attr in name_candidates:
                target_slide = slide
                break
            if shape.has_text_frame:
                text_lower = (shape.text or "").strip().lower()
                if text_lower in [p.lower() for p in pattern_candidates]:
                    target_slide = slide
                    break
        if target_slide:
            break

    if target_slide is None:
        print("[FX Weekly HTML] ERROR: Slide not found")
        return prs

    # Insert chart image with exact hardcoded dimensions
    left = Cm(FX_PPT_LEFT_CM)
    top = Cm(FX_PPT_TOP_CM)
    width = Cm(FX_PPT_WIDTH_CM)
    height = Cm(FX_PPT_HEIGHT_CM)

    stream = io.BytesIO(image_bytes)
    picture = target_slide.shapes.add_picture(stream, left, top, width=width, height=height)

    # Send picture to back
    spTree = target_slide.shapes._spTree
    pic_element = picture._element
    spTree.remove(pic_element)
    spTree.insert(2, pic_element)

    print(f"[FX Weekly HTML] Chart inserted at ({FX_PPT_LEFT_CM}, {FX_PPT_TOP_CM}) cm")

    # Update source placeholder if date available
    if used_date is not None:
        date_str = used_date.strftime("%d/%m/%Y")
        source_text = f"Source: Bloomberg, Herculis Group. Data as of {date_str}"
        source_candidates = ["fx_1w_source", "fx_perf_1w_source"]
        source_patterns = [f"[{n}]" for n in source_candidates]

        for shape in target_slide.shapes:
            name_attr = getattr(shape, "name", "").lower()
            if name_attr in [n.lower() for n in source_candidates]:
                if shape.has_text_frame:
                    runs = shape.text_frame.paragraphs[0].runs
                    attrs = _capture_font_attrs(runs[0]) if runs else (None, None, None, None, None, None)
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    new_run = p.add_run()
                    new_run.text = source_text
                    _apply_font_attrs(new_run, *attrs)
                break
            if shape.has_text_frame:
                for pattern in source_patterns:
                    if pattern.lower() in (shape.text or "").lower():
                        runs = shape.text_frame.paragraphs[0].runs
                        attrs = _capture_font_attrs(runs[0]) if runs else (None, None, None, None, None, None)
                        shape.text_frame.clear()
                        p = shape.text_frame.paragraphs[0]
                        new_run = p.add_run()
                        new_run.text = source_text
                        _apply_font_attrs(new_run, *attrs)
                        break
                else:
                    continue
                break

    return prs


# =============================================================================
# HTML-BASED CURRENCY HISTORICAL PERFORMANCE HEATMAP
# =============================================================================

# Chart dimensions for historical heatmap (hardcoded)
FX_HIST_PNG_WIDTH_PX = 1930
FX_HIST_PNG_HEIGHT_PX = 1200
FX_HIST_PPT_WIDTH_CM = 17.02
FX_HIST_PPT_HEIGHT_CM = 10.58
FX_HIST_PPT_LEFT_CM = 3.35
FX_HIST_PPT_TOP_CM = 5.5
FX_HIST_HTML_SCALE = 3


def _get_color_class_for_value(value: float) -> str:
    """Get the CSS color class for a value using fixed thresholds.

    Thresholds: 3%, 8%, 15%, 25%
    """
    abs_val = abs(value)

    if abs_val < 0.5:
        return "neutral"

    if value > 0:
        if abs_val >= 25:
            return "positive-5"
        elif abs_val >= 15:
            return "positive-4"
        elif abs_val >= 8:
            return "positive-3"
        elif abs_val >= 3:
            return "positive-2"
        else:
            return "positive-1"
    else:
        if abs_val >= 25:
            return "negative-5"
        elif abs_val >= 15:
            return "negative-4"
        elif abs_val >= 8:
            return "negative-3"
        elif abs_val >= 3:
            return "negative-2"
        else:
            return "negative-1"


def _format_percentage_with_sign(value: float) -> str:
    """Format percentage value with sign."""
    if pd.isna(value):
        return "N/A"
    if value >= 0:
        return f"+{value:.1f}%"
    else:
        return f"{value:.1f}%"


def create_historical_html_performance_chart(
    excel_path: Union[str, pathlib.Path],
    *,
    price_mode: str = "Last Price",
) -> Tuple[bytes, Optional[pd.Timestamp]]:
    """Generate HTML-based historical currency performance heatmap.

    Returns:
        Tuple of (PNG bytes, effective date used for computation)
    """
    # Get all tickers from configuration
    tickers = [c["ticker"] for c in CURRENCY_CONFIG]

    # Load and adjust price data
    df = _load_price_data(excel_path, tickers)
    df_adj, used_date = adjust_prices_for_mode(df, price_mode)
    today = df_adj["Date"].max()

    # Build row data with returns for all horizons
    rows_data = []
    for config in CURRENCY_CONFIG:
        ticker = config["ticker"]

        # Calculate returns for all horizons
        ret_ytd = float("nan")
        start_of_year = pd.Timestamp(year=today.year, month=1, day=1)
        past_series = df_adj.loc[df_adj["Date"] <= start_of_year, ticker]
        if len(past_series) > 0:
            past_price = past_series.iloc[-1]
            current_price = df_adj[ticker].iloc[-1]
            if past_price and not pd.isna(past_price):
                ret_ytd = (current_price - past_price) / past_price * 100.0

        ret_1m = _compute_horizon_returns(df_adj, ticker, today, 30)
        ret_3m = _compute_horizon_returns(df_adj, ticker, today, 90)
        ret_6m = _compute_horizon_returns(df_adj, ticker, today, 180)
        ret_12m = _compute_horizon_returns(df_adj, ticker, today, 365)

        rows_data.append({
            "name": config["name"],
            "flag": config["flag"],
            "flag_html": get_flag_html(config["flag"]),
            "ytd_value": ret_ytd if not pd.isna(ret_ytd) else 0.0,
            "m1_value": ret_1m if not pd.isna(ret_1m) else 0.0,
            "m3_value": ret_3m if not pd.isna(ret_3m) else 0.0,
            "m6_value": ret_6m if not pd.isna(ret_6m) else 0.0,
            "m12_value": ret_12m if not pd.isna(ret_12m) else 0.0,
        })

    # Sort by YTD descending (best performers first)
    rows_data.sort(key=lambda x: x["ytd_value"], reverse=True)

    # Prepare rows for template
    prepared_rows = []
    for row in rows_data:
        prepared_rows.append({
            "name": row["name"],
            "flag": row["flag"],
            "flag_html": get_flag_html(row["flag"]),
            "ytd_formatted": _format_percentage_with_sign(row["ytd_value"]),
            "ytd_class": _get_color_class_for_value(row["ytd_value"]),
            "m1_formatted": _format_percentage_with_sign(row["m1_value"]),
            "m1_class": _get_color_class_for_value(row["m1_value"]),
            "m3_formatted": _format_percentage_with_sign(row["m3_value"]),
            "m3_class": _get_color_class_for_value(row["m3_value"]),
            "m6_formatted": _format_percentage_with_sign(row["m6_value"]),
            "m6_class": _get_color_class_for_value(row["m6_value"]),
            "m12_formatted": _format_percentage_with_sign(row["m12_value"]),
            "m12_class": _get_color_class_for_value(row["m12_value"]),
        })

    # Render HTML template
    template = Template(CURRENCY_HISTORICAL_HTML_TEMPLATE)
    html_content = template.render(
        scale=FX_HIST_HTML_SCALE,
        width=FX_HIST_PNG_WIDTH_PX,
        height=FX_HIST_PNG_HEIGHT_PX,
        rows=prepared_rows,
    )

    # Convert HTML to PNG
    hti = Html2Image(size=(FX_HIST_PNG_WIDTH_PX, FX_HIST_PNG_HEIGHT_PX))
    hti.screenshot(html_str=html_content, save_as="fx_historical_temp.png")

    # Read the generated PNG
    with open("fx_historical_temp.png", "rb") as f:
        png_bytes = f.read()

    # Clean up temp file
    import os
    try:
        os.remove("fx_historical_temp.png")
    except Exception:
        pass

    return png_bytes, used_date


def insert_fx_historical_html_slide(
    prs: Presentation,
    image_bytes: bytes,
    used_date: Optional[pd.Timestamp] = None,
    price_mode: str = "Last Price",
) -> Presentation:
    """Insert the HTML-based currency historical performance heatmap into PowerPoint."""
    if not image_bytes:
        return prs

    # Find target slide by placeholder name
    target_slide = None
    placeholder_names = ["fx_perf_histo"]
    name_candidates = [n.lower() for n in placeholder_names]
    pattern_candidates = [f"[{n}]" for n in name_candidates]

    for slide in prs.slides:
        for shape in slide.shapes:
            name_attr = getattr(shape, "name", "").lower()
            if name_attr in name_candidates:
                target_slide = slide
                break
            if shape.has_text_frame:
                text_lower = (shape.text or "").strip().lower()
                if text_lower in [p.lower() for p in pattern_candidates]:
                    target_slide = slide
                    break
        if target_slide:
            break

    if target_slide is None:
        print("[FX Historical HTML] ERROR: Slide not found")
        return prs

    # Insert chart image with exact hardcoded dimensions
    left = Cm(FX_HIST_PPT_LEFT_CM)
    top = Cm(FX_HIST_PPT_TOP_CM)
    width = Cm(FX_HIST_PPT_WIDTH_CM)
    height = Cm(FX_HIST_PPT_HEIGHT_CM)

    stream = io.BytesIO(image_bytes)
    picture = target_slide.shapes.add_picture(stream, left, top, width=width, height=height)

    # Send picture to back
    spTree = target_slide.shapes._spTree
    pic_element = picture._element
    spTree.remove(pic_element)
    spTree.insert(2, pic_element)

    print(f"[FX Historical HTML] Heatmap inserted at ({FX_HIST_PPT_LEFT_CM}, {FX_HIST_PPT_TOP_CM}) cm")

    # Update source placeholder if date available
    if used_date is not None:
        date_str = used_date.strftime("%d/%m/%Y")
        source_text = f"Source: Bloomberg, Herculis Group. Data as of {date_str}"
        source_candidates = ["fx_1w_source2", "fx_hist_source"]
        source_patterns = [f"[{n}]" for n in source_candidates]

        for shape in target_slide.shapes:
            name_attr = getattr(shape, "name", "").lower()
            if name_attr in [n.lower() for n in source_candidates]:
                if shape.has_text_frame:
                    runs = shape.text_frame.paragraphs[0].runs
                    attrs = _capture_font_attrs(runs[0]) if runs else (None, None, None, None, None, None)
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    new_run = p.add_run()
                    new_run.text = source_text
                    _apply_font_attrs(new_run, *attrs)
                break
            if shape.has_text_frame:
                for pattern in source_patterns:
                    if pattern.lower() in (shape.text or "").lower():
                        runs = shape.text_frame.paragraphs[0].runs
                        attrs = _capture_font_attrs(runs[0]) if runs else (None, None, None, None, None, None)
                        shape.text_frame.clear()
                        p = shape.text_frame.paragraphs[0]
                        new_run = p.add_run()
                        new_run.text = source_text
                        _apply_font_attrs(new_run, *attrs)
                        break
                else:
                    continue
                break

    return prs
