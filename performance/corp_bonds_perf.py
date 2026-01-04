"""Corporate Bonds Weekly Performance chart generation.

This module generates a weekly performance bar chart for corporate bonds,
showing Investment Grade (IG) and High Yield (HY) bonds sorted by performance.
Each bond displays a flag, credit badge (IG/HY), and performance bar.

Functions
---------
``create_weekly_performance_chart``
    Generate the corporate bonds weekly performance PNG.
``insert_corp_bonds_performance_slide``
    Insert the chart into a PowerPoint slide.
"""

from __future__ import annotations

import io
import tempfile
from dataclasses import dataclass
from pathlib import Path
from typing import Dict, List, Optional, Tuple, Union

import pandas as pd
from jinja2 import Template
from html2image import Html2Image
from pptx import Presentation
from pptx.util import Cm

from utils import adjust_prices_for_mode
from market_compass.weekly_performance.html_template import (
    CORP_BONDS_WEEKLY_HTML_TEMPLATE,
    CORP_BONDS_HISTORICAL_HTML_TEMPLATE,
)


# Scale factor for PNG generation
SCALE_FACTOR = 3

# Corporate bond configuration - uses Bloomberg index tickers from data_prices sheet
# Maps ticker to (display_name, flag, credit_type)
CORP_BOND_CONFIG = {
    # Investment Grade
    "LUACTRUU Index": {"name": "US IG Corp", "flag": "\U0001F1FA\U0001F1F8", "credit": "IG"},
    "LECPTREU Index": {"name": "EUR IG Corp", "flag": "\U0001F1EA\U0001F1FA", "credit": "IG"},
    "JPEIESGE Index": {"name": "EM IG Corp", "flag": "\U0001F30D", "credit": "IG"},
    # High Yield
    "LF98TRUU Index": {"name": "US HY Corp", "flag": "\U0001F1FA\U0001F1F8", "credit": "HY"},
    "IBOXXMJA Index": {"name": "EUR HY Corp", "flag": "\U0001F1EA\U0001F1FA", "credit": "HY"},
    "JPEIEMHY Index": {"name": "EM HY Corp", "flag": "\U0001F30D", "credit": "HY"},
}


@dataclass
class CorpBondRow:
    """Data class for a corporate bond row."""
    name: str           # e.g., "US IG Corp"
    flag: str           # e.g., "🇺🇸"
    credit_type: str    # "IG" or "HY"
    value: float        # Weekly return as decimal (e.g., 0.0032 for 0.32%)


@dataclass
class CorpBondHistoricalRow:
    """Data class for a corporate bond historical row."""
    name: str           # e.g., "US IG Corp"
    flag: str           # e.g., "🇺🇸"
    credit_type: str    # "IG" or "HY"
    ytd: float          # YTD return as decimal
    m1: float           # 1M return
    m3: float           # 3M return
    m6: float           # 6M return
    m12: float          # 12M return


def _get_color_class(value: float) -> str:
    """Determine color class based on value magnitude.

    Thresholds (for bonds - typically smaller moves):
    - level 1: 0-2%
    - level 2: 2-5%
    - level 3: 5-10%
    - level 4: 10-15%
    - level 5: >15%
    """
    abs_val = abs(value)
    prefix = "positive" if value >= 0 else "negative"

    if abs_val <= 0.02:
        level = 1
    elif abs_val <= 0.05:
        level = 2
    elif abs_val <= 0.10:
        level = 3
    elif abs_val <= 0.15:
        level = 4
    else:
        level = 5

    return f"{prefix}-{level}"


def _format_percentage(value: float) -> str:
    """Format value as percentage string."""
    if pd.isna(value):
        return "N/A"
    return f"{value * 100:.1f}%"


def _load_price_data(
    excel_path: Union[str, Path],
    tickers: List[str],
) -> pd.DataFrame:
    """Load price data for the specified tickers.

    Parameters
    ----------
    excel_path : str or Path
        Path to the Excel workbook containing ``data_prices`` sheet.
    tickers : list of str
        Column names for the desired bonds.

    Returns
    -------
    DataFrame
        A DataFrame with ``Date`` column and one column per ticker.
    """
    df = pd.read_excel(excel_path, sheet_name="data_prices")
    df = df.drop(index=0)
    df = df[df[df.columns[0]] != "DATES"].copy()
    df = df.rename(columns={df.columns[0]: "Date"})
    df["Date"] = pd.to_datetime(df["Date"], errors="coerce")

    # Get available tickers
    available_tickers = [t for t in tickers if t in df.columns]
    if not available_tickers:
        return pd.DataFrame(columns=["Date"])

    out = df[["Date"] + available_tickers].copy()
    for t in available_tickers:
        out[t] = pd.to_numeric(out[t], errors="coerce")

    return out.dropna(subset=["Date"]).sort_values("Date").reset_index(drop=True)


def _compute_weekly_return(
    df: pd.DataFrame,
    ticker: str,
) -> float:
    """Compute 1-week return for a ticker.

    Parameters
    ----------
    df : DataFrame
        Price data with Date column.
    ticker : str
        Column name for the ticker.

    Returns
    -------
    float
        Weekly return as decimal, or NaN if not computable.
    """
    if ticker not in df.columns:
        return float("nan")

    today = df["Date"].max()
    past_date = today - pd.Timedelta(days=7)
    past_series = df.loc[df["Date"] <= past_date, ticker]

    if len(past_series) == 0:
        return float("nan")

    past_price = past_series.iloc[-1]
    current_price = df[ticker].iloc[-1]

    if past_price == 0 or pd.isna(past_price) or pd.isna(current_price):
        return float("nan")

    return (current_price - past_price) / past_price


def create_weekly_performance_chart(
    excel_path: Union[str, Path],
    *,
    price_mode: str = "Last Price",
    width_cm: float = 17.31,
    height_cm: float = 10.0,
) -> Tuple[bytes, Optional[pd.Timestamp]]:
    """Generate the Corporate Bonds Weekly Performance chart.

    Parameters
    ----------
    excel_path : str or Path
        Path to the Excel workbook.
    price_mode : str, default "Last Price"
        Either "Last Price" or "Last Close".
    width_cm, height_cm : float
        Dimensions for the chart in centimeters.

    Returns
    -------
    tuple
        (png_bytes, used_date) where png_bytes is the chart image and
        used_date is the effective date used.
    """
    tickers = list(CORP_BOND_CONFIG.keys())

    # Load and adjust prices
    df = _load_price_data(excel_path, tickers)
    df_adj, used_date = adjust_prices_for_mode(df, price_mode)

    # Build rows with weekly returns
    rows: List[CorpBondRow] = []

    for ticker, config in CORP_BOND_CONFIG.items():
        if ticker not in df_adj.columns:
            print(f"[Corp Bonds Weekly] DEBUG: Ticker {ticker} not found in data")
            continue

        weekly_ret = _compute_weekly_return(df_adj, ticker)
        if pd.isna(weekly_ret):
            print(f"[Corp Bonds Weekly] DEBUG: Ticker {ticker} has NaN weekly return")
            continue

        rows.append(CorpBondRow(
            name=config["name"],
            flag=config["flag"],
            credit_type=config["credit"],
            value=weekly_ret,
        ))
        print(f"[Corp Bonds Weekly] DEBUG: {config['name']}: {weekly_ret * 100:.2f}%")

    print(f"[Corp Bonds Weekly] DEBUG: Total bonds found: {len(rows)}")

    if not rows:
        # Return empty image if no data
        return b"", used_date

    # Sort by value descending (best to worst)
    rows.sort(key=lambda x: x.value, reverse=True)

    # Calculate max for bar scaling
    max_abs_value = max(abs(r.value) for r in rows)
    max_abs_value = max(max_abs_value, 0.005)  # At least 0.5%

    # Prepare rows for template
    prepared_rows = []
    for i, row in enumerate(rows):
        # Bar width as percentage of half the chart
        bar_width = abs(row.value) / max_abs_value * 50
        bar_width = min(bar_width, 48)

        # Determine highlight class
        if i == 0 and row.value > 0:
            highlight_class = "top-performer"
        elif i == len(rows) - 1 and row.value < 0:
            highlight_class = "worst-performer"
        else:
            highlight_class = ""

        prepared_rows.append({
            "name": row.name,
            "flag": row.flag,
            "credit_type": row.credit_type,
            "credit_class": "ig" if row.credit_type == "IG" else "hy",
            "value": row.value,
            "bar_class": "positive" if row.value >= 0 else "negative",
            "bar_width": bar_width,
            "value_class": "positive" if row.value >= 0 else "negative",
            "formatted_value": f"+{row.value * 100:.2f}%" if row.value >= 0 else f"{row.value * 100:.2f}%",
            "highlight_class": highlight_class,
        })

    # Calculate scale values
    scale_pct = max_abs_value * 100
    if scale_pct <= 0.5:
        scale_max = 0.5
    elif scale_pct <= 1:
        scale_max = 1
    elif scale_pct <= 2:
        scale_max = 2
    elif scale_pct <= 5:
        scale_max = 5
    else:
        scale_max = round(scale_pct) + 1

    scale_values = {
        "scale_min": f"-{scale_max:.1f}%",
        "scale_mid_low": f"-{scale_max / 2:.1f}%",
        "scale_mid_high": f"+{scale_max / 2:.1f}%",
        "scale_max": f"+{scale_max:.1f}%",
    }

    # Generate HTML - use EXACT hardcoded dimensions
    # DO NOT recalculate these values
    PNG_WIDTH_PX = 1963
    PNG_HEIGHT_PX = 1134

    template = Template(CORP_BONDS_WEEKLY_HTML_TEMPLATE)
    html = template.render(
        rows=prepared_rows,
        width=PNG_WIDTH_PX,
        height=PNG_HEIGHT_PX,
        scale=SCALE_FACTOR,
        **scale_values,
    )

    # Convert to PNG at EXACT size
    with tempfile.TemporaryDirectory() as tmpdir:
        hti = Html2Image(output_path=tmpdir, size=(1963, 1134))
        hti.screenshot(html_str=html, save_as="corp_bonds_weekly.png")

        img_path = Path(tmpdir) / "corp_bonds_weekly.png"
        with open(img_path, "rb") as f:
            png_bytes = f.read()

    return png_bytes, used_date


def insert_corp_bonds_performance_slide(
    prs: Presentation,
    image_bytes: bytes,
    used_date: Optional[pd.Timestamp] = None,
    price_mode: str = "Last Price",
    *,
    left_cm: float = 3.35,
    top_cm: float = 4.6,
    width_cm: float = 17.31,
    height_cm: float = 10.0,
) -> Presentation:
    """Insert the Corporate Bonds Weekly Performance chart into a slide.

    The slide is identified by a shape named ``corp_bonds_perf_1w`` or
    containing ``[corp_bonds_perf_1w]``.

    Parameters
    ----------
    prs : Presentation
        PowerPoint presentation to modify.
    image_bytes : bytes
        PNG data for the chart.
    used_date : Timestamp, optional
        Effective date for the source footnote.
    price_mode : str
        Either "Last Price" or "Last Close".
    left_cm, top_cm, width_cm, height_cm : float
        Position and size for the chart.

    Returns
    -------
    Presentation
        The modified presentation.
    """
    # Find target slide
    target_slide = None
    placeholder_names = ["credit_perf_1w", "corp_bonds_perf_1w", "corp_bonds_weekly"]
    name_candidates = [n.lower() for n in placeholder_names]
    pattern_candidates = [f"[{n}]" for n in name_candidates]

    for slide in prs.slides:
        for shape in slide.shapes:
            name_attr = getattr(shape, "name", "").lower()
            if name_attr in name_candidates:
                target_slide = slide
                break
            if shape.has_text_frame:
                text_lower = (shape.text or "").strip().lower()
                if text_lower in [p.lower() for p in pattern_candidates]:
                    target_slide = slide
                    break
        if target_slide:
            break

    if target_slide is None:
        print("[Corp Bonds Weekly] WARNING: Slide not found")
        return prs

    # Insert picture at EXACT hardcoded PowerPoint dimensions
    # Corporate Bonds - adjusted position (fewer rows = lower on slide)
    # DO NOT modify these values
    stream = io.BytesIO(image_bytes)
    pic = target_slide.shapes.add_picture(
        stream,
        left=Cm(3.21),
        top=Cm(5.81),
        width=Cm(17.31),
    )

    # Send to back
    spTree = target_slide.shapes._spTree
    sp = pic._element
    spTree.remove(sp)
    spTree.insert(2, sp)

    # Insert source footnote if date available
    if used_date is not None:
        date_str = used_date.strftime("%d/%m/%Y")
        suffix = " Close" if price_mode.lower() == "last close" else ""
        source_text = f"Source: Bloomberg, Herculis Group, Data as of {date_str}{suffix}"

        source_names = ["credit_1w_source", "corp_bonds_1w_source"]
        source_candidates = [n.lower() for n in source_names]
        source_patterns = [f"[{n}]" for n in source_candidates]

        for shape in target_slide.shapes:
            name_attr = getattr(shape, "name", "").lower()
            if name_attr in source_candidates and shape.has_text_frame:
                tf = shape.text_frame
                if tf.paragraphs:
                    tf.paragraphs[0].runs[0].text = source_text if tf.paragraphs[0].runs else source_text
                break
            if shape.has_text_frame:
                for pattern in source_patterns:
                    if pattern.lower() in (shape.text or "").lower():
                        shape.text_frame.paragraphs[0].runs[0].text = source_text
                        break

    print(f"[Corp Bonds Weekly] Chart inserted")
    return prs


# =============================================================================
# CORPORATE BONDS HISTORICAL PERFORMANCE HEATMAP
# =============================================================================


def _compute_horizon_return(
    df: pd.DataFrame,
    ticker: str,
    today: pd.Timestamp,
    days: int,
) -> float:
    """Compute return over a specific horizon.

    Parameters
    ----------
    df : DataFrame
        Price data with Date column.
    ticker : str
        Column name for the ticker.
    today : Timestamp
        Current date.
    days : int
        Number of days to look back.

    Returns
    -------
    float
        Return as decimal, or NaN if not computable.
    """
    if ticker not in df.columns:
        return float("nan")

    past_date = today - pd.Timedelta(days=days)
    past_series = df.loc[df["Date"] <= past_date, ticker]

    if len(past_series) == 0:
        return float("nan")

    past_price = past_series.iloc[-1]
    current_price = df.loc[df["Date"] == today, ticker]

    if len(current_price) == 0:
        current_price = df[ticker].iloc[-1]
    else:
        current_price = current_price.iloc[0]

    if past_price == 0 or pd.isna(past_price) or pd.isna(current_price):
        return float("nan")

    return (current_price - past_price) / past_price


def create_historical_performance_chart(
    excel_path: Union[str, Path],
    *,
    price_mode: str = "Last Price",
) -> Tuple[bytes, Optional[pd.Timestamp]]:
    """Generate the Corporate Bonds Historical Performance heatmap.

    Parameters
    ----------
    excel_path : str or Path
        Path to the Excel workbook.
    price_mode : str, default "Last Price"
        Either "Last Price" or "Last Close".

    Returns
    -------
    tuple
        (png_bytes, used_date) where png_bytes is the chart image and
        used_date is the effective date used.
    """
    tickers = list(CORP_BOND_CONFIG.keys())

    # Load and adjust prices
    df = _load_price_data(excel_path, tickers)
    df_adj, used_date = adjust_prices_for_mode(df, price_mode)

    if df_adj.empty or "Date" not in df_adj.columns:
        return b"", used_date

    today = df_adj["Date"].max()

    # Get YTD start date (use data year, not system year)
    data_year = today.year
    ytd_start = pd.Timestamp(year=data_year, month=1, day=1)

    # Build rows with historical returns
    rows: List[CorpBondHistoricalRow] = []

    for ticker, config in CORP_BOND_CONFIG.items():
        if ticker not in df_adj.columns:
            print(f"[Corp Bonds Historical] DEBUG: Ticker {ticker} not found in data")
            continue

        # Compute returns for each horizon
        m1 = _compute_horizon_return(df_adj, ticker, today, 30)
        m3 = _compute_horizon_return(df_adj, ticker, today, 90)
        m6 = _compute_horizon_return(df_adj, ticker, today, 180)
        m12 = _compute_horizon_return(df_adj, ticker, today, 365)

        # Compute YTD
        ytd_series = df_adj.loc[df_adj["Date"] <= ytd_start, ticker]
        if len(ytd_series) > 0:
            ytd_price = ytd_series.iloc[-1]
            current_price = df_adj[ticker].iloc[-1]
            if ytd_price and not pd.isna(ytd_price) and not pd.isna(current_price):
                ytd = (current_price - ytd_price) / ytd_price
            else:
                ytd = float("nan")
        else:
            ytd = float("nan")

        rows.append(CorpBondHistoricalRow(
            name=config["name"],
            flag=config["flag"],
            credit_type=config["credit"],
            ytd=ytd,
            m1=m1,
            m3=m3,
            m6=m6,
            m12=m12,
        ))
        print(f"[Corp Bonds Historical] DEBUG: {config['name']}: YTD={ytd * 100:.1f}%")

    print(f"[Corp Bonds Historical] DEBUG: Total bonds found: {len(rows)}")

    if not rows:
        return b"", used_date

    # Sort by YTD descending (best to worst)
    rows.sort(key=lambda x: x.ytd if not pd.isna(x.ytd) else float("-inf"), reverse=True)

    # Prepare rows for template
    prepared_rows = []
    for row in rows:
        prepared_rows.append({
            "name": row.name,
            "flag": row.flag,
            "credit_type": row.credit_type,
            "credit_class": "ig" if row.credit_type == "IG" else "hy",
            "ytd_formatted": _format_percentage(row.ytd),
            "ytd_class": _get_color_class(row.ytd) if not pd.isna(row.ytd) else "neutral",
            "m1_formatted": _format_percentage(row.m1),
            "m1_class": _get_color_class(row.m1) if not pd.isna(row.m1) else "neutral",
            "m3_formatted": _format_percentage(row.m3),
            "m3_class": _get_color_class(row.m3) if not pd.isna(row.m3) else "neutral",
            "m6_formatted": _format_percentage(row.m6),
            "m6_class": _get_color_class(row.m6) if not pd.isna(row.m6) else "neutral",
            "m12_formatted": _format_percentage(row.m12),
            "m12_class": _get_color_class(row.m12) if not pd.isna(row.m12) else "neutral",
        })

    # Generate HTML - use hardcoded dimensions for 6 rows
    # Shorter aspect ratio for fewer rows
    PNG_WIDTH_PX = 1930
    PNG_HEIGHT_PX = 800

    template = Template(CORP_BONDS_HISTORICAL_HTML_TEMPLATE)
    html = template.render(
        rows=prepared_rows,
        width=PNG_WIDTH_PX,
        height=PNG_HEIGHT_PX,
        scale=SCALE_FACTOR,
    )

    # Convert to PNG
    with tempfile.TemporaryDirectory() as tmpdir:
        hti = Html2Image(output_path=tmpdir, size=(PNG_WIDTH_PX, PNG_HEIGHT_PX))
        hti.screenshot(html_str=html, save_as="corp_bonds_historical.png")

        img_path = Path(tmpdir) / "corp_bonds_historical.png"
        with open(img_path, "rb") as f:
            png_bytes = f.read()

    return png_bytes, used_date


def insert_corp_bonds_historical_slide(
    prs: Presentation,
    image_bytes: bytes,
    used_date: Optional[pd.Timestamp] = None,
    price_mode: str = "Last Price",
) -> Presentation:
    """Insert the Corporate Bonds Historical Performance heatmap into a slide.

    The slide is identified by a shape named ``credit_perf_histo``.

    Parameters
    ----------
    prs : Presentation
        PowerPoint presentation to modify.
    image_bytes : bytes
        PNG data for the chart.
    used_date : Timestamp, optional
        Effective date for the source footnote.
    price_mode : str
        Either "Last Price" or "Last Close".

    Returns
    -------
    Presentation
        The modified presentation.
    """
    if not image_bytes:
        print("[Corp Bonds Historical] WARNING: No image data to insert")
        return prs

    # Find target slide
    target_slide = None
    placeholder_names = ["credit_perf_histo", "corp_bonds_historical", "credit_historical"]
    name_candidates = [n.lower() for n in placeholder_names]
    pattern_candidates = [f"[{n}]" for n in name_candidates]

    for slide in prs.slides:
        for shape in slide.shapes:
            name_attr = getattr(shape, "name", "").lower()
            if name_attr in name_candidates:
                target_slide = slide
                break
            if shape.has_text_frame:
                text_lower = (shape.text or "").strip().lower()
                if text_lower in [p.lower() for p in pattern_candidates]:
                    target_slide = slide
                    break
        if target_slide:
            break

    if target_slide is None:
        print("[Corp Bonds Historical] WARNING: Slide not found")
        return prs

    # Insert picture at hardcoded PowerPoint dimensions
    # width=17.02cm, left=3.35cm, top=4.6cm
    stream = io.BytesIO(image_bytes)
    pic = target_slide.shapes.add_picture(
        stream,
        left=Cm(3.35),
        top=Cm(4.6),
        width=Cm(17.02),
    )

    # Send to back
    spTree = target_slide.shapes._spTree
    sp = pic._element
    spTree.remove(sp)
    spTree.insert(2, sp)

    # Insert source footnote if date available
    if used_date is not None:
        date_str = used_date.strftime("%d/%m/%Y")
        suffix = " Close" if price_mode.lower() == "last close" else ""
        source_text = f"Source: Bloomberg, Herculis Group, Data as of {date_str}{suffix}"

        source_names = ["credit_histo_source", "corp_bonds_histo_source"]
        source_candidates = [n.lower() for n in source_names]
        source_patterns = [f"[{n}]" for n in source_candidates]

        for shape in target_slide.shapes:
            name_attr = getattr(shape, "name", "").lower()
            if name_attr in source_candidates and shape.has_text_frame:
                tf = shape.text_frame
                if tf.paragraphs:
                    tf.paragraphs[0].runs[0].text = source_text if tf.paragraphs[0].runs else source_text
                break
            if shape.has_text_frame:
                for pattern in source_patterns:
                    if pattern.lower() in (shape.text or "").lower():
                        shape.text_frame.paragraphs[0].runs[0].text = source_text
                        break

    print(f"[Corp Bonds Historical] Heatmap inserted")
    return prs
