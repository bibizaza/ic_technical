"""Rates performance dashboard generation with price‑mode awareness and source footnotes.

This module produces charts summarising recent yield movements for a selection
of government bond rates.  Unlike price‑based instruments, yield increases
correspond to price declines (and vice versa).  Therefore, positive changes
in yields are displayed as negative performance (red bars/cells) and
negative changes as positive performance (green).  Users can choose
to base calculations on either the most recent intraday yield ("Last
Price") or the previous day's closing yield ("Last Close").  The
resulting figures include a weekly change bar chart and a heatmap of
longer horizons (1M, 3M, 6M, 12M and YTD).  When the charts are
inserted into a PowerPoint presentation, a source footnote is added to
the designated text boxes on each slide, reflecting the effective
date used in the computations and the chosen price mode.  Text
formatting from the placeholders is preserved.

Key functions
-------------

* ``create_weekly_performance_chart`` – Build the weekly bar chart and
  return both the image bytes and the effective date used for
  computation.
* ``create_historical_performance_table`` – Build the heatmap of
  yield changes for multiple horizons and return the image bytes and
  the effective date.
* ``insert_rates_performance_bar_slide`` – Insert the weekly bar
  chart and source footnote into its slide.
* ``insert_rates_performance_histo_slide`` – Insert the historical
  performance heatmap and source footnote into its slide.

Usage example::

    from performance.rates_perf import (
        create_weekly_performance_chart,
        create_historical_performance_table,
        insert_rates_performance_bar_slide,
        insert_rates_performance_histo_slide,
    )

    bar_bytes, used_date = create_weekly_performance_chart(excel_path, price_mode="Last Close")
    histo_bytes, _ = create_historical_performance_table(excel_path, price_mode="Last Close")
    prs = insert_rates_performance_bar_slide(prs, bar_bytes, used_date, price_mode)
    prs = insert_rates_performance_histo_slide(prs, histo_bytes, used_date, price_mode)

"""

from __future__ import annotations

import io
import pathlib
import shutil
import tempfile
from pathlib import Path
from typing import Dict, List, Tuple, Union, Optional

import matplotlib.pyplot as plt
import matplotlib.colors as mcolors
import numpy as np
import pandas as pd
from jinja2 import Template
from html2image import Html2Image
from utils.flag_utils import get_flag_html
from pptx import Presentation
from pptx.util import Cm

from utils import adjust_prices_for_mode
from market_compass.weekly_performance.html_template import BONDS_RATES_HTML_TEMPLATE, BONDS_HISTORICAL_HTML_TEMPLATE

# Scale factor for high-resolution rendering
SCALE_FACTOR = 4

# Country configuration for bonds rates chart
# Flag codes are ISO 3166-1 alpha-2 for flagcdn.com
COUNTRY_CONFIG = [
    {
        "name": "United States",
        "flag": "us",
        "tickers": {
            "USGG2YR Index": "2Y",
            "USGG10YR Index": "10Y",
            "USGG30YR Index": "30Y",
        }
    },
    {
        "name": "Eurozone",
        "flag": "eu",
        "tickers": {
            "GECU2YR Index": "2Y",
            "GECU10YR Index": "10Y",
            "GECU30YR Index": "30Y",
        }
    },
    {
        "name": "Japan",
        "flag": "jp",
        "tickers": {
            "GJGB2 Index": "2Y",
            "GJGB10 Index": "10Y",
            "GJGB30 Index": "30Y",
        }
    },
    {
        "name": "China",
        "flag": "cn",
        "tickers": {
            "GCNY2YR Index": "2Y",
            "GCNY10YR Index": "10Y",
            "GCNY30YR Index": "30Y",
        }
    },
]

try:
    # Reuse font attribute helpers from the SPX module if available
    from technical_analysis.equity.spx import (
        _get_run_font_attributes as _capture_font_attrs,  # type: ignore
        _apply_run_font_attributes as _apply_font_attrs,  # type: ignore
    )
except Exception:
    # Fallback helpers: only support basic size and RGB; ignore theme colours
    def _capture_font_attrs(run):  # type: ignore
        if run is None:
            return (None, None, None, None, None, None)
        size = run.font.size
        colour = run.font.color
        rgb = None
        try:
            rgb = colour.rgb
        except Exception:
            rgb = None
        return (size, rgb, None, None, run.font.bold, run.font.italic)

    def _apply_font_attrs(new_run, size, rgb, theme_color, brightness, bold, italic):  # type: ignore
        if size is not None:
            new_run.font.size = size
        if rgb is not None:
            try:
                new_run.font.color.rgb = rgb
            except Exception:
                pass
        if bold is not None:
            new_run.font.bold = bold
        if italic is not None:
            new_run.font.italic = italic

###############################################################################
# Data loading and return calculations
###############################################################################

def _load_price_data(excel_path: Union[str, pathlib.Path], tickers: List[str]) -> pd.DataFrame:
    """Load a DataFrame of dates and yields for the specified rate tickers.

    Parameters
    ----------
    excel_path : str or pathlib.Path
        Path to the Excel workbook containing a sheet named ``data_prices``.
    tickers : list of str
        Column names in the sheet corresponding to the desired rate indices.

    Returns
    -------
    DataFrame
        A tidy DataFrame with columns ``Date`` and one column per ticker,
        sorted by date and with numeric yields.
    """
    df = pd.read_excel(excel_path, sheet_name="data_prices")
    df = df.drop(index=0)
    df = df[df[df.columns[0]] != "DATES"].copy()
    df = df.rename(columns={df.columns[0]: "Date"})
    df["Date"] = pd.to_datetime(df["Date"], errors="coerce")
    out = df[["Date"] + tickers].copy()
    for t in tickers:
        out[t] = pd.to_numeric(out[t], errors="coerce")
    return out.dropna(subset=["Date"]).sort_values("Date").reset_index(drop=True)


def _compute_horizon_returns(
    df: pd.DataFrame,
    ticker: str,
    today: pd.Timestamp,
    days: int,
) -> float:
    """Compute the change in yield over the given lookback, expressed in basis points.

    Unlike price series, yield changes are computed as simple differences
    rather than percentage changes.  For example, a move from 1.00 %
    to 1.25 % over a one‑week horizon is reported as +25 bps (positive
    values correspond to rising yields and thus negative price action).

    Parameters
    ----------
    df : DataFrame
        A tidy DataFrame with ``Date`` and yield columns.
    ticker : str
        Column name for which to compute the return.
    today : Timestamp
        Anchor date; typically the last date in the dataset.
    days : int
        Lookback period in days.

    Returns
    -------
    float
        Yield change over the lookback horizon in basis points (bps), or
        NaN if not computable.
    """
    past_date = today - pd.Timedelta(days=days)
    past_series = df.loc[df["Date"] <= past_date, ticker]
    if len(past_series) == 0:
        return float("nan")
    past_yield = past_series.iloc[-1]
    current_yield = df[ticker].iloc[-1]
    if pd.isna(past_yield) or pd.isna(current_yield):
        return float("nan")
    # The raw yield values in the Excel file are expressed in percentage points
    # (e.g. 1.941485 represents 1.941485%).  A one‑basis‑point move equals
    # 0.01 percentage points.  Therefore, to convert the difference between
    # two percentage yields into basis points, multiply by 100.  This yields
    # e.g. (1.94 – 1.84) * 100 ≈ 9.7 bps.  Do not multiply by 10,000, which
    # would treat the inputs as decimals rather than percentages.
    return (current_yield - past_yield) * 100.0  # convert to basis points


def _build_returns_table(
    df: pd.DataFrame,
    mapping: Dict[str, str],
) -> pd.DataFrame:
    """Calculate yield changes (in bps) for each ticker across multiple horizons.

    Returns a DataFrame indexed by the ticker code with the following
    columns: ``Name``, ``1W``, ``1M``, ``3M``, ``6M``, ``12M`` and ``YTD``.

    Parameters
    ----------
    df : DataFrame
        A tidy DataFrame as returned by ``_load_price_data``.
    mapping : dict
        Mapping from ticker codes to human‑friendly names.

    Returns
    -------
    DataFrame
        Table of yield changes in basis points.
    """
    today = df["Date"].max()
    results: Dict[str, Dict[str, float]] = {}
    for ticker in mapping:
        results[ticker] = {
            "1W": _compute_horizon_returns(df, ticker, today, 7),
            "1M": _compute_horizon_returns(df, ticker, today, 30),
            "3M": _compute_horizon_returns(df, ticker, today, 90),
            "6M": _compute_horizon_returns(df, ticker, today, 180),
            "12M": _compute_horizon_returns(df, ticker, today, 365),
        }
        # Year‑to‑date: yield change since the start of the calendar year
        start_of_year = pd.Timestamp(year=today.year, month=1, day=1)
        past_series = df.loc[df["Date"] <= start_of_year, ticker]
        if len(past_series) > 0:
            past_yield = past_series.iloc[-1]
            current_yield = df[ticker].iloc[-1]
            if pd.notna(past_yield) and pd.notna(current_yield):
                # Convert the yield difference (expressed in percentage points) to bps.
                # See _compute_horizon_returns for rationale: multiply by 100.
                results[ticker]["YTD"] = (current_yield - past_yield) * 100.0
            else:
                results[ticker]["YTD"] = float("nan")
        else:
            results[ticker]["YTD"] = float("nan")
    table = pd.DataFrame.from_dict(results, orient="index")
    table["Name"] = table.index.map(mapping.get)
    return table[["Name", "1W", "1M", "3M", "6M", "12M", "YTD"]]


###############################################################################
# Figure generation functions
###############################################################################

def _calculate_rates_scale(max_abs_value: float) -> dict:
    """Calculate nice scale values for rates chart (in bps)."""
    if max_abs_value <= 5:
        scale_max = 5
    elif max_abs_value <= 10:
        scale_max = 10
    elif max_abs_value <= 20:
        scale_max = 20
    elif max_abs_value <= 50:
        scale_max = 50
    else:
        scale_max = int(max_abs_value / 10) * 10 + 10

    return {
        "scale_min": f"-{scale_max} bps",
        "scale_mid_low": f"-{scale_max // 2} bps",
        "scale_mid_high": f"+{scale_max // 2} bps",
        "scale_max": f"+{scale_max} bps",
    }


def create_weekly_performance_chart(
    excel_path: Union[str, pathlib.Path],
    ticker_mapping: Dict[str, str] | None = None,
    *,
    width_cm: float = 17.02,
    height_cm: float = 13.0,
    price_mode: str = "Last Price",
) -> Tuple[bytes, Optional[pd.Timestamp]]:
    """Generate a bar chart of 1‑week yield changes using HTML template.

    This helper reads the specified Excel file, optionally adjusts the yield
    history according to the selected ``price_mode`` (``"Last Price"`` vs
    ``"Last Close"``), computes 1‑week changes for the configured rate
    tickers grouped by country (US, EUR, JP, CN).

    Uses INVERTED color logic for bond rates:
    - Positive change (rising yields) = RED (bearish)
    - Negative change (falling yields) = GREEN (bullish)

    Parameters
    ----------
    excel_path : str or pathlib.Path
        Path to the Excel workbook containing the yield series.
    ticker_mapping : dict, optional
        Not used in new implementation, kept for API compatibility.
    width_cm, height_cm : float
        Dimensions of the generated figure in centimeters.
    price_mode : str, default ``"Last Price"``
        One of ``"Last Price"`` or ``"Last Close"``.  Determines whether
        intraday yields are used or the most recent row is dropped if it
        corresponds to today's date.

    Returns
    -------
    tuple
        A two‑tuple ``(image_bytes, used_date)`` where ``image_bytes`` is
        the PNG data for the bar chart and ``used_date`` is the effective
        date in the adjusted yield series.
    """
    # Build flat ticker mapping from country config
    flat_mapping = {}
    for country in COUNTRY_CONFIG:
        for ticker, label in country["tickers"].items():
            flat_mapping[ticker] = f"{country['name']} - {label}"

    tickers = list(flat_mapping.keys())
    df = _load_price_data(excel_path, tickers)
    df_adj, used_date = adjust_prices_for_mode(df, price_mode)
    perf = _build_returns_table(df_adj, flat_mapping)

    # Calculate max absolute value for bar scaling
    all_changes = []
    for country in COUNTRY_CONFIG:
        for ticker in country["tickers"].keys():
            if ticker in perf.index:
                change = perf.loc[ticker, "1W"]
                if pd.notna(change):
                    all_changes.append(abs(change))

    max_abs_value = max(all_changes) if all_changes else 10
    max_abs_value = max(max_abs_value, 5)  # Minimum 5 bps scale

    # Build country data for template
    countries = []
    for country in COUNTRY_CONFIG:
        tenors = []
        for ticker, label in country["tickers"].items():
            if ticker in perf.index:
                change = perf.loc[ticker, "1W"]
                if pd.notna(change):
                    # Bar width as percentage of half the chart
                    bar_width = abs(change) / max_abs_value * 48
                    bar_width = min(bar_width, 48)

                    # Determine bar class (inverted: up = red, down = green)
                    if change > 0:
                        bar_class = "rates-up"
                        value_class = "rates-up"
                        formatted = f"+{change:.1f} bps"
                    elif change < 0:
                        bar_class = "rates-down"
                        value_class = "rates-down"
                        formatted = f"{change:.1f} bps"
                    else:
                        bar_class = ""
                        value_class = ""
                        formatted = "0.0 bps"

                    tenors.append({
                        "label": label,
                        "change": change,
                        "bar_class": bar_class,
                        "bar_width": bar_width,
                        "value_class": value_class,
                        "formatted_change": formatted,
                    })

        if tenors:
            countries.append({
                "name": country["name"],
                "flag": country["flag"],
                "flag_html": get_flag_html(country["flag"]),
                "tenors": tenors,
            })

    # Calculate pixel dimensions (37.8 px per cm at 96 DPI)
    width_px = int(width_cm * 37.8 * SCALE_FACTOR)
    height_px = int(height_cm * 37.8 * SCALE_FACTOR)

    # Calculate scale
    scale_values = _calculate_rates_scale(max_abs_value)

    # Generate HTML
    template = Template(BONDS_RATES_HTML_TEMPLATE)
    html = template.render(
        countries=countries,
        width=width_px,
        height=height_px,
        scale=SCALE_FACTOR,
        **scale_values,
    )

    # Convert to PNG
    with tempfile.TemporaryDirectory() as tmpdir:
        hti = Html2Image(output_path=tmpdir, size=(width_px, height_px))
        hti.screenshot(html_str=html, save_as="rates_chart.png")
        with open(Path(tmpdir) / "rates_chart.png", "rb") as f:
            image_bytes = f.read()

    print(f"[Rates Performance] Generated chart: {width_px}x{height_px}px")
    return image_bytes, used_date


def _get_rates_color_class(value_bps: float) -> str:
    """Determine color class based on yield change in bps.

    INVERTED: rates down = green (good), rates up = red (bad)

    Thresholds (in bps):
    - level 1: 0-10 bps
    - level 2: 10-25 bps
    - level 3: 25-50 bps
    - level 4: 50-75 bps
    - level 5: >75 bps
    """
    abs_val = abs(value_bps)

    # INVERTED: negative change (rates down) = green
    if value_bps <= 0:
        prefix = "green"
    else:
        prefix = "red"

    if abs_val <= 10:
        level = 1
    elif abs_val <= 25:
        level = 2
    elif abs_val <= 50:
        level = 3
    elif abs_val <= 75:
        level = 4
    else:
        level = 5

    return f"{prefix}-{level}"


def _format_bps(value_bps: float) -> str:
    """Format value as bps string with sign."""
    if value_bps >= 0:
        return f"+{value_bps:.0f}"
    else:
        return f"{value_bps:.0f}"


def create_historical_performance_table(
    excel_path: Union[str, pathlib.Path],
    ticker_mapping: Dict[str, str] | None = None,
    *,
    width_cm: float = 17.02,
    height_cm: float = 13.0,
    price_mode: str = "Last Price",
) -> Tuple[bytes, Optional[pd.Timestamp]]:
    """Generate a heatmap table of yield changes using HTML template.

    This helper reads the yield data, optionally adjusts it according to
    ``price_mode``, computes changes for multiple horizons, grouped by
    country (US, EUR, JP, CN).

    Uses INVERTED color logic for bond rates:
    - Negative change (rates down) = GREEN (bullish)
    - Positive change (rates up) = RED (bearish)

    Parameters
    ----------
    excel_path : str or pathlib.Path
        Path to the Excel workbook containing the yield series.
    ticker_mapping : dict, optional
        Not used in new implementation, kept for API compatibility.
    width_cm, height_cm : float
        Dimensions of the generated figure in centimeters.
    price_mode : str, default ``"Last Price"``
        Either ``"Last Price"`` or ``"Last Close"``.  Determines how
        yield data is adjusted prior to computing changes and which
        effective date appears in the source footnote.

    Returns
    -------
    tuple
        A two‑tuple ``(image_bytes, used_date)`` where ``image_bytes`` is
        the PNG data for the heatmap and ``used_date`` is the effective
        date in the adjusted yield series.
    """
    # Build flat ticker mapping from country config
    flat_mapping = {}
    for country in COUNTRY_CONFIG:
        for ticker, label in country["tickers"].items():
            flat_mapping[ticker] = f"{country['name']} - {label}"

    tickers = list(flat_mapping.keys())
    df = _load_price_data(excel_path, tickers)
    df_adj, used_date = adjust_prices_for_mode(df, price_mode)
    perf = _build_returns_table(df_adj, flat_mapping)

    # Build country data for template
    countries = []
    for country in COUNTRY_CONFIG:
        tenors = []
        for ticker, label in country["tickers"].items():
            if ticker in perf.index:
                row = perf.loc[ticker]
                ytd = row.get("YTD", 0) if pd.notna(row.get("YTD")) else 0
                m1 = row.get("1M", 0) if pd.notna(row.get("1M")) else 0
                m3 = row.get("3M", 0) if pd.notna(row.get("3M")) else 0
                m6 = row.get("6M", 0) if pd.notna(row.get("6M")) else 0
                m12 = row.get("12M", 0) if pd.notna(row.get("12M")) else 0

                tenors.append({
                    "label": label,
                    "ytd_formatted": _format_bps(ytd),
                    "ytd_class": _get_rates_color_class(ytd),
                    "m1_formatted": _format_bps(m1),
                    "m1_class": _get_rates_color_class(m1),
                    "m3_formatted": _format_bps(m3),
                    "m3_class": _get_rates_color_class(m3),
                    "m6_formatted": _format_bps(m6),
                    "m6_class": _get_rates_color_class(m6),
                    "m12_formatted": _format_bps(m12),
                    "m12_class": _get_rates_color_class(m12),
                })

        if tenors:
            countries.append({
                "name": country["name"],
                "flag": country["flag"],
                "flag_html": get_flag_html(country["flag"]),
                "tenors": tenors,
            })

    # Calculate pixel dimensions (37.8 px per cm at 96 DPI)
    width_px = int(width_cm * 37.8 * SCALE_FACTOR)
    height_px = int(height_cm * 37.8 * SCALE_FACTOR)

    # Generate HTML
    template = Template(BONDS_HISTORICAL_HTML_TEMPLATE)
    html = template.render(
        countries=countries,
        width=width_px,
        height=height_px,
        scale=SCALE_FACTOR,
    )

    # Convert to PNG
    with tempfile.TemporaryDirectory() as tmpdir:
        hti = Html2Image(output_path=tmpdir, size=(width_px, height_px))
        hti.screenshot(html_str=html, save_as="rates_historical.png")
        with open(Path(tmpdir) / "rates_historical.png", "rb") as f:
            image_bytes = f.read()

    print(f"[Rates Historical] Generated heatmap: {width_px}x{height_px}px")
    return image_bytes, used_date


###############################################################################
# Colour helper for rates (invert green/red)
###############################################################################

def _colour_for_value_rates(val: float, max_pos: float, min_neg: float) -> Tuple[float, float, float, float]:
    """Return an RGBA colour for a single yield change value.

    Positive changes (rising yields) produce a gradient from white to red;
    negative changes (falling yields) produce a gradient from white to green.
    """
    white = np.array([1.0, 1.0, 1.0, 1.0])
    red = np.array(mcolors.to_rgba("#C70039"))
    green = np.array(mcolors.to_rgba("#2E8B57"))
    if val > 0 and max_pos > 0:
        ratio = float(val) / max_pos
        return tuple(white + ratio * (red - white))
    if val < 0 and min_neg < 0:
        ratio = float(val) / min_neg  # both negative → positive ratio
        return tuple(white + ratio * (green - white))
    return tuple(white)


###############################################################################
# Slide insertion helpers
###############################################################################

def _insert_dashboard_to_placeholder(
    prs: Presentation,
    image_bytes: bytes,
    placeholder_names: List[str],
    *,
    left_cm: float,
    top_cm: float,
    width_cm: float,
    height_cm: Optional[float] = None,
    used_date: Optional[pd.Timestamp],
    price_mode: str,
    source_placeholder_names: List[str],
) -> Presentation:
    # Reuse the helper from other modules but local to this file for clarity.
    target_slide = None
    placeholder_box = None
    name_candidates = [n.lower() for n in placeholder_names]
    pattern_candidates = [f"[{n}]" for n in name_candidates]
    for slide in prs.slides:
        for shape in slide.shapes:
            name_attr = getattr(shape, "name", "").lower()
            if name_attr in name_candidates:
                target_slide = slide
                placeholder_box = shape
                break
            if shape.has_text_frame:
                text_lower = (shape.text or "").strip().lower()
                if text_lower in [p.lower() for p in pattern_candidates]:
                    target_slide = slide
                    placeholder_box = shape
                    break
        if target_slide:
            break
    if target_slide is None:
        target_slide = prs.slides[min(11, len(prs.slides) - 1)]
    left = Cm(left_cm)
    top = Cm(top_cm)
    width = Cm(width_cm)
    stream = io.BytesIO(image_bytes)
    # Only specify width; let height auto-scale to maintain aspect ratio
    if height_cm is not None:
        pic = target_slide.shapes.add_picture(stream, left, top, width=width, height=Cm(height_cm))
    else:
        pic = target_slide.shapes.add_picture(stream, left, top, width=width)

    # Send picture to back (behind other elements like footnote)
    spTree = target_slide.shapes._spTree
    sp = pic._element
    spTree.remove(sp)
    spTree.insert(2, sp)  # Index 2 = back (0 and 1 are reserved)

    if used_date is not None:
        date_str = used_date.strftime("%d/%m/%Y")
        suffix = " Close" if price_mode.lower() == "last close" else ""
        source_text = f"Source: Bloomberg, Herculis Group, Data as of {date_str}{suffix}"
        source_candidates = [n.lower() for n in source_placeholder_names]
        source_patterns = [f"[{n}]" for n in source_candidates]
        for shape in target_slide.shapes:
            name_attr2 = getattr(shape, "name", "").lower()
            if name_attr2 in source_candidates:
                if shape.has_text_frame:
                    runs = shape.text_frame.paragraphs[0].runs
                    attrs = _capture_font_attrs(runs[0]) if runs else (None, None, None, None, None, None)
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    new_run = p.add_run()
                    new_run.text = source_text
                    _apply_font_attrs(new_run, *attrs)
                break
            if shape.has_text_frame:
                for pattern in source_patterns:
                    if pattern.lower() in (shape.text or "").lower():
                        runs = shape.text_frame.paragraphs[0].runs
                        attrs = _capture_font_attrs(runs[0]) if runs else (None, None, None, None, None, None)
                        try:
                            new_text = shape.text.replace(pattern, source_text)
                        except Exception:
                            new_text = source_text
                        shape.text_frame.clear()
                        p = shape.text_frame.paragraphs[0]
                        new_run = p.add_run()
                        new_run.text = new_text
                        _apply_font_attrs(new_run, *attrs)
                        break
                else:
                    continue
                break
    return prs


def insert_rates_performance_bar_slide(
    prs: Presentation,
    image_bytes: bytes,
    used_date: Optional[pd.Timestamp] = None,
    price_mode: str = "Last Price",
    *,
    left_cm: float = 3.35,
    top_cm: float = 4.6,
    width_cm: float = 17.02,
) -> Presentation:
    """Insert the weekly rates bar chart and source footnote into its designated slide.

    The slide is identified by a shape named ``rates_perf_1week`` or
    ``rates_perf_1w`` (or containing ``[rates_perf_1week]``).  The chart
    is inserted at the specified coordinates and a source footnote is
    written into ``rates_1w_source`` or a shape containing ``[rates_1w_source]``.

    Height is auto-calculated to maintain aspect ratio.
    """
    return _insert_dashboard_to_placeholder(
        prs,
        image_bytes,
        placeholder_names=["rates_perf_1w", "rates_perf_1week"],
        left_cm=left_cm,
        top_cm=top_cm,
        width_cm=width_cm,
        height_cm=None,  # Auto-scale to maintain aspect ratio
        used_date=used_date,
        price_mode=price_mode,
        source_placeholder_names=["rates_1w_source"],
    )


def insert_rates_performance_histo_slide(
    prs: Presentation,
    image_bytes: bytes,
    used_date: Optional[pd.Timestamp] = None,
    price_mode: str = "Last Price",
    *,
    left_cm: float = 3.35,
    top_cm: float = 4.6,
    width_cm: float = 17.02,
) -> Presentation:
    """Insert the rates historical performance heatmap and source footnote into its slide.

    The slide is identified by a shape named ``rates_perf_histo`` (or
    containing ``[rates_perf_histo]``).  A source footnote is written
    into ``rates_1w_source2`` or a shape containing ``[rates_1w_source2]``.

    Height is auto-calculated to maintain aspect ratio.
    """
    return _insert_dashboard_to_placeholder(
        prs,
        image_bytes,
        placeholder_names=["rates_perf_histo"],
        left_cm=left_cm,
        top_cm=top_cm,
        width_cm=width_cm,
        height_cm=None,  # Auto-scale to maintain aspect ratio
        used_date=used_date,
        price_mode=price_mode,
        source_placeholder_names=["rates_1w_source2"],
    )
