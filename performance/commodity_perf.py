# === Windows Playwright Fix - MUST BE FIRST ===
import sys
import asyncio
if sys.platform == 'win32':
    asyncio.set_event_loop_policy(asyncio.WindowsProactorEventLoopPolicy())

try:
    import nest_asyncio
    nest_asyncio.apply()
except ImportError:
    pass
# === End Fix ===

"""Commodity performance dashboard generation with price‑mode awareness and source footnotes.

This module produces charts summarising recent performance for a selection
of commodity indices grouped into Energy, Green Transition and Precious Metals.
Users can choose to base calculations on either the most recent intraday
price ("Last Price") or the previous day's closing price ("Last Close").
The resulting figures include a weekly returns bar chart and a heatmap
of longer horizons.  When the charts are inserted into a PowerPoint
presentation, a source footnote is added to the designated text boxes
on each slide, reflecting the effective date used in the computations
and the chosen price mode.  Text formatting from the placeholders
is preserved.

The commodities are grouped into three categories:

* **Energy:** Oil (WTI), Oil (Brent), Natural Gas (TTF)
* **Green Transition:** Uranium, Lithium, Manganese, Copper, Cobalt
* **Precious Metals:** Gold, Silver, Platinum

Key functions
-------------

* ``create_weekly_performance_chart`` – Build the weekly bar chart and
  return both the image bytes and the effective date used for
  computation.  Bars are ordered within each category and separated
  visually by inserting a small gap between categories.
* ``create_historical_performance_table`` – Build the heatmap of
  returns for multiple horizons and return the image bytes and the
  effective date.
* ``insert_commodity_performance_bar_slide`` – Insert the weekly bar
  chart and source footnote into its slide.
* ``insert_commodity_performance_histo_slide`` – Insert the historical
  performance heatmap and source footnote into its slide.
"""

from __future__ import annotations

import io
import json
import os
import pathlib
from typing import Dict, List, Tuple, Union, Optional

import matplotlib.pyplot as plt
import matplotlib.colors as mcolors
import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.util import Cm
from jinja2 import Template, Environment
from html2image import Html2Image
from playwright.sync_api import sync_playwright

from utils import adjust_prices_for_mode
from market_compass.weekly_performance.html_template import (
    COMMODITIES_WEEKLY_HTML_TEMPLATE,
    COMMODITIES_HISTORICAL_HTML_TEMPLATE,
    COMMODITY_YTD_EVOLUTION_HTML_TEMPLATE,
    YTD_INSUFFICIENT_DATA_HTML_TEMPLATE,
)

try:
    # Reuse font attribute helpers from the SPX module if available
    from technical_analysis.equity.spx import (
        _get_run_font_attributes as _capture_font_attrs,  # type: ignore
        _apply_run_font_attributes as _apply_font_attrs,  # type: ignore
    )
except Exception:
    # Fallback helpers: only support basic size and RGB; ignore theme colours
    def _capture_font_attrs(run):  # type: ignore
        if run is None:
            return (None, None, None, None, None, None)
        size = run.font.size
        colour = run.font.color
        rgb = None
        try:
            rgb = colour.rgb
        except Exception:
            rgb = None
        return (size, rgb, None, None, run.font.bold, run.font.italic)

    def _apply_font_attrs(new_run, size, rgb, theme_color, brightness, bold, italic):  # type: ignore
        if size is not None:
            new_run.font.size = size
        if rgb is not None:
            try:
                new_run.font.color.rgb = rgb
            except Exception:
                pass
        if bold is not None:
            new_run.font.bold = bold
        if italic is not None:
            new_run.font.italic = italic

###############################################################################
# Data loading and return calculations
###############################################################################

def _load_price_data(excel_path: Union[str, pathlib.Path], tickers: List[str]) -> pd.DataFrame:
    """Load a DataFrame of dates and prices for the specified commodity tickers."""
    df = pd.read_excel(excel_path, sheet_name="data_prices")
    df = df.drop(index=0)
    df = df[df[df.columns[0]] != "DATES"].copy()
    df = df.rename(columns={df.columns[0]: "Date"})
    df["Date"] = pd.to_datetime(df["Date"], errors="coerce")
    out = df[["Date"] + tickers].copy()
    for t in tickers:
        out[t] = pd.to_numeric(out[t], errors="coerce")
    return out.dropna(subset=["Date"]).sort_values("Date").reset_index(drop=True)


def _compute_horizon_returns(
    df: pd.DataFrame,
    ticker: str,
    today: pd.Timestamp,
    days: int,
) -> float:
    """Compute the percentage return of a series over the given lookback."""
    past_date = today - pd.Timedelta(days=days)
    past_series = df.loc[df["Date"] <= past_date, ticker]
    if len(past_series) == 0:
        return float("nan")
    past_price = past_series.iloc[-1]
    current_price = df[ticker].iloc[-1]
    if past_price == 0 or pd.isna(past_price):
        return float("nan")
    return (current_price - past_price) / past_price * 100.0


def _build_returns_table(
    df: pd.DataFrame,
    mapping: Dict[str, str],
) -> pd.DataFrame:
    """Calculate performance metrics for each commodity ticker."""
    today = df["Date"].max()
    results: Dict[str, Dict[str, float]] = {}
    for ticker in mapping:
        results[ticker] = {
            "1W": _compute_horizon_returns(df, ticker, today, 7),
            "1M": _compute_horizon_returns(df, ticker, today, 30),
            "3M": _compute_horizon_returns(df, ticker, today, 90),
            "6M": _compute_horizon_returns(df, ticker, today, 180),
            "12M": _compute_horizon_returns(df, ticker, today, 365),
        }
        start_of_year = pd.Timestamp(year=today.year, month=1, day=1)
        past_series = df.loc[df["Date"] <= start_of_year, ticker]
        if len(past_series) > 0:
            past_price = past_series.iloc[-1]
            current_price = df[ticker].iloc[-1]
            if past_price and not pd.isna(past_price):
                results[ticker]["YTD"] = (current_price - past_price) / past_price * 100.0
            else:
                results[ticker]["YTD"] = float("nan")
        else:
            results[ticker]["YTD"] = float("nan")
    table = pd.DataFrame.from_dict(results, orient="index")
    table["Name"] = table.index.map(mapping.get)
    return table[["Name", "1W", "1M", "3M", "6M", "12M", "YTD"]]


###############################################################################
# Colour helper (green/red)
###############################################################################

def _colour_for_value(val: float, max_pos: float, min_neg: float) -> Tuple[float, float, float, float]:
    """Return an RGBA colour for a single value using independent scales."""
    white = np.array([1.0, 1.0, 1.0, 1.0])
    red = np.array(mcolors.to_rgba("#C70039"))
    green = np.array(mcolors.to_rgba("#2E8B57"))
    if val > 0 and max_pos > 0:
        ratio = float(val) / max_pos
        return tuple(white + ratio * (green - white))
    if val < 0 and min_neg < 0:
        ratio = float(val) / min_neg  # both negative → positive ratio
        return tuple(white + ratio * (red - white))
    return tuple(white)


###############################################################################
# Figure generation functions
###############################################################################

def create_weekly_performance_chart(
    excel_path: Union[str, pathlib.Path],
    ticker_mapping: Dict[str, str] | None = None,
    *,
    width: float = 14.0,
    height: float = 5.0,
    price_mode: str = "Last Price",
) -> Tuple[bytes, Optional[pd.Timestamp]]:
    """Generate a bar chart of 1‑week returns grouped by commodity category."""
    default_mapping = {
        # Energy
        "CL1 Comdty": "Oil (WTI)",
        "CO1 Comdty": "Oil (Brent)",
        "TTFG1MON OECM Index": "Natural Gas (TTF)",
        # Green Transition
        "UXA1 Comdty": "Uranium",
        "LMY1 Comdty": "Lithium",
        "IMRU5 Comdty": "Manganese",
        "LP1 Comdty": "Copper",
        "CVT1 Comdty": "Cobalt",
        # Precious Metals
        "GCA Comdty": "Gold",
        "SIA Comdty": "Silver",
        "XPT Comdty": "Platinum",
        "XPD Curncy": "Palladium",
    }
    mapping = ticker_mapping or default_mapping
    tickers = list(mapping.keys())
    df = _load_price_data(excel_path, tickers)
    df_adj, used_date = adjust_prices_for_mode(df, price_mode)
    today = df_adj["Date"].max()
    # Compute 1W returns directly for each ticker
    returns_dict = {}
    for t in tickers:
        returns_dict[t] = _compute_horizon_returns(df_adj, t, today, 7)
    # Define category order and mapping ticker → category
    category_map = {
        "CL1 Comdty": "Energy",
        "CO1 Comdty": "Energy",
        "TTFG1MON OECM Index": "Energy",
        "UXA1 Comdty": "Green Transition",
        "LMY1 Comdty": "Green Transition",
        "IMRU5 Comdty": "Green Transition",
        "LP1 Comdty": "Green Transition",
        "CVT1 Comdty": "Green Transition",
        "GCA Comdty": "Precious Metals",
        "SIA Comdty": "Precious Metals",
        "XPT Comdty": "Precious Metals",
        "XPD Curncy": "Precious Metals",
    }
    categories = ["Energy", "Green Transition", "Precious Metals"]
    # Build a list of (category, human name, return) and sort within categories by return descending
    perf_list: List[Tuple[str, str, float]] = []
    for t in tickers:
        cat = category_map.get(t, "Other")
        perf_list.append((cat, mapping[t], returns_dict.get(t, float("nan"))))
    # Sort by category order then return descending (NaNs last)
    perf_list.sort(key=lambda x: (categories.index(x[0]) if x[0] in categories else 99, float("inf") if pd.isna(x[2]) else -x[2]))
    names = [item[1] for item in perf_list]
    returns = [item[2] for item in perf_list]
    bar_df = pd.DataFrame({"Name": names, "Return": returns})
    fig, ax = plt.subplots(figsize=(width, height))
    # Colour bars based on sign
    bar_colors = ["#2E8B57" if x > 0 else "#C70039" for x in bar_df["Return"]]
    bars = ax.barh(bar_df["Name"], bar_df["Return"], color=bar_colors)
    # Compute dynamic offset for labels
    min_val = min(bar_df["Return"].min(), 0)
    max_val = bar_df["Return"].max()
    range_val = max_val - min_val if max_val > min_val else abs(max_val)
    label_offset = range_val * 0.05
    for bar, val in zip(bars, bar_df["Return"]):
        width_val = bar.get_width()
        x_pos = width_val + (label_offset if width_val > 0 else -label_offset)
        ax.text(
            x_pos,
            bar.get_y() + bar.get_height() / 2.0,
            f"{val:+.1f}%",
            va="center",
            ha="left" if width_val > 0 else "right",
            fontweight="bold",
            color="black",
        )
    # Add category separators: draw horizontal lines between categories
    # Determine index boundaries for categories
    cat_indices: List[int] = []
    idx_count = 0
    for cat in categories:
        count_in_cat = sum(1 for item in perf_list if item[0] == cat)
        if count_in_cat == 0:
            continue
        idx_count += count_in_cat
        cat_indices.append(idx_count)
    # Draw horizontal lines to separate categories
    for boundary in cat_indices[:-1]:  # skip last
        ax.axhline(boundary - 0.5, color="#BFBFBF", linewidth=0.8, linestyle="--")
    ax.axvline(0.0, color="grey", linewidth=0.8, linestyle="--")
    ax.set_xticks([])
    ax.set_yticks(range(len(bar_df)))
    ax.set_yticklabels(bar_df["Name"], fontsize=10)
    ax.invert_yaxis()
    for spine in ax.spines.values():
        spine.set_visible(False)
    ax.tick_params(axis="y", length=0)
    margin = (max_val - min_val) * 0.2
    ax.set_xlim(min_val - margin, max_val + margin)
    fig.subplots_adjust(left=0.4)
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=300, bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf.getvalue(), used_date


def create_historical_performance_table(
    excel_path: Union[str, pathlib.Path],
    ticker_mapping: Dict[str, str] | None = None,
    *,
    width: float = 14.0,
    height: float = 6.0,
    price_mode: str = "Last Price",
) -> Tuple[bytes, Optional[pd.Timestamp]]:
    """Generate a heatmap table of returns for commodities."""
    default_mapping = {
        # Energy
        "CL1 Comdty": "Oil (WTI)",
        "CO1 Comdty": "Oil (Brent)",
        "TTFG1MON OECM Index": "Natural Gas (TTF)",
        # Green Transition
        "UXA1 Comdty": "Uranium",
        "LMY1 Comdty": "Lithium",
        "IMRU5 Comdty": "Manganese",
        "LP1 Comdty": "Copper",
        "CVT1 Comdty": "Cobalt",
        # Precious Metals
        "GCA Comdty": "Gold",
        "SIA Comdty": "Silver",
        "XPT Comdty": "Platinum",
        "XPD Curncy": "Palladium",
    }
    mapping = ticker_mapping or default_mapping
    tickers = list(mapping.keys())
    df = _load_price_data(excel_path, tickers)
    df_adj, used_date = adjust_prices_for_mode(df, price_mode)
    perf = _build_returns_table(df_adj, mapping)
    heat_df = perf[["Name", "YTD", "1M", "3M", "6M", "12M"]].dropna().sort_values("YTD", ascending=False).reset_index(drop=True)
    n_rows = len(heat_df)
    cols = ["YTD", "1M", "3M", "6M", "12M"]
    color_array = np.zeros((n_rows, len(cols), 4))
    for j, col in enumerate(cols):
        col_vals = heat_df[col].astype(float).values
        pos_vals = col_vals[col_vals > 0]
        neg_vals = col_vals[col_vals < 0]
        max_pos = pos_vals.max() if len(pos_vals) > 0 else 0.0
        min_neg = neg_vals.min() if len(neg_vals) > 0 else 0.0
        for i, val in enumerate(col_vals):
            color_array[i, j] = _colour_for_value(float(val), max_pos, min_neg)
    fig, ax = plt.subplots(figsize=(width, height))
    ax.imshow(color_array, aspect="auto")
    # Add text
    for i in range(n_rows):
        for j, col in enumerate(cols):
            val = heat_df.iloc[i][col]
            ax.text(
                j,
                i,
                f"{val:+.1f}%",
                ha="center",
                va="center",
                fontsize=9,
                fontweight="bold",
                color="black",
            )
    ax.set_xticks(range(len(cols)))
    ax.set_xticklabels(cols, fontsize=12, fontweight="bold")
    ax.xaxis.set_ticks_position("top")
    ax.xaxis.set_label_position("top")
    ax.set_yticks(range(n_rows))
    ax.set_yticklabels(heat_df["Name"], fontsize=9)
    ax.tick_params(axis="y", length=0)
    ax.tick_params(axis="x", length=0)
    for spine in ax.spines.values():
        spine.set_visible(False)
    fig.subplots_adjust(left=0.5)
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=300, bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf.getvalue(), used_date


###############################################################################
# Slide insertion helpers
###############################################################################

def _insert_dashboard_to_placeholder(
    prs: Presentation,
    image_bytes: bytes,
    placeholder_names: List[str],
    *,
    left_cm: float,
    top_cm: float,
    width_cm: float,
    height_cm: float,
    used_date: Optional[pd.Timestamp],
    price_mode: str,
    source_placeholder_names: List[str],
) -> Presentation:
    target_slide = None
    placeholder_box = None
    name_candidates = [n.lower() for n in placeholder_names]
    pattern_candidates = [f"[{n}]" for n in name_candidates]
    for slide in prs.slides:
        for shape in slide.shapes:
            name_attr = getattr(shape, "name", "").lower()
            if name_attr in name_candidates:
                target_slide = slide
                placeholder_box = shape
                break
            if shape.has_text_frame:
                text_lower = (shape.text or "").strip().lower()
                if text_lower in [p.lower() for p in pattern_candidates]:
                    target_slide = slide
                    placeholder_box = shape
                    break
        if target_slide:
            break
    if target_slide is None:
        target_slide = prs.slides[min(11, len(prs.slides) - 1)]
    left = Cm(left_cm)
    top = Cm(top_cm)
    width = Cm(width_cm)
    height = Cm(height_cm)
    stream = io.BytesIO(image_bytes)
    target_slide.shapes.add_picture(stream, left, top, width=width, height=height)
    # Insert source footnote if a date is available
    if used_date is not None:
        date_str = used_date.strftime("%d/%m/%Y")
        suffix = " Close" if price_mode.lower() == "last close" else ""
        source_text = f"Source: Bloomberg, Herculis Group, Data as of {date_str}{suffix}"
        source_candidates = [n.lower() for n in source_placeholder_names]
        source_patterns = [f"[{n}]" for n in source_candidates]
        for shape in target_slide.shapes:
            name_attr2 = getattr(shape, "name", "").lower()
            if name_attr2 in source_candidates:
                if shape.has_text_frame:
                    runs = shape.text_frame.paragraphs[0].runs
                    attrs = _capture_font_attrs(runs[0]) if runs else (None, None, None, None, None, None)
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    new_run = p.add_run()
                    new_run.text = source_text
                    _apply_font_attrs(new_run, *attrs)
                break
            if shape.has_text_frame:
                for pattern in source_patterns:
                    if pattern.lower() in (shape.text or "").lower():
                        runs = shape.text_frame.paragraphs[0].runs
                        attrs = _capture_font_attrs(runs[0]) if runs else (None, None, None, None, None, None)
                        try:
                            new_text = shape.text.replace(pattern, source_text)
                        except Exception:
                            new_text = source_text
                        shape.text_frame.clear()
                        p = shape.text_frame.paragraphs[0]
                        new_run = p.add_run()
                        new_run.text = new_text
                        _apply_font_attrs(new_run, *attrs)
                        break
                else:
                    continue
                break
    return prs


def insert_commodity_performance_bar_slide(
    prs: Presentation,
    image_bytes: bytes,
    used_date: Optional[pd.Timestamp] = None,
    price_mode: str = "Last Price",
    *,
    left_cm: float = 0.0,
    top_cm: float = 3.22,
    width_cm: float = 25.0,
    height_cm: float = 11.66,
) -> Presentation:
    """Insert the weekly commodity bar chart and source footnote into its slide."""
    return _insert_dashboard_to_placeholder(
        prs,
        image_bytes,
        # Use only commo placeholders for consistency
        placeholder_names=["commo_perf_1w", "commo_perf_1week"],
        left_cm=left_cm,
        top_cm=top_cm,
        width_cm=width_cm,
        height_cm=height_cm,
        used_date=used_date,
        price_mode=price_mode,
        source_placeholder_names=["commo_1w_source"],
    )


def insert_commodity_performance_histo_slide(
    prs: Presentation,
    image_bytes: bytes,
    used_date: Optional[pd.Timestamp] = None,
    price_mode: str = "Last Price",
    *,
    left_cm: float = 0.0,
    top_cm: float = 3.22,
    width_cm: float = 25.0,
    height_cm: float = 11.66,
) -> Presentation:
    """Insert the commodity historical performance heatmap and source footnote."""
    return _insert_dashboard_to_placeholder(
        prs,
        image_bytes,
        # Use only commo placeholder for consistency
        placeholder_names=["commo_perf_histo"],
        left_cm=left_cm,
        top_cm=top_cm,
        width_cm=width_cm,
        height_cm=height_cm,
        used_date=used_date,
        price_mode=price_mode,
        source_placeholder_names=["commo_1w_source2"],
    )


###############################################################################
# HTML-based Commodities Weekly Performance Chart
###############################################################################

# Commodity configuration with categories, icons, and tickers
COMMODITY_CATEGORIES = [
    {
        "name": "Energy",
        "icon": "⚡",
        "items": [
            {"ticker": "CL1 Comdty", "name": "Oil (WTI)", "icon": "🛢️"},
            {"ticker": "CO1 Comdty", "name": "Oil (Brent)", "icon": "🛢️"},
            {"ticker": "TTFG1MON OECM Index", "name": "Natural Gas (TTF)", "icon": "🔥"},
        ]
    },
    {
        "name": "Green Transition",
        "icon": "🌱",
        "items": [
            {"ticker": "UXA1 Comdty", "name": "Uranium", "icon": "☢️"},
            {"ticker": "LMY1 Comdty", "name": "Lithium", "icon": "🔋"},
            {"ticker": "IMRU5 Comdty", "name": "Manganese", "icon": "⚡"},
            {"ticker": "LP1 Comdty", "name": "Copper", "icon": "🔶"},
            {"ticker": "CVT1 Comdty", "name": "Cobalt", "icon": "💎"},
        ]
    },
    {
        "name": "Precious Metals",
        "icon": "💰",
        "items": [
            {"ticker": "GCA Comdty", "name": "Gold", "icon": "🏆"},
            {"ticker": "SIA Comdty", "name": "Silver", "icon": "🥈"},
            {"ticker": "XPT Comdty", "name": "Platinum", "icon": "💠"},
            {"ticker": "XPD Curncy", "name": "Palladium", "icon": "🔘"},
        ]
    },
]

# Chart dimensions (hardcoded, no calculations)
PNG_WIDTH_PX = 1963
PNG_HEIGHT_PX = 1134
PPT_WIDTH_CM = 17.31
PPT_HEIGHT_CM = 10.0
PPT_LEFT_CM = 3.35
PPT_TOP_CM = 4.6
HTML_SCALE = 3


def _format_commodity_percentage(value: float) -> str:
    """Format percentage with sign."""
    return f"{value:+.1f}%"


def create_weekly_html_performance_chart(
    excel_path: Union[str, pathlib.Path],
    *,
    price_mode: str = "Last Price",
) -> Tuple[bytes, Optional[pd.Timestamp]]:
    """Generate HTML-based weekly commodity performance bar chart.

    Returns:
        Tuple of (PNG bytes, effective date used for computation)
    """
    # Get all tickers from configuration
    all_tickers = []
    for cat in COMMODITY_CATEGORIES:
        for item in cat["items"]:
            all_tickers.append(item["ticker"])

    # Load and adjust price data
    df = _load_price_data(excel_path, all_tickers)
    df_adj, used_date = adjust_prices_for_mode(df, price_mode)
    today = df_adj["Date"].max()

    # Compute 1W returns for all commodities
    returns_dict = {}
    for ticker in all_tickers:
        try:
            ret_val = _compute_horizon_returns(df_adj, ticker, today, 7)
            returns_dict[ticker] = ret_val
        except Exception:
            returns_dict[ticker] = float("nan")

    # Calculate max absolute value for bar scaling (minimum 5%)
    valid_returns = [r for r in returns_dict.values() if not pd.isna(r)]
    max_abs_value = max((abs(r) for r in valid_returns), default=5.0)
    max_abs_value = max(max_abs_value, 5.0)  # At least 5%

    # Build category data for template
    categories_data = []
    for cat in COMMODITY_CATEGORIES:
        items_data = []
        for item in cat["items"]:
            ticker = item["ticker"]
            value = returns_dict.get(ticker, 0.0)
            if pd.isna(value):
                value = 0.0

            # Calculate bar width as percentage of max (48% max to leave margin)
            bar_width = abs(value) / max_abs_value * 48

            items_data.append({
                "name": item["name"],
                "icon": item["icon"],
                "value": value,
                "formatted_value": _format_commodity_percentage(value),
                "bar_class": "positive" if value >= 0 else "negative",
                "bar_width": bar_width,
                "value_class": "positive" if value >= 0 else "negative",
            })

        categories_data.append({
            "name": cat["name"],
            "icon": cat["icon"],
            "commodities": items_data,
        })

    # Calculate scale labels
    scale_max = max_abs_value
    scale_min = -max_abs_value
    scale_mid_high = max_abs_value / 2
    scale_mid_low = -max_abs_value / 2

    # Render HTML template
    template = Template(COMMODITIES_WEEKLY_HTML_TEMPLATE)
    html_content = template.render(
        scale=HTML_SCALE,
        width=PNG_WIDTH_PX,
        height=PNG_HEIGHT_PX,
        categories=categories_data,
        scale_max=f"+{scale_max:.0f}%",
        scale_min=f"{scale_min:.0f}%",
        scale_mid_high=f"+{scale_mid_high:.0f}%",
        scale_mid_low=f"{scale_mid_low:.0f}%",
    )

    # Convert HTML to PNG
    hti = Html2Image(size=(PNG_WIDTH_PX, PNG_HEIGHT_PX))
    hti.screenshot(html_str=html_content, save_as="commodity_weekly_temp.png")

    # Read the generated PNG
    with open("commodity_weekly_temp.png", "rb") as f:
        png_bytes = f.read()

    # Clean up temp file
    import os
    try:
        os.remove("commodity_weekly_temp.png")
    except Exception:
        pass

    return png_bytes, used_date


def insert_commodity_weekly_html_slide(
    prs: Presentation,
    image_bytes: bytes,
    used_date: Optional[pd.Timestamp] = None,
    price_mode: str = "Last Price",
) -> Presentation:
    """Insert the HTML-based commodity weekly performance chart into PowerPoint."""
    if not image_bytes:
        return prs

    # Find target slide by placeholder name
    target_slide = None
    placeholder_names = ["commo_perf_1w", "commo_perf_1week"]
    name_candidates = [n.lower() for n in placeholder_names]
    pattern_candidates = [f"[{n}]" for n in name_candidates]

    for slide in prs.slides:
        for shape in slide.shapes:
            name_attr = getattr(shape, "name", "").lower()
            if name_attr in name_candidates:
                target_slide = slide
                break
            if shape.has_text_frame:
                text_lower = (shape.text or "").strip().lower()
                if text_lower in [p.lower() for p in pattern_candidates]:
                    target_slide = slide
                    break
        if target_slide:
            break

    if target_slide is None:
        target_slide = prs.slides[min(11, len(prs.slides) - 1)]

    # Insert chart image with exact hardcoded dimensions
    left = Cm(PPT_LEFT_CM)
    top = Cm(PPT_TOP_CM)
    width = Cm(PPT_WIDTH_CM)
    height = Cm(PPT_HEIGHT_CM)

    stream = io.BytesIO(image_bytes)
    picture = target_slide.shapes.add_picture(stream, left, top, width=width, height=height)

    # Send picture to back
    spTree = target_slide.shapes._spTree
    pic_element = picture._element
    spTree.remove(pic_element)
    spTree.insert(2, pic_element)

    print(f"[Commodity Weekly HTML] Chart inserted at ({PPT_LEFT_CM}, {PPT_TOP_CM}) cm, size: {PPT_WIDTH_CM}x{PPT_HEIGHT_CM} cm")

    # Update source placeholder if date available
    if used_date is not None:
        date_str = used_date.strftime("%d/%m/%Y")
        suffix = " Close" if price_mode.lower() == "last close" else ""
        source_text = f"Source: Bloomberg, Herculis Group, Data as of {date_str}{suffix}"
        source_candidates = ["commo_1w_source"]
        source_patterns = [f"[{n}]" for n in source_candidates]

        for shape in target_slide.shapes:
            name_attr = getattr(shape, "name", "").lower()
            if name_attr in [n.lower() for n in source_candidates]:
                if shape.has_text_frame:
                    runs = shape.text_frame.paragraphs[0].runs
                    attrs = _capture_font_attrs(runs[0]) if runs else (None, None, None, None, None, None)
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    new_run = p.add_run()
                    new_run.text = source_text
                    _apply_font_attrs(new_run, *attrs)
                break
            if shape.has_text_frame:
                for pattern in source_patterns:
                    if pattern.lower() in (shape.text or "").lower():
                        runs = shape.text_frame.paragraphs[0].runs
                        attrs = _capture_font_attrs(runs[0]) if runs else (None, None, None, None, None, None)
                        shape.text_frame.clear()
                        p = shape.text_frame.paragraphs[0]
                        new_run = p.add_run()
                        new_run.text = source_text
                        _apply_font_attrs(new_run, *attrs)
                        break
                else:
                    continue
                break

    return prs


# =============================================================================
# COMMODITIES HISTORICAL PERFORMANCE HEATMAP
# =============================================================================

# Historical chart dimensions
HIST_PNG_WIDTH_PX = 1930
HIST_PNG_HEIGHT_PX = 1200
HIST_PPT_WIDTH_CM = 17.02
HIST_PPT_LEFT_CM = 3.35
HIST_PPT_TOP_CM = 4.6
HIST_HTML_SCALE = 3


def _get_commodity_historical_color_class(value: float) -> str:
    """
    Determine color class based on value magnitude.

    Thresholds for commodities (higher volatility than bonds):
    - level 1: 0-5%
    - level 2: 5-15%
    - level 3: 15-30%
    - level 4: 30-50%
    - level 5: >50%
    """
    abs_val = abs(value)
    prefix = "positive" if value >= 0 else "negative"

    if abs_val <= 5:
        level = 1
    elif abs_val <= 15:
        level = 2
    elif abs_val <= 30:
        level = 3
    elif abs_val <= 50:
        level = 4
    else:
        level = 5

    return f"{prefix}-{level}"


def _format_historical_percentage(value: float) -> str:
    """Format value as percentage string with sign."""
    if value >= 0:
        return f"+{value:.1f}%"
    else:
        return f"{value:.1f}%"


def create_historical_html_performance_chart(
    excel_path: Union[str, pathlib.Path],
    *,
    price_mode: str = "Last Price",
) -> Tuple[bytes, Optional[pd.Timestamp]]:
    """Generate HTML-based historical commodity performance heatmap.

    Returns:
        Tuple of (PNG bytes, effective date used for computation)
    """
    # Get all tickers from configuration
    all_tickers = []
    for cat in COMMODITY_CATEGORIES:
        for item in cat["items"]:
            all_tickers.append(item["ticker"])

    # Load and adjust price data
    df = _load_price_data(excel_path, all_tickers)
    df_adj, used_date = adjust_prices_for_mode(df, price_mode)
    today = df_adj["Date"].max()

    # Compute returns for all horizons
    horizons = {
        "ytd": None,  # Special case for YTD
        "m1": 30,
        "m3": 90,
        "m6": 180,
        "m12": 365,
    }

    returns_dict = {}
    for ticker in all_tickers:
        returns_dict[ticker] = {}
        for horizon_name, days in horizons.items():
            try:
                if horizon_name == "ytd":
                    # YTD: from Jan 1 of current year
                    year_start = pd.Timestamp(year=today.year, month=1, day=1)
                    df_ticker = df_adj[df_adj["Date"] >= year_start]
                    if ticker in df_ticker.columns and len(df_ticker) >= 2:
                        first_price = df_ticker[ticker].dropna().iloc[0]
                        last_price = df_ticker[ticker].dropna().iloc[-1]
                        ret_val = ((last_price / first_price) - 1) * 100
                    else:
                        ret_val = float("nan")
                else:
                    ret_val = _compute_horizon_returns(df_adj, ticker, today, days)
                returns_dict[ticker][horizon_name] = ret_val
            except Exception:
                returns_dict[ticker][horizon_name] = float("nan")

    # Build category data for template
    categories_data = []
    for cat in COMMODITY_CATEGORIES:
        items_data = []
        for item in cat["items"]:
            ticker = item["ticker"]
            ticker_returns = returns_dict.get(ticker, {})

            ytd = ticker_returns.get("ytd", 0.0)
            m1 = ticker_returns.get("m1", 0.0)
            m3 = ticker_returns.get("m3", 0.0)
            m6 = ticker_returns.get("m6", 0.0)
            m12 = ticker_returns.get("m12", 0.0)

            # Replace NaN with 0
            ytd = 0.0 if pd.isna(ytd) else ytd
            m1 = 0.0 if pd.isna(m1) else m1
            m3 = 0.0 if pd.isna(m3) else m3
            m6 = 0.0 if pd.isna(m6) else m6
            m12 = 0.0 if pd.isna(m12) else m12

            items_data.append({
                "name": item["name"],
                "icon": item["icon"],
                "ytd_value": ytd,  # Keep raw value for sorting
                "ytd_formatted": _format_historical_percentage(ytd),
                "ytd_class": _get_commodity_historical_color_class(ytd),
                "m1_formatted": _format_historical_percentage(m1),
                "m1_class": _get_commodity_historical_color_class(m1),
                "m3_formatted": _format_historical_percentage(m3),
                "m3_class": _get_commodity_historical_color_class(m3),
                "m6_formatted": _format_historical_percentage(m6),
                "m6_class": _get_commodity_historical_color_class(m6),
                "m12_formatted": _format_historical_percentage(m12),
                "m12_class": _get_commodity_historical_color_class(m12),
            })

        # Sort items by YTD descending (best performers first)
        items_data.sort(key=lambda x: x["ytd_value"], reverse=True)

        categories_data.append({
            "name": cat["name"],
            "icon": cat["icon"],
            "commodities": items_data,  # Using "commodities" to avoid conflict with dict.items()
        })

    # Render HTML template
    template = Template(COMMODITIES_HISTORICAL_HTML_TEMPLATE)
    html_content = template.render(
        scale=HIST_HTML_SCALE,
        width=HIST_PNG_WIDTH_PX,
        height=HIST_PNG_HEIGHT_PX,
        categories=categories_data,
    )

    # Convert HTML to PNG
    hti = Html2Image(size=(HIST_PNG_WIDTH_PX, HIST_PNG_HEIGHT_PX))
    hti.screenshot(html_str=html_content, save_as="commodity_historical_temp.png")

    # Read the generated PNG
    with open("commodity_historical_temp.png", "rb") as f:
        png_bytes = f.read()

    # Clean up temp file
    import os
    try:
        os.remove("commodity_historical_temp.png")
    except Exception:
        pass

    return png_bytes, used_date


def insert_commodity_historical_html_slide(
    prs: Presentation,
    image_bytes: bytes,
    used_date: Optional[pd.Timestamp] = None,
    price_mode: str = "Last Price",
) -> Presentation:
    """Insert the HTML-based commodity historical performance heatmap into PowerPoint."""
    if not image_bytes:
        return prs

    # Find target slide by placeholder name
    target_slide = None
    placeholder_names = ["commo_perf_histo", "commo_perf_hist"]
    name_candidates = [n.lower() for n in placeholder_names]
    pattern_candidates = [f"[{n}]" for n in name_candidates]

    for slide in prs.slides:
        for shape in slide.shapes:
            name_attr = getattr(shape, "name", "").lower()
            if name_attr in name_candidates:
                target_slide = slide
                break
            if shape.has_text_frame:
                text_lower = (shape.text or "").strip().lower()
                if text_lower in [p.lower() for p in pattern_candidates]:
                    target_slide = slide
                    break
        if target_slide:
            break

    if target_slide is None:
        print("[Commodity Historical HTML] ERROR: Slide not found")
        return prs

    # Insert chart image with exact hardcoded dimensions
    left = Cm(HIST_PPT_LEFT_CM)
    top = Cm(HIST_PPT_TOP_CM)
    width = Cm(HIST_PPT_WIDTH_CM)

    stream = io.BytesIO(image_bytes)
    picture = target_slide.shapes.add_picture(stream, left, top, width=width)

    # Send picture to back
    spTree = target_slide.shapes._spTree
    pic_element = picture._element
    spTree.remove(pic_element)
    spTree.insert(2, pic_element)

    print(f"[Commodity Historical HTML] Heatmap inserted at ({HIST_PPT_LEFT_CM}, {HIST_PPT_TOP_CM}) cm")

    # Update source placeholder if date available
    if used_date is not None:
        date_str = used_date.strftime("%d/%m/%Y")
        suffix = " Close" if price_mode.lower() == "last close" else ""
        source_text = f"Source: Bloomberg, Herculis Group, Data as of {date_str}{suffix}"
        source_candidates = ["commo_hist_source", "commo_1w_source2"]
        source_patterns = [f"[{n}]" for n in source_candidates]

        for shape in target_slide.shapes:
            name_attr = getattr(shape, "name", "").lower()
            if name_attr in [n.lower() for n in source_candidates]:
                if shape.has_text_frame:
                    runs = shape.text_frame.paragraphs[0].runs
                    attrs = _capture_font_attrs(runs[0]) if runs else (None, None, None, None, None, None)
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    new_run = p.add_run()
                    new_run.text = source_text
                    _apply_font_attrs(new_run, *attrs)
                break
            if shape.has_text_frame:
                for pattern in source_patterns:
                    if pattern.lower() in (shape.text or "").lower():
                        runs = shape.text_frame.paragraphs[0].runs
                        attrs = _capture_font_attrs(runs[0]) if runs else (None, None, None, None, None, None)
                        shape.text_frame.clear()
                        p = shape.text_frame.paragraphs[0]
                        new_run = p.add_run()
                        new_run.text = source_text
                        _apply_font_attrs(new_run, *attrs)
                        break
                else:
                    continue
                break

    return prs

# =============================================================================
# COMMODITY YTD EVOLUTION LINE CHART (Chart.js + Playwright)
# =============================================================================

# Commodity configuration with real commodity colors
# Using same tickers as COMMODITY_CATEGORIES for weekly/historical charts
COMMODITY_YTD_CONFIG = {
    "GCA Comdty": {"name": "Gold", "color": "#DAA520"},
    "SIA Comdty": {"name": "Silver", "color": "#8C8C8C"},
    "XPT Comdty": {"name": "Platinum", "color": "#4682B4"},
    "XPD Curncy": {"name": "Palladium", "color": "#E8B4B8"},
    "LP1 Comdty": {"name": "Copper", "color": "#B87333"},
    "UXA1 Comdty": {"name": "Uranium", "color": "#6B8E23"},
    "CL1 Comdty": {"name": "Oil (WTI)", "color": "#1C1C1C"},
}

# Chart dimensions
COMMODITY_YTD_PNG_WIDTH_PX = 2700
COMMODITY_YTD_PNG_HEIGHT_PX = 1350
COMMODITY_YTD_PPT_WIDTH_CM = 23.0
COMMODITY_YTD_PPT_LEFT_CM = 0.5
COMMODITY_YTD_PPT_TOP_CM = 4.2
COMMODITY_YTD_HTML_SCALE = 3


def _get_commodity_ytd_year(data_end_date: pd.Timestamp) -> int:
    """Determine which year to show as YTD based on data end date."""
    return data_end_date.year


def _get_commodity_chart_title(data_end_date: pd.Timestamp) -> str:
    """Generate chart title with correct year."""
    year = _get_commodity_ytd_year(data_end_date)
    return f"{year} YTD Evolution"


def _load_commodity_price_data(excel_path, tickers):
    """Load a DataFrame of dates and prices for the specified commodity tickers."""
    df = pd.read_excel(excel_path, sheet_name="data_prices")
    df = df.drop(index=0)
    df = df[df[df.columns[0]] != "DATES"].copy()
    df = df.rename(columns={df.columns[0]: "Date"})
    df["Date"] = pd.to_datetime(df["Date"], errors="coerce")
    out = df[["Date"] + [t for t in tickers if t in df.columns]].copy()
    for t in tickers:
        if t in out.columns:
            out[t] = pd.to_numeric(out[t], errors="coerce")
    return out.dropna(subset=["Date"]).sort_values("Date").reset_index(drop=True)


def _compute_commodity_ytd_series(df, ticker):
    """Compute weekly YTD performance series for a commodity."""
    if ticker not in df.columns:
        return [], []

    last_date = df["Date"].max()
    year = last_date.year
    start_of_year = pd.Timestamp(year=year, month=1, day=1)
    df_year = df[df["Date"] >= start_of_year].copy().sort_values("Date")

    if len(df_year) == 0:
        return [], []

    first_valid_idx = df_year[ticker].first_valid_index()
    if first_valid_idx is None:
        return [], []

    baseline_price = df_year.loc[first_valid_idx, ticker]
    if pd.isna(baseline_price) or baseline_price == 0:
        return [], []

    month_names = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun',
                   'Jul', 'Aug', 'Sep', 'Oct', 'Nov', 'Dec']

    daily_labels = []
    daily_values = []

    for _, row in df_year.iterrows():
        price = row[ticker]
        date = row["Date"]
        if not pd.isna(price):
            ytd_return = (price - baseline_price) / baseline_price * 100
            daily_labels.append(month_names[date.month - 1])
            daily_values.append(round(ytd_return, 1))

    # If fewer than 22 data points, use daily frequency instead of weekly
    # This handles beginning of year when there's only a few days of data
    if len(daily_values) < 22:
        labels = daily_labels
        values = daily_values
    else:
        # Sample to weekly (every 5th trading day)
        weekly_indices = list(range(0, len(daily_values), 5))
        if len(daily_values) > 0 and weekly_indices[-1] != len(daily_values) - 1:
            weekly_indices.append(len(daily_values) - 1)

        labels = [daily_labels[i] for i in weekly_indices]
        values = [daily_values[i] for i in weekly_indices]

    return labels, values


def create_commodity_ytd_evolution_chart(excel_path, *, price_mode="Last Price"):
    """Generate HTML-based Commodity YTD Evolution line chart."""
    tickers = list(COMMODITY_YTD_CONFIG.keys())
    df = _load_commodity_price_data(excel_path, tickers)
    df_adj, used_date = adjust_prices_for_mode(df, price_mode)

    if used_date is None:
        used_date = df_adj["Date"].max()

    chart_title = _get_commodity_chart_title(used_date)

    all_labels = []
    datasets = []

    for ticker, config in COMMODITY_YTD_CONFIG.items():
        labels, values = _compute_commodity_ytd_series(df_adj, ticker)
        if labels and values:
            if len(labels) > len(all_labels):
                all_labels = labels
            datasets.append({
                "label": config["name"],
                "data": values,
                "borderColor": config["color"],
                "backgroundColor": "transparent",
            })

    max_len = len(all_labels)
    for dataset in datasets:
        while len(dataset["data"]) < max_len:
            dataset["data"].append(None)

    # Check for insufficient data (< 3 data points)
    if len(all_labels) < 3:
        print(f"[DEBUG] Commodity: Insufficient data ({len(all_labels)} points), rendering placeholder")
        env = Environment()
        template = env.from_string(YTD_INSUFFICIENT_DATA_HTML_TEMPLATE)
        html_content = template.render(
            width=COMMODITY_YTD_PNG_WIDTH_PX,
            height=COMMODITY_YTD_PNG_HEIGHT_PX,
            scale=COMMODITY_YTD_HTML_SCALE,
            chart_title=chart_title,
        )
        png_bytes = None
        try:
            with sync_playwright() as p:
                browser = p.chromium.launch(headless=True)
                page = browser.new_page(viewport={
                    'width': COMMODITY_YTD_PNG_WIDTH_PX,
                    'height': COMMODITY_YTD_PNG_HEIGHT_PX
                })
                page.set_content(html_content, wait_until='networkidle')
                page.wait_for_timeout(500)
                png_bytes = page.screenshot()
                browser.close()
        except Exception as e:
            print(f"[ERROR] Playwright failed for placeholder: {e}")
        return png_bytes, used_date

    all_values = []
    for ds in datasets:
        all_values.extend([v for v in ds["data"] if v is not None])

    if all_values:
        data_min = min(all_values)
        data_max = max(all_values)
        y_min = (int(data_min / 5) - 1) * 5
        y_max = (int(data_max / 5) + 2) * 5
        if data_min < 0 < data_max:
            y_min = min(y_min, -5)
            y_max = max(y_max, 5)
    else:
        y_min = -20
        y_max = 40

    env = Environment()
    env.filters['tojson'] = json.dumps
    template = env.from_string(COMMODITY_YTD_EVOLUTION_HTML_TEMPLATE)
    html_content = template.render(
        width=COMMODITY_YTD_PNG_WIDTH_PX,
        height=COMMODITY_YTD_PNG_HEIGHT_PX,
        scale=COMMODITY_YTD_HTML_SCALE,
        chart_title=chart_title,
        labels=all_labels,
        datasets=datasets,
        y_min=y_min,
        y_max=y_max,
    )

    png_bytes = None
    try:
        with sync_playwright() as p:
            browser = p.chromium.launch(headless=True)
            page = browser.new_page(viewport={
                'width': COMMODITY_YTD_PNG_WIDTH_PX,
                'height': COMMODITY_YTD_PNG_HEIGHT_PX
            })
            page.set_content(html_content, wait_until='networkidle')
            try:
                page.wait_for_selector('body[data-chart-ready="true"]', timeout=10000)
            except Exception:
                page.wait_for_timeout(3000)
            png_bytes = page.screenshot()
            browser.close()
    except Exception as e:
        print(f"[ERROR] Playwright failed for commodity chart: {e}")

    return png_bytes, used_date


def insert_commodity_ytd_evolution_slide(prs, image_bytes, used_date=None, price_mode="Last Price", subtitle=None):
    """Insert the Commodity YTD Evolution chart into PowerPoint.

    Parameters
    ----------
    prs : Presentation
        The PowerPoint presentation to modify.
    image_bytes : bytes
        PNG data for the YTD evolution chart.
    used_date : pandas.Timestamp or None, optional
        Effective date used for performance calculations.
    price_mode : str, default "Last Price"
        Either "Last Price" or "Last Close".
    subtitle : str or None, optional
        Subtitle text to insert into the ytd_commo_subtitle placeholder.
    """
    if not image_bytes:
        return prs

    target_slide = None
    placeholder_names = ["ytd_commo_perf"]

    for slide in prs.slides:
        for shape in slide.shapes:
            name_attr = getattr(shape, "name", "").lower()
            if name_attr in [n.lower() for n in placeholder_names]:
                target_slide = slide
                break
            if shape.has_text_frame:
                text_lower = (shape.text or "").strip().lower()
                if text_lower == "[ytd_commo_perf]":
                    target_slide = slide
                    break
        if target_slide:
            break

    if target_slide is None:
        print(f"[Commodity YTD Evolution] ERROR: Slide not found")
        return prs

    # ------------------------------------------------------------------
    # Insert subtitle into the 'ytd_commo_subtitle' placeholder
    # ------------------------------------------------------------------
    if subtitle:
        for shape in target_slide.shapes:
            name_attr = getattr(shape, "name", "")
            if name_attr and name_attr.lower() == "ytd_commo_subtitle" and shape.has_text_frame:
                tf = shape.text_frame
                paragraph = tf.paragraphs[0]
                runs = paragraph.runs
                attrs = _capture_font_attrs(runs[0]) if runs else (None, None, None, None, None, None)
                # Determine replacement: replace tokens "XXX" or "[ytd_commo_subtitle]"
                original_text = "".join(run.text for run in runs) if runs else ""
                new_text = original_text
                if "XXX" in new_text:
                    new_text = new_text.replace("XXX", subtitle)
                elif "[ytd_commo_subtitle]" in new_text:
                    new_text = new_text.replace("[ytd_commo_subtitle]", subtitle)
                else:
                    new_text = subtitle
                # Clear and insert new run
                tf.clear()
                p = tf.paragraphs[0]
                new_run = p.add_run()
                new_run.text = new_text
                _apply_font_attrs(new_run, *attrs)
                break

    # ------------------------------------------------------------------
    # Insert chart image with exact hardcoded dimensions
    # ------------------------------------------------------------------
    from io import BytesIO
    image_stream = BytesIO(image_bytes)

    left = Cm(COMMODITY_YTD_PPT_LEFT_CM)
    top = Cm(COMMODITY_YTD_PPT_TOP_CM)
    width = Cm(COMMODITY_YTD_PPT_WIDTH_CM)

    pic = target_slide.shapes.add_picture(image_stream, left, top, width=width)

    spTree = target_slide.shapes._spTree
    sp = pic._element
    spTree.remove(sp)
    spTree.insert(2, sp)

    print(f"[Commodity YTD Evolution] Chart inserted")

    # ------------------------------------------------------------------
    # Insert data source footnote
    # ------------------------------------------------------------------
    if used_date is not None:
        date_str = used_date.strftime("%d/%m/%Y")
        suffix = " Close" if price_mode.lower() == "last close" else ""
        source_text = f"Source: Bloomberg, Herculis Group, Data as of {date_str}{suffix}"
        placeholder_name = "ytd_commo_source"
        placeholder_patterns = ["[ytd_commo_source]", "ytd_commo_source"]
        inserted = False
        for shape in target_slide.shapes:
            name_attr = getattr(shape, "name", "")
            if name_attr and name_attr.lower() == placeholder_name:
                if shape.has_text_frame:
                    runs = shape.text_frame.paragraphs[0].runs
                    attrs = _capture_font_attrs(runs[0]) if runs else (None, None, None, None, None, None)
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    new_run = p.add_run()
                    new_run.text = source_text
                    _apply_font_attrs(new_run, *attrs)
                inserted = True
                break
            if shape.has_text_frame:
                for pattern in placeholder_patterns:
                    if pattern.lower() in (shape.text or "").lower():
                        runs = shape.text_frame.paragraphs[0].runs
                        attrs = _capture_font_attrs(runs[0]) if runs else (None, None, None, None, None, None)
                        try:
                            new_text = shape.text.replace(pattern, source_text)
                        except Exception:
                            new_text = source_text
                        shape.text_frame.clear()
                        p = shape.text_frame.paragraphs[0]
                        new_run = p.add_run()
                        new_run.text = new_text
                        _apply_font_attrs(new_run, *attrs)
                        inserted = True
                        break
                if inserted:
                    break

    return prs
