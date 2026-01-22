"""Equity performance dashboard generation with price‐mode awareness and source footnotes.

This module produces charts summarising recent performance for a selection of major
equity indices.  Users can choose to base calculations on either the most
recent intraday price ("Last Price") or the previous day's closing price
("Last Close").  The resulting figures include a weekly returns bar chart and
a heatmap of longer horizons.  When the charts are inserted into a PowerPoint
presentation, a source footnote is added to the designated text boxes on
each slide, reflecting the effective date used in the computations and the
chosen price mode.  Text formatting from the placeholders is preserved.

Key functions
-------------

* ``create_weekly_performance_chart`` – Build the weekly bar chart and return
  both the image bytes and the effective date used for computation.
* ``create_historical_performance_table`` – Build the heatmap of returns for
  multiple horizons and return the image bytes and the effective date.
* ``insert_equity_performance_bar_slide`` – Insert the weekly bar chart and
  source footnote into its slide.
* ``insert_equity_performance_histo_slide`` – Insert the historical performance
  heatmap and source footnote into its slide.

Usage example::

    from performance.equity_perf import (
        create_weekly_performance_chart,
        create_historical_performance_table,
        insert_equity_performance_bar_slide,
        insert_equity_performance_histo_slide,
    )

    # Generate charts with price-mode awareness
    bar_bytes, used_date = create_weekly_performance_chart(excel_path, price_mode="Last Close")
    histo_bytes, _ = create_historical_performance_table(excel_path, price_mode="Last Close")
    # Insert into PPT, supplying used_date and price_mode for source footnote
    prs = insert_equity_performance_bar_slide(prs, bar_bytes, used_date, price_mode)
    prs = insert_equity_performance_histo_slide(prs, histo_bytes, used_date, price_mode)

"""

from __future__ import annotations

# === Windows Playwright Fix - MUST BE EARLY ===
import sys
import asyncio
if sys.platform == 'win32':
    asyncio.set_event_loop_policy(asyncio.WindowsProactorEventLoopPolicy())

try:
    import nest_asyncio
    nest_asyncio.apply()
except ImportError:
    pass
# === End Fix ===

import io
import os
import json
import pathlib
from datetime import datetime
from typing import Dict, List, Tuple, Union, Optional

import matplotlib.pyplot as plt
import matplotlib.colors as mcolors
import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.util import Cm
from jinja2 import Template, Environment
from playwright.sync_api import sync_playwright

from utils import adjust_prices_for_mode
from helpers.flag_utils import get_flag_html
from helpers.html_to_image import render_html_to_image
from market_compass.weekly_performance.html_template import EQUITY_YTD_EVOLUTION_HTML_TEMPLATE, YTD_INSUFFICIENT_DATA_HTML_TEMPLATE, FX_IMPACT_ANALYSIS_EUR_HTML_TEMPLATE

try:
    # Reuse font attribute helpers from the SPX module if available
    from technical_analysis.equity.spx import (
        _get_run_font_attributes as _capture_font_attrs,  # type: ignore
        _apply_run_font_attributes as _apply_font_attrs,  # type: ignore
    )
except Exception:
    # Fallback helpers: only support basic size and RGB; ignore theme colours
    def _capture_font_attrs(run):  # type: ignore
        if run is None:
            return (None, None, None, None, None, None)
        size = run.font.size
        colour = run.font.color
        rgb = None
        try:
            rgb = colour.rgb
        except Exception:
            rgb = None
        return (size, rgb, None, None, run.font.bold, run.font.italic)

    def _apply_font_attrs(new_run, size, rgb, theme_color, brightness, bold, italic):  # type: ignore
        if size is not None:
            new_run.font.size = size
        if rgb is not None:
            try:
                new_run.font.color.rgb = rgb
            except Exception:
                pass
        if bold is not None:
            new_run.font.bold = bold
        if italic is not None:
            new_run.font.italic = italic

###############################################################################
# Data loading and return calculations
###############################################################################

def _load_price_data(excel_path: Union[str, pathlib.Path], tickers: List[str]) -> pd.DataFrame:
    """Load a DataFrame of dates and prices for the specified tickers.

    Parameters
    ----------
    excel_path : str or pathlib.Path
        Path to the Excel workbook containing a sheet named ``data_prices``.
    tickers : list of str
        Column names in the sheet corresponding to the desired indices.

    Returns
    -------
    DataFrame
        A tidy DataFrame with columns ``Date`` and one column per ticker,
        sorted by date and with numeric prices.
    """
    df = pd.read_excel(excel_path, sheet_name="data_prices")
    # Drop the first row which usually contains metadata and rows with the
    # sentinel ``DATES`` value.
    df = df.drop(index=0)
    df = df[df[df.columns[0]] != "DATES"].copy()
    df = df.rename(columns={df.columns[0]: "Date"})
    df["Date"] = pd.to_datetime(df["Date"], errors="coerce")
    # Restrict to requested tickers and coerce to numeric
    out = df[["Date"] + tickers].copy()
    for t in tickers:
        out[t] = pd.to_numeric(out[t], errors="coerce")
    return out.dropna(subset=["Date"]).sort_values("Date").reset_index(drop=True)

def _compute_horizon_returns(
    df: pd.DataFrame,
    ticker: str,
    today: pd.Timestamp,
    days: int,
) -> float:
    """Compute the percentage return of a series over the given lookback.

    This helper finds the last available price on or before ``today`` minus
    ``days`` and compares it to the most recent price.  If insufficient
    history exists, NaN is returned.

    Parameters
    ----------
    df : DataFrame
        A tidy DataFrame with ``Date`` and price columns.
    ticker : str
        Column name for which to compute the return.
    today : Timestamp
        Anchor date; typically the last date in the dataset.
    days : int
        Lookback period in days.

    Returns
    -------
    float
        Percent return over the lookback horizon, or NaN if not computable.
    """
    past_date = today - pd.Timedelta(days=days)
    past_series = df.loc[df["Date"] <= past_date, ticker]
    if len(past_series) == 0:
        return float("nan")
    past_price = past_series.iloc[-1]
    current_price = df[ticker].iloc[-1]
    if past_price == 0 or pd.isna(past_price):
        return float("nan")
    return (current_price - past_price) / past_price * 100.0

def _build_returns_table(
    df: pd.DataFrame,
    mapping: Dict[str, str],
) -> pd.DataFrame:
    """Calculate performance metrics for each ticker.

    Returns a DataFrame indexed by the ticker code with the following
    columns: ``Name``, ``1W``, ``1M``, ``3M``, ``6M``, ``12M`` and ``YTD``.

    Parameters
    ----------
    df : DataFrame
        A tidy DataFrame as returned by ``_load_price_data``.
    mapping : dict
        Mapping from ticker codes to human‑friendly names.

    Returns
    -------
    DataFrame
        Performance table.
    """
    today = df["Date"].max()
    results: Dict[str, Dict[str, float]] = {}
    for ticker in mapping:
        results[ticker] = {
            "1W": _compute_horizon_returns(df, ticker, today, 7),
            "1M": _compute_horizon_returns(df, ticker, today, 30),
            "3M": _compute_horizon_returns(df, ticker, today, 90),
            "6M": _compute_horizon_returns(df, ticker, today, 180),
            "12M": _compute_horizon_returns(df, ticker, today, 365),
        }
        # Year‑to‑date: look back to the start of the calendar year
        start_of_year = pd.Timestamp(year=today.year, month=1, day=1)
        past_series = df.loc[df["Date"] <= start_of_year, ticker]
        if len(past_series) > 0:
            past_price = past_series.iloc[-1]
            current_price = df[ticker].iloc[-1]
            if past_price and not pd.isna(past_price):
                results[ticker]["YTD"] = (current_price - past_price) / past_price * 100.0
            else:
                results[ticker]["YTD"] = float("nan")
        else:
            results[ticker]["YTD"] = float("nan")
    table = pd.DataFrame.from_dict(results, orient="index")
    table["Name"] = table.index.map(mapping.get)
    return table[["Name", "1W", "1M", "3M", "6M", "12M", "YTD"]]

###############################################################################
# Figure generation functions
###############################################################################

def create_weekly_performance_chart(
    excel_path: Union[str, pathlib.Path],
    ticker_mapping: Dict[str, str] | None = None,
    *,
    width: float = 14.0,
    height: float = 5.0,
    price_mode: str = "Last Price",
) -> Tuple[bytes, Optional[pd.Timestamp]]:
    """Generate a bar chart of 1‑week returns with price‑mode adjustment.

    This helper reads the specified Excel file, optionally adjusts the price
    history according to the selected ``price_mode`` (``"Last Price"`` vs
    ``"Last Close"``), computes 1‑week returns for the configured tickers,
    sorts them in descending order and constructs a horizontal bar chart.
    The effective date used for the calculations is returned along with
    the PNG data.

    Parameters
    ----------
    excel_path : str or pathlib.Path
        Path to the Excel workbook containing the price series.
    ticker_mapping : dict, optional
        Mapping from ticker codes to display names.  If not provided,
        sensible defaults are used.
    width, height : float
        Dimensions of the generated figure in inches.
    price_mode : str, default ``"Last Price"``
        One of ``"Last Price"`` or ``"Last Close"``.  Determines whether
        intraday prices are used or the most recent row is dropped if it
        corresponds to today's date.

    Returns
    -------
    tuple
        A two‑tuple ``(image_bytes, used_date)`` where ``image_bytes`` is
        the PNG data for the bar chart and ``used_date`` is the effective
        date in the adjusted price series.
    """
    import tempfile
    from pathlib import Path
    from jinja2 import Template
    from market_compass.weekly_performance.html_template import WEEKLY_PERFORMANCE_HTML_TEMPLATE

    # Index mapping with display names and flag country codes (for flagcdn.com)
    INDEX_CONFIG = {
        "SPX Index": {"name": "S&P 500", "flag": "us"},
        "DAX Index": {"name": "Dax", "flag": "de"},
        "SMI Index": {"name": "SMI", "flag": "ch"},
        "NKY Index": {"name": "Nikkei 225", "flag": "jp"},
        "SHSZ300 Index": {"name": "CSI 300", "flag": "cn"},
        "SENSEX Index": {"name": "Sensex", "flag": "in"},
        "IBOV Index": {"name": "Bovespa", "flag": "br"},
        "MEXBOL Index": {"name": "Mexbol", "flag": "mx"},
        "SASEIDX Index": {"name": "TASI", "flag": "sa"},
    }

    tickers = list(INDEX_CONFIG.keys())

    # Load data and adjust according to price mode
    df = _load_price_data(excel_path, tickers)
    df_adj, used_date = adjust_prices_for_mode(df, price_mode)

    # Build returns
    default_mapping = {k: v["name"] for k, v in INDEX_CONFIG.items()}
    perf = _build_returns_table(df_adj, default_mapping)

    # Get 1W returns and sort descending
    bar_df = perf[["Name", "1W"]].dropna().sort_values("1W", ascending=False).reset_index(drop=True)

    # Get latest prices for each ticker (for price column)
    last_prices = {}
    for ticker in INDEX_CONFIG.keys():
        if ticker in df_adj.columns:
            last_prices[ticker] = df_adj[ticker].iloc[-1]

    # Prepare rows for HTML template
    rows = []
    max_abs_value = max(abs(bar_df["1W"].max()), abs(bar_df["1W"].min())) if len(bar_df) > 0 else 1
    max_abs_value = max(max_abs_value, 1.0)  # At least 1%

    for i, row in bar_df.iterrows():
        name = row["Name"]
        value = row["1W"] / 100  # Convert percentage to decimal

        # Find flag and ticker for this name
        flag = ""
        ticker = ""
        for t, config in INDEX_CONFIG.items():
            if config["name"] == name:
                flag = config["flag"]
                ticker = t
                break

        # Get and format price (integer with thousand separator)
        price = last_prices.get(ticker, 0)
        formatted_price = f"{int(price):,}" if price else ""

        # Determine highlight class
        highlight_class = ""
        if i == 0 and value > 0:
            highlight_class = "top-performer"
        elif i == len(bar_df) - 1 and value < 0:
            highlight_class = "worst-performer"

        # Bar width as percentage of half the chart
        bar_width = abs(value) / (max_abs_value / 100) * 50
        bar_width = min(bar_width, 48)

        # Value classes
        if value > 0:
            bar_class = "positive"
            value_class = "positive"
            formatted_value = f"+{value * 100:.1f}%"
        elif value < 0:
            bar_class = "negative"
            value_class = "negative"
            formatted_value = f"{value * 100:.1f}%"
        else:
            bar_class = ""
            value_class = "zero"
            formatted_value = "0.0%"

        rows.append({
            "name": name,
            "flag": flag,
            "flag_html": get_flag_html(flag),
            "value": value,
            "formatted_price": formatted_price,
            "highlight_class": highlight_class,
            "bar_class": bar_class,
            "bar_width": bar_width,
            "value_class": value_class,
            "formatted_value": formatted_value,
        })

    # Calculate scale
    if max_abs_value <= 1:
        scale_max = 1
    elif max_abs_value <= 2:
        scale_max = 2
    elif max_abs_value <= 5:
        scale_max = 5
    elif max_abs_value <= 10:
        scale_max = 10
    else:
        scale_max = int(max_abs_value) + 5

    scale_values = {
        "scale_min": f"-{scale_max}%",
        "scale_mid_low": f"-{scale_max // 2}%",
        "scale_mid_high": f"+{scale_max // 2}%",
        "scale_max": f"+{scale_max}%",
    }

    # Generate HTML - PNG must be 3x scale for crisp rendering
    # Target: 17.31cm × 10cm in PowerPoint
    SCALE_FACTOR = 4
    width_px = int(17.31 * 37.8 * SCALE_FACTOR)  # = 1963 px
    height_px = int(10 * 37.8 * SCALE_FACTOR)     # = 1134 px

    template = Template(WEEKLY_PERFORMANCE_HTML_TEMPLATE)
    html = template.render(
        rows=rows,
        width=width_px,
        height=height_px,
        scale=SCALE_FACTOR,
        **scale_values,
    )

    # Convert to PNG using Playwright
    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as f:
        img_path = f.name

    print(f"[Equity Weekly] Rendering image ({width_px}x{height_px}px)...")
    render_html_to_image(
        html_content=html,
        output_path=img_path,
        size=(width_px, height_px),
        device_scale_factor=1  # CSS already scaled by SCALE_FACTOR
    )

    with open(img_path, "rb") as f:
        png_bytes = f.read()
    Path(img_path).unlink(missing_ok=True)

    return png_bytes, used_date

def create_historical_performance_table(
    excel_path: Union[str, pathlib.Path],
    ticker_mapping: Dict[str, str] | None = None,
    *,
    width: float = 14.0,
    height: float = 6.0,
    price_mode: str = "Last Price",
) -> Tuple[bytes, Optional[pd.Timestamp]]:
    """Generate a heatmap table of historical returns with price‑mode adjustment.

    This helper reads the price data, optionally adjusts it according to
    ``price_mode``, computes returns for multiple horizons, sorts the
    table by YTD performance (descending) and constructs a heatmap where
    each column is coloured independently.  The effective date used
    for the computations is returned along with the PNG data.

    Parameters
    ----------
    excel_path : str or pathlib.Path
        Path to the Excel workbook containing the price series.
    ticker_mapping : dict, optional
        Mapping from ticker codes to display names.  If not provided,
        sensible defaults are used.
    width, height : float
        Dimensions of the generated figure in inches.
    price_mode : str, default ``"Last Price"``
        Either ``"Last Price"`` or ``"Last Close"``.  Determines how
        price data is adjusted prior to computing returns and which
        effective date appears in the source footnote.

    Returns
    -------
    tuple
        A two‑tuple ``(image_bytes, used_date)`` where ``image_bytes`` is
        the PNG data for the heatmap and ``used_date`` is the effective
        date in the adjusted price series.
    """
    import tempfile
    from pathlib import Path
    from jinja2 import Template
    from market_compass.weekly_performance.html_template import HISTORICAL_PERFORMANCE_HTML_TEMPLATE

    # Index mapping with display names and flag country codes (for flagcdn.com)
    INDEX_CONFIG = {
        "SPX Index": {"name": "S&P 500", "flag": "us"},
        "DAX Index": {"name": "Dax", "flag": "de"},
        "SMI Index": {"name": "SMI", "flag": "ch"},
        "NKY Index": {"name": "Nikkei 225", "flag": "jp"},
        "SHSZ300 Index": {"name": "CSI 300", "flag": "cn"},
        "SENSEX Index": {"name": "Sensex", "flag": "in"},
        "IBOV Index": {"name": "Bovespa", "flag": "br"},
        "MEXBOL Index": {"name": "Mexbol", "flag": "mx"},
        "SASEIDX Index": {"name": "TASI", "flag": "sa"},
    }

    tickers = list(INDEX_CONFIG.keys())

    # Load data and adjust according to price mode
    df = _load_price_data(excel_path, tickers)
    df_adj, used_date = adjust_prices_for_mode(df, price_mode)

    # Build returns
    default_mapping = {k: v["name"] for k, v in INDEX_CONFIG.items()}
    perf = _build_returns_table(df_adj, default_mapping)

    # Get returns and sort by YTD descending
    heat_df = perf[["Name", "YTD", "1M", "3M", "6M", "12M"]].dropna().sort_values("YTD", ascending=False).reset_index(drop=True)

    def get_color_class(value: float) -> str:
        """Determine color class based on value magnitude. All values show green or red."""
        val = value / 100  # Convert from percentage to decimal
        # Always green (positive) or red (negative) - no neutral
        prefix = "positive" if val >= 0 else "negative"
        abs_val = abs(val)
        if abs_val <= 0.05:
            level = 1
        elif abs_val <= 0.10:
            level = 2
        elif abs_val <= 0.20:
            level = 3
        elif abs_val <= 0.30:
            level = 4
        else:
            level = 5
        return f"{prefix}-{level}"

    def format_percentage(value: float) -> str:
        """Format value as percentage string."""
        return f"{value:.1f}%"

    # Prepare rows for HTML template
    rows = []
    for _, row in heat_df.iterrows():
        name = row["Name"]

        # Find flag for this name
        flag = ""
        for ticker, config in INDEX_CONFIG.items():
            if config["name"] == name:
                flag = config["flag"]
                break

        rows.append({
            "name": name,
            "flag": flag,
            "flag_html": get_flag_html(flag),
            "ytd_formatted": format_percentage(row["YTD"]),
            "ytd_class": get_color_class(row["YTD"]),
            "m1_formatted": format_percentage(row["1M"]),
            "m1_class": get_color_class(row["1M"]),
            "m3_formatted": format_percentage(row["3M"]),
            "m3_class": get_color_class(row["3M"]),
            "m6_formatted": format_percentage(row["6M"]),
            "m6_class": get_color_class(row["6M"]),
            "m12_formatted": format_percentage(row["12M"]),
            "m12_class": get_color_class(row["12M"]),
        })

    # Generate HTML - PNG must match PowerPoint aspect ratio
    # Target: 25.0cm × 11.66cm in PowerPoint (from insert_equity_performance_histo_slide)
    SCALE_FACTOR = 4
    width_px = int(25.0 * 37.8 * SCALE_FACTOR)   # = 3780 px
    height_px = int(11.66 * 37.8 * SCALE_FACTOR) # = 1762 px

    template = Template(HISTORICAL_PERFORMANCE_HTML_TEMPLATE)
    html = template.render(
        rows=rows,
        width=width_px,
        height=height_px,
        scale=SCALE_FACTOR,
    )

    # Convert to PNG using Playwright
    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as f:
        img_path = f.name

    print(f"[Equity Historical] Rendering image ({width_px}x{height_px}px)...")
    render_html_to_image(
        html_content=html,
        output_path=img_path,
        size=(width_px, height_px),
        device_scale_factor=1  # CSS already scaled by SCALE_FACTOR
    )

    with open(img_path, "rb") as f:
        png_bytes = f.read()
    Path(img_path).unlink(missing_ok=True)

    return png_bytes, used_date

###############################################################################
# Colour helper
###############################################################################

def _colour_for_value(val: float, max_pos: float, min_neg: float) -> Tuple[float, float, float, float]:
    """Return an RGBA colour for a single value using independent scales.

    Positive values produce a gradient from white to green; negative
    values from white to red.  Zeros (or undefined scales) return
    white.  The input ``max_pos`` should be the maximum positive value
    in the column, and ``min_neg`` the minimum (most negative) value.
    """
    white = np.array([1.0, 1.0, 1.0, 1.0])
    red = np.array(mcolors.to_rgba("#C70039"))
    green = np.array(mcolors.to_rgba("#2E8B57"))
    if val > 0 and max_pos > 0:
        ratio = float(val) / max_pos
        return tuple(white + ratio * (green - white))
    if val < 0 and min_neg < 0:
        ratio = float(val) / min_neg  # both negative → positive ratio
        return tuple(white + ratio * (red - white))
    return tuple(white)

###############################################################################
# Slide insertion helpers
###############################################################################

def _insert_dashboard_to_placeholder(
    prs: Presentation,
    image_bytes: bytes,
    placeholder_names: List[str],
    *,
    left_cm: float,
    top_cm: float,
    width_cm: float,
    height_cm: Optional[float] = None,
    used_date: Optional[pd.Timestamp],
    price_mode: str,
    source_placeholder_names: List[str],
) -> Presentation:
    """Helper to insert an image and source footnote into a slide identified by placeholders.

    This internal helper searches for shapes whose names match any of
    ``placeholder_names`` or whose text contains any of those names in
    square brackets.  Once found, it uses the slide but ignores the
    shape's dimensions, inserting the provided image at the specified
    coordinates and size.  It also searches for a text box named
    ``source_placeholder_names`` (or containing the pattern) and
    replaces its contents with a source footnote based on ``used_date``
    and ``price_mode`` while preserving the original formatting.  If no
    matching placeholder is found, the image is inserted on the last
    slide.  Footnote insertion is skipped if ``used_date`` is ``None``.

    Parameters
    ----------
    prs : Presentation
        The presentation into which the image should be inserted.
    image_bytes : bytes
        The PNG data to insert.
    placeholder_names : list of str
        Names of shapes to look for (lowercased).  The function also
        recognises text placeholders like ``[equity_perf_1week]``.
    left_cm, top_cm, width_cm, height_cm : float
        Desired position and size in centimetres.
    used_date : pandas.Timestamp or None
        The effective date for the source footnote.  If ``None`` the
        source is not inserted.
    price_mode : str
        Either ``"Last Price"`` or ``"Last Close"``, used to suffix
        ``" Close"`` in the footnote.
    source_placeholder_names : list of str
        Names or patterns (without brackets) for the source text box.

    Returns
    -------
    Presentation
        The modified presentation.
    """
    target_slide = None
    placeholder_box = None
    # Normalise names for comparison
    name_candidates = [n.lower() for n in placeholder_names]
    pattern_candidates = [f"[{n}]" for n in name_candidates]
    print(f"[_insert_dashboard DEBUG] Searching for placeholders: {placeholder_names}")
    print(f"[_insert_dashboard DEBUG] Name candidates: {name_candidates}")
    print(f"[_insert_dashboard DEBUG] Pattern candidates: {pattern_candidates}")
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            name_attr = getattr(shape, "name", "").lower()
            if name_attr in name_candidates:
                print(f"[_insert_dashboard DEBUG] FOUND by name on slide {slide_idx}: '{name_attr}'")
                target_slide = slide
                placeholder_box = shape
                break
            if shape.has_text_frame:
                text_lower = (shape.text or "").strip().lower()
                if text_lower in [p.lower() for p in pattern_candidates]:
                    print(f"[_insert_dashboard DEBUG] FOUND by text on slide {slide_idx}: '{text_lower}'")
                    target_slide = slide
                    placeholder_box = shape
                    break
        if target_slide:
            break
    if target_slide is None:
        print(f"[_insert_dashboard DEBUG] NOT FOUND - using fallback slide")
        target_slide = prs.slides[min(11, len(prs.slides) - 1)]
    else:
        print(f"[_insert_dashboard DEBUG] Target slide found at index {slide_idx}")
    # Do not modify or remove the placeholder box; it serves only to locate
    # the slide.  Leave its text intact so the placeholder remains in
    # the template for future reference.
    left = Cm(left_cm)
    top = Cm(top_cm)
    width = Cm(width_cm)
    stream = io.BytesIO(image_bytes)
    print(f"[_insert_dashboard DEBUG] Inserting image: left={left_cm}cm, top={top_cm}cm, width={width_cm}cm, height={height_cm}cm")
    # Only specify width; let height auto-scale to maintain aspect ratio
    if height_cm is not None:
        pic = target_slide.shapes.add_picture(stream, left, top, width=width, height=Cm(height_cm))
    else:
        pic = target_slide.shapes.add_picture(stream, left, top, width=width)
    print(f"[_insert_dashboard DEBUG] Image inserted successfully")

    # Send to back (behind other elements like footnote)
    spTree = target_slide.shapes._spTree
    sp = pic._element
    spTree.remove(sp)
    spTree.insert(2, sp)  # Index 2 = back (0 and 1 are reserved)

    # Insert source footnote if a date is available
    if used_date is not None:
        date_str = used_date.strftime("%d/%m/%Y")
        source_text = f"Source: Bloomberg, Herculis Group. Data as of {date_str}"
        # Look for source placeholder
        source_candidates = [n.lower() for n in source_placeholder_names]
        source_patterns = [f"[{n}]" for n in source_candidates]
        for shape in target_slide.shapes:
            name_attr2 = getattr(shape, "name", "").lower()
            if name_attr2 in source_candidates:
                if shape.has_text_frame:
                    runs = shape.text_frame.paragraphs[0].runs
                    attrs = _capture_font_attrs(runs[0]) if runs else (None, None, None, None, None, None)
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    new_run = p.add_run()
                    new_run.text = source_text
                    _apply_font_attrs(new_run, *attrs)
                break
            if shape.has_text_frame:
                for pattern in source_patterns:
                    if pattern.lower() in (shape.text or "").lower():
                        runs = shape.text_frame.paragraphs[0].runs
                        attrs = _capture_font_attrs(runs[0]) if runs else (None, None, None, None, None, None)
                        try:
                            new_text = shape.text.replace(pattern, source_text)
                        except Exception:
                            new_text = source_text
                        shape.text_frame.clear()
                        p = shape.text_frame.paragraphs[0]
                        new_run = p.add_run()
                        new_run.text = new_text
                        _apply_font_attrs(new_run, *attrs)
                        break
                else:
                    continue
                break
    return prs

def insert_equity_performance_bar_slide(
    prs: Presentation,
    image_bytes: bytes,
    used_date: Optional[pd.Timestamp] = None,
    price_mode: str = "Last Price",
    *,
    left_cm: float = 3.47,
    top_cm: float = 5.28,
    width_cm: float = 17.31,
    height_cm: Optional[float] = 10.0,
) -> Presentation:
    """Insert the weekly performance bar chart and source footnote into its designated slide.

    The slide is identified by a shape named ``equity_perf_1week`` or
    ``equity_perf_1w`` (or containing ``[equity_perf_1week]`` etc.).  The
    chart is inserted at the specified coordinates regardless of the
    placeholder's original size.  A source footnote is written into
    ``equity_1w_source`` or a text box containing ``[equity_1w_source]``
    reflecting the ``used_date`` and ``price_mode``.

    Parameters
    ----------
    prs : Presentation
        The PowerPoint presentation to modify.
    image_bytes : bytes
        PNG data for the weekly bar chart.
    used_date : pandas.Timestamp or None, optional
        Effective date used for performance calculations.  If ``None``,
        the source footnote is not inserted.
    price_mode : str, default ``"Last Price"``
        Either ``"Last Price"`` or ``"Last Close"``.  Determines the
        suffix appended to the date in the source footnote.
    left_cm, top_cm, width_cm, height_cm : float
        Position and size for inserting the chart.

    Returns
    -------
    Presentation
        The modified presentation.
    """
    return _insert_dashboard_to_placeholder(
        prs,
        image_bytes,
        placeholder_names=["equity_perf_1w", "equity_perf_1week"],
        left_cm=left_cm,
        top_cm=top_cm,
        width_cm=width_cm,
        height_cm=height_cm,
        used_date=used_date,
        price_mode=price_mode,
        source_placeholder_names=["equity_1w_source"],
    )

def insert_equity_performance_histo_slide(
    prs: Presentation,
    image_bytes: bytes,
    used_date: Optional[pd.Timestamp] = None,
    price_mode: str = "Last Price",
    *,
    left_cm: float = 0.0,
    top_cm: float = 3.22,
    width_cm: float = 25.0,
    height_cm: float = 11.66,
) -> Presentation:
    """Insert the historical performance heatmap and source footnote into its designated slide.

    The slide is identified by a shape named ``equity_perf_histo`` (or
    containing ``[equity_perf_histo]``).  The heatmap is inserted at the
    specified coordinates.  A source footnote is written into
    ``equity_1w_source2`` or a text box containing ``[equity_1w_source2]``
    based on the provided ``used_date`` and ``price_mode``.

    Parameters
    ----------
    prs : Presentation
        The PowerPoint presentation to modify.
    image_bytes : bytes
        PNG data for the heatmap.
    used_date : pandas.Timestamp or None, optional
        Effective date used for performance calculations.  If ``None``,
        the source footnote is not inserted.
    price_mode : str, default ``"Last Price"``
        Either ``"Last Price"`` or ``"Last Close"``.  Determines the
        suffix appended to the date in the source footnote.
    left_cm, top_cm, width_cm, height_cm : float
        Position and size for inserting the heatmap.

    Returns
    -------
    Presentation
        The modified presentation.
    """
    return _insert_dashboard_to_placeholder(
        prs,
        image_bytes,
        placeholder_names=["equity_perf_histo"],
        left_cm=left_cm,
        top_cm=top_cm,
        width_cm=width_cm,
        height_cm=height_cm,
        used_date=used_date,
        price_mode=price_mode,
        source_placeholder_names=["equity_1w_source2"],
    )

# =============================================================================
# EQUITY YTD EVOLUTION CHART (Chart.js Line Chart)
# =============================================================================

# Index configuration with Bloomberg tickers and colors from directive
EQUITY_YTD_CONFIG = {
    "IBOV Index": {"name": "Ibov", "color": "#9333EA"},
    "MEXBOL Index": {"name": "Mexbol", "color": "#06B6D4"},
    "NKY Index": {"name": "Nikkei 225", "color": "#10B981"},
    "DAX Index": {"name": "Dax", "color": "#DC2626"},
    "SHSZ300 Index": {"name": "CSI 300", "color": "#F97316"},
    "SPX Index": {"name": "S&P 500", "color": "#3B82F6"},
    "SMI Index": {"name": "SMI", "color": "#1B3A5A"},
    "SENSEX Index": {"name": "Sensex", "color": "#22D3EE"},
    "SASEIDX Index": {"name": "TASI", "color": "#9CA3AF"},
}

# Chart dimensions (hardcoded)
EQUITY_YTD_PNG_WIDTH_PX = 2700
EQUITY_YTD_PNG_HEIGHT_PX = 1350
EQUITY_YTD_PPT_WIDTH_CM = 23.0
EQUITY_YTD_PPT_LEFT_CM = 0.5
EQUITY_YTD_PPT_TOP_CM = 4.2
EQUITY_YTD_HTML_SCALE = 3

def _get_ytd_year(data_end_date: pd.Timestamp) -> int:
    """Determine which year to show as YTD based on data end date.

    The YTD year is simply the year of the last data point.
    """
    return data_end_date.year

def _get_chart_title(data_end_date: pd.Timestamp) -> str:
    """Generate chart title with correct year."""
    year = _get_ytd_year(data_end_date)
    return f"{year} YTD Evolution"

def _compute_ytd_series(
    df: pd.DataFrame,
    ticker: str,
) -> Tuple[List[str], List[float]]:
    """Compute weekly YTD performance series for a ticker.

    Uses daily data from 01/01/YYYY, sampled to weekly (every 5th trading day).
    Baseline is the FIRST available price in the year.

    Returns:
        Tuple of (date labels, YTD performance values as cumulative returns)
    """
    if ticker not in df.columns:
        print(f"[DEBUG] {ticker} not in columns")
        return [], []

    # Get the year of the last data point
    last_date = df["Date"].max()
    year = last_date.year

    print(f"[DEBUG] {ticker}: last_date={last_date}, year={year}")

    # Filter to current year data ONLY (from 01/01/YYYY onwards)
    start_of_year = pd.Timestamp(year=year, month=1, day=1)
    df_year = df[df["Date"] >= start_of_year].copy().sort_values("Date")

    if len(df_year) == 0:
        print(f"[DEBUG] {ticker}: no data for year {year}")
        return [], []

    # Baseline = FIRST available price in the year
    first_valid_idx = df_year[ticker].first_valid_index()
    if first_valid_idx is None:
        print(f"[DEBUG] {ticker}: no valid prices in year {year}")
        return [], []

    baseline_price = df_year.loc[first_valid_idx, ticker]
    baseline_date = df_year.loc[first_valid_idx, "Date"]

    if pd.isna(baseline_price) or baseline_price == 0:
        print(f"[DEBUG] {ticker}: invalid baseline price {baseline_price}")
        return [], []

    print(f"[DEBUG] {ticker}: baseline_price={baseline_price} on {baseline_date}")

    # Compute YTD return for EVERY daily data point
    month_names = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun',
                   'Jul', 'Aug', 'Sep', 'Oct', 'Nov', 'Dec']

    daily_labels = []
    daily_values = []

    for _, row in df_year.iterrows():
        price = row[ticker]
        date = row["Date"]
        if not pd.isna(price):
            ytd_return = (price - baseline_price) / baseline_price * 100
            # Use month name as label
            month_label = month_names[date.month - 1]
            daily_labels.append(month_label)
            daily_values.append(round(ytd_return, 1))

    # If fewer than 22 data points, use daily frequency instead of weekly
    # This handles beginning of year when there's only a few days of data
    if len(daily_values) < 22:
        print(f"[DEBUG] {ticker}: {len(daily_values)} daily points (< 22), using daily frequency")
        labels = daily_labels
        values = daily_values
    else:
        # Sample to weekly (every 5th trading day)
        weekly_indices = list(range(0, len(daily_values), 5))
        # Always include the last point
        if weekly_indices[-1] != len(daily_values) - 1:
            weekly_indices.append(len(daily_values) - 1)

        labels = [daily_labels[i] for i in weekly_indices]
        values = [daily_values[i] for i in weekly_indices]

    print(f"[DEBUG] {ticker}: {len(daily_values)} daily -> {len(values)} final points, first 5={values[:5]}, last 5={values[-5:]}")
    return labels, values

def create_equity_ytd_evolution_chart(
    excel_path: Union[str, pathlib.Path],
    *,
    price_mode: str = "Last Price",
) -> Tuple[bytes, Optional[pd.Timestamp]]:
    """Generate HTML-based Equity YTD Evolution line chart.

    Features:
    - 9 equity indices with smooth lines
    - Dynamic year in title based on data
    - Label collision avoidance at end of lines
    - Emphasized zero line
    - Connector lines for adjusted labels

    Returns:
        Tuple of (PNG bytes, effective date used for computation)
    """
    # Get all tickers from configuration
    tickers = list(EQUITY_YTD_CONFIG.keys())

    # Load and adjust price data
    df = _load_price_data(excel_path, tickers)
    df_adj, used_date = adjust_prices_for_mode(df, price_mode)

    print(f"[DEBUG] Raw data date range: {df['Date'].min()} to {df['Date'].max()}, rows={len(df)}")
    print(f"[DEBUG] Adjusted data date range: {df_adj['Date'].min()} to {df_adj['Date'].max()}, rows={len(df_adj)}")
    print(f"[DEBUG] used_date from adjust_prices_for_mode: {used_date}")

    if used_date is None:
        used_date = df_adj["Date"].max()

    # Get chart title with correct year
    chart_title = _get_chart_title(used_date)

    # Compute YTD series for each index
    all_labels = []
    datasets = []

    for ticker, config in EQUITY_YTD_CONFIG.items():
        labels, values = _compute_ytd_series(df_adj, ticker)
        if labels and values:
            # Track the longest label list
            if len(labels) > len(all_labels):
                all_labels = labels

            datasets.append({
                "label": config["name"],
                "data": values,
                "borderColor": config["color"],
                "backgroundColor": "transparent",
            })

    # Ensure all datasets have the same length (pad with None if needed)
    max_len = len(all_labels)
    for dataset in datasets:
        while len(dataset["data"]) < max_len:
            dataset["data"].append(None)

    # Debug output
    print(f"[DEBUG] chart_title={chart_title}")
    print(f"[DEBUG] all_labels={all_labels} (len={len(all_labels)})")
    for ds in datasets:
        print(f"[DEBUG] {ds['label']}: data={ds['data'][:5]}... (len={len(ds['data'])})")

    # Check for insufficient data (< 3 data points)
    if len(all_labels) < 3:
        print(f"[DEBUG] Insufficient data ({len(all_labels)} points), rendering placeholder")
        env = Environment()
        template = env.from_string(YTD_INSUFFICIENT_DATA_HTML_TEMPLATE)
        html_content = template.render(
            width=EQUITY_YTD_PNG_WIDTH_PX,
            height=EQUITY_YTD_PNG_HEIGHT_PX,
            scale=EQUITY_YTD_HTML_SCALE,
            chart_title=chart_title,
        )
        # Convert placeholder HTML to PNG
        png_bytes = None
        try:
            with sync_playwright() as p:
                browser = p.chromium.launch(headless=True)
                page = browser.new_page(viewport={
                    'width': EQUITY_YTD_PNG_WIDTH_PX,
                    'height': EQUITY_YTD_PNG_HEIGHT_PX
                })
                page.set_content(html_content, wait_until='networkidle')
                page.wait_for_timeout(500)
                png_bytes = page.screenshot()
                browser.close()
        except Exception as e:
            print(f"[ERROR] Playwright failed for placeholder: {e}")
        return png_bytes, used_date

    # Compute y-axis min and max from all data values
    all_values = []
    for ds in datasets:
        all_values.extend([v for v in ds["data"] if v is not None])

    if all_values:
        data_min = min(all_values)
        data_max = max(all_values)
        # Round to nice values with some padding
        y_min = (int(data_min / 5) - 1) * 5  # Round down to nearest 5 with margin
        y_max = (int(data_max / 5) + 2) * 5  # Round up to nearest 5 with margin
        # Ensure zero is visible if data crosses zero
        if data_min < 0 < data_max:
            y_min = min(y_min, -5)
            y_max = max(y_max, 5)
    else:
        y_min = -20
        y_max = 40

    print(f"[DEBUG] y_min={y_min}, y_max={y_max}")

    # Render HTML template with proper JSON serialization
    # Create environment with tojson filter
    env = Environment()
    env.filters['tojson'] = json.dumps
    template = env.from_string(EQUITY_YTD_EVOLUTION_HTML_TEMPLATE)
    html_content = template.render(
        width=EQUITY_YTD_PNG_WIDTH_PX,
        height=EQUITY_YTD_PNG_HEIGHT_PX,
        scale=EQUITY_YTD_HTML_SCALE,
        chart_title=chart_title,
        labels=all_labels,
        datasets=datasets,
        y_min=y_min,
        y_max=y_max,
    )

    # Convert HTML to PNG using Playwright (waits for Chart.js to render)
    png_bytes = None
    try:
        with sync_playwright() as p:
            browser = p.chromium.launch(headless=True)
            page = browser.new_page(viewport={
                'width': EQUITY_YTD_PNG_WIDTH_PX,
                'height': EQUITY_YTD_PNG_HEIGHT_PX
            })

            # Load HTML content - use data URL to ensure proper loading
            page.set_content(html_content, wait_until='networkidle')

            # Wait for Chart.js to signal rendering complete
            try:
                page.wait_for_selector('body[data-chart-ready="true"]', timeout=10000)
                print("[DEBUG] Chart.js rendering complete")
            except Exception as e:
                print(f"[DEBUG] Chart ready selector timeout, using fallback wait: {e}")
                page.wait_for_timeout(3000)

            # Take screenshot
            png_bytes = page.screenshot()
            print(f"[DEBUG] Screenshot taken, size: {len(png_bytes)} bytes")
            browser.close()
    except Exception as e:
        print(f"[ERROR] Playwright failed: {e}")
        import traceback
        traceback.print_exc()

    return png_bytes, used_date

def insert_equity_ytd_evolution_slide(
    prs: Presentation,
    image_bytes: bytes,
    used_date: Optional[pd.Timestamp] = None,
    price_mode: str = "Last Price",
    subtitle: Optional[str] = None,
) -> Presentation:
    """Insert the Equity YTD Evolution chart into PowerPoint.

    Parameters
    ----------
    prs : Presentation
        The PowerPoint presentation to modify.
    image_bytes : bytes
        PNG data for the YTD evolution chart.
    used_date : pandas.Timestamp or None, optional
        Effective date used for performance calculations.
    price_mode : str, default "Last Price"
        Either "Last Price" or "Last Close".
    subtitle : str or None, optional
        Subtitle text to insert into the ytd_eq_subtitle placeholder.
    """
    if not image_bytes:
        return prs

    # Find target slide by placeholder name
    target_slide = None
    placeholder_names = ["ytd_eq_perf"]
    name_candidates = [n.lower() for n in placeholder_names]
    pattern_candidates = [f"[{n}]" for n in name_candidates]

    for slide in prs.slides:
        for shape in slide.shapes:
            name_attr = getattr(shape, "name", "").lower()
            if name_attr in name_candidates:
                target_slide = slide
                break
            if shape.has_text_frame:
                text_lower = (shape.text or "").strip().lower()
                if text_lower in [p.lower() for p in pattern_candidates]:
                    target_slide = slide
                    break
        if target_slide:
            break

    if target_slide is None:
        print("[Equity YTD Evolution] ERROR: Slide not found")
        return prs

    # ------------------------------------------------------------------
    # Insert subtitle into the 'ytd_eq_subtitle' placeholder
    # ------------------------------------------------------------------
    if subtitle:
        for shape in target_slide.shapes:
            name_attr = getattr(shape, "name", "")
            if name_attr and name_attr.lower() == "ytd_eq_subtitle" and shape.has_text_frame:
                tf = shape.text_frame
                paragraph = tf.paragraphs[0]
                runs = paragraph.runs
                attrs = _capture_font_attrs(runs[0]) if runs else (None, None, None, None, None, None)
                # Determine replacement: replace tokens "XXX" or "[ytd_eq_subtitle]"
                original_text = "".join(run.text for run in runs) if runs else ""
                new_text = original_text
                if "XXX" in new_text:
                    new_text = new_text.replace("XXX", subtitle)
                elif "[ytd_eq_subtitle]" in new_text:
                    new_text = new_text.replace("[ytd_eq_subtitle]", subtitle)
                else:
                    new_text = subtitle
                # Clear and insert new run
                tf.clear()
                p = tf.paragraphs[0]
                new_run = p.add_run()
                new_run.text = new_text
                _apply_font_attrs(new_run, *attrs)
                break

    # ------------------------------------------------------------------
    # Insert chart image with exact hardcoded dimensions
    # ------------------------------------------------------------------
    left = Cm(EQUITY_YTD_PPT_LEFT_CM)
    top = Cm(EQUITY_YTD_PPT_TOP_CM)
    width = Cm(EQUITY_YTD_PPT_WIDTH_CM)

    stream = io.BytesIO(image_bytes)
    picture = target_slide.shapes.add_picture(stream, left, top, width=width)

    # Send picture to back
    spTree = target_slide.shapes._spTree
    pic_element = picture._element
    spTree.remove(pic_element)
    spTree.insert(2, pic_element)

    print(f"[Equity YTD Evolution] Chart inserted at ({EQUITY_YTD_PPT_LEFT_CM}, {EQUITY_YTD_PPT_TOP_CM}) cm")

    # ------------------------------------------------------------------
    # Insert data source footnote
    # ------------------------------------------------------------------
    if used_date is not None:
        date_str = used_date.strftime("%d/%m/%Y")
        source_text = f"Source: Bloomberg, Herculis Group. Data as of {date_str}"
        placeholder_name = "ytd_eq_source"
        placeholder_patterns = ["[ytd_eq_source]", "ytd_eq_source"]
        inserted = False
        for shape in target_slide.shapes:
            name_attr = getattr(shape, "name", "")
            if name_attr and name_attr.lower() == placeholder_name:
                if shape.has_text_frame:
                    runs = shape.text_frame.paragraphs[0].runs
                    attrs = _capture_font_attrs(runs[0]) if runs else (None, None, None, None, None, None)
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    new_run = p.add_run()
                    new_run.text = source_text
                    _apply_font_attrs(new_run, *attrs)
                inserted = True
                break
            if shape.has_text_frame:
                for pattern in placeholder_patterns:
                    if pattern.lower() in (shape.text or "").lower():
                        runs = shape.text_frame.paragraphs[0].runs
                        attrs = _capture_font_attrs(runs[0]) if runs else (None, None, None, None, None, None)
                        try:
                            new_text = shape.text.replace(pattern, source_text)
                        except Exception:
                            new_text = source_text
                        shape.text_frame.clear()
                        p = shape.text_frame.paragraphs[0]
                        new_run = p.add_run()
                        new_run.text = new_text
                        _apply_font_attrs(new_run, *attrs)
                        inserted = True
                        break
                if inserted:
                    break

    return prs


# =============================================================================
# FX IMPACT ANALYSIS (EUR) - YTD Performance Decomposition
# =============================================================================

# Index configuration with local currency and FX pair for EUR conversion
FX_IMPACT_EUR_CONFIG = {
    "SPX Index": {
        "name": "S&P 500",
        "flag": "us",
        "currency": "USD",
        "fx_pair": "EURUSD Curncy",  # EUR per USD
    },
    "DAX Index": {
        "name": "Dax",
        "flag": "de",
        "currency": "EUR",
        "fx_pair": None,  # Already in EUR
    },
    "SMI Index": {
        "name": "SMI",
        "flag": "ch",
        "currency": "CHF",
        "fx_pair": "EURCHF Curncy",  # EUR per CHF
    },
    "NKY Index": {
        "name": "Nikkei 225",
        "flag": "jp",
        "currency": "JPY",
        "fx_pair": "EURJPY Curncy",  # EUR per JPY
    },
    "SHSZ300 Index": {
        "name": "CSI 300",
        "flag": "cn",
        "currency": "CNH",
        "fx_pair": "EURCNH Curncy",  # EUR per CNH
    },
    "SENSEX Index": {
        "name": "Sensex",
        "flag": "in",
        "currency": "INR",
        "fx_pair": "EURINR Curncy",  # EUR per INR
    },
    "IBOV Index": {
        "name": "Bovespa",
        "flag": "br",
        "currency": "BRL",
        "fx_pair": "EURBRL Curncy",  # EUR per BRL
    },
    "MEXBOL Index": {
        "name": "Mexbol",
        "flag": "mx",
        "currency": "MXN",
        "fx_pair": "EURMXN Curncy",  # EUR per MXN
    },
    "SASEIDX Index": {
        "name": "TASI",
        "flag": "sa",
        "currency": "SAR",
        "fx_pair": "EURSAR Curncy",  # EUR per SAR
    },
}

# Chart dimensions
FX_IMPACT_PNG_WIDTH_PX = 2400
FX_IMPACT_PNG_HEIGHT_PX = 1400
FX_IMPACT_HTML_SCALE = 3
FX_IMPACT_PPT_WIDTH_CM = 20.0
FX_IMPACT_PPT_LEFT_CM = 2.5
FX_IMPACT_PPT_TOP_CM = 4.0


def _compute_ytd_return(
    df: pd.DataFrame,
    ticker: str,
) -> float:
    """Compute YTD return for a ticker.

    Returns:
        YTD return as a percentage, or NaN if not computable.
    """
    if ticker not in df.columns:
        return float("nan")

    last_date = df["Date"].max()
    year = last_date.year

    # Start of year baseline
    start_of_year = pd.Timestamp(year=year, month=1, day=1)
    past_series = df.loc[df["Date"] <= start_of_year, ticker]

    if len(past_series) == 0:
        return float("nan")

    baseline_price = past_series.iloc[-1]
    current_price = df[ticker].iloc[-1]

    if pd.isna(baseline_price) or baseline_price == 0:
        return float("nan")

    return (current_price - baseline_price) / baseline_price * 100.0


def _generate_fx_insight_text(rows_data: list) -> str:
    """Generate insight text highlighting biggest FX tailwind and headwind.

    Parameters
    ----------
    rows_data : list
        List of row dictionaries with 'name', 'fx_effect', 'local_return', 'eur_return' keys.

    Returns
    -------
    str
        HTML-formatted insight text.
    """
    if not rows_data:
        return "Insufficient data to generate insights."

    # Find biggest tailwind (most positive FX effect) and headwind (most negative)
    sorted_by_fx = sorted(rows_data, key=lambda x: x.get("fx_effect", 0), reverse=True)

    best = sorted_by_fx[0] if sorted_by_fx else None
    worst = sorted_by_fx[-1] if sorted_by_fx else None

    parts = []

    if best and best.get("fx_effect", 0) > 0:
        # Determine currency from flag
        currency_map = {
            "us": "USD strength",
            "jp": "JPY movement",
            "ch": "CHF movement",
            "cn": "CNY movement",
            "br": "BRL movement",
            "mx": "MXN movement",
            "in": "INR movement",
            "sa": "SAR movement",
        }
        currency = currency_map.get(best.get("flag", ""), "Currency movement")
        local_ret = best.get("local_return", 0)
        eur_ret = best.get("eur_return", 0)
        fx_eff = best.get("fx_effect", 0)

        local_sign = "+" if local_ret >= 0 else ""
        eur_sign = "+" if eur_ret >= 0 else ""

        parts.append(
            f'<span class="highlight-positive">{currency}</span> has been the dominant FX tailwind YTD, '
            f'turning a {local_sign}{local_ret:.1f}% {best["name"]} local return into a '
            f'{eur_sign}{eur_ret:.1f}% gain for EUR investors (<span class="highlight-positive">+{fx_eff:.1f}% FX contribution</span>).'
        )

    if worst and worst.get("fx_effect", 0) < 0:
        currency_map = {
            "us": "USD weakness",
            "jp": "JPY weakness",
            "ch": "CHF weakness",
            "cn": "CNY weakness",
            "br": "BRL weakness",
            "mx": "MXN weakness",
            "in": "INR weakness",
            "sa": "SAR weakness",
        }
        currency = currency_map.get(worst.get("flag", ""), "Currency weakness")
        fx_eff = worst.get("fx_effect", 0)

        parts.append(
            f'Conversely, <span class="highlight-negative">{currency}</span> has eroded {worst["name"]} gains by '
            f'<span class="highlight-negative">{abs(fx_eff):.1f}%</span> for EUR-based investors.'
        )

    if not parts:
        return "FX effects have been relatively balanced across major indices YTD."

    return " ".join(parts)


def create_fx_impact_analysis_chart_eur(
    excel_path: Union[str, pathlib.Path],
    *,
    price_mode: str = "Last Price",
) -> Tuple[bytes, Optional[pd.Timestamp]]:
    """Generate FX Impact Analysis chart for EUR investors.

    Shows YTD performance decomposition: Local return, EUR return, and FX effect
    with visual bars showing FX tailwind (green) or headwind (red).

    Uses compounded calculation:
    - EUR_Return = (price_end / price_start) * (EURXXX_start / EURXXX_end) - 1
    - FX_Effect = EUR_Return - Local_Return

    Parameters
    ----------
    excel_path : str or pathlib.Path
        Path to the Excel workbook containing price data.
    price_mode : str, default "Last Price"
        Either "Last Price" or "Last Close".

    Returns
    -------
    tuple
        (PNG bytes, effective date used for computation)
    """
    print(f"[FX Impact EUR DEBUG] Starting chart generation...")
    print(f"[FX Impact EUR DEBUG] Excel path: {excel_path}")
    print(f"[FX Impact EUR DEBUG] Price mode: {price_mode}")

    # First, load the Excel to see which tickers are available
    df_raw = pd.read_excel(excel_path, sheet_name="data_prices")
    available_columns = set(df_raw.columns)
    print(f"[FX Impact EUR DEBUG] Available columns count: {len(available_columns)}")

    # Collect all tickers needed (indices + FX pairs)
    tickers_to_load = []
    for index_ticker, config in FX_IMPACT_EUR_CONFIG.items():
        if index_ticker in available_columns:
            tickers_to_load.append(index_ticker)
            print(f"[FX Impact EUR DEBUG] Found index: {index_ticker}")
        else:
            print(f"[FX Impact EUR DEBUG] MISSING index: {index_ticker}")
        fx_pair = config.get("fx_pair")
        if fx_pair and fx_pair in available_columns:
            tickers_to_load.append(fx_pair)
            print(f"[FX Impact EUR DEBUG] Found FX pair: {fx_pair}")
        elif fx_pair:
            print(f"[FX Impact EUR DEBUG] MISSING FX pair: {fx_pair}")

    tickers_to_load = list(set(tickers_to_load))
    print(f"[FX Impact EUR DEBUG] Total tickers to load: {len(tickers_to_load)}")

    if not tickers_to_load:
        print("[FX Impact Analysis EUR] ERROR: No valid tickers found in data")
        return None, None

    # Load and adjust price data
    df = _load_price_data(excel_path, tickers_to_load)
    df_adj, used_date = adjust_prices_for_mode(df, price_mode)
    print(f"[FX Impact EUR DEBUG] Data loaded. Rows: {len(df_adj)}, Used date: {used_date}")

    # Get YTD start date
    last_date = df_adj["Date"].max()
    year = last_date.year
    start_of_year = pd.Timestamp(year=year, month=1, day=1)
    print(f"[FX Impact EUR DEBUG] Last date: {last_date}, Year: {year}, Start of year: {start_of_year}")

    # Build row data
    rows_data = []
    for index_ticker, config in FX_IMPACT_EUR_CONFIG.items():
        print(f"[FX Impact EUR DEBUG] Processing {config['name']} ({index_ticker})...")

        # Get price at start and end of year
        if index_ticker in df_adj.columns:
            past_series = df_adj.loc[df_adj["Date"] <= start_of_year, index_ticker]
            if len(past_series) > 0:
                price_start = past_series.iloc[-1]
            else:
                price_start = float("nan")
            price_end = df_adj[index_ticker].iloc[-1]

            if pd.notna(price_start) and price_start != 0 and pd.notna(price_end):
                local_return = (price_end / price_start - 1) * 100
            else:
                local_return = 0.0
            print(f"[FX Impact EUR DEBUG]   Price: {price_start:.2f} -> {price_end:.2f}, Local return: {local_return:.2f}%")
        else:
            local_return = 0.0
            price_start = float("nan")
            price_end = float("nan")
            print(f"[FX Impact EUR DEBUG]   Index not in data, local_return = 0")

        # Compute EUR return using compounded FX calculation
        fx_pair = config.get("fx_pair")
        if fx_pair is None:
            # Already in EUR (e.g., DAX)
            eur_return = local_return
            fx_effect = 0.0
            print(f"[FX Impact EUR DEBUG]   No FX pair (EUR asset), EUR return = local return")
        elif fx_pair in df_adj.columns:
            # Get FX rate at start and end of year
            fx_past_series = df_adj.loc[df_adj["Date"] <= start_of_year, fx_pair]
            if len(fx_past_series) > 0:
                fx_start = fx_past_series.iloc[-1]
            else:
                fx_start = float("nan")
            fx_end = df_adj[fx_pair].iloc[-1]

            if pd.notna(fx_start) and fx_start != 0 and pd.notna(fx_end) and fx_end != 0:
                # Compounded EUR return:
                # EUR_Return = (price_end / price_start) * (fx_start / fx_end) - 1
                # For EURXXX pairs, if EURXXX falls (foreign currency strengthens vs EUR),
                # fx_start / fx_end > 1, which is positive for EUR investors
                if pd.notna(price_start) and price_start != 0 and pd.notna(price_end):
                    eur_return = ((price_end / price_start) * (fx_start / fx_end) - 1) * 100
                else:
                    eur_return = 0.0
                fx_effect = eur_return - local_return
                print(f"[FX Impact EUR DEBUG]   FX: {fx_start:.4f} -> {fx_end:.4f}, EUR return: {eur_return:.2f}%, FX effect: {fx_effect:.2f}%")
            else:
                eur_return = local_return
                fx_effect = 0.0
                print(f"[FX Impact EUR DEBUG]   FX data invalid, EUR return = local return")
        else:
            # FX pair not available
            eur_return = local_return
            fx_effect = 0.0
            print(f"[FX Impact EUR DEBUG]   FX pair {fx_pair} not in data, EUR return = local return")

        rows_data.append({
            "name": config["name"],
            "flag": config["flag"],
            "local_return": local_return,
            "eur_return": eur_return,
            "fx_effect": fx_effect,
        })

    print(f"[FX Impact EUR DEBUG] Built {len(rows_data)} rows")

    # Sort by FX effect descending (biggest tailwind first)
    rows_data.sort(key=lambda x: x["fx_effect"], reverse=True)

    # Calculate max absolute FX effect for bar scaling
    max_abs_fx = max((abs(r["fx_effect"]) for r in rows_data), default=5.0)
    max_abs_fx = max(max_abs_fx, 2.0)  # At least 2%
    print(f"[FX Impact EUR DEBUG] Max absolute FX effect: {max_abs_fx:.2f}%")

    # Prepare rows for template
    prepared_rows = []
    for i, row in enumerate(rows_data):
        fx_effect = row["fx_effect"]
        bar_width = abs(fx_effect) / max_abs_fx * 48
        bar_width = min(bar_width, 48)

        # Determine highlight class
        highlight_class = ""
        if i == 0 and fx_effect > 0:
            highlight_class = "fx-tailwind"
        elif i == len(rows_data) - 1 and fx_effect < 0:
            highlight_class = "fx-headwind"

        # Format values
        local_sign = "+" if row["local_return"] >= 0 else ""
        eur_sign = "+" if row["eur_return"] >= 0 else ""
        fx_sign = "+" if fx_effect >= 0 else ""

        prepared_rows.append({
            "name": row["name"],
            "flag": row["flag"],
            "flag_html": get_flag_html(row["flag"]),
            "local_formatted": f"{local_sign}{row['local_return']:.1f}%",
            "eur_formatted": f"{eur_sign}{row['eur_return']:.1f}%",
            "fx_formatted": f"{fx_sign}{fx_effect:.1f}%",
            "fx_effect": fx_effect,
            "fx_class": "positive" if fx_effect >= 0 else "negative",
            "bar_class": "positive" if fx_effect >= 0 else "negative",
            "bar_width": bar_width,
            "highlight_class": highlight_class,
        })

    # Generate insight text
    insight_text = _generate_fx_insight_text(rows_data)
    print(f"[FX Impact EUR DEBUG] Insight text generated: {insight_text[:100]}...")

    # Calculate scale values
    scale_max = int(max_abs_fx) + 2
    scale_values = {
        "scale_min": f"-{scale_max}%",
        "scale_mid_low": f"-{scale_max // 2}%",
        "scale_mid_high": f"+{scale_max // 2}%",
        "scale_max": f"+{scale_max}%",
    }

    # Render HTML template
    print(f"[FX Impact EUR DEBUG] Rendering HTML template...")
    template = Template(FX_IMPACT_ANALYSIS_EUR_HTML_TEMPLATE)
    html_content = template.render(
        width=FX_IMPACT_PNG_WIDTH_PX,
        height=FX_IMPACT_PNG_HEIGHT_PX,
        scale=FX_IMPACT_HTML_SCALE,
        rows=prepared_rows,
        insight_text=insight_text,
        **scale_values,
    )
    print(f"[FX Impact EUR DEBUG] HTML generated, length: {len(html_content)} chars")

    # Convert HTML to PNG using Playwright
    print(f"[FX Impact EUR DEBUG] Starting Playwright rendering ({FX_IMPACT_PNG_WIDTH_PX}x{FX_IMPACT_PNG_HEIGHT_PX}px)...")
    png_bytes = None
    try:
        with sync_playwright() as p:
            print(f"[FX Impact EUR DEBUG] Launching browser...")
            browser = p.chromium.launch(headless=True)
            page = browser.new_page(viewport={
                'width': FX_IMPACT_PNG_WIDTH_PX,
                'height': FX_IMPACT_PNG_HEIGHT_PX
            })
            print(f"[FX Impact EUR DEBUG] Setting page content...")
            page.set_content(html_content, wait_until='networkidle')
            page.wait_for_timeout(500)
            print(f"[FX Impact EUR DEBUG] Taking screenshot...")
            png_bytes = page.screenshot()
            browser.close()
            print(f"[FX Impact EUR DEBUG] Screenshot taken, size: {len(png_bytes)} bytes")
    except Exception as e:
        print(f"[FX Impact Analysis EUR] ERROR rendering: {e}")
        import traceback
        traceback.print_exc()

    if png_bytes:
        print(f"[FX Impact EUR DEBUG] SUCCESS: Generated chart {FX_IMPACT_PNG_WIDTH_PX}x{FX_IMPACT_PNG_HEIGHT_PX}px, {len(png_bytes)} bytes")
    else:
        print(f"[FX Impact EUR DEBUG] FAILED: No PNG bytes generated")

    return png_bytes, used_date


def insert_fx_impact_analysis_slide_eur(
    prs: Presentation,
    image_bytes: bytes,
    used_date: Optional[pd.Timestamp] = None,
    price_mode: str = "Last Price",
    *,
    left_cm: float = 2.5,
    top_cm: float = 4.0,
    width_cm: float = 20.0,
    height_cm: Optional[float] = None,
) -> Presentation:
    """Insert the FX Impact Analysis (EUR) chart into PowerPoint.

    The slide is identified by a shape named ``equity_perf_fx_eur`` (or
    containing ``[equity_perf_fx_eur]``). The chart is inserted at the
    specified coordinates.

    Parameters
    ----------
    prs : Presentation
        The PowerPoint presentation to modify.
    image_bytes : bytes
        PNG data for the FX impact chart.
    used_date : pandas.Timestamp or None, optional
        Effective date used for calculations.
    price_mode : str, default "Last Price"
        Either "Last Price" or "Last Close".
    left_cm, top_cm, width_cm, height_cm : float
        Position and size for inserting the chart.

    Returns
    -------
    Presentation
        The modified presentation.
    """
    print(f"[FX Impact EUR DEBUG] insert_fx_impact_analysis_slide_eur called")
    print(f"[FX Impact EUR DEBUG] image_bytes: {len(image_bytes) if image_bytes else 'None'} bytes")
    print(f"[FX Impact EUR DEBUG] used_date: {used_date}, price_mode: {price_mode}")

    if not image_bytes:
        print("[FX Impact Analysis EUR] ERROR: No image bytes to insert")
        return prs

    print(f"[FX Impact EUR DEBUG] Calling _insert_dashboard_to_placeholder with placeholder_names=['equity_perf_fx_eur']")
    return _insert_dashboard_to_placeholder(
        prs,
        image_bytes,
        placeholder_names=["equity_perf_fx_eur"],
        left_cm=left_cm,
        top_cm=top_cm,
        width_cm=width_cm,
        height_cm=height_cm,
        used_date=used_date,
        price_mode=price_mode,
        source_placeholder_names=["fx_impact_eur_source"],
    )
