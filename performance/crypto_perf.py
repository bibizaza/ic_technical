# === Windows Playwright Fix - MUST BE FIRST ===
import sys
import asyncio
if sys.platform == 'win32':
    asyncio.set_event_loop_policy(asyncio.WindowsProactorEventLoopPolicy())

try:
    import nest_asyncio
    nest_asyncio.apply()
except ImportError:
    pass
# === End Fix ===

"""Cryptocurrency performance dashboard generation with price‑mode awareness and source footnotes.

This module produces charts summarising recent performance for a selection
of major cryptocurrencies.  Users can choose to base calculations on
either the most recent intraday price ("Last Price") or the previous
day's closing price ("Last Close").  The resulting figures include a
weekly returns bar chart and a heatmap of longer horizons.  When the
charts are inserted into a PowerPoint presentation, a source footnote
is added to the designated text boxes on each slide, reflecting the
effective date used in the computations and the chosen price mode.
Text formatting from the placeholders is preserved.

Key functions
-------------

* ``create_weekly_performance_chart`` – Build the weekly bar chart and
  return both the image bytes and the effective date used for
  computation.
* ``create_historical_performance_table`` – Build the heatmap of
  returns for multiple horizons and return the image bytes and the
  effective date.
* ``insert_crypto_performance_bar_slide`` – Insert the weekly bar
  chart and source footnote into its slide.
* ``insert_crypto_performance_histo_slide`` – Insert the historical
  performance heatmap and source footnote into its slide.

Usage example::

    from performance.crypto_perf import (
        create_weekly_performance_chart,
        create_historical_performance_table,
        insert_crypto_performance_bar_slide,
        insert_crypto_performance_histo_slide,
    )

    # Generate charts with price-mode awareness
    bar_bytes, used_date = create_weekly_performance_chart(excel_path, price_mode="Last Close")
    histo_bytes, _ = create_historical_performance_table(excel_path, price_mode="Last Close")
    # Insert into PPT, supplying used_date and price_mode for source footnote
    prs = insert_crypto_performance_bar_slide(prs, bar_bytes, used_date, price_mode)
    prs = insert_crypto_performance_histo_slide(prs, histo_bytes, used_date, price_mode)

"""

from __future__ import annotations

import io
import json
import pathlib
from typing import Dict, List, Tuple, Union, Optional

import matplotlib.pyplot as plt
import matplotlib.colors as mcolors
import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.util import Cm
from jinja2 import Template, Environment
from html2image import Html2Image
from playwright.sync_api import sync_playwright

from utils import adjust_prices_for_mode
from market_compass.weekly_performance.html_template import (
    CRYPTO_WEEKLY_HTML_TEMPLATE,
    CRYPTO_HISTORICAL_HTML_TEMPLATE,
    CRYPTO_YTD_EVOLUTION_HTML_TEMPLATE,
    YTD_INSUFFICIENT_DATA_HTML_TEMPLATE,
)

try:
    # Reuse font attribute helpers from the SPX module if available
    from technical_analysis.equity.spx import (
        _get_run_font_attributes as _capture_font_attrs,  # type: ignore
        _apply_run_font_attributes as _apply_font_attrs,  # type: ignore
    )
except Exception:
    # Fallback helpers: only support basic size and RGB; ignore theme colours
    def _capture_font_attrs(run):  # type: ignore
        if run is None:
            return (None, None, None, None, None, None)
        size = run.font.size
        colour = run.font.color
        rgb = None
        try:
            rgb = colour.rgb
        except Exception:
            rgb = None
        return (size, rgb, None, None, run.font.bold, run.font.italic)

    def _apply_font_attrs(new_run, size, rgb, theme_color, brightness, bold, italic):  # type: ignore
        if size is not None:
            new_run.font.size = size
        if rgb is not None:
            try:
                new_run.font.color.rgb = rgb
            except Exception:
                pass
        if bold is not None:
            new_run.font.bold = bold
        if italic is not None:
            new_run.font.italic = italic

###############################################################################
# Data loading and return calculations
###############################################################################

def _load_price_data(excel_path: Union[str, pathlib.Path], tickers: List[str]) -> pd.DataFrame:
    """Load a DataFrame of dates and prices for the specified cryptocurrency tickers.

    Parameters
    ----------
    excel_path : str or pathlib.Path
        Path to the Excel workbook containing a sheet named ``data_prices``.
    tickers : list of str
        Column names in the sheet corresponding to the desired crypto indices.

    Returns
    -------
    DataFrame
        A tidy DataFrame with columns ``Date`` and one column per ticker,
        sorted by date and with numeric prices.
    """
    df = pd.read_excel(excel_path, sheet_name="data_prices")
    # Drop the first row which usually contains metadata and rows with the
    # sentinel ``DATES`` value.
    df = df.drop(index=0)
    df = df[df[df.columns[0]] != "DATES"].copy()
    df = df.rename(columns={df.columns[0]: "Date"})
    df["Date"] = pd.to_datetime(df["Date"], errors="coerce")
    # Filter tickers to only include those that actually exist in the Excel file
    available_tickers = [t for t in tickers if t in df.columns]
    missing_tickers = [t for t in tickers if t not in df.columns]
    if missing_tickers:
        print(f"Warning: The following tickers are not in the Excel file and will be skipped: {missing_tickers}")
    # Restrict to requested tickers and coerce to numeric
    out = df[["Date"] + available_tickers].copy()
    for t in available_tickers:
        out[t] = pd.to_numeric(out[t], errors="coerce")
    return out.dropna(subset=["Date"]).sort_values("Date").reset_index(drop=True)


def _compute_horizon_returns(
    df: pd.DataFrame,
    ticker: str,
    today: pd.Timestamp,
    days: int,
) -> float:
    """Compute the percentage return of a series over the given lookback.

    This helper finds the last available price on or before ``today`` minus
    ``days`` and compares it to the most recent price.  If insufficient
    history exists, NaN is returned.

    Parameters
    ----------
    df : DataFrame
        A tidy DataFrame with ``Date`` and price columns.
    ticker : str
        Column name for which to compute the return.
    today : Timestamp
        Anchor date; typically the last date in the dataset.
    days : int
        Lookback period in days.

    Returns
    -------
    float
        Percent return over the lookback horizon, or NaN if not computable.
    """
    past_date = today - pd.Timedelta(days=days)
    past_series = df.loc[df["Date"] <= past_date, ticker]
    if len(past_series) == 0:
        return float("nan")
    past_price = past_series.iloc[-1]
    current_price = df[ticker].iloc[-1]
    if past_price == 0 or pd.isna(past_price):
        return float("nan")
    return (current_price - past_price) / past_price * 100.0


def _build_returns_table(
    df: pd.DataFrame,
    mapping: Dict[str, str],
) -> pd.DataFrame:
    """Calculate performance metrics for each cryptocurrency ticker.

    Returns a DataFrame indexed by the ticker code with the following
    columns: ``Name``, ``1W``, ``1M``, ``3M``, ``6M``, ``12M`` and ``YTD``.

    Parameters
    ----------
    df : DataFrame
        A tidy DataFrame as returned by ``_load_price_data``.
    mapping : dict
        Mapping from ticker codes to human‑friendly names.

    Returns
    -------
    DataFrame
        Performance table.
    """
    today = df["Date"].max()
    results: Dict[str, Dict[str, float]] = {}
    for ticker in mapping:
        results[ticker] = {
            "1W": _compute_horizon_returns(df, ticker, today, 7),
            "1M": _compute_horizon_returns(df, ticker, today, 30),
            "3M": _compute_horizon_returns(df, ticker, today, 90),
            "6M": _compute_horizon_returns(df, ticker, today, 180),
            "12M": _compute_horizon_returns(df, ticker, today, 365),
        }
        # Year‑to‑date: look back to the start of the calendar year
        start_of_year = pd.Timestamp(year=today.year, month=1, day=1)
        past_series = df.loc[df["Date"] <= start_of_year, ticker]
        if len(past_series) > 0:
            past_price = past_series.iloc[-1]
            current_price = df[ticker].iloc[-1]
            if past_price and not pd.isna(past_price):
                results[ticker]["YTD"] = (current_price - past_price) / past_price * 100.0
            else:
                results[ticker]["YTD"] = float("nan")
        else:
            results[ticker]["YTD"] = float("nan")
    table = pd.DataFrame.from_dict(results, orient="index")
    table["Name"] = table.index.map(mapping.get)
    return table[["Name", "1W", "1M", "3M", "6M", "12M", "YTD"]]


###############################################################################
# Figure generation functions
###############################################################################

def create_weekly_performance_chart(
    excel_path: Union[str, pathlib.Path],
    ticker_mapping: Dict[str, str] | None = None,
    *,
    width: float = 14.0,
    height: float = 5.0,
    price_mode: str = "Last Price",
) -> Tuple[bytes, Optional[pd.Timestamp]]:
    """Generate a bar chart of 1‑week cryptocurrency returns with price‑mode adjustment.

    This helper reads the specified Excel file, optionally adjusts the price
    history according to the selected ``price_mode`` (``"Last Price"`` vs
    ``"Last Close"``), computes 1‑week returns for the configured crypto
    tickers, sorts them in descending order and constructs a horizontal
    bar chart.  The effective date used for the calculations is returned
    along with the PNG data.

    Parameters
    ----------
    excel_path : str or pathlib.Path
        Path to the Excel workbook containing the price series.
    ticker_mapping : dict, optional
        Mapping from ticker codes to display names.  If not provided,
        sensible defaults are used.
    width, height : float
        Dimensions of the generated figure in inches.
    price_mode : str, default ``"Last Price"``
        One of ``"Last Price"`` or ``"Last Close"``.  Determines whether
        intraday prices are used or the most recent row is dropped if it
        corresponds to today's date.

    Returns
    -------
    tuple
        A two‑tuple ``(image_bytes, used_date)`` where ``image_bytes`` is
        the PNG data for the bar chart and ``used_date`` is the effective
        date in the adjusted price series.
    """
    default_mapping = {
        "XBTUSD Curncy": "Bitcoin",
        "XETUSD Curncy": "Ethereum",
        "XRPUSD Curncy": "Ripple",
        "XSOUSD Curncy": "Solana",
        "XBIUSD Curncy": "Binance",
        "XDOUSD Curncy": "Polkadot",
        "XVV Curncy": "AAVE",
        "CHYPEE Index":"HyperLiquid",
        "CCTON Curncy": "Ton",
        "BGCI Index": "Bloomberg Galaxy Crypto",
    }
    mapping = ticker_mapping or default_mapping
    tickers = list(mapping.keys())
    # Load data and adjust according to price mode
    df = _load_price_data(excel_path, tickers)
    df_adj, used_date = adjust_prices_for_mode(df, price_mode)
    perf = _build_returns_table(df_adj, mapping)
    # Build bar chart data sorted by performance descending
    bar_df = perf[["Name", "1W"]].dropna().sort_values("1W", ascending=False).reset_index(drop=True)
    # Create figure and axis
    fig, ax = plt.subplots(figsize=(width, height))
    bar_colors = ["#2E8B57" if x > 0 else "#C70039" for x in bar_df["1W"]]
    bars = ax.barh(bar_df["Name"], bar_df["1W"], color=bar_colors)
    # Add labels
    for bar in bars:
        width_val = bar.get_width()
        x_pos = width_val + (0.1 if width_val > 0 else -0.1)
        ax.text(
            x_pos,
            bar.get_y() + bar.get_height() / 2.0,
            f"{width_val:.1f}%",
            va="center",
            ha="left" if width_val > 0 else "right",
            fontweight="bold",
            color="black",
        )
    # Style chart
    ax.axvline(0.0, color="grey", linewidth=0.8, linestyle="--")
    ax.set_xticks([])
    ax.set_yticks(range(len(bar_df)))
    ax.set_yticklabels(bar_df["Name"], fontsize=10)
    ax.invert_yaxis()
    for spine in ax.spines.values():
        spine.set_visible(False)
    ax.tick_params(axis="y", length=0)
    # Provide extra horizontal space so bars do not overlap names
    min_val = min(bar_df["1W"].min(), 0)
    max_val = bar_df["1W"].max()
    margin = (max_val - min_val) * 0.2
    ax.set_xlim(min_val - margin, max_val + margin)
    fig.subplots_adjust(left=0.4)
    # Export to PNG
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=300, bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf.getvalue(), used_date


def create_historical_performance_table(
    excel_path: Union[str, pathlib.Path],
    ticker_mapping: Dict[str, str] | None = None,
    *,
    width: float = 14.0,
    height: float = 6.0,
    price_mode: str = "Last Price",
) -> Tuple[bytes, Optional[pd.Timestamp]]:
    """Generate a heatmap table of cryptocurrency returns with price‑mode adjustment.

    This helper reads the price data, optionally adjusts it according to
    ``price_mode``, computes returns for multiple horizons, sorts the
    table by YTD performance (descending) and constructs a heatmap where
    each column is coloured independently.  The effective date used
    for the computations is returned along with the PNG data.

    Parameters
    ----------
    excel_path : str or pathlib.Path
        Path to the Excel workbook containing the price series.
    ticker_mapping : dict, optional
        Mapping from ticker codes to display names.  If not provided,
        sensible defaults are used.
    width, height : float
        Dimensions of the generated figure in inches.
    price_mode : str, default ``"Last Price"``
        Either ``"Last Price"`` or ``"Last Close"``.  Determines how
        price data is adjusted prior to computing returns and which
        effective date appears in the source footnote.

    Returns
    -------
    tuple
        A two‑tuple ``(image_bytes, used_date)`` where ``image_bytes`` is
        the PNG data for the heatmap and ``used_date`` is the effective
        date in the adjusted price series.
    """
    default_mapping = {
        "XBTUSD Curncy": "Bitcoin",
        "XETUSD Curncy": "Ethereum",
        "XRPUSD Curncy": "Ripple",
        "XSOUSD Curncy": "Solana",
        "XBIUSD Curncy": "Binance",
        "XDOUSD Curncy": "Polkadot",
        "XVV Curncy": "AAVE",
        "CHYPEE Index":"HyperLiquid",
        "CCTON Curncy": "Ton",
        "BGCI Index": "Bloomberg Galaxy Crypto",
    }
    mapping = ticker_mapping or default_mapping
    tickers = list(mapping.keys())
    # Load data and adjust according to price mode
    df = _load_price_data(excel_path, tickers)
    df_adj, used_date = adjust_prices_for_mode(df, price_mode)
    perf = _build_returns_table(df_adj, mapping)
    # Keep all cryptos, even if they have missing data for some time periods
    # Only drop rows where the Name itself is missing
    heat_df = perf[["Name", "YTD", "1M", "3M", "6M", "12M"]].dropna(subset=["Name"]).sort_values("YTD", ascending=False).reset_index(drop=True)
    n_rows = len(heat_df)
    cols = ["YTD", "1M", "3M", "6M", "12M"]
    color_array = np.zeros((n_rows, len(cols), 4))
    for j, col in enumerate(cols):
        col_vals = heat_df[col].astype(float).values
        # Handle NaN values when computing color scale
        valid_vals = col_vals[~np.isnan(col_vals)]
        if len(valid_vals) == 0:
            continue
        pos_vals = valid_vals[valid_vals > 0]
        neg_vals = valid_vals[valid_vals < 0]
        max_pos = pos_vals.max() if len(pos_vals) > 0 else 0.0
        min_neg = neg_vals.min() if len(neg_vals) > 0 else 0.0
        for i, val in enumerate(col_vals):
            if not np.isnan(val):
                color_array[i, j] = _colour_for_value(float(val), max_pos, min_neg)
            else:
                # Use light gray for missing data
                color_array[i, j] = [0.9, 0.9, 0.9, 1.0]
    # Create figure and axis
    fig, ax = plt.subplots(figsize=(width, height))
    ax.imshow(color_array, aspect="auto")
    # Add text
    for i in range(n_rows):
        for j, col in enumerate(cols):
            val = heat_df.iloc[i][col]
            # Display "N/A" for missing data instead of trying to format NaN
            if pd.isna(val):
                text = "N/A"
            else:
                text = f"{val:.1f}%"
            ax.text(
                j,
                i,
                text,
                ha="center",
                va="center",
                fontsize=9,
                fontweight="bold",
                color="black",
            )
    ax.set_xticks(range(len(cols)))
    ax.set_xticklabels(cols, fontsize=12, fontweight="bold")
    ax.xaxis.set_ticks_position("top")
    ax.xaxis.set_label_position("top")
    ax.set_yticks(range(n_rows))
    ax.set_yticklabels(heat_df["Name"], fontsize=9)
    ax.tick_params(axis="y", length=0)
    ax.tick_params(axis="x", length=0)
    for spine in ax.spines.values():
        spine.set_visible(False)
    fig.subplots_adjust(left=0.5)
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=300, bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf.getvalue(), used_date


###############################################################################
# Colour helper
###############################################################################

def _colour_for_value(val: float, max_pos: float, min_neg: float) -> Tuple[float, float, float, float]:
    """Return an RGBA colour for a single value using independent scales.

    Positive values produce a gradient from white to green; negative
    values from white to red.  Zeros (or undefined scales) return
    white.  The input ``max_pos`` should be the maximum positive value
    in the column, and ``min_neg`` the minimum (most negative) value.
    """
    white = np.array([1.0, 1.0, 1.0, 1.0])
    red = np.array(mcolors.to_rgba("#C70039"))
    green = np.array(mcolors.to_rgba("#2E8B57"))
    if val > 0 and max_pos > 0:
        ratio = float(val) / max_pos
        return tuple(white + ratio * (green - white))
    if val < 0 and min_neg < 0:
        ratio = float(val) / min_neg  # both negative → positive ratio
        return tuple(white + ratio * (red - white))
    return tuple(white)


###############################################################################
# Slide insertion helpers
###############################################################################

def _insert_dashboard_to_placeholder(
    prs: Presentation,
    image_bytes: bytes,
    placeholder_names: List[str],
    *,
    left_cm: float,
    top_cm: float,
    width_cm: float,
    height_cm: float,
    used_date: Optional[pd.Timestamp],
    price_mode: str,
    source_placeholder_names: List[str],
) -> Presentation:
    """Helper to insert an image and source footnote into a slide identified by placeholders.

    This internal helper searches for shapes whose names match any of
    ``placeholder_names`` or whose text contains any of those names in
    square brackets.  Once found, it uses the slide but ignores the
    shape's dimensions, inserting the provided image at the specified
    coordinates and size.  It also searches for a text box named
    ``source_placeholder_names`` (or containing the pattern) and
    replaces its contents with a source footnote based on ``used_date``
    and ``price_mode`` while preserving the original formatting.  If no
    matching placeholder is found, the image is inserted on the last
    slide.  Footnote insertion is skipped if ``used_date`` is ``None``.

    Parameters
    ----------
    prs : Presentation
        The presentation into which the image should be inserted.
    image_bytes : bytes
        The PNG data to insert.
    placeholder_names : list of str
        Names of shapes to look for (lowercased).  The function also
        recognises text placeholders like ``[crypto_perf_1week]``.
    left_cm, top_cm, width_cm, height_cm : float
        Desired position and size in centimetres.
    used_date : pandas.Timestamp or None
        The effective date for the source footnote.  If ``None`` the
        source is not inserted.
    price_mode : str
        Either ``"Last Price"`` or ``"Last Close"``, used to suffix
        ``" Close"`` in the footnote.
    source_placeholder_names : list of str
        Names or patterns (without brackets) for the source text box.

    Returns
    -------
    Presentation
        The modified presentation.
    """
    target_slide = None
    placeholder_box = None
    # Normalise names for comparison
    name_candidates = [n.lower() for n in placeholder_names]
    pattern_candidates = [f"[{n}]" for n in name_candidates]
    for slide in prs.slides:
        for shape in slide.shapes:
            name_attr = getattr(shape, "name", "").lower()
            if name_attr in name_candidates:
                target_slide = slide
                placeholder_box = shape
                break
            if shape.has_text_frame:
                text_lower = (shape.text or "").strip().lower()
                if text_lower in [p.lower() for p in pattern_candidates]:
                    target_slide = slide
                    placeholder_box = shape
                    break
        if target_slide:
            break
    if target_slide is None:
        target_slide = prs.slides[min(11, len(prs.slides) - 1)]
    # Do not modify or remove the placeholder box; it serves only to locate
    # the slide.  Leave its text intact so the placeholder remains in
    # the template for future reference.
    left = Cm(left_cm)
    top = Cm(top_cm)
    width = Cm(width_cm)
    height = Cm(height_cm)
    stream = io.BytesIO(image_bytes)
    target_slide.shapes.add_picture(stream, left, top, width=width, height=height)
    # Insert source footnote if a date is available
    if used_date is not None:
        date_str = used_date.strftime("%d/%m/%Y")
        suffix = " Close" if price_mode.lower() == "last close" else ""
        source_text = f"Source: Bloomberg, Herculis Group, Data as of {date_str}{suffix}"
        # Look for source placeholder
        source_candidates = [n.lower() for n in source_placeholder_names]
        source_patterns = [f"[{n}]" for n in source_candidates]
        for shape in target_slide.shapes:
            name_attr2 = getattr(shape, "name", "").lower()
            if name_attr2 in source_candidates:
                if shape.has_text_frame:
                    runs = shape.text_frame.paragraphs[0].runs
                    attrs = _capture_font_attrs(runs[0]) if runs else (None, None, None, None, None, None)
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    new_run = p.add_run()
                    new_run.text = source_text
                    _apply_font_attrs(new_run, *attrs)
                break
            if shape.has_text_frame:
                for pattern in source_patterns:
                    if pattern.lower() in (shape.text or "").lower():
                        runs = shape.text_frame.paragraphs[0].runs
                        attrs = _capture_font_attrs(runs[0]) if runs else (None, None, None, None, None, None)
                        try:
                            new_text = shape.text.replace(pattern, source_text)
                        except Exception:
                            new_text = source_text
                        shape.text_frame.clear()
                        p = shape.text_frame.paragraphs[0]
                        new_run = p.add_run()
                        new_run.text = new_text
                        _apply_font_attrs(new_run, *attrs)
                        break
                else:
                    continue
                break
    return prs


def insert_crypto_performance_bar_slide(
    prs: Presentation,
    image_bytes: bytes,
    used_date: Optional[pd.Timestamp] = None,
    price_mode: str = "Last Price",
    *,
    left_cm: float = 0.0,
    top_cm: float = 3.22,
    width_cm: float = 25.0,
    height_cm: float = 11.66,
) -> Presentation:
    """Insert the weekly cryptocurrency performance bar chart and source footnote into its designated slide.

    The slide is identified by a shape named ``crypto_perf_1week`` or
    ``crypto_perf_1w`` (or containing ``[crypto_perf_1week]`` etc.).  The chart
    is inserted at the specified coordinates regardless of the
    placeholder's original size.  A source footnote is written into
    ``crypto_1w_source`` or a text box containing ``[crypto_1w_source]`` reflecting
    the ``used_date`` and ``price_mode``.

    Parameters
    ----------
    prs : Presentation
        The PowerPoint presentation to modify.
    image_bytes : bytes
        PNG data for the weekly bar chart.
    used_date : pandas.Timestamp or None, optional
        Effective date used for performance calculations.  If ``None``,
        the source footnote is not inserted.
    price_mode : str, default ``"Last Price"``
        Either ``"Last Price"`` or ``"Last Close"``.  Determines the
        suffix appended to the date in the source footnote.
    left_cm, top_cm, width_cm, height_cm : float
        Position and size for inserting the chart.

    Returns
    -------
    Presentation
        The modified presentation.
    """
    return _insert_dashboard_to_placeholder(
        prs,
        image_bytes,
        placeholder_names=["crypto_perf_1w", "crypto_perf_1week"],
        left_cm=left_cm,
        top_cm=top_cm,
        width_cm=width_cm,
        height_cm=height_cm,
        used_date=used_date,
        price_mode=price_mode,
        source_placeholder_names=["crypto_1w_source"],
    )


def insert_crypto_performance_histo_slide(
    prs: Presentation,
    image_bytes: bytes,
    used_date: Optional[pd.Timestamp] = None,
    price_mode: str = "Last Price",
    *,
    left_cm: float = 0.0,
    top_cm: float = 3.22,
    width_cm: float = 25.0,
    height_cm: float = 11.66,
) -> Presentation:
    """Insert the cryptocurrency historical performance heatmap and source footnote into its designated slide.

    The slide is identified by a shape named ``crypto_perf_histo`` (or
    containing ``[crypto_perf_histo]``).  The heatmap is inserted at the
    specified coordinates.  A source footnote is written into
    ``crypto_1w_source2`` or a text box containing ``[crypto_1w_source2]`` based
    on the provided ``used_date`` and ``price_mode``.

    Parameters
    ----------
    prs : Presentation
        The PowerPoint presentation to modify.
    image_bytes : bytes
        PNG data for the heatmap.
    used_date : pandas.Timestamp or None, optional
        Effective date used for performance calculations.  If ``None``,
        the source footnote is not inserted.
    price_mode : str, default ``"Last Price"``
        Either ``"Last Price"`` or ``"Last Close"``.  Determines the
        suffix appended to the date in the source footnote.
    left_cm, top_cm, width_cm, height_cm : float
        Position and size for inserting the heatmap.

    Returns
    -------
    Presentation
        The modified presentation.
    """
    return _insert_dashboard_to_placeholder(
        prs,
        image_bytes,
        placeholder_names=["crypto_perf_histo"],
        left_cm=left_cm,
        top_cm=top_cm,
        width_cm=width_cm,
        height_cm=height_cm,
        used_date=used_date,
        price_mode=price_mode,
        source_placeholder_names=["crypto_1w_source2"],
    )


# =============================================================================
# HTML-BASED CRYPTO WEEKLY PERFORMANCE CHART
# =============================================================================

# Crypto configuration with display names
CRYPTO_CONFIG = [
    {"ticker": "CCTON Curncy", "name": "Ton"},
    {"ticker": "CHYPEE Index", "name": "HyperLiquid"},
    {"ticker": "XDOUSD Curncy", "name": "Polkadot"},
    {"ticker": "XSOUSD Curncy", "name": "Solana"},
    {"ticker": "XBIUSD Curncy", "name": "Binance"},
    {"ticker": "XBTUSD Curncy", "name": "Bitcoin"},
    {"ticker": "XETUSD Curncy", "name": "Ethereum"},
    {"ticker": "BGCI Index", "name": "Bloomberg Galaxy Crypto"},
    {"ticker": "XVV Curncy", "name": "AAVE"},
    {"ticker": "XRPUSD Curncy", "name": "Ripple"},
]

# Chart dimensions (hardcoded)
CRYPTO_PNG_WIDTH_PX = 1963
CRYPTO_PNG_HEIGHT_PX = 1134
CRYPTO_PPT_WIDTH_CM = 17.31
CRYPTO_PPT_HEIGHT_CM = 10.0
CRYPTO_PPT_LEFT_CM = 3.35
CRYPTO_PPT_TOP_CM = 4.6
CRYPTO_HTML_SCALE = 3

# Historical chart dimensions (hardcoded)
CRYPTO_HIST_PNG_WIDTH_PX = 1930
CRYPTO_HIST_PNG_HEIGHT_PX = 1100
CRYPTO_HIST_PPT_WIDTH_CM = 17.02
CRYPTO_HIST_PPT_HEIGHT_CM = 9.7
CRYPTO_HIST_PPT_LEFT_CM = 3.35
CRYPTO_HIST_PPT_TOP_CM = 4.6
CRYPTO_HIST_HTML_SCALE = 3


def _format_crypto_percentage(value: float) -> str:
    """Format percentage with sign."""
    if value >= 0:
        return f"+{value:.1f}%"
    else:
        return f"{value:.1f}%"


def create_weekly_html_performance_chart(
    excel_path: Union[str, pathlib.Path],
    *,
    price_mode: str = "Last Price",
) -> Tuple[bytes, Optional[pd.Timestamp]]:
    """Generate HTML-based weekly crypto performance bar chart.

    Returns:
        Tuple of (PNG bytes, effective date used for computation)
    """
    # Get all tickers from configuration
    tickers = [c["ticker"] for c in CRYPTO_CONFIG]

    # Load and adjust price data
    df = _load_price_data(excel_path, tickers)
    df_adj, used_date = adjust_prices_for_mode(df, price_mode)
    today = df_adj["Date"].max()

    # Compute 1W returns for all cryptos
    rows_data = []
    for config in CRYPTO_CONFIG:
        ticker = config["ticker"]
        # Skip if ticker not in dataframe
        if ticker not in df_adj.columns:
            continue
        try:
            ret_val = _compute_horizon_returns(df_adj, ticker, today, 7)
            if pd.isna(ret_val):
                ret_val = 0.0
        except Exception:
            ret_val = 0.0

        rows_data.append({
            "name": config["name"],
            "value": ret_val,
        })

    # Sort by value descending (best performers first)
    rows_data.sort(key=lambda x: x["value"], reverse=True)

    # Calculate max absolute value for bar scaling (minimum 5%)
    valid_values = [abs(r["value"]) for r in rows_data if r["value"] != 0]
    max_abs_value = max(valid_values) if valid_values else 5.0
    max_abs_value = max(max_abs_value, 5.0)  # At least 5%

    # Prepare rows for template with highlight classes
    prepared_rows = []
    for i, row in enumerate(rows_data):
        value = row["value"]
        bar_width = abs(value) / max_abs_value * 48
        bar_width = min(bar_width, 48)

        # Determine highlight class
        if i == 0:
            highlight_class = "top-performer"
        elif i == len(rows_data) - 1:
            highlight_class = "worst-performer"
        else:
            highlight_class = ""

        prepared_rows.append({
            "name": row["name"],
            "value": value,
            "bar_class": "positive" if value >= 0 else "negative",
            "bar_width": bar_width,
            "value_class": "positive" if value >= 0 else "negative",
            "formatted_value": _format_crypto_percentage(value),
            "highlight_class": highlight_class,
        })

    # Calculate scale values
    scale_values = {
        "scale_min": f"-{max_abs_value:.0f}%",
        "scale_mid_low": f"-{max_abs_value/2:.0f}%",
        "scale_mid_high": f"+{max_abs_value/2:.0f}%",
        "scale_max": f"+{max_abs_value:.0f}%",
    }

    # Render HTML template
    template = Template(CRYPTO_WEEKLY_HTML_TEMPLATE)
    html_content = template.render(
        scale=CRYPTO_HTML_SCALE,
        width=CRYPTO_PNG_WIDTH_PX,
        height=CRYPTO_PNG_HEIGHT_PX,
        rows=prepared_rows,
        **scale_values,
    )

    # Convert HTML to PNG
    hti = Html2Image(size=(CRYPTO_PNG_WIDTH_PX, CRYPTO_PNG_HEIGHT_PX))
    hti.screenshot(html_str=html_content, save_as="crypto_weekly_temp.png")

    # Read the generated PNG
    with open("crypto_weekly_temp.png", "rb") as f:
        png_bytes = f.read()

    # Clean up temp file
    import os
    try:
        os.remove("crypto_weekly_temp.png")
    except Exception:
        pass

    return png_bytes, used_date


def insert_crypto_weekly_html_slide(
    prs: Presentation,
    image_bytes: bytes,
    used_date: Optional[pd.Timestamp] = None,
    price_mode: str = "Last Price",
) -> Presentation:
    """Insert the HTML-based crypto weekly performance chart into PowerPoint."""
    if not image_bytes:
        return prs

    # Find target slide by placeholder name
    target_slide = None
    placeholder_names = ["crypto_perf_1w", "crypto_perf_1week"]
    name_candidates = [n.lower() for n in placeholder_names]
    pattern_candidates = [f"[{n}]" for n in name_candidates]

    for slide in prs.slides:
        for shape in slide.shapes:
            name_attr = getattr(shape, "name", "").lower()
            if name_attr in name_candidates:
                target_slide = slide
                break
            if shape.has_text_frame:
                text_lower = (shape.text or "").strip().lower()
                if text_lower in [p.lower() for p in pattern_candidates]:
                    target_slide = slide
                    break
        if target_slide:
            break

    if target_slide is None:
        print("[Crypto Weekly HTML] ERROR: Slide not found")
        return prs

    # Insert chart image with exact hardcoded dimensions
    left = Cm(CRYPTO_PPT_LEFT_CM)
    top = Cm(CRYPTO_PPT_TOP_CM)
    width = Cm(CRYPTO_PPT_WIDTH_CM)
    height = Cm(CRYPTO_PPT_HEIGHT_CM)

    stream = io.BytesIO(image_bytes)
    picture = target_slide.shapes.add_picture(stream, left, top, width=width, height=height)

    # Send picture to back
    spTree = target_slide.shapes._spTree
    pic_element = picture._element
    spTree.remove(pic_element)
    spTree.insert(2, pic_element)

    print(f"[Crypto Weekly HTML] Chart inserted at ({CRYPTO_PPT_LEFT_CM}, {CRYPTO_PPT_TOP_CM}) cm")

    # Update source placeholder if date available
    if used_date is not None:
        date_str = used_date.strftime("%d/%m/%Y")
        suffix = " Close" if price_mode.lower() == "last close" else ""
        source_text = f"Source: Bloomberg, Herculis Group, Data as of {date_str}{suffix}"
        source_candidates = ["crypto_1w_source"]
        source_patterns = [f"[{n}]" for n in source_candidates]

        for shape in target_slide.shapes:
            name_attr = getattr(shape, "name", "").lower()
            if name_attr in [n.lower() for n in source_candidates]:
                if shape.has_text_frame:
                    runs = shape.text_frame.paragraphs[0].runs
                    attrs = _capture_font_attrs(runs[0]) if runs else (None, None, None, None, None, None)
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    new_run = p.add_run()
                    new_run.text = source_text
                    _apply_font_attrs(new_run, *attrs)
                break
            if shape.has_text_frame:
                for pattern in source_patterns:
                    if pattern.lower() in (shape.text or "").lower():
                        runs = shape.text_frame.paragraphs[0].runs
                        attrs = _capture_font_attrs(runs[0]) if runs else (None, None, None, None, None, None)
                        shape.text_frame.clear()
                        p = shape.text_frame.paragraphs[0]
                        new_run = p.add_run()
                        new_run.text = source_text
                        _apply_font_attrs(new_run, *attrs)
                        break
                else:
                    continue
                break

    return prs


# =============================================================================
# HTML-BASED CRYPTO HISTORICAL PERFORMANCE CHART
# =============================================================================

def _get_crypto_historical_color_class(value: float) -> str:
    """Get color class for crypto historical heatmap with wider thresholds.

    Crypto-specific thresholds (wider for volatility):
    - Level 1: 0-10%
    - Level 2: 10-50%
    - Level 3: 50-100%
    - Level 4: 100-200%
    - Level 5: 200%+
    """
    if pd.isna(value):
        return "neutral"

    abs_val = abs(value)
    prefix = "positive" if value >= 0 else "negative"

    if abs_val >= 200:
        return f"{prefix}-5"
    elif abs_val >= 100:
        return f"{prefix}-4"
    elif abs_val >= 50:
        return f"{prefix}-3"
    elif abs_val >= 10:
        return f"{prefix}-2"
    elif abs_val > 0:
        return f"{prefix}-1"
    else:
        return "neutral"


def _format_crypto_large_percentage(value: float) -> str:
    """Format large crypto percentages without decimal for values >= 100%."""
    if pd.isna(value):
        return "N/A"

    sign = "+" if value >= 0 else ""
    abs_val = abs(value)

    # Large percentages without decimal
    if abs_val >= 100:
        return f"{sign}{value:.0f}%"
    else:
        return f"{sign}{value:.1f}%"


def create_historical_html_performance_chart(
    excel_path: Union[str, pathlib.Path],
    *,
    price_mode: str = "Last Price",
) -> Tuple[bytes, Optional[pd.Timestamp]]:
    """Generate HTML-based crypto historical performance heatmap.

    Features:
    - 10 cryptocurrencies sorted by YTD
    - Wider color thresholds for crypto volatility (10%, 50%, 100%, 200%)
    - Large percentages without decimal (e.g., +520%)
    - YTD column emphasized with larger font and shadow

    Returns:
        Tuple of (PNG bytes, effective date used for computation)
    """
    # Get all tickers from configuration
    tickers = [c["ticker"] for c in CRYPTO_CONFIG]
    ticker_to_name = {c["ticker"]: c["name"] for c in CRYPTO_CONFIG}

    # Load and adjust price data
    df = _load_price_data(excel_path, tickers)
    df_adj, used_date = adjust_prices_for_mode(df, price_mode)

    # Build returns table
    perf = _build_returns_table(df_adj, ticker_to_name)

    # Sort by YTD descending
    perf_sorted = perf.sort_values("YTD", ascending=False).reset_index(drop=True)

    # Prepare rows for template
    rows_data = []
    for _, row in perf_sorted.iterrows():
        ytd_val = row["YTD"] if not pd.isna(row["YTD"]) else 0.0
        m1_val = row["1M"] if not pd.isna(row["1M"]) else 0.0
        m3_val = row["3M"] if not pd.isna(row["3M"]) else 0.0
        m6_val = row["6M"] if not pd.isna(row["6M"]) else 0.0
        m12_val = row["12M"] if not pd.isna(row["12M"]) else 0.0

        rows_data.append({
            "name": row["Name"],
            "ytd_class": _get_crypto_historical_color_class(ytd_val),
            "ytd_formatted": _format_crypto_large_percentage(ytd_val),
            "m1_class": _get_crypto_historical_color_class(m1_val),
            "m1_formatted": _format_crypto_large_percentage(m1_val),
            "m3_class": _get_crypto_historical_color_class(m3_val),
            "m3_formatted": _format_crypto_large_percentage(m3_val),
            "m6_class": _get_crypto_historical_color_class(m6_val),
            "m6_formatted": _format_crypto_large_percentage(m6_val),
            "m12_class": _get_crypto_historical_color_class(m12_val),
            "m12_formatted": _format_crypto_large_percentage(m12_val),
        })

    # Render HTML template
    template = Template(CRYPTO_HISTORICAL_HTML_TEMPLATE)
    html_content = template.render(
        scale=CRYPTO_HIST_HTML_SCALE,
        width=CRYPTO_HIST_PNG_WIDTH_PX,
        height=CRYPTO_HIST_PNG_HEIGHT_PX,
        rows=rows_data,
    )

    # Convert HTML to PNG
    hti = Html2Image(size=(CRYPTO_HIST_PNG_WIDTH_PX, CRYPTO_HIST_PNG_HEIGHT_PX))
    hti.screenshot(html_str=html_content, save_as="crypto_hist_temp.png")

    # Read the generated PNG
    with open("crypto_hist_temp.png", "rb") as f:
        png_bytes = f.read()

    # Clean up temp file
    import os
    try:
        os.remove("crypto_hist_temp.png")
    except Exception:
        pass

    return png_bytes, used_date


def insert_crypto_historical_html_slide(
    prs: Presentation,
    image_bytes: bytes,
    used_date: Optional[pd.Timestamp] = None,
    price_mode: str = "Last Price",
) -> Presentation:
    """Insert the HTML-based crypto historical performance chart into PowerPoint."""
    if not image_bytes:
        return prs

    # Find target slide by placeholder name
    target_slide = None
    placeholder_names = ["crypto_perf_histo"]
    name_candidates = [n.lower() for n in placeholder_names]
    pattern_candidates = [f"[{n}]" for n in name_candidates]

    for slide in prs.slides:
        for shape in slide.shapes:
            name_attr = getattr(shape, "name", "").lower()
            if name_attr in name_candidates:
                target_slide = slide
                break
            if shape.has_text_frame:
                text_lower = (shape.text or "").strip().lower()
                if text_lower in [p.lower() for p in pattern_candidates]:
                    target_slide = slide
                    break
        if target_slide:
            break

    if target_slide is None:
        print("[Crypto Historical HTML] ERROR: Slide not found")
        return prs

    # Insert chart image with exact hardcoded dimensions
    left = Cm(CRYPTO_HIST_PPT_LEFT_CM)
    top = Cm(CRYPTO_HIST_PPT_TOP_CM)
    width = Cm(CRYPTO_HIST_PPT_WIDTH_CM)
    height = Cm(CRYPTO_HIST_PPT_HEIGHT_CM)

    stream = io.BytesIO(image_bytes)
    picture = target_slide.shapes.add_picture(stream, left, top, width=width, height=height)

    # Send picture to back
    spTree = target_slide.shapes._spTree
    pic_element = picture._element
    spTree.remove(pic_element)
    spTree.insert(2, pic_element)

    print(f"[Crypto Historical HTML] Chart inserted at ({CRYPTO_HIST_PPT_LEFT_CM}, {CRYPTO_HIST_PPT_TOP_CM}) cm")

    # Update source placeholder if date available
    if used_date is not None:
        date_str = used_date.strftime("%d/%m/%Y")
        suffix = " Close" if price_mode.lower() == "last close" else ""
        source_text = f"Source: Bloomberg, Herculis Group, Data as of {date_str}{suffix}"
        source_candidates = ["crypto_1w_source2"]
        source_patterns = [f"[{n}]" for n in source_candidates]

        for shape in target_slide.shapes:
            name_attr = getattr(shape, "name", "").lower()
            if name_attr in [n.lower() for n in source_candidates]:
                if shape.has_text_frame:
                    runs = shape.text_frame.paragraphs[0].runs
                    attrs = _capture_font_attrs(runs[0]) if runs else (None, None, None, None, None, None)
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    new_run = p.add_run()
                    new_run.text = source_text
                    _apply_font_attrs(new_run, *attrs)
                break
            if shape.has_text_frame:
                for pattern in source_patterns:
                    if pattern.lower() in (shape.text or "").lower():
                        runs = shape.text_frame.paragraphs[0].runs
                        attrs = _capture_font_attrs(runs[0]) if runs else (None, None, None, None, None, None)
                        shape.text_frame.clear()
                        p = shape.text_frame.paragraphs[0]
                        new_run = p.add_run()
                        new_run.text = source_text
                        _apply_font_attrs(new_run, *attrs)
                        break
                else:
                    continue
                break

    return prs


# =============================================================================
# CRYPTO YTD EVOLUTION LINE CHART (Chart.js + Playwright)
# =============================================================================

# 6 cryptos only with official brand colors
CRYPTO_YTD_CONFIG = {
    "XBTUSD Curncy": {"name": "Bitcoin", "color": "#F7931A"},
    "XETUSD Curncy": {"name": "Ethereum", "color": "#627EEA"},
    "XBIUSD Curncy": {"name": "Binance", "color": "#F3BA2F"},
    "XSOUSD Curncy": {"name": "Solana", "color": "#9945FF"},
    "XRPUSD Curncy": {"name": "Ripple", "color": "#0085C0"},
    "CCTON Curncy": {"name": "Ton", "color": "#00CED1"},
}

# Chart dimensions
CRYPTO_YTD_PNG_WIDTH_PX = 2700
CRYPTO_YTD_PNG_HEIGHT_PX = 1350
CRYPTO_YTD_PPT_WIDTH_CM = 23.0
CRYPTO_YTD_PPT_LEFT_CM = 0.5
CRYPTO_YTD_PPT_TOP_CM = 4.2
CRYPTO_YTD_HTML_SCALE = 3


def _get_crypto_ytd_year(data_end_date: pd.Timestamp) -> int:
    """Determine which year to show as YTD based on data end date."""
    return data_end_date.year


def _get_crypto_chart_title(data_end_date: pd.Timestamp) -> str:
    """Generate chart title with correct year."""
    year = _get_crypto_ytd_year(data_end_date)
    return f"{year} YTD Evolution"


def _compute_crypto_ytd_series(df, ticker):
    """Compute weekly YTD performance series for a crypto."""
    if ticker not in df.columns:
        return [], []

    last_date = df["Date"].max()
    year = last_date.year
    start_of_year = pd.Timestamp(year=year, month=1, day=1)
    df_year = df[df["Date"] >= start_of_year].copy().sort_values("Date")

    if len(df_year) == 0:
        return [], []

    first_valid_idx = df_year[ticker].first_valid_index()
    if first_valid_idx is None:
        return [], []

    baseline_price = df_year.loc[first_valid_idx, ticker]
    if pd.isna(baseline_price) or baseline_price == 0:
        return [], []

    month_names = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun',
                   'Jul', 'Aug', 'Sep', 'Oct', 'Nov', 'Dec']

    daily_labels = []
    daily_values = []

    for _, row in df_year.iterrows():
        price = row[ticker]
        date = row["Date"]
        if not pd.isna(price):
            ytd_return = (price - baseline_price) / baseline_price * 100
            daily_labels.append(month_names[date.month - 1])
            daily_values.append(round(ytd_return, 1))

    # If fewer than 22 data points, use daily frequency instead of weekly
    # This handles beginning of year when there's only a few days of data
    if len(daily_values) < 22:
        labels = daily_labels
        values = daily_values
    else:
        # Sample to weekly (every 5th trading day)
        weekly_indices = list(range(0, len(daily_values), 5))
        if len(daily_values) > 0 and weekly_indices[-1] != len(daily_values) - 1:
            weekly_indices.append(len(daily_values) - 1)

        labels = [daily_labels[i] for i in weekly_indices]
        values = [daily_values[i] for i in weekly_indices]

    return labels, values


def create_crypto_ytd_evolution_chart(excel_path, *, price_mode="Last Price"):
    """Generate HTML-based Crypto YTD Evolution line chart."""
    tickers = list(CRYPTO_YTD_CONFIG.keys())
    df = _load_price_data(excel_path, tickers)
    df_adj, used_date = adjust_prices_for_mode(df, price_mode)

    if used_date is None:
        used_date = df_adj["Date"].max()

    chart_title = _get_crypto_chart_title(used_date)

    all_labels = []
    datasets = []

    for ticker, config in CRYPTO_YTD_CONFIG.items():
        labels, values = _compute_crypto_ytd_series(df_adj, ticker)
        if labels and values:
            if len(labels) > len(all_labels):
                all_labels = labels
            datasets.append({
                "label": config["name"],
                "data": values,
                "borderColor": config["color"],
                "backgroundColor": "transparent",
            })

    max_len = len(all_labels)
    for dataset in datasets:
        while len(dataset["data"]) < max_len:
            dataset["data"].append(None)

    # Check for insufficient data (< 3 data points)
    if len(all_labels) < 3:
        print(f"[DEBUG] Crypto: Insufficient data ({len(all_labels)} points), rendering placeholder")
        env = Environment()
        template = env.from_string(YTD_INSUFFICIENT_DATA_HTML_TEMPLATE)
        html_content = template.render(
            width=CRYPTO_YTD_PNG_WIDTH_PX,
            height=CRYPTO_YTD_PNG_HEIGHT_PX,
            scale=CRYPTO_YTD_HTML_SCALE,
            chart_title=chart_title,
        )
        png_bytes = None
        try:
            with sync_playwright() as p:
                browser = p.chromium.launch(headless=True)
                page = browser.new_page(viewport={
                    'width': CRYPTO_YTD_PNG_WIDTH_PX,
                    'height': CRYPTO_YTD_PNG_HEIGHT_PX
                })
                page.set_content(html_content, wait_until='networkidle')
                page.wait_for_timeout(500)
                png_bytes = page.screenshot()
                browser.close()
        except Exception as e:
            print(f"[ERROR] Playwright failed for placeholder: {e}")
        return png_bytes, used_date

    all_values = []
    for ds in datasets:
        all_values.extend([v for v in ds["data"] if v is not None])

    if all_values:
        data_min = min(all_values)
        data_max = max(all_values)
        y_min = (int(data_min / 5) - 1) * 5
        y_max = (int(data_max / 5) + 2) * 5
        if data_min < 0 < data_max:
            y_min = min(y_min, -5)
            y_max = max(y_max, 5)
    else:
        y_min = -20
        y_max = 40

    env = Environment()
    env.filters['tojson'] = json.dumps
    template = env.from_string(CRYPTO_YTD_EVOLUTION_HTML_TEMPLATE)
    html_content = template.render(
        width=CRYPTO_YTD_PNG_WIDTH_PX,
        height=CRYPTO_YTD_PNG_HEIGHT_PX,
        scale=CRYPTO_YTD_HTML_SCALE,
        chart_title=chart_title,
        labels=all_labels,
        datasets=datasets,
        y_min=y_min,
        y_max=y_max,
    )

    png_bytes = None
    try:
        with sync_playwright() as p:
            browser = p.chromium.launch(headless=True)
            page = browser.new_page(viewport={
                'width': CRYPTO_YTD_PNG_WIDTH_PX,
                'height': CRYPTO_YTD_PNG_HEIGHT_PX
            })
            page.set_content(html_content, wait_until='networkidle')
            try:
                page.wait_for_selector('body[data-chart-ready="true"]', timeout=10000)
            except Exception:
                page.wait_for_timeout(3000)
            png_bytes = page.screenshot()
            browser.close()
    except Exception as e:
        print(f"[ERROR] Playwright failed for crypto chart: {e}")

    return png_bytes, used_date


def insert_crypto_ytd_evolution_slide(prs, image_bytes, used_date=None, price_mode="Last Price", subtitle=None):
    """Insert the Crypto YTD Evolution chart into PowerPoint.

    Parameters
    ----------
    prs : Presentation
        The PowerPoint presentation to modify.
    image_bytes : bytes
        PNG data for the YTD evolution chart.
    used_date : pandas.Timestamp or None, optional
        Effective date used for performance calculations.
    price_mode : str, default "Last Price"
        Either "Last Price" or "Last Close".
    subtitle : str or None, optional
        Subtitle text to insert into the ytd_crypto_subtitle placeholder.
    """
    if not image_bytes:
        return prs

    target_slide = None
    placeholder_names = ["ytd_crypto_perf"]

    for slide in prs.slides:
        for shape in slide.shapes:
            name_attr = getattr(shape, "name", "").lower()
            if name_attr in [n.lower() for n in placeholder_names]:
                target_slide = slide
                break
            if shape.has_text_frame:
                text_lower = (shape.text or "").strip().lower()
                if text_lower == "[ytd_crypto_perf]":
                    target_slide = slide
                    break
        if target_slide:
            break

    if target_slide is None:
        print(f"[Crypto YTD Evolution] ERROR: Slide not found")
        return prs

    # ------------------------------------------------------------------
    # Insert subtitle into the 'ytd_crypto_subtitle' placeholder
    # ------------------------------------------------------------------
    if subtitle:
        for shape in target_slide.shapes:
            name_attr = getattr(shape, "name", "")
            if name_attr and name_attr.lower() == "ytd_crypto_subtitle" and shape.has_text_frame:
                tf = shape.text_frame
                paragraph = tf.paragraphs[0]
                runs = paragraph.runs
                attrs = _capture_font_attrs(runs[0]) if runs else (None, None, None, None, None, None)
                # Determine replacement: replace tokens "XXX" or "[ytd_crypto_subtitle]"
                original_text = "".join(run.text for run in runs) if runs else ""
                new_text = original_text
                if "XXX" in new_text:
                    new_text = new_text.replace("XXX", subtitle)
                elif "[ytd_crypto_subtitle]" in new_text:
                    new_text = new_text.replace("[ytd_crypto_subtitle]", subtitle)
                else:
                    new_text = subtitle
                # Clear and insert new run
                tf.clear()
                p = tf.paragraphs[0]
                new_run = p.add_run()
                new_run.text = new_text
                _apply_font_attrs(new_run, *attrs)
                break

    # ------------------------------------------------------------------
    # Insert chart image with exact hardcoded dimensions
    # ------------------------------------------------------------------
    from io import BytesIO
    image_stream = BytesIO(image_bytes)

    left = Cm(CRYPTO_YTD_PPT_LEFT_CM)
    top = Cm(CRYPTO_YTD_PPT_TOP_CM)
    width = Cm(CRYPTO_YTD_PPT_WIDTH_CM)

    pic = target_slide.shapes.add_picture(image_stream, left, top, width=width)

    spTree = target_slide.shapes._spTree
    sp = pic._element
    spTree.remove(sp)
    spTree.insert(2, sp)

    print(f"[Crypto YTD Evolution] Chart inserted")

    # ------------------------------------------------------------------
    # Insert data source footnote
    # ------------------------------------------------------------------
    if used_date is not None:
        date_str = used_date.strftime("%d/%m/%Y")
        suffix = " Close" if price_mode.lower() == "last close" else ""
        source_text = f"Source: Bloomberg, Herculis Group, Data as of {date_str}{suffix}"
        placeholder_name = "ytd_crypto_source"
        placeholder_patterns = ["[ytd_crypto_source]", "ytd_crypto_source"]
        inserted = False
        for shape in target_slide.shapes:
            name_attr = getattr(shape, "name", "")
            if name_attr and name_attr.lower() == placeholder_name:
                if shape.has_text_frame:
                    runs = shape.text_frame.paragraphs[0].runs
                    attrs = _capture_font_attrs(runs[0]) if runs else (None, None, None, None, None, None)
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    new_run = p.add_run()
                    new_run.text = source_text
                    _apply_font_attrs(new_run, *attrs)
                inserted = True
                break
            if shape.has_text_frame:
                for pattern in placeholder_patterns:
                    if pattern.lower() in (shape.text or "").lower():
                        runs = shape.text_frame.paragraphs[0].runs
                        attrs = _capture_font_attrs(runs[0]) if runs else (None, None, None, None, None, None)
                        try:
                            new_text = shape.text.replace(pattern, source_text)
                        except Exception:
                            new_text = source_text
                        shape.text_frame.clear()
                        p = shape.text_frame.paragraphs[0]
                        new_run = p.add_run()
                        new_run.text = new_text
                        _apply_font_attrs(new_run, *attrs)
                        inserted = True
                        break
                if inserted:
                    break

    return prs
